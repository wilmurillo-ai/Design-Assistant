#!/usr/bin/env python3
"""
Deck to HTML Converter by UniqueClub
Converts PPT/PDF pitch decks into responsive, self-contained HTML presentations.

Usage:
    python convert_deck.py --input pitch_deck.pptx --output presentation.html
    python convert_deck.py --input document.pdf

Features:
    - Extracts text, images, and structure from .pptx or .pdf files
    - Generates single self-contained HTML file (no external dependencies)
    - Keyboard navigation (arrow keys), touch/swipe support
    - Progress bar, slide counter, fullscreen toggle
    - Mobile responsive design
"""

import argparse
import base64
import json
import os
import sys
from pathlib import Path


def extract_from_pptx(filepath):
    """Extract slide content from a .pptx file."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        print("Error: python-pptx not installed. Run: pip install python-pptx")
        sys.exit(1)
    
    prs = Presentation(filepath)
    slides = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_data = {
            "number": slide_num,
            "title": "",
            "content": [],
            "images": [],
            "layout": "content"
        }
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                # Detect title (largest font or top position)
                is_title = False
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.size and run.font.size > Pt(24):
                                is_title = True
                                break
                
                if is_title and not slide_data["title"]:
                    slide_data["title"] = text
                else:
                    slide_data["content"].append(text)
            
            # Extract images
            if shape.shape_type == 13:  # Picture type
                try:
                    image = shape.image
                    image_bytes = image.blob
                    image_b64 = base64.b64encode(image_bytes).decode('utf-8')
                    slide_data["images"].append({
                        "data": f"data:{image.content_type};base64,{image_b64}",
                        "ext": image.ext
                    })
                except Exception:
                    pass
        
        # Detect slide type based on content
        title_lower = slide_data["title"].lower()
        if slide_num == 1:
            slide_data["layout"] = "cover"
        elif any(word in title_lower for word in ["fund", "invest", "raise", "融资"]):
            slide_data["layout"] = "fundraising"
        elif any(word in title_lower for word in ["team", "团队"]):
            slide_data["layout"] = "team"
        elif any(word in title_lower for word in ["competitor", "竞争", "竞品"]):
            slide_data["layout"] = "comparison"
        elif any(word in title_lower for word in ["roadmap", "路线图"]):
            slide_data["layout"] = "roadmap"
        
        slides.append(slide_data)
    
    return slides


def extract_from_pdf(filepath):
    """Extract page content from a .pdf file."""
    try:
        import fitz  # pymupdf
    except ImportError:
        print("Error: pymupdf not installed. Run: pip install pymupdf")
        sys.exit(1)
    
    doc = fitz.open(filepath)
    slides = []
    
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        slide_data = {
            "number": page_num + 1,
            "title": f"Page {page_num + 1}",
            "content": [],
            "images": [],
            "layout": "content"
        }
        
        # Extract text
        text = page.get_text()
        lines = [line.strip() for line in text.split('\n') if line.strip()]
        
        if lines:
            # First substantial line as title
            slide_data["title"] = lines[0][:100]
            slide_data["content"] = lines[1:]
        
        # Extract images
        image_list = page.get_images()
        for img_index, img in enumerate(image_list):
            try:
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_b64 = base64.b64encode(image_bytes).decode('utf-8')
                ext = base_image["ext"]
                mime_type = f"image/{ext}" if ext else "image/png"
                slide_data["images"].append({
                    "data": f"data:{mime_type};base64,{image_b64}",
                    "ext": ext
                })
            except Exception:
                pass
        
        # Detect layout
        title_lower = slide_data["title"].lower()
        if page_num == 0:
            slide_data["layout"] = "cover"
        elif any(word in title_lower for word in ["fund", "invest", "team", "roadmap", "competitor"]):
            if "fund" in title_lower or "invest" in title_lower:
                slide_data["layout"] = "fundraising"
            elif "team" in title_lower:
                slide_data["layout"] = "team"
            elif "roadmap" in title_lower:
                slide_data["layout"] = "roadmap"
            elif "competitor" in title_lower:
                slide_data["layout"] = "comparison"
        
        slides.append(slide_data)
    
    doc.close()
    return slides


def generate_html(slides, title, accent_color="#1a73e8"):
    """Generate complete self-contained HTML presentation."""
    
    slides_html = []
    for slide in slides:
        layout_class = f"slide-{slide['layout']}"
        
        # Build content HTML
        content_html = ""
        
        # Title
        if slide["title"]:
            content_html += f'<h1 class="slide-title">{escape_html(slide["title"])}</h1>'
        
        # Images
        for img in slide["images"][:2]:  # Limit to 2 images per slide
            content_html += f'<img src="{img["data"]}" class="slide-image" alt="Slide image">'
        
        # Content paragraphs
        for content in slide["content"][:8]:  # Limit content items
            if content != slide["title"] and len(content) > 5:
                content_html += f'<p class="slide-text">{escape_html(content)}</p>'
        
        slides_html.append(f'''
        <div class="slide {layout_class}" data-index="{slide["number"] - 1}">
            <div class="slide-content">
                {content_html}
            </div>
        </div>
        ''')
    
    slides_container = '\n'.join(slides_html)
    dots_html = '\n'.join([f'<span class="dot" data-index="{i}"></span>' for i in range(len(slides))])
    
    html = f'''<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{escape_html(title)} — Business Plan</title>
    <style>
        * {{
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }}
        
        body {{
            font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", "Microsoft YaHei", sans-serif;
            background: linear-gradient(135deg, #0f0f1a, #1a1a2e);
            min-height: 100vh;
            overflow: hidden;
        }}
        
        /* Progress bar */
        .progress-bar {{
            position: fixed;
            top: 0;
            left: 0;
            width: 100%;
            height: 4px;
            background: rgba(255,255,255,0.1);
            z-index: 1000;
        }}
        .progress-fill {{
            height: 100%;
            background: {accent_color};
            width: 0%;
            transition: width 0.3s ease;
        }}
        
        /* Slides container */
        .slides-container {{
            position: relative;
            width: 100%;
            height: 100vh;
            display: flex;
            align-items: center;
            justify-content: center;
            padding: 40px 20px;
        }}
        
        /* Individual slide */
        .slide {{
            position: absolute;
            width: 100%;
            max-width: 1200px;
            aspect-ratio: 16/9;
            background: #ffffff;
            border-radius: 12px;
            box-shadow: 0 20px 60px rgba(0,0,0,0.3);
            padding: 60px;
            opacity: 0;
            transform: translateX(100px);
            transition: all 0.5s ease;
            pointer-events: none;
            overflow-y: auto;
        }}
        
        .slide.active {{
            opacity: 1;
            transform: translateX(0);
            pointer-events: auto;
        }}
        
        .slide.prev {{
            transform: translateX(-100px);
        }}
        
        /* Slide layouts */
        .slide-cover {{
            background: linear-gradient(135deg, {accent_color}, #1a1a2e);
            color: white;
            text-align: center;
            display: flex;
            flex-direction: column;
            justify-content: center;
        }}
        
        .slide-cover .slide-title {{
            font-size: 3.5rem;
            margin-bottom: 1rem;
            color: white;
        }}
        
        .slide-cover .slide-text {{
            color: rgba(255,255,255,0.9);
            font-size: 1.3rem;
        }}
        
        /* Content styling */
        .slide-content {{
            height: 100%;
        }}
        
        .slide-title {{
            font-size: 2.5rem;
            color: #202124;
            margin-bottom: 2rem;
            padding-bottom: 0.5rem;
            border-bottom: 3px solid {accent_color};
        }}
        
        .slide-text {{
            font-size: 1.2rem;
            line-height: 1.8;
            color: #5f6368;
            margin-bottom: 1rem;
        }}
        
        .slide-image {{
            max-width: 100%;
            max-height: 400px;
            border-radius: 8px;
            margin: 1rem 0;
            box-shadow: 0 4px 12px rgba(0,0,0,0.1);
        }}
        
        /* Navigation */
        .nav-dots {{
            position: fixed;
            bottom: 30px;
            left: 50%;
            transform: translateX(-50%);
            display: flex;
            gap: 10px;
            z-index: 100;
        }}
        
        .dot {{
            width: 10px;
            height: 10px;
            border-radius: 50%;
            background: rgba(255,255,255,0.3);
            cursor: pointer;
            transition: all 0.3s ease;
        }}
        
        .dot.active {{
            background: {accent_color};
            transform: scale(1.2);
        }}
        
        .slide-counter {{
            position: fixed;
            bottom: 30px;
            right: 30px;
            color: rgba(255,255,255,0.7);
            font-size: 14px;
            z-index: 100;
        }}
        
        .nav-arrow {{
            position: fixed;
            top: 50%;
            transform: translateY(-50%);
            background: rgba(255,255,255,0.1);
            border: none;
            color: white;
            font-size: 2rem;
            width: 50px;
            height: 50px;
            border-radius: 50%;
            cursor: pointer;
            z-index: 100;
            transition: all 0.3s ease;
            display: flex;
            align-items: center;
            justify-content: center;
        }}
        
        .nav-arrow:hover {{
            background: rgba(255,255,255,0.2);
        }}
        
        .nav-arrow.prev {{
            left: 20px;
        }}
        
        .nav-arrow.next {{
            right: 20px;
        }}
        
        .fullscreen-btn {{
            position: fixed;
            top: 20px;
            right: 20px;
            background: rgba(255,255,255,0.1);
            border: none;
            color: white;
            font-size: 1.5rem;
            padding: 10px;
            border-radius: 8px;
            cursor: pointer;
            z-index: 100;
            transition: all 0.3s ease;
        }}
        
        .fullscreen-btn:hover {{
            background: rgba(255,255,255,0.2);
        }}
        
        /* Responsive */
        @media (max-width: 768px) {{
            .slide {{
                padding: 30px;
                aspect-ratio: auto;
                min-height: 80vh;
            }}
            
            .slide-title {{
                font-size: 1.8rem;
            }}
            
            .slide-cover .slide-title {{
                font-size: 2rem;
            }}
            
            .slide-text {{
                font-size: 1rem;
            }}
            
            .nav-arrow {{
                display: none;
            }}
        }}
        
        /* Print styles */
        @media print {{
            body {{
                background: white;
            }}
            .slide {{
                position: relative;
                opacity: 1;
                transform: none;
                page-break-after: always;
                box-shadow: none;
                border: 1px solid #ddd;
            }}
            .nav-dots, .nav-arrow, .fullscreen-btn, .slide-counter, .progress-bar {{
                display: none;
            }}
        }}
    </style>
</head>
<body>
    <!-- Progress bar -->
    <div class="progress-bar">
        <div class="progress-fill" id="progressFill"></div>
    </div>
    
    <!-- Fullscreen button -->
    <button class="fullscreen-btn" onclick="toggleFullscreen()" title="Fullscreen (F)">⛶</button>
    
    <!-- Slides -->
    <div class="slides-container">
        {slides_container}
    </div>
    
    <!-- Navigation -->
    <div class="nav-dots">
        {dots_html}
    </div>
    
    <div class="slide-counter">
        <span id="currentSlide">1</span> / {len(slides)}
    </div>
    
    <button class="nav-arrow prev" onclick="prevSlide()" title="Previous (←)">‹</button>
    <button class="nav-arrow next" onclick="nextSlide()" title="Next (→)">›</button>
    
    <script>
        let currentSlide = 0;
        const slides = document.querySelectorAll('.slide');
        const dots = document.querySelectorAll('.dot');
        const totalSlides = slides.length;
        
        function updateSlide() {{
            slides.forEach((slide, index) => {{
                slide.classList.remove('active', 'prev');
                if (index === currentSlide) {{
                    slide.classList.add('active');
                }} else if (index < currentSlide) {{
                    slide.classList.add('prev');
                }}
            }});
            
            dots.forEach((dot, index) => {{
                dot.classList.toggle('active', index === currentSlide);
            }});
            
            document.getElementById('currentSlide').textContent = currentSlide + 1;
            document.getElementById('progressFill').style.width = 
                ((currentSlide + 1) / totalSlides * 100) + '%';
        }}
        
        function nextSlide() {{
            if (currentSlide < totalSlides - 1) {{
                currentSlide++;
                updateSlide();
            }}
        }}
        
        function prevSlide() {{
            if (currentSlide > 0) {{
                currentSlide--;
                updateSlide();
            }}
        }}
        
        function goToSlide(index) {{
            currentSlide = index;
            updateSlide();
        }}
        
        function toggleFullscreen() {{
            if (!document.fullscreenElement) {{
                document.documentElement.requestFullscreen();
            }} else {{
                document.exitFullscreen();
            }}
        }}
        
        // Dot navigation
        dots.forEach((dot, index) => {{
            dot.addEventListener('click', () => goToSlide(index));
        }});
        
        // Keyboard navigation
        document.addEventListener('keydown', (e) => {{
            if (e.key === 'ArrowRight' || e.key === ' ') {{
                nextSlide();
            }} else if (e.key === 'ArrowLeft') {{
                prevSlide();
            }} else if (e.key === 'f' || e.key === 'F') {{
                toggleFullscreen();
            }} else if (e.key === 'Escape') {{
                if (document.fullscreenElement) {{
                    document.exitFullscreen();
                }}
            }}
        }});
        
        // Touch/Swipe support
        let touchStartX = 0;
        let touchEndX = 0;
        
        document.addEventListener('touchstart', (e) => {{
            touchStartX = e.changedTouches[0].screenX;
        }});
        
        document.addEventListener('touchend', (e) => {{
            touchEndX = e.changedTouches[0].screenX;
            handleSwipe();
        }});
        
        function handleSwipe() {{
            if (touchEndX < touchStartX - 50) {{
                nextSlide();
            }}
            if (touchEndX > touchStartX + 50) {{
                prevSlide();
            }}
        }}
        
        // Initialize
        updateSlide();
    </script>
</body>
</html>
'''
    return html


def escape_html(text):
    """Escape HTML special characters."""
    return (text
            .replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
            .replace('"', "&quot;")
            .replace("'", "&#x27;"))


def main():
    parser = argparse.ArgumentParser(description="Convert PPT/PDF to HTML presentation")
    parser.add_argument("--input", "-i", required=True, help="Input PPTX or PDF file")
    parser.add_argument("--output", "-o", help="Output HTML file (optional)")
    parser.add_argument("--color", "-c", default="#1a73e8", help="Accent color (hex)")
    args = parser.parse_args()
    
    if not os.path.exists(args.input):
        print(f"Error: Input file not found: {args.input}")
        sys.exit(1)
    
    # Detect file type
    ext = os.path.splitext(args.input)[1].lower()
    
    if ext == ".pptx":
        print("Extracting content from PowerPoint...")
        slides = extract_from_pptx(args.input)
    elif ext == ".pdf":
        print("Extracting content from PDF...")
        slides = extract_from_pdf(args.input)
    else:
        print(f"Error: Unsupported format: {ext}. Use .pptx or .pdf")
        sys.exit(1)
    
    print(f"Extracted {len(slides)} slides")
    
    # Generate output filename
    if args.output:
        output_file = args.output
    else:
        base_name = os.path.splitext(args.input)[0]
        output_file = f"{base_name}_presentation.html"
    
    # Generate HTML
    title = os.path.splitext(os.path.basename(args.input))[0]
    html = generate_html(slides, title, args.color)
    
    # Save
    with open(output_file, 'w', encoding='utf-8') as f:
        f.write(html)
    
    file_size = os.path.getsize(output_file) / 1024
    print(f"✅ HTML presentation generated: {os.path.abspath(output_file)}")
    print(f"   Slides: {len(slides)} | Size: {file_size:.1f} KB")
    print(f"\nControls:")
    print(f"   ← → Arrow keys: Navigate slides")
    print(f"   F: Toggle fullscreen")
    print(f"   Swipe: Mobile navigation")


if __name__ == "__main__":
    main()
