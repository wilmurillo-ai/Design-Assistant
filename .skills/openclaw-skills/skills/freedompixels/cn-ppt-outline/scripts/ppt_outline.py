#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT大纲生成助手
主题 → 大纲 → PPT文件
"""

import os
import sys
from typing import List, Dict
from dataclasses import dataclass

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


@dataclass
class Slide:
    """单页幻灯片"""
    title: str
    bullets: List[str]
    layout: str = "content"  # title/content/section/end


@dataclass
class PPTOutline:
    """PPT大纲"""
    topic: str
    slides: List[Slide]
    design_tips: List[str]


# 场景模板
SCENE_TEMPLATES = {
    "report": {
        "name": "工作汇报",
        "structure": ["封面", "目录", "回顾", "成果", "问题", "计划", "总结"],
        "slides": [
            Slide("封面", ["{topic}", "汇报人：XXX", "日期：XXXX年XX月"], "title"),
            Slide("目录", ["1. 工作回顾", "2. 核心成果", "3. 问题分析", "4. 下阶段计划"]),
            Slide("工作回顾", ["时间周期：XXXX年XX月-XX月", "主要任务：XXX", "完成情况：XXX", "关键节点：XXX"]),
            Slide("核心成果", ["成果1：XXX（数据支撑）", "成果2：XXX（数据支撑）", "成果3：XXX（数据支撑）", "亮点：XXX"]),
            Slide("问题分析", ["问题1：XXX", "原因：XXX", "影响：XXX", "改进方向：XXX"]),
            Slide("下阶段计划", ["目标：XXX", "关键任务：XXX", "时间节点：XXX", "资源需求：XXX"]),
            Slide("总结", ["核心观点", "感谢语", "联系方式"], "end")
        ]
    },
    "product": {
        "name": "产品发布",
        "structure": ["封面", "痛点", "方案", "产品", "优势", "应用", "行动"],
        "slides": [
            Slide("封面", ["{topic}", "Slogan：一句话卖点"], "title"),
            Slide("市场痛点", ["现状问题1：XXX", "现状问题2：XXX", "用户痛点：XXX", "市场机会：XXX"]),
            Slide("解决方案", ["核心思路：XXX", "解决路径：XXX", "创新点：XXX"]),
            Slide("产品介绍", ["产品名称：XXX", "核心功能：XXX", "技术亮点：XXX", "使用场景：XXX"]),
            Slide("竞争优势", ["对比维度1：我们 vs 竞品A vs 竞品B", "对比维度2：...", "核心优势总结"]),
            Slide("应用场景", ["场景1：XXX", "场景2：XXX", "场景3：XXX", "客户案例：XXX"]),
            Slide("行动号召", ["购买/试用方式", "优惠政策", "联系方式", "二维码"], "end")
        ]
    },
    "training": {
        "name": "培训课件",
        "structure": ["封面", "目标", "理论", "案例", "练习", "总结"],
        "slides": [
            Slide("封面", ["{topic}", "讲师：XXX", "时长：X小时"], "title"),
            Slide("培训目标", ["学完本课，你将能够：", "• 掌握XXX技能", "• 理解XXX概念", "• 应用XXX方法"]),
            Slide("课程大纲", ["1. 基础概念", "2. 核心方法", "3. 实战案例", "4. 练习巩固"]),
            Slide("基础概念", ["概念1：定义+解释", "概念2：定义+解释", "概念3：定义+解释"]),
            Slide("核心方法", ["方法1：步骤说明", "方法2：步骤说明", "注意事项：XXX"]),
            Slide("实战案例", ["案例背景：XXX", "问题分析：XXX", "解决方案：XXX", "效果展示：XXX"]),
            Slide("练习环节", ["练习1：XXX", "练习2：XXX", "参考答案：XXX"]),
            Slide("课程总结", ["核心要点回顾", "延伸学习资源", "答疑时间", "感谢"], "end")
        ]
    },
    "business": {
        "name": "商业计划",
        "structure": ["封面", "市场", "方案", "模式", "团队", "融资"],
        "slides": [
            Slide("封面", ["{topic}", "一句话定位", "公司名称+Logo"], "title"),
            Slide("市场机会", ["市场规模：XXX亿", "增长趋势：年增长率XX%", "痛点分析：XXX", "市场空白：XXX"]),
            Slide("解决方案", ["产品/服务：XXX", "核心价值：XXX", "差异化优势：XXX"]),
            Slide("商业模式", ["盈利模式：XXX", "收入来源：XXX", "成本结构：XXX", "盈利预测：XXX"]),
            Slide("竞争优势", ["技术壁垒：XXX", "资源优势：XXX", "团队背景：XXX", "先发优势：XXX"]),
            Slide("运营数据", ["用户数据：XXX", "营收数据：XXX", "增长数据：XXX", "关键指标：XXX"]),
            Slide("团队介绍", ["创始人：XXX（背景）", "核心团队：XXX", "顾问团队：XXX"]),
            Slide("融资计划", ["融资额度：XXX万", "出让股份：XX%", "资金用途：XXX", "里程碑：XXX"], "end")
        ]
    },
    "proposal": {
        "name": "方案建议",
        "structure": ["封面", "背景", "问题", "方案", "收益", "风险"],
        "slides": [
            Slide("封面", ["{topic}", "提案单位：XXX", "日期：XXXX年XX月"], "title"),
            Slide("项目背景", ["背景介绍：XXX", "现状分析：XXX", "相关方：XXX"]),
            Slide("问题分析", ["核心问题：XXX", "问题影响：XXX", "紧迫性：XXX"]),
            Slide("解决方案", ["方案概述：XXX", "实施步骤：", "  1. XXX", "  2. XXX", "  3. XXX"]),
            Slide("预期收益", ["收益1：XXX（量化）", "收益2：XXX（量化）", "收益3：XXX（量化）"]),
            Slide("风险评估", ["风险1：XXX（应对措施）", "风险2：XXX（应对措施）"]),
            Slide("实施计划", ["阶段1：XXX（时间+负责人）", "阶段2：XXX", "阶段3：XXX"]),
            Slide("结语", ["方案价值总结", "期待合作", "联系方式"], "end")
        ]
    }
}


def generate_outline(topic: str, scene_type: str = "report") -> PPTOutline:
    """生成PPT大纲"""
    template = SCENE_TEMPLATES.get(scene_type, SCENE_TEMPLATES["report"])

    slides = []
    for slide_template in template["slides"]:
        # 替换主题占位符
        title = slide_template.title
        bullets = [b.replace("{topic}", topic) for b in slide_template.bullets]
        slides.append(Slide(title, bullets, slide_template.layout))

    design_tips = [
        f"配色建议：商务蓝(#1E3A5F)适合{template['name']}场景",
        "排版原则：每页不超过6行，每行不超过15字",
        "图表建议：数据用柱状图，趋势用折线图，占比用饼图",
        "字体建议：标题用微软雅黑Bold，正文用微软雅黑",
        "动画建议：简洁为主，避免过多花哨动画"
    ]

    return PPTOutline(topic, slides, design_tips)


def export_markdown(outline: PPTOutline) -> str:
    """导出Markdown格式"""
    lines = [f"# {outline.topic}", ""]

    for i, slide in enumerate(outline.slides, 1):
        lines.append(f"## 第{i}页：{slide.title}")
        for bullet in slide.bullets:
            lines.append(f"- {bullet}")
        lines.append("")

    lines.append("## 设计建议")
    for tip in outline.design_tips:
        lines.append(f"- {tip}")

    return "\n".join(lines)


def export_pptx(outline: PPTOutline, output_path: str):
    """导出PowerPoint文件"""
    if not HAS_PPTX:
        print("请先安装python-pptx: pip install python-pptx")
        return False

    prs = Presentation()

    for slide_data in outline.slides:
        # 选择版式
        if slide_data.layout == "title":
            layout = prs.slide_layouts[0]  # 标题页
        elif slide_data.layout == "end":
            layout = prs.slide_layouts[6]  # 空白页
        else:
            layout = prs.slide_layouts[1]  # 标题和内容

        slide = prs.slides.add_slide(layout)

        # 设置标题
        if slide.shapes.title:
            slide.shapes.title.text = slide_data.title

        # 设置内容
        if len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.text = slide_data.bullets[0] if slide_data.bullets else ""

            for bullet in slide_data.bullets[1:]:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

    prs.save(output_path)
    return True


def list_templates():
    """列出所有模板"""
    print("可用PPT场景模板：")
    print()
    for key, template in SCENE_TEMPLATES.items():
        print(f"【{key}】{template['name']}")
        print(f"  结构：{' → '.join(template['structure'])}")
        print(f"  页数：{len(template['slides'])}页")
        print()


def main():
    import argparse

    parser = argparse.ArgumentParser(description="PPT大纲生成助手")
    subparsers = parser.add_subparsers(dest='command', help='命令')

    # generate命令
    gen_parser = subparsers.add_parser('generate', help='生成PPT大纲')
    gen_parser.add_argument('topic', help='PPT主题')
    gen_parser.add_argument('--type', '-t', default='report',
        choices=['report', 'product', 'training', 'business', 'proposal'],
        help='场景类型')
    gen_parser.add_argument('--export', '-e', choices=['markdown', 'ppt'], default='markdown',
        help='导出格式')
    gen_parser.add_argument('--output', '-o', help='输出文件路径')

    # templates命令
    subparsers.add_parser('templates', help='列出所有模板')

    args = parser.parse_args()

    if not args.command:
        parser.print_help()
        return

    if args.command == 'generate':
        outline = generate_outline(args.topic, args.type)

        if args.export == 'markdown':
            output = export_markdown(outline)
            if args.output:
                with open(args.output, 'w', encoding='utf-8') as f:
                    f.write(output)
                print(f"已保存到：{args.output}")
            else:
                print(output)
        elif args.export == 'ppt':
            output_path = args.output or f"{args.topic}.pptx"
            if export_pptx(outline, output_path):
                print(f"PPT已生成：{output_path}")

    elif args.command == 'templates':
        list_templates()


if __name__ == "__main__":
    main()
