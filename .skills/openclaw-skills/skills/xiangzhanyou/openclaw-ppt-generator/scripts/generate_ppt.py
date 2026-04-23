#!/usr/bin/env python3
"""
OpenClaw PPT Generator
Generate PPT documents using python-pptx library (open source, no API calls)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
import argparse
import json
import os
import sys
import time


def parse_content(content_str):
    """Parse content string into slide data
    
    Format: "Slide1|Slide2:Item1,Item2|Slide3"
    - Use | to separate slides
    - Use : to separate slide title from bullet list
    - Use , to separate bullet items
    """
    slides = []
    slide_parts = content_str.split('|')
    
    for part in slide_parts:
        part = part.strip()
        if not part:
            continue
            
        if ':' in part:
            # Has bullet list
            title, items_str = part.split(':', 1)
            items = [item.strip() for item in items_str.split(',') if item.strip()]
            slides.append({
                'title': title.strip(),
                'content': items
            })
        else:
            # Just text content
            slides.append({
                'title': part,
                'content': []
            })
    
    return slides


def create_ppt(title, slides_data, output_path=None):
    """Create PPT presentation
    
    Args:
        title: Presentation title
        slides_data: List of slide dictionaries with 'title' and 'content'
        output_path: Output file path (optional)
    
    Returns:
        Path to generated PPT file
    """
    # Create presentation
    prs = Presentation()
    
    # Add title slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    
    # Add content slides
    for slide_data in slides_data:
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set slide title
        slide.shapes.title.text = slide_data['title']
        
        # Set content
        content_placeholder = slide.shapes.placeholders[1]
        if slide_data['content']:
            # Add bullet list
            text_frame = content_placeholder.text_frame
            text_frame.clear()
            for item in slide_data['content']:
                p = text_frame.add_paragraph()
                p.text = item
                p.level = 0
        else:
            # Simple text
            content_placeholder.text = slide_data['title']
    
    # Determine output path
    if not output_path:
        timestamp = int(time.time())
        output_path = f"output_{timestamp}.pptx"
    
    # Ensure output directory exists
    output_dir = os.path.dirname(output_path)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    # Save presentation
    prs.save(output_path)
    return output_path


def main():
    parser = argparse.ArgumentParser(
        description="Generate PPT using python-pptx (open source, no API calls)"
    )
    parser.add_argument(
        "--title", "-t",
        type=str,
        required=True,
        help="Presentation title"
    )
    parser.add_argument(
        "--content", "-c",
        type=str,
        required=True,
        help="Content string (use | to separate slides, : for bullet list, , for items)"
    )
    parser.add_argument(
        "--output", "-o",
        type=str,
        default=None,
        help="Output file path (default: output_<timestamp>.pptx)"
    )
    
    args = parser.parse_args()
    
    try:
        # Parse content
        slides_data = parse_content(args.content)
        
        if not slides_data:
            print(json.dumps({
                "status": "error",
                "message": "No valid slides content provided"
            }), file=sys.stderr)
            sys.exit(1)
        
        # Generate PPT
        ppt_path = create_ppt(args.title, slides_data, args.output)
        
        # Output result
        result = {
            "status": "success",
            "ppt_path": ppt_path
        }
        print(json.dumps(result, ensure_ascii=False, indent=2))
        
    except Exception as e:
        print(json.dumps({
            "status": "error",
            "message": str(e)
        }), file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()