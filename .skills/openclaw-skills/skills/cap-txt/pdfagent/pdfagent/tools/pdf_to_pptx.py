from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from .render import render_pdf_to_images


def pdf_to_pptx(input_path: Path, output_path: Path, dpi: int = 150) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    temp_dir = output_path.parent / f".{output_path.stem}_images"
    images = render_pdf_to_images(input_path, temp_dir, dpi=dpi, fmt="png")

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    for image_path in images:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)

    prs.save(str(output_path))
