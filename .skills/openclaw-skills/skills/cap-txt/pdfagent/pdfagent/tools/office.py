from __future__ import annotations

from pathlib import Path
import os

from ..core.deps import require_bin
from ..core.exec import ExternalCommandError, run_cmd


def office_to_pdf(input_path: Path, output_dir: Path, timeout_sec: int = 60) -> Path:
    env_timeout = None
    try:
        env_timeout = int(os.environ.get("PDFAGENT_SOFFICE_TIMEOUT", "0"))
    except ValueError:
        env_timeout = None
    if env_timeout and env_timeout > 0:
        timeout_sec = env_timeout
    soffice = require_bin("soffice", "Install LibreOffice for Office conversions")
    output_dir.mkdir(parents=True, exist_ok=True)
    profile_dir = (output_dir / ".lo-profile").resolve()
    profile_dir.mkdir(parents=True, exist_ok=True)
    profile_uri = profile_dir.as_uri()
    cmd = [
        soffice,
        "--headless",
        "--nologo",
        "--nodefault",
        "--nolockcheck",
        "--norestore",
        "--invisible",
        "--nofirststartwizard",
        f"-env:UserInstallation={profile_uri}",
        "--convert-to",
        "pdf",
        "--outdir",
        str(output_dir),
        str(input_path),
    ]
    try:
        run_cmd(cmd, timeout=timeout_sec)
    except ExternalCommandError:
        return _fallback_office_to_pdf(input_path, output_dir)

    out_path = output_dir / f"{input_path.stem}.pdf"
    if not out_path.exists():
        return _fallback_office_to_pdf(input_path, output_dir)
    return out_path


def _fallback_office_to_pdf(input_path: Path, output_dir: Path) -> Path:
    suffix = input_path.suffix.lower()
    if suffix in {".docx", ".doc"}:
        return _docx_to_pdf_fallback(input_path, output_dir)
    if suffix in {".xlsx", ".xls"}:
        return _xlsx_to_pdf_fallback(input_path, output_dir)
    if suffix in {".pptx", ".ppt"}:
        return _pptx_to_pdf_fallback(input_path, output_dir)
    raise RuntimeError("Unsupported Office file type for fallback conversion")


def _docx_to_pdf_fallback(input_path: Path, output_dir: Path) -> Path:
    try:
        import docx
    except Exception as exc:
        raise RuntimeError("python-docx is required for DOCX fallback") from exc
    from reportlab.pdfgen import canvas

    output_path = output_dir / f"{input_path.stem}.pdf"
    doc = docx.Document(str(input_path))
    c = canvas.Canvas(str(output_path))
    y = 800
    for para in doc.paragraphs:
        if y < 40:
            c.showPage()
            y = 800
        c.setFont("Helvetica", 11)
        c.drawString(40, y, para.text[:120])
        y -= 14
    c.save()
    return output_path


def _xlsx_to_pdf_fallback(input_path: Path, output_dir: Path) -> Path:
    try:
        import openpyxl
    except Exception as exc:
        raise RuntimeError("openpyxl is required for XLSX fallback") from exc
    from reportlab.pdfgen import canvas

    output_path = output_dir / f"{input_path.stem}.pdf"
    wb = openpyxl.load_workbook(str(input_path))
    ws = wb.active
    c = canvas.Canvas(str(output_path))
    y = 800
    for row in ws.iter_rows(values_only=True):
        if y < 40:
            c.showPage()
            y = 800
        line = "  ".join("" if v is None else str(v) for v in row)
        c.setFont("Helvetica", 10)
        c.drawString(40, y, line[:150])
        y -= 14
    c.save()
    return output_path


def _pptx_to_pdf_fallback(input_path: Path, output_dir: Path) -> Path:
    try:
        from pptx import Presentation
    except Exception as exc:
        raise RuntimeError("python-pptx is required for PPTX fallback") from exc
    from reportlab.pdfgen import canvas

    output_path = output_dir / f"{input_path.stem}.pdf"
    prs = Presentation(str(input_path))
    c = canvas.Canvas(str(output_path))
    for slide in prs.slides:
        y = 760
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                c.setFont("Helvetica", 12)
                c.drawString(40, y, shape.text[:120])
                y -= 16
        c.showPage()
    c.save()
    return output_path
