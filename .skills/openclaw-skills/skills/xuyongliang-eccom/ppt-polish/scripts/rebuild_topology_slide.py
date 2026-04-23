#!/usr/bin/env python3
"""Create a polished rebuilt copy of a single-slide PPT topology/diagram deck.

This is intentionally heuristic-driven rather than template-perfect.
It is useful when a user asks to optimize an already-generated PPT diagram
into a cleaner presentation slide.
"""

from __future__ import annotations

import argparse
from pathlib import Path
from zipfile import ZipFile
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


TITLE_FILL = (31, 45, 78)
TITLE_LINE = (82, 112, 188)
BG = (13, 20, 36)
LAYER_STYLES = {
    "external": ((23, 31, 54), (79, 104, 168)),
    "core": ((29, 63, 84), (73, 160, 186)),
    "data": ((56, 38, 80), (160, 109, 205)),
    "ops": ((98, 63, 38), (211, 145, 91)),
}


def extract_texts(pptx_path: Path) -> list[str]:
    texts: list[str] = []
    with ZipFile(pptx_path) as zf:
        slide_names = sorted(
            n for n in zf.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")
        )
        if not slide_names:
            return texts
        data = zf.read(slide_names[0]).decode("utf-8", "ignore")
        texts = re.findall(r"<a:t>(.*?)</a:t>", data)
    return [t.strip() for t in texts if t.strip()]


def classify_items(texts: list[str]) -> dict[str, list[str]]:
    buckets = {"external": [], "core": [], "data": [], "ops": []}
    for t in texts:
        if t in {"外部接入层", "核心平台层", "数据资源层", "运维保障", "图例"}:
            continue
        if any(k in t for k in ["MySQL", "Redis", "文件存储", "MQ", "日志", "ES", "BI"]):
            buckets["data"].append(t)
        elif any(k in t for k in ["监控", "备份", "运维", "巡检", "告警", "报表"]):
            buckets["ops"].append(t)
        elif any(k in t for k in ["用户", "第三方", "公网", "Internet", "门户", "移动", "管理端", "外部"]):
            buckets["external"].append(t)
        else:
            buckets["core"].append(t)
    for k in buckets:
        dedup = []
        seen = set()
        for item in buckets[k]:
            if item not in seen:
                seen.add(item)
                dedup.append(item)
        buckets[k] = dedup
    return buckets


def add_box(slide, x, y, w, h, text, fill, line, font_size=12, bold=False):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill)
    shape.line.color.rgb = RGBColor(*line)
    shape.line.width = Pt(1.1)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.name = "Microsoft YaHei"
    r.font.color.rgb = RGBColor(240, 245, 255)
    return shape


def add_line(slide, x1, y1, x2, y2, color, width=1.5):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = RGBColor(*color)
    c.line.width = Pt(width)
    return c


def build_slide(texts: list[str], output: Path, title: str):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(*BG)

    add_box(slide, Inches(0.35), Inches(0.22), Inches(4.2), Inches(0.55), title, TITLE_FILL, TITLE_LINE, 24, True)
    add_box(
        slide,
        Inches(4.8),
        Inches(0.24),
        Inches(5.2),
        Inches(0.46),
        "自动优化版：层次归一 / 连线降噪 / 更适合汇报",
        (20, 29, 51),
        (58, 83, 143),
        11,
        False,
    )

    headers = [
        ("外部接入层", 0.45, 1.02, 3.0, 0.40, (38, 56, 96), (91, 122, 194)),
        ("核心平台层", 3.72, 1.02, 5.95, 0.40, (28, 70, 87), (66, 156, 179)),
        ("数据资源层", 10.0, 1.02, 1.45, 0.40, (73, 51, 97), (154, 101, 198)),
        ("运维保障", 11.65, 1.02, 1.2, 0.40, (90, 58, 36), (194, 129, 78)),
    ]
    for text, x, y, w, h, fill, line in headers:
        add_box(slide, Inches(x), Inches(y), Inches(w), Inches(h), text, fill, line, 14, True)

    buckets = classify_items(texts)

    # external
    ext_x_positions = [0.55, 1.52, 2.49, 0.55, 1.95, 0.55, 1.62]
    ext_y_positions = [1.72, 1.72, 1.72, 2.46, 2.46, 3.20, 3.20]
    ext_widths = [0.84, 0.84, 0.95, 1.28, 1.38, 0.95, 0.82]
    external = buckets["external"][:7] or ["用户终端", "管理端", "第三方平台", "Internet / 公网", "外部业务系统", "移动应用", "门户"]
    external_shapes = []
    for i, t in enumerate(external):
        s = add_box(
            slide,
            Inches(ext_x_positions[i]),
            Inches(ext_y_positions[i]),
            Inches(ext_widths[i]),
            Inches(0.56),
            t,
            *LAYER_STYLES["external"],
            font_size=11,
        )
        external_shapes.append(s)

    gateway = add_box(slide, Inches(0.98), Inches(4.18), Inches(2.22), Inches(0.72), "统一入口 / API 网关 / 负载均衡", (45, 74, 140), (121, 162, 255), 14, True)

    # core area
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(3.78), Inches(1.60), Inches(5.8), Inches(4.45))
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(17, 31, 46)
    frame.line.color.rgb = RGBColor(49, 113, 135)
    frame.line.width = Pt(1.4)

    core_items = buckets["core"][:11] or [
        "认证服务", "应用服务 A", "应用服务 B", "应用服务 C", "应用服务 D", "管理后台", "支撑服务", "开放服务",
        "调度中心 / 服务编排", "消息总线 / 中间件", "规则引擎 / 任务中心",
    ]
    top_row = core_items[:4]
    mid_row = core_items[4:8]
    bottom_row = core_items[8:11]

    top_x = [4.05, 5.3, 6.6, 7.9]
    mid_x = [4.7, 6.0, 7.3, 8.55]
    bottom_x = [4.42, 6.34, 8.10]
    top_shapes = []
    for i, t in enumerate(top_row):
        top_shapes.append(add_box(slide, Inches(top_x[i]), Inches(1.88), Inches(1.08), Inches(0.60), t, *LAYER_STYLES["core"], font_size=12, bold=True))
    mid_shapes = []
    for i, t in enumerate(mid_row):
        mid_shapes.append(add_box(slide, Inches(mid_x[i]), Inches(2.72), Inches(1.02 if i < 3 else 0.88), Inches(0.60), t, *LAYER_STYLES["core"], font_size=12))
    bottom_shapes = []
    for i, t in enumerate(bottom_row):
        bottom_shapes.append(add_box(slide, Inches(bottom_x[i]), Inches(3.72), Inches(1.68 if i == 0 else 1.54 if i == 1 else 1.2), Inches(0.62), t, (24, 88, 103), (79, 192, 198), font_size=11, bold=True))

    # data area
    data_items = buckets["data"][:7] or ["MySQL", "Redis", "文件存储", "MQ", "日志系统", "ES", "BI"]
    data_shapes = []
    data_y = [1.82, 2.56, 3.30, 4.04, 4.78, 5.52, 5.52]
    data_x = [10.12, 10.12, 10.12, 10.12, 10.12, 10.12, 10.88]
    data_w = [0.92, 0.92, 1.02, 0.82, 0.95, 0.72, 0.62]
    for i, t in enumerate(data_items):
        data_shapes.append(add_box(slide, Inches(data_x[i]), Inches(data_y[i]), Inches(data_w[i]), Inches(0.52), t, *LAYER_STYLES["data"], font_size=11))

    # ops area
    ops_items = buckets["ops"][:3] or ["监控平台", "备份 / 运维工具", "巡检 / 告警 / 报表"]
    ops_shapes = []
    ops_specs = [(11.8, 1.82, 0.85), (11.72, 2.72, 1.02), (11.68, 3.92, 1.08)]
    for i, t in enumerate(ops_items):
        x, y, w = ops_specs[i]
        ops_shapes.append(add_box(slide, Inches(x), Inches(y), Inches(w), Inches(0.62), t, *LAYER_STYLES["ops"], font_size=11))

    for s in external_shapes:
        add_line(slide, s.left + s.width / 2, s.top + s.height, gateway.left + gateway.width / 2, gateway.top, (121, 162, 255), 1.4)

    for s in top_shapes + mid_shapes:
        add_line(slide, gateway.left + gateway.width, gateway.top + gateway.height / 2, s.left, s.top + s.height / 2, (90, 172, 198), 1.2)

    for s in top_shapes + mid_shapes:
        add_line(slide, s.left + s.width / 2, s.top + s.height, s.left + s.width / 2, Inches(3.72), (88, 183, 200), 0.9)

    for s in bottom_shapes:
        add_line(slide, s.left + s.width, s.top + s.height / 2, Inches(10.0), s.top + s.height / 2, (165, 121, 222), 1.0)

    for s in ops_shapes:
        add_line(slide, Inches(9.3), Inches(4.03), s.left, s.top + s.height / 2, (212, 151, 86), 1.0)

    note = slide.shapes.add_textbox(Inches(4.20), Inches(6.45), Inches(8.6), Inches(0.4))
    p = note.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "说明：本页由规则化重构生成，目标是提升结构清晰度、层次感和汇报可读性，同时保留原始图中的核心信息。"
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(11)
    r.font.color.rgb = RGBColor(194, 205, 229)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("input", help="Source PPTX file")
    parser.add_argument("-o", "--output", required=True, help="Output PPTX file")
    parser.add_argument("--title", default="系统拓扑架构图", help="Slide title")
    args = parser.parse_args()

    input_path = Path(args.input)
    output_path = Path(args.output)
    texts = extract_texts(input_path)
    build_slide(texts, output_path, args.title)
    print(output_path)


if __name__ == "__main__":
    main()
