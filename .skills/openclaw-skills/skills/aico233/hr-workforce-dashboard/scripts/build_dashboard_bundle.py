#!/usr/bin/env python3
from __future__ import annotations

import argparse
import base64
import math
import re
import shutil
import textwrap
import zipfile
from datetime import datetime
from html import escape as html_escape
from io import BytesIO
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import pandas as pd
from matplotlib.patches import Rectangle

plt.rcParams["font.sans-serif"] = ["PingFang SC", "Hiragino Sans GB", "Arial Unicode MS", "DejaVu Sans"]
plt.rcParams["axes.unicode_minus"] = False
from openpyxl import Workbook, load_workbook
from openpyxl.drawing.image import Image as XLImage
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches, Pt

REGION_ORDER = ["Americas", "APAC", "EMEA"]
REGION_ALL_ORDER = ["Americas", "APAC", "EMEA", "Greater China"]
REGION_COLORS = {
    "Americas": "#1F6B8F",
    "APAC": "#ED7D31",
    "EMEA": "#2E7D32",
}
ROW_FILL = {
    "Americas": "F5E1D6",
    "APAC": "DCE8F4",
    "EMEA": "F3E0D4",
    "Greater China": "D8E3EE",
    "subtotal": "7AC2E3",
    "grand": "4FA9D5",
    "missing": "F6E3A1",
}
BG_DISPLAY_ORDER = ["合计", "IEG", "CSIG", "WXG", "TEG", "CDG", "PCG", "OFS", "S1", "S2", "S3"]
BG_PATTERNS = {
    "IEG": ["ieg - interactive entertainment group"],
    "CSIG": ["csig - cloud & smart industries group"],
    "WXG": ["wxg - weixin group"],
    "TEG": ["teg - technology & engineering group"],
    "CDG": ["cdg - corporate development group"],
    "PCG": ["pcg - platform & content group"],
    "OFS": ["overseas functional system"],
    "S1": ["s1 - functional line"],
    "S2": ["s2 - finance line"],
    "S3": ["s3 - hr & management line"],
}
TERM_CATEGORY_MAP = {
    "terminate employee > voluntary": "Voluntary",
    "terminate employee > involuntary": "Involuntary",
    "terminate employee > others": "Others",
    "terminate employee > statutory termination": "Statutory",
}
COMMON_TEXT_COLS = ["Employee Type", "Region", "WD Employee ID", "Country/Territory", "BG", "Termination Category"]

# --- Contingent Worker (看板5) constants ---
# Country mapping: Chinese name -> (English name, Region)
CONTINGENT_COUNTRY_MAP: dict[str, tuple[str, str]] = {
    "爱尔兰": ("Ireland", "EMEA"),
    "俄罗斯": ("Russia", "EMEA"),
    "菲律宾": ("Philippines", "APAC"),
    "哈萨克斯坦": ("Kazakhstan", "EMEA"),
    "罗马尼亚": ("Romania", "EMEA"),
    "孟加拉国": ("Bangladesh", "APAC"),
    "挪威": ("Norway", "EMEA"),
    "葡萄牙": ("Portugal", "EMEA"),
    "西班牙": ("Spain", "EMEA"),
    "越南": ("Vietnam", "APAC"),
    "美国": ("United States of America", "Americas"),
    "墨西哥": ("Mexico", "Americas"),
    "智利": ("Chile", "Americas"),
    "加拿大": ("Canada", "Americas"),
    "巴西": ("Brazil", "Americas"),
    "泰国": ("Thailand", "APAC"),
    "新加坡": ("Singapore", "APAC"),
    "巴基斯坦": ("Pakistan", "APAC"),
    "新西兰": ("New Zealand", "APAC"),
    "马来西亚": ("Malaysia", "APAC"),
    "韩国": ("Korea, Republic of", "APAC"),
    "日本": ("Japan", "APAC"),
    "印度尼西亚": ("Indonesia", "APAC"),
    "印度": ("India", "APAC"),
    "澳大利亚": ("Australia", "APAC"),
    "英国": ("United Kingdom", "EMEA"),
    "阿联酋": ("United Arab Emirates", "EMEA"),
    "阿拉伯联合酋长国": ("United Arab Emirates", "EMEA"),
    "土耳其": ("Türkiye", "EMEA"),
    "瑞典": ("Sweden", "EMEA"),
    "沙特阿拉伯": ("Saudi Arabia", "EMEA"),
    "波兰": ("Poland", "EMEA"),
    "荷兰": ("Netherlands", "EMEA"),
    "意大利": ("Italy", "EMEA"),
    "德国": ("Germany", "EMEA"),
    "法国": ("France", "EMEA"),
    "埃及": ("Egypt", "EMEA"),
    "乌兹别克斯坦": ("Uzbekistan", "APAC"),
    "柬埔寨": ("Cambodia", "APAC"),
}

WORKER_TYPE_PREFIX = {"v_": "Contractor", "p_": "Partner"}


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build workforce dashboard bundle from Excel detail files.")
    parser.add_argument("--files", nargs="+", required=True, help="Input Excel files")
    parser.add_argument("--output-dir", required=True, help="Output directory")
    parser.add_argument("--bundle-name", default="workforce_dashboard", help="Final zip and workbook bundle name")
    parser.add_argument("--title", default="人力数据看板", help="Dashboard title")
    return parser.parse_args()


def clean_text(value) -> str:
    if value is None:
        return ""
    if isinstance(value, float) and math.isnan(value):
        return ""
    return str(value).strip()


def normalize_key(value) -> str:
    text = clean_text(value).lower()
    text = re.sub(r"[^a-z0-9]+", "_", text)
    return text.strip("_")


def classify_dataset_type(value: object) -> str:
    text = clean_text(value).lower()
    if any(token in text for token in ["terminat", "attrition", "离职"]):
        return "termination"
    if any(token in text for token in ["contingent", "contractor", "partner", "外包"]):
        return "contingent"
    if any(token in text for token in ["active", "在职", "headcount"]):
        return "active"
    raise ValueError(f"无法识别 dataset_type: {value}")


def parse_any_date(value: object) -> pd.Timestamp | pd.NaT:
    if value is None or clean_text(value) == "":
        return pd.NaT
    try:
        return pd.to_datetime(value, errors="coerce")
    except Exception:
        return pd.NaT


def _read_pipe_row(ws, row_num: int) -> tuple[str, str]:
    """Read a metadata row in 'Label | Value' format. Supports two layouts:
    - Two cells: A=label, B=value (value may be a datetime object)
    - One merged cell: 'Label | Value' text"""
    vals = []
    for cell in next(ws.iter_rows(min_row=row_num, max_row=row_num)):
        if cell.value is not None:
            vals.append(cell.value)
    if len(vals) >= 2:
        return clean_text(vals[0]), vals[1]  # keep raw value for date parsing
    if len(vals) == 1:
        text = clean_text(vals[0])
        if "|" in text:
            parts = text.split("|", 1)
            return parts[0].strip(), parts[1].strip()
    return "", ""


def normalize_region(value: object) -> str:
    text = clean_text(value)
    lower = text.lower()
    if lower == "":
        return "Unknown"
    if "america" in lower:
        return "Americas"
    if lower == "apac":
        return "APAC"
    if "emea" in lower:
        return "EMEA"
    if lower in {"greater china", "greater_china"}:
        return "Greater China"
    if lower in {"mainland china", "hong kong"}:
        return text
    return text


def is_greater_china(region: object, country: object) -> bool:
    region_text = clean_text(region).lower()
    country_text = clean_text(country).lower()
    return region_text == "greater china" or country_text in {"mainland china", "hong kong"}


def map_bg(value: object) -> str | None:
    text = clean_text(value).lower()
    if not text:
        return None
    for display, patterns in BG_PATTERNS.items():
        if any(pattern in text for pattern in patterns):
            return display
    return None


def _resolve_employee_id_column(df: pd.DataFrame) -> pd.DataFrame:
    """Normalize employee ID column name to 'WD Employee ID'."""
    if "WD Employee ID" in df.columns:
        return df
    id_candidates = [c for c in df.columns if "employee" in c.lower() and "id" in c.lower()]
    if id_candidates:
        df = df.rename(columns={id_candidates[0]: "WD Employee ID"})
    return df


def read_dataset(file_path: Path) -> tuple[pd.DataFrame, str, dict[str, str]]:
    """New-mode only: parse metadata from structured header rows.
    Active files: rows 1-6 metadata, row 7 header, row 8+ data.
    Termination files: rows 1-7 metadata, row 8 header, row 9+ data.
    """
    wb = load_workbook(file_path, read_only=True, data_only=True)
    ws = wb[wb.sheetnames[0]]

    # Row 1 = report name -> determines dataset_type
    row1_vals = [c.value for c in next(ws.iter_rows(min_row=1, max_row=1))]
    report_name = clean_text(row1_vals[0]) if row1_vals else ""
    dataset_type = classify_dataset_type(report_name)

    metadata: dict[str, str] = {"dataset_type": dataset_type, "report_name": report_name}

    if dataset_type == "active":
        # Active: row 4 = Effective Date | <date>  ->  snapshot_date
        _label, raw_val = _read_pipe_row(ws, 4)
        snapshot_date = parse_any_date(raw_val)
        if pd.isna(snapshot_date):
            raise ValueError(f"文件 {file_path.name} 第4行 Effective Date 无法解析: {raw_val}")
        metadata["snapshot_date"] = str(snapshot_date.date())
        header_row = 7  # 0-indexed for pandas: header=6
    else:
        # Termination: row 4 = Termination Date From | <start_date>
        #              row 5 = Termination Date To   | <end_date>
        _label4, raw_start = _read_pipe_row(ws, 4)
        _label5, raw_end = _read_pipe_row(ws, 5)
        period_start = parse_any_date(raw_start)
        period_end = parse_any_date(raw_end)
        if pd.isna(period_start):
            raise ValueError(f"文件 {file_path.name} 第4行 Termination Date From 无法解析: {raw_start}")
        if pd.isna(period_end):
            raise ValueError(f"文件 {file_path.name} 第5行 Termination Date To 无法解析: {raw_end}")
        metadata["period_start"] = str(period_start.date())
        metadata["period_end"] = str(period_end.date())
        # Use period_end as snapshot_date for termination files
        snapshot_date = period_end
        metadata["snapshot_date"] = str(period_end.date())
        header_row = 8  # 0-indexed for pandas: header=7
    wb.close()

    snapshot_ts = parse_any_date(metadata["snapshot_date"])
    if pd.isna(snapshot_ts):
        raise ValueError(f"文件 {file_path.name} 的 snapshot_date 无法解析: {metadata['snapshot_date']}")

    # header_row is 1-indexed; pandas header param is 0-indexed
    df = pd.read_excel(file_path, header=header_row - 1)
    df = df.dropna(how="all")
    df.columns = [clean_text(col) for col in df.columns]
    df = _resolve_employee_id_column(df)
    df["snapshot_date"] = snapshot_ts.normalize()
    df["dataset_type"] = dataset_type
    df["__source_file__"] = file_path.name
    for col in COMMON_TEXT_COLS:
        if col in df.columns:
            df[col] = df[col].apply(clean_text)
    if "Region" in df.columns:
        df["Region_std"] = df["Region"].apply(normalize_region)
    else:
        df["Region_std"] = "Unknown"
    if "Country/Territory" in df.columns:
        df["Country/Territory"] = df["Country/Territory"].replace("", "Unknown")
    df["snapshot_month"] = df["snapshot_date"].dt.to_period("M")
    if "Employee Type" in df.columns:
        df["Employee Type"] = df["Employee Type"].str.strip().replace({"regular": "Regular", "intern": "Intern"}, regex=False)
    if "WD Employee ID" in df.columns:
        df["WD Employee ID"] = df["WD Employee ID"].astype(str).str.strip()
    if "BG" in df.columns:
        df["BG_mapped"] = df["BG"].apply(map_bg)
    else:
        df["BG_mapped"] = None
    for date_col in ["Hire Date", "Termination Date"]:
        if date_col in df.columns:
            df[date_col] = pd.to_datetime(df[date_col], errors="coerce")
    required = ["Employee Type", "Region", "WD Employee ID", "Country/Territory", "BG"]
    if dataset_type == "termination":
        required += ["Termination Date", "Termination Category"]
    missing = [col for col in required if col not in df.columns]
    if missing:
        raise ValueError(f"文件 {file_path.name} 缺少字段: {', '.join(missing)}")
    return df, dataset_type, metadata


def read_contingent_dataset(file_path: Path) -> tuple[pd.DataFrame, dict[str, str]]:
    """Read Contingent Worker report. Row 1 = report name, Row 2 = header, Row 3+ = data.
    No structured metadata rows; the file is always the latest snapshot."""
    wb = load_workbook(file_path, read_only=True, data_only=True)
    ws = wb[wb.sheetnames[0]]
    row1_vals = [c.value for c in next(ws.iter_rows(min_row=1, max_row=1))]
    report_name = clean_text(row1_vals[0]) if row1_vals else ""
    wb.close()

    df = pd.read_excel(file_path, header=1)  # Row 2 = header (0-indexed: 1)
    df = df.dropna(how="all")
    df.columns = [clean_text(col) for col in df.columns]

    # Resolve employee ID column
    if "HQ Employee ID" in df.columns:
        df = df.rename(columns={"HQ Employee ID": "Worker_ID"})
    else:
        id_candidates = [c for c in df.columns if "employee" in c.lower() and "id" in c.lower()]
        if id_candidates:
            df = df.rename(columns={id_candidates[0]: "Worker_ID"})

    df["Worker_ID"] = df["Worker_ID"].astype(str).str.strip()

    # Determine worker type from WeCom Name prefix (v_ = Contractor, p_ = Partner)
    if "WeCom Name" in df.columns:
        def _worker_type(name):
            name_str = clean_text(name).lower()
            for prefix, wtype in WORKER_TYPE_PREFIX.items():
                if name_str.startswith(prefix):
                    return wtype
            return "Contractor"  # default
        df["Worker_Type"] = df["WeCom Name"].apply(_worker_type)
    else:
        df["Worker_Type"] = "Contractor"

    # Map Country (Chinese -> English) and assign Region
    raw_country = df["Country/Territory"].apply(clean_text) if "Country/Territory" in df.columns else pd.Series(["Unknown"] * len(df))
    df["Country_EN"] = raw_country.map(lambda c: CONTINGENT_COUNTRY_MAP.get(c, (c, "Unknown"))[0])
    df["Region_std"] = raw_country.map(lambda c: CONTINGENT_COUNTRY_MAP.get(c, (c, "Unknown"))[1])

    # Map BG using existing map_bg
    if "BG" in df.columns:
        df["BG"] = df["BG"].apply(clean_text)
        df["BG_mapped"] = df["BG"].apply(map_bg)
    else:
        df["BG_mapped"] = None

    metadata = {"dataset_type": "contingent", "report_name": report_name}
    return df, metadata


def unique_employee_count(df: pd.DataFrame) -> int:
    if df.empty:
        return 0
    return df["WD Employee ID"].dropna().astype(str).nunique()


def filter_regular(df: pd.DataFrame) -> pd.DataFrame:
    if df.empty:
        return df.copy()
    return df[df["Employee Type"].eq("Regular")].copy()


def filter_intern(df: pd.DataFrame) -> pd.DataFrame:
    if df.empty:
        return df.copy()
    return df[df["Employee Type"].eq("Intern")].copy()


def filter_overseas(df: pd.DataFrame) -> pd.DataFrame:
    if df.empty:
        return df.copy()
    return df[df["Region_std"].isin(REGION_ORDER)].copy()


def month_label(period: pd.Period) -> str:
    return period.to_timestamp().strftime("%Y %b")


def compact_date_label(value: pd.Timestamp | None) -> str:
    if value is None or pd.isna(value):
        return "N/A"
    return f"{value.year}-{value.month}-{value.day}"


def month_day_label(value: pd.Timestamp | None) -> str:
    if value is None or pd.isna(value):
        return "N/A"
    return f"{value.month}-{value.day}"


def build_dashboard1(active_df: pd.DataFrame, notes: list[str]) -> tuple[dict, pd.DataFrame]:
    regular = filter_overseas(filter_regular(active_df))
    if regular.empty:
        return {
            "latest_month": None,
            "yoy_months": [],
            "mom_months": [],
            "incomplete": True,
            "notes": ["缺少可用于看板1的海外正式员工在职数据"],
        }, pd.DataFrame(columns=["section", "month", *REGION_ORDER])

    monthly = (
        regular.drop_duplicates(["snapshot_month", "Region_std", "WD Employee ID"])
        .groupby(["snapshot_month", "Region_std"])["WD Employee ID"]
        .nunique()
        .reset_index(name="headcount")
    )
    latest_month = monthly["snapshot_month"].max()
    yoy_month = latest_month - 12
    # Build candidate MoM months: up to 3 months ending at latest
    mom_candidates = [latest_month - 2, latest_month - 1, latest_month]
    available_months = set(monthly["snapshot_month"].tolist())

    # Filter to actually available MoM months
    mom_periods = [period for period in mom_candidates if period in available_months]
    mom_count = len(mom_periods)

    missing_local = []
    if yoy_month not in available_months:
        missing_local.append(f"看板1缺少同比月份 {month_label(yoy_month)}")

    if mom_count < 2:
        missing_local.append("看板1环比图需至少2个月数据，当前数据不足")
        mom_periods = []  # Not enough data for MoM chart
    notes.extend(missing_local)

    def rows_for(periods: list[pd.Period], section: str) -> list[dict]:
        records = []
        for period in periods:
            row = {"section": section, "month": month_label(period)}
            month_slice = monthly[monthly["snapshot_month"].eq(period)]
            for region in REGION_ORDER:
                region_slice = month_slice[month_slice["Region_std"].eq(region)]
                row[region] = int(region_slice["headcount"].sum()) if not region_slice.empty else 0
            records.append(row)
        return records

    yoy_periods = [period for period in [yoy_month, latest_month] if period in available_months]
    data_rows = rows_for(yoy_periods, "同比") + rows_for(mom_periods, "环比")
    data_df = pd.DataFrame(data_rows)
    return {
        "latest_month": month_label(latest_month),
        "yoy_months": [month_label(period) for period in yoy_periods],
        "mom_months": [month_label(period) for period in mom_periods],
        "mom_count": mom_count if mom_count >= 2 else 0,
        "incomplete": bool(missing_local),
        "notes": missing_local,
    }, data_df


def build_dashboard2(active_df: pd.DataFrame, term_df: pd.DataFrame, notes: list[str]) -> tuple[dict, pd.DataFrame]:
    regular = filter_overseas(filter_regular(active_df))
    empty = pd.DataFrame(
        [{"Metric": metric, "Region": region, "Headcount": 0, "Share": 0.0} for metric in ["Active", "AttritionYTD"] for region in REGION_ORDER]
    )
    if regular.empty:
        notes.append("看板2缺少海外正式员工在职数据")
        return {"latest_month": None, "active_snapshot_date": None, "attrition_cutoff_date": None, "attrition_year": None}, empty

    latest_snapshot = regular["snapshot_date"].max().normalize()
    latest_month = regular["snapshot_month"].max()
    latest_df = regular[regular["snapshot_date"].eq(latest_snapshot)].drop_duplicates(["Region_std", "WD Employee ID"])
    active_total = unique_employee_count(latest_df)

    rows = []
    for region in REGION_ORDER:
        count = unique_employee_count(latest_df[latest_df["Region_std"].eq(region)])
        share = round((count / active_total) * 100, 1) if active_total else 0.0
        rows.append({"Metric": "Active", "Region": region, "Headcount": count, "Share": share})

    attrition_cutoff = latest_snapshot
    attrition_year = latest_snapshot.year
    term_regular = filter_overseas(filter_regular(term_df)).copy() if not term_df.empty else pd.DataFrame()
    attrition_total = 0
    if term_regular.empty or "Termination Date" not in term_regular.columns:
        notes.append(f"看板2缺少截至 {compact_date_label(attrition_cutoff)} 的海外正式员工离职数据")
        ytd_term = pd.DataFrame(columns=term_regular.columns if not term_regular.empty else [])
    else:
        year_start = pd.Timestamp(year=attrition_year, month=1, day=1)
        ytd_term = term_regular[
            term_regular["Termination Date"].between(year_start, attrition_cutoff, inclusive="both")
        ].copy()
        attrition_total = unique_employee_count(ytd_term)
        if attrition_total == 0:
            notes.append(f"看板2缺少截至 {compact_date_label(attrition_cutoff)} 的海外正式员工离职数据")

    for region in REGION_ORDER:
        count = unique_employee_count(ytd_term[ytd_term["Region_std"].eq(region)]) if not ytd_term.empty else 0
        share = round((count / attrition_total) * 100, 1) if attrition_total else 0.0
        rows.append({"Metric": "AttritionYTD", "Region": region, "Headcount": count, "Share": share})

    return {
        "latest_month": month_label(latest_month),
        "active_snapshot_date": latest_snapshot,
        "attrition_cutoff_date": attrition_cutoff,
        "attrition_year": attrition_year,
    }, pd.DataFrame(rows)


def bg_pair_counts(df: pd.DataFrame, bg: str) -> tuple[int, int]:
    subset = df[df["BG_mapped"].eq(bg)]
    return unique_employee_count(filter_regular(subset)), unique_employee_count(filter_intern(subset))


def summary_row(region_label: str, country_label: str, subset: pd.DataFrame, row_kind: str) -> dict:
    row = {
        "Region": region_label,
        "Country/Territory": country_label,
        "Total R": unique_employee_count(filter_regular(subset)),
        "Total I": unique_employee_count(filter_intern(subset)),
        "row_kind": row_kind,
    }
    for bg in BG_DISPLAY_ORDER[1:]:
        reg_count, int_count = bg_pair_counts(subset, bg)
        row[f"{bg} R"] = reg_count
        row[f"{bg} I"] = int_count
    return row


def build_people_summary_lines(subset: pd.DataFrame) -> list[str]:
    regular = unique_employee_count(filter_regular(subset))
    intern = unique_employee_count(filter_intern(subset))
    return [f"Regular {regular}", f"Intern {intern}"]


def build_dashboard3(active_df: pd.DataFrame, notes: list[str]) -> tuple[dict, pd.DataFrame]:
    if active_df.empty:
        columns = ["Region", "Country/Territory", "Total R", "Total I"] + [f"{bg} {suffix}" for bg in BG_DISPLAY_ORDER[1:] for suffix in ["R", "I"]] + ["row_kind"]
        return {"latest_month": None, "region_annotations": {}}, pd.DataFrame(columns=columns)

    latest_month = active_df["snapshot_month"].max()
    latest_df = active_df[active_df["snapshot_month"].eq(latest_month)].drop_duplicates(["WD Employee ID", "snapshot_month"])
    unknown_bgs = sorted({value for value in latest_df[latest_df["BG_mapped"].isna()]["BG"].dropna().tolist() if clean_text(value)})
    if unknown_bgs:
        notes.append("未识别 BG: " + ", ".join(unknown_bgs))

    overseas_df = latest_df[latest_df["Region_std"].isin(REGION_ORDER)].copy()
    gc_df = latest_df[latest_df.apply(lambda row: is_greater_china(row.get("Region_std"), row.get("Country/Territory")), axis=1)].copy()

    rows: list[dict] = []
    region_annotations: dict[str, list[str]] = {}
    for region in REGION_ORDER:
        region_df = overseas_df[overseas_df["Region_std"].eq(region)].copy()
        if region_df.empty:
            continue
        region_annotations[region] = build_people_summary_lines(region_df)
        country_order = (
            region_df.groupby("Country/Territory")["WD Employee ID"]
            .nunique()
            .sort_values(ascending=False)
            .index.tolist()
        )
        for country in country_order:
            country_df = region_df[region_df["Country/Territory"].eq(country)]
            rows.append(summary_row(region, country, country_df, region))

    rows.append(summary_row("海外合计", "", overseas_df, "grand"))

    if not gc_df.empty:
        gc_country_order = (
            gc_df.groupby("Country/Territory")["WD Employee ID"]
            .nunique()
            .sort_values(ascending=False)
            .index.tolist()
        )
        for country in gc_country_order:
            rows.append(summary_row("Greater China", country, gc_df[gc_df["Country/Territory"].eq(country)], "Greater China"))

    combined_df = pd.concat([overseas_df, gc_df], ignore_index=True)
    rows.append(summary_row("海外合计（含试点国内）", "", combined_df, "grand"))
    return {"latest_month": month_label(latest_month), "region_annotations": region_annotations}, pd.DataFrame(rows)


def choose_attrition_period(active_df: pd.DataFrame, term_df: pd.DataFrame) -> tuple[pd.Timestamp, pd.Timestamp, pd.Timestamp]:
    end_date = active_df["snapshot_date"].max()
    if term_df.empty or term_df["Termination Date"].dropna().empty:
        previous = sorted([d for d in active_df["snapshot_date"].drop_duplicates().tolist() if d < end_date])
        if previous:
            start_anchor = previous[-1]
        else:
            start_anchor = end_date
    else:
        term_start = term_df["Termination Date"].min().normalize()
        candidates = sorted([d for d in active_df["snapshot_date"].drop_duplicates().tolist() if d <= term_start])
        if candidates:
            start_anchor = candidates[-1]
        else:
            previous = sorted([d for d in active_df["snapshot_date"].drop_duplicates().tolist() if d < end_date])
            start_anchor = previous[-1] if previous else end_date
    return start_anchor, end_date, end_date


def calculate_attrition_metrics(start_df: pd.DataFrame, end_df: pd.DataFrame, term_df: pd.DataFrame) -> dict[str, float]:
    start_hc = unique_employee_count(start_df)
    end_hc = unique_employee_count(end_df)
    denominator = (start_hc + end_hc) / 2
    no_attrition = unique_employee_count(term_df)
    has_reason = "Term Reason Std" in term_df.columns
    voluntary_attrition = unique_employee_count(term_df[term_df["Term Reason Std"].eq("Voluntary")]) if has_reason else 0
    involuntary_attrition = unique_employee_count(term_df[term_df["Term Reason Std"].eq("Involuntary")]) if has_reason else 0
    others_attrition = unique_employee_count(term_df[term_df["Term Reason Std"].eq("Others")]) if has_reason else 0
    statutory_attrition = unique_employee_count(term_df[term_df["Term Reason Std"].eq("Statutory")]) if has_reason else 0

    def calc_rate(count: int) -> float:
        return round((count / denominator) * 100, 1) if denominator else 0.0

    return {
        "start_hc": start_hc,
        "end_hc": end_hc,
        "no_attrition": no_attrition,
        "rate": calc_rate(no_attrition),
        "voluntary_attrition": voluntary_attrition,
        "voluntary_rate": calc_rate(voluntary_attrition),
        "involuntary_attrition": involuntary_attrition,
        "involuntary_rate": calc_rate(involuntary_attrition),
        "others_attrition": others_attrition,
        "others_rate": calc_rate(others_attrition),
        "statutory_attrition": statutory_attrition,
        "statutory_rate": calc_rate(statutory_attrition),
    }


def rebuild_year_start_population(active_df: pd.DataFrame, term_df: pd.DataFrame, period_start: pd.Timestamp, first_snapshot: pd.Timestamp) -> pd.DataFrame:
    first_snapshot_df = active_df[active_df["snapshot_date"].eq(first_snapshot)].copy()
    if "Hire Date" in first_snapshot_df.columns:
        first_snapshot_df = first_snapshot_df[first_snapshot_df["Hire Date"].isna() | (first_snapshot_df["Hire Date"] <= period_start)]
    if term_df.empty:
        return first_snapshot_df.drop_duplicates(["WD Employee ID"]) if not first_snapshot_df.empty else first_snapshot_df

    bridge_term = term_df[
        term_df["Termination Date"].between(period_start.normalize(), first_snapshot.normalize(), inclusive="both")
    ].copy()
    if "Hire Date" in bridge_term.columns:
        bridge_term = bridge_term[bridge_term["Hire Date"].isna() | (bridge_term["Hire Date"] <= period_start)]

    combined = pd.concat([first_snapshot_df, bridge_term], ignore_index=True, sort=False)
    return combined.drop_duplicates(["WD Employee ID"]) if not combined.empty else combined


def build_attrition_row(region_label: str, country_label: str, start_df: pd.DataFrame, end_df: pd.DataFrame, term_df: pd.DataFrame, row_kind: str) -> dict:
    metrics = calculate_attrition_metrics(start_df, end_df, term_df)
    return {
        "Region": region_label,
        "Country/Territory": country_label,
        "Active Regular Start": int(metrics["start_hc"]),
        "Active Regular End": int(metrics["end_hc"]),
        "Overall No. of Attrition": int(metrics["no_attrition"]),
        "Overall Attrition %": metrics["rate"],
        "Voluntary No. of Attrition": int(metrics["voluntary_attrition"]),
        "Voluntary Attrition %": metrics["voluntary_rate"],
        "Involuntary No. of Attrition": int(metrics["involuntary_attrition"]),
        "Involuntary Attrition %": metrics["involuntary_rate"],
        "Others No. of Attrition": int(metrics["others_attrition"]),
        "Others Attrition %": metrics["others_rate"],
        "Statutory No. of Attrition": int(metrics["statutory_attrition"]),
        "Statutory Attrition %": metrics["statutory_rate"],
        "row_kind": row_kind,
    }


def build_dashboard4(active_df: pd.DataFrame, term_df: pd.DataFrame, notes: list[str]) -> tuple[dict, pd.DataFrame]:
    active_regular = filter_overseas(filter_regular(active_df)).copy()
    active_regular = active_regular[active_regular["Country/Territory"].ne("Mainland China")]
    if active_regular.empty:
        placeholder = pd.DataFrame([
            {
                "Region": "N/A",
                "Country/Territory": "缺少海外正式员工在职数据",
                "Active Regular Start": 0,
                "Active Regular End": 0,
                "Overall No. of Attrition": 0,
                "Overall Attrition %": 0.0,
                "Voluntary No. of Attrition": 0,
                "Voluntary Attrition %": 0.0,
                "Involuntary No. of Attrition": 0,
                "Involuntary Attrition %": 0.0,
                "Others No. of Attrition": 0,
                "Others Attrition %": 0.0,
                "Statutory No. of Attrition": 0,
                "Statutory Attrition %": 0.0,
                "row_kind": "missing",
            }
        ])
        notes.append("看板4缺少海外正式员工在职数据")
        return {"period_start": None, "period_end": None, "header_start": None, "header_end": None, "incomplete": True}, placeholder

    end_snapshot = active_regular["snapshot_date"].max().normalize()
    report_year = end_snapshot.year
    header_start = pd.Timestamp(year=report_year, month=1, day=1)
    header_end = end_snapshot

    if term_df.empty:
        placeholder = pd.DataFrame([
            {
                "Region": "N/A",
                "Country/Territory": "缺少离职明细文件",
                "Active Regular Start": 0,
                "Active Regular End": 0,
                "Overall No. of Attrition": 0,
                "Overall Attrition %": 0.0,
                "Voluntary No. of Attrition": 0,
                "Voluntary Attrition %": 0.0,
                "Involuntary No. of Attrition": 0,
                "Involuntary Attrition %": 0.0,
                "Others No. of Attrition": 0,
                "Others Attrition %": 0.0,
                "Statutory No. of Attrition": 0,
                "Statutory Attrition %": 0.0,
                "row_kind": "missing",
            }
        ])
        notes.append("看板4缺少离职明细文件")
        return {
            "period_start": header_start.strftime("%Y-%m-%d"),
            "period_end": header_end.strftime("%Y-%m-%d"),
            "header_start": header_start.strftime("%Y-%m-%d"),
            "header_end": header_end.strftime("%Y-%m-%d"),
            "incomplete": True,
        }, placeholder

    term_regular = filter_overseas(filter_regular(term_df)).copy()
    term_regular = term_regular[term_regular["Country/Territory"].ne("Mainland China")]
    term_regular["Term Reason Std"] = term_regular["Termination Category"].str.lower().map(TERM_CATEGORY_MAP).fillna("Others")

    year_snapshots = sorted({date.normalize() for date in active_regular["snapshot_date"].dropna().tolist() if date.year == report_year})
    first_snapshot = year_snapshots[0] if year_snapshots else header_end
    start_df = rebuild_year_start_population(active_regular, term_regular, header_start, first_snapshot)
    end_df = active_regular[active_regular["snapshot_date"].eq(header_end)].copy()
    period_term = term_regular[
        term_regular["Termination Date"].between(header_start.normalize(), header_end.normalize(), inclusive="both")
    ].copy()

    if first_snapshot.normalize() != header_start.normalize():
        notes.append(
            f"看板4期初人数按 {first_snapshot.strftime('%Y-%m-%d')} 首个可用快照结合年初至该日离职记录回推至 {header_start.strftime('%Y-%m-%d')}"
        )

    rows: list[dict] = []
    for region in REGION_ORDER:
        region_start = start_df[start_df["Region_std"].eq(region)]
        region_end = end_df[end_df["Region_std"].eq(region)]
        region_term = period_term[period_term["Region_std"].eq(region)]
        countries = sorted(
            set(region_start["Country/Territory"].tolist())
            | set(region_end["Country/Territory"].tolist())
            | set(region_term["Country/Territory"].tolist())
        )
        if not countries:
            continue
        country_order = region_end.groupby("Country/Territory")["WD Employee ID"].nunique().sort_values(ascending=False).index.tolist()
        countries = sorted(countries, key=lambda c: (country_order.index(c) if c in country_order else 9999, c))
        for country in countries:
            rows.append(
                build_attrition_row(
                    region,
                    country,
                    region_start[region_start["Country/Territory"].eq(country)],
                    region_end[region_end["Country/Territory"].eq(country)],
                    region_term[region_term["Country/Territory"].eq(country)],
                    region,
                )
            )

    total_row = build_attrition_row(
        "海外合计",
        "",
        start_df,
        end_df,
        period_term,
        "grand",
    )
    rows.append(total_row)
    return {
        "period_start": header_start.strftime("%Y-%m-%d"),
        "period_end": header_end.strftime("%Y-%m-%d"),
        "header_start": header_start.strftime("%Y-%m-%d"),
        "header_end": header_end.strftime("%Y-%m-%d"),
        "incomplete": False,
    }, pd.DataFrame(rows)


def _contingent_pair_counts(df: pd.DataFrame, bg: str | None = None) -> tuple[int, int]:
    """Count unique Contractor and Partner workers, optionally filtered by BG."""
    subset = df if bg is None else df[df["BG_mapped"].eq(bg)]
    contractor = subset[subset["Worker_Type"].eq("Contractor")]["Worker_ID"].nunique() if not subset.empty else 0
    partner = subset[subset["Worker_Type"].eq("Partner")]["Worker_ID"].nunique() if not subset.empty else 0
    return contractor, partner


def _contingent_summary_row(region_label: str, country_label: str, subset: pd.DataFrame, row_kind: str) -> dict:
    """Build one row for the dashboard5 table."""
    total_c, total_p = _contingent_pair_counts(subset)
    row: dict = {
        "Region": region_label,
        "Country/Territory": country_label,
        "Total C": total_c,
        "Total P": total_p,
    }
    for bg in BG_DISPLAY_ORDER[1:]:
        c, p = _contingent_pair_counts(subset, bg)
        row[f"{bg} C"] = c
        row[f"{bg} P"] = p
    row["row_kind"] = row_kind
    return row


def build_dashboard5(contingent_df: pd.DataFrame, notes: list[str]) -> tuple[dict, pd.DataFrame]:
    """Build Dashboard 5 – Active Contractor & Partner breakdown."""
    columns = (
        ["Region", "Country/Territory", "Total C", "Total P"]
        + [f"{bg} {suffix}" for bg in BG_DISPLAY_ORDER[1:] for suffix in ["C", "P"]]
        + ["row_kind"]
    )
    if contingent_df.empty:
        return {"has_data": False}, pd.DataFrame(columns=columns)

    overseas_df = contingent_df[contingent_df["Region_std"].isin(REGION_ORDER)].copy()

    rows: list[dict] = []
    for region in REGION_ORDER:
        region_df = overseas_df[overseas_df["Region_std"].eq(region)].copy()
        if region_df.empty:
            continue
        country_order = (
            region_df.groupby("Country_EN")["Worker_ID"]
            .nunique()
            .sort_values(ascending=False)
            .index.tolist()
        )
        for country in country_order:
            country_df = region_df[region_df["Country_EN"].eq(country)]
            rows.append(_contingent_summary_row(region, country, country_df, region))

    rows.append(_contingent_summary_row("海外合计", "", overseas_df, "grand"))

    # Build the bottom total line (Contractor / Partner per BG)
    total_c, total_p = _contingent_pair_counts(overseas_df)
    bottom_row: dict = {
        "Region": "",
        "Country/Territory": "",
        "Total C": total_c,
        "Total P": total_p,
        "row_kind": "grand",
    }
    for bg in BG_DISPLAY_ORDER[1:]:
        c, p = _contingent_pair_counts(overseas_df, bg)
        bottom_row[f"{bg} C"] = c
        bottom_row[f"{bg} P"] = p
    # The second grand-total row shows column sums (same as 海外合计 since no GC)
    # Already captured by 海外合计 row, skip duplicate

    return {"has_data": True}, pd.DataFrame(rows)


def format_percentage(value: object) -> str:
    try:
        if value is None or pd.isna(value):
            return ""
    except Exception:
        pass
    return f"{float(value):.1f}%"


def format_table_value(column_name: str, value: object) -> str:
    try:
        if value is None or pd.isna(value):
            return ""
    except Exception:
        pass
    if column_name == "Country/Territory":
        return textwrap.fill(clean_text(value), width=18, break_long_words=False)
    if "Attrition %" in column_name:
        return format_percentage(value)
    if isinstance(value, float):
        return str(int(value)) if value.is_integer() else f"{value:.1f}"
    return clean_text(value)


def add_inside_bar_labels(ax, bars) -> None:
    for bar in bars:
        height = bar.get_height()
        if height <= 0:
            continue
        y = height * 0.52 if height >= 40 else height + max(1, height * 0.05)
        va = "center" if height >= 40 else "bottom"
        color = "white" if height >= 40 else "#4d4d4d"
        ax.text(
            bar.get_x() + bar.get_width() / 2,
            y,
            f"{int(height)}",
            ha="center",
            va=va,
            fontsize=10.5,
            fontweight="bold",
            color=color,
        )


def add_title_badge(ax, text: str) -> None:
    ax.text(
        0.5,
        1.07,
        text,
        transform=ax.transAxes,
        ha="center",
        va="center",
        fontsize=14,
        fontweight="bold",
        color="#4f4f4f",
        bbox={"boxstyle": "square,pad=0.28", "facecolor": "#d9d9d9", "edgecolor": "none"},
    )


def style_chart_axis(ax) -> None:
    ax.set_facecolor("#f3f3f3")
    ax.spines[["top", "left", "right"]].set_visible(False)
    ax.spines["bottom"].set_color("#b9b9b9")
    ax.spines["bottom"].set_linewidth(1.1)
    ax.tick_params(axis="y", left=False, labelleft=False)
    ax.tick_params(axis="x", length=0, colors="#5a5a5a", labelsize=10.5)
    ax.grid(False)


def plot_grouped_bars(ax, data: pd.DataFrame, title: str, highlight_last: bool):
    style_chart_axis(ax)
    add_title_badge(ax, title)
    if data.empty:
        ax.text(0.5, 0.46, "数据缺失", ha="center", va="center", fontsize=13, color="#7a7a7a")
        ax.set_xticks([])
        return []

    months = data["month"].tolist()
    x = list(range(len(months)))
    width = 0.30
    handles = []
    max_height = 0
    for idx, region in enumerate(REGION_ORDER):
        heights = data[region].tolist()
        max_height = max(max_height, max(heights) if heights else 0)
        offset = (idx - 1) * width
        bars = ax.bar(
            [pos + offset for pos in x],
            heights,
            width=width,
            label=region,
            color=REGION_COLORS[region],
            edgecolor="none",
        )
        handles.append(bars[0])
        add_inside_bar_labels(ax, bars)
    ax.set_xticks(x)
    ax.set_xticklabels(months)
    ax.set_ylim(0, max(max_height * 1.16, 10))
    ax.margins(x=0.04)
    if highlight_last and months:
        last_index = len(months) - 1
        ymax = ax.get_ylim()[1]
        rect = Rectangle(
            (last_index - 0.5, 0),
            1.0,
            ymax * 0.96,
            fill=False,
            linestyle=(0, (3, 2)),
            linewidth=1.0,
            edgecolor="#ababab",
        )
        ax.add_patch(rect)
    return handles


def figure_to_data_uri(fig) -> str:
    buffer = BytesIO()
    fig.savefig(buffer, dpi=180, bbox_inches="tight", facecolor=fig.get_facecolor())
    buffer.seek(0)
    encoded = base64.b64encode(buffer.read()).decode("ascii")
    return f"data:image/png;base64,{encoded}"


def build_dashboard1_figure(data: pd.DataFrame, meta: dict):
    fig, axes = plt.subplots(1, 2, figsize=(13.8, 4.2), gridspec_kw={"width_ratios": [1.0, 1.45]})
    fig.patch.set_facecolor("#f3f3f3")
    fig.subplots_adjust(top=0.82, left=0.035, right=0.98, bottom=0.16, wspace=0.16)

    yoy = data[data["section"].eq("同比")].copy()
    mom = data[data["section"].eq("环比")].copy()
    handles = plot_grouped_bars(axes[0], yoy, "① 同比去年同期正式员工人数", True)
    mom_count = meta.get("mom_count", 3)
    mom_title = f"② 环比近{mom_count}个月正式员工人数" if mom_count >= 2 else "② 环比正式员工人数（数据不足）"
    plot_grouped_bars(axes[1], mom, mom_title, True)

    if handles:
        fig.legend(
            handles,
            REGION_ORDER,
            loc="upper center",
            bbox_to_anchor=(0.52, 0.985),
            ncol=3,
            frameon=False,
            fontsize=11,
            handlelength=0.8,
            handletextpad=0.35,
        )
    fig.text(0.975, 0.985, f"最新月份: {meta['latest_month'] or 'N/A'}", ha="right", va="top", fontsize=10.5, color="#4f4f4f", fontweight="bold")
    if meta.get("notes"):
        fig.text(0.035, 0.03, "；".join(meta["notes"]), ha="left", va="bottom", fontsize=9.3, color="#b76e00")
    return fig


def render_dashboard1_png(data: pd.DataFrame, meta: dict, output_path: Path) -> None:
    fig = build_dashboard1_figure(data, meta)
    fig.savefig(output_path, dpi=180, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


def render_dashboard1_image(data: pd.DataFrame, meta: dict) -> str:
    fig = build_dashboard1_figure(data, meta)
    uri = figure_to_data_uri(fig)
    plt.close(fig)
    return uri


def add_panel_title_badge(ax, text: str) -> None:
    ax.text(
        0.03,
        0.97,
        text,
        transform=ax.transAxes,
        ha="left",
        va="top",
        fontsize=12.5,
        fontweight="bold",
        color="#4f4f4f",
        bbox={"boxstyle": "square,pad=0.24", "facecolor": "#d9d9d9", "edgecolor": "none"},
        zorder=5,
    )


def plot_pie_panel(ax, data: pd.DataFrame, title: str, background: str, title_parts: list[tuple[str, str]] | None = None) -> None:
    from matplotlib.offsetbox import AnchoredOffsetbox, HPacker, TextArea
    ax.set_facecolor(background)
    if title_parts:
        text_areas = []
        for seg_text, seg_color in title_parts:
            ta = TextArea(
                seg_text,
                textprops=dict(fontsize=12.5, fontweight="bold", color=seg_color),
            )
            text_areas.append(ta)
        hbox = HPacker(children=text_areas, align="baseline", pad=0, sep=0)
        anchored = AnchoredOffsetbox(
            loc="upper left",
            child=hbox,
            pad=0.3,
            borderpad=0.4,
            bbox_to_anchor=(0.03, 0.99),
            bbox_transform=ax.transAxes,
            frameon=True,
        )
        anchored.patch.set_boxstyle("square,pad=0.24")
        anchored.patch.set_facecolor("#d9d9d9")
        anchored.patch.set_edgecolor("none")
        anchored.zorder = 5
        ax.add_artist(anchored)
    else:
        add_panel_title_badge(ax, title)
    values = data["Headcount"].tolist() if not data.empty else []
    if not values or sum(values) == 0:
        ax.text(0.5, 0.48, "数据缺失", ha="center", va="center", fontsize=14, color="#7a7a7a")
        ax.set_xticks([])
        ax.set_yticks([])
        for spine in ax.spines.values():
            spine.set_visible(False)
        return

    wedges, _ = ax.pie(
        values,
        colors=[REGION_COLORS[row.Region] for row in data.itertuples()],
        startangle=90,
        counterclock=False,
        wedgeprops={"linewidth": 1.6, "edgecolor": "white"},
        radius=0.9,
    )
    for wedge, row in zip(wedges, data.itertuples()):
        angle = math.radians((wedge.theta1 + wedge.theta2) / 2)
        radius = 0.53
        share_label = int(round(row.Share))
        ax.text(
            radius * math.cos(angle),
            radius * math.sin(angle),
            f"{row.Region}\n{int(row.Headcount)}人 {share_label}%",
            ha="center",
            va="center",
            fontsize=12,
            fontweight="bold",
            color="white",
        )
    ax.set_aspect("equal")
    ax.set_xticks([])
    ax.set_yticks([])
    for spine in ax.spines.values():
        spine.set_visible(False)


def build_dashboard2_figure(data: pd.DataFrame, meta: dict):
    fig, axes = plt.subplots(1, 2, figsize=(12.2, 5.3))
    fig.patch.set_facecolor("white")
    fig.subplots_adjust(left=0.03, right=0.97, top=0.96, bottom=0.06, wspace=0.02)

    active_df = data[data["Metric"].eq("Active")].copy()
    attrition_df = data[data["Metric"].eq("AttritionYTD")].copy()
    active_title = f"③ {compact_date_label(meta.get('active_snapshot_date'))} 正式员工在职人数"
    active_title_parts = [
        (f"③ {compact_date_label(meta.get('active_snapshot_date'))} 正式员工", "#4f4f4f"),
        ("在职", "#1565C0"),
        ("人数", "#4f4f4f"),
    ]
    attrition_title = f"③ {meta.get('attrition_year') or ''}截止{month_day_label(meta.get('attrition_cutoff_date'))}正式员工离职人数"
    attrition_title_parts = [
        (f"③ {meta.get('attrition_year') or ''}截止{month_day_label(meta.get('attrition_cutoff_date'))}正式员工", "#4f4f4f"),
        ("离职", "#C62828"),
        ("人数", "#4f4f4f"),
    ]
    plot_pie_panel(
        axes[0],
        active_df,
        active_title,
        "#dfe8f1",
        title_parts=active_title_parts,
    )
    plot_pie_panel(
        axes[1],
        attrition_df,
        attrition_title,
        "#eadccf",
        title_parts=attrition_title_parts,
    )

    for ax in axes:
        ax.add_patch(Rectangle((0.01, 0.01), 0.98, 0.98, transform=ax.transAxes, fill=False, edgecolor="#cfd3d8", linewidth=1.0))
    return fig


def render_dashboard2_png(data: pd.DataFrame, meta: dict, output_path: Path) -> None:
    fig = build_dashboard2_figure(data, meta)
    fig.savefig(output_path, dpi=180, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


def render_dashboard2_image(data: pd.DataFrame, meta: dict) -> str:
    fig = build_dashboard2_figure(data, meta)
    uri = figure_to_data_uri(fig)
    plt.close(fig)
    return uri


def build_grouped_header_rows(columns: list[str]) -> tuple[list[str], list[str]]:
    top_row: list[str] = []
    second_row: list[str] = []
    previous_group = None
    for column in columns:
        if column in {"Region", "Country/Territory"}:
            top_row.append(column)
            second_row.append("")
            previous_group = None
            continue
        group, suffix = column.rsplit(" ", 1)
        display_group = "合计" if group == "Total" else group
        top_row.append(display_group if display_group != previous_group else "")
        second_row.append("Regular" if suffix == "R" else "Intern")
        previous_group = display_group
    return top_row, second_row


def apply_table_widths(table, total_rows: int, col_widths: list[float]) -> None:
    total = sum(col_widths)
    normalized = [width / total for width in col_widths]
    for col_idx, width in enumerate(normalized):
        for row_idx in range(total_rows):
            if (row_idx, col_idx) in table._cells:
                table[(row_idx, col_idx)].set_width(width)


def build_region_merge_groups(display_df: pd.DataFrame) -> tuple[pd.DataFrame, list[tuple[int, int]]]:
    if display_df.empty or "Region" not in display_df.columns:
        return display_df, []
    display_df = display_df.copy()
    groups: list[tuple[int, int]] = []
    regions = display_df["Region"].tolist()
    start = 0
    while start < len(regions):
        region = clean_text(regions[start])
        end = start
        while end + 1 < len(regions) and clean_text(regions[end + 1]) == region:
            end += 1
        groups.append((start, end))
        if end > start:
            mid = start + (end - start) // 2
            for idx in range(start, end + 1):
                if idx != mid:
                    display_df.at[display_df.index[idx], "Region"] = ""
        start = end + 1
    return display_df, groups


def apply_region_merge_edges(table, groups: list[tuple[int, int]], header_rows: int, region_col_idx: int = 0) -> None:
    for start, end in groups:
        if start == end:
            continue
        for offset, row_idx in enumerate(range(start + header_rows, end + header_rows + 1)):
            cell = table[(row_idx, region_col_idx)]
            if offset == 0:
                cell.visible_edges = "TLR"
            elif row_idx == end + header_rows:
                cell.visible_edges = "BLR"
            else:
                cell.visible_edges = "LR"
            cell.set_edgecolor("#505050")
            cell.set_linewidth(0.8)


def render_table_png(
    df: pd.DataFrame,
    title: str,
    output_path: Path,
    min_width: float = 16,
    row_height: float = 0.35,
    grouped_header: bool = False,
    col_widths: list[float] | None = None,
    font_size: float = 8.0,
    merge_region_cells: bool = False,
) -> None:
    display_df = df.copy()
    row_kinds = display_df.pop("row_kind") if "row_kind" in display_df.columns else pd.Series(["data"] * len(display_df))
    region_groups: list[tuple[int, int]] = []
    if merge_region_cells:
        display_df, region_groups = build_region_merge_groups(display_df)
    columns = display_df.columns.tolist()
    body_rows = [[format_table_value(column, row[column]) for column in columns] for _, row in display_df.iterrows()]
    header_rows = 2 if grouped_header else 1
    fig_width = max(min_width, len(columns) * 0.95)
    fig_height = max(4.4, 0.95 + (len(display_df) + header_rows) * row_height)

    fig, ax = plt.subplots(figsize=(fig_width, fig_height))
    fig.patch.set_facecolor("#f3f3f3")
    ax.set_facecolor("#f3f3f3")
    ax.axis("off")

    if grouped_header:
        top_header, second_header = build_grouped_header_rows(columns)
        cell_text = [top_header, second_header] + body_rows
        table = ax.table(cellText=cell_text, cellLoc="center", bbox=[0.0, 0.0, 1.0, 0.92])
    else:
        table = ax.table(cellText=body_rows, colLabels=columns, cellLoc="center", bbox=[0.0, 0.0, 1.0, 0.92])

    table.auto_set_font_size(False)
    table.set_fontsize(font_size)
    table.scale(1, 1.15 if grouped_header else 1.08)

    total_rows = header_rows + len(display_df)
    if col_widths:
        apply_table_widths(table, total_rows, col_widths)

    for row_idx in range(header_rows):
        for col_idx in range(len(columns)):
            cell = table[(row_idx, col_idx)]
            cell.set_facecolor("#B7D7E8")
            cell.set_text_props(weight="bold", color="#243746")
            cell.set_edgecolor("#505050")
            cell.set_linewidth(0.8)

    for row_idx, kind in enumerate(row_kinds.tolist(), start=header_rows):
        fill = ROW_FILL.get(kind, ROW_FILL.get(clean_text(df.iloc[row_idx - header_rows, 0]), "FFFFFF"))
        for col_idx in range(len(columns)):
            cell = table[(row_idx, col_idx)]
            cell.set_facecolor(f"#{fill}")
            cell.set_edgecolor("#505050")
            cell.set_linewidth(0.8)
            cell.set_text_props(color="#333333")

    if merge_region_cells and region_groups:
        apply_region_merge_edges(table, region_groups, header_rows)

    ax.text(0.5, 0.985, title, transform=ax.transAxes, ha="center", va="top", fontsize=15.5, fontweight="bold", color="#4f4f4f")
    fig.add_artist(Rectangle((0.004, 0.02), 0.992, 0.96, transform=fig.transFigure, fill=False, edgecolor="#d0d0d0", linewidth=1.1))
    fig.savefig(output_path, dpi=180, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


def load_css() -> str:
    css_path = Path(__file__).resolve().parent.parent / "assets" / "dashboard_style.css"
    return css_path.read_text(encoding="utf-8") if css_path.exists() else ""


def format_html_value(column_name: str, value: object) -> str:
    try:
        if value is None or pd.isna(value):
            return ""
    except Exception:
        pass
    if "Attrition %" in column_name:
        return format_percentage(value)
    if isinstance(value, float):
        return str(int(value)) if value.is_integer() else f"{value:.1f}"
    return clean_text(value)


def row_fill_hex(kind: object) -> str:
    normalized = clean_text(kind)
    fill = ROW_FILL.get(normalized, ROW_FILL.get(normalized.strip(), "FFFFFF"))
    return f"#{fill}"


def build_region_rowspan_map(display_df: pd.DataFrame) -> dict[int, int]:
    if display_df.empty or "Region" not in display_df.columns:
        return {}
    spans: dict[int, int] = {}
    regions = display_df["Region"].tolist()
    start = 0
    while start < len(regions):
        region = clean_text(regions[start])
        end = start
        while end + 1 < len(regions) and clean_text(regions[end + 1]) == region:
            end += 1
        spans[start] = end - start + 1
        start = end + 1
    return spans


def build_region_cell_html(region_label: str, annotation_lines: list[str] | None = None) -> str:
    label_html = f"<div class='region-label'>{html_escape(region_label)}</div>"
    if not annotation_lines:
        return label_html

    metric_blocks: list[str] = []
    for line in annotation_lines:
        parts = clean_text(line).rsplit(" ", 1)
        if len(parts) == 2:
            metric_label, metric_value = parts
            metric_blocks.append(
                "<div class='region-metric'>"
                f"<div class='region-meta-label'>{html_escape(metric_label)}</div>"
                f"<div class='region-meta-value'>{html_escape(metric_value)}</div>"
                "</div>"
            )
        else:
            metric_blocks.append(f"<div class='region-meta-line'>{html_escape(line)}</div>")
    return f"{label_html}<div class='region-meta'>{''.join(metric_blocks)}</div>"


def build_table_body_rows(display_df: pd.DataFrame, row_kinds: pd.Series, region_annotations: dict[str, list[str]] | None = None) -> str:
    columns = display_df.columns.tolist()
    rowspan_map = build_region_rowspan_map(display_df)
    # Labels whose Region+Country cells should be merged into a single colspan=2 cell
    MERGE_LABELS = {"海外合计", "海外合计（含试点国内）"}
    body_rows: list[str] = []

    for row_pos, (_, row) in enumerate(display_df.iterrows()):
        row_kind = row_kinds.iloc[row_pos]
        row_style = f" style='background:{row_fill_hex(row_kind)}'"
        cells: list[str] = []
        region_label = clean_text(row.get("Region", ""))
        is_merge_row = region_label in MERGE_LABELS
        for col_idx, column in enumerate(columns):
            value = format_html_value(column, row[column])
            escaped = html_escape(value)
            if column == "Region":
                if is_merge_row:
                    # Merge Region + Country/Territory into one colspan=2 cell
                    region_html = build_region_cell_html(region_label, (region_annotations or {}).get(region_label))
                    cells.append(f"<td class='region-cell' colspan='2'>{region_html}</td>")
                elif row_pos in rowspan_map:
                    rowspan = rowspan_map[row_pos]
                    region_html = build_region_cell_html(region_label, (region_annotations or {}).get(region_label))
                    cells.append(f"<td class='region-cell' rowspan='{rowspan}'>{region_html}</td>")
                continue
            if column == "Country/Territory" and is_merge_row:
                # Skip Country column for merge rows (already covered by colspan)
                continue
            cells.append(f"<td>{escaped}</td>")
        body_rows.append(f"<tr{row_style}>{''.join(cells)}</tr>")
    return ''.join(body_rows)


def build_table_html(df: pd.DataFrame, grouped_header: bool = False, region_annotations: dict[str, list[str]] | None = None) -> str:
    display_df = df.copy()
    row_kinds = display_df.pop("row_kind") if "row_kind" in display_df.columns else pd.Series(["data"] * len(display_df))
    columns = display_df.columns.tolist()
    body_html = build_table_body_rows(display_df, row_kinds, region_annotations)

    if grouped_header:
        header_cells = ["<th rowspan='2'>Region</th>", "<th rowspan='2'>Country/Territory</th>"]
        second_row: list[str] = []
        metric_columns = columns[2:]
        for idx in range(0, len(metric_columns), 2):
            group = metric_columns[idx].rsplit(" ", 1)[0]
            display_group = "合计" if group == "Total" else group
            header_cells.append(f"<th colspan='2'>{html_escape(display_group)}</th>")
            second_row.extend(["<th>Regular</th>", "<th>Intern</th>"])
        thead = f"<thead><tr>{''.join(header_cells)}</tr><tr>{''.join(second_row)}</tr></thead>"
    else:
        head_cells = "".join(f"<th>{html_escape(column)}</th>" for column in columns)
        thead = f"<thead><tr>{head_cells}</tr></thead>"

    return f"<div class='table-wrap'><table class='dashboard-table{' grouped' if grouped_header else ''}'>{thead}<tbody>{body_html}</tbody></table></div>"


def build_dashboard4_table_html(df: pd.DataFrame, meta: dict) -> str:
    display_df = df.copy()
    row_kinds = display_df.pop("row_kind") if "row_kind" in display_df.columns else pd.Series(["data"] * len(display_df))

    # Determine whether Others / Statutory columns have any non-zero data
    has_others = (
        "Others No. of Attrition" in display_df.columns
        and display_df["Others No. of Attrition"].sum() > 0
    )
    has_statutory = (
        "Statutory No. of Attrition" in display_df.columns
        and display_df["Statutory No. of Attrition"].sum() > 0
    )

    # Drop empty optional columns from display
    if not has_others:
        display_df = display_df.drop(columns=["Others No. of Attrition", "Others Attrition %"], errors="ignore")
    if not has_statutory:
        display_df = display_df.drop(columns=["Statutory No. of Attrition", "Statutory Attrition %"], errors="ignore")

    body_html = build_table_body_rows(display_df, row_kinds)
    header_start = html_escape(meta.get("header_start") or "Start")
    header_end = html_escape(meta.get("header_end") or "End")

    # Build dynamic header
    top_row_parts = [
        "<th rowspan='2'>Region</th>",
        "<th rowspan='2'>Country/Territory</th>",
        "<th colspan='2'>Active Regular</th>",
        "<th colspan='2'>Overall</th>",
        "<th colspan='2'>Voluntary</th>",
        "<th colspan='2'>Involuntary</th>",
    ]
    second_row_parts = [
        f"<th>{header_start}</th><th>{header_end}</th>",
        "<th>No.</th><th>% of Involuntary</th>",
        "<th>No.</th><th>% of Involuntary</th>",
        "<th>No.</th><th>% of Involuntary</th>",
    ]
    if has_others:
        top_row_parts.append("<th colspan='2'>Others</th>")
        second_row_parts.append("<th>No.</th><th>% of Involuntary</th>")
    if has_statutory:
        top_row_parts.append("<th colspan='2'>Statutory</th>")
        second_row_parts.append("<th>No.</th><th>% of Involuntary</th>")

    thead = (
        f"<thead><tr>{''.join(top_row_parts)}</tr>"
        f"<tr>{''.join(second_row_parts)}</tr></thead>"
    )
    return f"<div class='table-wrap'><table class='dashboard-table grouped dashboard4-table'>{thead}<tbody>{body_html}</tbody></table></div>"


def build_dashboard5_table_html(df: pd.DataFrame) -> str:
    """Build HTML table for Dashboard 5 – Active Contractor & Partner.
    Uses the same grouped-header style as Dashboard 3, but with Contractor/Partner sub-columns."""
    display_df = df.copy()
    row_kinds = display_df.pop("row_kind") if "row_kind" in display_df.columns else pd.Series(["data"] * len(display_df))
    body_html = build_table_body_rows(display_df, row_kinds)

    # Build grouped header: Region | Country | 合计(C/P) | IEG(C/P) | ... | S3(C/P)
    header_cells = ["<th rowspan='2'>Region</th>", "<th rowspan='2'>Country/Territory</th>"]
    second_row: list[str] = []
    metric_columns = display_df.columns.tolist()[2:]  # skip Region, Country
    for idx in range(0, len(metric_columns), 2):
        col_name = metric_columns[idx]
        group = col_name.rsplit(" ", 1)[0]
        display_group = "合计" if group == "Total" else group
        header_cells.append(f"<th colspan='2'>{html_escape(display_group)}</th>")
        second_row.extend(["<th>Contractor</th>", "<th>Partner</th>"])
    thead = f"<thead><tr>{''.join(header_cells)}</tr><tr>{''.join(second_row)}</tr></thead>"

    return f"<div class='table-wrap'><table class='dashboard-table grouped'>{thead}<tbody>{body_html}</tbody></table></div>"


def _es_bold(val: object) -> str:
    """Wrap a value in <strong> for Executive Summary highlight."""
    return f"<strong>{html_escape(str(val))}</strong>"


def build_executive_summary_defaults(meta: dict, d1_meta: dict, d2_meta: dict, d2_df: pd.DataFrame, d3_meta: dict, d3_df: pd.DataFrame, d4_meta: dict, d4_df: pd.DataFrame, d5_meta: dict | None = None, d5_df: pd.DataFrame | None = None) -> dict[str, str]:
    """Generate default Executive Summary HTML (unordered list) for each category.
    Returns dict with keys: active, termination, others.  Values are HTML strings."""

    # --- 在职情况 (based on 看板1 + 看板2-active + 看板3) ---
    active_items: list[str] = []

    active_data = d2_df[d2_df["Metric"].eq("Active")] if not d2_df.empty else pd.DataFrame()
    if not active_data.empty:
        total_active = int(active_data["Headcount"].sum())
        snapshot_label = compact_date_label(d2_meta.get("active_snapshot_date"))
        region_parts = []
        for _, row in active_data.iterrows():
            r_name = html_escape(row['Region'])
            r_hc = _es_bold(str(int(row['Headcount'])) + '人')
            r_share = _es_bold(str(int(round(row['Share']))) + '%')
            region_parts.append(f"{r_name} {r_hc}({r_share})")
        active_items.append(f"截至{_es_bold(snapshot_label)}，海外正式员工在职共{_es_bold(str(total_active) + '人')}，其中{'、'.join(region_parts)}。")

    if d1_meta.get("yoy_months") and len(d1_meta["yoy_months"]) >= 2:
        active_items.append(f"同比（{_es_bold(d1_meta['yoy_months'][0])} vs {_es_bold(d1_meta['yoy_months'][-1])}）正式员工人数变化趋势见看板1。")
    if d1_meta.get("mom_months") and len(d1_meta["mom_months"]) >= 2:
        mom_count = d1_meta.get('mom_count', len(d1_meta['mom_months']))
        active_items.append(f"环比近{_es_bold(str(mom_count) + '个月')}趋势见看板1。")

    if not d3_df.empty and "row_kind" in d3_df.columns:
        grand_rows = d3_df[d3_df["row_kind"].eq("grand")]
        if not grand_rows.empty:
            first_grand = grand_rows.iloc[0]
            total_r = int(first_grand.get("Total R", 0))
            total_i = int(first_grand.get("Total I", 0))
            active_items.append(f"看板3明细汇总：海外合计正式员工{_es_bold(str(total_r) + '人')}、实习生{_es_bold(str(total_i) + '人')}。")

    # --- 离职情况 (based on 看板2-attrition + 看板4) ---
    term_items: list[str] = []

    attrition_data = d2_df[d2_df["Metric"].eq("AttritionYTD")] if not d2_df.empty else pd.DataFrame()
    if not attrition_data.empty:
        total_attrition = int(attrition_data["Headcount"].sum())
        attrition_year = d2_meta.get("attrition_year", "")
        cutoff_label = month_day_label(d2_meta.get("attrition_cutoff_date"))
        if total_attrition > 0:
            region_parts = []
            for _, row in attrition_data.iterrows():
                r_name = html_escape(row['Region'])
                r_hc = _es_bold(str(int(row['Headcount'])) + '人')
                r_share = _es_bold(str(int(round(row['Share']))) + '%')
                region_parts.append(f"{r_name} {r_hc}({r_share})")
            term_items.append(f"{_es_bold(str(attrition_year) + '年')}截止{_es_bold(cutoff_label)}，海外正式员工离职共{_es_bold(str(total_attrition) + '人')}，其中{'、'.join(region_parts)}。")

    if not d4_df.empty and "row_kind" in d4_df.columns:
        grand_rows = d4_df[d4_df["row_kind"].eq("grand")]
        if not grand_rows.empty:
            grand = grand_rows.iloc[0]
            overall_rate = grand.get("Overall Attrition %", 0)
            vol_rate = grand.get("Voluntary Attrition %", 0)
            invol_rate = grand.get("Involuntary Attrition %", 0)
            others_rate = grand.get("Others Attrition %", 0)
            statutory_rate = grand.get("Statutory Attrition %", 0)
            period_start = html_escape(d4_meta.get('header_start', ''))
            period_end = html_escape(d4_meta.get('header_end', ''))
            period_label = period_start + '至' + period_end
            rate_parts = f"整体离职率{_es_bold(format_percentage(overall_rate))}，主动离职率{_es_bold(format_percentage(vol_rate))}，被动离职率{_es_bold(format_percentage(invol_rate))}"
            if others_rate and float(others_rate) > 0:
                rate_parts += f"，其他离职率{_es_bold(format_percentage(others_rate))}"
            if statutory_rate and float(statutory_rate) > 0:
                rate_parts += f"，法定离职率{_es_bold(format_percentage(statutory_rate))}"
            term_items.append(f"看板4离职分析（{_es_bold(period_label)}）：{rate_parts}。")

    # --- 其他人员 (based on 看板5) ---
    others_items: list[str] = []

    if d5_df is not None and not d5_df.empty and d5_meta and d5_meta.get("has_data"):
        grand_rows = d5_df[d5_df["row_kind"].eq("grand")] if "row_kind" in d5_df.columns else pd.DataFrame()
        if not grand_rows.empty:
            grand = grand_rows.iloc[0]
            total_c = int(grand.get("Total C", 0))
            total_p = int(grand.get("Total P", 0))
            others_items.append(f"海外外包人员（Contractor）共{_es_bold(str(total_c) + '人')}，合作伙伴（Partner）共{_es_bold(str(total_p) + '人')}。详见看板5。")
    if not others_items:
        others_items.append("暂无外包/合作伙伴数据。")

    def _to_inline(items: list[str]) -> str:
        if not items:
            return "暂无数据。"
        return "".join(items)

    return {
        "active": _to_inline(active_items),
        "termination": _to_inline(term_items),
        "others": _to_inline(others_items),
    }


def build_executive_summary_html(summary: dict[str, str]) -> str:
    """Build the editable Executive Summary HTML section.
    Each category (在职/离职/其他) is rendered as a single inline paragraph
    with the title prefix in bold, all within one unified block."""
    # summary values are HTML (ul/li with <strong>); convert to inline text
    return f"""<section class='section executive-summary' id='executiveSummary'>
  <h2>Executive Summary
    <button type='button' class='es-btn es-edit-btn' id='esEditBtn' title='编辑'>✏️ 编辑</button>
    <button type='button' class='es-btn es-save-btn' id='esSaveBtn' title='保存' style='display:none'>💾 保存</button>
    <button type='button' class='es-btn es-cancel-btn' id='esCancelBtn' title='取消' style='display:none'>❌ 取消</button>
    <button type='button' class='es-btn es-copy-btn' id='esCopyBtn' title='复制'>📋 复制</button>
    <button type='button' class='es-btn es-mail-btn' id='esMailBtn' title='生成邮件' style='margin-left: 8px;'>✉️ 生成邮件</button>
  </h2>
  <div class='es-block' id='esGrid'>
    <div class='es-row' id='esActive' data-field='active'>
      <span class='es-row-title'>【在职情况】</span><span class='es-row-body'>{summary["active"]}</span>
    </div>
    <div class='es-row' id='esTermination' data-field='termination'>
      <span class='es-row-title'>【离职情况】</span><span class='es-row-body'>{summary["termination"]}</span>
    </div>
    <div class='es-row' id='esOthers' data-field='others'>
      <span class='es-row-title'>【其他人员】</span><span class='es-row-body'>{summary["others"]}</span>
    </div>
  </div>
</section>"""


def build_footnotes_html(meta: dict, d4_meta: dict) -> str:
    """Build the footnotes / disclaimer section at the bottom of the dashboard.
    Date variables are filled from metadata."""
    # Determine dates for the footnote
    snapshot_date = meta.get("period_end") or "N/A"
    period_start = d4_meta.get("header_start") or meta.get("period_start") or "N/A"
    period_end = d4_meta.get("header_end") or meta.get("period_end") or "N/A"

    # Each line: (label_before_colon, rest_after_colon)
    # We build HTML directly to support bold labels, red highlights, and clean dates.
    lines_html = ""

    def _fn_line(label: str, body: str) -> str:
        return f"<p>▲ <strong>{html_escape(label)}</strong>：{body}</p>\n"

    lines_html += _fn_line("中英对照", html_escape("Regular = 正式员工; Intern = 实习生; Contractor = 外包; Partner = 合作伙伴"))
    lines_html += _fn_line("数据定义", html_escape("正式员工、实习生均为腾讯雇员，即自有员工；外包、合作伙伴均通过第三方雇佣，即非自有员工。"))
    lines_html += _fn_line("数据范围", html_escape("正式员工、实习生人数统计包括海外各区域与国内NHS试点团队；外包、合作伙伴人数统计仅包括海外各区域，区域范围外国家未纳入统计。"))
    lines_html += _fn_line("数据源", html_escape("正式员工、实习生数据统计源自Workday；外包、合作伙伴数据统计源自总部系统，外包统计字段采用员工子类型为人力服务外包、人力服务外包驻场、项目外包。合作伙伴定义未明。"))
    # 数据健康度: "外包、合作伙伴人数统计不完全准确" 标红
    health_body = (
        html_escape("正式员工、实习生人数统计运作正常；")
        + "<span class='fn-red'>外包、合作伙伴人数统计不完全准确</span>"
        + html_escape("，如需应用，请联系当地第三方服务提供商接口人。")
    )
    lines_html += _fn_line("数据健康度", health_body)
    lines_html += _fn_line("离职率计算", html_escape("统计周期离职人数 / 期初与期末在职人数均值 x 100%；Mainland China 数据不在统计范围内。"))
    # 数据统计日期: no curly braces
    date_body = f"{html_escape(snapshot_date)} 在职统计；{html_escape(period_start)}至{html_escape(snapshot_date)}离职统计。"
    lines_html += _fn_line("数据统计日期", date_body)

    return f"""<section class='section footnotes-section' id='footnotesSection'>
  <h2>看板备注<button type='button' class='copy-table-btn' id='copyFootnotesBtn' title='复制备注'>📋 复制</button></h2>
  <div class='footnotes-content' id='footnotesContent'>
    {lines_html}
  </div>
</section>"""


def build_chart_card_html(title: str, image_src: str) -> str:
    safe_title = html_escape(title)
    return (
        f"<section class='section'><h2>{safe_title}</h2>"
        f"<button type='button' class='zoomable-card' data-modal-title='{safe_title}'>"
        f"<img class='card-image' src='{image_src}' alt='{safe_title}'>"
        f"<span class='zoom-hint'>点击查看大图</span></button></section>"
    )


def _build_echarts_dashboard1_html(d1_df: pd.DataFrame, d1_meta: dict) -> str:
    """Build ECharts interactive grouped bar charts for Dashboard 1."""
    import json as _json

    yoy_df = d1_df[d1_df["section"].eq("同比")] if not d1_df.empty else pd.DataFrame()
    mom_df = d1_df[d1_df["section"].eq("环比")] if not d1_df.empty else pd.DataFrame()

    yoy_months = yoy_df["month"].tolist() if not yoy_df.empty else []
    mom_months = mom_df["month"].tolist() if not mom_df.empty else []

    def region_data(df: pd.DataFrame, region: str) -> list:
        return df[region].tolist() if not df.empty and region in df.columns else []

    yoy_data = {r: region_data(yoy_df, r) for r in REGION_ORDER}
    mom_data = {r: region_data(mom_df, r) for r in REGION_ORDER}

    mom_count = d1_meta.get("mom_count", 3)
    mom_title = f"② 环比近{mom_count}个月正式员工人数" if mom_count >= 2 else "② 环比正式员工人数（数据不足）"
    latest_month = d1_meta.get("latest_month") or "N/A"
    notes_text = "；".join(d1_meta.get("notes", []))

    region_colors = _json.dumps([REGION_COLORS[r] for r in REGION_ORDER])
    region_names = _json.dumps(REGION_ORDER)

    return f"""<section class='section'>
  <h2>看板1：正式员工趋势图</h2>
  <div class='charts-inline-row' style='display:flex;gap:8px;flex-wrap:nowrap;'>
    <div id='echartsD1Yoy' style='flex:1;min-width:0;height:320px;'></div>
    <div id='echartsD1Mom' style='flex:1.45;min-width:0;height:320px;'></div>
  </div>
</section>
<script>
(function() {{
  var regionNames = {region_names};
  var regionColors = {region_colors};
  var yoyMonths = {_json.dumps(yoy_months)};
  var momMonths = {_json.dumps(mom_months)};
  var yoyData = {_json.dumps({r: yoy_data[r] for r in REGION_ORDER})};
  var momData = {_json.dumps({r: mom_data[r] for r in REGION_ORDER})};
  var latestMonth = {_json.dumps(latest_month)};
  var notesText = {_json.dumps(notes_text)};

  function makeBarOption(title, months, data, highlightLast) {{
    var series = regionNames.map(function(r, i) {{
      return {{
        name: r,
        type: 'bar',
        barGap: '8%',
        data: data[r] || [],
        itemStyle: {{ color: regionColors[i] }},
        label: {{
          show: true,
          position: 'top',
          fontWeight: 'bold',
          fontSize: 11,
          color: '#4d4d4d',
          formatter: function(p) {{ return p.value > 0 ? p.value : ''; }}
        }}
      }};
    }});
    var markAreaData = [];
    if (highlightLast && months.length > 0) {{
      var lastMonth = months[months.length - 1];
      markAreaData = [[ {{ xAxis: months.length - 1.5 }}, {{ xAxis: months.length - 0.5 }} ]];
    }}
    if (markAreaData.length > 0) {{
      series[0].markArea = {{
        silent: true,
        itemStyle: {{ color: 'rgba(180,180,180,0.08)', borderColor: '#ababab', borderWidth: 1, borderType: 'dashed' }},
        data: markAreaData
      }};
    }}
    return {{
      title: {{
        text: title,
        left: 'center',
        top: 8,
        textStyle: {{ fontSize: 14, fontWeight: 'bold', color: '#4f4f4f',
          backgroundColor: '#d9d9d9', padding: [6, 12], borderRadius: 0 }}
      }},
      legend: {{
        data: regionNames,
        top: 36,
        left: 'center',
        textStyle: {{ fontSize: 11 }},
        itemWidth: 14,
        itemHeight: 10,
        itemGap: 16
      }},
      tooltip: {{
        trigger: 'axis',
        axisPointer: {{ type: 'shadow' }}
      }},
      grid: {{ left: 40, right: 20, top: 72, bottom: 36 }},
      xAxis: {{
        type: 'category',
        data: months,
        axisLine: {{ lineStyle: {{ color: '#b9b9b9' }} }},
        axisTick: {{ show: false }},
        axisLabel: {{ color: '#5a5a5a', fontSize: 11 }}
      }},
      yAxis: {{
        type: 'value',
        show: false,
        splitLine: {{ show: false }}
      }},
      series: series,
      backgroundColor: '#f3f3f3'
    }};
  }}

  var d1yoy = echarts.init(document.getElementById('echartsD1Yoy'));
  var d1mom = echarts.init(document.getElementById('echartsD1Mom'));
  var yoyOpt = makeBarOption('① 同比去年同期正式员工人数', yoyMonths, yoyData, true);
  var momOpt = makeBarOption({_json.dumps(mom_title)}, momMonths, momData, true);
  /* Add "最新月份" annotation to the MoM chart */
  momOpt.graphic = [{{
    type: 'text',
    right: 16,
    top: 8,
    style: {{ text: '最新月份: ' + latestMonth, fontSize: 11, fontWeight: 'bold', fill: '#4f4f4f' }}
  }}];
  /* Add notes if data is incomplete */
  if (notesText) {{
    yoyOpt.graphic = [{{
      type: 'text',
      left: 10,
      bottom: 4,
      style: {{ text: notesText, fontSize: 10, fill: '#b76e00' }}
    }}];
  }}
  d1yoy.setOption(yoyOpt);
  d1mom.setOption(momOpt);

  /* responsive resize */
  window.addEventListener('resize', function() {{ d1yoy.resize(); d1mom.resize(); }});
}})();
</script>"""


def _build_echarts_dashboard2_html(d2_df: pd.DataFrame, d2_meta: dict) -> str:
    """Build ECharts interactive pie charts for Dashboard 2."""
    import json as _json

    active_df = d2_df[d2_df["Metric"].eq("Active")] if not d2_df.empty else pd.DataFrame()
    attrition_df = d2_df[d2_df["Metric"].eq("AttritionYTD")] if not d2_df.empty else pd.DataFrame()

    def pie_data(df: pd.DataFrame) -> list:
        result = []
        for _, row in df.iterrows():
            result.append({
                "name": row["Region"],
                "value": int(row["Headcount"]),
                "share": round(float(row["Share"]), 1),
            })
        return result

    active_data = pie_data(active_df)
    attrition_data = pie_data(attrition_df)

    snapshot_label = compact_date_label(d2_meta.get("active_snapshot_date"))
    attrition_year = d2_meta.get("attrition_year") or ""
    cutoff_label = month_day_label(d2_meta.get("attrition_cutoff_date"))

    region_colors = {r: REGION_COLORS[r] for r in REGION_ORDER}

    return f"""<section class='section'>
  <h2>看板2：期末在离职分析</h2>
  <div class='charts-inline-row' style='display:flex;gap:8px;flex-wrap:nowrap;'>
    <div id='echartsD2Active' style='flex:1;min-width:0;height:320px;border-radius:12px;'></div>
    <div id='echartsD2Attrition' style='flex:1;min-width:0;height:320px;border-radius:12px;'></div>
  </div>
</section>
<script>
(function() {{
  var regionColors = {_json.dumps(region_colors)};
  var activeData = {_json.dumps(active_data)};
  var attritionData = {_json.dumps(attrition_data)};

  function makePieOption(titleParts, data, bgColor) {{
    var totalHC = data.reduce(function(s, d) {{ return s + d.value; }}, 0);
    return {{
      title: {{
        text: titleParts.map(function(p) {{ return p[0]; }}).join(''),
        left: 14,
        top: 10,
        textStyle: {{
          fontSize: 13,
          fontWeight: 'bold',
          color: '#4f4f4f',
          backgroundColor: '#d9d9d9',
          padding: [5, 10],
          borderRadius: 0,
          rich: titleParts.reduce(function(acc, p, i) {{
            acc['s' + i] = {{ color: p[1], fontWeight: 'bold', fontSize: 13 }};
            return acc;
          }}, {{}})
        }},
        textStyle: {{ fontSize: 13, fontWeight: 'bold', color: '#4f4f4f',
          backgroundColor: '#d9d9d9', padding: [5, 10] }}
      }},
      backgroundColor: bgColor,
      tooltip: {{
        trigger: 'item',
        formatter: function(p) {{
          return p.name + '<br/>人数: <b>' + p.value + '人</b> (' + p.data.share + '%)';
        }}
      }},
      series: [{{
        type: 'pie',
        radius: ['0%', '72%'],
        center: ['50%', '55%'],
        startAngle: 90,
        clockwise: false,
        data: data.map(function(d) {{
          return {{
            name: d.name,
            value: d.value,
            share: d.share,
            itemStyle: {{ color: regionColors[d.name] || '#999' }},
            label: {{
              show: true,
              position: 'inside',
              formatter: function(p) {{
                return p.name + '\\n' + p.value + '人 ' + p.data.share + '%';
              }},
              color: '#fff',
              fontWeight: 'bold',
              fontSize: 13,
              lineHeight: 20
            }}
          }};
        }}),
        emphasis: {{
          itemStyle: {{ shadowBlur: 10, shadowOffsetX: 0, shadowColor: 'rgba(0,0,0,0.2)' }}
        }},
        itemStyle: {{
          borderColor: '#fff',
          borderWidth: 2
        }}
      }}]
    }};
  }}

  /* Build rich-text titles */
  var activeTitle = '③ {snapshot_label} 正式员工在职人数';
  var attritionTitle = '③ {attrition_year}截止{cutoff_label}正式员工离职人数';

  var d2a = echarts.init(document.getElementById('echartsD2Active'));
  var d2b = echarts.init(document.getElementById('echartsD2Attrition'));

  var activeOpt = makePieOption(
    [['③ {snapshot_label} 正式员工', '#4f4f4f'], ['在职', '#1565C0'], ['人数', '#4f4f4f']],
    activeData,
    '#dfe8f1'
  );
  /* Use rich text for active title */
  activeOpt.title = {{
    left: 14, top: 10,
    text: '{{s0|③ {snapshot_label} 正式员工}}{{s1|在职}}{{s2|人数}}',
    textStyle: {{
      rich: {{
        s0: {{ color: '#4f4f4f', fontWeight: 'bold', fontSize: 13 }},
        s1: {{ color: '#1565C0', fontWeight: 'bold', fontSize: 13 }},
        s2: {{ color: '#4f4f4f', fontWeight: 'bold', fontSize: 13 }}
      }},
      backgroundColor: '#d9d9d9',
      padding: [5, 10]
    }}
  }};

  var attritionOpt = makePieOption(
    [['③ {attrition_year}截止{cutoff_label}正式员工', '#4f4f4f'], ['离职', '#C62828'], ['人数', '#4f4f4f']],
    attritionData,
    '#eadccf'
  );
  attritionOpt.title = {{
    left: 14, top: 10,
    text: '{{s0|③ {attrition_year}截止{cutoff_label}正式员工}}{{s1|离职}}{{s2|人数}}',
    textStyle: {{
      rich: {{
        s0: {{ color: '#4f4f4f', fontWeight: 'bold', fontSize: 13 }},
        s1: {{ color: '#C62828', fontWeight: 'bold', fontSize: 13 }},
        s2: {{ color: '#4f4f4f', fontWeight: 'bold', fontSize: 13 }}
      }},
      backgroundColor: '#d9d9d9',
      padding: [5, 10]
    }}
  }};

  d2a.setOption(activeOpt);
  d2b.setOption(attritionOpt);

  window.addEventListener('resize', function() {{ d2a.resize(); d2b.resize(); }});
}})();
</script>"""


def _build_combined_d1d2_html(d1_df: pd.DataFrame, d1_meta: dict, d2_df: pd.DataFrame, d2_meta: dict) -> str:
    """Build a single section combining Dashboard 1 (bar charts) and Dashboard 2 (pie charts)
    with all 4 charts displayed in one row."""
    import json as _json

    # --- Dashboard 1 data ---
    yoy_df = d1_df[d1_df["section"].eq("同比")] if not d1_df.empty else pd.DataFrame()
    mom_df = d1_df[d1_df["section"].eq("环比")] if not d1_df.empty else pd.DataFrame()
    yoy_months = yoy_df["month"].tolist() if not yoy_df.empty else []
    mom_months = mom_df["month"].tolist() if not mom_df.empty else []
    def region_data(df: pd.DataFrame, region: str) -> list:
        return df[region].tolist() if not df.empty and region in df.columns else []
    yoy_data = {r: region_data(yoy_df, r) for r in REGION_ORDER}
    mom_data = {r: region_data(mom_df, r) for r in REGION_ORDER}
    mom_count = d1_meta.get("mom_count", 3)
    mom_title = f"② 环比近{mom_count}个月正式员工人数" if mom_count >= 2 else "② 环比正式员工人数（数据不足）"
    latest_month = d1_meta.get("latest_month") or "N/A"
    notes_text = "；".join(d1_meta.get("notes", []))
    d1_region_colors = _json.dumps([REGION_COLORS[r] for r in REGION_ORDER])
    d1_region_names = _json.dumps(REGION_ORDER)

    # --- Dashboard 2 data ---
    active_df = d2_df[d2_df["Metric"].eq("Active")] if not d2_df.empty else pd.DataFrame()
    attrition_df = d2_df[d2_df["Metric"].eq("AttritionYTD")] if not d2_df.empty else pd.DataFrame()
    def pie_data(df: pd.DataFrame) -> list:
        result = []
        for _, row in df.iterrows():
            result.append({"name": row["Region"], "value": int(row["Headcount"]), "share": round(float(row["Share"]), 1)})
        return result
    active_data = pie_data(active_df)
    attrition_data = pie_data(attrition_df)
    snapshot_label = compact_date_label(d2_meta.get("active_snapshot_date"))
    attrition_year = d2_meta.get("attrition_year") or ""
    cutoff_label = month_day_label(d2_meta.get("attrition_cutoff_date"))
    d2_region_colors = {r: REGION_COLORS[r] for r in REGION_ORDER}

    # Pre-compute JSON strings to avoid f-string {{ conflicts
    yoy_data_json = _json.dumps({r: yoy_data[r] for r in REGION_ORDER})
    mom_data_json = _json.dumps({r: mom_data[r] for r in REGION_ORDER})
    yoy_months_json = _json.dumps(yoy_months)
    mom_months_json = _json.dumps(mom_months)
    latest_month_json = _json.dumps(latest_month)
    notes_text_json = _json.dumps(notes_text)
    mom_title_json = _json.dumps(mom_title)
    active_data_json = _json.dumps(active_data)
    attrition_data_json = _json.dumps(attrition_data)
    d2_region_colors_json = _json.dumps(d2_region_colors)

    return f"""<section class='section'>
  <h2>看板1&amp;2：正式员工趋势 &amp; 期末在离职分析</h2>
  <div class='charts-inline-row' style='display:flex;gap:6px;flex-wrap:nowrap;'>
    <div id='echartsD1Yoy' style='flex:1;min-width:0;height:300px;'></div>
    <div id='echartsD1Mom' style='flex:1.3;min-width:0;height:300px;'></div>
    <div id='echartsD2Active' style='flex:0.8;min-width:0;height:300px;border-radius:10px;'></div>
    <div id='echartsD2Attrition' style='flex:0.8;min-width:0;height:300px;border-radius:10px;'></div>
  </div>
</section>
<script>
(function() {{
  // === Dashboard 1: Bar charts ===
  var regionNames = {d1_region_names};
  var regionColors = {d1_region_colors};
  var yoyMonths = {yoy_months_json};
  var momMonths = {mom_months_json};
  var yoyData = {yoy_data_json};
  var momData = {mom_data_json};
  var latestMonth = {latest_month_json};
  var notesText = {notes_text_json};

  function makeBarOption(title, months, data, highlightLast) {{
    var series = regionNames.map(function(r, i) {{
      return {{
        name: r,
        type: 'bar',
        barGap: '8%',
        data: data[r] || [],
        itemStyle: {{ color: regionColors[i] }},
        label: {{
          show: true,
          position: 'top',
          fontWeight: 'bold',
          fontSize: 10,
          color: '#4d4d4d',
          formatter: function(p) {{ return p.value > 0 ? p.value : ''; }}
        }}
      }};
    }});
    var markAreaData = [];
    if (highlightLast && months.length > 0) {{
      markAreaData = [[ {{ xAxis: months.length - 1.5 }}, {{ xAxis: months.length - 0.5 }} ]];
    }}
    if (markAreaData.length > 0) {{
      series[0].markArea = {{
        silent: true,
        itemStyle: {{ color: 'rgba(180,180,180,0.08)', borderColor: '#ababab', borderWidth: 1, borderType: 'dashed' }},
        data: markAreaData
      }};
    }}
    return {{
      title: {{
        text: title,
        left: 'center',
        top: 4,
        textStyle: {{ fontSize: 12, fontWeight: 'bold', color: '#4f4f4f',
          backgroundColor: '#d9d9d9', padding: [4, 8], borderRadius: 0 }}
      }},
      legend: {{
        data: regionNames,
        top: 28,
        left: 'center',
        textStyle: {{ fontSize: 10 }},
        itemWidth: 12,
        itemHeight: 8,
        itemGap: 10
      }},
      tooltip: {{
        trigger: 'axis',
        axisPointer: {{ type: 'shadow' }}
      }},
      grid: {{ left: 36, right: 12, top: 56, bottom: 28 }},
      xAxis: {{
        type: 'category',
        data: months,
        axisLine: {{ lineStyle: {{ color: '#b9b9b9' }} }},
        axisTick: {{ show: false }},
        axisLabel: {{ color: '#5a5a5a', fontSize: 10 }}
      }},
      yAxis: {{
        type: 'value',
        show: false,
        splitLine: {{ show: false }}
      }},
      series: series,
      backgroundColor: '#f3f3f3'
    }};
  }}

  var d1yoy = echarts.init(document.getElementById('echartsD1Yoy'));
  var d1mom = echarts.init(document.getElementById('echartsD1Mom'));
  var yoyOpt = makeBarOption('① 同比去年同期正式员工人数', yoyMonths, yoyData, true);
  var momOpt = makeBarOption({mom_title_json}, momMonths, momData, true);
  momOpt.graphic = [{{
    type: 'text',
    right: 8,
    top: 4,
    style: {{ text: '最新月份: ' + latestMonth, fontSize: 10, fontWeight: 'bold', fill: '#4f4f4f' }}
  }}];
  if (notesText) {{
    yoyOpt.graphic = [{{
      type: 'text',
      left: 6,
      bottom: 2,
      style: {{ text: notesText, fontSize: 9, fill: '#b76e00' }}
    }}];
  }}
  d1yoy.setOption(yoyOpt);
  d1mom.setOption(momOpt);

  // === Dashboard 2: Pie charts ===
  var d2RegionColors = {d2_region_colors_json};
  var activeData = {active_data_json};
  var attritionData = {attrition_data_json};

  function makePieOption(titleParts, data, bgColor) {{
    return {{
      title: {{
        text: titleParts[0],
        left: 8,
        top: 4,
        textStyle: {{ fontSize: 11, fontWeight: 'bold', color: '#4f4f4f',
          backgroundColor: '#d9d9d9', padding: [3, 8] }}
      }},
      backgroundColor: bgColor,
      tooltip: {{
        trigger: 'item',
        formatter: function(p) {{
          return p.name + '<br/>人数: <b>' + p.value + '人</b> (' + p.data.share + '%)';
        }}
      }},
      series: [{{
        type: 'pie',
        radius: ['0%', '68%'],
        center: ['50%', '56%'],
        startAngle: 90,
        clockwise: false,
        data: data.map(function(d) {{
          return {{
            name: d.name,
            value: d.value,
            share: d.share,
            itemStyle: {{ color: d2RegionColors[d.name] || '#999' }},
            label: {{
              show: true,
              position: 'inside',
              formatter: function(p) {{
                return p.name + '\\n' + p.value + '人 ' + p.data.share + '%';
              }},
              color: '#fff',
              fontWeight: 'bold',
              fontSize: 11,
              lineHeight: 16
            }}
          }};
        }}),
        emphasis: {{
          itemStyle: {{ shadowBlur: 10, shadowOffsetX: 0, shadowColor: 'rgba(0,0,0,0.2)' }}
        }},
        itemStyle: {{
          borderColor: '#fff',
          borderWidth: 2
        }}
      }}]
    }};
  }}

  var d2a = echarts.init(document.getElementById('echartsD2Active'));
  var d2b = echarts.init(document.getElementById('echartsD2Attrition'));

  var activeOpt = makePieOption(
    ['③ {snapshot_label} 在职'],
    activeData,
    '#dfe8f1'
  );
  activeOpt.title = {{
    left: 8, top: 4,
    text: '{{s0|③ {snapshot_label} }}{{s1|在职}}{{s2|人数}}',
    textStyle: {{
      rich: {{
        s0: {{ color: '#4f4f4f', fontWeight: 'bold', fontSize: 11 }},
        s1: {{ color: '#1565C0', fontWeight: 'bold', fontSize: 11 }},
        s2: {{ color: '#4f4f4f', fontWeight: 'bold', fontSize: 11 }}
      }},
      backgroundColor: '#d9d9d9',
      padding: [3, 8]
    }}
  }};

  var attritionOpt = makePieOption(
    ['③ {attrition_year}截止{cutoff_label} 离职'],
    attritionData,
    '#eadccf'
  );
  attritionOpt.title = {{
    left: 8, top: 4,
    text: '{{s0|③ {attrition_year}截止{cutoff_label}}}{{s1|离职}}{{s2|人数}}',
    textStyle: {{
      rich: {{
        s0: {{ color: '#4f4f4f', fontWeight: 'bold', fontSize: 11 }},
        s1: {{ color: '#C62828', fontWeight: 'bold', fontSize: 11 }},
        s2: {{ color: '#4f4f4f', fontWeight: 'bold', fontSize: 11 }}
      }},
      backgroundColor: '#d9d9d9',
      padding: [3, 8]
    }}
  }};

  d2a.setOption(activeOpt);
  d2b.setOption(attritionOpt);

  window.addEventListener('resize', function() {{ d1yoy.resize(); d1mom.resize(); d2a.resize(); d2b.resize(); }});
}})();
</script>"""


def build_html(meta: dict, output_path: Path, dashboard1_image: str, dashboard2_image: str, d3_df: pd.DataFrame, d4_df: pd.DataFrame, d3_meta: dict, d4_meta: dict, d5_df: pd.DataFrame | None = None, executive_summary: dict[str, str] | None = None, d1_df: pd.DataFrame | None = None, d1_meta: dict | None = None, d2_df: pd.DataFrame | None = None, d2_meta: dict | None = None) -> None:
    css = load_css()
    # Use ECharts interactive charts when data is available; fall back to static images
    # When both D1 and D2 use ECharts, merge them into a single combined section
    if d1_df is not None and d1_meta is not None and d2_df is not None and d2_meta is not None:
        d1_card = _build_echarts_dashboard1_html(d1_df, d1_meta)
        d2_card = _build_echarts_dashboard2_html(d2_df, d2_meta)
        combined_card = _build_combined_d1d2_html(d1_df, d1_meta, d2_df, d2_meta)
        cards = [combined_card]
    else:
        if d1_df is not None and d1_meta is not None:
            d1_card = _build_echarts_dashboard1_html(d1_df, d1_meta)
        else:
            d1_card = build_chart_card_html("看板1：正式员工趋势图", dashboard1_image)
        if d2_df is not None and d2_meta is not None:
            d2_card = _build_echarts_dashboard2_html(d2_df, d2_meta)
        else:
            d2_card = build_chart_card_html("看板2：期末在离职分析", dashboard2_image)
        cards = [d1_card, d2_card]
    cards.extend([
        f"<section class='section'><h2>看板3：Active Regular & Intern 明细汇总表<button type='button' class='copy-table-btn' data-target='dashboard3Table' title='复制表格'>📋 复制</button></h2><div id='dashboard3Table'>{build_table_html(d3_df, grouped_header=True, region_annotations=d3_meta.get('region_annotations'))}</div></section>",
        f"<section class='section'><h2>看板4：Attrition Regular 离职分析表<button type='button' class='copy-table-btn' data-target='dashboard4Table' title='复制表格'>📋 复制</button></h2><div id='dashboard4Table'>{build_dashboard4_table_html(d4_df, d4_meta)}</div></section>",
    ])
    if d5_df is not None and not d5_df.empty:
        cards.append(
            f"<section class='section'><h2>看板5：Active Contractor &amp; Partner<button type='button' class='copy-table-btn' data-target='dashboard5Table' title='复制表格'>📋 复制</button></h2><div id='dashboard5Table'>{build_dashboard5_table_html(d5_df)}</div></section>"
        )

    # Executive Summary section (editable)
    es_html = ""
    if executive_summary:
        es_html = build_executive_summary_html(executive_summary)

    # Footnotes section
    footnotes_html = build_footnotes_html(meta, d4_meta)

    html = f"""<!DOCTYPE html>
<html lang='zh-CN'>
<head>
<meta charset='utf-8'>
<meta name='viewport' content='width=device-width, initial-scale=1'>
<title>{html_escape(meta['title'])}</title>
<script src='https://cdn.jsdelivr.net/npm/echarts@5/dist/echarts.min.js'></script>
<style>{css}</style>
</head>
<body>
<div class='container'>
  {es_html}
  {''.join(cards)}
  {footnotes_html}
  <div class='footer'>Generated by hr-workforce-dashboard skill</div>
</div>
<div class='image-modal' id='imageModal' hidden>
  <button type='button' class='modal-close' id='modalClose' aria-label='关闭'>×</button>
  <div class='modal-content'>
    <div class='modal-title' id='modalTitle'></div>
    <img id='modalImage' alt='大图预览'>
  </div>
</div>
<script>
const modal = document.getElementById('imageModal');
const modalImage = document.getElementById('modalImage');
const modalTitle = document.getElementById('modalTitle');
const closeModal = () => {{
  modal.hidden = true;
  modalImage.src = '';
}};
document.querySelectorAll('.zoomable-card').forEach((card) => {{
  card.addEventListener('click', () => {{
    const img = card.querySelector('img');
    modalImage.src = img.src;
    modalTitle.textContent = card.dataset.modalTitle || img.alt || '图表预览';
    modal.hidden = false;
  }});
}});
document.getElementById('modalClose').addEventListener('click', closeModal);
modal.addEventListener('click', (event) => {{
  if (event.target === modal) closeModal();
}});
document.addEventListener('keydown', (event) => {{
  if (event.key === 'Escape' && !modal.hidden) closeModal();
}});
function showCopied(b) {{
  const orig = b.textContent;
  b.textContent = '✅ 已复制';
  b.classList.add('copied');
  setTimeout(() => {{ b.textContent = orig; b.classList.remove('copied'); }}, 2000);
}}
function fallbackCopyText(text, btn) {{
  const ta = document.createElement('textarea');
  ta.value = text;
  ta.style.cssText = 'position:fixed;left:-9999px;top:-9999px;opacity:0';
  document.body.appendChild(ta);
  ta.select();
  try {{ document.execCommand('copy'); showCopied(btn); }} catch(e) {{}}
  document.body.removeChild(ta);
}}
function fallbackCopyHTML(table, btn) {{
  const range = document.createRange();
  range.selectNode(table);
  const sel = window.getSelection();
  sel.removeAllRanges();
  sel.addRange(range);
  try {{ document.execCommand('copy'); showCopied(btn); }} catch(e) {{}}
  sel.removeAllRanges();
}}
document.querySelectorAll('.copy-table-btn').forEach((btn) => {{
  btn.addEventListener('click', (e) => {{
    e.stopPropagation();
    const targetId = btn.dataset.target;
    const container = document.getElementById(targetId);
    if (!container) return;
    const table = container.querySelector('table');
    if (!table) return;
    if (typeof ClipboardItem !== 'undefined' && navigator.clipboard && navigator.clipboard.write) {{
      const htmlContent = table.outerHTML;
      const textContent = table.innerText;
      const item = new ClipboardItem({{
        'text/html': new Blob([htmlContent], {{ type: 'text/html' }}),
        'text/plain': new Blob([textContent], {{ type: 'text/plain' }})
      }});
      navigator.clipboard.write([item]).then(() => {{
        showCopied(btn);
      }}).catch(() => {{
        fallbackCopyHTML(table, btn);
      }});
    }} else {{
      fallbackCopyHTML(table, btn);
    }}
  }});
}});
// --- Executive Summary edit/save/download/copy ---
(function() {{
  const editBtn = document.getElementById('esEditBtn');
  const saveBtn = document.getElementById('esSaveBtn');
  const cancelBtn = document.getElementById('esCancelBtn');
  const downloadBtn = document.getElementById('esDownloadBtn');
  const copyBtn = document.getElementById('esCopyBtn');
  const fields = document.querySelectorAll('.es-row');
  if (!editBtn || !copyBtn || !fields.length) return;
  let savedState = {{}};
  function captureState() {{
    savedState = {{}};
    fields.forEach(f => {{ savedState[f.id] = f.innerHTML; }});
  }}
  captureState();
  editBtn.addEventListener('click', () => {{
    captureState();
    fields.forEach(f => {{
      f.contentEditable = 'true';
      f.classList.add('editing');
    }});
    editBtn.style.display = 'none';
    saveBtn.style.display = '';
    cancelBtn.style.display = '';
  }});
  saveBtn.addEventListener('click', () => {{
    fields.forEach(f => {{
      f.contentEditable = 'false';
      f.classList.remove('editing');
    }});
    captureState();
    editBtn.style.display = '';
    saveBtn.style.display = 'none';
    cancelBtn.style.display = 'none';
    saveBtn.textContent = '✅ 已保存';
    setTimeout(() => {{ saveBtn.textContent = '💾 保存'; }}, 1500);
  }});
  cancelBtn.addEventListener('click', () => {{
    fields.forEach(f => {{
      f.innerHTML = savedState[f.id] || '';
      f.contentEditable = 'false';
      f.classList.remove('editing');
    }});
    editBtn.style.display = '';
    saveBtn.style.display = 'none';
    cancelBtn.style.display = 'none';
  }});
  if (downloadBtn) {{
    downloadBtn.addEventListener('click', () => {{
      let text = 'Executive Summary\\n\\n';
      fields.forEach(f => {{
        text += f.innerText.trim() + '\\n';
      }});
      const blob = new Blob([text], {{ type: 'text/plain;charset=utf-8' }});
      const url = URL.createObjectURL(blob);
      const a = document.createElement('a');
      a.href = url;
      a.download = 'executive_summary.txt';
      a.click();
      URL.revokeObjectURL(url);
    }});
  }}
  copyBtn.addEventListener('click', () => {{
    let text = 'Executive Summary\\n\\n';
    fields.forEach(f => {{
      text += f.innerText.trim() + '\\n';
    }});
    if (navigator.clipboard && navigator.clipboard.writeText) {{
      navigator.clipboard.writeText(text).then(() => {{
        showCopied(copyBtn);
      }}).catch(() => {{
        fallbackCopyText(text, copyBtn);
      }});
    }} else {{
      fallbackCopyText(text, copyBtn);
    }}
  }});
}})();
// === 一键生成邮件功能（轻量版：无背景格式，方便微调） ===
(function() {{
  var mailBtn = document.getElementById('esMailBtn');
  if (!mailBtn) return;

  // --- 工具函数：提取日期字符串 YYYYMMDD ---
  function extractDateStr() {{
    var esActiveEl = document.getElementById('esActive');
    if (esActiveEl) {{
      var text = esActiveEl.textContent || esActiveEl.innerText || '';
      var m = text.match(/截至\\s*(\\d{{4}})\\s*[-\\/年]\\s*(\\d{{1,2}})\\s*[-\\/月]\\s*(\\d{{1,2}})/);
      if (m) return m[1] + m[2].padStart(2, '0') + m[3].padStart(2, '0');
    }}
    var now = new Date();
    return now.getFullYear() + String(now.getMonth() + 1).padStart(2, '0') + String(now.getDate()).padStart(2, '0');
  }}

  // --- 工具函数：生成纯文本版正文（用于 mailto body） ---
  function buildPlainBody() {{
    var lines = [];
    lines.push('Executive Summary');
    lines.push('');

    var fields = [
      {{ id: 'esActive',      title: '【在职情况】' }},
      {{ id: 'esTermination', title: '【离职情况】' }},
      {{ id: 'esOthers',      title: '【其他人员】' }}
    ];
    fields.forEach(function(f) {{
      var el = document.getElementById(f.id);
      if (el) {{
        lines.push(el.innerText.trim());
      }}
    }});

    // 提取备注
    var fnEl = document.getElementById('footnotesContent');
    if (fnEl) {{
      lines.push('');
      lines.push('看板备注');
      lines.push(fnEl.innerText.trim());
    }}

    lines.push('');
    lines.push('---');
    lines.push('（完整看板含图表已复制到剪贴板，请在正文区域 Ctrl+V / Cmd+V 粘贴即可获得图表版本）');
    return lines.join('\\n');
  }}

  // --- 工具函数：克隆 + 清理 DOM（保留表格颜色，去除非表格背景） ---
  function buildCleanClone() {{
    var container = document.querySelector('.container');
    if (!container) throw new Error('未找到 .container 容器');

    var clone = container.cloneNode(true);
    // 移除所有按钮
    clone.querySelectorAll('button').forEach(function(b) {{ b.remove(); }});
    // 移除 .zoom-hint
    clone.querySelectorAll('.zoom-hint').forEach(function(e) {{ e.remove(); }});
    // 移除 contentEditable
    clone.querySelectorAll('[contenteditable]').forEach(function(e) {{ e.removeAttribute('contenteditable'); }});
    // 移除 .editing
    clone.querySelectorAll('.editing').forEach(function(e) {{ e.classList.remove('editing'); }});

    // 仅移除 Executive Summary 的 h2 标题（看板标题和备注标题保留）
    var esSection = clone.querySelector('.executive-summary');
    if (esSection) {{
      var esH2 = esSection.querySelector('h2');
      if (esH2) esH2.remove();
    }}

    // 仅清除非表格元素的背景样式；保留表格行/单元格的 inline background
    clone.querySelectorAll('*').forEach(function(el) {{
      if (el.style) {{
        // 保留 <tr>、<td>、<th> 和 table 内部元素的背景颜色
        var tag = el.tagName.toUpperCase();
        var isTableElement = (tag === 'TR' || tag === 'TD' || tag === 'TH' || tag === 'TABLE' || tag === 'THEAD' || tag === 'TBODY');
        if (!isTableElement) {{
          el.style.background = '';
          el.style.backgroundColor = '';
          el.style.backgroundImage = '';
        }}
        el.style.boxShadow = '';
      }}
    }});
    return clone;
  }}

  // --- 工具函数：ECharts → 静态 <img>，用 table 横向排布（邮件客户端不支持 flex） ---
  function chartsToImages(clone) {{
    // 收集每个 .charts-inline-row 下的图表，转为 table 横排
    clone.querySelectorAll('.charts-inline-row').forEach(function(row) {{
      var chartDivs = row.querySelectorAll('[id^="echarts"]');
      if (chartDivs.length === 0) return;

      // 计算每个 chart 的宽度百分比（按原始 flex 比例）
      var totalFlex = 0;
      var flexValues = [];
      chartDivs.forEach(function(div) {{
        var flexVal = 1;
        var flexMatch = (div.style.cssText || '').match(/flex\s*:\s*([\d.]+)/);
        if (flexMatch) flexVal = parseFloat(flexMatch[1]);
        flexValues.push(flexVal);
        totalFlex += flexVal;
      }});

      var tds = [];
      chartDivs.forEach(function(div, idx) {{
        var chartId = div.id;
        var originalEl = document.getElementById(chartId);
        if (!originalEl) {{
          tds.push('<td></td>');
          return;
        }}
        var inst = echarts.getInstanceByDom(originalEl);
        if (!inst) {{
          tds.push('<td></td>');
          return;
        }}
        var dataURL = inst.getDataURL({{ type: 'png', pixelRatio: 2, backgroundColor: '#fff' }});
        var pct = Math.round(flexValues[idx] / totalFlex * 100);
        tds.push(
          '<td style="width:' + pct + '%;padding:2px;vertical-align:top;">'
          + '<img src="' + dataURL + '" style="width:100%;height:auto;display:block;" alt="' + chartId + '">'
          + '</td>'
        );
      }});

      // 用 table 替换整个 .charts-inline-row
      var tableHTML = '<table style="width:100%;border-collapse:collapse;border:none;table-layout:fixed;">'
        + '<tr>' + tds.join('') + '</tr></table>';
      row.innerHTML = tableHTML;
      // 清除 flex 样式，避免干扰
      row.style.cssText = '';
    }});
  }}

  // --- 工具函数：生成干净 HTML（保留表格颜色，图表用 table 横排） ---
  function buildCleanHTML(clone) {{
    var minimalCSS = '<style>'
      + 'table {{ border-collapse: collapse; font-size: 13px; }}'
      + '.dashboard-table {{ width: 100%; }}'
      + '.dashboard-table th, .dashboard-table td {{ border: 1px solid #8a9199; padding: 6px 8px; text-align: center; white-space: nowrap; }}'
      + '.dashboard-table th {{ background: #b7d7e8; color: #243746; font-weight: bold; }}'
      + 'td.region-cell {{ text-align: center; vertical-align: middle; font-weight: bold; color: #1f2f46; padding: 10px 12px; }}'
      + '.region-label {{ color: #1f2f46; font-size: 14px; font-weight: bold; }}'
      + '.region-meta-label {{ color: #3c78a8; font-size: 12px; font-weight: bold; }}'
      + '.region-meta-value {{ color: #1f2f46; font-size: 13px; font-weight: bold; }}'
      + '.region-meta-line {{ color: #1f2f46; font-size: 13px; font-weight: bold; }}'
      + '.region-metric {{ margin-top: 4px; }}'
      + '.fn-red {{ color: #e53e3e; font-weight: bold; }}'
      + 'strong {{ color: #0f4c81; }}'
      + 'img {{ max-width: 100%; height: auto; }}'
      + '.es-row-title {{ font-weight: bold; color: #0f4c81; }}'
      + '.es-row-body strong {{ color: #0f4c81; }}'
      + '</style>';
    return minimalCSS + clone.innerHTML;
  }}

  // --- 主流程 ---
  mailBtn.addEventListener('click', async function() {{
    try {{
      mailBtn.disabled = true;
      mailBtn.textContent = '⏳ 生成中...';

      // 1. 生成邮件主题
      var dateStr = extractDateStr();
      var mailSubject = dateStr + '在职&累计离职人员数据统计看板 - 海外各区域及国内NHS试点';

      // 2. 生成纯文本正文（用于 mailto body 预填）
      var plainBody = buildPlainBody();

      // 3. 克隆 DOM + 图表转图片 → 干净 HTML（无背景格式）
      var clone = buildCleanClone();
      chartsToImages(clone);
      var htmlContent = buildCleanHTML(clone);
      var plainFallback = document.querySelector('.container').innerText;

      // 4. 写入系统剪贴板（轻量 HTML，无背景格式，粘贴后可自由微调）
      if (typeof ClipboardItem !== 'undefined' && navigator.clipboard && navigator.clipboard.write) {{
        await navigator.clipboard.write([
          new ClipboardItem({{
            'text/html': new Blob([htmlContent], {{ type: 'text/html' }}),
            'text/plain': new Blob([plainFallback], {{ type: 'text/plain' }})
          }})
        ]);
      }} else {{
        var tempDiv = document.createElement('div');
        tempDiv.innerHTML = htmlContent;
        tempDiv.style.cssText = 'position:fixed;left:-9999px;top:-9999px;opacity:0;';
        document.body.appendChild(tempDiv);
        var range = document.createRange();
        range.selectNodeContents(tempDiv);
        var sel = window.getSelection();
        sel.removeAllRanges();
        sel.addRange(range);
        document.execCommand('copy');
        sel.removeAllRanges();
        document.body.removeChild(tempDiv);
      }}

      // 5. 唤起企微/系统邮件客户端，自动填充主题 + 纯文本正文
      var mailtoURL = 'mailto:?subject=' + encodeURIComponent(mailSubject)
                    + '&body=' + encodeURIComponent(plainBody);
      window.location.href = mailtoURL;

      // 6. UI 反馈
      mailBtn.innerHTML = '✅ 已复制 — 请在邮件正文中 <kbd style="padding:1px 6px;border-radius:4px;border:1px solid #ccc;font-size:12px;">Ctrl+V</kbd> 粘贴看板';
      mailBtn.classList.add('mail-success');
      mailBtn.disabled = false;

      setTimeout(function() {{
        mailBtn.textContent = '✉️ 生成邮件';
        mailBtn.classList.remove('mail-success');
      }}, 5000);

    }} catch (err) {{
      console.error('生成邮件失败:', err);
      mailBtn.textContent = '❌ 生成失败，请重试';
      mailBtn.disabled = false;
      setTimeout(function() {{
        mailBtn.textContent = '✉️ 生成邮件';
      }}, 3000);
    }}
  }});
}})();
// --- Footnotes copy ---
(function() {{
  const btn = document.getElementById('copyFootnotesBtn');
  const content = document.getElementById('footnotesContent');
  if (!btn || !content) return;
  btn.addEventListener('click', () => {{
    const text = content.innerText;
    if (navigator.clipboard && navigator.clipboard.writeText) {{
      navigator.clipboard.writeText(text).then(() => {{
        showCopied(btn);
      }}).catch(() => {{
        fallbackCopyText(text, btn);
      }});
    }} else {{
      fallbackCopyText(text, btn);
    }}
  }});
}})();
</script>
</body>
</html>
"""
    output_path.write_text(html, encoding="utf-8")


def cleanup_output_dir(output_dir: Path) -> None:
    for folder_name in ["png", "excel", "ppt"]:
        folder = output_dir / folder_name
        if folder.exists():
            shutil.rmtree(folder)
    for file_path in output_dir.iterdir() if output_dir.exists() else []:
        if file_path.is_file() and file_path.name != "dashboard.html":
            file_path.unlink()


def excel_safe_value(value):
    if isinstance(value, pd.Period):
        return str(value)
    if isinstance(value, pd.Timestamp):
        return value.to_pydatetime()
    try:
        if pd.isna(value):
            return ""
    except Exception:
        pass
    return value


def write_dataframe_sheet(ws, df: pd.DataFrame, percent_cols: set[str] | None = None) -> None:
    percent_cols = percent_cols or set()
    thin = Side(style="thin", color="666666")
    header_fill = PatternFill("solid", fgColor="B7D7E8")
    for row in df.itertuples(index=False):
        ws.append([excel_safe_value(value) for value in row])
    if ws.max_row == 1 and ws.max_column == 1 and ws["A1"].value is None:
        ws.delete_rows(1)
    ws.insert_rows(1)
    for col_idx, col_name in enumerate(df.columns, start=1):
        cell = ws.cell(row=1, column=col_idx, value=col_name)
        cell.font = Font(bold=True)
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center", vertical="center")
    for row in ws.iter_rows():
        for cell in row:
            cell.border = Border(left=thin, right=thin, top=thin, bottom=thin)
            cell.alignment = Alignment(horizontal="center", vertical="center")
    for col_idx, col_name in enumerate(df.columns, start=1):
        letter = get_column_letter(col_idx)
        max_len = max(len(clean_text(cell.value)) for cell in ws[letter]) if ws[letter] else len(col_name)
        ws.column_dimensions[letter].width = min(max(max_len + 2, 12), 28)
        if col_name in percent_cols:
            for row_idx in range(2, ws.max_row + 1):
                ws.cell(row=row_idx, column=col_idx).number_format = '0.0%'


def apply_row_fills(ws, row_kind_col: int | None = None) -> None:
    if row_kind_col is None:
        return
    for row_idx in range(2, ws.max_row + 1):
        kind = clean_text(ws.cell(row=row_idx, column=row_kind_col).value)
        fill = ROW_FILL.get(kind)
        if fill:
            for col_idx in range(1, ws.max_column + 1):
                ws.cell(row=row_idx, column=col_idx).fill = PatternFill("solid", fgColor=fill)


def build_excel_bundle(output_dir: Path, bundle_name: str, meta: dict, cleaned_active: pd.DataFrame, cleaned_term: pd.DataFrame, d1_df: pd.DataFrame, d2_df: pd.DataFrame, d3_df: pd.DataFrame, d4_df: pd.DataFrame) -> Path:
    excel_dir = output_dir / "excel"
    excel_dir.mkdir(parents=True, exist_ok=True)
    workbook_path = excel_dir / f"{bundle_name}.xlsx"
    wb = Workbook()
    wb.remove(wb.active)

    summary_ws = wb.create_sheet("Summary")
    summary_rows = [
        ["标题", meta["title"]],
        ["生成时间", meta["generated_at"]],
        ["最新月份", meta["latest_month"] or "N/A"],
        ["Attrition Start", meta["period_start"] or "N/A"],
        ["Attrition End", meta["period_end"] or "N/A"],
        ["Active Files", ", ".join(meta["active_files"])],
        ["Termination Files", ", ".join(meta["termination_files"])],
        ["Notes", " | ".join(meta["notes"]) or "None"],
    ]
    for row in summary_rows:
        summary_ws.append(row)
    summary_ws.column_dimensions["A"].width = 20
    summary_ws.column_dimensions["B"].width = 90

    active_sheet = wb.create_sheet("Cleaned_Active")
    write_dataframe_sheet(active_sheet, cleaned_active.fillna(""))
    term_sheet = wb.create_sheet("Cleaned_Termination")
    write_dataframe_sheet(term_sheet, cleaned_term.fillna(""))

    d1_sheet = wb.create_sheet("Dashboard1")
    write_dataframe_sheet(d1_sheet, d1_df.fillna(""))
    d1_sheet.add_image(XLImage(str(output_dir / "png" / "dashboard_1.png")), "H2")

    d2_sheet = wb.create_sheet("Dashboard2")
    write_dataframe_sheet(d2_sheet, d2_df.assign(Share=d2_df["Share"] / 100).fillna(""), percent_cols={"Share"})
    d2_sheet.add_image(XLImage(str(output_dir / "png" / "dashboard_2.png")), "F2")

    d3_sheet = wb.create_sheet("Dashboard3")
    write_dataframe_sheet(d3_sheet, d3_df.fillna(""))
    apply_row_fills(d3_sheet, d3_df.columns.get_loc("row_kind") + 1)

    d4_export = d4_df.copy()
    d4_percent_cols = {column for column in d4_export.columns if "Attrition %" in column}
    for column in d4_percent_cols:
        d4_export[column] = d4_export[column] / 100
    d4_sheet = wb.create_sheet("Dashboard4")
    write_dataframe_sheet(d4_sheet, d4_export.fillna(""), percent_cols=d4_percent_cols)
    apply_row_fills(d4_sheet, d4_df.columns.get_loc("row_kind") + 1)

    wb.save(workbook_path)
    return workbook_path


def build_ppt_bundle(output_dir: Path, bundle_name: str, meta: dict) -> Path:
    ppt_dir = output_dir / "ppt"
    ppt_dir.mkdir(parents=True, exist_ok=True)
    ppt_path = ppt_dir / f"{bundle_name}.pptx"
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    titles = [
        "看板1：正式员工趋势图",
        "看板2：期末在离职分析",
        "看板3：Active Regular & Intern 明细汇总表",
        "看板4：Attrition Regular 离职分析表",
    ]
    for idx, title in enumerate(titles, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.4), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        subtitle = slide.shapes.add_textbox(Inches(0.4), Inches(0.7), Inches(12.4), Inches(0.3))
        subtitle.text_frame.text = f"{meta['title']} | {meta['generated_at']}"
        subtitle.text_frame.paragraphs[0].font.size = Pt(11)
        slide.shapes.add_picture(str(output_dir / "png" / f"dashboard_{idx}.png"), Inches(0.35), Inches(1.0), width=Inches(12.6))
    prs.save(ppt_path)
    return ppt_path


def write_summary(output_dir: Path, meta: dict, bundle_name: str) -> Path:
    summary_path = output_dir / "summary.md"
    lines = [
        f"# {meta['title']}",
        "",
        f"- **生成时间**: {meta['generated_at']}",
        f"- **最新月份**: {meta['latest_month'] or 'N/A'}",
        f"- **Attrition 周期**: {meta['period_start'] or 'N/A'} ~ {meta['period_end'] or 'N/A'}",
        f"- **Active 文件**: {', '.join(meta['active_files']) or '无'}",
        f"- **Termination 文件**: {', '.join(meta['termination_files']) or '无'}",
        "",
        "## 数据完整性说明",
        "",
    ]
    if meta["notes"]:
        lines.extend([f"- {note}" for note in meta["notes"]])
    else:
        lines.append("- 数据完整，可直接使用。")
    lines.extend([
        "",
        "## 输出文件",
        "",
        "- `dashboard.html`",
        "- `png/dashboard_1.png` ~ `png/dashboard_4.png`",
        f"- `excel/{bundle_name}.xlsx`",
        f"- `ppt/{bundle_name}.pptx`",
        f"- `{bundle_name}.zip`",
    ])
    summary_path.write_text("\n".join(lines), encoding="utf-8")
    return summary_path


def package_zip(output_dir: Path, bundle_name: str) -> Path:
    zip_path = output_dir / f"{bundle_name}.zip"
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for relative in [
            Path("dashboard.html"),
            Path("summary.md"),
            Path("png/dashboard_1.png"),
            Path("png/dashboard_2.png"),
            Path("png/dashboard_3.png"),
            Path("png/dashboard_4.png"),
            Path(f"excel/{bundle_name}.xlsx"),
            Path(f"ppt/{bundle_name}.pptx"),
        ]:
            absolute = output_dir / relative
            if absolute.exists():
                zf.write(absolute, arcname=str(relative))
    return zip_path


def main() -> None:
    args = parse_args()
    output_dir = Path(args.output_dir).resolve()
    output_dir.mkdir(parents=True, exist_ok=True)
    cleanup_output_dir(output_dir)
    active_frames = []
    term_frames = []
    contingent_frames = []
    active_files = []
    term_files = []
    contingent_files = []
    notes: list[str] = []

    for file_name in args.files:
        file_path = Path(file_name).expanduser().resolve()
        # Peek at report name to decide type before full parse
        wb_peek = load_workbook(file_path, read_only=True, data_only=True)
        ws_peek = wb_peek[wb_peek.sheetnames[0]]
        row1_vals = [c.value for c in next(ws_peek.iter_rows(min_row=1, max_row=1))]
        report_name = clean_text(row1_vals[0]) if row1_vals else ""
        wb_peek.close()
        dtype = classify_dataset_type(report_name)

        if dtype == "contingent":
            cdf, _cmeta = read_contingent_dataset(file_path)
            contingent_frames.append(cdf)
            contingent_files.append(file_path.name)
        else:
            df, dataset_type, _metadata = read_dataset(file_path)
            if dataset_type == "active":
                active_frames.append(df)
                active_files.append(file_path.name)
            else:
                term_frames.append(df)
                term_files.append(file_path.name)

    if not active_frames:
        raise SystemExit("至少需要 1 个 active 文件才能生成看板")

    active_df = pd.concat(active_frames, ignore_index=True)
    term_df = pd.concat(term_frames, ignore_index=True) if term_frames else pd.DataFrame()
    contingent_df = pd.concat(contingent_frames, ignore_index=True) if contingent_frames else pd.DataFrame()

    d1_meta, d1_df = build_dashboard1(active_df, notes)
    d2_meta, d2_df = build_dashboard2(active_df, term_df, notes)
    d3_meta, d3_df = build_dashboard3(active_df, notes)
    d4_meta, d4_df = build_dashboard4(active_df, term_df, notes)

    d5_meta, d5_df = build_dashboard5(contingent_df, notes) if not contingent_df.empty else ({"has_data": False}, None)

    meta = {
        "title": args.title,
        "generated_at": datetime.now().strftime("%Y-%m-%d %H:%M"),
        "latest_month": d2_meta.get("latest_month") or d3_meta.get("latest_month"),
        "period_start": d4_meta.get("period_start"),
        "period_end": d4_meta.get("period_end"),
        "active_files": active_files,
        "termination_files": term_files,
        "contingent_files": contingent_files,
        "notes": notes,
    }

    dashboard1_image = render_dashboard1_image(d1_df, d1_meta)
    dashboard2_image = render_dashboard2_image(d2_df, d2_meta)

    # Generate Executive Summary defaults
    executive_summary = build_executive_summary_defaults(
        meta, d1_meta, d2_meta, d2_df, d3_meta, d3_df, d4_meta, d4_df,
        d5_meta=d5_meta if not contingent_df.empty else None,
        d5_df=d5_df,
    )

    html_path = output_dir / "dashboard.html"
    build_html(meta, html_path, dashboard1_image, dashboard2_image, d3_df, d4_df, d3_meta, d4_meta, d5_df=d5_df, executive_summary=executive_summary, d1_df=d1_df, d1_meta=d1_meta, d2_df=d2_df, d2_meta=d2_meta)
    cleanup_output_dir(output_dir)
    print(html_path)


if __name__ == "__main__":
    main()
