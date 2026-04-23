"""
extract_text.py - 通用文档文本提取工具
支持：PDF / Word (.docx, .doc) / Excel (.xlsx, .xls) / TXT / PPT
依赖：fitz (PyMuPDF), python-docx, openpyxl, pdfplumber, xlrd, python-pptx
"""

import sys
import os
import argparse
import logging

logging.basicConfig(level=logging.INFO, format='%(levelname)s: %(message)s')
logger = logging.getLogger(__name__)


def extract_pdf(filepath: str) -> str:
    """提取 PDF 文本内容（优先 PyMuPDF，降级 pdfplumber）"""
    text_parts = []
    try:
        import fitz
        doc = fitz.open(filepath)
        logger.info(f"  [PyMuPDF] 读取 PDF，共 {doc.page_count} 页")
        for i, page in enumerate(doc):
            t = page.get_text()
            if t.strip():
                text_parts.append(f"\n===== 第 {i+1} 页 =====\n{t}")
            else:
                # 尝试提取图片中的文字（OCR 备选）
                logger.info(f"  第 {i+1} 页无文本，尝试表格/图片提取...")
        doc.close()
    except Exception as e:
        logger.warning(f"  PyMuPDF 失败: {e}，尝试 pdfplumber...")
        try:
            import pdfplumber
            with pdfplumber.open(filepath) as pdf:
                logger.info(f"  [pdfplumber] 读取 PDF，共 {len(pdf.pages)} 页")
                for i, page in enumerate(pdf.pages):
                    t = page.extract_text() or ""
                    if t.strip():
                        text_parts.append(f"\n===== 第 {i+1} 页 =====\n{t}")
        except Exception as e2:
            logger.error(f"  pdfplumber 也失败: {e2}")
    return "\n".join(text_parts)


def extract_docx(filepath: str) -> str:
    """提取 Word .docx 文本内容"""
    import docx
    logger.info(f"  [python-docx] 读取 Word 文档")
    doc = docx.Document(filepath)
    paragraphs = []
    for i, para in enumerate(doc.paragraphs):
        if para.text.strip():
            paragraphs.append(para.text)
    # 提取表格
    for i, table in enumerate(doc.tables):
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(f"[表格行] {row_text}")
    return "\n".join(paragraphs)


def extract_doc(filepath: str) -> str:
    """提取 Word .doc 文本内容（通过 python-docx 兼容模式或 textutil）"""
    # 尝试 python-docx 直接读（仅限 OOXML 格式的 .doc）
    try:
        import docx
        doc = docx.Document(filepath)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        if paragraphs:
            logger.info(f"  [python-docx] 读取 .doc 成功")
            return "\n".join(paragraphs)
    except Exception:
        pass
    # 备选：macOS/Windows 系统工具
    import platform
    system = platform.system()
    if system == "Darwin":
        import subprocess
        try:
            result = subprocess.run(
                ["textutil", "-convert", "txt", "-stdout", filepath],
                capture_output=True, text=True, timeout=30
            )
            if result.returncode == 0:
                logger.info(f"  [textutil] 读取 .doc 成功")
                return result.stdout
        except Exception as e:
            logger.warning(f"  textutil 失败: {e}")
    elif system == "Windows":
        # 尝试使用 pywin32 或直接读取
        try:
            # python-docx 本身不支持 .doc，但可以通过 antiword 备选
            import subprocess
            # 尝试 antiword（需安装）
            result = subprocess.run(
                ["antiword", filepath], capture_output=True, text=True, timeout=30
            )
            if result.returncode == 0:
                logger.info(f"  [antiword] 读取 .doc 成功")
                return result.stdout
        except FileNotFoundError:
            logger.warning("  antiword 未安装，.doc 文件需要手动转换为 .docx")
        except Exception as e:
            logger.warning(f"  antiword 失败: {e}")
    logger.info(f"  .doc 文件尝试多种方式均失败，建议手动另存为 .docx")
    return ""


def extract_xlsx(filepath: str) -> str:
    """提取 Excel .xlsx / .xls 文本内容"""
    import openpyxl
    logger.info(f"  [openpyxl] 读取 Excel 文档")
    wb = openpyxl.load_workbook(filepath, data_only=True, read_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"\n===== Sheet: {sheet_name} =====")
        for row in ws.iter_rows(values_only=True):
            row_vals = [str(c) if c is not None else "" for c in row]
            line = " | ".join(v for v in row_vals if v)
            if line:
                parts.append(line)
    wb.close()
    return "\n".join(parts)


def extract_xls(filepath: str) -> str:
    """提取 Excel .xls 文本内容（xlrd）"""
    try:
        import xlrd
        logger.info(f"  [xlrd] 读取 Excel .xls 文档")
        wb = xlrd.open_workbook(filepath)
        parts = []
        for idx in range(wb.nsheets):
            sheet = wb.sheet_by_index(idx)
            parts.append(f"\n===== Sheet: {sheet.name} =====")
            for row_idx in range(sheet.nrows):
                row_vals = [str(sheet.cell_value(row_idx, col_idx)) for col_idx in range(sheet.ncols)]
                line = " | ".join(v for v in row_vals if v and v != "")
                if line:
                    parts.append(line)
        return "\n".join(parts)
    except ImportError:
        logger.warning("  xlrd 未安装，尝试 pandas 读取 .xls...")
        try:
            import pandas as pd
            logger.info(f"  [pandas] 读取 Excel .xls 文档")
            df = pd.read_excel(filepath, sheet_name=None, engine="xlrd")
            parts = []
            for sheet_name, sheet_df in df.items():
                parts.append(f"\n===== Sheet: {sheet_name} =====")
                parts.append(sheet_df.to_string(index=False))
            return "\n".join(parts)
        except Exception as e:
            logger.error(f"  pandas 读取 .xls 也失败: {e}")
            return ""


def extract_pptx(filepath: str) -> str:
    """提取 PPT .pptx 文本内容"""
    try:
        from pptx import Presentation
        logger.info(f"  [python-pptx] 读取 PPT 文档")
        prs = Presentation(filepath)
        parts = []
        for i, slide in enumerate(prs.slides):
            parts.append(f"\n===== 第 {i+1} 页幻灯片 =====")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)
        return "\n".join(parts)
    except ImportError:
        logger.warning("  python-pptx 未安装，无法提取 PPT 内容")
        return ""


def extract_txt(filepath: str) -> str:
    """提取纯文本文件（自动检测编码）"""
    encodings = ["utf-8", "gbk", "gb2312", "gb18030", "utf-16", "latin-1"]
    for enc in encodings:
        try:
            with open(filepath, "r", encoding=enc) as f:
                content = f.read()
            logger.info(f"  [TXT] 成功以 {enc} 编码读取")
            return content
        except (UnicodeDecodeError, UnicodeError):
            continue
    logger.error(f"  无法识别文本文件编码")
    return ""


def extract_file(filepath: str) -> str:
    """
    根据文件扩展名自动选择提取方法
    返回提取的纯文本内容
    """
    ext = os.path.splitext(filepath)[1].lower()
    logger.info(f"正在提取文件: {os.path.basename(filepath)} (.{ext})")

    extractor_map = {
        ".pdf": extract_pdf,
        ".docx": extract_docx,
        ".doc": extract_doc,
        ".xlsx": extract_xlsx,
        ".xls": extract_xls,
        ".pptx": extract_pptx,
        ".txt": extract_txt,
        ".md": extract_txt,
        ".csv": extract_txt,
    }

    extractor = extractor_map.get(ext)
    if not extractor:
        logger.error(f"  不支持的文件类型: {ext}")
        return ""

    text = extractor(filepath)
    logger.info(f"  提取完成，共 {len(text)} 个字符")
    return text


def main():
    parser = argparse.ArgumentParser(description="通用文档文本提取工具")
    parser.add_argument("file", help="要提取的文件路径")
    parser.add_argument("-o", "--output", help="输出文件路径（默认为标准输出）")
    parser.add_argument("--encoding", default="utf-8", help="输出编码（默认 utf-8）")
    args = parser.parse_args()

    if not os.path.exists(args.file):
        print(f"错误: 文件不存在: {args.file}", file=sys.stderr)
        sys.exit(1)

    text = extract_file(args.file)

    if args.output:
        with open(args.output, "w", encoding=args.encoding) as f:
            f.write(text)
        logger.info(f"内容已保存到: {args.output}")
    else:
        print(text)


if __name__ == "__main__":
    main()
