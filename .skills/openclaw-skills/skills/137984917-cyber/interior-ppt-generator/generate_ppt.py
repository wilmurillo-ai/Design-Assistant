#!/usr/bin/env python3
"""
设计方案PPT自动生成器 - 隐室空间设计版 v3.0
==============================================
大幅升级设计品质：
- 深色+橙色品牌配色体系
- 全局字体统一处理
- 封面：全屏大图+渐变遮罩+居中标题
- 内容页：左图右文/右图左文多种版式
- 效果页：沉浸式全屏图+底部标题条
- 材质/打卡点：用色块网格代替纯文字
- 页码/品牌角标统一
"""

import os
import sys
import math
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.enum.shapes import MSO_SHAPE

# ─── 品牌配色 ────────────────────────────────────────────────
C = {
    'deep':    RGBColor(18,  22,  28),   # 深邃黑 #12161c
    'ink':     RGBColor(28,  35,  45),   # 墨色 #1c232d
    'slate':   RGBColor(55,  65,  85),   # 石板蓝灰 #374055
    'muted':   RGBColor(100, 115, 140),  # 灰蓝 #64648c
    'warm':    RGBColor(220, 175, 120),  # 暖金 #dcb478
    'accent':  RGBColor(195, 145,  80),  # 品牌金 #c39150
    'cream':   RGBColor(245, 240, 230),  # 米白 #f5f0e6
    'white':   RGBColor(255, 255, 255),
    'light':   RGBColor(242, 242, 245),   # 浅灰背景
}

# ─── 字体处理 ────────────────────────────────────────────────
FONT_CN = 'PingFang SC'
FONT_EN = 'Avenir Next'
FONT_MONO = 'SF Mono'

def apply_font(run, font_name=FONT_CN, size=Pt(14), bold=False, color=None, italic=False):
    run.font.name = font_name
    try:
        run._element.rPr.rFonts.set(qn('w:eastAsia'), font_name)
    except Exception:
        pass
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color

def add_para(tf, text, font=FONT_CN, size=Pt(14), bold=False,
             color=None, align=PP_ALIGN.LEFT, space_before=Pt(0), space_after=Pt(8)):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    p.space_after = space_after
    run = p.add_run()
    apply_font(run, font, size, bold, color)
    run.text = text
    return p

def set_tf_style(tf, font=FONT_CN, size=Pt(13), color=None, bold=False,
                 align=PP_ALIGN.LEFT, line_spacing=1.35):
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(6)
        for r in p.runs:
            apply_font(r, font, size, bold, color)

# ─── 形状工具 ────────────────────────────────────────────────
def fill_shape(slide, shape_type, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, left, top, width, height,
                font=FONT_CN, size=Pt(13), bold=False, color=None,
                align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    if text:
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        apply_font(r, font, size, bold, color)
        r.text = text
    return tb, tf

# ─── 页面框架 ────────────────────────────────────────────────
def new_slide(prs, bg=C['white']):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 全屏背景
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg
    bg_shape.line.fill.background()
    return slide

def add_page_number(slide, num, total, color=C['muted']):
    """右下角页码"""
    num_text = f"{num:02d} / {total:02d}"
    tb = slide.shapes.add_textbox(
        Inches(14.5), Inches(8.6), Inches(1.2), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    apply_font(r, FONT_EN, Pt(9), False, color)

def add_brand_corner(slide, name="温州隐室空间设计"):
    """左上角品牌角标"""
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.35), Inches(3), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    apply_font(r, FONT_CN, Pt(8), False, C['muted'])
    r.text = name

def add_bottom_bar(slide, height=Inches(0.6), color=C['ink']):
    """底部色条"""
    return fill_shape(slide, MSO_SHAPE.RECTANGLE,
                      0, Inches(9) - height,
                      prs.slide_width, height, color)

# ─── 幻灯片类型 ──────────────────────────────────────────────

def slide_cover(prs, project_name, style, area, image_path=None):
    """封面：全屏图+渐变遮罩+居中标题"""
    slide = new_slide(prs, bg=C['deep'])

    # 插入图片（全屏覆盖）
    if image_path and os.path.exists(image_path):
        try:
            pic = slide.shapes.add_picture(image_path, 0, 0,
                                           width=prs.slide_width, height=prs.slide_height)
        except Exception:
            pass

    # 渐变遮罩（用深色半透明矩形叠加）
    overlay = fill_shape(slide, MSO_SHAPE.RECTANGLE,
                         0, 0, prs.slide_width, prs.slide_height,
                         RGBColor(18, 22, 28))
    # 用图片亮度模拟渐变：上方深一些
    grad_top = fill_shape(slide, MSO_SHAPE.RECTANGLE,
                          0, 0, prs.slide_width, Inches(4),
                          RGBColor(10, 12, 16))

    # 金色细线装饰
    line_top = fill_shape(slide, MSO_SHAPE.RECTANGLE,
                          Inches(5), Inches(3.3), Inches(6), Pt(1.5),
                          C['warm'])

    # 主标题
    add_textbox(slide, project_name,
                Inches(1), Inches(3.5), Inches(14), Inches(1.2),
                font=FONT_CN, size=Pt(48), bold=True,
                color=C['cream'], align=PP_ALIGN.CENTER)

    # 副标题
    add_textbox(slide, f"{style}  ·  {area}㎡",
                Inches(1), Inches(4.7), Inches(14), Inches(0.5),
                font=FONT_CN, size=Pt(20), bold=False,
                color=C['warm'], align=PP_ALIGN.CENTER)

    # 金色细线
    line_bot = fill_shape(slide, MSO_SHAPE.RECTANGLE,
                          Inches(5), Inches(5.4), Inches(6), Pt(1.5),
                          C['warm'])

    # 品牌信息
    add_textbox(slide, "温州隐室空间设计",
                Inches(1), Inches(7.8), Inches(14), Inches(0.4),
                font=FONT_CN, size=Pt(11), bold=False,
                color=C['muted'], align=PP_ALIGN.CENTER)

    add_brand_corner(slide)
    return slide

def slide_chapter(prs, number, title, subtitle=None):
    """章节页：深色背景+大号数字+标题"""
    slide = new_slide(prs, bg=C['ink'])

    # 左侧大号章节编号
    add_textbox(slide, number,
                Inches(0.8), Inches(2.5), Inches(3), Inches(2),
                font=FONT_EN, size=Pt(120), bold=True,
                color=RGBColor(40, 50, 65), align=PP_ALIGN.LEFT)

    # 金色竖线
    fill_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(4.2), Inches(2.8), Pt(3), Inches(2.8),
               C['warm'])

    # 标题
    add_textbox(slide, title,
                Inches(4.6), Inches(3.0), Inches(10), Inches(1.2),
                font=FONT_CN, size=Pt(38), bold=True,
                color=C['cream'], align=PP_ALIGN.LEFT)

    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(4.6), Inches(4.4), Inches(10), Inches(0.5),
                    font=FONT_CN, size=Pt(16), bold=False,
                    color=C['muted'], align=PP_ALIGN.LEFT)

    add_brand_corner(slide)
    return slide

def slide_text_page(prs, page_num, total, chapter, title, body_lines,
                    accent_lines=None, layout='left'):
    """内容页：左文右图 or 左图右文 版式"""
    slide = new_slide(prs, bg=C['light'])

    # 顶部深色条带
    top_bar = fill_shape(slide, MSO_SHAPE.RECTANGLE,
                         0, 0, prs.slide_width, Inches(1.5), C['ink'])

    # 章节标签
    add_textbox(slide, chapter,
                Inches(0.5), Inches(0.55), Inches(3), Inches(0.4),
                font=FONT_CN, size=Pt(9), bold=False,
                color=C['muted'], align=PP_ALIGN.LEFT)

    # 页面标题
    add_textbox(slide, title,
                Inches(0.5), Inches(0.95), Inches(13), Inches(0.45),
                font=FONT_CN, size=Pt(20), bold=True,
                color=C['cream'], align=PP_ALIGN.LEFT)

    # 主体内容区
    content_top = Inches(1.8)
    content_height = Inches(6.5)

    if layout == 'left':
        # 左侧文字区
        content_bg = fill_shape(slide, MSO_SHAPE.RECTANGLE,
                               Inches(0.5), content_top, Inches(7.5), content_height,
                               C['white'])
        # 文字
        tb, tf = add_textbox(slide, "",
                             Inches(0.7), content_top + Inches(0.2),
                             Inches(7.1), content_height - Inches(0.4),
                             font=FONT_CN, size=Pt(13), color=C['slate'])
        # 逐行写入
        for i, line in enumerate(body_lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(0)
            p.space_after = Pt(10)
            r = p.add_run()
            is_accent = accent_lines and i in accent_lines
            apply_font(r, FONT_CN, Pt(13), is_accent,
                      C['ink'] if is_accent else C['slate'])
            r.text = line

        # 右侧色块装饰
        accent_bg = fill_shape(slide, MSO_SHAPE.RECTANGLE,
                              Inches(8.3), content_top, Inches(7.2), content_height,
                              C['ink'])
        # 装饰性大号数字
        add_textbox(slide, f"{page_num:02d}",
                    Inches(8.5), Inches(2.5), Inches(6.8), Inches(3),
                    font=FONT_EN, size=Pt(160), bold=True,
                    color=RGBColor(35, 45, 58), align=PP_ALIGN.RIGHT)
        # 装饰横线
        fill_shape(slide, MSO_SHAPE.RECTANGLE,
                   Inches(8.5), Inches(5.8), Inches(4), Pt(2),
                   C['warm'])
    else:
        # 右图左文：左侧色块+右侧文字
        accent_bg = fill_shape(slide, MSO_SHAPE.RECTANGLE,
                              Inches(0.5), content_top, Inches(7.5), content_height,
                              C['ink'])
        add_textbox(slide, f"{page_num:02d}",
                    Inches(0.8), Inches(2.5), Inches(6.8), Inches(3),
                    font=FONT_EN, size=Pt(160), bold=True,
                    color=RGBColor(35, 45, 58), align=PP_ALIGN.LEFT)
        fill_shape(slide, MSO_SHAPE.RECTANGLE,
                   Inches(4), Inches(5.8), Inches(4), Pt(2),
                   C['warm'])

        content_bg = fill_shape(slide, MSO_SHAPE.RECTANGLE,
                               Inches(8.3), content_top, Inches(7.2), content_height,
                               C['white'])
        tb, tf = add_textbox(slide, "",
                             Inches(8.5), content_top + Inches(0.2),
                             Inches(6.8), content_height - Inches(0.4),
                             font=FONT_CN, size=Pt(13), color=C['slate'])
        for i, line in enumerate(body_lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(0)
            p.space_after = Pt(10)
            r = p.add_run()
            is_accent = accent_lines and i in accent_lines
            apply_font(r, FONT_CN, Pt(13), is_accent,
                      C['ink'] if is_accent else C['slate'])
            r.text = line

    add_page_number(slide, page_num, total)
    return slide

def slide_image_full(prs, page_num, total, space_name, image_path):
    """效果页：沉浸式全屏图+底部半透明标题条"""
    slide = new_slide(prs, bg=C['deep'])

    if image_path and os.path.exists(image_path):
        try:
            # 全屏图片
            pic = slide.shapes.add_picture(image_path, 0, 0,
                                           width=prs.slide_width, height=prs.slide_height)
            # 底部渐变遮罩
            mask = fill_shape(slide, MSO_SHAPE.RECTANGLE,
                             0, Inches(5.8), prs.slide_width, Inches(3.2),
                             RGBColor(10, 12, 18))
            # 标题条
            bar = fill_shape(slide, MSO_SHAPE.RECTANGLE,
                           0, Inches(7.6), prs.slide_width, Inches(1.4),
                           RGBColor(18, 22, 28))
            # 金色装饰线
            fill_shape(slide, MSO_SHAPE.RECTANGLE,
                       Inches(0.5), Inches(7.7), Inches(0.5), Pt(2),
                       C['warm'])
            # 空间名称
            add_textbox(slide, space_name,
                       Inches(1.2), Inches(7.75), Inches(8), Inches(0.5),
                       font=FONT_CN, size=Pt(18), bold=True,
                       color=C['cream'], align=PP_ALIGN.LEFT)
            # 页码
            add_textbox(slide, f"{page_num:02d} / {total:02d}",
                       Inches(13.5), Inches(8.2), Inches(2), Inches(0.4),
                       font=FONT_EN, size=Pt(10), bold=False,
                       color=C['muted'], align=PP_ALIGN.RIGHT)
        except Exception as e:
            print(f"⚠️ 插入图片失败: {e}")
            # 回退到纯色页
            add_textbox(slide, space_name,
                       Inches(2), Inches(4), Inches(12), Inches(1),
                       font=FONT_CN, size=Pt(32), bold=True,
                       color=C['cream'], align=PP_ALIGN.CENTER)
    else:
        add_textbox(slide, space_name,
                   Inches(2), Inches(4), Inches(12), Inches(1),
                   font=FONT_CN, size=Pt(32), bold=True,
                   color=C['cream'], align=PP_ALIGN.CENTER)

    add_brand_corner(slide)
    return slide

def slide_grid_cards(prs, page_num, total, chapter, title, cards):
    """卡片网格页：2×2 或 1×2 色块卡片"""
    slide = new_slide(prs, bg=C['light'])

    # 顶部深色条
    fill_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.5), C['ink'])
    add_textbox(slide, chapter, Inches(0.5), Inches(0.55), Inches(3), Inches(0.4),
                font=FONT_CN, size=Pt(9), color=C['muted'])
    add_textbox(slide, title, Inches(0.5), Inches(0.95), Inches(13), Inches(0.45),
                font=FONT_CN, size=Pt(20), bold=True, color=C['cream'])

    # 卡片布局
    card_w = Inches(7.2)
    card_h = Inches(3.0)
    gap = Inches(0.35)
    start_x = Inches(0.5)
    start_y = Inches(1.8)

    card_colors = [C['ink'], C['slate'], C['ink'], C['slate']]
    accent_colors = [C['warm'], C['accent'], C['accent'], C['warm']]

    for i, (card_title, card_body) in enumerate(cards[:4]):
        col = i % 2
        row = i // 2
        left = start_x + col * (card_w + gap)
        top = start_y + row * (card_h + gap)

        # 卡片背景
        fill_shape(slide, MSO_SHAPE.RECTANGLE,
                  left, top, card_w, card_h, card_colors[i])

        # 金色顶部细线
        fill_shape(slide, MSO_SHAPE.RECTANGLE,
                  left, top, card_w, Pt(3), accent_colors[i])

        # 卡片标题
        add_textbox(slide, card_title,
                   left + Inches(0.3), top + Inches(0.25),
                   card_w - Inches(0.6), Inches(0.45),
                   font=FONT_CN, size=Pt(15), bold=True, color=C['cream'])

        # 卡片内容
        tb, tf = add_textbox(slide, "",
                            left + Inches(0.3), top + Inches(0.8),
                            card_w - Inches(0.6), card_h - Inches(1.1),
                            font=FONT_CN, size=Pt(12), color=C['cream'])
        lines = card_body if isinstance(card_body, list) else [card_body]
        for j, line in enumerate(lines):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.space_before = Pt(0)
            p.space_after = Pt(6)
            r = p.add_run()
            apply_font(r, FONT_CN, Pt(12), False, RGBColor(190, 200, 215))
            r.text = f"· {line}" if not line.startswith("·") else line

    add_page_number(slide, page_num, total)
    return slide

def slide_material_board(prs, page_num, total, materials):
    """材质板页：色块+文字标签"""
    slide = new_slide(prs, bg=C['deep'])

    # 标题区
    fill_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.5), C['ink'])
    add_textbox(slide, "06 材质说明", Inches(0.5), Inches(0.55), Inches(3), Inches(0.4),
                font=FONT_CN, size=Pt(9), color=C['muted'])
    add_textbox(slide, "MATERIAL BOARD", Inches(0.5), Inches(0.95), Inches(10), Inches(0.45),
                font=FONT_EN, size=Pt(20), bold=True, color=C['cream'])
    add_textbox(slide, "主材选型与色彩体系", Inches(10.5), Inches(0.95), Inches(5), Inches(0.45),
                font=FONT_CN, size=Pt(12), color=C['warm'], align=PP_ALIGN.RIGHT)

    # 材质色块
    colors = [
        RGBColor(210, 195, 170), RGBColor(160, 140, 115),
        RGBColor(235, 230, 220), RGBColor(180, 165, 145),
        RGBColor(245, 240, 230), RGBColor(130, 120, 105),
        RGBColor(200, 185, 160), RGBColor(150, 135, 115),
    ]
    swatch_w = Inches(3.6)
    swatch_h = Inches(1.8)
    gap = Inches(0.25)
    start_x = Inches(0.5)
    start_y = Inches(1.8)

    mat_names = [
        "地面  750×1500mm 仿古砖", "墙面  淡灰色乳胶漆",
        "木作  原木色饰面板", "石材  雅士白门槛石",
        "厨卫  300×600mm墙砖", "门窗  黑色窄边铝合金",
        "涂料  低VOC净味墙漆", "灯具  防眩光深杯射灯",
    ]

    for i, (name, rgb) in enumerate(zip(mat_names, colors)):
        col = i % 4
        row = i // 4
        left = start_x + col * (swatch_w + gap)
        top = start_y + row * (swatch_h + gap + Inches(0.5))

        # 色块
        fill_shape(slide, MSO_SHAPE.RECTANGLE,
                  left, top, swatch_w, swatch_h, rgb)
        # 标签
        fill_shape(slide, MSO_SHAPE.RECTANGLE,
                  left, top + swatch_h, swatch_w, Inches(0.45), C['ink'])
        add_textbox(slide, name,
                   left + Inches(0.15), top + swatch_h + Inches(0.05),
                   swatch_w - Inches(0.3), Inches(0.35),
                   font=FONT_CN, size=Pt(10), color=C['cream'])

    # 右侧配色说明
    fill_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0.5), Inches(5.6), Inches(15), Inches(2.6), C['ink'])
    add_textbox(slide, "色彩比例",
               Inches(0.8), Inches(5.8), Inches(2), Inches(0.4),
               font=FONT_CN, size=Pt(11), bold=True, color=C['warm'])
    palette = [
        ("主色调", "40%", "视觉基础", C['slate']),
        ("辅色调", "30%", "层次过渡", C['muted']),
        ("点缀色", "30%", "视觉焦点", C['warm']),
    ]
    for j, (name, pct, desc, color) in enumerate(palette):
        x = Inches(0.8) + j * Inches(5)
        fill_shape(slide, MSO_SHAPE.RECTANGLE,
                  x, Inches(6.4), Inches(0.3), Inches(0.3), color)
        add_textbox(slide, f"{name}  {pct}",
                   x + Inches(0.45), Inches(6.35), Inches(2), Inches(0.35),
                   font=FONT_CN, size=Pt(12), bold=True, color=C['cream'])
        add_textbox(slide, desc,
                   x + Inches(0.45), Inches(6.7), Inches(2), Inches(0.3),
                   font=FONT_CN, size=Pt(10), color=C['muted'])

    add_page_number(slide, page_num, total)
    return slide

def slide_closing(prs):
    """封底：简洁深色+品牌信息"""
    slide = new_slide(prs, bg=C['deep'])

    # 金色装饰线
    fill_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(6), Inches(3.5), Inches(4), Pt(2), C['warm'])

    add_textbox(slide, "感谢聆听",
               Inches(1), Inches(3.8), Inches(14), Inches(1.0),
               font=FONT_CN, size=Pt(44), bold=True,
               color=C['cream'], align=PP_ALIGN.CENTER)

    add_textbox(slide, "期待项目落地",
               Inches(1), Inches(4.9), Inches(14), Inches(0.6),
               font=FONT_CN, size=Pt(22), bold=False,
               color=C['muted'], align=PP_ALIGN.CENTER)

    fill_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(6), Inches(5.7), Inches(4), Pt(1.5), C['warm'])

    add_textbox(slide, "温州隐室空间设计  |  老冷",
               Inches(1), Inches(6.4), Inches(14), Inches(0.4),
               font=FONT_CN, size=Pt(13), bold=False,
               color=C['muted'], align=PP_ALIGN.CENTER)

    add_brand_corner(slide)
    return slide

# ─── 主程序 ──────────────────────────────────────────────────

DEFAULT_SPACES = [
    "客厅", "餐厅", "主卧", "次卧", "书房", "厨房",
    "主卫", "公卫", "阳台", "茶室", "衣帽间", "玄关",
    "多功能厅", "楼梯间", "储藏间"
]

def get_image_files(image_dir):
    exts = ('.png', '.jpg', '.jpeg', '.webp', '.bmp')
    files = sorted([f for f in os.listdir(image_dir)
                    if f.lower().endswith(exts)])
    valid = []
    for f in files:
        try:
            with Image.open(os.path.join(image_dir, f)) as img:
                img.verify()
            valid.append(f)
        except Exception:
            pass
    return valid

def main():
    if len(sys.argv) < 3:
        print("用法: python3 generate_ppt.py <图片文件夹> <输出.pptx> [项目名] [风格] [面积]")
        sys.exit(1)

    image_dir = sys.argv[1]
    output_path = sys.argv[2]
    project = sys.argv[3] if len(sys.argv) > 3 else "设计方案"
    style = sys.argv[4] if len(sys.argv) > 4 else "现代简约"
    area = sys.argv[5] if len(sys.argv) > 5 else ""

    images = get_image_files(image_dir)
    print(f"📁 找到 {len(images)} 张效果图")

    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    total_pages = 6 + len(images)  # 封面+章节页×2+材质+封底+效果图
    if total_pages < 10:
        total_pages = 10

    # ── 1. 封面 ──────────────────────────────────────────────
    cover_img = os.path.join(image_dir, images[0]) if images else None
    slide_cover(prs, project, style,
                f"{area}㎡" if area else "")

    # ── 2. 任务回顾 ─────────────────────────────────────────
    slide_chapter(prs, "01", "任务回顾", "项目需求与设计目标")

    body = []
    if area:
        body.append(f"建筑面积：{area} ㎡")
    body.append(f"风格定位：{style}")
    body.append("核心目标：")
    body.append("满足功能使用需求")
    body.append("打造符合业主审美品味的空间氛围")
    body.append("兼顾品质与预算的平衡")
    body.append("确保施工可落地性")
    slide_text_page(prs, 3, total_pages, "01 任务回顾",
                    "客户需求 · 设计目标", body, {2, 3}, layout='left')

    # ── 3. 设计概念 ─────────────────────────────────────────
    slide_chapter(prs, "02", "设计概念", "灵感来源与方案推导")

    concept_body = [
        "设计风格：",
        f"提取{style}风格核心元素进行设计表达",
        f"结合业主生活习惯与功能需求进行重组优化",
        "用材质、色彩与光影共同营造空间氛围",
        "",
        "概念推导：",
        "从功能布局出发，层层推进，形成完整设计方案",
        "注重空间的叙事性与居住者的情感体验",
        "在满足实用性的基础上追求精神层面的满足",
    ]
    slide_text_page(prs, 5, total_pages, "02 设计概念",
                    "灵感来源 · 概念推导", concept_body,
                    {0, 5}, layout='right')

    # ── 4. 效果展示 ─────────────────────────────────────────
    effect_images = images[1:] if len(images) > 1 else images
    for i, img_name in enumerate(effect_images):
        idx = 6 + i
        space = DEFAULT_SPACES[i] if i < len(DEFAULT_SPACES) else f"空间{i+1}"
        img_path = os.path.join(image_dir, img_name)
        slide_image_full(prs, idx, total_pages, space, img_path)

    # ── 5. 材质板 ───────────────────────────────────────────
    slide_material_board(prs, total_pages - 2, total_pages, [])

    # ── 6. 封底 ─────────────────────────────────────────────
    slide_closing(prs)

    prs.save(output_path)
    print(f"✅ PPT生成完成！")
    print(f"📁 {output_path}")
    print(f"📄 共 {len(prs.slides)} 页")

if __name__ == "__main__":
    main()
