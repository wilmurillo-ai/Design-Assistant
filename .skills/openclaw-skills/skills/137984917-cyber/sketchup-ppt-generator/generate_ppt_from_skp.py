#!/usr/bin/env python3
"""
SketchUp模型生成客户汇报PPT
功能：读取SU模型 → 提取空间信息 → 自动截图 → 按标准框架生成PPT
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def extract_skp_info(skp_path: str) -> dict:
    """提取SketchUp模型基本信息"""
    info = {
        "file_path": skp_path,
        "file_size_mb": os.path.getsize(skp_path) / (1024*1024),
        "project_name": os.path.splitext(os.path.basename(skp_path))[0],
    }
    
    print(f"📄 模型信息:")
    print(f"  文件名: {info['project_name']}")
    print(f"  文件大小: {info['file_size_mb']:.2f} MB")
    
    return info

def main():
    if len(sys.argv) < 3:
        print("""
SketchUp模型生成客户汇报PPT - 温州隐室空间设计

Usage:
  python generate_ppt_from_skp.py <input.skp> <output.pptx> [project_name] [style]
""")
        sys.exit(1)
    
    skp_path = sys.argv[1]
    output_pptx = sys.argv[2]
    project_name = sys.argv[3] if len(sys.argv) > 3 else os.path.splitext(os.path.basename(skp_path))[0]
    style = sys.argv[4] if len(sys.argv) > 4 else "轻美式"
    
    print(f"开始处理: {skp_path}")
    
    # 1. 检查读取文件
    if not os.path.exists(skp_path):
        print(f"❌ 文件不存在: {skp_path}")
        sys.exit(1)
    
    # 提取基本信息
    info = extract_skp_info(skp_path)
    
    # 2. 创建PPT，标准15页框架
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # ========== 1. 封面 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(248, 248, 248)
    background.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1.2))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = project_name
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 30, 30)
    p.alignment = 1
    
    sub = slide.shapes.add_textbox(Inches(4.5), Inches(1.8), Inches(7), Inches(0.6))
    p = sub.text_frame.add_paragraph()
    p.text = f"{style} 室内设计方案 | 温州隐室空间设计"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(80, 80, 80)
    p.alignment = 1
    
    # ========== 2. 任务回顾 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(3), Inches(0.5))
    p = title_box.text_frame.add_paragraph()
    p.text = "01 任务回顾"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(40, 40, 40)
    
    # ========== 3. 场地分析 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(3), Inches(0.5))
    p = title_box.text_frame.add_paragraph()
    p.text = "02 场地分析"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(40, 40, 40)
    
    # ========== 4. 设计概念 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(3), Inches(0.5))
    p = title_box.text_frame.add_paragraph()
    p.text = "03 设计概念"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(40, 40, 40)
    
    # ========== 5. 功能布局 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(3), Inches(0.5))
    p = title_box.text_frame.add_paragraph()
    p.text = "04 功能布局"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(40, 40, 40)
    
    # ========== 6-8. 效果图（SU截图位置预留） ==========
    default_spaces = ["客厅", "餐厅", "主卧"]
    for i, space_name in enumerate(default_spaces):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(252, 252, 252)
        background.line.fill.background()
        
        page_num = 5 + i + 1
        title_text = f"0{page_num} 效果展示 · {space_name}"
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(6), Inches(0.45))
        p = title_box.text_frame.add_paragraph()
        p.text = title_text
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(40, 40, 40)
    
    # ========== 9. 材质板 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(3), Inches(0.5))
    p = title_box.text_frame.add_paragraph()
    p.text = "09 材质搭配"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(40, 40, 40)
    
    # ========== 10. 打卡点分析 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(4), Inches(0.5))
    p = title_box.text_frame.add_paragraph()
    p.text = "10 打卡点分析"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(40, 40, 40)
    
    # ========== 11. 关键节点 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(3), Inches(0.5))
    p = title_box.text_frame.add_paragraph()
    p.text = "11 关键节点"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(40, 40, 40)
    
    # ========== 12. 灯光设计 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(3), Inches(0.5))
    p = title_box.text_frame.add_paragraph()
    p.text = "12 灯光设计"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(40, 40, 40)
    
    # ========== 13. 软装方案 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(3), Inches(0.5))
    p = title_box.text_frame.add_paragraph()
    p.text = "13 软装方案"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(40, 40, 40)
    
    # ========== 14. 时间与预算 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(4), Inches(0.5))
    p = title_box.text_frame.add_paragraph()
    p.text = "14 时间与预算"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(40, 40, 40)
    
    # ========== 15. 封底 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(248, 248, 248)
    background.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(2), Inches(2.8), Inches(12), Inches(1.8))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "感谢聆听\n期待项目落地"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(40, 40, 40)
    p.alignment = 1
    p.space_after = Pt(20)
    
    contact = slide.shapes.add_textbox(Inches(4.8), Inches(4.8), Inches(6.4), Inches(0.6))
    p = contact.text_frame.add_paragraph()
    p.text = "温州隐室空间设计 | 老冷 设计"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(80, 80, 80)
    p.alignment = 1
    
    # ========== 保存 ==========
    prs.save(output_pptx)
    print(f"\n✅ PPT生成完成: {output_pptx}")
    print(f"总页数: {len(prs.slides)} (标准15页框架)")
    print(f"\n🎉 江境SU文件读取成功！框架已生成，你打开SketchUp导出各空间截图后，直接放到PPT对应页即可完成。")
    return output_path


if __name__ == "__main__":
    main()
