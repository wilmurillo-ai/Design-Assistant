#!/usr/bin/env python3
"""
室内设计方案提案PPT自动生成器 v4.1 - 隐室空间设计
====================================================
修复内容：
- 所有布局严格控制在16寸页面内（左右各0.5寸留白，内容区15寸）
- 修复timeline/home_scenes等4栏布局溢出问题
- 深化商业/住宅两套内容质量
"""

import os, sys, argparse
from datetime import datetime
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE

# ─── 页面常量 ─────────────────────────────────────────────────
PW = Inches(16)    # 页面宽度
PH = Inches(9)     # 页面高度
ML = Inches(0.5)   # 左边距
MR = Inches(0.5)   # 右边距
CW = PW - ML - MR  # 内容宽度 = 15 inch
CH = PH - Inches(1.5) - Inches(0.5)  # 内容高度(去掉顶栏和底边距)

# ─── 配色 ────────────────────────────────────────────────────
C = dict(
    deep=RGBColor(18,22,28), ink=RGBColor(28,35,45),
    slate=RGBColor(55,65,85), muted=RGBColor(100,115,140),
    warm=RGBColor(220,175,120), accent=RGBColor(195,145,80),
    cream=RGBColor(245,240,230), white=RGBColor(255,255,255),
    light=RGBColor(242,242,245), fog=RGBColor(200,205,215),
    biz=RGBColor(30,50,80), biz_accent=RGBColor(80,180,160),
    biz_warn=RGBColor(220,100,80), home=RGBColor(50,40,35),
    home_accent=RGBColor(180,140,100),
)
FCN = 'PingFang SC'
FEN = 'Avenir Next'

# ─── 工具函数 ────────────────────────────────────────────────
def fr(run, fn=FCN, sz=Pt(12), b=False, c=None, i=False):
    run.font.name = fn
    try: run._element.rPr.rFonts.set(qn('w:eastAsia'), fn)
    except: pass
    run.font.size, run.font.bold, run.font.italic = sz, b, i
    if c: run.font.color.rgb = c

def pp(tf, t='', fn=FCN, sz=Pt(12), b=False, c=None,
      al=PP_ALIGN.LEFT, sb=Pt(0), sa=Pt(6)):
    p = tf.add_paragraph(); p.alignment, p.space_before, p.space_after = al, sb, sa
    if t:
        r = p.add_run(); fr(r, fn, sz, b, c); r.text = t
    return p

def tx(s, t, l, t1, w, h, fn=FCN, sz=Pt(12), b=False, c=None, al=PP_ALIGN.LEFT):
    box = s.shapes.add_textbox(l, t1, w, h)
    tf = box.text_frame; tf.word_wrap = True
    if t:
        p = tf.paragraphs[0]; p.alignment = al
        r = p.add_run(); fr(r, fn, sz, b, c); r.text = t
    return box, tf

def fl(st, s, l, t1, w, h, fc, lc=None):
    sh = s.shapes.add_shape(st, l, t1, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fc
    sh.line.color.rgb = lc if lc else RGBColor(0,0,0)
    sh.line.fill.background() if not lc else None
    return sh

def new(prs, bg=C['light']):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fl(MSO_SHAPE.RECTANGLE, sl, 0, 0, PW, PH, bg)
    return sl

def brand(sl):
    tx(sl, "温州隐室空间设计", Inches(0.4), Inches(0.3), Inches(3), Inches(0.3),
       sz=Pt(8), c=C['muted'])

def pg(sl, n, tot):
    tx(sl, f"{n:02d} / {tot:02d}", Inches(14.5), Inches(8.55),
       Inches(1.0), Inches(0.3), fn=FEN, sz=Pt(9), c=C['muted'], al=PP_ALIGN.RIGHT)

def glv(sl, l, t1, h, col=None):
    col = col or C['warm']
    fl(MSO_SHAPE.RECTANGLE, sl, l, t1, Pt(3), h, col)

def glh(sl, l, t1, w, col=None):
    col = col or C['warm']
    fl(MSO_SHAPE.RECTANGLE, sl, l, t1, w, Pt(2), col)

def topbar(sl, ch, title, acc=None):
    acc = acc or C['warm']
    fl(MSO_SHAPE.RECTANGLE, sl, 0, 0, PW, Inches(1.5), C['ink'])
    tx(sl, ch, Inches(0.5), Inches(0.52), Inches(3), Inches(0.35), sz=Pt(9), c=C['muted'])
    tx(sl, title, Inches(0.5), Inches(0.9), Inches(13), Inches(0.45), sz=Pt(20), b=True, c=C['cream'])
    glh(sl, Inches(0.5), Inches(1.42), Inches(1.5), acc)

# ─── 幻灯片 ─────────────────────────────────────────────────

def cover(prs, proj, style, area, ptype):
    sl = new(prs, bg=C['deep'])
    acc = C['biz_accent'] if ptype=='commercial' else C['warm']
    fl(MSO_SHAPE.RECTANGLE, sl, 0, 0, PW, Inches(5), RGBColor(10,12,18))
    glh(sl, Inches(5.5), Inches(3.3), Inches(5), acc)
    tx(sl, proj, Inches(1), Inches(3.5), Inches(14), Inches(1.2),
       sz=Pt(50), b=True, c=C['cream'], al=PP_ALIGN.CENTER)
    tx(sl, f"{style}  ·  {area}㎡", Inches(1), Inches(4.75), Inches(14), Inches(0.5),
       sz=Pt(20), c=acc, al=PP_ALIGN.CENTER)
    glh(sl, Inches(5.5), Inches(5.45), Inches(5), acc)
    lbl = "商业空间设计提案" if ptype=='commercial' else "住宅空间设计提案"
    tx(sl, lbl, Inches(1), Inches(6.25), Inches(14), Inches(0.4),
       sz=Pt(11), c=C['muted'], al=PP_ALIGN.CENTER)
    tx(sl, "温州隐室空间设计  |  老冷 设计",
       Inches(1), Inches(7.7), Inches(14), Inches(0.4),
       sz=Pt(12), c=C['muted'], al=PP_ALIGN.CENTER)
    tx(sl, datetime.now().strftime('%Y年%m月%d日'),
       Inches(1), Inches(8.2), Inches(14), Inches(0.35),
       sz=Pt(10), c=C['muted'], al=PP_ALIGN.CENTER)
    brand(sl); return sl

def toc(prs, ptype):
    sl = new(prs, bg=C['ink'])
    acc = C['biz_accent'] if ptype=='commercial' else C['warm']
    tx(sl, "CONTENTS", Inches(0.5), Inches(0.5), Inches(4), Inches(0.45),
       fn=FEN, sz=Pt(11), c=C['muted'])
    tx(sl, "目录", Inches(0.5), Inches(0.9), Inches(4), Inches(0.7),
       sz=Pt(32), b=True, c=C['cream'])
    glh(sl, Inches(0.5), Inches(1.7), Inches(2), acc)
    items = [
        ("01","项目定位","投资定位与目标客群","家庭结构与设计愿景"),
        ("02","投资分析","坪效测算与回报周期","家庭场景与生活方式"),
        ("03","消费者心理","行为动线与触点设计","生活美学与情感设计"),
        ("04","空间设计","功能布局与品牌语言","功能分区与动线优化"),
        ("05","效果展示","效果图呈现","效果图呈现"),
        ("06","材质与品牌","主材选型与品牌调性","材质美学与触感设计"),
        ("07","投资预算","三档预算与收益预测","收纳系统与预算方案"),
        ("08","设计流程","时间节点与服务内容","时间节点与服务内容"),
        ("09","感谢聆听","期待合作","期待合作"),
    ]
    for i,(num,t1,sub_c,sub_h) in enumerate(items):
        col = i%2; row=i//2
        lx = ML + col*Inches(7.5); ty = Inches(2.2)+row*Inches(1.55)
        tx(sl, num, lx, ty, Inches(1), Inches(0.6), fn=FEN,
           sz=Pt(34), b=True, c=RGBColor(40,50,65))
        glv(sl, lx+Inches(1.15), ty+Inches(0.05), Inches(0.5), acc)
        sub = sub_c if ptype=='commercial' else sub_h
        tx(sl, t1, lx+Inches(1.3), ty, Inches(4), Inches(0.45),
           sz=Pt(16), b=True, c=C['cream'])
        tx(sl, sub, lx+Inches(1.3), ty+Inches(0.5), Inches(5), Inches(0.35),
           sz=Pt(10), c=C['muted'])
    brand(sl); return sl

# ── 商业：项目定位 ──────────────────────────────────────────
def comm_positioning(prs, n, tot, proj, style, area, budget, target, competitors):
    sl = new(prs, bg=C['light'])
    topbar(sl, "01 项目定位", "投资定位与目标客群", C['biz_accent'])
    # 左：项目信息
    fl(MSO_SHAPE.RECTANGLE, sl, ML, Inches(1.8), Inches(7.0), Inches(6.4), C['white'])
    fl(MSO_SHAPE.RECTANGLE, sl, ML, Inches(1.8), Inches(7.0), Pt(3), C['biz_accent'])
    data = [("项目名称",proj),("建筑面积",f"{area} ㎡"),("风格定位",style),("预算区间",budget or "待确认")]
    for j,(lbl,val) in enumerate(data):
        y = Inches(2.05)+j*Inches(1.5)
        tx(sl, lbl, Inches(0.8), y, Inches(2.2), Inches(0.32), sz=Pt(10), c=C['muted'])
        tx(sl, val, Inches(0.8), y+Inches(0.32), Inches(6.0), Inches(0.5),
           sz=Pt(16), b=True, c=C['ink'])
    # 右：目标客群
    fl(MSO_SHAPE.RECTANGLE, sl, Inches(7.8), Inches(1.8), Inches(7.7), Inches(3.0), C['biz'])
    tx(sl, "目标客群", Inches(8.1), Inches(1.95), Inches(3), Inches(0.38),
       sz=Pt(12), b=True, c=C['biz_accent'])
    glh(sl, Inches(8.1), Inches(2.42), Inches(1.5), C['biz_accent'])
    for j,t in enumerate(target):
        tx(sl, f"· {t}", Inches(8.1), Inches(2.6)+j*Inches(0.6),
           Inches(7.0), Inches(0.38), sz=Pt(12), c=C['cream'])
    # 右下：竞争分析
    fl(MSO_SHAPE.RECTANGLE, sl, Inches(7.8), Inches(5.0), Inches(7.7), Inches(3.2), C['ink'])
    tx(sl, "差异化竞争点", Inches(8.1), Inches(5.15), Inches(4), Inches(0.38),
       sz=Pt(12), b=True, c=C['biz_accent'])
    glh(sl, Inches(8.1), Inches(5.62), Inches(1.5), C['biz_accent'])
    for j,c in enumerate(competitors):
        tx(sl, f"· {c}", Inches(8.1), Inches(5.8)+j*Inches(0.6),
           Inches(7.0), Inches(0.38), sz=Pt(12), c=C['cream'])
    pg(sl, n, tot); brand(sl); return sl

# ── 家装：项目定位 ──────────────────────────────────────────
def home_positioning(prs, n, tot, proj, style, area, family_struct, vision):
    sl = new(prs, bg=C['light'])
    topbar(sl, "01 项目定位", "家庭结构与设计愿景", C['home_accent'])
    fl(MSO_SHAPE.RECTANGLE, sl, ML, Inches(1.8), Inches(7.0), Inches(6.4), C['white'])
    fl(MSO_SHAPE.RECTANGLE, sl, ML, Inches(1.8), Inches(7.0), Pt(3), C['home_accent'])
    data = [("项目名称",proj),("建筑面积",f"{area} ㎡"),
            ("风格定位",style),("家庭结构",family_struct or "待了解")]
    for j,(lbl,val) in enumerate(data):
        y = Inches(2.05)+j*Inches(1.5)
        tx(sl, lbl, Inches(0.8), y, Inches(2.2), Inches(0.32), sz=Pt(10), c=C['muted'])
        tx(sl, val, Inches(0.8), y+Inches(0.32), Inches(6.0), Inches(0.5),
           sz=Pt(16), b=True, c=C['ink'])
    fl(MSO_SHAPE.RECTANGLE, sl, Inches(7.8), Inches(1.8), Inches(7.7), Inches(6.4), C['home'])
    tx(sl, "设计愿景", Inches(8.1), Inches(2.0), Inches(3), Inches(0.38),
       sz=Pt(12), b=True, c=C['home_accent'])
    glh(sl, Inches(8.1), Inches(2.47), Inches(1.5), C['home_accent'])
    vs = vision or ["打造三代同堂的温馨居所","让孩子有独立成长的空间",
                    "让老人日常生活更安全便捷","让夫妻拥有私密放松的角落"]
    for j,v in enumerate(vs):
        tx(sl, f"· {v}", Inches(8.1), Inches(2.65)+j*Inches(1.0),
           Inches(7.0), Inches(0.42), sz=Pt(12), c=C['cream'])
    pg(sl, n, tot); brand(sl); return sl

# ── 商业：投资分析 ──────────────────────────────────────────
def comm_investment(prs, n, tot, area, expected_rev, customer_price):
    sl = new(prs, bg=C['ink'])
    tx(sl, "02 投资分析", Inches(0.5), Inches(0.52), Inches(3), Inches(0.35),
       sz=Pt(9), c=C['muted'])
    tx(sl, "INVESTMENT ANALYSIS", Inches(0.5), Inches(0.9), Inches(10), Inches(0.45),
       fn=FEN, sz=Pt(20), b=True, c=C['cream'])
    tx(sl, "坪效测算与回报周期", Inches(12), Inches(0.9), Inches(3.5), Inches(0.45),
       sz=Pt(12), c=C['biz_accent'], al=PP_ALIGN.RIGHT)
    # 三大KPI
    kw = (CW - Inches(0.6)) / 3  # 每列宽度 ≈ 4.8 inch
    kpis = [
        ("目标坪效", f"¥{expected_rev or '—'}/㎡/月", "高于区域均值30%"),
        ("客单价目标", f"¥{customer_price or '—'}/单", "提升客单价15%以上"),
        ("预期回本周期", "2-3年", "稳健投资模型测算"),
    ]
    for j,(lbl,val,note) in enumerate(kpis):
        x = ML + j*(kw + Inches(0.3))
        fl(MSO_SHAPE.RECTANGLE, sl, x, Inches(1.8), kw, Inches(3.4), C['slate'])
        fl(MSO_SHAPE.RECTANGLE, sl, x, Inches(1.8), kw, Pt(3), C['biz_accent'])
        tx(sl, lbl, x+Inches(0.25), Inches(2.0), kw-Inches(0.5), Inches(0.38),
           sz=Pt(11), c=C['muted'])
        tx(sl, val, x+Inches(0.25), Inches(2.4), kw-Inches(0.5), Inches(0.9),
           fn=FEN, sz=Pt(30), b=True, c=C['cream'])
        glh(sl, x+Inches(0.25), Inches(3.4), Inches(1.5), C['biz_accent'])
        tx(sl, note, x+Inches(0.25), Inches(3.6), kw-Inches(0.5), Inches(0.4),
           sz=Pt(11), c=C['fog'])
    # 回报周期表
    fl(MSO_SHAPE.RECTANGLE, sl, ML, Inches(5.45), CW, Inches(2.75), C['slate'])
    tx(sl, "投资回报测算（按套内面积估算）", Inches(0.75), Inches(5.6), Inches(6), Inches(0.38),
       sz=Pt(12), b=True, c=C['biz_accent'])
    tiers = [
        ("基础档", "1500-2000元/㎡", "预计回本周期3-4年，适合稳健型投资者"),
        ("品质档", "2500-3500元/㎡", "预计回本周期2-3年，平衡品质与投资回报"),
        ("高端档", "4500+元/㎡", "预计回本周期1.5-2年，追求溢价空间"),
    ]
    tw = (CW - Inches(0.6)) / 3
    for j,(tier,price,desc) in enumerate(tiers):
        x = ML + j*(tw + Inches(0.3))
        tx(sl, tier, x, Inches(6.15), tw, Inches(0.4),
           sz=Pt(13), b=True, c=C['cream'])
        tx(sl, price, x, Inches(6.6), tw, Inches(0.4),
           fn=FEN, sz=Pt(18), b=True, c=C['biz_accent'])
        tx(sl, desc, x, Inches(7.05), tw, Inches(0.8),
           sz=Pt(10), c=C['fog'])
    pg(sl, n, tot); brand(sl); return sl

# ── 家装：家庭场景 ──────────────────────────────────────────
def home_scenes(prs, n, tot, members):
    sl = new(prs, bg=C['ink'])
    tx(sl, "02 家庭场景", Inches(0.5), Inches(0.52), Inches(3), Inches(0.35),
       sz=Pt(9), c=C['muted'])
    tx(sl, "FAMILY SCENES", Inches(0.5), Inches(0.9), Inches(10), Inches(0.45),
       fn=FEN, sz=Pt(20), b=True, c=C['cream'])
    tx(sl, "每位成员的生活地图", Inches(12), Inches(0.9), Inches(3.5), Inches(0.45),
       sz=Pt(12), c=C['home_accent'], al=PP_ALIGN.RIGHT)
    # 4卡布局，每卡宽度 = (15 - 3*0.3) / 4 = 3.525 inch
    cw = (CW - Inches(0.9)) / 4  # ≈ 3.525
    scenes = members or [
        ("长辈","安全·便捷",["卫生间扶手+防滑地面","起夜感应灯","沙发坐高适老化","储物弯腰不费力"]),
        ("孩子","成长·陪伴",["独立学习角","安全插座保护","涂鸦墙设计","户外活动动线"]),
        ("主人","私密·放松",["主卧套房设计","独立书房","步入式衣帽间","南向采光主卫"]),
        ("全家","情感·连接",["开放式客餐厅","阳台亲子区","节日团聚餐桌","客厅阅读角"]),
    ]
    for j,(person,theme,details) in enumerate(scenes[:4]):
        x = ML + j*(cw + Inches(0.3))
        fl(MSO_SHAPE.RECTANGLE, sl, x, Inches(1.8), cw, Inches(6.4), C['slate'])
        fl(MSO_SHAPE.RECTANGLE, sl, x, Inches(1.8), cw, Pt(3), C['home_accent'])
        tx(sl, person, x+Inches(0.2), Inches(2.0), cw-Inches(0.4), Inches(0.5),
           sz=Pt(16), b=True, c=C['home_accent'])
        tx(sl, theme, x+Inches(0.2), Inches(2.5), cw-Inches(0.4), Inches(0.38),
           sz=Pt(11), c=C['fog'])
        glh(sl, x+Inches(0.2), Inches(3.0), Inches(1.2), C['home_accent'])
        for k,d in enumerate(details):
            tx(sl, f"· {d}", x+Inches(0.2), Inches(3.2)+k*Inches(1.05),
               cw-Inches(0.4), Inches(0.42), sz=Pt(12), c=C['cream'])
    pg(sl, n, tot); brand(sl); return sl

# ── 共享：空间设计 ──────────────────────────────────────────
def space_design(prs, n, tot, ch, space, points, ptype):
    sl = new(prs, bg=C['light'])
    acc = C['biz_accent'] if ptype=='commercial' else C['home_accent']
    topbar(sl, ch, f"空间设计 · {space}", acc)
    # 左文字区
    fl(MSO_SHAPE.RECTANGLE, sl, ML, Inches(1.8), Inches(7.0), Inches(6.4), C['white'])
    fl(MSO_SHAPE.RECTANGLE, sl, ML, Inches(1.8), Inches(7.0), Pt(3), acc)
    for j,(title,detail) in enumerate(points):
        y = Inches(2.1)+j*Inches(1.9)
        tx(sl, title, Inches(0.8), y, Inches(6.4), Inches(0.4),
           sz=Pt(14), b=True, c=C['ink'])
        tx(sl, detail, Inches(0.8), y+Inches(0.45), Inches(6.4), Inches(0.85),
           sz=Pt(12), c=C['slate'])
    # 右装饰
    fl(MSO_SHAPE.RECTANGLE, sl, Inches(7.8), Inches(1.8), Inches(7.7), Inches(6.4), C['ink'])
    tx(sl, f"{n:02d}", Inches(8.0), Inches(2.2), Inches(7.2), Inches(4),
       fn=FEN, sz=Pt(200), b=True, c=RGBColor(35,45,58), al=PP_ALIGN.RIGHT)
    glh(sl, Inches(8.0), Inches(6.3), Inches(4), acc)
    pg(sl, n, tot); brand(sl); return sl

# ── 共享：效果展示 ──────────────────────────────────────────
def renders(prs, n, tot, images, img_dir, ptype):
    sl = new(prs, bg=C['deep'])
    acc = C['biz_accent'] if ptype=='commercial' else C['home_accent']
    fl(MSO_SHAPE.RECTANGLE, sl, 0, 0, PW, Inches(1.5), C['ink'])
    tx(sl, "04 效果展示", Inches(0.5), Inches(0.52), Inches(3), Inches(0.35),
       sz=Pt(9), c=C['muted'])
    tx(sl, "RENDERINGS", Inches(0.5), Inches(0.9), Inches(10), Inches(0.45),
       fn=FEN, sz=Pt(20), b=True, c=C['cream'])
    tx(sl, "效果图呈现", Inches(12), Inches(0.9), Inches(3.5), Inches(0.45),
       sz=Pt(12), c=acc, al=PP_ALIGN.RIGHT)
    cols = 2; iw = Inches(7.3); ih = Inches(3.5)
    gx = Inches(0.45); gy = Inches(0.3)
    for i,(fname,fpath) in enumerate(images[:4]):
        col = i%cols; row=i//cols
        x = ML + col*(iw+gx); y = Inches(1.8)+row*(ih+gy)
        try: sl.shapes.add_picture(fpath, x, y, width=iw, height=ih)
        except:
            fl(MSO_SHAPE.RECTANGLE, sl, x, y, iw, ih, C['slate'])
            tx(sl, fname, x+Inches(0.2), y+Inches(1.3), iw-Inches(0.4), Inches(0.5),
               sz=Pt(11), c=C['muted'], al=PP_ALIGN.CENTER)
    pg(sl, n, tot); brand(sl); return sl

# ── 共享：材质 ───────────────────────────────────────────────
def materials(prs, n, tot, ptype):
    sl = new(prs, bg=C['deep'])
    acc = C['biz_accent'] if ptype=='commercial' else C['home_accent']
    fl(MSO_SHAPE.RECTANGLE, sl, 0, 0, PW, Inches(1.5), C['ink'])
    tx(sl, "05 材质说明", Inches(0.5), Inches(0.52), Inches(3), Inches(0.35),
       sz=Pt(9), c=C['muted'])
    tx(sl, "MATERIAL BOARD", Inches(0.5), Inches(0.9), Inches(10), Inches(0.45),
       fn=FEN, sz=Pt(20), b=True, c=C['cream'])
    tx(sl, "主材选型与品牌调性", Inches(12), Inches(0.9), Inches(3.5), Inches(0.45),
       sz=Pt(12), c=acc, al=PP_ALIGN.RIGHT)
    if ptype=='commercial':
        swatch = [
            ("地面","750×1500mm 防滑地砖",RGBColor(160,155,145)),
            ("墙面","品牌墙纸/高光板",RGBColor(218,213,203)),
            ("吊顶","轻钢龙骨+石膏板",RGBColor(235,230,220)),
            ("灯光","品牌轨道灯 显指>90",RGBColor(80,85,95)),
            ("吧台","岩板台面+钢结构",RGBColor(50,50,55)),
            ("隔断","磨砂玻璃/金属网",RGBColor(170,170,175)),
            ("外立面","铝板/穿孔板",RGBColor(130,130,135)),
            ("软装","品牌定制家具",RGBColor(195,185,170)),
        ]
    else:
        swatch = [
            ("地面","实木复合地板 温润脚感",RGBColor(195,175,150)),
            ("墙面","蛋壳光乳胶漆 柔和护眼",RGBColor(240,235,225)),
            ("木材","北美白橡木 自然纹理",RGBColor(175,150,120)),
            ("石材","爵士白/雅士白 点缀",RGBColor(240,240,235)),
            ("布艺","棉麻/羊毛 触感温暖",RGBColor(200,190,175)),
            ("皮革","头层牛皮 历久弥新",RGBColor(130,90,70)),
            ("金属","黄铜/拉丝不锈钢",RGBColor(180,175,165)),
            ("植物","散尾葵/龟背竹",RGBColor(100,130,100)),
        ]
    sw = (CW - Inches(0.75)) / 4  # 4列色块
    sh = Inches(1.8); sy = Inches(1.8)
    for i,(name,desc,rgb) in enumerate(swatch):
        col = i%4; row=i//4
        x = ML + col*(sw+Inches(0.25)); y = sy+row*(sh+Inches(0.45))
        fl(MSO_SHAPE.RECTANGLE, sl, x, y, sw, sh, rgb)
        fl(MSO_SHAPE.RECTANGLE, sl, x, y+sh, sw, Inches(0.5), C['ink'])
        tx(sl, name, x+Inches(0.15), y+sh+Inches(0.05),
           sw-Inches(0.3), Inches(0.3), sz=Pt(10), b=True, c=C['cream'])
        tx(sl, desc, x+Inches(0.15), y+sh+Inches(0.32),
           sw-Inches(0.3), Inches(0.45), sz=Pt(8), c=C['fog'])
    # 色彩比例
    fl(MSO_SHAPE.RECTANGLE, sl, ML, Inches(5.85), CW, Inches(2.35), C['slate'])
    tx(sl, "色彩比例", Inches(0.75), Inches(6.05), Inches(2), Inches(0.38),
       sz=Pt(11), b=True, c=acc)
    palette = [("主色调","40%",C['slate']),("辅色调","30%",C['muted']),("点缀色","30%",acc)]
    pw = (CW - Inches(0.6)) / 3
    for j,(name,pct,col) in enumerate(palette):
        x = ML + j*(pw+Inches(0.3))
        fl(MSO_SHAPE.RECTANGLE, sl, x, Inches(6.6), Inches(0.3), Inches(0.3), col)
        tx(sl, f"{name}  {pct}", x+Inches(0.45), Inches(6.55), Inches(2.5), Inches(0.38),
           sz=Pt(13), b=True, c=C['cream'])
        tx(sl, "视觉基础" if j==0 else ("层次过渡" if j==1 else "视觉焦点"),
           x+Inches(0.45), Inches(6.9), Inches(3), Inches(0.3),
           sz=Pt(10), c=C['muted'])
    pg(sl, n, tot); brand(sl); return sl

# ── 商业：消费者心理 ─────────────────────────────────────────
def comm_psychology(prs, n, tot, journey, touch_pts, pain_pts):
    sl = new(prs, bg=C['ink'])
    tx(sl, "03 消费者心理", Inches(0.5), Inches(0.52), Inches(3), Inches(0.35),
       sz=Pt(9), c=C['muted'])
    tx(sl, "CONSUMER PSYCHOLOGY", Inches(0.5), Inches(0.9), Inches(10), Inches(0.45),
       fn=FEN, sz=Pt(20), b=True, c=C['cream'])
    tx(sl, "行为动线与触点设计", Inches(12), Inches(0.9), Inches(3.5), Inches(0.45),
       sz=Pt(12), c=C['biz_accent'], al=PP_ALIGN.RIGHT)
    # 消费者旅程
    fl(MSO_SHAPE.RECTANGLE, sl, ML, Inches(1.8), CW, Inches(2.5), C['slate'])
    tx(sl, "消费者决策旅程", Inches(0.75), Inches(1.95), Inches(3), Inches(0.38),
       sz=Pt(12), b=True, c=C['biz_accent'])
    stages = journey or ["认知","兴趣","欲望","行动","复购","推荐"]
    sw = (CW - Inches(0.6)) / len(stages)  # 动态宽度
    for j,stage in enumerate(stages):
        x = ML + j*(sw+Inches(0.12))
        fl(MSO_SHAPE.RECTANGLE, sl, x, Inches(2.5), sw, Inches(1.5), C['ink'])
        fl(MSO_SHAPE.RECTANGLE, sl, x, Inches(2.5), sw, Pt(2), C['biz_accent'])
        tx(sl, stage, x, Inches(2.85), sw, Inches(0.5),
           sz=Pt(14), b=True, c=C['cream'], al=PP_ALIGN.CENTER)
        if j < len(stages)-1:
            tx(sl, "→", x+sw, Inches(3.0), Inches(0.25), Inches(0.5),
               fn=FEN, sz=Pt(16), b=True, c=C['biz_accent'], al=PP_ALIGN.CENTER)
    # 触点设计
    fl(MSO_SHAPE.RECTANGLE, sl, ML, Inches(4.5), Inches(7.2), Inches(3.7), C['slate'])
    tx(sl, "关键触点设计", Inches(0.75), Inches(4.65), Inches(3), Inches(0.38),
       sz=Pt(12), b=True, c=C['biz_accent'])
    tps = touch_pts or ["入口吸引：橱窗+音乐","等待区：舒适座椅+茶水",
                         "决策区：明档/样品展示","支付区：便捷+积分"]
    for j,tp in enumerate(tps):
        tx(sl, f"· {tp}", Inches(0.75), Inches(5.15)+j*Inches(0.7),
           Inches(6.5), Inches(0.38), sz=Pt(12), c=C['cream'])
    # 痛点规避
    fl(MSO_SHAPE.RECTANGLE, sl, Inches(7.8), Inches(4.5), Inches(7.7), Inches(3.7), C['deep'])
    tx(sl, "心理痛点规避", Inches(8.1), Inches(4.65), Inches(3), Inches(0.38),
       sz=Pt(12), b=True, c=C['biz_accent'])
    pps = pain_pts or ["等位焦虑：提前预约+实时通知","选择困难：精简菜单+套餐推荐",
                        "价格顾虑：透明定价+价值对比","体验平庸：拍照打卡点设计"]
    for j,pp in enumerate(pps):
        tx(sl, f"· {pp}", Inches(8.1), Inches(5.15)+j*Inches(0.7),
           Inches(7.0), Inches(0.38), sz=Pt(12), c=C['cream'])
    pg(sl, n, tot); brand(sl); return sl

# ── 家装：生活美学 ──────────────────────────────────────────
def home_aesthetics(prs, n, tot, aesthetics, light_design, emotion_scenes):
    sl = new(prs, bg=C['ink'])
    tx(sl, "03 生活美学", Inches(0.5), Inches(0.52), Inches(3), Inches(0.35),
       sz=Pt(9), c=C['muted'])
    tx(sl, "AESTHETICS OF LIVING", Inches(0.5), Inches(0.9), Inches(10), Inches(0.45),
       fn=FEN, sz=Pt(20), b=True, c=C['cream'])
    tx(sl, "材质、光影与情感设计", Inches(12), Inches(0.9), Inches(3.5), Inches(0.45),
       sz=Pt(12), c=C['home_accent'], al=PP_ALIGN.RIGHT)
    # 3列支柱
    pw = (CW - Inches(0.6)) / 3  # ≈ 4.8 inch
    pillars = [
        ("材质美学","触感即情感",aesthetics or ["天然木材的温润感","石材的沉稳大气",
                   "织物的柔软亲和","金属的精致细节"]),
        ("光影设计","光是空间的灵魂",light_design or ["自然光：引光入室，借景自然",
                   "氛围光：调光设计，分时分景","重点光：突出焦点，聚焦美感",
                   "夜灯：柔和指引，不扰睡眠"]),
        ("情感场景","家的意义",emotion_scenes or ["晨光中的早餐角","雨天的阅读窗边",
                    "周末的家庭影院","深夜的一盏小灯"]),
    ]
    for j,(title,subtitle,items) in enumerate(pillars):
        x = ML + j*(pw+Inches(0.3))
        fl(MSO_SHAPE.RECTANGLE, sl, x, Inches(1.8), pw, Inches(6.4), C['slate'])
        fl(MSO_SHAPE.RECTANGLE, sl, x, Inches(1.8), pw, Pt(3), C['home_accent'])
        tx(sl, title, x+Inches(0.3), Inches(2.0), pw-Inches(0.6), Inches(0.5),
           sz=Pt(17), b=True, c=C['cream'])
        tx(sl, subtitle, x+Inches(0.3), Inches(2.55), pw-Inches(0.6), Inches(0.38),
           sz=Pt(11), c=C['home_accent'])
        glh(sl, x+Inches(0.3), Inches(3.05), Inches(1.2), C['home_accent'])
        for k,item in enumerate(items):
            tx(sl, f"· {item}", x+Inches(0.3), Inches(3.25)+k*Inches(1.1),
               pw-Inches(0.6), Inches(0.42), sz=Pt(12), c=C['cream'])
    pg(sl, n, tot); brand(sl); return sl

# ── 共享：预算 ───────────────────────────────────────────────
def budget(prs, n, tot, ptype):
    sl = new(prs, bg=C['light'])
    acc = C['biz_accent'] if ptype=='commercial' else C['home_accent']
    topbar(sl, "06 投资预算" if ptype=='commercial' else "07 预算方案",
           "三档预算参考", acc)
    tiers = [
        ("基础档","1500-2000","元/㎡",
         "适合稳健型投资" if ptype=='commercial' else "适合刚需改善",
         ["国内品牌主材","标准施工工艺",
          "基础软装" if ptype=='commercial' else "实用家具配置",
          "基础保修" if ptype=='commercial' else "一年质保"]),
        ("品质档","2500-3500","元/㎡",
         "平衡品质与回报" if ptype=='commercial' else "适合改善型需求",
         ["国内一线+部分进口","精细化施工工艺",
          "设计师全程跟踪" if ptype=='commercial' else "全屋定制",
          "两年质保" if ptype=='commercial' else "两年质保"]),
        ("豪华档","4500+","元/㎡",
         "追求溢价空间" if ptype=='commercial' else "适合终极改善",
         ["进口品牌/全屋定制","顶级工艺+管家服务",
          "拎包入住" if ptype=='commercial' else "全案设计",
          "五年质保" if ptype=='commercial' else "五年质保"]),
    ]
    tw = (CW - Inches(0.6)) / 3
    tc = [C['slate'],C['ink'],C['deep']]
    for j,(name,price,unit,tagline,feats) in enumerate(tiers):
        x = ML + j*(tw+Inches(0.3))
        fl(MSO_SHAPE.RECTANGLE, sl, x, Inches(1.8), tw, Inches(6.4), tc[j])
        fl(MSO_SHAPE.RECTANGLE, sl, x, Inches(1.8), tw, Pt(3), acc)
        tx(sl, name, x+Inches(0.3), Inches(2.0), tw-Inches(0.6), Inches(0.45),
           sz=Pt(16), b=True, c=acc)
        tx(sl, price, x+Inches(0.3), Inches(2.5), tw-Inches(0.6), Inches(0.75),
           fn=FEN, sz=Pt(32), b=True, c=C['cream'])
        tx(sl, unit, x+Inches(0.3), Inches(3.3), tw-Inches(0.6), Inches(0.35),
           sz=Pt(12), c=C['muted'])
        tx(sl, tagline, x+Inches(0.3), Inches(3.75), tw-Inches(0.6), Inches(0.45),
           sz=Pt(11), c=C['fog'])
        for k,feat in enumerate(feats):
            tx(sl, f"· {feat}", x+Inches(0.3), Inches(4.35)+k*Inches(0.75),
               tw-Inches(0.6), Inches(0.42), sz=Pt(12), c=C['cream'])
    pg(sl, n, tot); brand(sl); return sl

# ── 共享：设计流程 ──────────────────────────────────────────
def timeline(prs, n, tot, ptype):
    sl = new(prs, bg=C['ink'])
    acc = C['biz_accent'] if ptype=='commercial' else C['home_accent']
    tx(sl, "07 设计流程" if ptype=='commercial' else "08 设计流程",
       Inches(0.5), Inches(0.52), Inches(3), Inches(0.35),
       sz=Pt(9), c=C['muted'])
    tx(sl, "PROJECT TIMELINE", Inches(0.5), Inches(0.9), Inches(10), Inches(0.45),
       fn=FEN, sz=Pt(20), b=True, c=C['cream'])
    tx(sl, "时间节点与推进计划", Inches(12), Inches(0.9), Inches(3.5), Inches(0.45),
       sz=Pt(12), c=acc, al=PP_ALIGN.RIGHT)
    if ptype=='commercial':
        steps = [
            ("01","品牌定位","5-7天","品牌梳理→定位分析\n客群画像→风格方向"),
            ("02","方案设计","10-15天","平面方案→效果图\n客户确认→方案定稿"),
            ("03","施工图深化","10-15天","水电点位→立面图\n材料表→预算清单"),
            ("04","施工落地","45-90天","拆除→基础→安装\n消防/空调/灯光"),
        ]
    else:
        steps = [
            ("01","需求沟通","5-7天","家庭结构→生活方式\n收纳需求→风格偏好"),
            ("02","方案设计","10-15天","平面布局→效果预览\n设计说明→确认签字"),
            ("03","施工图","10-15天","水电点位→吊顶图\n立面尺寸→材料清单"),
            ("04","施工落地","60-90天","拆改→水电→泥木\n油漆→安装→保洁"),
            ("05","软装摆场","7-10天","家具采购→窗帘灯饰\n绿植装饰→入住摄影"),
        ]
    nsteps = len(steps)
    sw = (CW - Inches(0.5) - (nsteps-1)*Inches(0.3)) / nsteps  # 动态宽度
    sy = Inches(2.0)
    sh = PH - sy - Inches(0.5)
    for j,(num,title,days,detail) in enumerate(steps):
        x = ML + j*(sw+Inches(0.3))
        fl(MSO_SHAPE.RECTANGLE, sl, x, sy, sw, sh, C['slate'])
        fl(MSO_SHAPE.RECTANGLE, sl, x, sy, sw, Pt(3), acc)
        tx(sl, num, x+Inches(0.2), sy+Inches(0.15), Inches(1.5), Inches(0.7),
           fn=FEN, sz=Pt(34), b=True, c=acc)
        tx(sl, title, x+Inches(0.2), sy+Inches(1.0), sw-Inches(0.4), Inches(0.5),
           sz=Pt(14), b=True, c=C['cream'])
        tx(sl, days, x+Inches(0.2), sy+Inches(1.55), sw-Inches(0.4), Inches(0.38),
           sz=Pt(12), c=acc)
        glh(sl, x+Inches(0.2), sy+Inches(2.1), Inches(1.5), acc)
        tx(sl, detail, x+Inches(0.2), sy+Inches(2.3), sw-Inches(0.4), Inches(1.3),
           sz=Pt(11), c=C['fog'])
        if j < nsteps-1:
            tx(sl, "→", x+sw+Inches(0.05), sy+Inches(2.3), Inches(0.25), Inches(0.5),
               fn=FEN, sz=Pt(22), b=True, c=acc, al=PP_ALIGN.CENTER)
    pg(sl, n, tot); brand(sl); return sl

# ── 家装：收纳系统 ──────────────────────────────────────────
def home_storage(prs, n, tot):
    sl = new(prs, bg=C['light'])
    topbar(sl, "06 收纳系统", "全屋收纳规划", C['home_accent'])
    storage = [
        ("玄关","鞋柜+换鞋凳+挂衣+快递桌","容量按家庭人口×20双计算"),
        ("客厅","电视柜+茶几收纳+阳台柜","隐藏式收纳为主，视觉整洁"),
        ("厨房","地柜+吊柜+高柜+转角柜","抽屉式拿取比层板方便30%"),
        ("卧室","衣柜+床头柜+床底收纳","挂衣:叠放=7:3比例最实用"),
        ("卫生间","镜柜+浴室柜+壁龛","干区收纳最大化，湿区防潮"),
        ("阳台","家政柜+洗衣台+储物架","家政工具集中管理，效率提升"),
    ]
    # 3列×2行布局
    cw = (CW - Inches(0.6)) / 3  # ≈ 4.8 inch
    ch = Inches(3.5)
    for j,(space,content,tip) in enumerate(storage):
        col = j%3; row=j//3
        x = ML + col*(cw+Inches(0.3)); y = Inches(1.8)+row*(ch+Inches(0.25))
        fl(MSO_SHAPE.RECTANGLE, sl, x, y, cw, ch, C['home'])
        fl(MSO_SHAPE.RECTANGLE, sl, x, y, cw, Pt(3), C['home_accent'])
        tx(sl, space, x+Inches(0.3), y+Inches(0.2), cw-Inches(0.6), Inches(0.45),
           sz=Pt(16), b=True, c=C['cream'])
        tx(sl, content, x+Inches(0.3), y+Inches(0.75), cw-Inches(0.6), Inches(1.2),
           sz=Pt(11), c=C['fog'])
        tx(sl, f"💡 {tip}", x+Inches(0.3), y+Inches(2.05), cw-Inches(0.6), Inches(0.8),
           sz=Pt(10), c=C['home_accent'])
    pg(sl, n, tot); brand(sl); return sl

# ── 共享：封底 ───────────────────────────────────────────────
def closing(prs, proj, ptype, company="温州隐室空间设计", designer="老冷"):
    sl = new(prs, bg=C['deep'])
    acc = C['biz_accent'] if ptype=='commercial' else C['warm']
    fl(MSO_SHAPE.RECTANGLE, sl, 0, 0, PW, Inches(5), RGBColor(10,12,18))
    glh(sl, Inches(6), Inches(3.3), Inches(4), acc)
    tx(sl, "感谢聆听", Inches(1), Inches(3.5), Inches(14), Inches(1.0),
       sz=Pt(44), b=True, c=C['cream'], al=PP_ALIGN.CENTER)
    suffix = "打造高回报商业空间" if ptype=='commercial' else "期待为您的家创造美好生活"
    tx(sl, f"期待 {suffix}", Inches(1), Inches(4.65), Inches(14), Inches(0.5),
       sz=Pt(18), c=C['muted'], al=PP_ALIGN.CENTER)
    glh(sl, Inches(6), Inches(5.35), Inches(4), acc)
    tx(sl, f"{company}  |  {designer}", Inches(1), Inches(5.9), Inches(14), Inches(0.4),
       sz=Pt(13), c=C['muted'], al=PP_ALIGN.CENTER)
    brand(sl); return sl

# ─── 主程序 ──────────────────────────────────────────────────
DEFAULT_SPACES_C = ["接待区","展示区","用餐区"]
DEFAULT_SPACES_H = ["客厅","餐厅","主卧"]

def get_images(img_dir):
    exts = ('.png','.jpg','.jpeg','.webp','.bmp')
    files = sorted([f for f in os.listdir(img_dir) if f.lower().endswith(exts)])
    valid = []
    for f in files:
        try:
            with Image.open(os.path.join(img_dir,f)) as img: img.verify()
            valid.append((f, os.path.join(img_dir,f)))
        except: pass
    return valid

def main():
    p = argparse.ArgumentParser(description='室内设计方案提案PPT生成器 v4.1')
    p.add_argument('--type','-t',required=True,choices=['commercial','residential'])
    p.add_argument('--project','-p',required=True)
    p.add_argument('--area','-a',required=True)
    p.add_argument('--style','-s',required=True)
    p.add_argument('--output','-o',required=True)
    p.add_argument('--images','-i',default='')
    p.add_argument('--budget','-b',default='')
    p.add_argument('--spaces','-sp',default='')
    p.add_argument('--company','-c',default='温州隐室空间设计')
    p.add_argument('--designer','-d',default='老冷')
    args = p.parse_args()

    os.makedirs(os.path.dirname(args.output) or '.', exist_ok=True)
    ptype = args.type
    images = get_images(args.images) if args.images and os.path.isdir(args.images) else []
    print(f"📁 {len(images)}张图 | 类型: {ptype}")

    prs = Presentation()
    prs.slide_width = PW; prs.slide_height = PH

    # 1 封面
    cover(prs, args.project, args.style, f"{args.area}㎡", ptype)
    # 2 目录
    toc(prs, ptype)

    if ptype=='commercial':
        total = 12
        comm_positioning(prs, 3, total, args.project, args.style, args.area, args.budget,
                         target=["高净值消费客群","年轻家庭客群","商务宴请客群"],
                         competitors=["品牌差异化定位","空间体验独特性","服务溢价能力"])
        comm_investment(prs, 4, total, args.area, "280", "380")
        comm_psychology(prs, 5, total,
                        journey=["认知","兴趣","欲望","行动","复购","推荐"],
                        touch_pts=["入口吸引：橱窗+音乐","等待区：舒适座椅+茶水",
                                   "决策区：明档/样品展示","支付区：便捷+积分"],
                        pain_pts=["等位焦虑：提前预约+实时通知","选择困难：精简菜单+套餐推荐",
                                  "价格顾虑：透明定价+价值对比","体验平庸：拍照打卡点设计"])
        spaces = [s.strip() for s in args.spaces.split(',')] if args.spaces else DEFAULT_SPACES_C
        space_pts = [
            [("功能布局","接待→展示→体验→消费的完整动线设计"),
             ("坪效优化","高频使用区靠窗，低频区退居内侧，最大化商业价值")],
            [("空间叙事","从进入→停留→消费的情感递进设计"),
             ("打卡点","设置2-3个拍照打卡点，驱动社交传播")],
            [("座位设计","不同场景的座位配比：单人/双人/多人灵活配置"),
             ("动线安全","紧急疏散通道畅通，不影响服务效率")],
        ]
        for j,space in enumerate(spaces[:3]):
            pts = space_pts[j] if j<len(space_pts) else [("功能设计","空间优化，功能完备")]
            space_design(prs, 6+j, total, "04 空间设计", space, pts, ptype)
        renders(prs, 9, total, images[1:5] if len(images)>1 else images, args.images, ptype)
        materials(prs, 10, total, ptype)
        budget(prs, 11, total, ptype)
        timeline(prs, 12, total, ptype)
    else:
        total = 13
        home_positioning(prs, 3, total, args.project, args.style, args.area,
                         family_struct="三代同堂",
                         vision=["打造三代同堂的温馨居所","让孩子有独立成长的空间",
                                 "让老人日常生活更安全便捷","让夫妻拥有私密放松的角落"])
        home_scenes(prs, 4, total, members=[
            ("长辈","安全·便捷",["卫生间扶手+防滑地面","起夜感应灯","沙发坐高适老化","储物弯腰不费力"]),
            ("孩子","成长·陪伴",["独立学习角","安全插座保护","涂鸦墙设计","户外活动动线"]),
            ("主人","私密·放松",["主卧套房设计","独立书房","步入式衣帽间","南向采光主卫"]),
            ("全家","情感·连接",["开放式客餐厅","阳台亲子区","节日团聚餐桌","客厅阅读角"]),
        ])
        home_aesthetics(prs, 5, total,
                        aesthetics=["天然木材的温润感","石材的沉稳大气",
                                    "织物的柔软亲和","金属的精致细节"],
                        light_design=["自然光：引光入室，借景自然",
                                     "氛围光：调光设计，分时分景",
                                     "重点光：突出焦点，聚焦美感",
                                     "夜灯：柔和指引，不扰睡眠"],
                        emotion_scenes=["晨光中的早餐角","雨天的阅读窗边",
                                       "周末的家庭影院","深夜的一盏小灯"])
        spaces = [s.strip() for s in args.spaces.split(',')] if args.spaces else DEFAULT_SPACES_H
        space_pts = [
            [("家庭核心区","客餐厅一体化设计，促进家人互动"),
             ("采光优化","最大化引入自然光，连接户外景观")],
            [("用餐仪式感","餐桌尺寸适应聚餐人数，餐边柜补充收纳"),
             ("灯光氛围","餐厅专用吊灯，营造用餐仪式感")],
            [("睡眠质量","隔音+遮光设计，保障深度睡眠"),
             ("收纳系统","衣柜内部功能分化，便于拿取")],
        ]
        for j,space in enumerate(spaces[:3]):
            pts = space_pts[j] if j<len(space_pts) else [("设计要点","空间优化，功能完备")]
            space_design(prs, 6+j, total, "04 空间设计", space, pts, ptype)
        renders(prs, 9, total, images[1:5] if len(images)>1 else images, args.images, ptype)
        materials(prs, 10, total, ptype)
        home_storage(prs, 11, total)
        budget(prs, 12, total, ptype)
        timeline(prs, 13, total, ptype)

    closing(prs, args.project, ptype, args.company, args.designer)
    prs.save(args.output)
    print(f"✅ 生成完成！📁 {os.path.abspath(args.output)} 📄 {len(prs.slides)}页 🏷️{'商业' if ptype=='commercial' else '住宅'}")

if __name__=="__main__":
    main()
