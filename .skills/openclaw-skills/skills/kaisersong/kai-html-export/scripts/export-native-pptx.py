#!/usr/bin/env python3
"""
export-native-pptx-v5.py — 改进版 v5

v5 修复：
1. slideSize 坐标修复（critical）：使用 PX_PER_IN 代替 96 → 幻灯片宽度从 15 英寸修正为 13.33 英寸，元素位置正确居中
2. 渐变背景提取：从 body/slide 提取两端颜色，输出为 PPTX 原生渐变填充
3. 字体映射补全：添加 PingFang SC、system-ui、-apple-system → 中文字体正确映射至 Microsoft YaHei

v4 修复：
1. extractSegments 追踪 fontWeight（bold）→ <strong>/<b> 标签在 PPTX 中正确加粗
2. 修复卡片边框处理：区分全边框 vs 单侧边框
3. 合并逻辑同时考虑 color 和 fontWeight
"""

import sys
import argparse
import io
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
import re

def check_deps():
    missing = []
    try:
        from playwright.sync_api import sync_playwright
    except ImportError:
        missing.append("playwright")
    try:
        from pptx import Presentation
    except ImportError:
        missing.append("python-pptx")
    if missing:
        print(f"Install: pip install {' '.join(missing)}")
        sys.exit(1)

check_deps()

from playwright.sync_api import sync_playwright, Page
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree as _etree


def set_roundrect_adj(shape, radius_px: float, width_in: float, height_in: float):
    """Set rounded rectangle corner radius via OOXML adj value."""
    from lxml import etree
    NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    radius_in = radius_px / 108.0
    shorter = min(width_in, height_in)
    if shorter <= 0:
        return
    adj = int(radius_in / (shorter / 2) * 100000)
    adj = max(0, min(50000, adj))
    prstGeom = shape._element.spPr.find(f'{{{NS}}}prstGeom')
    if prstGeom is None:
        return
    avLst = prstGeom.find(f'{{{NS}}}avLst')
    if avLst is None:
        avLst = etree.SubElement(prstGeom, f'{{{NS}}}avLst')
    for gd in avLst.findall(f'{{{NS}}}gd'):
        avLst.remove(gd)
    gd = etree.SubElement(avLst, f'{{{NS}}}gd')
    gd.set('name', 'adj')
    gd.set('fmla', f'val {adj}')


def suppress_line(shape):
    """Write <a:ln><a:noFill/></a:ln> directly into spPr to remove shape outline."""
    from lxml import etree
    NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    spPr = shape._element.spPr
    ln_tag = f'{{{NS}}}ln'
    ln = spPr.find(ln_tag)
    if ln is not None:
        spPr.remove(ln)
    ln = etree.SubElement(spPr, ln_tag)
    etree.SubElement(ln, f'{{{NS}}}noFill')


def set_light_shadow(shape):
    """Replace heavy theme shadow (alpha≈35%) with a subtle light drop shadow (alpha≈8%).
    Applies to both spPr/effectLst and resets effectRef idx=0 so theme doesn't override."""
    from lxml import etree
    NP = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    NA = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    spPr = shape._element.spPr
    eff_tag = f'{{{NA}}}effectLst'
    existing = spPr.find(eff_tag)
    if existing is not None:
        spPr.remove(existing)
    effectLst = etree.fromstring(
        f'<a:effectLst xmlns:a="{NA}">'
        f'<a:outerShdw blurRad="25000" dist="8000" dir="5400000" rotWithShape="0">'
        f'<a:srgbClr val="000000"><a:alpha val="8000"/></a:srgbClr>'
        f'</a:outerShdw></a:effectLst>'
    )
    spPr.append(effectLst)
    style_elem = shape._element.find(f'{{{NP}}}style')
    if style_elem is not None:
        eff_ref = style_elem.find(f'{{{NA}}}effectRef')
        if eff_ref is not None:
            eff_ref.set('idx', '0')


def parse_color(css_color: str, bg: Tuple[int,int,int] = (255, 255, 255)) -> Optional[Tuple[int, int, int]]:
    """Parse a CSS color string, blending rgba() alpha over the given bg color (default white)."""
    if not css_color or css_color in ('transparent',) or 'rgba(0, 0, 0, 0)' in css_color:
        return None
    m = re.search(r'rgba?\((\d+),\s*(\d+),\s*(\d+)(?:,\s*([\d.]+))?\)', css_color)
    if m:
        r, g, b = int(m.group(1)), int(m.group(2)), int(m.group(3))
        a = float(m.group(4)) if m.group(4) else 1.0
        if a <= 0:
            return None
        if a < 1.0:
            r = int(a * r + (1 - a) * bg[0])
            g = int(a * g + (1 - a) * bg[1])
            b = int(a * b + (1 - a) * bg[2])
        return (r, g, b)
    m = re.search(r'#([0-9a-fA-F]{6}|[0-9a-fA-F]{3})', css_color)
    if m:
        h = m.group(1)
        if len(h) == 3:
            h = ''.join([c*2 for c in h])
        return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    return None


def px_to_pt(px_value: str) -> float:
    m = re.search(r'([\d.]+)px', str(px_value))
    if m:
        return round(float(m.group(1)) * 0.75, 1)
    return 12.0


def inject_visible(page: Page):
    page.evaluate(
        'document.querySelectorAll(".slide").forEach(s => s.classList.add("visible"));'
        'document.querySelectorAll(".reveal").forEach(el => {'
        '  el.style.opacity = "1"; el.style.transform = "none";'
        '});'
    )


# JS extraction code (stored as a string, careful with escaping)
_EXTRACT_JS = r"""
(slideIndex) => {
    const slide = document.querySelectorAll('.slide')[slideIndex];
    if (!slide) return {background: null, elements: [], slideSize: null};

    // Fallback: if no .slide-content wrapper, use the slide itself as content root
    const content = slide.querySelector('.slide-content') || slide;

    const slideRect = slide.getBoundingClientRect();
    // Correct px→inch scale: PPTX slide is 13.33in wide, HTML slide is slideRect.width px wide
    // Using 96 (CSS dpi) causes 12.5% layout error. Use actual viewport-to-PPTX ratio instead.
    const PX_PER_IN = slideRect.width / 13.33;
    const bodyStyle = window.getComputedStyle(document.body);
    let bgColor = bodyStyle.backgroundColor;
    if (!bgColor || bgColor === 'transparent' || bgColor === 'rgba(0, 0, 0, 0)') {
        // Body may use a CSS gradient (background shorthand) — extract first color stop
        const bgImg = bodyStyle.backgroundImage || '';
        if (bgImg.includes('gradient')) {
            const cm = bgImg.match(/rgba?\([^)]+\)/g);
            bgColor = (cm && cm.length > 0) ? cm[0] : null;
        } else {
            bgColor = null;
        }
    }

    const slideW = slideRect.width;
    const slideH = slideRect.height;
    let exportCounter = 0;

    function safeInnerText(node) {
        if (!node) return '';
        if (typeof node.innerText === 'string') return node.innerText;
        if (typeof node.textContent === 'string') return node.textContent;
        return '';
    }

    // Helper: resolve border-radius to px string.
    // getComputedStyle may return "50%" for border-radius:50% class rules — convert to px.
    function resolveBorderRadius(style, rect) {
        const br = style.borderTopLeftRadius || style.borderRadius || '';
        if (!br || br === '0px') return '0px';
        if (br.endsWith('%')) {
            const px = parseFloat(br) / 100 * Math.min(rect.width, rect.height);
            return px.toFixed(2) + 'px';
        }
        return br;
    }

    // Helper: extract inline text segments with colors
    function extractSegments(el) {
        const elStyle = window.getComputedStyle(el);
        const bgImage = elStyle.backgroundImage || '';
        const bgClip = elStyle.webkitBackgroundClip || elStyle.backgroundClip || '';
        const isGradient = bgImage.includes('gradient') && bgClip === 'text';

        let gradientColor = null;
        if (isGradient) {
            const cm = bgImage.match(/rgba?\([^)]+\)/g);
            if (cm) gradientColor = cm[0];  // use first (darkest) stop — last stop is too light
        }

        let gradientColors = null;
        if (isGradient) {
            const cm2 = bgImage.match(/rgba?\([^)]+\)/g);
            if (cm2 && cm2.length >= 2) gradientColors = [cm2[0], cm2[cm2.length-1]];
        }

        const segments = [];
        // Track the Chrome bounding rect of the nearest ancestor element that introduced
        // an inline background. Used by Python to draw a tight background shape (instead of
        // <a:highlight> which fills the full line cell and bleeds into adjacent lines).
        let currentInlineBgBounds = null;
        function walk(node, color, bold, fontSize, strike, bgColor) {
            if (node.nodeType === 3) {
                const t = node.textContent;
                // Preserve internal whitespace (e.g. " text" after <strong>Note:</strong>)
                // but skip nodes that are purely whitespace
                if (t && t.trim()) segments.push({text: t, color: color, bold: bold, fontSize: fontSize, strike: strike, bgColor: bgColor || null, inlineBgBounds: currentInlineBgBounds || null});
            } else if (node.nodeType === 1) {
                const tag = node.tagName;
                if (tag === 'BR') { segments.push({text: '\n', color: color, bold: bold, fontSize: fontSize, strike: strike, bgColor: null, inlineBgBounds: null}); return; }
                const s2 = window.getComputedStyle(node);
                const bi = s2.backgroundImage || '';
                const bc = s2.webkitBackgroundClip || s2.backgroundClip || '';
                let c = color;
                if (bi.includes('gradient') && bc === 'text') {
                    const cm = bi.match(/rgba?\([^)]+\)/g);
                    if (cm) c = cm[0];  // first (darkest) stop
                } else {
                    const sc = s2.color;
                    if (sc && sc !== 'rgba(0, 0, 0, 0)') c = sc;
                }
                let b = bold;
                const fw = s2.fontWeight;
                if (fw === 'bold' || fw === '700' || fw === '800' || fw === '900' || parseInt(fw) >= 600) b = true;
                const fs = s2.fontSize || fontSize;
                let sk = strike;
                const td = s2.textDecoration || s2.textDecorationLine || '';
                if (td.includes('line-through')) sk = true;
                // Propagate inline background color (e.g. <span style="background:yellow">)
                const childBg = s2.backgroundColor;
                const childHasBg = childBg && childBg !== 'transparent' && childBg !== 'rgba(0, 0, 0, 0)';
                const newBgColor = childHasBg ? childBg : bgColor;
                // Capture Chrome bounding rect for elements that introduce a new background.
                // Python uses this to draw a precisely-positioned background shape.
                const prevInlineBgBounds = currentInlineBgBounds;
                if (childHasBg) {
                    const r = node.getBoundingClientRect();
                    const st2 = window.getComputedStyle(node);
                    const hPad = (parseFloat(st2.paddingLeft || '0') + parseFloat(st2.paddingRight || '0')) / 2;
                    currentInlineBgBounds = {
                        x: (r.left - slideRect.left) / PX_PER_IN,
                        y: (r.top - slideRect.top) / PX_PER_IN,
                        w: r.width / PX_PER_IN,
                        h: r.height / PX_PER_IN,
                        vPad: hPad / PX_PER_IN,
                    };
                }
                for (const child of node.childNodes) walk(child, c, b, fs, sk, newBgColor);
                currentInlineBgBounds = prevInlineBgBounds;  // restore after leaving element
            }
        }

        const baseColor = isGradient ? (gradientColor || elStyle.color) : elStyle.color;
        const baseBold = parseInt(elStyle.fontWeight) >= 600;
        const baseFontSize = elStyle.fontSize;
        const baseStrike = (elStyle.textDecoration || elStyle.textDecorationLine || '').includes('line-through');
        for (const child of el.childNodes) walk(child, baseColor, baseBold, baseFontSize, baseStrike, null);

        // Merge consecutive same-color+bold+fontSize+strike+bgColor+inlineBgBounds segments
        const merged = [];
        for (const seg of segments) {
            const prevIbb = merged.length > 0 ? merged[merged.length-1].inlineBgBounds : null;
            const curIbb = seg.inlineBgBounds;
            const sameIbb = (!prevIbb && !curIbb) ||
                (prevIbb && curIbb && Math.abs(prevIbb.x - curIbb.x) < 0.001 && Math.abs(prevIbb.w - curIbb.w) < 0.001);
            if (merged.length > 0 && merged[merged.length-1].color === seg.color &&
                merged[merged.length-1].bold === seg.bold &&
                merged[merged.length-1].fontSize === seg.fontSize &&
                merged[merged.length-1].strike === seg.strike &&
                merged[merged.length-1].bgColor === seg.bgColor &&
                sameIbb &&
                seg.text !== '\n' && merged[merged.length-1].text !== '\n') {
                merged[merged.length-1].text += seg.text;
            } else {
                merged.push({text: seg.text, color: seg.color, bold: seg.bold, fontSize: seg.fontSize, strike: seg.strike, bgColor: seg.bgColor || null, inlineBgBounds: seg.inlineBgBounds || null});
            }
        }

        return {segments: merged, gradientColors: gradientColors};
    }

    const TEXT_TAGS = new Set(['h1','h2','h3','h4','h5','h6','p','li','span','a']);
    const RASTER_TAGS = new Set(['img', 'svg', 'canvas']);

    // CJK width correction factor: PingFang SC / Source Han is ~15% wider than Chrome's
    // Latin fallback. Apply to bordered/backgrounded shapes that contain CJK text so they
    // don't clip the text in Keynote/PowerPoint.
    const CJK_BOX_FACTOR = 1.15;
    const CJK_V_FACTOR = 1.30;
    function hasCJK(text) {
        return /[\u2E80-\u9FFF\uF900-\uFAFF\uFE10-\uFE6F\uFF00-\uFFEF]/.test(text);
    }

    // Flat recursive traversal
    function flatExtract(el) {
        const rect = el.getBoundingClientRect();
        if (rect.width < 1 || rect.height < 1) return [];

        const tag = el.tagName.toLowerCase();
        const style = window.getComputedStyle(el);

        const bounds = {
            x: (rect.left - slideRect.left) / PX_PER_IN,
            y: (rect.top - slideRect.top) / PX_PER_IN,
            width: rect.width / PX_PER_IN,
            height: rect.height / PX_PER_IN
        };

        const bgColor = style.backgroundColor;
        const hasBg = bgColor && bgColor !== 'transparent' && bgColor !== 'rgba(0, 0, 0, 0)';

        // Check all border sides
        const borderLeft = style.borderLeft || '';
        const borderRight = style.borderRight || '';
        const borderTop = style.borderTop || '';
        const borderBottom = style.borderBottom || '';
        const borderStr = style.border || '';

        const hasLeftBorder = borderLeft && !borderLeft.includes('none') && !borderLeft.startsWith('0px');
        const hasRightBorder = borderRight && !borderRight.includes('none') && !borderRight.startsWith('0px');
        const hasTopBorder = borderTop && !borderTop.includes('none') && !borderTop.startsWith('0px');
        const hasBottomBorder = borderBottom && !borderBottom.includes('none') && !borderBottom.startsWith('0px');
        const hasGeneralBorder = borderStr && !borderStr.includes('none') && !borderStr.startsWith('0px');
        const hasBorder = hasLeftBorder || hasRightBorder || hasTopBorder || hasBottomBorder || hasGeneralBorder;

        const results = [];

        function registerRaster(kind, source = '') {
            const exportId = `kai-export-${slideIndex}-${++exportCounter}`;
            el.setAttribute('data-kai-export-id', exportId);
            return {
                type: 'image',
                tag: tag,
                imageKind: kind,
                exportId: exportId,
                source: source,
                bounds: bounds,
                styles: {
                    borderRadius: resolveBorderRadius(style, rect),
                    objectFit: style.objectFit || ''
                }
            };
        }

        if (RASTER_TAGS.has(tag)) {
            const source = tag === 'img' ? (el.currentSrc || el.src || '') : '';
            results.push(registerRaster(tag, source));
            return results;
        }

        if (TEXT_TAGS.has(tag)) {
            // If element has visible background/border (e.g. .pill badge span), emit a shape first
            if (hasBg || hasBorder) {
                const elText = safeInnerText(el);
                const elFontStr = style.fontFamily || '';
                const elIsCondensed = /condensed|narrow|compressed/i.test(elFontStr);
                // Condensed fonts (e.g. Barlow Condensed) get a 30% width boost to compensate for
                // system font substitution in PPTX. CJK expansion only applies to small boxes (<3in)
                // to avoid blowing up wide containers like list items with border-left.
                const cjkShapeW = hasBorder ? (
                    elIsCondensed ? bounds.width * 1.50 :
                    (hasCJK(elText) && bounds.width < 3.0 ? bounds.width * CJK_BOX_FACTOR : bounds.width)
                ) : bounds.width;
                results.push({
                    type: 'shape', tag: tag,
                    bounds: { x: bounds.x, y: bounds.y, width: cjkShapeW, height: bounds.height },
                    styles: {
                        backgroundColor: bgColor,
                        backgroundImage: '',
                        border: borderStr, borderLeft: borderLeft, borderRight: borderRight,
                        borderTop: borderTop, borderBottom: borderBottom,
                        borderRadius: resolveBorderRadius(style, rect)
                    }
                });
            }
            // Inline child bg colors are captured inside extractSegments (bgColor on each segment)
            // and rendered as <a:highlight> text run properties — no separate shape needed.
            // Text element - render as text box
            const {segments, gradientColors} = extractSegments(el);
            const text = safeInnerText(el).trim();
            if (text || segments.length > 0) {
                const elFontStr2 = style.fontFamily || '';
                const elIsCondensed2 = /condensed|narrow|compressed/i.test(elFontStr2);
                // Condensed fonts: expand text box without requiring hasBorder
                // CJK: expand only for small bordered boxes to avoid wide list-item overflow
                const cjkTxtFactor = elIsCondensed2 ? 1.50 :
                    (hasBorder && hasCJK(text) && bounds.width < 3.0 ? CJK_BOX_FACTOR : 1.0);
                const txtBounds = cjkTxtFactor === 1.0 ? bounds : {
                    x: bounds.x, y: bounds.y,
                    width: bounds.width * cjkTxtFactor, height: bounds.height
                };
                results.push({
                    type: 'text',
                    tag: tag,
                    text: text,
                    segments: segments,
                    gradientColors: gradientColors,
                    textTransform: style.textTransform,
                    naturalHeight: (function() {
                        // scrollHeight == clientHeight for overflow:visible — doesn't capture
                        // content flowing past the element's CSS height boundary.
                        // Temporarily remove the height constraint to measure true content height.
                        // Safe: height is restored synchronously before any sibling is measured.
                        const savedH = el.style.height;
                        const savedMaxH = el.style.maxHeight;
                        el.style.height = 'auto';
                        el.style.maxHeight = 'none';
                        const naturalH = el.getBoundingClientRect().height;
                        el.style.height = savedH;
                        el.style.maxHeight = savedMaxH;
                        return Math.max(rect.height, naturalH) / PX_PER_IN;
                    })(),
                    bounds: txtBounds,
                    styles: {
                        fontSize: style.fontSize,
                        fontWeight: style.fontWeight,
                        fontFamily: style.fontFamily,
                        letterSpacing: style.letterSpacing,
                        color: style.color,
                        textAlign: style.textAlign,
                        lineHeight: style.lineHeight,
                        listStyleType: style.listStyleType,
                        paddingLeft: style.paddingLeft,
                        paddingRight: style.paddingRight,
                        paddingTop: style.paddingTop,
                        paddingBottom: style.paddingBottom
                    }
                });
            }
            // Don't recurse into text elements
            return results;
        }

        if (tag === 'table') {
            // Extract table structure: all rows and cells with bounds/styles
            const tableRows = [];
            const allRows = el.querySelectorAll('tr');
            allRows.forEach((row) => {
                const isHeader = !!row.closest('thead');
                const cells = row.querySelectorAll('th, td');
                const rowCells = [];
                cells.forEach(cell => {
                    const cellRect = cell.getBoundingClientRect();
                    const cellStyle = window.getComputedStyle(cell);
                    const {segments} = extractSegments(cell);
                    rowCells.push({
                        bounds: {
                            x: (cellRect.left - slideRect.left) / PX_PER_IN,
                            y: (cellRect.top - slideRect.top) / PX_PER_IN,
                            width: cellRect.width / PX_PER_IN,
                            height: cellRect.height / PX_PER_IN
                        },
                        text: safeInnerText(cell).trim(),
                        segments: segments,
                        isHeader: isHeader,
                        styles: {
                            fontSize: cellStyle.fontSize,
                            fontWeight: cellStyle.fontWeight,
                            color: cellStyle.color,
                            backgroundColor: cellStyle.backgroundColor,
                            textAlign: cellStyle.textAlign,
                            paddingLeft: cellStyle.paddingLeft,
                            paddingRight: cellStyle.paddingRight,
                            paddingTop: cellStyle.paddingTop,
                            paddingBottom: cellStyle.paddingBottom,
                            fontFamily: cellStyle.fontFamily,
                            letterSpacing: cellStyle.letterSpacing,
                            borderBottom: cellStyle.borderBottom,
                            borderRight: cellStyle.borderRight
                        }
                    });
                });
                if (rowCells.length > 0) tableRows.push({isHeader: isHeader, cells: rowCells});
            });
            results.push({type: 'table', bounds: bounds, rows: tableRows});
            return results;  // Don't recurse into table children
        }

        if (tag === 'div' || tag === 'section' || tag === 'article' || tag === 'ul' || tag === 'ol') {
            const bgImage = style.backgroundImage || 'none';
            const hasGradientBg = bgImage !== 'none' && bgImage.includes('gradient');
            const hasUrlBg = bgImage !== 'none' && bgImage.includes('url(');
            const totalText = safeInnerText(el).trim();

            if (hasUrlBg && !totalText) {
                results.push(registerRaster('background-image', bgImage));
                return results;
            }

            // Filter out decorative elements with highly transparent backgrounds and no text
            // (e.g. ambient orb/cloud divs in blue-sky style: rgba(x,x,x,0.3) blobs)
            const alphaMatch = (bgColor || '').match(/rgba\([^)]+,\s*([\d.]+)\s*\)/);
            const bgAlpha = alphaMatch ? parseFloat(alphaMatch[1]) : 1.0;
            const hasChildRasters = el.querySelectorAll('img, svg, canvas').length > 0;
            if (bgAlpha < 0.5 && !totalText && !hasChildRasters && bounds.width > 1.5 && bounds.height > 1.5) return [];

            // Detect "leaf-text container": a div whose entire visible content is text
            // Case A: no child elements at all — e.g. <div class="chapter-num">01</div>
            // Case B: only inline child elements + sibling text nodes — e.g. <div class="co"><strong>Note:</strong> more text</div>
            const INLINE_TAGS = new Set(['STRONG','EM','B','I','SPAN','A','MARK','CODE','SMALL','KBD','VAR','ABBR','TIME','SUP','SUB','BR']);
            const allChildInline = el.children.length > 0 &&
                Array.from(el.children).every(c => INLINE_TAGS.has(c.tagName));
            const childrenTextLen = Array.from(el.children)
                .map(c => safeInnerText(c).trim()).join('').replace(/\s+/g, '').length;
            const totalTextLen = totalText.replace(/\s+/g, '').length;
            // "has direct text" if total text is notably more than what children account for
            const hasDirectText = totalText && (
                el.children.length === 0 ||
                (allChildInline && totalTextLen > childrenTextLen + 1)
            );

            if (hasDirectText) {
                // First: preserve background/border styling (e.g. callout amber bg + left border)
                const bgImageLeaf = style.backgroundImage || 'none';
                const bgClipLeaf = style.webkitBackgroundClip || style.backgroundClip || '';
                // background-clip:text means gradient is used as text fill, not a visible background
                const hasGradientBgLeaf = bgImageLeaf !== 'none' && bgImageLeaf.includes('gradient') && bgClipLeaf !== 'text';
                if (hasBg || hasBorder || hasGradientBgLeaf) {
                    const cjkLeafW = (hasBorder && hasCJK(totalText)) ? bounds.width * CJK_BOX_FACTOR : bounds.width;
                    const leafBounds = cjkLeafW === bounds.width ? bounds : {
                        x: bounds.x, y: bounds.y, width: cjkLeafW, height: bounds.height
                    };
                    results.push({
                        type: 'shape', tag: tag, bounds: leafBounds,
                        styles: {
                            backgroundColor: bgColor,
                            backgroundImage: hasGradientBgLeaf ? bgImageLeaf : '',
                            border: borderStr, borderLeft: borderLeft, borderRight: borderRight,
                            borderTop: borderTop, borderBottom: borderBottom,
                            borderRadius: resolveBorderRadius(style, rect)
                        }
                    });
                }
                // Inline bg colors are captured via extractSegments bgColor field → text run highlight
                // Then: render the whole container as a text box (captures direct text + inline formatting)
                const {segments, gradientColors} = extractSegments(el);
                // Condensed fonts: expand text box to compensate for system font substitution
                const hd_fontStr = style.fontFamily || '';
                const hd_isCondensed = /condensed|narrow|compressed/i.test(hd_fontStr);
                const textBounds = hd_isCondensed
                    ? {x: bounds.x, y: bounds.y, width: bounds.width * 1.50, height: bounds.height}
                    : bounds;
                results.push({
                    type: 'text', tag: tag, text: totalText,
                    segments: segments, gradientColors: gradientColors,
                    textTransform: style.textTransform,
                    bounds: textBounds,
                    styles: {
                        fontSize: style.fontSize, fontWeight: style.fontWeight,
                        fontFamily: style.fontFamily, letterSpacing: style.letterSpacing,
                        color: style.color, textAlign: style.textAlign,
                        lineHeight: style.lineHeight, listStyleType: style.listStyleType,
                        paddingLeft: style.paddingLeft, paddingRight: style.paddingRight,
                        paddingTop: style.paddingTop, paddingBottom: style.paddingBottom,
                        alignItems: style.alignItems, justifyContent: style.justifyContent
                    }
                });
                return results;
            }

            // Standard container: maybe has background shape, then recurse
            if (hasBg || hasBorder || hasGradientBg) {
                const containerText = safeInnerText(el);
                // Condensed fonts: check self AND direct children (parent div may not inherit
                // the font-family set on a child .headline element).
                const containerFontStr = style.fontFamily || '';
                const isCondensedContainer = /condensed|narrow|compressed/i.test(containerFontStr) ||
                    Array.from(el.children).some(ch =>
                        /condensed|narrow|compressed/i.test(window.getComputedStyle(ch).fontFamily));
                // Condensed: expand always (title blocks can be wide — no width gate)
                // CJK: only expand small boxes (<3in) to avoid blowing up wide list containers
                const cjkContainerW = hasBorder ? (
                    isCondensedContainer ? bounds.width * 1.50 :
                    (hasCJK(containerText) && bounds.width < 3.0 ? bounds.width * CJK_BOX_FACTOR : bounds.width)
                ) : bounds.width;
                const containerBounds = cjkContainerW === bounds.width ? bounds : {
                    x: bounds.x, y: bounds.y, width: cjkContainerW, height: bounds.height
                };
                results.push({
                    type: 'shape',
                    tag: tag,
                    bounds: containerBounds,
                    styles: {
                        backgroundColor: bgColor,
                        backgroundImage: hasGradientBg ? bgImage : '',
                        border: borderStr,
                        borderLeft: borderLeft,
                        borderRight: borderRight,
                        borderTop: borderTop,
                        borderBottom: borderBottom,
                        borderRadius: resolveBorderRadius(style, rect)
                    }
                });
            }

            // Recurse into children
            for (const child of el.children) {
                results.push(...flatExtract(child));
            }
        }

        return results;
    }

    // Only process children of .slide-content (not the entire slide)
    const elements = [];
    for (const child of content.children) {
        elements.push(...flatExtract(child));
    }

    // Also try to extract gradient stops for slide background
    let bgGradient = null;
    let gridBg = null;
    const bodyBgImg = window.getComputedStyle(document.body).backgroundImage || '';
    // Detect grid pattern: two crossing linear-gradients (one at 90deg) = grid
    // Check body first (grid is often on body), then fall back to slide element
    const slideBgImgFull = window.getComputedStyle(slide).backgroundImage || '';
    const gridSourceImg = (bodyBgImg.match(/linear-gradient/g) || []).length >= 2 && bodyBgImg.includes('90deg')
        ? bodyBgImg : slideBgImgFull;
    const gridSourceEl = (bodyBgImg.match(/linear-gradient/g) || []).length >= 2 && bodyBgImg.includes('90deg')
        ? document.body : slide;
    const gradientCount = (gridSourceImg.match(/linear-gradient/g) || []).length;
    if (gradientCount >= 2 && gridSourceImg.includes('90deg')) {
        // Grid pattern detected
        const colorMatch = gridSourceImg.match(/rgba?\([^)]+\)/);
        const sizeStr = window.getComputedStyle(gridSourceEl).backgroundSize || '';
        const sizeMatch = sizeStr.match(/([\d.]+)px/);
        const gridSizePx = sizeMatch ? parseFloat(sizeMatch[1]) : 24;
        if (colorMatch) gridBg = {color: colorMatch[0], sizePx: gridSizePx};
    } else {
        if (bodyBgImg.includes('gradient') && (bodyBgImg.match(/linear-gradient/g) || []).length < 2) {
            const stops = bodyBgImg.match(/rgba?\([^)]+\)/g);
            if (stops && stops.length >= 2) bgGradient = [stops[0], stops[stops.length - 1]];
        }
        if (!bgGradient) {
            const slideBgImg = slideBgImgFull;
            if (slideBgImg.includes('gradient') && (slideBgImg.match(/linear-gradient/g) || []).length < 2) {
                const stops = slideBgImg.match(/rgba?\([^)]+\)/g);
                if (stops && stops.length >= 2) bgGradient = [stops[0], stops[stops.length - 1]];
            }
        }
    }

    const exportProgressSetting =
        (document.body && document.body.dataset.exportProgress) ||
        document.documentElement.dataset.exportProgress ||
        'true';
    const shouldExportChrome = !['false', '0', 'off', 'no'].includes(String(exportProgressSetting).toLowerCase());

    // Detect if the template already provides its own navigation chrome
    // (.nav-dots, .slide-counter, etc.) — if so, skip PPTX chrome injection
    const hasOwnChrome = !!(document.querySelector('.nav-dots') ||
                            document.querySelector('.slide-counter') ||
                            document.querySelector('.page-counter'));

    // Extract fixed-position chrome: nav-dots + progress-bar
    let fixedChrome = null;
    const navDotsEl = document.querySelector('.nav-dots');
    const progressBarEl = document.querySelector('.progress-bar');
    if (navDotsEl || progressBarEl) {
        fixedChrome = {};
        if (navDotsEl && shouldExportChrome) {
            const buttons = navDotsEl.querySelectorAll('button');
            const navRect = navDotsEl.getBoundingClientRect();
            const dots = [];
            buttons.forEach(btn => {
                const r = btn.getBoundingClientRect();
                const s = window.getComputedStyle(btn);
                dots.push({
                    x: r.left / PX_PER_IN,
                    y: r.top / PX_PER_IN,
                    w: r.width / PX_PER_IN,
                    h: r.height / PX_PER_IN,
                    bg: s.backgroundColor,
                    border: s.borderColor,
                    active: btn.classList.contains('active')
                });
            });
            fixedChrome.navDots = dots;
        }
        if (progressBarEl && shouldExportChrome) {
            const r = progressBarEl.getBoundingClientRect();
            const s = window.getComputedStyle(progressBarEl);
            if (r.width > 2) {  // only capture if visible (width > 0)
                fixedChrome.progressBar = {
                    x: 0, y: 0,
                    w: r.width / PX_PER_IN,
                    h: r.height / PX_PER_IN,
                    bg: s.backgroundColor
                };
            }
        }
    }

    return {
        background: bgColor,
        bgGradient: bgGradient,
        gridBg: gridBg,
        elements: elements,
        hasOwnChrome: hasOwnChrome,
        fixedChrome: fixedChrome,
        slideSize: slideW && slideH ? {width: slideW/PX_PER_IN, height: slideH/PX_PER_IN} : null
    };
}
"""


def extract_slide_elements(page: Page, slide_index: int) -> Dict[str, Any]:
    result = page.evaluate(_EXTRACT_JS, slide_index)
    bg_rgb = parse_color(result['background']) if result['background'] else None
    bg_gradient = result.get('bgGradient')  # [colorStr1, colorStr2] or None
    bg_gradient_rgb = None
    if bg_gradient and len(bg_gradient) >= 2:
        c1 = parse_color(bg_gradient[0])
        c2 = parse_color(bg_gradient[1])
        if c1 and c2:
            bg_gradient_rgb = (c1, c2)
    return {
        'background': bg_rgb,
        'bgGradient': bg_gradient_rgb,
        'gridBg': result.get('gridBg'),   # {color: "rgba(...)", sizePx: 24} or None
        'elements': result['elements'],
        'hasOwnChrome': result.get('hasOwnChrome', False),
        'fixedChrome': result.get('fixedChrome'),  # {navDots: [...], progressBar: {...}} or None
        'slideSize': result.get('slideSize')
    }


# CSS font → PPTX font mapping
# Key: substring to match in CSS fontFamily; Value: (latin_font, ea_font)
_FONT_MAP = {
    'Clash Display': ('Calibri Light', 'Microsoft YaHei'),
    'Satoshi':       ('Calibri',       'Microsoft YaHei'),
    'Microsoft YaHei': ('Microsoft YaHei', 'Microsoft YaHei'),
    '微软雅黑':          ('Microsoft YaHei', 'Microsoft YaHei'),
    'PingFang SC':      ('PingFang SC',     'PingFang SC'),
    'Noto Sans CJK SC': ('Noto Sans CJK SC','Noto Sans CJK SC'),
    'Source Han Sans':  ('Source Han Sans', 'Source Han Sans'),
    'system-ui':        ('Microsoft YaHei', 'Microsoft YaHei'),
    '-apple-system':    ('PingFang SC',     'PingFang SC'),
}
_DEFAULT_FONTS = ('Microsoft YaHei', 'Microsoft YaHei')


def map_font(css_font_family: str):
    """Map CSS fontFamily string to (latin_font, ea_font) tuple."""
    if css_font_family:
        for css_name, fonts in _FONT_MAP.items():
            if css_name in css_font_family:
                return fonts
    return _DEFAULT_FONTS


def set_run_fonts(run, latin_font, ea_font):
    """Set latin, ea, cs fonts on a run element."""
    from lxml import etree
    NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    run.font.name = latin_font
    rPr = run._r.get_or_add_rPr()
    for tag, typeface in [('ea', ea_font), ('cs', ea_font)]:
        el = rPr.find(f'{{{NS}}}{tag}')
        if el is None:
            el = etree.SubElement(rPr, f'{{{NS}}}{tag}')
        el.set('typeface', typeface)


def set_letter_spacing(run, css_letter_spacing: str):
    """Set character spacing from CSS letterSpacing (e.g. '1.65px')."""
    if not css_letter_spacing or css_letter_spacing in ('normal', '0px'):
        return
    m = re.search(r'([\d.]+)px', css_letter_spacing)
    if m:
        px = float(m.group(1))
        # OOXML spc unit = 1/100 pt; 1px = 0.75pt → spc = px * 75
        spc = int(px * 75)
        if spc > 0:
            run._r.get_or_add_rPr().set('spc', str(spc))


def segments_to_lines(segments):
    """将 segments 列表拆分成行列表 (list of list)"""
    # Strip leading/trailing whitespace from text segments (HTML formatting whitespace)
    # But preserve \n (from BR tags) as line separators
    cleaned = []
    for s in segments:
        t = s['text']
        if t == '\n':
            cleaned.append(s)  # Keep newlines as-is
        elif t.strip():
            # Preserve original text (including leading/trailing spaces between inline elements
            # e.g. "<strong>Note:</strong> text" produces " text" which must keep its leading space)
            cleaned.append({'text': t, 'color': s['color'], 'bold': s.get('bold', False),
                            'fontSize': s.get('fontSize', ''), 'strike': s.get('strike', False),
                            'bgColor': s.get('bgColor'), 'inlineBgBounds': s.get('inlineBgBounds')})
    segments = cleaned
    lines = []
    current_line = []
    for seg in segments:
        text = seg['text']
        color = seg['color']
        bold = seg.get('bold', False)
        fontSize = seg.get('fontSize', '')
        strike = seg.get('strike', False)
        bg_color = seg.get('bgColor')
        if '\n' in text:
            parts = text.split('\n')
            for i, part in enumerate(parts):
                if part:
                    current_line.append({'text': part, 'color': color, 'bold': bold, 'fontSize': fontSize, 'strike': strike, 'bgColor': bg_color, 'inlineBgBounds': seg.get('inlineBgBounds')})
                if i < len(parts) - 1:
                    lines.append(current_line)
                    current_line = []
        else:
            current_line.append({'text': text, 'color': color, 'bold': bold, 'fontSize': fontSize, 'strike': strike, 'bgColor': bg_color, 'inlineBgBounds': seg.get('inlineBgBounds')})
    lines.append(current_line)
    # Strip leading/trailing empty lines, but preserve internal empty lines (from <BR><BR>)
    result = lines
    while result and not any(s['text'].strip() for s in result[0]):
        result = result[1:]
    while result and not any(s['text'].strip() for s in result[-1]):
        result = result[:-1]
    return result


def apply_run(run, text, color_str, font_size_pt, font_weight,
              text_transform='none', font_family='', letter_spacing='', strike=False, bg_color_str=None):
    if text_transform == 'uppercase':
        text = text.upper()
    run.text = text
    # Preserve leading/trailing spaces in runs (OOXML strips them without xml:space="preserve")
    if text and (text[0] == ' ' or text[-1] == ' '):
        _nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        _t_elem = run._r.find('.//a:t', _nsmap)
        if _t_elem is not None:
            _t_elem.set('{http://www.w3.org/XML/1998/namespace}space', 'preserve')
    # P2: font mapping
    latin_font, ea_font = map_font(font_family)
    set_run_fonts(run, latin_font, ea_font)
    if font_size_pt:
        run.font.size = Pt(font_size_pt)
    try:
        if font_weight == 'bold':
            run.font.bold = True
        else:
            run.font.bold = int(font_weight) >= 600
    except Exception:
        pass
    rgb = parse_color(color_str)
    if rgb:
        run.font.color.rgb = RGBColor(*rgb)
    if strike:
        # python-pptx 1.0.x: Font.strike is not a real property; set via XML
        _ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        _rPr = run._r.find(f'{{{_ns}}}rPr')
        if _rPr is None:
            _rPr = _etree.SubElement(run._r, f'{{{_ns}}}rPr')
            run._r.insert(0, _rPr)
        _rPr.set('strike', 'sngStrike')
    # P1: letter-spacing
    set_letter_spacing(run, letter_spacing)


def add_grid_background(slide, slide_w_in: float, slide_h_in: float,
                        grid_color_str: str, grid_size_px: float):
    """Render a CSS repeating-linear-gradient grid as a full-slide PNG picture."""
    import io
    try:
        from PIL import Image, ImageDraw
    except ImportError:
        return  # skip if Pillow not available
    scale = 3  # 3x resolution for crisp lines
    w = int(slide_w_in * 96 * scale)
    h = int(slide_h_in * 96 * scale)
    grid_px = max(1, int(grid_size_px * scale))
    img = Image.new('RGBA', (w, h), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    # Parse rgba(r,g,b,a)
    m = re.match(r'rgba?\((\d+),\s*(\d+),\s*(\d+)(?:,\s*([\d.]+))?\)', grid_color_str.strip())
    if m:
        r, g, b = int(m.group(1)), int(m.group(2)), int(m.group(3))
        a = float(m.group(4) or '1.0')
        line_color = (r, g, b, int(a * 255))
    else:
        line_color = (80, 100, 170, 25)
    for y in range(0, h, grid_px):
        draw.line([(0, y), (w - 1, y)], fill=line_color, width=1)
    for x in range(0, w, grid_px):
        draw.line([(x, 0), (x, h - 1)], fill=line_color, width=1)
    buf = io.BytesIO()
    img.save(buf, format='PNG')
    buf.seek(0)
    from pptx.util import Inches
    pic = slide.shapes.add_picture(buf, Inches(0), Inches(0),
                                   Inches(slide_w_in), Inches(slide_h_in))
    # Move grid picture to just above background (z-index 2) so content stays on top
    sp_tree = slide.shapes._spTree
    sp_tree.remove(pic._element)
    sp_tree.insert(2, pic._element)


def apply_para_format(p, s):
    lh = s.get('lineHeight', 'normal')
    if lh == 'normal':
        p.line_spacing = 1.2
    else:
        try:
            if 'px' in lh:
                lh_px = float(re.search(r'([\d.]+)', lh).group(1))
                # Use exact Pt line spacing so inline runs with larger fonts
                # don't inflate the line height beyond the CSS-measured value
                p.line_spacing = Pt(round(lh_px * 0.75, 1))
            else:
                p.line_spacing = float(lh)
        except Exception:
            p.line_spacing = 1.2
    align = s.get('textAlign', 'left')
    if align == 'center':
        p.alignment = PP_ALIGN.CENTER
    elif align == 'right':
        p.alignment = PP_ALIGN.RIGHT


def gradient_to_solid(bg_image, slide_bg=(13, 17, 23)):
    """Approximate a CSS gradient with a solid color by blending the first stop over slide bg."""
    if not bg_image or 'gradient' not in bg_image:
        return None
    rgba_matches = re.findall(r'rgba?\((\d+),\s*(\d+),\s*(\d+)(?:,\s*([\d.]+))?\)', bg_image)
    if not rgba_matches:
        return None
    r, g, b = int(rgba_matches[0][0]), int(rgba_matches[0][1]), int(rgba_matches[0][2])
    a = float(rgba_matches[0][3]) if rgba_matches[0][3] else 1.0
    if a <= 0:
        return None
    if a < 1.0:
        r = int(a * r + (1 - a) * slide_bg[0])
        g = int(a * g + (1 - a) * slide_bg[1])
        b = int(a * b + (1 - a) * slide_bg[2])
    return (r, g, b)


def export_shape_background(slide, elem, slide_bg=(255, 255, 255)):
    """仅创建背景形状（无文字），用于 type=shape 的容器"""
    b = elem['bounds']
    s = elem['styles']

    border_radius = s.get('borderRadius', '')
    radius_px = 0.0
    if border_radius and border_radius != '0px':
        m = re.search(r'([\d.]+)px', border_radius)
        if m:
            radius_px = float(m.group(1))

    # [X] 非徽章形状使用最小圆角（精致感）。
    # 徽章/胶囊判定：radius >= 40% of height（近乎圆形）→ 保留完整圆角。
    height_px = b['height'] * 108.0
    if radius_px > 0 and radius_px < height_px * 0.4:
        radius_px = min(radius_px, 6.0)

    # 提前解析边框（在画色条之前需要）
    def parse_border_side(bs):
        if not bs or 'none' in bs or bs.startswith('0px'):
            return None
        m = re.search(r'([\d.]+)px.*?rgb\((\d+),\s*(\d+),\s*(\d+)\)', bs)
        if m:
            return {'width': float(m.group(1)), 'rgb': (int(m.group(2)), int(m.group(3)), int(m.group(4)))}
        return None

    bl = parse_border_side(s.get('borderLeft', ''))
    br = parse_border_side(s.get('borderRight', ''))
    bt = parse_border_side(s.get('borderTop', ''))
    bb = parse_border_side(s.get('borderBottom', ''))

    sides_with_border = sum(1 for x in [bl, br, bt, bb] if x is not None)
    borders = [x for x in [bl, br, bt, bb] if x is not None]
    all_uniform = (len(borders) >= 3 and
                   all(bd['rgb'] == borders[0]['rgb'] and bd['width'] == borders[0]['width']
                       for bd in borders))

    # [Y2] 左侧色条"融为一体"技法 v2：色条本身是圆角矩形，整体往左移 bar_visible_px，
    # 卡片盖住右侧部分，左边圆角自然露出，效果与 HTML 原版一致。
    BAR_VISIBLE_PX = 4.0   # 色条在卡片左侧露出的宽度（px）
    bl_handled = False
    if bl and not all_uniform and radius_px > 0:
        bar_visible_in = BAR_VISIBLE_PX / 108.0
        bar_total_width_px = BAR_VISIBLE_PX + radius_px * 2 + 4  # 右侧足够被卡片覆盖
        bar_total_width_in = bar_total_width_px / 108.0
        bar_x_in = b['x'] - bar_visible_in  # 整体左移，使左侧 BAR_VISIBLE_PX 露出
        bar_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(bar_x_in), Inches(b['y']),
            Inches(bar_total_width_in), Inches(b['height'])
        )
        set_roundrect_adj(bar_shape, radius_px, bar_total_width_in, b['height'])
        bar_shape.fill.solid()
        bar_shape.fill.fore_color.rgb = RGBColor(*bl['rgb'])
        suppress_line(bar_shape)
        set_light_shadow(bar_shape)
        bl_handled = True

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius_px > 0 else MSO_SHAPE.RECTANGLE,
        Inches(b['x']), Inches(b['y']),
        Inches(b['width']), Inches(b['height'])
    )
    if radius_px > 0:
        set_roundrect_adj(shape, radius_px, b['width'], b['height'])

    bg_rgb = parse_color(s.get('backgroundColor', ''), bg=slide_bg)
    if bg_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*bg_rgb)
    else:
        grad_fill = gradient_to_solid(s.get('backgroundImage', ''), slide_bg=slide_bg)
        if grad_fill:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*grad_fill)
        else:
            shape.fill.background()

    if all_uniform:
        # 所有边框颜色/宽度一致 → 使用 shape.line
        shape.line.color.rgb = RGBColor(*borders[0]['rgb'])
        shape.line.width = Pt(max(0.5, borders[0]['width'] * 0.75))
    elif sides_with_border >= 1:
        suppress_line(shape)
        if bl and not bl_handled:
            # 无圆角时的回退：色条画在卡片上方（全高，无需缩进）
            border_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(b['x']), Inches(b['y']),
                Inches(bl['width'] / 108.0), Inches(b['height'])
            )
            border_shape.fill.solid()
            border_shape.fill.fore_color.rgb = RGBColor(*bl['rgb'])
            suppress_line(border_shape)
        if br:
            border_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(b['x'] + b['width'] - br['width']/108.0), Inches(b['y']),
                Inches(br['width'] / 108.0), Inches(b['height'])
            )
            border_shape.fill.solid()
            border_shape.fill.fore_color.rgb = RGBColor(*br['rgb'])
            suppress_line(border_shape)
        if bt:
            border_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(b['x']), Inches(b['y']),
                Inches(b['width']), Inches(bt['width'] / 108.0)
            )
            border_shape.fill.solid()
            border_shape.fill.fore_color.rgb = RGBColor(*bt['rgb'])
            suppress_line(border_shape)
        if bb:
            border_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(b['x']), Inches(b['y'] + b['height'] - bb['width']/108.0),
                Inches(b['width']), Inches(bb['width'] / 108.0)
            )
            border_shape.fill.solid()
            border_shape.fill.fore_color.rgb = RGBColor(*bb['rgb'])
            suppress_line(border_shape)
    else:
        suppress_line(shape)

    set_light_shadow(shape)

    # Remove default text frame content
    tf = shape.text_frame
    if tf.paragraphs:
        for para in tf.paragraphs:
            for run in para.runs:
                run.text = ''

    return shape


def _download_img_direct(source: str, b: Dict, object_fit: str, pptx_slide) -> bool:
    """
    For <img> elements with http(s)/file sources and object-fit: cover: download the
    original image, PIL-crop to the exact visible pixel region, and embed the cropped
    image. Only handles cover — other fits fall back to Playwright screenshot so that
    CSS opacity and rendering are captured correctly.
    Returns True on success, False to fall back to screenshot.
    """
    # Handle cover, contain, fill, and empty (default) — other values fall back to screenshot.
    if object_fit not in ('cover', 'contain', 'fill', ''):
        return False

    try:
        import urllib.request
        from PIL import Image as _PILImg
        import ssl
        if source.startswith('file://'):
            img_path = source[len('file://'):]
            with open(img_path, 'rb') as f:
                img_bytes = f.read()
        else:
            req = urllib.request.Request(source, headers={'User-Agent': 'Mozilla/5.0'})
            try:
                import certifi
                _ssl_ctx = ssl.create_default_context(cafile=certifi.where())
            except ImportError:
                _ssl_ctx = ssl._create_unverified_context()
            with urllib.request.urlopen(req, context=_ssl_ctx, timeout=15) as resp:
                img_bytes = resp.read()

        pil_img = _PILImg.open(io.BytesIO(img_bytes))
        w_img, h_img = pil_img.size
        if w_img < 1 or h_img < 1:
            return False

        w_box = b['width']   # inches
        h_box = b['height']  # inches
        aspect_box = w_box / h_box
        aspect_img = w_img / h_img

        # Compute crop fractions for object-fit: cover
        # For contain/fill/'', embed the image without cropping.
        crop_l = crop_r = crop_t = crop_b = 0.0
        if object_fit == 'cover':
            if aspect_box > aspect_img:
                # Fill width, crop top/bottom symmetrically
                crop_t = crop_b = 0.5 * (1.0 - aspect_img / aspect_box)
            else:
                # Fill height, crop left/right symmetrically
                crop_l = crop_r = 0.5 * (1.0 - aspect_box / aspect_img)

        # PIL-crop to the exact visible pixel region instead of using PPTX srcRect.
        # PPTX srcRect is rendered inconsistently across Keynote/LibreOffice/PowerPoint,
        # causing visible position shifts. Pre-cropping the image is viewer-agnostic.
        x0 = int(round(crop_l * w_img))
        x1 = int(round((1.0 - crop_r) * w_img))
        y0 = int(round(crop_t * h_img))
        y1 = int(round((1.0 - crop_b) * h_img))
        cropped = pil_img.crop((x0, y0, x1, y1))
        out = io.BytesIO()
        cropped.save(out, format='PNG')
        out.seek(0)

        pptx_slide.shapes.add_picture(
            out,
            Inches(b['x']), Inches(b['y']),
            Inches(w_box), Inches(h_box)
        )
        return True
    except Exception:
        return False


def export_raster_element(page: Page, slide, elem: Dict[str, Any]):
    """Render DOM-backed raster elements (<img>, <svg>, url() backgrounds) as PPT pictures."""
    export_id = elem.get('exportId')
    if not export_id:
        raise ValueError("Raster element missing exportId")

    # For <img> elements with a resolvable URL: download the original and use PPTX-native
    # crop to simulate object-fit. This avoids Playwright screenshot cropping differences
    # and preserves PNG transparency without background bleed-through.
    if elem.get('tag') == 'img':
        source = elem.get('source', '')
        if source.startswith(('http://', 'https://', 'file://')):
            b = elem['bounds']
            object_fit = elem.get('styles', {}).get('objectFit', 'fill')
            if _download_img_direct(source, b, object_fit, slide):
                return

    locator = page.locator(f'[data-kai-export-id="{export_id}"]').first
    if locator.count() == 0:
        raise ValueError(f"Raster element not found: {export_id}")

    # Isolate the target element before screenshotting. For full-slide images, Playwright's
    # locator.screenshot() captures the final painted pixels inside the element's bounding box,
    # which includes overlapping sibling text. Hiding non-target nodes avoids duplicating text
    # when the PPT also gets native text boxes.
    page.evaluate(
        """
        (targetId) => {
            const target = document.querySelector(`[data-kai-export-id="${targetId}"]`);
            if (!target) return false;
            const slide = target.closest('.slide');
            if (!slide) return false;
            // Hide all sibling nodes so only the target is rendered.
            const nodes = slide.querySelectorAll('*');
            nodes.forEach((node) => {
                if (node === target || target.contains(node) || node.contains(target)) return;
                if (node.dataset.kaiPrevVisibility !== undefined) return;
                node.dataset.kaiPrevVisibility = node.style.visibility || '';
                node.style.visibility = 'hidden';
            });
            // Make ancestor backgrounds transparent so PNG transparency is preserved.
            // omit_background=True only removes the browser's default white; CSS
            // background-color on ancestor elements (e.g. .slide-title { background: white })
            // still bleeds through transparent areas of the <img>/<svg>/canvas.
            let anc = target.parentElement;
            while (anc && anc !== document.documentElement) {
                anc.dataset.kaiPrevCssText = anc.style.cssText;
                anc.style.setProperty('background-color', 'transparent', 'important');
                anc.style.setProperty('background-image', 'none', 'important');
                anc = anc.parentElement;
            }
            return true;
        }
        """,
        export_id,
    )
    try:
        png_bytes = locator.screenshot(type="png", animations="disabled", omit_background=True)
        b = elem['bounds']
        slide.shapes.add_picture(
            io.BytesIO(png_bytes),
            Inches(b['x']), Inches(b['y']),
            Inches(b['width']), Inches(b['height'])
        )
    finally:
        page.evaluate(
            """
            (targetId) => {
                const target = document.querySelector(`[data-kai-export-id="${targetId}"]`);
                if (!target) return false;
                const slide = target.closest('.slide');
                if (!slide) return false;
                // Restore sibling visibility.
                const nodes = slide.querySelectorAll('[data-kai-prev-visibility]');
                nodes.forEach((node) => {
                    node.style.visibility = node.dataset.kaiPrevVisibility;
                    delete node.dataset.kaiPrevVisibility;
                });
                // Restore ancestor backgrounds.
                document.querySelectorAll('[data-kai-prev-css-text]').forEach((a) => {
                    a.style.cssText = a.dataset.kaiPrevCssText;
                    delete a.dataset.kaiPrevCssText;
                });
                return true;
            }
            """,
            export_id,
        )


def interpolate_color(c1, c2, t):
    """在两色间插值，t=0返回c1，t=1返回c2"""
    return (
        int(c1[0] + (c2[0] - c1[0]) * t),
        int(c1[1] + (c2[1] - c1[1]) * t),
        int(c1[2] + (c2[2] - c1[2]) * t)
    )


def export_text_element(slide, elem: Dict[str, Any], bg_color=None):
    b = elem['bounds']
    s = elem['styles']
    segments = elem.get('segments', [])
    text_transform = elem.get('textTransform', 'none')
    font_size_pt = px_to_pt(s.get('fontSize', '16px'))
    font_weight = s.get('fontWeight', '400')
    font_family = s.get('fontFamily', '')
    letter_spacing = s.get('letterSpacing', '')

    # Use naturalHeight (scrollHeight-based) if it's taller than the clipped bounds.
    # This corrects overflow:visible containers where Chrome clips getBoundingClientRect
    # but the content naturally extends further.
    natural_h = elem.get('naturalHeight', b['height'])
    effective_h = max(b['height'], natural_h)

    if not segments:
        raw = (elem.get('text', '') or '').strip()
        segments = [{'text': raw, 'color': s.get('color', '')}]

    lines = segments_to_lines(segments)
    if not lines:
        lines = [[{'text': '', 'color': s.get('color', '')}]]

    # Pre-pass: draw inline background shapes BEFORE text box (text on top in z-order).
    #
    # Coordinate system note:
    #   JS element bounds use PX_PER_IN = slideRect.width/13.33 ≈ 108 px/in.
    #   Python CSS properties (font-size, line-height) use ×0.75 ≡ 96 px/in.
    #   These differ by 108/96 = 1.125×, so ibg['y'] (Chrome, 108-based) is 12.5%
    #   smaller per-line than the PPTX line position (96-based). Using ibg['y'] directly
    #   puts the block too high; using b['y'] + line_idx*lh_pptx (PPTX coords) is correct.
    #
    # Formula:
    #   bg_y = b['y'] + line_idx × lh_pptx − vPad    (PPTX line top, shift up by vPad)
    #   bg_h = font_size_in + 2 × vPad               (em-square ± vPad above and below)
    #   vPad = avg CSS horizontal padding of span / PX_PER_IN  (same proportion as h-pad)
    _lh = s.get('lineHeight', 'normal')
    if _lh != 'normal' and 'px' in str(_lh):
        _lh_m = re.search(r'([\d.]+)', str(_lh))
        _lh_pptx_in = float(_lh_m.group(1)) * 0.75 / 72 if _lh_m else font_size_pt / 72.0
    else:
        try:
            _lh_ratio = float(_lh) if _lh != 'normal' else 1.2
        except Exception:
            _lh_ratio = 1.2
        _lh_pptx_in = font_size_pt * _lh_ratio / 72.0

    for _li, _lsegs in enumerate(lines):
        _seen_x: set = set()
        for _seg in _lsegs:
            _ibg = _seg.get('inlineBgBounds')
            if _ibg and _seg.get('bgColor'):
                _bg_rgb = parse_color(_seg['bgColor'])
                if not _bg_rgb:
                    continue
                _xk = round(_ibg['x'], 3)
                if _xk in _seen_x:
                    continue
                _seen_x.add(_xk)
                _v_pad = _ibg.get('vPad', 0)
                _bg_y = b['y'] + _li * _lh_pptx_in - _v_pad   # PPTX line top − vPad
                _bg_h = font_size_pt / 72.0 + 2 * _v_pad       # em-square + vPad above+below
                _bg_shape = slide.shapes.add_shape(1,
                    Inches(b['x'] + (_ibg['x'] - b['x']) * 1.15),
                    Inches(_bg_y),
                    Inches(_ibg['w'] * 1.15),
                    Inches(_bg_h))
                _bg_shape.fill.solid()
                _bg_shape.fill.fore_color.rgb = RGBColor(*_bg_rgb)
                _bg_shape.line.fill.background()

    txBox = slide.shapes.add_textbox(
        Inches(b['x']), Inches(b['y']),
        Inches(b['width']), Inches(effective_h)
    )
    tf = txBox.text_frame
    # Disable word wrap for:
    # 1. Large/heading text (font_size_pt >= 40): PingFang SC metrics differ from Calibri;
    #    intra-paragraph wrapping causes overflow onto adjacent shapes. At 40pt+, each
    #    paragraph is a designed single line — the multi-line structure comes from separate
    #    paragraphs, not CSS word-wrap.
    # 2. Single-line elements: detected by estimated_lines < 1.5
    line_height_in = font_size_pt / 72.0 * 1.2
    estimated_lines = effective_h / max(line_height_in, 0.001)
    # word_wrap / auto_size strategy:
    # - Multi-line with background shape (pptx_height_corrected=True): word_wrap=True + NONE.
    #   Height was pre-corrected ×1.3 in the pre-pass; no auto-expansion needed. Using
    #   SHAPE_TO_FIT_TEXT here would expand the text box beyond the background shape in
    #   PowerPoint/Keynote because CJK character metrics differ between Chrome and PPTX renderers.
    # - Multi-line without background: word_wrap=True + SHAPE_TO_FIT_TEXT — box expands freely.
    # - Single-line (estimated_lines < 2): word_wrap=False + TEXT_TO_FIT_SHAPE — text shrinks
    #   slightly to fit box width, preventing horizontal overflow into adjacent elements.
    tf.word_wrap = estimated_lines >= 2.0
    from pptx.enum.text import MSO_AUTO_SIZE
    if estimated_lines >= 2.0 and elem.get('pptx_height_corrected'):
        tf.auto_size = MSO_AUTO_SIZE.NONE
    elif estimated_lines >= 2.0:
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    else:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    # Map CSS padding to tf.margin_* so text doesn't flush against card edges.
    # getBoundingClientRect() returns border-box (includes padding), so the text box
    # dimensions already include padding. We set margins = CSS padding so the text
    # starts at the correct visual inset. Falls back to 0 when no padding is set.
    def parse_px(val):
        if not val or val == '0px' or val == 'auto':
            return 0.0
        m = re.search(r'([\d.]+)px', str(val))
        return float(m.group(1)) if m else 0.0
    pad_l = parse_px(s.get('paddingLeft', ''))
    pad_r = parse_px(s.get('paddingRight', ''))
    pad_t = parse_px(s.get('paddingTop', ''))
    pad_b = parse_px(s.get('paddingBottom', ''))
    tf.margin_left = Inches(pad_l / 108.0) if pad_l > 0 else 0
    tf.margin_right = Inches(pad_r / 108.0) if pad_r > 0 else 0
    tf.margin_top = Inches(pad_t / 108.0) if pad_t > 0 else 0
    tf.margin_bottom = Inches(pad_b / 108.0) if pad_b > 0 else 0
    # Vertical centering: apply when parent flex container has align-items: center
    align_items = s.get('alignItems', '')
    justify_content = s.get('justifyContent', '')
    if align_items in ('center', 'middle'):
        from pptx.oxml.ns import qn
        tf._txBody.attrib['anchor'] = 'ctr'

    # H1 渐变近似：多行时按比例分配渐变色
    gradient_colors = elem.get('gradientColors') if elem.get('tag') == 'h1' else None
    gc_start = parse_color(gradient_colors[0]) if gradient_colors else None
    gc_end = parse_color(gradient_colors[1]) if gradient_colors else None
    total_lines = len(lines)

    # li 元素：添加蓝色 ▶ 前缀
    is_li = elem.get('tag') == 'li'
    li_bullet_color = 'rgb(56, 139, 253)'  # --accent-blue

    for line_idx, line_segs in enumerate(lines):
        p = tf.add_paragraph() if line_idx > 0 else tf.paragraphs[0]
        apply_para_format(p, s)
        # justifyContent: center → override paragraph alignment to CENTER
        # (e.g. numbered badge circles, centered flex containers)
        if justify_content in ('center', 'space-around', 'space-evenly'):
            p.alignment = PP_ALIGN.CENTER

        # li 前缀（只在第一行加）
        if is_li and line_idx == 0:
            bullet_run = p.add_run()
            apply_run(bullet_run, '▶ ', li_bullet_color, font_size_pt * 0.7, '400')

        # 计算当前行的渐变色
        if gc_start and gc_end and total_lines > 1:
            t = line_idx / (total_lines - 1)
            grad_rgb = interpolate_color(gc_start, gc_end, t)
            override_color = 'rgb({},{},{})'.format(*grad_rgb)
        elif gc_start and gc_end:
            override_color = gradient_colors[1]
        else:
            override_color = None

        for seg in line_segs:
            if not seg['text']:
                continue
            run = p.add_run()
            # 如果有渐变色覆盖，使用覆盖色；否则用 seg 原色
            color = override_color or seg['color']
            # 使用 segment 级别的 bold（如果有），否则用元素级别
            seg_weight = 'bold' if seg.get('bold') else font_weight
            seg_fs_raw = seg.get('fontSize', '')
            seg_font_size_pt = px_to_pt(seg_fs_raw) if seg_fs_raw and 'px' in str(seg_fs_raw) else font_size_pt
            apply_run(run, seg['text'], color, seg_font_size_pt, seg_weight, text_transform,
                      font_family=font_family, letter_spacing=letter_spacing,
                      strike=seg.get('strike', False), bg_color_str=seg.get('bgColor'))


def export_shape_with_text(slide, elem: Dict[str, Any], bg_color=None):
    b = elem['bounds']
    s = elem['styles']

    border_radius = s.get('borderRadius', '')
    radius_px = 0.0
    if border_radius and border_radius != '0px':
        m = re.search(r'([\d.]+)px', border_radius)
        if m:
            radius_px = float(m.group(1))

    # [X] 非徽章形状使用最小圆角
    height_px = b['height'] * 108.0
    if radius_px > 0 and radius_px < height_px * 0.4:
        radius_px = min(radius_px, 6.0)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius_px > 0 else MSO_SHAPE.RECTANGLE,
        Inches(b['x']), Inches(b['y']),
        Inches(b['width']), Inches(b['height'])
    )
    if radius_px > 0:
        set_roundrect_adj(shape, radius_px, b['width'], b['height'])

    bg_rgb = parse_color(s.get('backgroundColor', ''))
    if bg_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*bg_rgb)
    else:
        grad_fill = gradient_to_solid(s.get('backgroundImage', ''))
        if grad_fill:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*grad_fill)
        else:
            shape.fill.background()

    border_str = s.get('border', '')
    if border_str and 'none' not in border_str:
        m = re.search(r'rgb\((\d+),\s*(\d+),\s*(\d+)\)', border_str)
        if m:
            shape.line.color.rgb = RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3)))
    else:
        shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    from pptx.enum.text import MSO_AUTO_SIZE
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(10)
    tf.margin_bottom = Pt(10)

    segments = elem.get('segments', [])
    font_size_pt = px_to_pt(s.get('fontSize', '16px'))
    font_weight = s.get('fontWeight', '400')
    font_family = s.get('fontFamily', '')
    letter_spacing = s.get('letterSpacing', '')
    text_transform = elem.get('textTransform', 'none')

    if not segments:
        raw = (elem.get('text', '') or '').strip()
        segments = [{'text': raw, 'color': s.get('color', '')}]

    lines = segments_to_lines(segments)
    for line_idx, line_segs in enumerate(lines):
        p = tf.add_paragraph() if line_idx > 0 else tf.paragraphs[0]
        apply_para_format(p, s)
        for seg in line_segs:
            if not seg['text']:
                continue
            run = p.add_run()
            seg_weight = 'bold' if seg.get('bold') else font_weight
            apply_run(run, seg['text'], seg['color'], font_size_pt, seg_weight, text_transform,
                      font_family=font_family, letter_spacing=letter_spacing,
                      strike=seg.get('strike', False), bg_color_str=seg.get('bgColor'))


def export_table_element(slide, elem: Dict[str, Any]):
    """Render HTML table as individual cell rectangles + text frames."""
    rows = elem.get('rows', [])
    if not rows:
        return

    for row_data in rows:
        for cell in row_data['cells']:
            cb = cell['bounds']
            cs = cell['styles']

            # Skip zero-size cells
            if cb['width'] < 0.01 or cb['height'] < 0.01:
                continue

            # Cell background rectangle
            cell_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(cb['x']), Inches(cb['y']),
                Inches(cb['width']), Inches(cb['height'])
            )

            bg_rgb = parse_color(cs.get('backgroundColor', ''))
            if bg_rgb:
                cell_shape.fill.solid()
                cell_shape.fill.fore_color.rgb = RGBColor(*bg_rgb)
            else:
                cell_shape.fill.background()

            # No rectangle border (avoids vertical dividers)
            suppress_line(cell_shape)

            # border-bottom only: draw as a separate thin rectangle
            border_bottom = cs.get('borderBottom', '')
            if border_bottom and 'none' not in border_bottom and not border_bottom.startswith('0px'):
                m = re.search(r'rgb\((\d+),\s*(\d+),\s*(\d+)\)', border_bottom)
                if m:
                    divider = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(cb['x']), Inches(cb['y'] + cb['height'] - 0.005),
                        Inches(cb['width']), Inches(0.005)
                    )
                    divider.fill.solid()
                    divider.fill.fore_color.rgb = RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3)))
                    suppress_line(divider)

            # Text content
            segments = cell.get('segments', [])
            text = cell.get('text', '').strip()
            if not segments and text:
                segments = [{'text': text, 'color': cs.get('color', '')}]
            if not segments:
                continue

            font_size_pt = px_to_pt(cs.get('fontSize', '14px'))
            font_weight = cs.get('fontWeight', '400')
            font_family = cs.get('fontFamily', '')
            letter_spacing = cs.get('letterSpacing', '')
            if cell['isHeader']:
                font_weight = 'bold'

            tf = cell_shape.text_frame
            tf.word_wrap = True
            tf.margin_left = Pt(6)
            tf.margin_right = Pt(6)
            tf.margin_top = Pt(4)
            tf.margin_bottom = Pt(4)

            lines = segments_to_lines(segments)
            for line_idx, line_segs in enumerate(lines):
                p = tf.add_paragraph() if line_idx > 0 else tf.paragraphs[0]
                align = cs.get('textAlign', 'left')
                if align == 'center':
                    p.alignment = PP_ALIGN.CENTER
                elif align == 'right':
                    p.alignment = PP_ALIGN.RIGHT
                for seg in line_segs:
                    if not seg['text']:
                        continue
                    run = p.add_run()
                    seg_weight = 'bold' if seg.get('bold') else font_weight
                    apply_run(run, seg['text'], seg['color'], font_size_pt, seg_weight,
                              font_family=font_family, letter_spacing=letter_spacing,
                              strike=seg.get('strike', False))


def apply_slide_gradient_bg(slide, color1: tuple, color2: tuple, angle_deg: float = 135.0):
    """Apply a linear gradient fill to the slide background using python-pptx native API.
    angle_deg: CSS gradient angle (135 = top-left to bottom-right diagonal).
    """
    try:
        fill = slide.background.fill
        fill.gradient()
        fill.gradient_angle = angle_deg
        stops = fill.gradient_stops
        stops[0].position = 0.0
        stops[0].color.rgb = RGBColor(*color1)
        stops[1].position = 1.0
        stops[1].color.rgb = RGBColor(*color2)
    except Exception:
        # Fallback: use first color as solid
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*color1)


def add_slide_chrome(slide, slide_idx: int, slide_count: int,
                     slide_w_in: float, slide_h_in: float, px_per_in: float = 108.0):
    """Add page counter (top-left) and pagination dots (bottom-center) to each slide."""
    # ── Page counter ──────────────────────────────────────────────────────────
    # CSS: position:fixed; top:24px; left:36px; font-size:12px; font-weight:700
    counter_x = 36 / px_per_in
    counter_y = 24 / px_per_in
    txBox = slide.shapes.add_textbox(
        Inches(counter_x), Inches(counter_y), Inches(1.0), Inches(0.22)
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{slide_idx + 1:02d} / {slide_count:02d}"
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(100, 116, 139)   # --text-muted: #64748b
    run.font.bold = True

    # ── Pagination dots ───────────────────────────────────────────────────────
    # CSS: position:fixed; bottom:24px; left:50%; gap:8px
    # Dot: height:6px; inactive width:6px; active width:28px
    dot_h = 6 / px_per_in
    dot_inactive_w = 6 / px_per_in
    dot_active_w = 28 / px_per_in
    gap = 8 / px_per_in
    total_w = (dot_inactive_w * (slide_count - 1) + dot_active_w
               + gap * (slide_count - 1))
    start_x = slide_w_in / 2 - total_w / 2
    dot_y = slide_h_in - 24 / px_per_in - dot_h

    x = start_x
    for j in range(slide_count):
        is_active = (j == slide_idx)
        w = dot_active_w if is_active else dot_inactive_w
        dot_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(dot_y), Inches(w), Inches(dot_h)
        )
        dot_shape.fill.solid()
        if is_active:
            dot_shape.fill.fore_color.rgb = RGBColor(37, 99, 235)    # #2563eb
        else:
            dot_shape.fill.fore_color.rgb = RGBColor(147, 197, 253)  # inactive
        suppress_line(dot_shape)
        x += w + gap


def render_fixed_chrome(slide, fixed_chrome: dict, slide_idx: int = 1, slide_count: int = 1, slide_w_in: float = 13.33):
    """Render fixed-position chrome elements (nav-dots, progress-bar) extracted from the HTML DOM.
    slide_idx: 1-based current slide index (used to compute active dot and progress bar width).
    """
    if not fixed_chrome:
        return

    # ── Nav dots (right-side vertical dots) ──────────────────────────────────
    nav_dots = fixed_chrome.get('navDots', [])
    if nav_dots:
        # Detect active color from captured state (slide 1 has dot 0 active); fallback to yellow
        active_color = None
        border_color_rgb = (17, 17, 17)  # --border fallback
        for dot in nav_dots:
            bc = parse_color(dot.get('border', ''))
            if bc:
                border_color_rgb = bc
            bg = parse_color(dot.get('bg', ''))
            if bg and active_color is None:
                active_color = bg
        if active_color is None:
            active_color = (255, 225, 77)  # --yellow fallback

        active_idx = slide_idx - 1  # 0-based
        for dot_i, dot in enumerate(nav_dots):
            dot_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(dot['x']), Inches(dot['y']),
                Inches(max(dot['w'], 0.001)), Inches(max(dot['h'], 0.001))
            )
            if dot_i == active_idx:
                dot_shape.fill.solid()
                dot_shape.fill.fore_color.rgb = RGBColor(*active_color)
            else:
                dot_shape.fill.background()
            dot_shape.line.color.rgb = RGBColor(*border_color_rgb)
            dot_shape.line.width = Pt(1.5)

    # ── Progress bar (thin line at top) ──────────────────────────────────────
    # Always compute width from slide_idx/slide_count (HTML scroll state may not update in headless)
    pb = fixed_chrome.get('progressBar')
    if pb:
        pb_h = pb.get('h', 0.046)
        pb_color_rgb = parse_color(pb.get('bg', ''))
        if not pb_color_rgb:
            pb_color_rgb = (255, 60, 126)  # --pink fallback
        pb_w = slide_idx / slide_count * slide_w_in
        if pb_w > 0.01:
            pb_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0),
                Inches(pb_w), Inches(max(pb_h, 0.05))
            )
            pb_shape.fill.solid()
            pb_shape.fill.fore_color.rgb = RGBColor(*pb_color_rgb)
            suppress_line(pb_shape)


def _find_and_launch_browser(playwright):
    """
    Try system-installed browsers first (no download needed), fall back to
    Playwright's bundled Chromium as a last resort.

    Linux/Docker/CI: adds --no-sandbox because root/container environments
    block Chrome's kernel sandbox by default.
    """
    import platform
    is_linux = platform.system() == 'Linux'
    extra_args = ['--no-sandbox', '--disable-setuid-sandbox'] if is_linux else []

    for channel in ['chrome', 'msedge', 'chromium']:
        try:
            browser = playwright.chromium.launch(channel=channel, headless=True, args=extra_args)
            print(f"  Using browser: {channel}")
            return browser
        except Exception:
            continue

    # Last resort: Playwright's own Chromium (requires: playwright install chromium)
    try:
        browser = playwright.chromium.launch(headless=True, args=extra_args)
        print("  Using browser: playwright-chromium")
        return browser
    except Exception as e:
        print("\nNo browser found.")
        if is_linux:
            print("  Linux options (pick one):")
            print("    apt install chromium-browser      # system package")
            print("    playwright install chromium        # self-contained")
        else:
            print("  Install Google Chrome: https://www.google.com/chrome/")
            print("  Or run: playwright install chromium")
        raise SystemExit(1) from e


def _screenshot_slide(page, slide_index: int, width: int, height: int) -> Optional[bytes]:
    """Screenshot one slide via scrollIntoView + getBoundingClientRect clip."""
    try:
        page.evaluate(f"""
            () => {{
                const s = document.querySelectorAll('.slide')[{slide_index}];
                if (s) s.scrollIntoView({{behavior: 'instant', block: 'start'}});
            }}
        """)
        page.wait_for_timeout(80)
        box = page.evaluate(f"""
            () => {{
                const r = document.querySelectorAll('.slide')[{slide_index}].getBoundingClientRect();
                return {{x: r.left, y: r.top, w: r.width, h: r.height}};
            }}
        """)
        if not box or box['w'] <= 0:
            return None
        x, y, w, h = box['x'], box['y'], box['w'], box['h']
        clip = {'x': max(0.0, x), 'y': max(0.0, y),
                'width': min(w, width - max(0.0, x)),
                'height': min(h, height - max(0.0, y))}
        if clip['width'] < 10 or clip['height'] < 10:
            return None
        return page.screenshot(type='png', clip=clip)
    except Exception:
        return None


def _save_preview_grid(screenshots: List[Tuple[int, bytes]], output_path: Path) -> Optional[Path]:
    """Assemble slide thumbnails into a labeled grid PNG."""
    try:
        from PIL import Image, ImageDraw
        THUMB_W, LABEL_H, PAD = 320, 22, 4
        thumbs = []
        for label, png in screenshots:
            img = Image.open(io.BytesIO(png)).convert('RGB')
            ratio = img.height / img.width
            img = img.resize((THUMB_W, int(THUMB_W * ratio)), Image.LANCZOS)
            thumbs.append((label, img))
        if not thumbs:
            return None
        n = len(thumbs)
        tw, th = thumbs[0][1].size
        grid = Image.new('RGB', (n * THUMB_W + (n - 1) * PAD, th + LABEL_H), (32, 32, 32))
        draw = ImageDraw.Draw(grid)
        for j, (label, thumb) in enumerate(thumbs):
            x = j * (THUMB_W + PAD)
            grid.paste(thumb, (x, 0))
            draw.text((x + THUMB_W // 2, th + 3), f"Slide {label}",
                      fill=(200, 200, 200), anchor='mt')
        preview_path = output_path.with_name(output_path.stem + '-preview.png')
        grid.save(str(preview_path))
        return preview_path
    except Exception:
        return None


def _validate_pptx(pptx_path: Path, expected_slides: int) -> List[str]:
    """Quick structural validation after save. Returns list of warning strings."""
    issues = []
    try:
        from pptx import Presentation as _Prs
        prs2 = _Prs(str(pptx_path))
        actual = len(prs2.slides)
        if actual != expected_slides:
            issues.append(f"slide count mismatch: expected {expected_slides}, got {actual}")
        for i, slide in enumerate(prs2.slides):
            try:
                _ = len(slide.shapes)
            except Exception as e:
                issues.append(f"slide {i+1} unreadable: {e}")
    except Exception as e:
        issues.append(f"failed to open PPTX: {e}")
    return issues


def export_native(html_path, output_path=None, width=1440, height=810):
    html_path = Path(html_path).resolve()
    if not html_path.exists():
        print(f"Error: {html_path}")
        sys.exit(1)

    output_path = Path(output_path) if output_path else html_path.with_suffix('.pptx')
    print(f"导出（native v5）: {html_path.name}")

    with sync_playwright() as p:
        browser = _find_and_launch_browser(p)
        page = browser.new_page(viewport={"width": width, "height": height})
        page.goto(f"file://{html_path}", wait_until="networkidle")
        page.wait_for_timeout(500)

        slide_count = page.evaluate("document.querySelectorAll('.slide').length")
        if slide_count == 0:
            print("未找到 .slide 元素")
            browser.close()
            return

        print(f"找到 {slide_count} 张幻灯片")
        inject_visible(page)
        page.wait_for_timeout(200)

        _preview_indices = sorted({0, slide_count // 3, 2 * slide_count // 3, slide_count - 1})
        preview_screenshots: List[Tuple[int, bytes]] = []

        prs = Presentation()
        first = extract_slide_elements(page, 0)
        if first.get('slideSize'):
            prs.slide_width = Inches(first['slideSize']['width'])
            prs.slide_height = Inches(first['slideSize']['height'])
        else:
            prs.slide_width = Inches(15.0)
            prs.slide_height = Inches(9.375)

        blank_layout = prs.slide_layouts[6]
        max_w = prs.slide_width / 914400

        slide_h_in = prs.slide_height / 914400

        for i in range(slide_count):
            print(f"  [{i+1}/{slide_count}] 处理中...")
            data = extract_slide_elements(page, i)
            slide = prs.slides.add_slide(blank_layout)

            if data.get('bgGradient'):
                c1, c2 = data['bgGradient']
                try:
                    apply_slide_gradient_bg(slide, c1, c2, angle_deg=135.0)
                except Exception:
                    if data['background']:
                        r, g, b = data['background']
                        slide.background.fill.solid()
                        slide.background.fill.fore_color.rgb = RGBColor(r, g, b)
            elif data['background']:
                r, g, b = data['background']
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(r, g, b)

            # Grid background overlay (CSS double linear-gradient grid pattern)
            if data.get('gridBg'):
                grid_info = data['gridBg']
                add_grid_background(slide, max_w, slide_h_in,
                                    grid_info['color'], grid_info['sizePx'])

            # Pre-pass: sync background shape height with adjacent text's naturalHeight,
            # and apply a PPTX font-metrics correction factor for multi-line elements.
            #
            # Problem: for TEXT_TAG elements with background (callouts, info boxes), the JS
            # emits a shape (background/border bar) followed by a text element at the same
            # y-position and height. In the PPTX, the text box uses spAutoFit which expands
            # it at render time in PowerPoint/Keynote when CJK font metrics cause more line-
            # wrapping than Chrome. The background shape has no auto-sizing and stays fixed,
            # so text visually overflows below the orange/blue background box.
            #
            # Fix: detect these shape+text pairs, apply a 1.3x PPTX correction factor to
            # both shapes for multi-line elements, and mark the text element so that
            # export_text_element switches from SHAPE_TO_FIT_TEXT to NONE auto-size
            # (since the height is pre-corrected, auto-expansion is no longer needed).
            elems = data['elements']
            for _i in range(len(elems) - 1):
                _s, _t = elems[_i], elems[_i + 1]
                if (_s.get('type') == 'shape' and _t.get('type') == 'text' and
                        abs(_s['bounds']['y'] - _t['bounds']['y']) < 0.05 and
                        abs(_s['bounds']['height'] - _t['bounds']['height']) < 0.05):
                    _nat = _t.get('naturalHeight', _t['bounds']['height'])
                    _base = max(_nat, _s['bounds']['height'])
                    # Only correct multi-line elements (estimated_lines >= 2).
                    # Single-line badge/pill elements must keep their original height.
                    _t_font_pt = px_to_pt(_t['styles'].get('fontSize', '16px'))
                    _t_line_h = _t_font_pt / 72.0 * 1.2
                    _t_est_lines = _base / max(_t_line_h, 0.001)
                    if _t_est_lines >= 2.0:
                        # Apply 1.3x PPTX correction: accounts for wider CJK character
                        # metrics in PowerPoint/Keynote vs Chrome rendering, which causes
                        # additional line-wrapping and requires extra vertical space.
                        _corrected = _base * 1.3
                        _s['bounds']['height'] = _corrected
                        _t['bounds']['height'] = _corrected
                        _t['pptx_height_corrected'] = True
                    elif _nat > _s['bounds']['height'] * 1.05:
                        # True overflow (naturalHeight > bounds): sync heights without factor.
                        _s['bounds']['height'] = _nat
                        _t['bounds']['height'] = _nat

            # Second pre-pass: push elements immediately to the right of large single-line
            # text (titles/headings). In Keynote/PowerPoint, CJK characters are ~8% wider
            # than in Chrome, so large-font title text visually overflows into adjacent
            # badges/pills (e.g., "21 个主题适应所有场景" + "按内容类型自动匹配" pill on P4).
            # Fix: for each large (font > 24pt) single-line text element, detect any element
            # horizontally adjacent (gap 0–0.3in) in the same y-band, and push it right by
            # title_width × CJK_H_FACTOR so the PPTX visual gap matches the HTML gap.
            _CJK_H_FACTOR = 0.15
            for _i in range(len(elems)):
                _el = elems[_i]
                if _el.get('type') != 'text':
                    continue
                _fp = px_to_pt(_el['styles'].get('fontSize', '16px'))
                if _fp <= 24.0:
                    continue
                _lh = _fp / 72.0 * 1.2
                _est = _el['bounds']['height'] / max(_lh, 0.001)
                if _est >= 2.0:
                    continue
                # Large single-line text: compute extra PPTX horizontal space needed
                _orig_right = _el['bounds']['x'] + _el['bounds']['width']
                _extra = _el['bounds']['width'] * _CJK_H_FACTOR
                for _j in range(len(elems)):
                    if _j == _i:
                        continue
                    _other = elems[_j]
                    _gap = _other['bounds']['x'] - _orig_right
                    # Same y-band (vertically overlapping) and small horizontal gap
                    _y_overlap = (abs(_other['bounds']['y'] - _el['bounds']['y'])
                                  < _el['bounds']['height'])
                    if 0 <= _gap <= 0.3 and _y_overlap:
                        _other['bounds']['x'] += _extra

            for elem in data['elements']:
                try:
                    # Clamp width
                    if elem['bounds']['x'] < max_w and elem['bounds']['width'] > max_w - elem['bounds']['x']:
                        elem['bounds']['width'] = max_w - elem['bounds']['x']

                    elem_type = elem.get('type', 'text')

                    if elem_type == 'shape':
                        export_shape_background(slide, elem, slide_bg=data['background'] or (255, 255, 255))
                    elif elem_type == 'image':
                        export_raster_element(page, slide, elem)
                    elif elem_type == 'table':
                        export_table_element(slide, elem)
                    else:
                        export_text_element(slide, elem, data['background'])
                except Exception as e:
                    print(f"    警告: {e}")

            px_per_in = width / max_w
            if not data.get('hasOwnChrome'):
                add_slide_chrome(slide, i, slide_count, max_w, slide_h_in, px_per_in)
            render_fixed_chrome(slide, data.get('fixedChrome'), i, slide_count, max_w)

            # Collect thumbnail for preview grid
            if i in _preview_indices:
                ss = _screenshot_slide(page, i, width, height)
                if ss:
                    preview_screenshots.append((i + 1, ss))

        browser.close()

    prs.save(str(output_path))
    print(f"Saved: {output_path}  ({slide_count} 张幻灯片)")

    # Structural validation
    for issue in _validate_pptx(output_path, slide_count):
        print(f"  ⚠ {issue}")

    # Preview grid
    preview_path = _save_preview_grid(preview_screenshots, output_path)
    if preview_path:
        print(f"Preview: {preview_path}")

    return output_path


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("html")
    parser.add_argument("output", nargs="?")
    parser.add_argument("--width", type=int, default=1440)
    parser.add_argument("--height", type=int, default=810)
    args = parser.parse_args()
    export_native(args.html, args.output, args.width, args.height)


if __name__ == "__main__":
    main()
