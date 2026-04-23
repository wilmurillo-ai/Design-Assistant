#!/usr/bin/env python3
"""验证 PPTX 文件的样式是否正确"""

import sys
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

def verify_pptx(pptx_path):
    prs = Presentation(pptx_path)

    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n=== 幻灯片 {slide_idx + 1} ===")

        for shape_idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                print(f"\n形状 {shape_idx}:")
                print(f"  位置: ({shape.left}, {shape.top})")
                print(f"  大小: ({shape.width}, {shape.height})")
                print(f"  形状类型: {shape.shape_type}")

                for para_idx, para in enumerate(shape.text_frame.paragraphs):
                    print(f"\n  段落 {para_idx}:")
                    print(f"    对齐: {para.alignment}")

                    for run_idx, run in enumerate(para.runs):
                        text = run.text[:50] + "..." if len(run.text) > 50 else run.text
                        print(f"\n    文本块 {run_idx}: '{text}'")
                        print(f"      字体: {run.font.name}")
                        print(f"      字号: {run.font.size}")
                        print(f"      粗体: {run.font.bold}")
                        print(f"      斜体: {run.font.italic}")
                        try:
                            if run.font.color.type is not None:
                                print(f"      颜色: {run.font.color.rgb}")
                        except:
                            pass

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法: python verify-pptx.py <pptx文件>")
        sys.exit(1)

    verify_pptx(sys.argv[1])