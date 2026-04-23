"""
Tests for export-pptx.py

Run: python -m pytest tests/test_pptx.py -v
     python tests/run_tests.py
"""

import subprocess
import sys
import tempfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

FIXTURES = Path(__file__).parent / "fixtures"
SCRIPT = Path(__file__).parent.parent / "scripts" / "export-pptx.py"


def run_export(args: list, cwd=None) -> subprocess.CompletedProcess:
    return subprocess.run(
        [sys.executable, str(SCRIPT)] + args,
        capture_output=True, text=True, cwd=cwd
    )


# ─── Happy path ──────────────────────────────────────────────────────────────

class TestPptxExport:

    def test_slide_count_matches(self, tmp_path):
        """3 .slide elements → PPTX should have exactly 3 slides."""
        out = tmp_path / "out.pptx"
        r = run_export([str(FIXTURES / "simple_slides.html"), str(out)])
        assert r.returncode == 0, f"Script failed:\n{r.stderr}"
        assert out.exists(), "Output PPTX not created"
        prs = Presentation(str(out))
        assert len(prs.slides) == 3

    def test_output_file_has_content(self, tmp_path):
        """Output PPTX should be larger than 10 KB (proves real screenshots)."""
        out = tmp_path / "out.pptx"
        run_export([str(FIXTURES / "simple_slides.html"), str(out)])
        assert out.stat().st_size > 10_000, "PPTX is suspiciously small"

    def test_default_output_path(self, tmp_path):
        """No output arg → PPTX is saved alongside the HTML with same stem."""
        import shutil
        html = tmp_path / "my_deck.html"
        shutil.copy(FIXTURES / "simple_slides.html", html)
        r = run_export([str(html)])
        assert r.returncode == 0, f"Script failed:\n{r.stderr}"
        pptx = tmp_path / "my_deck.pptx"
        assert pptx.exists(), f"Expected {pptx} to be created"

    def test_default_aspect_ratio_matches_viewport(self, tmp_path):
        """Default 1440×900 viewport → slide ratio should match 1440/900 = 1.6 (8:5)."""
        out = tmp_path / "out.pptx"
        run_export([str(FIXTURES / "simple_slides.html"), str(out)])
        prs = Presentation(str(out))
        ratio = prs.slide_width / prs.slide_height
        expected = 1440 / 900  # 1.6 (8:5), not 16:9
        assert abs(ratio - expected) < 0.01, f"Expected {expected:.4f}, got {ratio:.4f}"

    def test_custom_dimensions(self, tmp_path):
        """--width 1920 --height 1080 → slide ratio should be 16:9."""
        out = tmp_path / "out.pptx"
        r = run_export([
            str(FIXTURES / "simple_slides.html"), str(out),
            "--width", "1920", "--height", "1080"
        ])
        assert r.returncode == 0, f"Script failed:\n{r.stderr}"
        prs = Presentation(str(out))
        assert len(prs.slides) == 3
        ratio = prs.slide_width / prs.slide_height
        assert abs(ratio - 16 / 9) < 0.01

    def test_custom_4x3_ratio(self, tmp_path):
        """--width 1024 --height 768 → slide ratio should be 4:3."""
        out = tmp_path / "out.pptx"
        run_export([
            str(FIXTURES / "simple_slides.html"), str(out),
            "--width", "1024", "--height", "768"
        ])
        prs = Presentation(str(out))
        ratio = prs.slide_width / prs.slide_height
        assert abs(ratio - 4 / 3) < 0.01, f"Expected 4:3, got {ratio:.4f}"

    def test_stdout_reports_slide_count(self, tmp_path):
        """Script stdout should mention the number of slides found."""
        out = tmp_path / "out.pptx"
        r = run_export([str(FIXTURES / "simple_slides.html"), str(out)])
        assert "3" in r.stdout, f"Expected slide count in stdout: {r.stdout}"

    def test_scroll_snap_not_misaligned(self, tmp_path):
        """HTML with scroll-snap-type: y mandatory should export without slide misalignment.

        Regression test for: PPTX slides misaligned when HTML uses scroll-snap.

        Bug cause: window.scrollTo() with scroll-snap snaps to wrong position,
        causing screenshots to capture wrong slides (e.g., slide 2 shows slide 6 content).

        Fix: Use locator().nth(i).screenshot() to capture each slide element directly,
        bypassing scroll positioning entirely.
        """
        out = tmp_path / "out.pptx"
        r = run_export([str(FIXTURES / "scroll_snap_slides.html"), str(out)])
        assert r.returncode == 0, f"Script failed:\n{r.stderr}"
        assert out.exists(), "Output PPTX not created"
        prs = Presentation(str(out))
        # 5 slides in scroll_snap_slides.html
        assert len(prs.slides) == 5, \
            f"Expected 5 slides, got {len(prs.slides)} — possible scroll-snap misalignment"


# ─── Error cases ─────────────────────────────────────────────────────────────

class TestPptxErrorCases:

    def test_no_slide_elements_exits_gracefully(self, tmp_path):
        """HTML with no .slide elements → script exits 0 with 'Nothing to export'."""
        out = tmp_path / "out.pptx"
        r = run_export([str(FIXTURES / "no_slides.html"), str(out)])
        # Should exit cleanly (not crash), output file should not exist
        assert r.returncode == 0
        assert not out.exists(), "No PPTX should be created when there are no slides"
        assert "Nothing" in r.stdout or "0 slides" in r.stdout.lower() or "No .slide" in r.stdout

    def test_missing_file_exits_nonzero(self, tmp_path):
        """Nonexistent input file → non-zero exit code."""
        out = tmp_path / "out.pptx"
        r = run_export([str(tmp_path / "ghost.html"), str(out)])
        assert r.returncode != 0, "Expected non-zero exit for missing file"


# ─── Native mode ─────────────────────────────────────────────────────────────

class TestNativeMode:
    """Tests for --mode native (editable text/shapes via BS4 parsing)."""

    def test_native_mode_creates_pptx(self, tmp_path):
        """--mode native should create a valid PPTX file."""
        out = tmp_path / "out.pptx"
        r = run_export([str(FIXTURES / "simple_slides.html"), str(out), "--mode", "native"])
        assert r.returncode == 0, f"Script failed:\n{r.stderr}"
        assert out.exists(), "PPTX not created in native mode"
        assert out.stat().st_size > 5_000, "PPTX is suspiciously small"

    def test_native_mode_slide_count_matches(self, tmp_path):
        """--mode native: 3 .slide elements → PPTX should have exactly 3 slides."""
        out = tmp_path / "out.pptx"
        r = run_export([str(FIXTURES / "simple_slides.html"), str(out), "--mode", "native"])
        assert r.returncode == 0, f"Script failed:\n{r.stderr}"
        prs = Presentation(str(out))
        assert len(prs.slides) == 3, f"Expected 3 slides, got {len(prs.slides)}"

    def test_native_mode_importlib_dispatch(self, tmp_path):
        """export-pptx.py must dispatch to export-native-pptx.py via importlib.

        Regression test for: native mode silently fell back to image mode
        because 'from export_native_pptx import ...' fails on hyphenated filenames.
        Fix: importlib.util.spec_from_file_location loads the file by path.
        """
        out = tmp_path / "out.pptx"
        r = run_export([str(FIXTURES / "simple_slides.html"), str(out), "--mode", "native"])
        # If importlib dispatch is broken, the script either crashes or
        # falls back to image mode (which would not print native-mode output).
        assert r.returncode == 0, f"importlib dispatch failed:\n{r.stderr}"
        # export-native-pptx.py prints output that doesn't contain "Capturing..."
        # (which is image-mode language). Native mode uses different log messages.
        assert "Capturing" not in r.stdout, \
            "Output looks like image mode — native mode may not have been dispatched"

    def test_native_mode_missing_file_exits_nonzero(self, tmp_path):
        """--mode native with nonexistent input → non-zero exit code."""
        out = tmp_path / "out.pptx"
        r = run_export([str(tmp_path / "ghost.html"), str(out), "--mode", "native"])
        assert r.returncode != 0, "Expected non-zero exit for missing file in native mode"

    def test_both_modes_same_slide_count(self, tmp_path):
        """image mode and native mode must produce the same slide count."""
        out_image = tmp_path / "image.pptx"
        out_native = tmp_path / "native.pptx"
        run_export([str(FIXTURES / "simple_slides.html"), str(out_image), "--mode", "image"])
        run_export([str(FIXTURES / "simple_slides.html"), str(out_native), "--mode", "native"])
        count_image = len(Presentation(str(out_image)).slides)
        count_native = len(Presentation(str(out_native)).slides)
        assert count_image == count_native, \
            f"Slide count mismatch: image={count_image}, native={count_native}"

    def test_native_mode_exports_img_svg_and_background_images(self, tmp_path):
        """Native mode should keep raster template elements instead of dropping or crashing on them."""
        out = tmp_path / "raster-elements.pptx"
        r = run_export([str(FIXTURES / "native_raster_elements.html"), str(out), "--mode", "native"])
        assert r.returncode == 0, f"Script failed:\nSTDOUT:\n{r.stdout}\nSTDERR:\n{r.stderr}"

        prs = Presentation(str(out))
        assert len(prs.slides) == 3

        picture_counts = [
            sum(1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
            for slide in prs.slides
        ]
        assert picture_counts[0] >= 1, "Expected native mode to export <img> as a picture"
        assert picture_counts[1] >= 1, "Expected native mode to export inline <svg> as a picture"
        assert picture_counts[2] >= 1, "Expected native mode to export background-image:url(...) containers as pictures"

    def test_native_mode_respects_export_progress_opt_out(self, tmp_path):
        """data-export-progress=false should suppress PPT progress-bar and nav-dot rendering."""
        out = tmp_path / "progress-opt-out.pptx"
        r = run_export([str(FIXTURES / "native_progress_opt_out.html"), str(out), "--mode", "native"])
        assert r.returncode == 0, f"Script failed:\nSTDOUT:\n{r.stdout}\nSTDERR:\n{r.stderr}"

        prs = Presentation(str(out))
        assert len(prs.slides) == 2

        for slide in prs.slides:
            top_thin_rects = [
                shape for shape in slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                and getattr(shape, "top", 0) == 0
                and getattr(shape, "height", 0) < Inches(0.08)
                and getattr(shape, "width", 0) > Inches(1.0)
            ]
            assert not top_thin_rects, "Expected no PPT progress bar when data-export-progress=false"

            nav_dot_shapes = [
                shape for shape in slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                and getattr(shape, "left", 0) > Inches(12.5)
                and Inches(0.05) <= getattr(shape, "width", 0) <= Inches(0.3)
                and Inches(0.05) <= getattr(shape, "height", 0) <= Inches(0.3)
            ]
            assert not nav_dot_shapes, "Expected no PPT nav dots when data-export-progress=false"
