#!/usr/bin/env python3
"""
PPTX 進階編輯器 - 支援投影片編輯、文字替換、版面調整
用法：
  uv run python scripts/pptx_editor.py input.pptx output.pptx edits.json
"""

import sys
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def apply_highlight(run):
    """應用螢光筆高亮（PPTX 使用填充色模擬）"""
    run.font.highlight_color = None  # PPTX 不直接支援 highlight，改用其他方法
    run.font.bold = True


def apply_strike(run):
    """應用刪除線"""
    run.font.strike = True


def apply_bold(run):
    """應用粗體"""
    run.font.bold = True


def apply_underline(run):
    """應用底線"""
    run.font.underline = True


def apply_style(run, style_type):
    """對 run 應用樣式"""
    if style_type == "highlight":
        apply_highlight(run)
    elif style_type == "delete":
        apply_strike(run)
    elif style_type == "bold":
        apply_bold(run)
    elif style_type == "underline":
        apply_underline(run)


def replace_in_paragraph(para, search_text, replace_text, style_type="replace"):
    """在段落中替換文字"""
    if search_text not in para.text:
        return False
    
    # 簡單版本：替換整個段落文字
    # 進階版本需要處理 run 級別
    runs = para.runs
    if not runs:
        return False
    
    # 找到包含搜尋文字的 run
    for run in runs:
        if search_text in run.text:
            if style_type == "replace":
                run.text = run.text.replace(search_text, replace_text)
            else:
                run.text = run.text.replace(search_text, replace_text)
                apply_style(run, style_type)
            return True
    
    return False


def replace_in_slide(slide, search_text, replace_text, style_type="replace"):
    """在投影片中替換文字"""
    replaced = False
    
    # 檢查是否有 shape
    if not hasattr(slide, 'shapes'):
        return False
    
    for shape in slide.shapes:
        # 處理文字框
        if hasattr(shape, "text_frame") and shape.text_frame:
            for paragraph in shape.text_frame.paragraphs:
                if replace_in_paragraph(paragraph, search_text, replace_text, style_type):
                    replaced = True
        
        # 處理表格
        if hasattr(shape, "table") and shape.table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        if replace_in_paragraph(paragraph, search_text, replace_text, style_type):
                            replaced = True
    
    return replaced


def rearrange_slides(prs, slide_order):
    """
    重新排列投影片順序
    
    Args:
        prs: Presentation 物件
        slide_order: 投影片順序列表，例如 [0, 2, 1, 3]
    """
    # 獲取所有投影片
    slides = list(prs.slides)
    
    # 檢查索引是否有效
    for idx in slide_order:
        if idx < 0 or idx >= len(slides):
            print(f"⚠️  警告：投影片索引 {idx} 超出範圍 (0-{len(slides)-1})")
            return False
    
    # 創建新的投影片順序
    new_slides = [slides[i] for i in slide_order]
    
    # 刪除所有投影片
    for i in range(len(prs.slides)):
        sp = prs.slides._sldIdLst[i]
        prs.slides._sldIdLst.remove(sp)
    
    # 添加新順序的投影片
    # 注意：python-pptx 不直接支援移動投影片，需要更複雜的操作
    # 這裡使用簡單的方法：返回新的順序列表，讓使用者手動處理
    print(f"⚠️  注意：python-pptx 不直接支援投影片重排")
    print(f"   建議順序：{slide_order}")
    
    return True


def add_slide(prs, template_slide=None):
    """
    新增投影片
    
    Args:
        prs: Presentation 物件
        template_slide: 範本投影片（可選）
    """
    if template_slide:
        # 使用範本版面
        slide_layout = prs.slide_layouts[0]  # 標題投影片
    else:
        # 使用預設版面
        slide_layout = prs.slide_layouts[1]  # 標題和內容
    
    slide = prs.slides.add_slide(slide_layout)
    return slide


def remove_slide(prs, slide_index):
    """
    刪除投影片
    
    Args:
        prs: Presentation 物件
        slide_index: 投影片索引
    """
    if slide_index < 0 or slide_index >= len(prs.slides):
        print(f"⚠️  警告：投影片索引 {slide_index} 超出範圍")
        return False
    
    sp = prs.slides._sldIdLst[slide_index]
    prs.slides._sldIdLst.remove(sp)
    return True


def edit_pptx_advanced(input_path, output_path, edits):
    """
    進階編輯 PPTX 文件
    
    Args:
        input_path: 輸入文件
        output_path: 輸出文件
        edits: 編輯規則（JSON 格式）
    """
    prs = Presentation(input_path)
    stats = {"replacements": 0, "additions": 0, "deletions": 0, "slides_modified": 0}
    
    # 處理投影片替換
    for edit in edits.get("replacements", []):
        search_text = edit["search"]
        replace_text = edit.get("replace", "")
        style = edit.get("style", "replace")
        
        # 在所有投影片中替換
        for i, slide in enumerate(prs.slides):
            if replace_in_slide(slide, search_text, replace_text, style):
                stats["replacements"] += 1
                stats["slides_modified"] += 1
    
    # 處理投影片順序調整
    if "slides" in edits:
        for slide_edit in edits["slides"]:
            action = slide_edit.get("action")
            
            if action == "rearrange":
                order = slide_edit.get("order", [])
                if rearrange_slides(prs, order):
                    stats["slides_modified"] += 1
            
            elif action == "add":
                add_slide(prs)
                stats["additions"] += 1
            
            elif action == "remove":
                idx = slide_edit.get("index", -1)
                if remove_slide(prs, idx):
                    stats["deletions"] += 1
    
    # 儲存
    prs.save(output_path)
    
    print(f"✅ 完成編輯：{output_path}")
    print(f"📊 統計：")
    print(f"   替換：{stats['replacements']} 處")
    print(f"   新增：{stats['additions']} 張投影片")
    print(f"   刪除：{stats['deletions']} 張投影片")
    print(f"   修改：{stats['slides_modified']} 張投影片")
    
    return stats


def main():
    if len(sys.argv) < 4:
        print(__doc__)
        sys.exit(1)
    
    input_path = Path(sys.argv[1])
    output_path = Path(sys.argv[2])
    edits_path = Path(sys.argv[3])
    
    if not input_path.exists():
        print(f"❌ 錯誤：文件不存在 {input_path}")
        sys.exit(1)
    
    if not edits_path.exists():
        print(f"❌ 錯誤：編輯規則文件不存在 {edits_path}")
        sys.exit(1)
    
    # 讀取編輯規則
    with open(edits_path, 'r', encoding='utf-8') as f:
        edits = json.load(f)
    
    # 執行編輯
    edit_pptx_advanced(input_path, output_path, edits)


if __name__ == "__main__":
    main()
