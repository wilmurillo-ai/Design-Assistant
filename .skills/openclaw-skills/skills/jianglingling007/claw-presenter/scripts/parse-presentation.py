#!/usr/bin/env python3
"""
Parse PPT/PPTX/PDF → per-page images + extract text/notes.

Usage:
  python3 parse-presentation.py <file> [--output-dir <dir>]

Output structure:
  <output-dir>/
    presentation.json    — metadata + per-page text/notes
    slides/
      001.png            — page 1 image
      002.png            — page 2 image
      ...
"""
import argparse
import json
import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path


def parse_pptx_text(filepath):
    """Extract per-slide text and notes from PPTX."""
    from pptx import Presentation
    prs = Presentation(filepath)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        title = ''
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        notes = ''
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        slides.append({
            'page': i,
            'title': title,
            'content': texts,
            'notes': notes,
            'has_notes': bool(notes),
            'image': f'slides/{i:03d}.png',
            'script': ''  # to be filled by agent
        })
    return slides


def parse_pdf_text(filepath):
    """Extract per-page text from PDF."""
    slides = []
    try:
        import pdfplumber
        with pdfplumber.open(filepath) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text() or ''
                lines = [l.strip() for l in text.split('\n') if l.strip()]
                title = lines[0] if lines else f'Page {i}'
                slides.append({
                    'page': i, 'title': title, 'content': lines,
                    'notes': '', 'has_notes': False,
                    'image': f'slides/{i:03d}.png', 'script': ''
                })
    except ImportError:
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(filepath)
            for i, page in enumerate(reader.pages, 1):
                text = page.extract_text() or ''
                lines = [l.strip() for l in text.split('\n') if l.strip()]
                title = lines[0] if lines else f'Page {i}'
                slides.append({
                    'page': i, 'title': title, 'content': lines,
                    'notes': '', 'has_notes': False,
                    'image': f'slides/{i:03d}.png', 'script': ''
                })
        except ImportError:
            print('Error: need pdfplumber or PyPDF2', file=sys.stderr)
            sys.exit(1)
    return slides


def pdf_to_images_pdftoppm(pdf_path, slides_dir, dpi=150):
    """Convert PDF to images using pdftoppm CLI (more reliable than pdf2image for large files)."""
    pdftoppm = shutil.which('pdftoppm')
    if not pdftoppm:
        return False
    try:
        prefix = os.path.join(slides_dir, 'slide')
        result = subprocess.run(
            [pdftoppm, '-png', '-r', str(dpi), pdf_path, prefix],
            capture_output=True, timeout=600
        )
        if result.returncode != 0:
            print(f'Warning: pdftoppm failed: {result.stderr.decode()}', file=sys.stderr)
            return False
        # Rename slide-01.png → 001.png, slide-1.png → 001.png, etc.
        for fname in sorted(os.listdir(slides_dir)):
            if not fname.startswith('slide-') or not fname.endswith('.png'):
                continue
            # Extract page number from slide-01.png or slide-1.png
            num_str = fname.replace('slide-', '').replace('.png', '')
            try:
                num = int(num_str)
            except ValueError:
                continue
            new_name = f'{num:03d}.png'
            os.rename(os.path.join(slides_dir, fname), os.path.join(slides_dir, new_name))
        return True
    except subprocess.TimeoutExpired:
        print('Warning: pdftoppm timed out', file=sys.stderr)
        return False
    except Exception as e:
        print(f'Warning: pdftoppm error: {e}', file=sys.stderr)
        return False


def pdf_to_images_pdf2image(pdf_path, slides_dir, dpi=150):
    """Convert PDF to images using pdf2image (python library)."""
    try:
        from pdf2image import convert_from_path
        images = convert_from_path(pdf_path, dpi=dpi, fmt='png')
        for i, img in enumerate(images, 1):
            img.save(os.path.join(slides_dir, f'{i:03d}.png'))
        return True
    except Exception as e:
        print(f'Warning: pdf2image failed: {e}', file=sys.stderr)
        return False


def pptx_to_pdf(filepath):
    """Convert PPTX to PDF via LibreOffice. Returns PDF path or None."""
    soffice = shutil.which('soffice') or shutil.which('libreoffice')
    if not soffice:
        return None
    try:
        # Use a persistent temp dir (not auto-cleaned) so the PDF survives
        tmpdir = tempfile.mkdtemp(prefix='claw-presenter-')
        print(f'Converting PPTX → PDF via LibreOffice...', file=sys.stderr)
        result = subprocess.run(
            [soffice, '--headless', '--convert-to', 'pdf', '--outdir', tmpdir, filepath],
            capture_output=True, timeout=300
        )
        pdf_path = os.path.join(tmpdir, Path(filepath).stem + '.pdf')
        if os.path.exists(pdf_path):
            print(f'PDF created: {pdf_path}', file=sys.stderr)
            return pdf_path
        else:
            print(f'Warning: LibreOffice did not produce PDF. stdout={result.stdout.decode()}, stderr={result.stderr.decode()}', file=sys.stderr)
            return None
    except subprocess.TimeoutExpired:
        print('Warning: LibreOffice conversion timed out (300s)', file=sys.stderr)
        return None
    except Exception as e:
        print(f'Warning: LibreOffice conversion failed: {e}', file=sys.stderr)
        return None


def convert_to_images(filepath, slides_dir, ext):
    """Convert file to per-page PNG images."""
    os.makedirs(slides_dir, exist_ok=True)

    pdf_path = filepath if ext == '.pdf' else None

    # Step 1: For PPTX/PPT, convert to PDF first
    if ext in ('.pptx', '.ppt'):
        pdf_path = pptx_to_pdf(filepath)

    # Step 2: Convert PDF to images (prefer pdftoppm CLI, fallback to pdf2image)
    if pdf_path:
        if pdf_to_images_pdftoppm(pdf_path, slides_dir):
            # Clean up temp PDF if we created it
            if pdf_path != filepath:
                try:
                    os.unlink(pdf_path)
                    os.rmdir(os.path.dirname(pdf_path))
                except OSError:
                    pass
            return True
        if pdf_to_images_pdf2image(pdf_path, slides_dir):
            if pdf_path != filepath:
                try:
                    os.unlink(pdf_path)
                    os.rmdir(os.path.dirname(pdf_path))
                except OSError:
                    pass
            return True

    # Step 3: Fallback for PPTX — simple text-rendered slides
    if ext in ('.pptx', '.ppt'):
        print('Warning: falling back to simple text rendering', file=sys.stderr)
        return fallback_render_pptx(filepath, slides_dir)

    return False


def fallback_render_pptx(filepath, slides_dir):
    """Fallback: render simple text-based slide images."""
    try:
        from pptx import Presentation
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        return False

    prs = Presentation(filepath)
    w, h = 1280, 720

    for i, slide in enumerate(prs.slides, 1):
        img = Image.new('RGB', (w, h), '#1e1e2e')
        draw = ImageDraw.Draw(img)

        title = slide.shapes.title.text.strip() if slide.shapes.title else f'Slide {i}'
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape != getattr(slide.shapes, 'title', None):
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)

        try:
            ft = ImageFont.truetype('/System/Library/Fonts/PingFang.ttc', 36)
            fb = ImageFont.truetype('/System/Library/Fonts/PingFang.ttc', 22)
        except:
            ft = ImageFont.load_default()
            fb = ft

        draw.text((60, 40), title, fill='#e8e8ed', font=ft)
        y = 120
        for t in texts[:10]:
            if len(t) > 80: t = t[:77] + '...'
            draw.text((60, y), t, fill='#8888a0', font=fb)
            y += 36
        draw.text((w - 80, h - 50), str(i), fill='#55556a', font=fb)
        img.save(os.path.join(slides_dir, f'{i:03d}.png'))

    return True


def main():
    parser = argparse.ArgumentParser(description='Parse PPT/PDF → images + text')
    parser.add_argument('file', help='Path to PPT/PPTX/PDF')
    parser.add_argument('--output-dir', '-o', default=None)
    args = parser.parse_args()

    filepath = os.path.expanduser(args.file)
    if not os.path.exists(filepath):
        print(f'Error: file not found: {filepath}', file=sys.stderr)
        sys.exit(1)

    ext = Path(filepath).suffix.lower()
    basename = Path(filepath).stem

    if args.output_dir:
        output_dir = args.output_dir
    else:
        workspace = os.environ.get('OPENCLAW_WORKSPACE', os.getcwd())
        output_dir = os.path.join(workspace, 'presentations', basename)

    os.makedirs(output_dir, exist_ok=True)
    slides_dir = os.path.join(output_dir, 'slides')

    # Step 1: Convert to images
    convert_to_images(filepath, slides_dir, ext)

    # Step 2: Extract text/notes
    if ext in ('.pptx', '.ppt'):
        slides = parse_pptx_text(filepath)
    elif ext == '.pdf':
        slides = parse_pdf_text(filepath)
    else:
        print(f'Error: unsupported: {ext}', file=sys.stderr)
        sys.exit(1)

    # Step 3: For slides with notes, pre-fill script
    for s in slides:
        if s['has_notes']:
            s['script'] = s['notes']

    # Step 4: Save
    presentation = {
        'source': os.path.basename(filepath),
        'total_pages': len(slides),
        'output_dir': output_dir,
        'slides': slides
    }
    out_path = os.path.join(output_dir, 'presentation.json')
    with open(out_path, 'w', encoding='utf-8') as f:
        json.dump(presentation, f, ensure_ascii=False, indent=2)

    print(json.dumps({
        'status': 'ok',
        'output_dir': output_dir,
        'json_path': out_path,
        'total_pages': len(slides),
        'pages_with_notes': sum(1 for s in slides if s['has_notes']),
        'pages_need_script': sum(1 for s in slides if not s['script']),
    }, ensure_ascii=False, indent=2))


if __name__ == '__main__':
    main()
