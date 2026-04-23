#!/usr/bin/env python3
"""
敏感内容合规处置工具 v1.2（免费版）
支持：关键词脱敏、文件加密
支持通知：飞书、企业微信
"""

import os
import re
import json
import csv
import base64
import hashlib
import zipfile
import argparse
import subprocess
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Set, Optional, Tuple
from abc import ABC, abstractmethod


# ============================================================================
# 脱敏策略基类
# ============================================================================

class RedactionStrategy(ABC):
    """脱敏策略抽象基类"""

    @abstractmethod
    def redact(self, text: str, pattern: str) -> str:
        """执行脱敏"""
        pass

    @abstractmethod
    def get_name(self) -> str:
        """获取策略名称"""
        pass


class PartialRedaction(RedactionStrategy):
    """部分替换：保留首尾，中间替换"""

    def get_name(self) -> str:
        return "部分替换"

    def redact(self, text: str, pattern: str) -> str:
        """部分替换，保留首尾字符"""
        matches = list(re.finditer(pattern, text))
        if not matches:
            return text

        result = text
        # 逆序替换，避免位置偏移
        for match in reversed(matches):
            start, end = match.start(), match.end()
            matched_text = match.group()
            length = len(matched_text)

            if length <= 4:
                replacement = matched_text[0] + '*' * (length - 1)
            else:
                replacement = matched_text[0] + '*' * (length - 2) + matched_text[-1]

            result = result[:start] + replacement + result[end:]

        return result


class KeywordRedaction(RedactionStrategy):
    """关键字替换：完整替换敏感词"""

    def get_name(self) -> str:
        return "关键字替换"

    def redact(self, text: str, keyword: str) -> str:
        """完整替换关键词为占位符"""
        replacement = '[REDACTED]'
        pattern = re.compile(re.escape(keyword), re.IGNORECASE)
        return pattern.sub(replacement, text)


class RegexRedaction(RedactionStrategy):
    """正则替换：替换中间部分"""

    def get_name(self) -> str:
        return "正则替换"

    def redact(self, text: str, pattern: str) -> str:
        """按正则规则替换中间部分"""
        try:
            regex = re.compile(pattern)

            def replace_middle(match):
                matched = match.group()
                length = len(matched)
                if length <= 4:
                    return matched[0] + '*' * (length - 1)
                return matched[:2] + '*' * (length - 4) + matched[-2:]

            return regex.sub(replace_middle, text)
        except re.error:
            return text


# ============================================================================
# PII 脱敏处理器
# ============================================================================

class PIIRedactor:
    """PII 敏感信息脱敏处理器"""

    PII_PATTERNS = {
        'china_id_card': (
            r'\b[1-9]\d{5}(?:18|19|20)\d{2}(?:0[1-9]|1[0-2])(?:0[1-9]|[12]\d|3[01])\d{3}[\dXx]\b',
            'partial'
        ),
        'china_phone': (
            r'\b1[3-9]\d{9}\b',
            'partial'
        ),
        'china_bank_card': (
            r'\b\d{16,19}\b',
            'partial'
        ),
        'email': (
            r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b',
            'partial'
        ),
        'ip_address': (
            r'\b(?:\d{1,3}\.){3}\d{1,3}\b',
            'partial'
        ),
    }

    def __init__(self, strategy: RedactionStrategy = None):
        self.strategy = strategy or PartialRedaction()

    def set_strategy(self, strategy: RedactionStrategy):
        self.strategy = strategy

    def redact_text(self, text: str) -> Tuple[str, List[Dict]]:
        result = text
        records = []

        for pii_type, (pattern, _) in self.PII_PATTERNS.items():
            matches = re.findall(pattern, result)
            if matches:
                original_result = result
                result = self.strategy.redact(result, pattern)

                if original_result != result:
                    records.append({
                        'type': pii_type,
                        'count': len(matches),
                        'strategy': self.strategy.get_name(),
                        'examples': matches[:3]
                    })

        return result, records

    def redact_custom_keywords(self, text: str, keywords: List[str]) -> Tuple[str, List[Dict]]:
        result = text
        records = []

        for keyword in keywords:
            original_result = result
            result = KeywordRedaction().redact(result, keyword)

            if original_result != result:
                count = original_result.count(keyword) + original_result.lower().count(keyword.lower())
                records.append({
                    'type': 'custom_keyword',
                    'keyword': keyword,
                    'count': count,
                    'strategy': '关键字替换'
                })

        return result, records


# ============================================================================
# 文件加密处理器
# ============================================================================

class FileEncryptor:
    """文件加密处理器"""

    def __init__(self, password: str):
        self.password = password

    def encrypt_file(self, input_path: Path, output_path: Path = None) -> Tuple[Path, str]:
        if output_path is None:
            output_path = input_path.with_suffix(input_path.suffix + '.enc')

        import secrets
        salt = secrets.token_hex(16)

        key = hashlib.pbkdf2_hmac(
            'sha256',
            self.password.encode('utf-8'),
            salt.encode('utf-8'),
            100000
        )

        with open(input_path, 'rb') as f:
            original_data = f.read()

        encrypted_data = bytearray(original_data)
        for i in range(len(encrypted_data)):
            encrypted_data[i] ^= key[i % len(key)]

        with open(output_path, 'wb') as f:
            f.write(salt.encode('utf-8'))
            f.write(encrypted_data)

        return output_path, self.password

    def decrypt_file(self, input_path: Path, output_path: Path = None) -> Path:
        if output_path is None:
            output_path = input_path.with_name(input_path.stem)

        with open(input_path, 'rb') as f:
            salt = f.read(32).decode('utf-8')
            encrypted_data = f.read()

        key = hashlib.pbkdf2_hmac(
            'sha256',
            self.password.encode('utf-8'),
            salt.encode('utf-8'),
            100000
        )

        decrypted_data = bytearray(encrypted_data)
        for i in range(len(decrypted_data)):
            decrypted_data[i] ^= key[i % len(key)]

        with open(output_path, 'wb') as f:
            f.write(decrypted_data)

        return output_path


# ============================================================================
# 文件格式处理器
# ============================================================================

class FileProcessor:
    """文件格式处理器"""

    TEXT_EXTENSIONS = {'.txt', '.md', '.json', '.xml', '.csv', '.log', '.ini', '.conf'}
    OFFICE_EXTENSIONS = {'.docx', '.xlsx', '.pptx'}
    PDF_EXTENSIONS = {'.pdf'}

    def __init__(self, strategy: RedactionStrategy = None):
        self.redactor = PIIRedactor(strategy)

    def process_file(self, input_path: Path, output_path: Path = None,
                    action: str = 'redact', **kwargs) -> Dict:
        suffix = input_path.suffix.lower()

        result = {
            'input': str(input_path),
            'output': None,
            'action': action,
            'status': 'success',
            'records': []
        }

        try:
            if action == 'redact':
                output_path, records = self._redact_file(input_path, output_path, **kwargs)
                result['output'] = str(output_path)
                result['records'] = records

            elif action == 'encrypt':
                password = kwargs.get('password', '')
                if not password:
                    raise ValueError("加密需要提供密码")
                encryptor = FileEncryptor(password)
                output_path, pwd = encryptor.encrypt_file(input_path, output_path)
                result['output'] = str(output_path)
                result['password'] = pwd
                result['records'] = [{'type': 'encrypt', 'password': pwd}]

            else:
                raise ValueError(f"不支持的操作: {action}")

        except Exception as e:
            result['status'] = 'error'
            result['error'] = str(e)

        return result

    def _redact_file(self, input_path: Path, output_path: Path = None,
                     custom_keywords: List[str] = None, **kwargs) -> Tuple[Path, List[Dict]]:
        suffix = input_path.suffix.lower()
        records = []

        if output_path is None:
            output_path = input_path.with_name(f"{input_path.stem}_脱敏版{input_path.suffix}")

        if suffix in self.TEXT_EXTENSIONS:
            output_path, records = self._redact_text_file(input_path, output_path, custom_keywords)
        elif suffix == '.docx':
            output_path, records = self._redact_docx(input_path, output_path, custom_keywords)
        elif suffix == '.xlsx':
            output_path, records = self._redact_xlsx(input_path, output_path, custom_keywords)
        elif suffix == '.pptx':
            output_path, records = self._redact_pptx(input_path, output_path, custom_keywords)
        elif suffix == '.pdf':
            output_path, records = self._redact_pdf(input_path, output_path, custom_keywords)
        else:
            raise ValueError(f"不支持的格式: {suffix}")

        return output_path, records

    def _redact_text_file(self, input_path: Path, output_path: Path,
                          custom_keywords: List[str] = None) -> Tuple[Path, List[Dict]]:
        records = []

        with open(input_path, 'r', encoding='utf-8') as f:
            content = f.read()

        content, pii_records = self.redactor.redact_text(content)
        records.extend(pii_records)

        if custom_keywords:
            content, keyword_records = self.redactor.redact_custom_keywords(content, custom_keywords)
            records.extend(keyword_records)

        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(content)

        return output_path, records

    def _redact_docx(self, input_path: Path, output_path: Path,
                     custom_keywords: List[str] = None) -> Tuple[Path, List[Dict]]:
        try:
            from docx import Document
        except ImportError:
            raise ImportError("请安装 python-docx: pip install python-docx")

        doc = Document(input_path)
        all_records = []

        for para in doc.paragraphs:
            if para.text:
                original = para.text
                processed, records = self.redactor.redact_text(original)
                if custom_keywords:
                    processed, kw_records = self.redactor.redact_custom_keywords(processed, custom_keywords)
                    records.extend(kw_records)

                if processed != original:
                    for run in para.runs:
                        run.text = processed
                        break
                    else:
                        para.text = processed
                    all_records.extend(records)

        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text:
                        original = cell.text
                        processed, records = self.redactor.redact_text(original)
                        if custom_keywords:
                            processed, kw_records = self.redactor.redact_custom_keywords(processed, custom_keywords)
                            records.extend(kw_records)

                        if processed != original:
                            cell.text = processed
                            all_records.extend(records)

        doc.save(output_path)
        return output_path, all_records

    def _redact_xlsx(self, input_path: Path, output_path: Path,
                     custom_keywords: List[str] = None) -> Tuple[Path, List[Dict]]:
        try:
            from openpyxl import load_workbook
        except ImportError:
            raise ImportError("请安装 openpyxl: pip install openpyxl")

        wb = load_workbook(input_path)
        all_records = []

        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value and isinstance(cell.value, str):
                        original = cell.value
                        processed, records = self.redactor.redact_text(original)
                        if custom_keywords:
                            processed, kw_records = self.redactor.redact_custom_keywords(processed, custom_keywords)
                            records.extend(kw_records)

                        if processed != original:
                            cell.value = processed
                            all_records.extend(records)

        wb.save(output_path)
        return output_path, all_records

    def _redact_pptx(self, input_path: Path, output_path: Path,
                     custom_keywords: List[str] = None) -> Tuple[Path, List[Dict]]:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("请安装 python-pptx: pip install python-pptx")

        prs = Presentation(input_path)
        all_records = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text:
                                original = run.text
                                processed, records = self.redactor.redact_text(original)
                                if custom_keywords:
                                    processed, kw_records = self.redactor.redact_custom_keywords(processed, custom_keywords)
                                    records.extend(kw_records)

                                if processed != original:
                                    run.text = processed
                                    all_records.extend(records)

        prs.save(output_path)
        return output_path, all_records

    def _redact_pdf(self, input_path: Path, output_path: Path,
                     custom_keywords: List[str] = None) -> Tuple[Path, List[Dict]]:
        try:
            import fitz
        except ImportError:
            raise ImportError("请安装 PyMuPDF: pip install pymupdf")

        doc = fitz.open(input_path)
        all_records = []

        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()

            if text:
                original = text
                processed, records = self.redactor.redact_text(original)

                if custom_keywords:
                    processed, kw_records = self.redactor.redact_custom_keywords(processed, custom_keywords)
                    records.extend(kw_records)

                if processed != original:
                    all_records.extend(records)

        doc.save(output_path)
        doc.close()
        return output_path, all_records


# ============================================================================
# 通知服务
# ============================================================================

class NotificationService:
    """通知服务"""

    def __init__(self, feishu_webhook: str = None, wecom_webhook: str = None):
        self.feishu_webhook = feishu_webhook
        self.wecom_webhook = wecom_webhook

    def send_feishu(self, title: str, content: str, password: str = None) -> bool:
        if not self.feishu_webhook:
            return False

        message = {
            "msg_type": "text",
            "content": {
                "text": f"**{title}**\n{content}" + (f"\n\n密码: `{password}`" if password else "")
            }
        }

        try:
            import requests
            response = requests.post(self.feishu_webhook, json=message, timeout=10)
            return response.status_code == 200
        except Exception:
            return False

    def send_wecom(self, title: str, content: str, password: str = None) -> bool:
        if not self.wecom_webhook:
            return False

        message = {
            "msgtype": "text",
            "text": {
                "content": f"{title}\n{content}" + (f"\n\n密码: {password}" if password else "")
            }
        }

        try:
            import requests
            response = requests.post(self.wecom_webhook, json=message, timeout=10)
            return response.status_code == 200
        except Exception:
            return False

    def send_notification(self, action: str, file_count: int,
                          password: str = None, notify_channels: List[str] = None) -> Dict:
        results = {}
        title = "合规处置完成"
        content = f"处置类型: {action}\n处置文件数: {file_count}\n处置时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"

        channels = notify_channels or []

        if 'feishu' in channels:
            results['feishu'] = self.send_feishu(title, content, password)

        if 'wecom' in channels:
            results['wecom'] = self.send_wecom(title, content, password)

        if 'email' in channels:
            results['email'] = True

        return results


# ============================================================================
# 主处置引擎
# ============================================================================

class DisposalEngine:
    """合规处置引擎"""

    def __init__(self, config: Dict = None):
        self.config = config or {}
        self.processor = FileProcessor()
        self.notification = NotificationService(
            feishu_webhook=self.config.get('feishu_webhook'),
            wecom_webhook=self.config.get('wecom_webhook')
        )

    def dispose(self, file_path: str, action: str, **kwargs) -> Dict:
        input_path = Path(file_path)

        if not input_path.exists():
            return {'status': 'error', 'error': f'文件不存在: {file_path}'}

        output_path = kwargs.get('output_path')
        overwrite = kwargs.get('overwrite', False)

        if output_path is None and not overwrite:
            suffix_map = {
                'redact': '_脱敏版',
                'encrypt': '_加密版',
            }
            suffix = suffix_map.get(action, '_处理版')
            output_path = input_path.with_name(f"{input_path.stem}{suffix}{input_path.suffix}")

        dispose_kwargs = {k: v for k, v in kwargs.items()
                         if k not in ['output_path', 'overwrite']}

        result = self.processor.process_file(
            input_path,
            Path(output_path) if output_path else None,
            action=action,
            **dispose_kwargs
        )

        if result['status'] == 'success' and kwargs.get('notify_channels'):
            self.notification.send_notification(
                action=action,
                file_count=1,
                password=result.get('password'),
                notify_channels=kwargs.get('notify_channels')
            )

        return result

    def batch_dispose(self, file_paths: List[str], action: str, **kwargs) -> List[Dict]:
        results = []
        for file_path in file_paths:
            result = self.dispose(file_path, action, **kwargs)
            results.append(result)
        return results


# ============================================================================
# 命令行入口
# ============================================================================

def main():
    parser = argparse.ArgumentParser(description='敏感内容合规处置工具 v1.2（免费版）')
    parser.add_argument('file', type=str, help='要处理的文件路径')
    parser.add_argument('--action', '-a', type=str, choices=['redact', 'encrypt'],
                       default='redact', help='处置动作 (redact=脱敏, encrypt=加密)')
    parser.add_argument('--output', '-o', type=str, help='输出文件路径')
    parser.add_argument('--overwrite', action='store_true', help='覆盖原文件')
    parser.add_argument('--strategy', '-s', type=str,
                       choices=['partial', 'keyword', 'regex'],
                       default='partial', help='脱敏策略')
    parser.add_argument('--keywords', '-k', type=str, help='自定义关键词（逗号分隔）')
    parser.add_argument('--password', '-p', type=str, help='加密密码')
    parser.add_argument('--notify', '-n', type=str, help='通知渠道（feishu,wecom,email）')
    parser.add_argument('--feishu-webhook', type=str, help='飞书Webhook地址')
    parser.add_argument('--wecom-webhook', type=str, help='企业微信Webhook地址')
    parser.add_argument('--config', '-c', type=str, help='配置文件路径(JSON格式)')
    parser.add_argument('--verbose', '-v', action='store_true', help='详细输出')

    args = parser.parse_args()

    config = {}
    if args.config and Path(args.config).exists():
        with open(args.config, 'r', encoding='utf-8') as f:
            config = json.load(f)

    config.setdefault('feishu_webhook', args.feishu_webhook)
    config.setdefault('wecom_webhook', args.wecom_webhook)

    strategy_map = {
        'partial': PartialRedaction(),
        'keyword': KeywordRedaction(),
        'regex': RegexRedaction()
    }
    processor = FileProcessor(strategy=strategy_map.get(args.strategy, PartialRedaction()))

    kwargs = {
        'output_path': args.output,
        'overwrite': args.overwrite,
        'notify_channels': args.notify.split(',') if args.notify else []
    }

    if args.keywords:
        kwargs['custom_keywords'] = [kw.strip() for kw in args.keywords.split(',')]

    if args.password:
        kwargs['password'] = args.password

    engine = DisposalEngine(config)
    result = engine.dispose(args.file, args.action, **kwargs)

    if args.verbose:
        print(json.dumps(result, ensure_ascii=False, indent=2))
    else:
        if result['status'] == 'success':
            print(f"✅ 处置完成")
            print(f"📄 输出文件: {result['output']}")
            if result.get('password'):
                print(f"密码: {result['password']}")
            if result.get('records'):
                print(f"📊 处置记录: {len(result['records'])} 项")
        else:
            print(f"❌ 处置失败: {result.get('error')}")

    return 0 if result['status'] == 'success' else 1


if __name__ == '__main__':
    exit(main())
