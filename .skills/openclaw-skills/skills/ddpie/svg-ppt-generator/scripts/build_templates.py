#!/usr/bin/env python3
"""Generate PPTX template files with custom Slide Layouts.

Each template has 8 layouts: cover, section, content, table_data,
two_column, stats, cards, closing.  Visual elements (card backgrounds,
accent bars, decorative circles) are baked into the layout so every
slide created from them already looks polished.
"""

import os
from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn

# ═══════════════════════════════════════════════════════════════════════════
#  Slide dimensions (16:9)
# ═══════════════════════════════════════════════════════════════════════════
W  = 12192000
H  = 6858000
M  = 450000        # side / top margin
MB = 380000        # bottom margin
CW = W - 2 * M    # usable content width

# Header zone (tag + title)
TAG_Y  = M
TAG_H  = 280000
TTL_Y  = TAG_Y + TAG_H + 60000
TTL_H  = 480000
CTN_Y  = TTL_Y + TTL_H + 130000        # content area top
CTN_H  = H - CTN_Y - MB                # content area height

# ═══════════════════════════════════════════════════════════════════════════
#  Theme colour palettes
# ═══════════════════════════════════════════════════════════════════════════
DARK = dict(
    bg='1A1D21', card='2A3A3A', border='4DB6AC',
    accent='4DB6AC', accent2='80CBC4',
    t1='FFFFFF', t2='B0BEC5', tag='4DB6AC',
    deco='4DB6AC', deco_a=12000,
    lw=19050, shadow=False,       # 1.5 pt border, no shadow
)

LIGHT = dict(
    bg='F5F5F5', card='FFFFFF', border='1976D2',
    accent='1976D2', accent2='42A5F5',
    t1='212121', t2='757575', tag='1976D2',
    deco='1976D2', deco_a=8000,
    lw=12700, shadow=True,        # 1 pt border + shadow
)

# ═══════════════════════════════════════════════════════════════════════════
#  XML shape factories
# ═══════════════════════════════════════════════════════════════════════════
_NS = (
    'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
    'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
    'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"'
)

_sid = 100

def _nid():
    global _sid; _sid += 1; return _sid


def _rect(nm, x, y, w, h, fill, lc=None, lw=None, r=None, a=None, sh=False):
    """Rectangle (or roundRect when *r* given)."""
    sid = _nid()
    prst = 'roundRect' if r else 'rect'
    av = f'<a:gd name="adj" fmla="val {r}"/>' if r else ''
    fc = f'<a:srgbClr val="{fill}">'
    if a is not None:
        fc += f'<a:alpha val="{a}"/>'
    fc += '</a:srgbClr>'
    ln = (f'<a:ln w="{lw}"><a:solidFill><a:srgbClr val="{lc}"/></a:solidFill></a:ln>'
          if lc else '<a:ln><a:noFill/></a:ln>')
    ef = ''
    if sh:
        ef = ('<a:effectLst><a:outerShdw blurRad="50800" dist="25400" '
              'dir="5400000" algn="tl" rotWithShape="0"><a:srgbClr val="000000">'
              '<a:alpha val="20000"/></a:srgbClr></a:outerShdw></a:effectLst>')
    return etree.fromstring(
        f'<p:sp {_NS}>'
        f'<p:nvSpPr><p:cNvPr id="{sid}" name="{nm}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr>'
        f'<a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{w}" cy="{h}"/></a:xfrm>'
        f'<a:prstGeom prst="{prst}"><a:avLst>{av}</a:avLst></a:prstGeom>'
        f'<a:solidFill>{fc}</a:solidFill>{ln}{ef}'
        f'</p:spPr></p:sp>')


def _circ(nm, x, y, d, fill, a=None):
    """Circle / ellipse."""
    sid = _nid()
    fc = f'<a:srgbClr val="{fill}">'
    if a is not None:
        fc += f'<a:alpha val="{a}"/>'
    fc += '</a:srgbClr>'
    return etree.fromstring(
        f'<p:sp {_NS}>'
        f'<p:nvSpPr><p:cNvPr id="{sid}" name="{nm}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr>'
        f'<a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{d}" cy="{d}"/></a:xfrm>'
        f'<a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom>'
        f'<a:solidFill>{fc}</a:solidFill><a:ln><a:noFill/></a:ln>'
        f'</p:spPr></p:sp>')


def _ph(nm, tp, idx, x, y, w, h, sz=None, clr=None, b=False, al='l', an='t'):
    """Placeholder shape (title / body / custom)."""
    sid = _nid()
    pa = (f' type="{tp}"' if tp else '') + (f' idx="{idx}"' if idx is not None else '')
    am = {'l': 'l', 'c': 'ctr', 'r': 'r'}.get(al, 'l')
    anm = {'t': 't', 'ctr': 'ctr', 'b': 'b'}.get(an, 't')
    rp = 'lang="zh-CN" altLang="en-US" dirty="0"'
    if sz:
        rp += f' sz="{sz}"'
    if b:
        rp += ' b="1"'
    fi = f'<a:solidFill><a:srgbClr val="{clr}"/></a:solidFill>' if clr else ''
    return etree.fromstring(
        f'<p:sp {_NS}>'
        f'<p:nvSpPr>'
        f'<p:cNvPr id="{sid}" name="{nm}"/>'
        f'<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
        f'<p:nvPr><p:ph{pa}/></p:nvPr>'
        f'</p:nvSpPr>'
        f'<p:spPr>'
        f'<a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{w}" cy="{h}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f'<a:noFill/><a:ln><a:noFill/></a:ln>'
        f'</p:spPr>'
        f'<p:txBody>'
        f'<a:bodyPr wrap="square" lIns="91440" tIns="45720" rIns="91440" bIns="45720" anchor="{anm}"/>'
        f'<a:lstStyle/>'
        f'<a:p><a:pPr algn="{am}"><a:buNone/></a:pPr><a:endParaRPr {rp}>{fi}</a:endParaRPr></a:p>'
        f'</p:txBody></p:sp>')


# ═══════════════════════════════════════════════════════════════════════════
#  Layout helpers
# ═══════════════════════════════════════════════════════════════════════════

def _clear(lay):
    """Strip every shape from a layout's shape tree."""
    tree = lay._element.find(qn('p:cSld')).find(qn('p:spTree'))
    for ch in list(tree):
        tag = ch.tag.split('}')[-1]
        if tag not in ('nvGrpSpPr', 'grpSpPr'):
            tree.remove(ch)


def _init(lay, name, bg):
    """Prepare a layout: clear, rename, set background, hide master."""
    _clear(lay)
    lay._element.set('showMasterSp', '0')
    lay._element.set('userDrawn', '1')
    cSld = lay._element.find(qn('p:cSld'))
    cSld.set('name', name)
    old = cSld.find(qn('p:bg'))
    if old is not None:
        cSld.remove(old)
    bg_el = etree.fromstring(
        f'<p:bg {_NS}><p:bgPr>'
        f'<a:solidFill><a:srgbClr val="{bg}"/></a:solidFill>'
        f'<a:effectLst/></p:bgPr></p:bg>')
    cSld.insert(0, bg_el)


def _add(lay, *elems):
    """Append shape elements to a layout's shape tree."""
    tree = lay._element.find(qn('p:cSld')).find(qn('p:spTree'))
    for e in elems:
        tree.append(e)


# ═══════════════════════════════════════════════════════════════════════════
#  Layout builders (one per slide type)
# ═══════════════════════════════════════════════════════════════════════════

def _cover(lay, t):
    """Cover — accent bar + big title + subtitle + date."""
    _init(lay, 'cover', t['bg'])
    _add(lay,
        # decorative
        _rect('bar_l', 0, 0, 100000, H, t['accent']),
        _rect('bar_b', 0, H - 70000, W, 70000, t['accent']),
        _circ('deco_tr', W - 900000, -500000, 1400000, t['deco'], t['deco_a']),
        _circ('deco_bl', -300000, H - 600000,  800000, t['deco'], t['deco_a']),
        # placeholders
        _ph('Title',    'ctrTitle', 0,  M + 300000, 2000000, CW - 600000, 1200000,
            sz=4400, clr=t['t1'], b=True, al='c', an='ctr'),
        _ph('Subtitle', 'subTitle', 1,  M + 300000, 3400000, CW - 600000,  600000,
            sz=2000, clr=t['t2'], al='c'),
        _ph('Date',     'body',     10, M + 300000, 4200000, CW - 600000,  400000,
            sz=1400, clr=t['t2'], al='c'),
    )


def _section(lay, t):
    """Section divider — large number + accent line + title."""
    _init(lay, 'section', t['bg'])
    _add(lay,
        _rect('bar_r', W - 100000, 0, 100000, H, t['accent']),
        _rect('line',  M, 3100000, 2500000, 50000, t['accent']),
        _circ('deco_tr', W - 700000, -400000, 1200000, t['deco'], t['deco_a']),
        _circ('deco_bl', -300000, H - 500000,  700000, t['deco'], t['deco_a']),
        _ph('Number',   'body',     10, M, 1500000, 3000000, 1400000,
            sz=8000, clr=t['accent'], b=True),
        _ph('Title',    'title',     0, M, 3300000, CW - 200000, 800000,
            sz=3600, clr=t['t1'], b=True),
        _ph('Subtitle', 'subTitle',  1, M, 4200000, CW - 200000, 500000,
            sz=1800, clr=t['t2']),
    )


def _content(lay, t):
    """Content — tag + title + card panel with body placeholder."""
    _init(lay, 'content', t['bg'])
    _add(lay,
        _circ('deco_br', W - 400000, H - 400000, 600000, t['deco'], t['deco_a']),
        _rect('card', M, CTN_Y, CW, CTN_H,
              t['card'], t['border'], t['lw'], r=3500, sh=t['shadow']),
        _rect('bar_l', M, CTN_Y, 55000, CTN_H, t['accent']),
        _ph('Tag',   'body',  10, M,           TAG_Y, 3000000, TAG_H,
            sz=1200, clr=t['tag'], b=True),
        _ph('Title', 'title',  0, M,           TTL_Y, CW,      TTL_H,
            sz=2800, clr=t['t1'], b=True),
        _ph('Body',  'body',   1, M + 210000,  CTN_Y + 130000,
            CW - 360000, CTN_H - 260000, sz=1600, clr=t['t1']),
    )


def _table_data(lay, t):
    """Table — tag + title + card area (table added at runtime)."""
    _init(lay, 'table_data', t['bg'])
    _add(lay,
        _circ('deco_br', W - 400000, H - 400000, 600000, t['deco'], t['deco_a']),
        _rect('card', M, CTN_Y, CW, CTN_H,
              t['card'], t['border'], t['lw'], r=3500, sh=t['shadow']),
        _rect('bar_t', M, CTN_Y, CW, 50000, t['accent']),
        _ph('Tag',   'body',  10, M, TAG_Y, 3000000, TAG_H,
            sz=1200, clr=t['tag'], b=True),
        _ph('Title', 'title',  0, M, TTL_Y, CW,      TTL_H,
            sz=2800, clr=t['t1'], b=True),
    )


def _two_column(lay, t):
    """Two-column comparison — left / right panels."""
    _init(lay, 'two_column', t['bg'])
    gap = 200000
    cw  = (CW - gap) // 2
    rx  = M + cw + gap
    _add(lay,
        # left panel
        _rect('panel_l', M,  CTN_Y, cw, CTN_H,
              t['card'], t['border'], t['lw'], r=3500, sh=t['shadow']),
        _rect('bar_lt',  M,  CTN_Y, cw, 50000, t['accent']),
        # right panel
        _rect('panel_r', rx, CTN_Y, cw, CTN_H,
              t['card'], t['border'], t['lw'], r=3500, sh=t['shadow']),
        _rect('bar_rt',  rx, CTN_Y, cw, 50000, t['accent2']),
        # placeholders
        _ph('Tag',   'body',  10, M,  TAG_Y, 3000000, TAG_H,
            sz=1200, clr=t['tag'], b=True),
        _ph('Title', 'title',  0, M,  TTL_Y, CW,      TTL_H,
            sz=2800, clr=t['t1'], b=True),
        _ph('Left',  'body',  11, M  + 100000, CTN_Y + 100000, cw - 200000, 400000,
            sz=2000, clr=t['t1'], b=True),
        _ph('Right', 'body',  12, rx + 100000, CTN_Y + 100000, cw - 200000, 400000,
            sz=2000, clr=t['t1'], b=True),
    )


def _stats(lay, t):
    """Stats / KPI — 4 metric cards."""
    _init(lay, 'stats', t['bg'])
    gap = 180000
    sw  = (CW - 3 * gap) // 4
    shapes = []
    for i in range(4):
        x = M + i * (sw + gap)
        shapes.append(_rect(f'stat_{i}', x, CTN_Y, sw, CTN_H,
                            t['card'], t['border'], t['lw'], r=5000, sh=t['shadow']))
        shapes.append(_rect(f'bar_{i}',  x, CTN_Y, sw, 50000, t['accent']))
    shapes += [
        _ph('Tag',   'body',  10, M, TAG_Y, 3000000, TAG_H,
            sz=1200, clr=t['tag'], b=True),
        _ph('Title', 'title',  0, M, TTL_Y, CW,      TTL_H,
            sz=2800, clr=t['t1'], b=True),
    ]
    _add(lay, *shapes)


def _cards(lay, t):
    """Cards — 3 card panels."""
    _init(lay, 'cards', t['bg'])
    gap = 200000
    cwd = (CW - 2 * gap) // 3
    shapes = []
    for i in range(3):
        x = M + i * (cwd + gap)
        shapes.append(_rect(f'card_{i}', x, CTN_Y, cwd, CTN_H,
                            t['card'], t['border'], t['lw'], r=5000, sh=t['shadow']))
        shapes.append(_rect(f'bar_{i}',  x, CTN_Y, cwd, 50000, t['accent']))
    shapes += [
        _ph('Tag',   'body',  10, M, TAG_Y, 3000000, TAG_H,
            sz=1200, clr=t['tag'], b=True),
        _ph('Title', 'title',  0, M, TTL_Y, CW,      TTL_H,
            sz=2800, clr=t['t1'], b=True),
    ]
    _add(lay, *shapes)


def _closing(lay, t):
    """Closing — centred title + subtitle."""
    _init(lay, 'closing', t['bg'])
    _add(lay,
        _rect('bar_b', 0, H - 70000, W, 70000, t['accent']),
        _rect('line',  W // 2 - 1200000, 2400000, 2400000, 50000, t['accent']),
        _circ('deco_tl', -400000, -400000, 1000000, t['deco'], t['deco_a']),
        _circ('deco_br', W - 500000, H - 500000, 800000, t['deco'], t['deco_a']),
        _ph('Title',    'ctrTitle', 0, M + 300000, 2700000, CW - 600000, 1200000,
            sz=4400, clr=t['t1'], b=True, al='c', an='ctr'),
        _ph('Subtitle', 'subTitle', 1, M + 300000, 4000000, CW - 600000,  600000,
            sz=2000, clr=t['t2'], al='c'),
    )


# ═══════════════════════════════════════════════════════════════════════════
#  Build entry-point
# ═══════════════════════════════════════════════════════════════════════════

_BUILDERS = [_cover, _section, _content, _table_data,
             _two_column, _stats, _cards, _closing]

def build(theme, path):
    """Create one template file."""
    global _sid
    _sid = 100
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    for i, fn in enumerate(_BUILDERS):
        fn(prs.slide_masters[0].slide_layouts[i], theme)
    os.makedirs(os.path.dirname(path), exist_ok=True)
    prs.save(path)
    print(f'  saved  {path}')


if __name__ == '__main__':
    d = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'templates')
    build(DARK,  os.path.join(d, 'dark_tech.pptx'))
    build(LIGHT, os.path.join(d, 'light_corporate.pptx'))
    print('Done — 2 templates generated.')
