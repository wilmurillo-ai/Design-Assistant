"""
Template-driven PPT engine.
Uses pre-defined visual parameters (matching dark_tech.pptx template) to generate
professional slides with visible cards, borders, and accent elements.

Usage:
    from template_engine import generate_ppt
    slides = [{"type": "cover", "title": "Hello", "subtitle": "World"}, ...]
    generate_ppt(slides, "/tmp/output.pptx", theme="dark_tech")
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================
# Constants
# ============================================================
SLIDE_W = 12192000
SLIDE_H = 6858000
MARGIN = 381000
CONTENT_W = SLIDE_W - 2 * MARGIN

# ============================================================
# Themes
# ============================================================
THEMES = {
    "dark_tech": {
        "bg": "1A1D21",
        "card_bg": "2A3A3A",
        "accent": "4DB6AC",
        "accent2": "FF7043",
        "text_white": "FFFFFF",
        "text_primary": "E0E0E0",
        "text_secondary": "9E9E9E",
        "border_w": 1.5,
        "card_right_bg": "3A2D2A",
        "deco1": "1F2A29",  # pre-blended accent on bg
        "deco2": "211F20",  # pre-blended accent2 on bg
        "deco3": "223230",  # pre-blended accent on bg (brighter)
    },
    "dark_blue": {
        "bg": "0D1B2A",
        "card_bg": "1B2D45",
        "accent": "FFC107",
        "accent2": "64B5F6",
        "text_white": "FFFFFF",
        "text_primary": "E0E0E0",
        "text_secondary": "90A4AE",
        "border_w": 1.5,
        "card_right_bg": "1A2A40",
        "deco1": "152535",
        "deco2": "1A2030",
        "deco3": "1A2838",
    },
    "dark_purple": {
        "bg": "1A1025",
        "card_bg": "2D1B4E",
        "accent": "BB86FC",
        "accent2": "03DAC5",
        "text_white": "FFFFFF",
        "text_primary": "E0E0E0",
        "text_secondary": "9E9E9E",
        "border_w": 1.5,
        "card_right_bg": "1A2D2D",
        "deco1": "221535",
        "deco2": "152225",
        "deco3": "251A38",
    },
    "dark_red": {
        "bg": "1A1215",
        "card_bg": "2D1A1E",
        "accent": "FF5252",
        "accent2": "FFD740",
        "text_white": "FFFFFF",
        "text_primary": "E0E0E0",
        "text_secondary": "9E9E9E",
        "border_w": 1.5,
        "card_right_bg": "2D2518",
        "deco1": "251518",
        "deco2": "252018",
        "deco3": "301820",
    },
    "dark_green": {
        "bg": "0D1A12",
        "card_bg": "1A2D1E",
        "accent": "66BB6A",
        "accent2": "FFA726",
        "text_white": "FFFFFF",
        "text_primary": "E0E0E0",
        "text_secondary": "9E9E9E",
        "border_w": 1.5,
        "card_right_bg": "2D2518",
        "deco1": "152518",
        "deco2": "1A2015",
        "deco3": "1A2A1A",
    },
    "dark_warm": {
        "bg": "F8FAFC",
        "card_bg": "FFFFFF",
        "accent": "D97757",
        "accent2": "4A90D9",
        "text_white": "1A1A2E",
        "text_primary": "334155",
        "text_secondary": "64748B",
        "border_w": 1.0,
        "card_right_bg": "EFF6FF",
        "deco1": "F5EDE8",
        "deco2": "E8EFF5",
        "deco3": "F0E5DF",
    },
    "consultant": {
        "bg": "FFFFFF",
        "card_bg": "ECF0F1",
        "accent": "005587",
        "accent2": "004D5C",
        "text_white": "2C3E50",
        "text_primary": "34495E",
        "text_secondary": "7F8C8D",
        "border_w": 1.0,
        "card_right_bg": "E8F4F8",
        "deco1": "E8EFF5",
        "deco2": "E5F0F0",
        "deco3": "D8E8F0",
    },
    "light_clean": {
        "bg": "FFFFFF",
        "card_bg": "F8F9FA",
        "accent": "4285F4",
        "accent2": "EA4335",
        "text_white": "202124",
        "text_primary": "3C4043",
        "text_secondary": "5F6368",
        "border_w": 1.0,
        "card_right_bg": "FEF7F6",
        "deco1": "EBF2FE",
        "deco2": "FDECEB",
        "deco3": "E3EDFC",
    },
    "light_corporate": {
        "bg": "F5F5F5",
        "card_bg": "FFFFFF",
        "accent": "1976D2",
        "accent2": "E65100",
        "text_white": "212121",
        "text_primary": "424242",
        "text_secondary": "757575",
        "border_w": 1.0,
        "card_right_bg": "FFF3E0",
        "deco1": "E3EDF7",
        "deco2": "F7E8E0",
        "deco3": "D0E4F2",
    },
    "cloud_orange": {
        "bg": "232F3E",
        "card_bg": "2A3B4D",
        "accent": "FF9900",
        "accent2": "527FFF",
        "text_white": "FFFFFF",
        "text_primary": "D5DBDB",
        "text_secondary": "A0AAB2",
        "border_w": 1.0,
        "card_right_bg": "1A2332",
        "deco1": "2E3F52",
        "deco2": "1E2D3D",
        "deco3": "354A5F",
    },
}

T = THEMES["dark_tech"]  # active theme


# ============================================================
# Helpers
# ============================================================
def hex_rgb(h):
    return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:], 16))


def add_rect(slide, left, top, w, h, fill, border=None, border_w=None, rounded=False):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(st, Emu(left), Emu(top), Emu(w), Emu(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(fill)
    if border and border_w:
        shape.line.color.rgb = hex_rgb(border)
        shape.line.width = Pt(border_w)
    else:
        shape.line.fill.background()
    return shape


def add_circle(slide, cx, cy, size, fill):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Emu(cx - size // 2), Emu(cy - size // 2), Emu(size), Emu(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(fill)
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, w, h, text, size=14, color="FFFFFF",
             bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(w), Emu(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(size)
    p.font.color.rgb = hex_rgb(color)
    p.font.bold = bold
    p.alignment = align
    return txBox


def add_decorations(slide):
    add_circle(slide, SLIDE_W - 500000, 400000, 800000, T["deco1"])
    add_circle(slide, 400000, SLIDE_H - 400000, 500000, T["deco2"])
    add_circle(slide, SLIDE_W - 800000, 250000, 100000, T["deco3"])


def add_accent_bar(slide, x, y, h, color=None, w=32000):
    return add_rect(slide, x, y, w, h, color or T["accent"])


def add_tag(slide, y, tag):
    add_circle(slide, MARGIN + 25000, y + 55000, 50000, T["accent"])
    add_text(slide, MARGIN + 70000, y + 30000, 2000000, 200000,
             tag, size=11, color=T["accent"], bold=True)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, T["bg"])
    return slide


# ============================================================
# Layout Builders
# ============================================================

def build_cover(prs, title="", subtitle="", author="", date="", **kw):
    slide = new_slide(prs)
    add_decorations(slide)
    add_rect(slide, 0, 0, SLIDE_W, 40000, T["accent"])
    add_text(slide, MARGIN, SLIDE_H // 2 - 600000, CONTENT_W, 500000,
             title, size=44, color=T["text_white"], bold=True, align=PP_ALIGN.CENTER)
    sub = subtitle or ""
    if author:
        sub = f"{sub}\n{author}" if sub else author
    add_text(slide, MARGIN, SLIDE_H // 2, CONTENT_W, 400000,
             sub, size=18, color=T["text_secondary"], align=PP_ALIGN.CENTER)
    add_rect(slide, MARGIN, SLIDE_H - 300000, CONTENT_W, 20000, T["accent"])
    if date:
        add_text(slide, MARGIN, SLIDE_H - 250000, CONTENT_W, 200000,
                 date, size=12, color=T["text_secondary"], align=PP_ALIGN.CENTER)
    return slide


def build_section(prs, number=1, title="", subtitle="", **kw):
    slide = new_slide(prs)
    add_decorations(slide)
    badge_size = 900000
    bx = MARGIN + 200000
    by = SLIDE_H // 2 - badge_size // 2 - 100000
    add_circle(slide, bx + badge_size // 2, by + badge_size // 2, badge_size, T["deco3"])
    add_circle(slide, bx + badge_size // 2, by + badge_size // 2,
               int(badge_size * 0.7), T["card_bg"])
    add_text(slide, bx, by, badge_size, badge_size,
             str(number), size=48, color=T["accent"], bold=True, align=PP_ALIGN.CENTER)
    tx = bx + badge_size + 200000
    add_text(slide, tx, SLIDE_H // 2 - 300000, SLIDE_W - tx - MARGIN, 350000,
             title, size=36, color=T["text_white"], bold=True)
    add_rect(slide, tx, SLIDE_H // 2 + 80000, 800000, 20000, T["accent"])
    if subtitle:
        add_text(slide, tx, SLIDE_H // 2 + 150000, SLIDE_W - tx - MARGIN, 250000,
                 subtitle, size=16, color=T["text_secondary"])
    return slide


def build_content(prs, tag="", title="", bullets=None, subtitle="", **kw):
    bullets = bullets or []
    slide = new_slide(prs)
    add_decorations(slide)

    y = MARGIN
    add_tag(slide, y, tag)
    y += 220000
    add_text(slide, MARGIN, y, CONTENT_W, 450000,
             title, size=28, color=T["text_white"], bold=True)
    y += 480000
    if subtitle:
        add_text(slide, MARGIN, y, CONTENT_W, 220000,
                 subtitle, size=12, color=T["text_secondary"])
        y += 250000

    n = len(bullets)
    if n == 0:
        return slide

    bottom = SLIDE_H - 350000
    card_gap = 50000
    card_h = (bottom - y - (n - 1) * card_gap) // n
    card_h = max(card_h, 250000)
    bfont = 16 if n <= 3 else (14 if n <= 5 else 12)

    for idx, bullet in enumerate(bullets):
        add_rect(slide, MARGIN, y, CONTENT_W, card_h,
                 T["card_bg"], border=T["accent"], border_w=T["border_w"], rounded=True)
        add_accent_bar(slide, MARGIN + 18000, y + int(card_h * 0.15), int(card_h * 0.7))
        badge_sz = min(int(card_h * 0.55), 280000)
        add_text(slide, MARGIN + 70000, y + (card_h - badge_sz) // 2, badge_sz, badge_sz,
                 str(idx + 1), size=18, color=T["accent"], bold=True, align=PP_ALIGN.CENTER)
        tl = MARGIN + 70000 + badge_sz + 20000
        tw = CONTENT_W - (tl - MARGIN) - 50000

        if isinstance(bullet, str):
            add_text(slide, tl, y + (card_h - int(card_h * 0.6)) // 2,
                     tw, int(card_h * 0.6), bullet, size=bfont, color=T["text_primary"])
        elif isinstance(bullet, dict):
            key = bullet.get("key", "")
            val = bullet.get("val", "")
            add_text(slide, tl, y + int(card_h * 0.12), tw, int(card_h * 0.38),
                     key, size=bfont, color=T["accent"], bold=True)
            add_text(slide, tl, y + int(card_h * 0.50), tw, int(card_h * 0.42),
                     val, size=bfont - 1, color=T["text_primary"])

        y += card_h + card_gap
    return slide


def build_table(prs, tag="", title="", headers=None, rows=None, subtitle="", **kw):
    headers = headers or []
    rows = rows or []
    slide = new_slide(prs)
    add_decorations(slide)

    y = MARGIN
    add_tag(slide, y, tag)
    y += 220000
    add_text(slide, MARGIN, y, CONTENT_W, 450000,
             title, size=28, color=T["text_white"], bold=True)
    y += 480000
    if subtitle:
        add_text(slide, MARGIN, y, CONTENT_W, 200000,
                 subtitle, size=12, color=T["text_secondary"])
        y += 230000

    n_cols = len(headers)
    if n_cols == 0:
        return slide
    col_w = CONTENT_W // n_cols

    bottom = SLIDE_H - 350000
    n_rows = len(rows)
    header_h = 340000
    avail = bottom - y - header_h - 20000
    row_h = min(avail // max(n_rows, 1), 500000)
    rfont = 13 if n_rows <= 6 else (12 if n_rows <= 8 else 11)

    # Header
    add_rect(slide, MARGIN, y, CONTENT_W, header_h, T["card_bg"],
             border=T["accent"], border_w=0.8, rounded=True)
    for i, h in enumerate(headers):
        add_text(slide, MARGIN + i * col_w + 60000, y + 40000,
                 col_w - 120000, header_h - 80000,
                 h, size=13, color=T["accent"], bold=True)
    y += header_h + 20000

    # Rows
    for ri, row in enumerate(rows):
        if ri % 2 == 0:
            add_rect(slide, MARGIN, y, CONTENT_W, int(row_h * 0.9),
                     T["card_bg"], rounded=True)
        add_accent_bar(slide, MARGIN + 10000, y + int(row_h * 0.15),
                      int(row_h * 0.6), w=22000)
        for i, cell in enumerate(row):
            add_text(slide, MARGIN + i * col_w + 60000, y + 25000,
                     col_w - 120000, row_h - 50000,
                     str(cell), size=rfont, color=T["text_primary"], bold=(i == 0))
        y += row_h
    return slide


def build_two_column(prs, tag="", title="", left_items=None, right_items=None,
                     left_title="", right_title="", **kw):
    left_items = left_items or []
    right_items = right_items or []
    slide = new_slide(prs)
    add_decorations(slide)

    y = MARGIN
    add_tag(slide, y, tag)
    y += 220000
    add_text(slide, MARGIN, y, CONTENT_W, 450000,
             title, size=28, color=T["text_white"], bold=True)
    y += 500000

    col_gap = 120000
    panel_w = (CONTENT_W - col_gap) // 2
    panel_h = SLIDE_H - y - 350000

    max_items = max(len(left_items), len(right_items))
    title_h = 380000 if (left_title or right_title) else 0
    inner_h = panel_h - 80000 - title_h
    item_slot = min(inner_h // max(max_items, 1), 550000)
    ifont = 15 if max_items <= 4 else (14 if max_items <= 5 else 13)

    # Left panel
    lx = MARGIN
    add_rect(slide, lx, y, panel_w, panel_h,
             T["card_bg"], border=T["accent"], border_w=T["border_w"], rounded=True)
    add_accent_bar(slide, lx + 18000, y + 30000, panel_h - 60000)
    ly = y + 40000
    if left_title:
        add_text(slide, lx + 70000, ly, panel_w - 100000, 300000,
                 left_title, size=16, color=T["accent"], bold=True)
        add_rect(slide, lx + 70000, ly + 280000, panel_w // 2, 16000, T["accent"])
        ly += title_h
    for item in left_items:
        add_circle(slide, lx + 90000, ly + 50000, 40000, T["accent"])
        add_text(slide, lx + 140000, ly + 30000, panel_w - 180000, int(item_slot * 0.7),
                 item, size=ifont, color=T["text_primary"])
        ly += item_slot

    # Right panel
    rx = MARGIN + panel_w + col_gap
    add_rect(slide, rx, y, panel_w, panel_h,
             T["card_right_bg"], border=T["accent2"], border_w=T["border_w"], rounded=True)
    add_accent_bar(slide, rx + 18000, y + 30000, panel_h - 60000, color=T["accent2"])
    ry = y + 40000
    if right_title:
        add_text(slide, rx + 70000, ry, panel_w - 100000, 300000,
                 right_title, size=16, color=T["accent2"], bold=True)
        add_rect(slide, rx + 70000, ry + 280000, panel_w // 2, 16000, T["accent2"])
        ry += title_h
    for item in right_items:
        add_circle(slide, rx + 90000, ry + 50000, 40000, T["accent2"])
        add_text(slide, rx + 140000, ry + 30000, panel_w - 180000, int(item_slot * 0.7),
                 item, size=ifont, color=T["text_primary"])
        ry += item_slot
    return slide


def build_stats(prs, tag="", title="", stats=None, subtitle="", **kw):
    stats = stats or []
    slide = new_slide(prs)
    add_decorations(slide)

    y = MARGIN
    add_tag(slide, y, tag)
    y += 220000
    add_text(slide, MARGIN, y, CONTENT_W, 450000,
             title, size=28, color=T["text_white"], bold=True)
    y += 480000
    if subtitle:
        add_text(slide, MARGIN, y, CONTENT_W, 200000,
                 subtitle, size=12, color=T["text_secondary"])
        y += 250000

    n = len(stats)
    if n == 0:
        return slide

    cols = min(n, 4)
    col_gap = 80000
    stat_w = (CONTENT_W - (cols - 1) * col_gap) // cols
    card_h = SLIDE_H - y - 350000
    card_h = min(card_h, 2400000)

    for i, s in enumerate(stats):
        sx = MARGIN + (i % cols) * (stat_w + col_gap)
        sy = y + (i // cols) * (card_h + 100000)
        add_rect(slide, sx, sy, stat_w, card_h,
                 T["card_bg"], border=T["accent"], border_w=1.2, rounded=True)
        add_text(slide, sx, sy + int(card_h * 0.10), stat_w, int(card_h * 0.35),
                 s.get("value", ""), size=48, color=T["accent"],
                 bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, sx + 50000, sy + int(card_h * 0.52), stat_w - 100000, int(card_h * 0.16),
                 s.get("label", ""), size=16, color=T["text_white"],
                 bold=True, align=PP_ALIGN.CENTER)
        desc = s.get("desc", "")
        if desc:
            add_text(slide, sx + 50000, sy + int(card_h * 0.70), stat_w - 100000, int(card_h * 0.20),
                     desc, size=11, color=T["text_secondary"], align=PP_ALIGN.CENTER)
    return slide


def build_cards(prs, tag="", title="", cards=None, **kw):
    cards = cards or []
    slide = new_slide(prs)
    add_decorations(slide)

    y = MARGIN
    add_tag(slide, y, tag)
    y += 220000
    add_text(slide, MARGIN, y, CONTENT_W, 450000,
             title, size=28, color=T["text_white"], bold=True)
    y += 500000

    n = len(cards)
    if n == 0:
        return slide

    card_gap = 100000
    card_w = (CONTENT_W - (n - 1) * card_gap) // n
    card_h = SLIDE_H - y - 350000

    for i, card in enumerate(cards):
        cx = MARGIN + i * (card_w + card_gap)
        add_rect(slide, cx, y, card_w, card_h,
                 T["card_bg"], border=T["accent"], border_w=T["border_w"], rounded=True)
        ctitle = card.get("title", "")
        add_text(slide, cx + 50000, y + 50000, card_w - 100000, 250000,
                 ctitle, size=16, color=T["accent"], bold=True)
        add_rect(slide, cx + 50000, y + 280000, card_w // 3, 16000, T["accent"])
        items = card.get("items", [])
        n_items = len(items)
        item_area = card_h - 350000
        item_slot = min(item_area // max(n_items, 1), 300000)
        for j, item in enumerate(items):
            iy = y + 350000 + j * item_slot
            add_circle(slide, cx + 70000, iy + 40000, 35000, T["accent"])
            add_text(slide, cx + 120000, iy + 20000, card_w - 180000, int(item_slot * 0.8),
                     item, size=13, color=T["text_primary"])
    return slide


def build_closing(prs, title="", subtitle="", **kw):
    slide = new_slide(prs)
    add_decorations(slide)
    add_rect(slide, 0, SLIDE_H - 40000, SLIDE_W, 40000, T["accent"])
    add_text(slide, MARGIN, SLIDE_H // 2 - 400000, CONTENT_W, 500000,
             title, size=44, color=T["text_white"], bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, MARGIN, SLIDE_H // 2 + 200000, CONTENT_W, 300000,
                 subtitle, size=16, color=T["text_secondary"], align=PP_ALIGN.CENTER)
    return slide


# ============================================================
# Layout dispatcher
# ============================================================
BUILDERS = {
    "cover": build_cover,
    "section": build_section,
    "content": build_content,
    "table": build_table,
    "two_column": build_two_column,
    "stats": build_stats,
    "cards": build_cards,
    "closing": build_closing,
}


# ============================================================
# Main entry point
# ============================================================
def generate_ppt(slides_data, output_path, theme="dark_tech", show_page_numbers=False):
    """Generate a PPTX from slides data using the template engine.

    Args:
        slides_data: list of dicts, each with "type" and layout-specific fields
        output_path: path to save the PPTX
        theme: "dark_tech" or "light_corporate"
        show_page_numbers: whether to add page numbers

    Returns:
        output_path
    """
    global T
    T = THEMES.get(theme, THEMES["dark_tech"])

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)

    for i, sd in enumerate(slides_data):
        stype = sd.get("type", "content")
        builder = BUILDERS.get(stype)
        if builder is None:
            print(f"Warning: unknown slide type '{stype}', skipping")
            continue
        builder(prs, **sd)

        if show_page_numbers and stype not in ("cover", "closing"):
            slide = prs.slides[-1]
            add_text(slide, SLIDE_W - 500000, SLIDE_H - 250000, 400000, 200000,
                     str(i + 1), size=10, color=T["text_secondary"], align=PP_ALIGN.RIGHT)

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    return output_path
