"""
PPT Validator - Detect and auto-fix layout issues in generated PPTX files.
Checks for: text overlap, boundary overflow, content density.
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.oxml.ns import qn
import math

SLIDE_W = 12192000  # EMU
SLIDE_H = 6858000
MARGIN = 381000

# Approximate characters per line based on font size and box width
def estimate_text_height(text, font_size_pt, box_width_emu, line_spacing=1.3):
    """Estimate the height needed for text in EMU."""
    if not text:
        return 0
    # Detect CJK characters for wider char width
    cjk_count = sum(1 for c in text if '\u4e00' <= c <= '\u9fff' or '\u3000' <= c <= '\u303f')
    cjk_ratio = cjk_count / max(1, len(text))
    # Digits and punctuation are narrower (~0.4), Latin ~0.55, CJK ~1.0
    digit_count = sum(1 for c in text if c.isdigit() or c in ' .,;:/-')
    digit_ratio = digit_count / max(1, len(text))
    char_width_factor = 0.4 * digit_ratio + 1.0 * cjk_ratio + 0.55 * (1 - digit_ratio - cjk_ratio)
    char_width_factor = max(0.35, min(1.0, char_width_factor))
    char_width_emu = int(font_size_pt * 12700 * char_width_factor)
    chars_per_line = max(1, box_width_emu // max(1, char_width_emu))
    num_lines = math.ceil(len(text) / chars_per_line)
    line_height_emu = int(font_size_pt * 12700 * line_spacing)
    return num_lines * line_height_emu


def get_shape_bounds(shape):
    """Get shape bounding box as (left, top, right, bottom) in EMU."""
    left = shape.left if shape.left is not None else 0
    top = shape.top if shape.top is not None else 0
    width = shape.width if shape.width is not None else 0
    height = shape.height if shape.height is not None else 0
    return (left, top, left + width, top + height)


def shapes_overlap(bounds1, bounds2, min_overlap_emu=50000):
    """Check if two bounding boxes overlap by more than min_overlap_emu."""
    l1, t1, r1, b1 = bounds1
    l2, t2, r2, b2 = bounds2
    
    overlap_x = max(0, min(r1, r2) - max(l1, l2))
    overlap_y = max(0, min(b1, b2) - max(t1, t2))
    
    if overlap_x > min_overlap_emu and overlap_y > min_overlap_emu:
        return (overlap_x, overlap_y)
    return None


def is_bg_or_overlay(shape, slide_w=SLIDE_W, slide_h=SLIDE_H):
    """Check if a shape is a background image or overlay (covers >80% of slide)."""
    if shape.width is None or shape.height is None:
        return False
    area = shape.width * shape.height
    slide_area = slide_w * slide_h
    return area > slide_area * 0.8


def get_text_content(shape):
    """Extract text content from a shape."""
    if not shape.has_text_frame:
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs if p.text)


def get_font_size(shape):
    """Get the primary font size from a shape in points."""
    if not shape.has_text_frame:
        return 14  # default
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                return run.font.size.pt
    return 14


def validate_slide(slide, slide_idx, slide_w=SLIDE_W, slide_h=SLIDE_H):
    """Validate a single slide for layout issues."""
    issues = []
    shapes = list(slide.shapes)
    
    # Filter out background/overlay shapes
    content_shapes = []
    for s in shapes:
        if not is_bg_or_overlay(s):
            content_shapes.append(s)
    
    # 1. Check boundary overflow
    for s in content_shapes:
        bounds = get_shape_bounds(s)
        l, t, r, b = bounds
        if r > slide_w + 50000:  # 50000 EMU tolerance
            issues.append({
                "slide": slide_idx + 1,
                "severity": "error",
                "type": "overflow_right",
                "shape": s.name,
                "detail": f"右边界溢出 {(r - slide_w) // 12700}pt",
                "bounds": bounds,
            })
        if b > slide_h + 50000:
            issues.append({
                "slide": slide_idx + 1,
                "severity": "error",
                "type": "overflow_bottom",
                "shape": s.name,
                "detail": f"底部溢出 {(b - slide_h) // 12700}pt",
                "bounds": bounds,
            })
    
    # 2. Check text-to-text overlaps (skip bg/overlay shapes)
    text_shapes = [s for s in content_shapes if s.has_text_frame and get_text_content(s)]
    
    for i in range(len(text_shapes)):
        for j in range(i + 1, len(text_shapes)):
            s1, s2 = text_shapes[i], text_shapes[j]
            b1 = get_shape_bounds(s1)
            b2 = get_shape_bounds(s2)
            overlap = shapes_overlap(b1, b2)
            if overlap:
                ox, oy = overlap
                # Only report if overlap is significant relative to shape sizes
                s1_area = (b1[2] - b1[0]) * (b1[3] - b1[1])
                s2_area = (b2[2] - b2[0]) * (b2[3] - b2[1])
                overlap_area = ox * oy
                min_area = min(s1_area, s2_area) if min(s1_area, s2_area) > 0 else 1
                overlap_pct = overlap_area / min_area * 100
                
                if overlap_pct > 10:  # >10% overlap is significant
                    issues.append({
                        "slide": slide_idx + 1,
                        "severity": "error" if overlap_pct > 30 else "warning",
                        "type": "text_overlap",
                        "shapes": (s1.name, s2.name),
                        "texts": (get_text_content(s1)[:30], get_text_content(s2)[:30]),
                        "overlap_pct": round(overlap_pct, 1),
                        "detail": f"文字重叠 {overlap_pct:.0f}%: '{get_text_content(s1)[:20]}' ↔ '{get_text_content(s2)[:20]}'",
                    })
    
    # 3. Check text truncation (estimated height vs actual box height)
    for s in text_shapes:
        text = get_text_content(s)
        font_size = get_font_size(s)
        box_width = s.width if s.width else SLIDE_W
        estimated_h = estimate_text_height(text, font_size, box_width)
        actual_h = s.height if s.height else 0
        
        if estimated_h > actual_h * 1.3:  # needs 30% more than available
            issues.append({
                "slide": slide_idx + 1,
                "severity": "warning",
                "type": "text_truncation",
                "shape": s.name,
                "detail": f"文字可能被截断: 需要 {estimated_h // 12700}pt 高度，实际 {actual_h // 12700}pt",
                "text_preview": text[:40],
            })
    
    # 4. Content density check
    text_count = len(text_shapes)
    if text_count > 20:
        issues.append({
            "slide": slide_idx + 1,
            "severity": "info",
            "type": "high_density",
            "detail": f"内容密度较高: {text_count} 个文字元素",
        })
    
    return issues


def validate_pptx(pptx_path):
    """Validate an entire PPTX file. Returns list of issues."""
    prs = Presentation(pptx_path)
    all_issues = []
    
    for idx, slide in enumerate(prs.slides):
        issues = validate_slide(slide, idx)
        all_issues.extend(issues)
    
    return all_issues


def format_report(issues):
    """Format issues into a readable report string."""
    if not issues:
        return "✅ 无排版问题"
    
    errors = [i for i in issues if i["severity"] == "error"]
    warnings = [i for i in issues if i["severity"] == "warning"]
    infos = [i for i in issues if i["severity"] == "info"]
    
    lines = [f"排版检查: {len(errors)} 错误, {len(warnings)} 警告, {len(infos)} 提示\n"]
    
    for issue in errors:
        lines.append(f"🔴 P{issue['slide']}: {issue['detail']}")
    for issue in warnings:
        lines.append(f"🟡 P{issue['slide']}: {issue['detail']}")
    for issue in infos:
        lines.append(f"🟢 P{issue['slide']}: {issue['detail']}")
    
    return "\n".join(lines)


if __name__ == "__main__":
    import sys
    path = sys.argv[1] if len(sys.argv) > 1 else "/tmp/output.pptx"
    issues = validate_pptx(path)
    print(format_report(issues))
