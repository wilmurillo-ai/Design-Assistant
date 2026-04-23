"""
PPT Generation Engine v2 — Enhanced Kimi-style PPTX generator
Supports multiple themes, 11+ layout types, page numbers, logo, and Chinese font auto-detection.
Backward-compatible with v1 slide data format.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os, json, glob, random, subprocess, shutil

# ============================================================
# Design System — multi-theme support
# ============================================================

THEMES = {
    "dark": {
        "bg": "1A1D21",
        "bg2": "263238",
        "accent": "4DB6AC",
        "accent2": "FF7043",
        "text_primary": "E0E0E0",
        "text_secondary": "9E9E9E",
        "text_white": "FFFFFF",
        "divider": "4DB6AC",
        "card_bg": "263238",
        "card_border": "37474F",
        "table_header_bg": "4DB6AC",
        "table_row_alt": "1E2428",
        "page_num": "616161",
    },
    "light": {
        "bg": "FFFFFF",
        "bg2": "F5F5F5",
        "accent": "1976D2",
        "accent2": "E65100",
        "text_primary": "212121",
        "text_secondary": "757575",
        "text_white": "FFFFFF",
        "divider": "1976D2",
        "card_bg": "F5F5F5",
        "card_border": "E0E0E0",
        "table_header_bg": "1976D2",
        "table_row_alt": "FAFAFA",
        "page_num": "BDBDBD",
    },
    "blue": {
        "bg": "0D1B2A",
        "bg2": "1B2838",
        "accent": "FFC107",
        "accent2": "FF5722",
        "text_primary": "E0E6ED",
        "text_secondary": "8899A6",
        "text_white": "FFFFFF",
        "divider": "FFC107",
        "card_bg": "1B2838",
        "card_border": "2C3E50",
        "table_header_bg": "FFC107",
        "table_row_alt": "112233",
        "page_num": "546E7A",
    },
    "green": {
        "bg": "1B2D1B",
        "bg2": "2E4A2E",
        "accent": "66BB6A",
        "accent2": "FFB74D",
        "text_primary": "E0E8E0",
        "text_secondary": "8DA68D",
        "text_white": "FFFFFF",
        "divider": "66BB6A",
        "card_bg": "2E4A2E",
        "card_border": "3E5A3E",
        "table_header_bg": "66BB6A",
        "table_row_alt": "223322",
        "page_num": "4E6B4E",
    },
}

# Active theme — set by generate_ppt() or default to dark
THEME = THEMES["dark"]

# ============================================================
# Font detection
# ============================================================

_FONT_CACHE = {}


def _detect_fonts():
    """Detect available Chinese fonts with fallback chain."""
    if _FONT_CACHE:
        return _FONT_CACHE

    # Preferred font chains
    zh_candidates = [
        "MiSans",
        "Microsoft YaHei",
        "WenQuanYi Zen Hei",
        "Noto Sans CJK SC",
        "SimHei",
        "PingFang SC",
        "Hiragino Sans GB",
    ]
    en_candidates = [
        "Segoe UI",
        "Inter",
        "Roboto",
        "Helvetica Neue",
        "Arial",
        "DejaVu Sans",
    ]

    available = set()
    try:
        result = subprocess.run(
            ["fc-list", "--format", "%{family}\n"],
            capture_output=True, text=True, timeout=5
        )
        for line in result.stdout.splitlines():
            for name in line.split(","):
                available.add(name.strip())
    except Exception:
        pass

    zh_font = "Microsoft YaHei"  # default fallback
    for f in zh_candidates:
        if f in available:
            zh_font = f
            break

    en_font = "Segoe UI"
    for f in en_candidates:
        if f in available:
            en_font = f
            break

    _FONT_CACHE["zh"] = zh_font
    _FONT_CACHE["en"] = en_font
    return _FONT_CACHE


def get_font(kind="zh"):
    fonts = _detect_fonts()
    return fonts.get(kind, fonts["zh"])


# ============================================================
# Constants
# ============================================================

SLIDE_W = 12192000  # EMU — 16:9 standard
SLIDE_H = 6858000
MARGIN = 381000  # ~0.42 inch
CONTENT_W = SLIDE_W - 2 * MARGIN

ASSETS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")


# ============================================================
# Utility helpers
# ============================================================

def hex_to_rgb(hex_str):
    return RGBColor(int(hex_str[:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def set_slide_bg(slide, color_hex):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color_hex)


def get_random_bg(category):
    """Pick a random background image from a category folder."""
    cat_dir = os.path.join(ASSETS_DIR, f"bg_{category}")
    if not os.path.isdir(cat_dir):
        return None
    imgs = glob.glob(os.path.join(cat_dir, "*.jpg")) + glob.glob(os.path.join(cat_dir, "*.png"))
    return random.choice(imgs) if imgs else None


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=None, bold=False, align=PP_ALIGN.LEFT, font_name=None,
                 line_spacing=1.3, anchor=None):
    """Add a text box with precise positioning."""
    if color is None:
        color = THEME["text_primary"]
    if font_name is None:
        font_name = get_font("zh")

    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    if anchor is not None:
        tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align

    # Line spacing
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
    spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
    lnSpc.append(spcPct)
    pPr.append(lnSpc)

    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = hex_to_rgb(color)
    run.font.bold = bold
    run.font.name = font_name

    # East Asian font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn('a:ea'), {'typeface': get_font("zh")})
    rPr.append(ea)

    return txBox


def add_shape_rect(slide, left, top, width, height, fill_color, alpha=100,
                   rounded=False, corner_radius=0.08, border_color=None, border_width=None):
    """Add a rectangle/rounded-rectangle shape with optional border."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Emu(left), Emu(top), Emu(width), Emu(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_color)

    if border_color and border_width:
        shape.line.color.rgb = hex_to_rgb(border_color)
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()

    # Set corner radius for rounded rect
    if rounded and hasattr(shape, 'adjustments') and len(shape.adjustments) > 0:
        shape.adjustments[0] = corner_radius

    if alpha < 100:
        _set_shape_alpha(shape, alpha)

    return shape


def _set_shape_alpha(shape, alpha_pct):
    """Set alpha transparency on a shape's solid fill."""
    spPr = shape._element.find(qn('p:spPr'))
    if spPr is None:
        return
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is not None:
        srgbClr = solidFill.find(qn('a:srgbClr'))
        if srgbClr is not None:
            # Remove existing alpha
            for old in srgbClr.findall(qn('a:alpha')):
                srgbClr.remove(old)
            alphaElem = srgbClr.makeelement(qn('a:alpha'), {'val': str(alpha_pct * 1000)})
            srgbClr.append(alphaElem)


def add_shadow(shape, blur=Pt(6), dist=Pt(3), direction=2700000, alpha=40):
    """Add outer shadow effect to a shape via XML."""
    spPr = shape._element.find(qn('p:spPr'))
    if spPr is None:
        return

    # Remove existing effectLst
    for old in spPr.findall(qn('a:effectLst')):
        spPr.remove(old)


# ============================================================
# Visual Enhancement Helpers
# ============================================================

def add_circle(slide, cx, cy, size, fill_color, alpha=100):
    """Add a circle shape centered at (cx, cy)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Emu(cx - size // 2), Emu(cy - size // 2),
        Emu(size), Emu(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
    shape.line.fill.background()
    if alpha < 100:
        _set_shape_alpha(shape, alpha)
    return shape


def add_bullet_dot(slide, x, y, color, size=60000):
    """Add a small colored circle as a bullet indicator."""
    return add_circle(slide, x + size // 2, y + size // 2, size, color, alpha=95)


def add_accent_bar(slide, x, y, height, color, width=36000, alpha=95):
    """Add a thin vertical accent bar (left-side indicator)."""
    return add_shape_rect(slide, x, y, width, height, color, alpha=alpha,
                          rounded=True, corner_radius=0.5)


def add_decorative_circles(slide, theme):
    """Add subtle decorative circles for visual depth. Higher contrast for dark themes."""
    accent = theme.get("accent", "4DB6AC")
    accent2 = theme.get("accent2", "FF7043")
    # Top-right area — large soft circle
    add_circle(slide, SLIDE_W - 500000, 400000, 800000, accent, alpha=12)
    # Bottom-left area — medium circle
    add_circle(slide, 400000, SLIDE_H - 400000, 500000, accent2, alpha=10)
    # Small accent dot top-right
    add_circle(slide, SLIDE_W - 800000, 250000, 100000, accent, alpha=20)


def add_section_number_badge(slide, x, y, number, theme):
    """Add a large styled number badge for section pages."""
    size = 700000
    # Outer ring
    add_circle(slide, x + size // 2, y + size // 2, size, theme["accent"], alpha=15)
    # Inner circle
    add_circle(slide, x + size // 2, y + size // 2, int(size * 0.7),
               theme["accent"], alpha=25)
    # Number text
    add_text_box(slide, x, y, size, size,
                 str(number), font_size=36, color=theme["accent"],
                 bold=True, align=PP_ALIGN.CENTER, font_name=get_font("en"))

    effectLst = spPr.makeelement(qn('a:effectLst'), {})
    outerShdw = effectLst.makeelement(qn('a:outerShdw'), {
        'blurRad': str(int(blur)),
        'dist': str(int(dist)),
        'dir': str(direction),
        'rotWithShape': '0',
    })
    srgbClr = outerShdw.makeelement(qn('a:srgbClr'), {'val': '000000'})
    alphaElem = srgbClr.makeelement(qn('a:alpha'), {'val': str(alpha * 1000)})
    srgbClr.append(alphaElem)
    outerShdw.append(srgbClr)
    effectLst.append(outerShdw)
    spPr.append(effectLst)


def add_image_bg(slide, image_path, alpha_pct=40):
    """Add a background image with transparency."""
    if not os.path.exists(image_path):
        return None
    pic = slide.shapes.add_picture(image_path, 0, 0, Emu(SLIDE_W), Emu(SLIDE_H))

    blip = pic._element.find('.//' + qn('a:blip'))
    if blip is not None:
        # Remove old alphaModFix
        for old in blip.findall(qn('a:alphaModFix')):
            blip.remove(old)
        alphaModFix = blip.makeelement(qn('a:alphaModFix'), {'amt': str(alpha_pct * 1000)})
        blip.append(alphaModFix)

    # Move to back
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)
    return pic


def add_gradient_overlay(slide, theme=None):
    """Add a dark gradient overlay for text readability."""
    if theme is None:
        theme = THEME
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Emu(SLIDE_W), Emu(SLIDE_H))
    shape.line.fill.background()

    spPr = shape._element.find(qn('p:spPr'))
    for child in list(spPr):
        if child.tag == qn('a:solidFill'):
            spPr.remove(child)

    gradFill = spPr.makeelement(qn('a:gradFill'), {'rotWithShape': '1'})
    gsLst = gradFill.makeelement(qn('a:gsLst'), {})

    bg_color = theme["bg"]
    bg2_color = theme.get("bg2", theme["bg"])
    for pos, color, alpha in [(0, bg_color, 80), (50000, bg_color, 85), (100000, bg2_color, 92)]:
        gs = gsLst.makeelement(qn('a:gs'), {'pos': str(pos)})
        srgbClr = gs.makeelement(qn('a:srgbClr'), {'val': color})
        alphaElem = srgbClr.makeelement(qn('a:alpha'), {'val': str(alpha * 1000)})
        srgbClr.append(alphaElem)
        gs.append(srgbClr)
        gsLst.append(gs)

    lin = gradFill.makeelement(qn('a:lin'), {'ang': '2700000', 'scaled': '1'})
    gradFill.append(gsLst)
    gradFill.append(lin)
    spPr.append(gradFill)
    return shape


def add_page_number(slide, page_num, total=None):
    """Add page number at bottom-right."""
    text = f"{page_num}" if total is None else f"{page_num} / {total}"
    add_text_box(
        slide, SLIDE_W - MARGIN - 600000, SLIDE_H - 300000, 600000, 200000,
        text, font_size=9, color=THEME.get("page_num", THEME["text_secondary"]),
        align=PP_ALIGN.RIGHT, font_name=get_font("en"),
    )


def add_logo(slide, logo_path, position="top-right", size=400000):
    """Add a logo image to the slide."""
    if not logo_path or not os.path.exists(logo_path):
        return None
    positions = {
        "top-right": (SLIDE_W - MARGIN - size, MARGIN // 2),
        "top-left": (MARGIN, MARGIN // 2),
        "bottom-right": (SLIDE_W - MARGIN - size, SLIDE_H - MARGIN // 2 - size),
        "bottom-left": (MARGIN, SLIDE_H - MARGIN // 2 - size),
    }
    left, top = positions.get(position, positions["top-right"])
    pic = slide.shapes.add_picture(logo_path, Emu(left), Emu(top), Emu(size), Emu(size))
    return pic


# ============================================================
# Slide Builders
# ============================================================

def build_cover(prs, title, subtitle, author="", date="", bg_image=None, **kw):
    """Build a cover slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, THEME["bg"])

    if bg_image:
        add_image_bg(slide, bg_image, alpha_pct=35)
        add_gradient_overlay(slide)

    y = SLIDE_H // 2 - 800000

    # Accent bar
    add_shape_rect(slide, MARGIN, y - 80000, 600000, 38100, THEME["accent"])

    # Title
    add_text_box(slide, MARGIN, y, CONTENT_W, 600000,
                 title, font_size=36, color=THEME["text_white"], bold=True)

    # Subtitle
    add_text_box(slide, MARGIN, y + 650000, CONTENT_W, 400000,
                 subtitle, font_size=18, color=THEME["text_secondary"])

    # Author & date
    if author or date:
        info = f"{author}    {date}" if author and date else (author or date)
        add_text_box(slide, MARGIN, SLIDE_H - 600000, CONTENT_W, 300000,
                     info, font_size=12, color=THEME["text_secondary"])

    return slide


def build_section(prs, number, title, subtitle="", bg_image=None, **kw):
    """Build a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, THEME["bg"])

    if bg_image:
        add_image_bg(slide, bg_image, alpha_pct=25)
        add_gradient_overlay(slide)

    y_base = SLIDE_H // 2 - 600000

    # Number (watermark style, above title with clear gap)
    add_text_box(slide, MARGIN, y_base - 900000, 1500000, 900000,
                 f"{number:02d}", font_size=72, color=THEME["accent"],
                 bold=True, font_name=get_font("en"))

    # Accent line
    add_shape_rect(slide, MARGIN, y_base + 50000, 500000, 30000, THEME["accent"])

    # Title (moved down to avoid overlap with number)
    add_text_box(slide, MARGIN, y_base + 150000, CONTENT_W, 500000,
                 title, font_size=32, color=THEME["text_white"], bold=True)

    if subtitle:
        add_text_box(slide, MARGIN, y_base + 700000, CONTENT_W, 500000,
                     subtitle, font_size=16, color=THEME["text_secondary"])

    return slide


def build_content(prs, tag, title, bullets, subtitle="", **kw):
    """Build a content slide. Each bullet in its own visible card, filling the page."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, THEME["bg"])
    add_decorative_circles(slide, THEME)

    y = MARGIN
    bottom_margin = 350000

    # Tag with accent dot
    add_bullet_dot(slide, MARGIN, y + 30000, THEME["accent"], size=50000)
    add_text_box(slide, MARGIN + 70000, y, CONTENT_W - 70000, 200000,
                 tag, font_size=11, color=THEME["accent"], bold=True,
                 font_name=get_font("en"))
    y += 220000

    # Title
    add_text_box(slide, MARGIN, y, CONTENT_W, 450000,
                 title, font_size=28, color=THEME["text_white"], bold=True)
    y += 480000

    if subtitle:
        add_text_box(slide, MARGIN, y, CONTENT_W, 220000,
                     subtitle, font_size=12, color=THEME["text_secondary"])
        y += 250000

    # Each bullet gets its own card — fill remaining space
    n_bullets = len(bullets)
    if n_bullets > 0:
        available_h = SLIDE_H - y - bottom_margin
        card_gap = 50000
        total_gaps = (n_bullets - 1) * card_gap
        card_h = (available_h - total_gaps) // n_bullets
        card_h = max(card_h, 250000)
        bfont = 16 if n_bullets <= 3 else (14 if n_bullets <= 5 else 12)

        for idx, bullet in enumerate(bullets):
            # Card with pre-blended fill + visible border
            card = add_shape_rect(slide, MARGIN, y, CONTENT_W, card_h,
                                 "2A4040", alpha=100, rounded=True, corner_radius=0.04,
                                 border_color=THEME["accent"], border_width=1.5)
            add_shadow(card)

            # Left accent bar
            bar_h = int(card_h * 0.65)
            add_accent_bar(slide, MARGIN + 18000, y + (card_h - bar_h) // 2,
                          bar_h, THEME["accent"], width=32000, alpha=95)

            # Number badge
            badge_size = min(int(card_h * 0.55), 280000)
            add_text_box(slide, MARGIN + 70000, y + (card_h - badge_size) // 2,
                        badge_size, badge_size,
                        str(idx + 1), font_size=18, color=THEME["accent"],
                        bold=True, align=PP_ALIGN.CENTER, font_name=get_font("en"))

            text_left = MARGIN + 70000 + badge_size + 20000
            text_w = CONTENT_W - (text_left - MARGIN) - 50000

            if isinstance(bullet, str):
                add_text_box(slide, text_left, y + (card_h - int(card_h * 0.65)) // 2,
                             text_w, int(card_h * 0.65),
                             bullet, font_size=bfont, color=THEME["text_primary"])
            elif isinstance(bullet, dict):
                key = bullet.get("key", "")
                val = bullet.get("val", "")
                add_text_box(slide, text_left, y + int(card_h * 0.12),
                             text_w, int(card_h * 0.38),
                             key, font_size=bfont, color=THEME["accent"], bold=True)
                add_text_box(slide, text_left, y + int(card_h * 0.50),
                             text_w, int(card_h * 0.42),
                             val, font_size=bfont - 1, color=THEME["text_primary"])

            y += card_h + card_gap

    return slide


def build_table_slide(prs, tag, title, headers, rows, subtitle="", **kw):
    """Build a table slide with card-style rows and alternating colors. Adaptive + visual."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, THEME["bg"])

    # Decorative background
    add_decorative_circles(slide, THEME)

    y = MARGIN
    bottom_margin = 400000

    # Tag with accent dot
    add_bullet_dot(slide, MARGIN, y + 30000, THEME["accent"], size=50000)
    add_text_box(slide, MARGIN + 70000, y, CONTENT_W - 70000, 200000,
                 tag, font_size=11, color=THEME["accent"], bold=True,
                 font_name=get_font("en"))
    y += 250000

    add_text_box(slide, MARGIN, y, CONTENT_W, 400000,
                 title, font_size=26, color=THEME["text_white"], bold=True)
    y += 500000

    if subtitle:
        add_text_box(slide, MARGIN, y, CONTENT_W, 200000,
                     subtitle, font_size=12, color=THEME["text_secondary"])
        y += 250000

    n_cols = len(headers)
    col_w = CONTENT_W // n_cols

    # Adaptive row sizing
    n_rows = len(rows)
    available_h = SLIDE_H - y - bottom_margin - 380000
    row_slot = min(available_h // max(n_rows, 1), 500000)
    row_h = int(row_slot * 0.85)
    header_h = min(int(row_slot * 0.95), 380000)
    row_font = 13 if n_rows <= 6 else (12 if n_rows <= 8 else 11)
    header_font = 13 if n_rows <= 6 else 12

    # Header — rounded card style
    header_shape = add_shape_rect(
        slide, MARGIN, y, CONTENT_W, header_h,
        THEME.get("table_header_bg", THEME["accent"]),
        alpha=40, rounded=True, corner_radius=0.04
    )

    for i, h in enumerate(headers):
        add_text_box(slide, MARGIN + i * col_w + 60000, y + 40000,
                     col_w - 120000, header_h - 80000,
                     h, font_size=header_font, color=THEME["accent"], bold=True)
    y += header_h + 20000

    # Rows with alternating bg + left accent indicator
    for row_idx, row in enumerate(rows):
        if row_idx % 2 == 0:
            add_shape_rect(
                slide, MARGIN, y, CONTENT_W, row_h,
                THEME.get("table_row_alt", THEME["card_bg"]), alpha=70,
                rounded=True, corner_radius=0.03
            )
        # Left accent bar on each row
        add_accent_bar(slide, MARGIN + 10000, y + int(row_h * 0.15),
                      int(row_h * 0.7), THEME["accent"], width=26000, alpha=55)

        for i, cell in enumerate(row):
            add_text_box(slide, MARGIN + i * col_w + 60000, y + 30000,
                         col_w - 120000, row_h - 60000,
                         str(cell), font_size=row_font,
                         color=THEME["text_primary"], bold=(i == 0))
        y += row_slot

    return slide


def build_two_column(prs, tag, title, left_items, right_items,
                     left_title="", right_title="", **kw):
    """Build a two-column comparison slide. Glass panels + adaptive layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, THEME["bg"])

    add_decorative_circles(slide, THEME)

    y = MARGIN
    bottom_margin = 400000

    # Tag with accent dot
    add_bullet_dot(slide, MARGIN, y + 30000, THEME["accent"], size=50000)
    add_text_box(slide, MARGIN + 70000, y, CONTENT_W - 70000, 200000,
                 tag, font_size=11, color=THEME["accent"], bold=True,
                 font_name=get_font("en"))
    y += 250000

    # Title — larger
    add_text_box(slide, MARGIN, y, CONTENT_W, 500000,
                 title, font_size=28, color=THEME["text_white"], bold=True)
    y += 550000

    col_gap = 120000
    half_w = (CONTENT_W - col_gap) // 2
    panel_h = SLIDE_H - y - bottom_margin

    # Adaptive item sizing
    max_items = max(len(left_items), len(right_items))
    title_h = 380000 if left_title or right_title else 0
    inner_h = panel_h - 100000 - title_h
    item_slot = min(inner_h // max(max_items, 1), 550000)
    item_h = int(item_slot * 0.75)
    item_font = 15 if max_items <= 4 else (14 if max_items <= 5 else 13)
    dot_size = 45000 if item_font >= 14 else 35000

    # Left panel — glass card
    lx = MARGIN
    left_panel = add_shape_rect(slide, lx, y, half_w, panel_h,
                               "2A4040", alpha=100, rounded=True, corner_radius=0.04, border_color=THEME["accent"], border_width=1.5)
    add_shadow(left_panel, blur=Pt(6), dist=Pt(3), alpha=25)
    # Left accent bar
    add_accent_bar(slide, lx + 15000, y + 30000, panel_h - 60000,
                  THEME["accent"], width=32000, alpha=80)

    ly = y + 50000
    if left_title:
        add_text_box(slide, lx + 60000, ly, half_w - 80000, 300000,
                     left_title, font_size=16, color=THEME["accent"], bold=True)
        add_shape_rect(slide, lx + 60000, ly + 280000, half_w // 2, 16000,
                      THEME["accent"], alpha=40)
        ly += title_h
    for item in left_items:
        add_bullet_dot(slide, lx + 70000, ly + (item_slot - dot_size) // 2,
                      THEME["accent"], size=dot_size)
        add_text_box(slide, lx + 140000, ly, half_w - 170000, item_h,
                     item, font_size=item_font, color=THEME["text_primary"])
        ly += item_slot

    # Right panel — glass card
    rx = MARGIN + half_w + col_gap
    right_panel = add_shape_rect(slide, rx, y, half_w, panel_h,
                                "3D2D2A", alpha=100, rounded=True, corner_radius=0.04, border_color=THEME["accent2"], border_width=1.5)
    add_shadow(right_panel, blur=Pt(6), dist=Pt(3), alpha=25)
    # Right accent bar
    add_accent_bar(slide, rx + 15000, y + 30000, panel_h - 60000,
                  THEME["accent2"], width=32000, alpha=80)

    ry = y + 50000
    if right_title:
        add_text_box(slide, rx + 60000, ry, half_w - 80000, 300000,
                     right_title, font_size=16, color=THEME["accent2"], bold=True)
        add_shape_rect(slide, rx + 60000, ry + 280000, half_w // 2, 16000,
                      THEME["accent2"], alpha=40)
        ry += title_h
    for item in right_items:
        add_bullet_dot(slide, rx + 70000, ry + (item_slot - dot_size) // 2,
                      THEME["accent2"], size=dot_size)
        add_text_box(slide, rx + 140000, ry, half_w - 170000, item_h,
                     item, font_size=item_font, color=THEME["text_primary"])
        ry += item_slot

    return slide


def build_cards(prs, tag, title, cards, **kw):
    """Build a slide with card-style layout (rounded + shadow)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, THEME["bg"])

    y = MARGIN

    add_text_box(slide, MARGIN, y, CONTENT_W, 200000,
                 tag, font_size=11, color=THEME["accent"], bold=True,
                 font_name=get_font("en"))
    y += 250000

    add_text_box(slide, MARGIN, y, CONTENT_W, 400000,
                 title, font_size=26, color=THEME["text_white"], bold=True)
    y += 500000

    n = len(cards)
    cols = min(n, 3)
    card_w = (CONTENT_W - (cols - 1) * 120000) // cols
    card_h = 2800000

    for i, card in enumerate(cards):
        col = i % cols
        row = i // cols
        cx = MARGIN + col * (card_w + 120000)
        cy = y + row * (card_h + 120000)

        # Rounded card with shadow
        card_shape = add_shape_rect(
            slide, cx, cy, card_w, card_h,
            THEME["card_bg"], alpha=85, rounded=True, corner_radius=0.06
        )
        add_shadow(card_shape)

        # Card accent top bar
        add_shape_rect(slide, cx + 50000, cy + 20000, card_w - 100000, 25000,
                       THEME["accent"], rounded=True, corner_radius=0.5)

        # Card title
        add_text_box(slide, cx + 100000, cy + 100000, card_w - 200000, 300000,
                     card.get("title", ""), font_size=16,
                     color=THEME["text_white"], bold=True)

        # Card items
        content_y = cy + 450000
        for line in card.get("items", []):
            add_text_box(slide, cx + 100000, content_y, card_w - 200000, 240000,
                         f"•  {line}", font_size=11, color=THEME["text_primary"])
            content_y += 260000

    return slide


def build_closing(prs, title, subtitle="", bg_image=None, **kw):
    """Build a closing/Q&A slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, THEME["bg"])

    if bg_image:
        add_image_bg(slide, bg_image, alpha_pct=25)
        add_gradient_overlay(slide)

    y_center = SLIDE_H // 2 - 300000

    add_text_box(slide, MARGIN, y_center, CONTENT_W, 500000,
                 title, font_size=40, color=THEME["text_white"], bold=True,
                 align=PP_ALIGN.CENTER)

    if subtitle:
        add_text_box(slide, MARGIN, y_center + 550000, CONTENT_W, 500000,
                     subtitle, font_size=16, color=THEME["text_secondary"],
                     align=PP_ALIGN.CENTER)

    line_w = 600000
    add_shape_rect(slide, (SLIDE_W - line_w) // 2, y_center - 100000,
                   line_w, 30000, THEME["accent"])

    return slide


# === NEW LAYOUTS ===

def build_timeline(prs, tag, title, events, subtitle="", **kw):
    """
    Build a horizontal timeline slide.
    events: list of {"date": "2024-Q1", "title": "Event", "desc": "Details"}
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, THEME["bg"])

    y = MARGIN

    add_text_box(slide, MARGIN, y, CONTENT_W, 200000,
                 tag, font_size=11, color=THEME["accent"], bold=True,
                 font_name=get_font("en"))
    y += 250000

    add_text_box(slide, MARGIN, y, CONTENT_W, 400000,
                 title, font_size=26, color=THEME["text_white"], bold=True)
    y += 450000

    if subtitle:
        add_text_box(slide, MARGIN, y, CONTENT_W, 200000,
                     subtitle, font_size=12, color=THEME["text_secondary"])
        y += 250000

    n = len(events)
    if n == 0:
        return slide

    # Timeline horizontal line
    line_y = y + 800000
    add_shape_rect(slide, MARGIN, line_y, CONTENT_W, 20000, THEME["accent"], alpha=60)

    # Events spread along the line
    segment_w = CONTENT_W // n
    for i, ev in enumerate(events):
        cx = MARGIN + i * segment_w + segment_w // 2
        is_above = (i % 2 == 0)

        # Dot on timeline
        dot_r = 60000
        add_shape_rect(slide, cx - dot_r, line_y - dot_r // 2 + 10000,
                       dot_r * 2, dot_r, THEME["accent"], rounded=True, corner_radius=0.5)

        # Date label
        date_str = ev.get("date", "")
        evt_title = ev.get("title", "")
        evt_desc = ev.get("desc", "")

        text_w = segment_w - 60000

        if is_above:
            # Content above the line
            add_text_box(slide, cx - text_w // 2, line_y - 700000, text_w, 200000,
                         date_str, font_size=10, color=THEME["accent"], bold=True,
                         align=PP_ALIGN.CENTER, font_name=get_font("en"))
            add_text_box(slide, cx - text_w // 2, line_y - 500000, text_w, 200000,
                         evt_title, font_size=13, color=THEME["text_white"], bold=True,
                         align=PP_ALIGN.CENTER)
            if evt_desc:
                add_text_box(slide, cx - text_w // 2, line_y - 320000, text_w, 260000,
                             evt_desc, font_size=10, color=THEME["text_secondary"],
                             align=PP_ALIGN.CENTER)
        else:
            # Content below the line
            add_text_box(slide, cx - text_w // 2, line_y + 100000, text_w, 200000,
                         date_str, font_size=10, color=THEME["accent"], bold=True,
                         align=PP_ALIGN.CENTER, font_name=get_font("en"))
            add_text_box(slide, cx - text_w // 2, line_y + 280000, text_w, 200000,
                         evt_title, font_size=13, color=THEME["text_white"], bold=True,
                         align=PP_ALIGN.CENTER)
            if evt_desc:
                add_text_box(slide, cx - text_w // 2, line_y + 460000, text_w, 260000,
                             evt_desc, font_size=10, color=THEME["text_secondary"],
                             align=PP_ALIGN.CENTER)

    return slide


def build_stats(prs, tag, title, stats, subtitle="", **kw):
    """Build a KPI/stats slide with big numbers in visible cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, THEME["bg"])
    add_decorative_circles(slide, THEME)

    y = MARGIN
    bottom_margin = 350000

    add_bullet_dot(slide, MARGIN, y + 30000, THEME["accent"], size=50000)
    add_text_box(slide, MARGIN + 70000, y, CONTENT_W - 70000, 200000,
                 tag, font_size=11, color=THEME["accent"], bold=True,
                 font_name=get_font("en"))
    y += 220000

    add_text_box(slide, MARGIN, y, CONTENT_W, 450000,
                 title, font_size=28, color=THEME["text_white"], bold=True)
    y += 480000

    if subtitle:
        add_text_box(slide, MARGIN, y, CONTENT_W, 200000,
                     subtitle, font_size=12, color=THEME["text_secondary"])
        y += 250000

    n = len(stats)
    if n == 0:
        return slide

    cols = min(n, 4)
    col_gap = 80000
    stat_w = (CONTENT_W - (cols - 1) * col_gap) // cols

    # Cards fill remaining vertical space
    available_h = SLIDE_H - y - bottom_margin
    card_h = min(available_h, 2400000)
    card_h = max(card_h, 1400000)
    # Center vertically
    card_y = y + (available_h - card_h) // 2

    for i, s in enumerate(stats):
        col = i % cols
        row = i // cols
        sx = MARGIN + col * (stat_w + col_gap)
        sy = card_y + row * (card_h + 100000)

        card = add_shape_rect(slide, sx, sy, stat_w, card_h,
                              "2A4040", alpha=100, rounded=True, corner_radius=0.06,
                              border_color=THEME["accent"], border_width=1.2)
        add_shadow(card)

        # Layout: number at 25%, label at 55%, desc at 72%
        num_y = sy + int(card_h * 0.10)
        num_h = int(card_h * 0.35)
        label_y = sy + int(card_h * 0.50)
        label_h = int(card_h * 0.16)
        desc_y = sy + int(card_h * 0.68)
        desc_h = int(card_h * 0.22)

        add_text_box(slide, sx, num_y, stat_w, num_h,
                     s.get("value", ""), font_size=48, color=THEME["accent"],
                     bold=True, align=PP_ALIGN.CENTER, font_name=get_font("en"))
        add_text_box(slide, sx + 50000, label_y, stat_w - 100000, label_h,
                     s.get("label", ""), font_size=16, color=THEME["text_white"],
                     bold=True, align=PP_ALIGN.CENTER)
        desc = s.get("desc", "")
        if desc:
            add_text_box(slide, sx + 50000, desc_y, stat_w - 100000, desc_h,
                         desc, font_size=11, color=THEME["text_secondary"],
                         align=PP_ALIGN.CENTER)

    return slide


def build_quote(prs, quote, author="", source="", bg_image=None, **kw):
    """Build a quote/citation slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, THEME["bg"])

    if bg_image:
        add_image_bg(slide, bg_image, alpha_pct=20)
        add_gradient_overlay(slide)

    y_center = SLIDE_H // 2 - 600000

    # Large quote mark
    add_text_box(slide, MARGIN + 200000, y_center - 200000, 500000, 500000,
                 "\u201C", font_size=80, color=THEME["accent"], bold=True,
                 align=PP_ALIGN.LEFT, font_name=get_font("en"))

    # Quote text
    add_text_box(slide, MARGIN + 400000, y_center + 100000, CONTENT_W - 800000, 800000,
                 quote, font_size=22, color=THEME["text_white"],
                 align=PP_ALIGN.LEFT, line_spacing=1.5)

    # Attribution
    if author:
        attr_text = f"— {author}"
        if source:
            attr_text += f"  ·  {source}"
        add_text_box(slide, MARGIN + 400000, y_center + 1000000,
                     CONTENT_W - 800000, 250000,
                     attr_text, font_size=14, color=THEME["accent"])

    # Accent line
    add_shape_rect(slide, MARGIN + 350000, y_center + 100000, 30000, 700000,
                   THEME["accent"], alpha=50)

    return slide


def build_image_text(prs, tag, title, text_content, image_path=None,
                     layout="left", bullets=None, **kw):
    """
    Build an image+text slide.
    layout: "left" (image left, text right) or "right" (text left, image right)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, THEME["bg"])

    y = MARGIN

    add_text_box(slide, MARGIN, y, CONTENT_W, 200000,
                 tag, font_size=11, color=THEME["accent"], bold=True,
                 font_name=get_font("en"))
    y += 250000

    add_text_box(slide, MARGIN, y, CONTENT_W, 400000,
                 title, font_size=26, color=THEME["text_white"], bold=True)
    y += 500000

    half_w = (CONTENT_W - 200000) // 2
    img_area_w = half_w
    img_area_h = SLIDE_H - y - MARGIN

    if layout == "left":
        img_x = MARGIN
        text_x = MARGIN + half_w + 200000
    else:
        text_x = MARGIN
        img_x = MARGIN + half_w + 200000

    # Image area
    if image_path and os.path.exists(image_path):
        # Add image with rounded clip (placeholder rect + image)
        pic = slide.shapes.add_picture(
            image_path, Emu(img_x), Emu(y),
            Emu(img_area_w), Emu(img_area_h)
        )
    else:
        # Placeholder box
        add_shape_rect(slide, img_x, y, img_area_w, img_area_h,
                       THEME["card_bg"], alpha=40, rounded=True, corner_radius=0.04)
        add_text_box(slide, img_x, y + img_area_h // 2 - 100000,
                     img_area_w, 200000, "[Image]",
                     font_size=14, color=THEME["text_secondary"],
                     align=PP_ALIGN.CENTER)

    # Text content area
    ty = y
    if text_content:
        add_text_box(slide, text_x, ty, half_w, 600000,
                     text_content, font_size=14, color=THEME["text_primary"],
                     line_spacing=1.5)
        ty += 650000

    if bullets:
        for b in bullets:
            add_text_box(slide, text_x + 100000, ty, half_w - 100000, 260000,
                         f"•  {b}", font_size=13, color=THEME["text_primary"])
            ty += 280000

    return slide


# ============================================================
# Main generator
# ============================================================

BUILDER_MAP = {
    "cover": build_cover,
    "section": build_section,
    "content": build_content,
    "table": build_table_slide,
    "two_column": build_two_column,
    "cards": build_cards,
    "closing": build_closing,
    "timeline": build_timeline,
    "stats": build_stats,
    "quote": build_quote,
    "image_text": build_image_text,
}


def generate_ppt(slides_data, output_path, bg_images=None, theme="dark",
                 logo_path=None, logo_position="top-right", show_page_numbers=True):
    """
    Generate a PPTX from structured slide data.

    Args:
        slides_data: list of dicts, each with 'type' and type-specific fields
        output_path: path to save the .pptx file
        bg_images: dict mapping image keys to file paths
        theme: "dark" | "light" | "blue" | "green" or a custom dict
        logo_path: path to a logo image (optional)
        logo_position: "top-right" | "top-left" | "bottom-right" | "bottom-left"
        show_page_numbers: whether to add page numbers

    Returns:
        output_path
    """
    global THEME

    # Set theme
    if isinstance(theme, dict):
        THEME = theme
    elif isinstance(theme, str) and theme in THEMES:
        THEME = THEMES[theme]
    else:
        THEME = THEMES["dark"]

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)

    if bg_images is None:
        bg_images = {}

    total_slides = len(slides_data)

    for idx, sd in enumerate(slides_data):
        t = sd["type"]
        bg_key = sd.get("bg_image")
        bg = bg_images.get(bg_key) if bg_key else None

        # Auto-select bg from category if specified
        if bg is None and sd.get("bg_category"):
            bg = get_random_bg(sd["bg_category"])

        builder = BUILDER_MAP.get(t)
        if builder is None:
            continue

        # Build kwargs from slide data
        kwargs = dict(sd)
        kwargs.pop("type", None)
        kwargs.pop("bg_image", None)
        kwargs.pop("bg_category", None)

        # Handle bg_image param for builders that accept it
        if t in ("cover", "section", "closing", "quote"):
            kwargs["bg_image"] = bg
        elif t == "image_text":
            # image_path comes from bg_images or direct path
            img_key = sd.get("image")
            if img_key and img_key in bg_images:
                kwargs["image_path"] = bg_images[img_key]
            elif img_key and os.path.exists(img_key):
                kwargs["image_path"] = img_key

        slide = builder(prs, **kwargs)

        # Add logo (skip cover and closing)
        if logo_path and t not in ("cover", "closing"):
            add_logo(slide, logo_path, logo_position)

        # Add page numbers (skip cover)
        if show_page_numbers and t != "cover":
            add_page_number(slide, idx + 1, total_slides)

    prs.save(output_path)

    # Auto-validate after generation
    try:
        from ppt_validator import validate_pptx, format_report
        issues = validate_pptx(output_path)
        if issues:
            errors = [i for i in issues if i["severity"] == "error"]
            if errors:
                print(f"⚠️ PPT 排版检查发现 {len(errors)} 个错误:")
                for e in errors:
                    print(f"  🔴 P{e['slide']}: {e['detail']}")
            warnings = [i for i in issues if i["severity"] == "warning"]
            if warnings:
                print(f"📋 PPT 排版检查发现 {len(warnings)} 个警告:")
                for w in warnings:
                    print(f"  🟡 P{w['slide']}: {w['detail']}")
    except ImportError:
        pass  # validator not available, skip

    return output_path


# ============================================================
# CLI test
# ============================================================

if __name__ == "__main__":
    test_slides = [
        {
            "type": "cover",
            "title": "PPT Engine v2 演示",
            "subtitle": "多主题 · 11种布局 · 页码 · Logo · 中文字体自适应",
            "author": "PPT Generator",
            "date": "2026-03",
        },
        {
            "type": "section",
            "number": 1,
            "title": "新增布局类型",
            "subtitle": "timeline · stats · quote · image_text",
        },
        {
            "type": "content",
            "tag": "OVERVIEW",
            "title": "支持的布局类型",
            "bullets": [
                "cover — 封面页",
                "section — 章节分隔页",
                "content — 内容列表页",
                "table — 表格展示页",
                "two_column — 双栏对比页",
                "cards — 卡片布局页",
                "timeline — 时间线页",
                "stats — 数字/KPI展示页",
                "quote — 引用/名言页",
                "image_text — 图文混排页",
                "closing — 结束页",
            ],
        },
        {
            "type": "timeline",
            "tag": "ROADMAP",
            "title": "项目里程碑",
            "events": [
                {"date": "2026-Q1", "title": "v1 发布", "desc": "基础7种布局"},
                {"date": "2026-Q2", "title": "v2 发布", "desc": "新增4种布局 + 主题"},
                {"date": "2026-Q3", "title": "v3 计划", "desc": "动画 + 图表"},
                {"date": "2026-Q4", "title": "v4 展望", "desc": "AI 自动排版"},
            ],
        },
        {
            "type": "stats",
            "tag": "METRICS",
            "title": "关键指标",
            "stats": [
                {"value": "11", "label": "布局类型", "desc": "覆盖常见演示场景"},
                {"value": "4", "label": "内置主题", "desc": "dark / light / blue / green"},
                {"value": "25+", "label": "背景素材", "desc": "5 个分类 × 5 张"},
                {"value": "99%", "label": "中文兼容", "desc": "自动检测系统字体"},
            ],
        },
        {
            "type": "quote",
            "quote": "设计不仅仅是外观和感觉。设计是它如何运作的。",
            "author": "Steve Jobs",
            "source": "Apple Inc.",
        },
        {
            "type": "table",
            "tag": "COMPARISON",
            "title": "主题配色方案",
            "headers": ["主题", "背景色", "强调色", "适用场景"],
            "rows": [
                ["Dark", "#1A1D21", "#4DB6AC", "技术/正式"],
                ["Light", "#FFFFFF", "#1976D2", "商务/教育"],
                ["Blue", "#0D1B2A", "#FFC107", "金融/高端"],
                ["Green", "#1B2D1B", "#66BB6A", "环保/健康"],
            ],
        },
        {
            "type": "two_column",
            "tag": "V1 vs V2",
            "title": "版本对比",
            "left_title": "v1 功能",
            "left_items": ["7 种布局", "深色主题", "基础表格", "无页码"],
            "right_title": "v2 新增",
            "right_items": ["11 种布局", "4 种主题", "圆角卡片表格", "页码 + Logo"],
        },
        {
            "type": "cards",
            "tag": "FEATURES",
            "title": "核心特性",
            "cards": [
                {"title": "多主题", "items": ["Dark / Light", "Blue / Green", "自定义配色"]},
                {"title": "新布局", "items": ["Timeline 时间线", "Stats KPI展示", "Quote 引用"]},
                {"title": "增强", "items": ["圆角卡片", "阴影效果", "页码/Logo"]},
            ],
        },
        {
            "type": "closing",
            "title": "谢谢观看",
            "subtitle": "PPT Engine v2 — 让演示更专业",
        },
    ]

    out = generate_ppt(test_slides, "/tmp/test_ppt_v2.pptx", theme="dark", show_page_numbers=True)
    print(f"Generated: {out}")

    # Also generate light theme version
    out2 = generate_ppt(test_slides, "/tmp/test_ppt_v2_light.pptx", theme="light")
    print(f"Generated: {out2}")
