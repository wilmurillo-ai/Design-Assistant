#!/usr/bin/env python3
"""slide_data_audit.py — every hard number in a compiled .pptx must have a
row in the companion data-provenance.md.

This closes the PPT-to-provenance gap: cn_typo_scan checks PPT text for
typos, provenance_verify checks analysis.md hard numbers against
provenance, but neither of them checks that the hard numbers that
actually end up on the PPT slides are each traceable back to a
provenance row. When an agent hand-edits slide-NN.js after analysis.md
is frozen (common in Phase 4), the numbers on the slides can drift away
from what's in data-provenance.md. This script is the gate for that.

Usage:
    python3 slide_data_audit.py <deck.pptx> <data-provenance.md>
    # exit 0 → every hard number in the extracted slide text has a
    #          matching row in data-provenance.md
    # exit 1 → one or more slide numbers are missing a provenance row

Implementation: extract PPT text via python-pptx (uvx fallback), then
reuse the hard-number + normalize_variants machinery from
provenance_verify.py so unit-equivalence (亿元 ↔ 万元) and
thousands-comma handling stay consistent.
"""
from __future__ import annotations

import pathlib
import subprocess
import sys

from provenance_verify import (
    HARD_NUMBER,
    extract_provenance_corpus,
    normalize_variants,
)

EXTRACT_SCRIPT = """\
import sys
try:
    from pptx import Presentation
except ImportError:
    print('ERR: python-pptx missing — install via pip or uvx', file=sys.stderr)
    sys.exit(2)
prs = Presentation(sys.argv[1])
for idx, s in enumerate(prs.slides, start=1):
    for sh in s.shapes:
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                t = para.text
                if t.strip():
                    print(f'[slide {idx}] {t}')
"""


def extract_slides(pptx_path: pathlib.Path) -> str:
    """Return one-line-per-paragraph text with slide prefixes. Raises
    ``RuntimeError`` when python-pptx is unavailable (even via uvx)."""
    for cmd in (
        ["python3", "-c", EXTRACT_SCRIPT, str(pptx_path)],
        ["uvx", "--with", "python-pptx", "python3", "-c", EXTRACT_SCRIPT, str(pptx_path)],
    ):
        try:
            r = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        except FileNotFoundError:
            continue
        if r.returncode == 0:
            return r.stdout
    raise RuntimeError(
        f"python-pptx unavailable and uvx fallback failed for {pptx_path.name}"
    )


def audit(pptx_text: str, provenance_text: str) -> tuple[int, list[str]]:
    """Return ``(missing_count, messages)``.

    Reuses ``normalize_variants`` from provenance_verify so the
    unit-equivalence behavior (亿元 ↔ 万元 etc.) matches the
    analysis-md gate — a number that passes provenance_verify on
    analysis.md will also pass here when it appears on a slide.
    """
    corpus = extract_provenance_corpus(provenance_text)
    hits: list[tuple[str, str, str, str]] = []
    for raw in pptx_text.splitlines():
        slide_tag = ""
        body = raw
        if raw.startswith("[slide "):
            end = raw.find("]")
            if end > 0:
                slide_tag = raw[1:end]
                body = raw[end + 1:].strip()
        for m in HARD_NUMBER.finditer(body):
            hits.append((slide_tag, m.group(1), m.group(2), body.strip()[:120]))

    missing: list[str] = []
    covered = 0
    for slide_tag, num, unit, snippet in hits:
        unit_variants, bare_variants = normalize_variants(num, unit)
        if any(v in corpus for v in unit_variants):
            covered += 1
            continue
        if unit in corpus and any(v in corpus for v in bare_variants):
            covered += 1
            continue
        missing.append(
            f"{slide_tag or 'slide ?'}: hard number '{num}{unit}' missing from data-provenance.md"
            f"\n       context: {snippet!r}"
        )

    summary = [
        f"scan: {len(hits)} hard numbers in pptx slides; "
        f"{covered} covered, {len(missing)} missing."
    ]
    return len(missing), summary + missing


def main(argv: list[str]) -> int:
    if len(argv) != 3:
        print(
            "usage: slide_data_audit.py <deck.pptx> <data-provenance.md>",
            file=sys.stderr,
        )
        return 2

    pptx_p = pathlib.Path(argv[1])
    prov_p = pathlib.Path(argv[2])
    for p in (pptx_p, prov_p):
        if not p.exists():
            print(f"file not found: {p}", file=sys.stderr)
            return 2

    try:
        slides_text = extract_slides(pptx_p)
    except RuntimeError as e:
        print(f"ERR: {e}", file=sys.stderr)
        return 2

    prov_text = prov_p.read_text(encoding="utf-8", errors="replace")
    missing, messages = audit(slides_text, prov_text)
    header = messages[0]
    details = messages[1:]

    if missing:
        print(f"FAIL: {header}", file=sys.stderr)
        for m in details:
            print(m, file=sys.stderr)
        return 1
    print(f"OK: {header}")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv))
