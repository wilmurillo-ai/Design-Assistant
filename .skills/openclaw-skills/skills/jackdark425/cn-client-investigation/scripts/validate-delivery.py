#!/usr/bin/env python3
"""
validate-delivery.py — run all three CN banker delivery gates on a deliverable
directory in one command.

Gates (each must exit 0 for overall PASS):

1. verify_intelligence.py  — upstream `*intelligence*.md` must have inline
                             source citation for every hard number.
                             Script lives in the paired `aigroup-lead-discovery-
                             openclaw` plugin, under cn-lead-safety/scripts/.
                             If that plugin isn't installed, skip this gate
                             with a warning (not a failure).

2. cn_typo_scan.py         — every `.pptx` in the dir extracted to text and
                             scanned for MiniMax escape-drift typos (宽厭谛79
                             family). Uses python-pptx (or uvx --with
                             python-pptx fallback) for extraction.

3. provenance_verify.py    — `analysis.md` ↔ `data-provenance.md` must have
                             every hard number in the analysis covered by a
                             row in the provenance table, including cross-
                             unit equivalence (亿元 ↔ 万元).

Usage:
    python3 validate-delivery.py <deliverable_dir>
    # exit 0 → all applicable gates passed
    # exit 1 → at least one gate failed

Design notes:
- A gate is "applicable" only if the input file it expects exists in the dir.
  e.g. no `*intelligence*.md` → skip gate 1 (with a note, not a failure).
- The output is bank-reviewer-readable: one summary line per gate plus any
  failure details copied from the underlying gate's stderr.
- This script is a thin orchestrator; it does NOT reimplement the gate logic
  — it just invokes the canonical per-gate scripts and aggregates results.
"""
from __future__ import annotations
import argparse
import glob
import os
import pathlib
import subprocess
import sys

# ---- Paths ------------------------------------------------------------------

# This script is in
#   skills/cn-client-investigation/scripts/validate-delivery.py
# so the other two gates live at:
SCRIPT_DIR = pathlib.Path(__file__).resolve().parent
CN_TYPO_SCAN = SCRIPT_DIR / "cn_typo_scan.py"
PROVENANCE_VERIFY = SCRIPT_DIR / "provenance_verify.py"
STYLE_SCAN = SCRIPT_DIR / "style_scan.py"
SLIDE_DATA_AUDIT = SCRIPT_DIR / "slide_data_audit.py"
SOURCE_AUTH_CHECK = SCRIPT_DIR / "source_authenticity_check.py"
RAW_DATA_CHECK = SCRIPT_DIR / "raw_data_check.py"
DOCX_DATA_AUDIT = SCRIPT_DIR / "docx_data_audit.py"

# verify_intelligence.py is in the paired lead-discovery plugin. Try a few
# known locations in order. If none resolves, that gate is reported as
# SKIPPED rather than failing — the deliverable may not include upstream
# intelligence markdown, and forcing this dependency would be wrong.
VERIFY_INTELLIGENCE_CANDIDATES = [
    # macmini ClawHub install
    pathlib.Path.home()
    / ".openclaw/extensions/aigroup-lead-discovery-openclaw/skills/cn-lead-safety/scripts/verify_intelligence.py",
    # local dev clone
    pathlib.Path.home()
    / "claude/ai/aigroup-lead-discovery-openclaw/skills/cn-lead-safety/scripts/verify_intelligence.py",
    # sibling dir (both repos cloned side by side)
    SCRIPT_DIR.parents[3]
    / "aigroup-lead-discovery-openclaw/skills/cn-lead-safety/scripts/verify_intelligence.py",
]


def find_verify_intelligence() -> pathlib.Path | None:
    for p in VERIFY_INTELLIGENCE_CANDIDATES:
        if p.exists():
            return p
    return None


# ---- PPTX text extraction (inline, matches compile template) -----------------

EXTRACT_SCRIPT = """\
import sys
try:
    from pptx import Presentation
except ImportError:
    print('ERR: python-pptx missing — install via pip or uvx', file=sys.stderr)
    sys.exit(2)
prs = Presentation(sys.argv[1])
for s in prs.slides:
    for sh in s.shapes:
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                t = para.text
                if t.strip():
                    print(t)
"""


def extract_pptx_to_txt(pptx_path: pathlib.Path, out_path: pathlib.Path) -> bool:
    """Run python-pptx extraction; if system python doesn't have pptx, try uvx."""
    for cmd in (
        ["python3", "-c", EXTRACT_SCRIPT, str(pptx_path)],
        ["uvx", "--with", "python-pptx", "python3", "-c", EXTRACT_SCRIPT, str(pptx_path)],
    ):
        try:
            r = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        except FileNotFoundError:
            continue
        except subprocess.TimeoutExpired:
            print(f"  ERR: extraction timeout on {pptx_path.name}", file=sys.stderr)
            return False
        if r.returncode == 0:
            out_path.write_text(r.stdout, encoding="utf-8")
            return True
    print(
        f"  ERR: python-pptx unavailable and uvx fallback failed for {pptx_path.name}",
        file=sys.stderr,
    )
    return False


# ---- Gate runners ------------------------------------------------------------


def run_gate(script: pathlib.Path, args: list[str]) -> tuple[int, str, str]:
    r = subprocess.run(
        ["python3", str(script), *args], capture_output=True, text=True, timeout=60
    )
    return r.returncode, r.stdout, r.stderr


def main() -> int:
    ap = argparse.ArgumentParser(
        description="Run all three CN banker delivery gates on a deliverable dir."
    )
    ap.add_argument("deliverable", help="Path to deliverable directory")
    ap.add_argument(
        "--strict",
        action="store_true",
        help="Pass --strict to provenance_verify (catches estimate-as-T1 smuggling).",
    )
    ap.add_argument(
        "--strict-mcp",
        action="store_true",
        help=("Pass --strict-mcp to source_authenticity_check (fails on media-only "
              "citations in addition to forbidden labels like Wind/同花顺)."),
    )
    ap.add_argument(
        "--style",
        action="store_true",
        help="Also run style_scan.py on analysis.md in warn-only mode (non-blocking).",
    )
    args = ap.parse_args()

    d = pathlib.Path(args.deliverable).resolve()
    if not d.is_dir():
        print(f"ERR: {d} is not a directory", file=sys.stderr)
        return 2

    print(f"validate-delivery on: {d}")
    print()

    overall_fail = 0
    summary: list[str] = []

    # -------- Gate 1: verify_intelligence on *intelligence*.md --------
    intel_files = sorted(glob.glob(str(d / "*intelligence*.md")))
    verify_script = find_verify_intelligence()
    if not intel_files:
        summary.append("[1/4] verify_intelligence  SKIPPED (no *intelligence*.md in dir)")
    elif not verify_script:
        summary.append(
            "[1/4] verify_intelligence  SKIPPED (aigroup-lead-discovery-openclaw not installed; "
            "paired plugin provides the script)"
        )
    else:
        all_ok = True
        for f in intel_files:
            rc, so, se = run_gate(verify_script, [f])
            if rc != 0:
                all_ok = False
                print(f"  [fail] verify_intelligence on {pathlib.Path(f).name}", file=sys.stderr)
                if se.strip():
                    print("    " + se.strip().replace("\n", "\n    "), file=sys.stderr)
            else:
                print(f"  [ok] verify_intelligence on {pathlib.Path(f).name}: {so.strip()}")
        summary.append(
            f"[1/4] verify_intelligence  {'OK' if all_ok else 'FAIL'} "
            f"({len(intel_files)} intelligence file(s))"
        )
        if not all_ok:
            overall_fail = 1

    # -------- Gate 2: cn_typo_scan on each *.pptx --------
    pptx_files = sorted(glob.glob(str(d / "*.pptx")))
    if not pptx_files:
        summary.append("[2/4] cn_typo_scan         SKIPPED (no .pptx in dir)")
    else:
        all_ok = True
        for p in pptx_files:
            pp = pathlib.Path(p)
            tmp = pathlib.Path("/tmp") / f"validate-delivery-{os.getpid()}-{pp.stem}.txt"
            if not extract_pptx_to_txt(pp, tmp):
                all_ok = False
                continue
            rc, so, se = run_gate(CN_TYPO_SCAN, [str(tmp)])
            if rc != 0:
                all_ok = False
                print(f"  [fail] cn_typo_scan on {pp.name}", file=sys.stderr)
                if se.strip():
                    print("    " + se.strip().replace("\n", "\n    ")[:2000], file=sys.stderr)
            else:
                print(f"  [ok] cn_typo_scan on {pp.name}: {so.strip()}")
            try:
                tmp.unlink()
            except OSError:
                pass
        summary.append(
            f"[2/4] cn_typo_scan         {'OK' if all_ok else 'FAIL'} "
            f"({len(pptx_files)} pptx file(s))"
        )
        if not all_ok:
            overall_fail = 1

    # -------- Gate 2b: slide_data_audit on each *.pptx ↔ data-provenance.md --------
    provenance_early = d / "data-provenance.md"
    if not pptx_files:
        summary.append("[2b]  slide_data_audit    SKIPPED (no .pptx in dir)")
    elif not provenance_early.exists():
        summary.append(
            "[2b]  slide_data_audit    SKIPPED (no data-provenance.md to cross-check against)"
        )
    else:
        all_ok = True
        for p in pptx_files:
            pp = pathlib.Path(p)
            rc, so, se = run_gate(SLIDE_DATA_AUDIT, [str(pp), str(provenance_early)])
            if rc != 0:
                all_ok = False
                print(f"  [fail] slide_data_audit on {pp.name}", file=sys.stderr)
                if se.strip():
                    print("    " + se.strip().replace("\n", "\n    ")[:2000], file=sys.stderr)
            else:
                print(f"  [ok] slide_data_audit on {pp.name}: {so.strip()}")
        summary.append(
            f"[2b]  slide_data_audit    {'OK' if all_ok else 'FAIL'} "
            f"({len(pptx_files)} pptx file(s))"
        )
        if not all_ok:
            overall_fail = 1

    # -------- Gate 3: provenance_verify on analysis.md ↔ data-provenance.md --------
    analysis = d / "analysis.md"
    provenance = d / "data-provenance.md"
    if not analysis.exists() or not provenance.exists():
        missing = []
        if not analysis.exists():
            missing.append("analysis.md")
        if not provenance.exists():
            missing.append("data-provenance.md")
        summary.append(
            f"[3/4] provenance_verify    SKIPPED (missing: {', '.join(missing)})"
        )
    else:
        prov_args = [str(analysis), str(provenance)]
        if args.strict:
            prov_args.insert(0, "--strict")
        rc, so, se = run_gate(PROVENANCE_VERIFY, prov_args)
        if rc != 0:
            overall_fail = 1
            print("  [fail] provenance_verify:", file=sys.stderr)
            if se.strip():
                print("    " + se.strip().replace("\n", "\n    ")[:2000], file=sys.stderr)
            summary.append(
                f"[3/4] provenance_verify    FAIL{' (--strict)' if args.strict else ''}"
            )
        else:
            print(f"  [ok] provenance_verify{' (--strict)' if args.strict else ''}: {so.strip()}")
            if se.strip():
                # strict mode may emit WARN on stderr even when exit 0
                print("    " + se.strip().replace("\n", "\n    "))
            summary.append(
                f"[3/4] provenance_verify    OK{' (--strict)' if args.strict else ''}"
            )

    # -------- Gate 3b: source_authenticity_check on data-provenance.md --------
    if not provenance.exists():
        summary.append("[3b]  source_authenticity SKIPPED (no data-provenance.md)")
    else:
        auth_args = [str(provenance)]
        if args.strict_mcp:
            auth_args.insert(0, "--strict-mcp")
        rc, so, se = run_gate(SOURCE_AUTH_CHECK, auth_args)
        mode_label = " (--strict-mcp)" if args.strict_mcp else ""
        if rc != 0:
            overall_fail = 1
            print(f"  [fail] source_authenticity_check{mode_label}:", file=sys.stderr)
            if se.strip():
                print("    " + se.strip().replace("\n", "\n    ")[:2000], file=sys.stderr)
            summary.append(f"[3b]  source_authenticity FAIL{mode_label}")
        else:
            print(f"  [ok] source_authenticity_check{mode_label}: {so.strip()}")
            if se.strip():
                print("    " + se.strip().replace("\n", "\n    "))
            summary.append(f"[3b]  source_authenticity OK{mode_label}")

    # -------- Gate 2c: docx_data_audit on each *.docx ↔ data-provenance.md --------
    docx_files = sorted(glob.glob(str(d / "*.docx")))
    if not docx_files:
        summary.append("[2c]  docx_data_audit     SKIPPED (no .docx in dir)")
    elif not provenance_early.exists():
        summary.append(
            "[2c]  docx_data_audit     SKIPPED (no data-provenance.md to cross-check against)"
        )
    else:
        all_ok = True
        for p in docx_files:
            pp = pathlib.Path(p)
            rc, so, se = run_gate(DOCX_DATA_AUDIT, [str(pp), str(provenance_early)])
            if rc != 0:
                all_ok = False
                print(f"  [fail] docx_data_audit on {pp.name}", file=sys.stderr)
                if se.strip():
                    print("    " + se.strip().replace("\n", "\n    ")[:2000], file=sys.stderr)
            else:
                print(f"  [ok] docx_data_audit on {pp.name}: {so.strip()}")
        summary.append(
            f"[2c]  docx_data_audit     {'OK' if all_ok else 'FAIL'} "
            f"({len(docx_files)} docx file(s))"
        )
        if not all_ok:
            overall_fail = 1

    # -------- Gate 3c: raw_data_check on deliverable-dir/raw-data/ --------
    rd_args = [str(d)]
    if args.strict_mcp:
        rd_args.insert(0, "--strict-mcp")
    rc, so, se = run_gate(RAW_DATA_CHECK, rd_args)
    rd_label = " (--strict-mcp)" if args.strict_mcp else ""
    if rc != 0:
        overall_fail = 1
        print(f"  [fail] raw_data_check{rd_label}:", file=sys.stderr)
        if se.strip():
            print("    " + se.strip().replace("\n", "\n    ")[:2000], file=sys.stderr)
        summary.append(f"[3c]  raw_data_check      FAIL{rd_label}")
    else:
        print(f"  [ok] raw_data_check{rd_label}: {so.strip()}")
        if se.strip():
            print("    " + se.strip().replace("\n", "\n    "))
        summary.append(f"[3c]  raw_data_check      OK{rd_label}")

    # -------- Optional Phase [4]: style_scan (warn-only) --------
    if args.style and analysis.exists():
        rc, so, se = run_gate(STYLE_SCAN, ["--warn-only", str(analysis)])
        out = (so + se).strip()
        if out:
            print(f"  [style] {out.replace(chr(10), chr(10) + '    ')}")
        summary.append("[4]   style_scan (opt)    WARN-ONLY (non-blocking)")

    # -------- Summary --------
    print()
    print("=" * 60)
    for line in summary:
        print(line)
    print("=" * 60)

    if overall_fail:
        print("OVERALL: FAIL — at least one gate failed. Do NOT ship.", file=sys.stderr)
        return 1
    print("OVERALL: PASS — all applicable gates clean. Shippable.")
    print(
        "Next step: run data-quality-audit skill for independent cross-source "
        "verification (strongly recommended before client delivery)."
    )
    return 0


if __name__ == "__main__":
    sys.exit(main())
