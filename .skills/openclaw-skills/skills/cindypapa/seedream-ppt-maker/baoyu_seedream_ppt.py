#!/usr/bin/env python3
"""
baoyu-seedream-ppt: 宝玉布局 × Seedream 图片PPT生成器

整合 baoyu-infographic 布局风格框架 + Seedream 5.0 文生图 + 自动输出PPT
"""

import argparse
import json
import os
import sys
import time
import shutil
from pathlib import Path
import requests

from pptx import Presentation
from pptx.util import Inches

__version__ = "1.0.0"

# 默认配置
DEFAULT_API_KEY = None
DEFAULT_MODEL = "doubao-seedream-5-0-260128"
BASE_URL = "https://ark.cn-beijing.volces.com/api/v3"

# 从配置文件读取API Key
def get_api_key():
    config_path = os.path.expanduser("~/.openclaw/config.json")
    if os.path.exists(config_path):
        try:
            with open(config_path, 'r', encoding='utf-8') as f:
                config = json.load(f)
                # 支持两种格式
                if 'volces' in config and 'apiKey' in config['volces']:
                    return config['volces']['apiKey']
                if 'volces_api_key' in config:
                    return config['volces_api_key']
                if 'volces' in config and 'api_key' in config['volces']:
                    return config['volces']['api_key']
        except Exception as e:
            print(f"⚠️  读取配置失败: {e}")
    return None

def generate_image(prompt, size, api_key, model_id):
    """调用Seedream API生成图片"""
    headers = {
        'Authorization': f'Bearer {api_key}',
        'Content-Type': 'application/json'
    }
    
    # 解析尺寸
    if size == "landscape":
        resolution = "2560x1440"
    elif size == "portrait":
        resolution = "1440x2560"
    elif size == "square":
        resolution = "1920x1920"
    else:
        # 自定义比例，使用最接近的分辨率
        resolution = "2560x1440"
    
    data = {
        "model": model_id,
        "prompt": prompt,
        "size": resolution
    }
    
    response = requests.post(
        f"{BASE_URL}/images/generations",
        headers=headers,
        json=data,
        timeout=180
    )
    
    if response.status_code != 200:
        print(f"❌ API请求失败: {response.status_code}")
        print(response.text[:200])
        return None
    
    result = response.json()
    if 'data' in result and len(result['data']) > 0:
        image_url = result['data'][0]['url']
        return image_url
    else:
        print(f"❌ 响应格式错误")
        print(json.dumps(result, indent=2, ensure_ascii=False)[:200])
        return None

def download_image(url, save_path):
    """下载生成的图片"""
    try:
        resp = requests.get(url, timeout=60)
        if resp.status_code == 200:
            with open(save_path, 'wb') as f:
                f.write(resp.content)
            file_size = os.path.getsize(save_path) / 1024 / 1024
            print(f"✅ 已保存: {os.path.basename(save_path)} ({file_size:.2f} MB)")
            return True
        else:
            print(f"❌ 下载失败: {resp.status_code}")
            return False
    except Exception as e:
        print(f"❌ 下载异常: {e}")
        return False

def create_ppt(image_dir, output_ppt_path, aspect="landscape"):
    """将所有图片整合为PPT"""
    # 设置PPT尺寸
    if aspect == "landscape":
        width = Inches(13.333)
        height = Inches(7.5)
    elif aspect == "portrait":
        width = Inches(7.5)
        height = Inches(13.333)
    else: # square
        width = Inches(10)
        height = Inches(10)
    
    prs = Presentation()
    prs.slide_width = width
    prs.slide_height = height
    
    # 获取所有图片文件，按名称排序
    image_files = sorted(list(Path(image_dir).glob("*.png")))
    
    for idx, img_path in enumerate(image_files, 1):
        print(f"  添加页面 {idx}/{len(image_files)}: {img_path.name}")
        blank_layout = prs.slide_layouts[6] # 空白布局
        slide = prs.slides.add_slide(blank_layout)
        
        # 图片铺满整页
        left = Inches(0)
        top = Inches(0)
        slide.shapes.add_picture(str(img_path), left, top, width, height)
    
    prs.save(output_ppt_path)
    file_size = os.path.getsize(output_ppt_path) / (1024*1024)
    print(f"\n📊 PPT生成完成:")
    print(f"   文件: {output_ppt_path}")
    print(f"   页数: {len(prs.slides)}")
    print(f"   大小: {file_size:.2f} MB")
    
    return True

def main():
    parser = argparse.ArgumentParser(
        description='baoyu-seedream-ppt: 宝玉布局 × Seedream 图片PPT生成器'
    )
    parser.add_argument('content_file', nargs='?', help='内容文件路径')
    parser.add_argument('--layout', default='bento-grid', help='布局类型 (默认: bento-grid)')
    parser.add_argument('--style', default='craft-handmade', help='风格类型 (默认: craft-handmade)')
    parser.add_argument('--aspect', default='landscape', 
                        choices=['landscape', 'portrait', 'square'], 
                        help='宽高比 (默认: landscape 16:9)')
    parser.add_argument('--api-key', help='火山方舟API Key (从config.json读取默认)')
    parser.add_argument('--model', default=DEFAULT_MODEL, help='Seedream模型ID')
    parser.add_argument('--output-dir', default='./output', help='输出目录')
    args = parser.parse_args()
    
    # 获取API Key
    api_key = args.api_key or get_api_key()
    if not api_key:
        print("❌ 未找到火山方舟API Key，请配置 ~/.openclaw/config.json 中的 volces.api_key")
        sys.exit(1)
    
    # 检查内容文件
    if not args.content_file or not os.path.exists(args.content_file):
        print("❌ 请提供内容文件路径")
        print(f"   用法: {sys.argv[0]} content.md [options]")
        sys.exit(1)
    
    # 读取布局和风格定义
    skill_root = os.path.dirname(os.path.abspath(__file__))
    baoyu_root = os.path.join(os.path.dirname(skill_root), "baoyu-infographic")
    
    layout_path = os.path.join(baoyu_root, "references", "layouts", f"{args.layout}.md")
    style_path = os.path.join(baoyu_root, "references", "styles", f"{args.style}.md")
    base_prompt_path = os.path.join(baoyu_root, "references", "base-prompt.md")
    
    if not os.path.exists(layout_path):
        print(f"❌ 布局不存在: {args.layout}")
        print(f"   请检查 {layout_path}")
        sys.exit(1)
    
    if not os.path.exists(style_path):
        print(f"❌ 风格不存在: {args.style}")
        print(f"   请检查 {style_path}")
        sys.exit(1)
    
    # 读取内容
    with open(args.content_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    with open(layout_path, 'r', encoding='utf-8') as f:
        layout_def = f.read()
    
    with open(style_path, 'r', encoding='utf-8') as f:
        style_def = f.read()
    
    with open(base_prompt_path, 'r', encoding='utf-8') as f:
        base_prompt = f.read()
    
    # 拆分内容为多页（每页一个#标题分隔）
    pages = []
    current_page = []
    for line in content.split('\n'):
        if line.startswith('# ') and current_page:
            pages.append('\n'.join(current_page))
            current_page = []
        current_page.append(line)
    if current_page:
        pages.append('\n'.join(current_page))
    
    print(f"📖 内容读取完成，共 {len(pages)} 页")
    
    # 创建输出目录
    project_name = os.path.splitext(os.path.basename(args.content_file))[0]
    output_dir = os.path.join(args.output_dir, project_name)
    images_dir = os.path.join(output_dir, "images")
    Path(output_dir).mkdir(parents=True, exist_ok=True)
    Path(images_dir).mkdir(parents=True, exist_ok=True)
    
    # 保存完整prompt信息
    prompt_info = {
        "project": project_name,
        "layout": args.layout,
        "style": args.style,
        "aspect": args.aspect,
        "model": args.model,
        "pages": len(pages)
    }
    with open(os.path.join(output_dir, "config.json"), 'w', encoding='utf-8') as f:
        json.dump(prompt_info, f, ensure_ascii=False, indent=2)
    
    # 逐页生成
    print(f"\n🎨 开始逐页生成图片 (模型: {args.model})...\n")
    
    success_count = 0
    for page_idx, page_content in enumerate(pages):
        page_num = page_idx + 1
        print(f"[{page_num}/{len(pages)}] 生成: {page_content.splitlines()[0][:60]}...")
        
        # 构建完整prompt
        full_prompt = base_prompt
        full_prompt = full_prompt.replace("{{LAYOUT}}", args.layout)
        full_prompt = full_prompt.replace("{{STYLE}}", args.style)
        full_prompt = full_prompt.replace("{{ASPECT_RATIO}}", args.aspect)
        full_prompt = full_prompt.replace("{{LANGUAGE}}", "Chinese")
        full_prompt = full_prompt.replace("{{LAYOUT_GUIDELINES}}", layout_def)
        full_prompt = full_prompt.replace("{{STYLE_GUIDELINES}}", style_def)
        full_prompt = full_prompt.replace("{{CONTENT}}", page_content)
        full_prompt = full_prompt.replace("{{TEXT_LABELS}}", "")
        
        # 保存prompt
        prompt_file = os.path.join(output_dir, f"prompts_page_{page_num:02d}.md")
        with open(prompt_file, 'w', encoding='utf-8') as f:
            f.write(full_prompt)
        
        # 构建完整prompt已经完成，full_prompt就是可以直接发给API的
        prompt_for_api = full_prompt
        
        # 生成图片
        image_url = generate_image(full_prompt, args.aspect, api_key, args.model)
        if not image_url:
            print(f"   ❌ 生成失败，跳过")
            continue
        
        # 下载图片
        image_path = os.path.join(images_dir, f"page_{page_num:02d}.png")
        if download_image(image_url, image_path):
            success_count += 1
        
        # 间隔避免限流
        if page_idx != len(pages) - 1:
            print(f"   ⏳ 等待10秒后生成下一页...")
            time.sleep(10)
    
    print(f"\n✅ 图片生成完成: {success_count}/{len(pages)} 页成功")
    
    # 整合PPT
    if success_count > 0:
        output_ppt = os.path.join(output_dir, f"{project_name}_{args.layout}_{args.style}.pptx")
        print(f"\n📦 整合图片到PPT...")
        create_ppt(images_dir, output_ppt, args.aspect)
        print(f"\n🎉 全部完成! 最终PPT: {output_ppt}")
    else:
        print(f"\n❌ 没有成功生成任何图片，请检查API配置")
        sys.exit(1)

if __name__ == "__main__":
    main()
