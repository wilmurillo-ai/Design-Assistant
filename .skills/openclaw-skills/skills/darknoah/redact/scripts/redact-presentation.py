#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Presentation Redaction Tool

Usage:
  redact-presentation.py <presentation.pptx> <rules.csv> <output.pptx>   # Redact presentation with privacy rules
  redact-presentation.py --help                                          # Show help

Rules CSV format:
  target_text,replacement_text    # Replace with new text
  target_text,                    # Empty replacement = cover with █ characters (or solid color for images)

Features:
  - Replace text in slides, tables, text boxes, shapes
  - Redact text in embedded images using PPStructureV3 OCR
  - Preserve presentation formatting
"""

import argparse
import csv
import io
import os
import re
import sys
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any, List, Optional, Sequence, Tuple

os.environ.setdefault("DISABLE_MODEL_SOURCE_CHECK", "True")
os.environ.setdefault("FLAGS_use_mkldnn", "0")

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Character used for redaction when replacement is empty
REDACTION_CHAR = "█"

# Default cover color for images
DEFAULT_COVER_COLOR = (50, 50, 50)


@dataclass
class TextRegion:
    text: str
    polygon: List[Tuple[float, float]]
    region_type: str = "ocr"
    confidence: float = 1.0


@dataclass
class ReplacementRule:
    target: str
    replacement: Optional[str]


@dataclass
class TextMatch:
    region: TextRegion
    target: str
    start_idx: int
    end_idx: int
    rule: ReplacementRule


def eprint(*args: Any) -> None:
    """Print to stderr."""
    print(*args, file=sys.stderr)


def load_rules(csv_path: Path) -> List[ReplacementRule]:
    """Load replacement rules from CSV file."""
    if not csv_path.exists():
        raise FileNotFoundError(f"CSV file not found: {csv_path}")
    
    rules: List[ReplacementRule] = []
    with csv_path.open("r", encoding="utf-8-sig", newline="") as f:
        reader = csv.reader(f)
        for row in reader:
            if not row:
                continue
            target = row[0].strip() if row[0] else ""
            if not target:
                continue
            
            if len(row) >= 2 and row[1]:
                replacement = row[1].strip()
            else:
                replacement = None
            
            rules.append(ReplacementRule(target=target, replacement=replacement))
    
    if not rules:
        raise ValueError(f"No valid replacement rules found in: {csv_path}")
    
    return rules


def replace_text_in_string(text: str, rules: Sequence[ReplacementRule]) -> Tuple[str, int]:
    """
    Replace text in a string according to rules.
    Returns (new_text, replacement_count).
    
    When replacement is None/empty, uses '█' characters to cover the target text.
    """
    count = 0
    new_text = text
    for rule in rules:
        occurrences = new_text.count(rule.target)
        if occurrences > 0:
            if rule.replacement is not None:
                # Replace with specified text
                new_text = new_text.replace(rule.target, rule.replacement)
                eprint(f"  Replaced '{rule.target}' -> '{rule.replacement}' ({occurrences} times)")
            else:
                # Replace with █ characters (same length as target)
                redacted = REDACTION_CHAR * len(rule.target)
                new_text = new_text.replace(rule.target, redacted)
                eprint(f"  Redacted '{rule.target}' -> '{redacted}' ({occurrences} times)")
            count += occurrences
    return new_text, count


def redact_text_frame(text_frame, rules: Sequence[ReplacementRule]) -> int:
    """Redact text in a text frame. Returns replacement count."""
    total_count = 0
    
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text:
                new_text, count = replace_text_in_string(run.text, rules)
                if count > 0:
                    run.text = new_text
                    total_count += count
    
    return total_count


def redact_table(table, rules: Sequence[ReplacementRule]) -> int:
    """Redact text in a table. Returns replacement count."""
    total_count = 0
    
    for row in table.rows:
        for cell in row.cells:
            if cell.text_frame:
                total_count += redact_text_frame(cell.text_frame, rules)
    
    return total_count


def redact_shape(shape, rules: Sequence[ReplacementRule]) -> int:
    """Redact text in a shape. Returns replacement count."""
    total_count = 0
    
    if shape.has_text_frame:
        total_count += redact_text_frame(shape.text_frame, rules)
    
    if shape.has_table:
        total_count += redact_table(shape.table, rules)
    
    # Handle group shapes recursively
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            total_count += redact_shape(sub_shape, rules)
    
    return total_count


# ============== Image Redaction Functions ==============

def polygon_bounds(polygon: Sequence[Tuple[float, float]]) -> Tuple[int, int, int, int]:
    """Get bounding box of polygon."""
    xs = [p[0] for p in polygon]
    ys = [p[1] for p in polygon]
    return int(min(xs)), int(min(ys)), int(max(xs)), int(max(ys))


def get_text_width(text: str, font: ImageFont.ImageFont, draw: ImageDraw.ImageDraw) -> int:
    """Get text width in pixels."""
    bbox = draw.textbbox((0, 0), text, font=font)
    return bbox[2] - bbox[0]


def estimate_font_size(region: TextRegion) -> int:
    """Estimate font size based on region height."""
    _, y1, _, y2 = polygon_bounds(region.polygon)
    height = y2 - y1
    return max(10, int(height * 0.8))


def resolve_font(font_size: int) -> ImageFont.ImageFont:
    """Find a suitable font for rendering across platforms."""
    candidates = [
        "/System/Library/Fonts/PingFang.ttc",
        "/System/Library/Fonts/STHeiti Light.ttc",
        "/System/Library/Fonts/Hiragino Sans GB.ttc",
        "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
        r"C:\Windows\Fonts\msyh.ttc",
        r"C:\Windows\Fonts\simhei.ttf",
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    for path in candidates:
        if os.path.exists(path):
            try:
                return ImageFont.truetype(path, max(10, font_size))
            except Exception:
                continue
    return ImageFont.load_default()


def sample_background_color(image: Image.Image, polygon: Sequence[Tuple[float, float]]) -> Tuple[int, int, int]:
    """Sample background color from the region corners and edges."""
    x1, y1, x2, y2 = polygon_bounds(polygon)
    width, height = x2 - x1, y2 - y1
    
    sample_points = []
    margin = max(2, min(width, height) // 10)
    
    for dx, dy in [(margin, margin), (width - margin, margin),
                   (margin, height - margin), (width - margin, height - margin)]:
        sample_points.append((x1 + dx, y1 + dy))
    
    if width > 20:
        sample_points.append((x1 + width // 2, y1 + margin))
        sample_points.append((x1 + width // 2, y2 - margin))
    if height > 20:
        sample_points.append((x1 + margin, y1 + height // 2))
        sample_points.append((x2 - margin, y1 + height // 2))
    
    pixels = image.load()
    colors = []
    img_width, img_height = image.size
    
    for px, py in sample_points:
        if 0 <= px < img_width and 0 <= py < img_height:
            pixel = pixels[px, py]
            colors.append(pixel[:3] if isinstance(pixel, (list, tuple)) and len(pixel) >= 3 else pixel)
    
    if not colors:
        return (255, 255, 255)
    
    r = sorted(c[0] for c in colors)[len(colors) // 2]
    g = sorted(c[1] for c in colors)[len(colors) // 2]
    b = sorted(c[2] for c in colors)[len(colors) // 2]
    return (r, g, b)


def calculate_sub_region(region: TextRegion, start_idx: int, end_idx: int,
                         font: ImageFont.ImageFont, draw: ImageDraw.ImageDraw) -> List[Tuple[float, float]]:
    """Calculate the polygon for a substring within a text region."""
    x1, y1, x2, y2 = polygon_bounds(region.polygon)
    total_width = x2 - x1
    
    full_text = region.text
    full_width = get_text_width(full_text, font, draw)
    
    prefix_text = full_text[:start_idx]
    target_text = full_text[start_idx:end_idx]
    
    prefix_width = get_text_width(prefix_text, font, draw) if prefix_text else 0
    target_width = get_text_width(target_text, font, draw)
    
    scale = total_width / full_width if full_width > 0 else 1
    
    start_x = x1 + prefix_width * scale
    end_x = x1 + (prefix_width + target_width) * scale
    
    padding = 1
    start_x = max(x1, start_x - padding)
    end_x = min(x2, end_x + padding)
    
    return [
        (start_x, float(y1)),
        (end_x, float(y1)),
        (end_x, float(y2)),
        (start_x, float(y2))
    ]


def cover_sub_region(draw: ImageDraw.ImageDraw, polygon: Sequence[Tuple[float, float]],
                     color: Tuple[int, int, int] = DEFAULT_COVER_COLOR) -> None:
    """Cover a sub-region with solid color block."""
    points = [(int(x), int(y)) for x, y in polygon]
    draw.polygon(points, fill=color)


def replace_text_in_sub_region(image: Image.Image, region: TextRegion,
                               sub_polygon: List[Tuple[float, float]],
                               new_text: str, font: ImageFont.ImageFont) -> Image.Image:
    """Replace text in a sub-region with new text."""
    x1, y1, x2, y2 = polygon_bounds(sub_polygon)
    width, height = x2 - x1, y2 - y1
    
    bg_color = sample_background_color(image, sub_polygon)
    draw = ImageDraw.Draw(image)
    
    points = [(int(x), int(y)) for x, y in sub_polygon]
    draw.polygon(points, fill=bg_color)
    
    test_font = font
    font_size = estimate_font_size(region)
    while font_size >= 10:
        test_font = resolve_font(font_size)
        text_width = get_text_width(new_text, test_font, draw)
        if text_width <= width - 2:
            break
        font_size -= 1
    
    text_width = get_text_width(new_text, test_font, draw)
    bbox = draw.textbbox((0, 0), new_text, font=test_font)
    text_height = bbox[3] - bbox[1]
    
    text_x = x1 + max(0, (width - text_width) // 2)
    text_y = y1 + max(0, (height - text_height) // 2)
    
    brightness = (bg_color[0] * 299 + bg_color[1] * 587 + bg_color[2] * 114) / 1000
    text_color = (0, 0, 0) if brightness > 128 else (255, 255, 255)
    
    draw.text((text_x, text_y), new_text, fill=text_color, font=test_font)
    
    return image


def extract_regions_from_json(data: Any) -> List[TextRegion]:
    """Extract text regions from PPStructureV3 JSON result."""
    regions = []
    seen = set()
    
    TEXT_KEYS = ("rec_text", "text", "ocr_text", "transcription", "label", "content", "markdown", "block_content")
    BOX_KEYS = ("text_region", "text_box", "dt_polys", "poly", "bbox", "box", "points", "position", "region", "block_bbox", "coordinate")
    
    def normalize_polygon(value):
        if value is None:
            return None
        if isinstance(value, dict):
            for key in BOX_KEYS:
                if key in value:
                    result = normalize_polygon(value[key])
                    if result:
                        return result
            return None
        if not isinstance(value, (list, tuple)):
            return None
        if len(value) == 4 and all(isinstance(v, (int, float)) for v in value):
            x1, y1, x2, y2 = value
            return [(float(x1), float(y1)), (float(x2), float(y1)),
                    (float(x2), float(y2)), (float(x1), float(y2))]
        if value and all(
            isinstance(p, (list, tuple)) and len(p) >= 2 and
            all(isinstance(v, (int, float)) for v in p[:2])
            for p in value
        ):
            return [(float(p[0]), float(p[1])) for p in value]
        return None
    
    def extract_text(value):
        if isinstance(value, str):
            return value.strip() or None
        if isinstance(value, dict):
            for key in TEXT_KEYS:
                if key in value:
                    result = extract_text(value[key])
                    if result:
                        return result
            for nested in value.values():
                result = extract_text(nested)
                if result:
                    return result
            return None
        if isinstance(value, list):
            parts = []
            for item in value:
                result = extract_text(item)
                if result:
                    parts.append(result)
            return " ".join(parts) if parts else None
        return None
    
    def add_region(text, polygon, region_type="ocr", confidence=1.0):
        if not text or not polygon:
            return
        key = (text, tuple(polygon))
        if key in seen:
            return
        seen.add(key)
        regions.append(TextRegion(text=text, polygon=polygon, region_type=region_type, confidence=confidence))
    
    def walk(node, default_type="unknown"):
        if isinstance(node, dict):
            if "block_label" in node and "block_content" in node:
                text = node.get("block_content")
                bbox = node.get("block_bbox")
                label = node.get("block_label", default_type)
                if label != "table" or not (text and text.startswith("<")):
                    add_region(extract_text(text), normalize_polygon(bbox), label)
            
            if "label" in node and "coordinate" in node:
                text = node.get("label")
                coord = node.get("coordinate")
                score = node.get("score", 1.0)
                add_region(text, normalize_polygon(coord), text, score)
            
            dt_polys = node.get("dt_polys")
            rec_texts = node.get("rec_texts")
            if dt_polys and rec_texts:
                for poly, text in zip(dt_polys, rec_texts):
                    add_region(extract_text(text), normalize_polygon(poly), "ocr")
            
            if "rec_text" in node or "text" in node:
                text = extract_text(node)
                box = None
                for key in BOX_KEYS:
                    if key in node:
                        box = normalize_polygon(node[key])
                        if box:
                            break
                rtype = node.get("type", node.get("block_label", "ocr"))
                add_region(text, box, rtype)
            
            if "cells" in node:
                for cell in node["cells"]:
                    walk(cell, "table_cell")
            
            for key, value in node.items():
                if key in ("dt_polys", "rec_texts", "boxes"):
                    continue
                if isinstance(value, (dict, list)):
                    walk(value, default_type)
        
        elif isinstance(node, list):
            for item in node:
                walk(item, default_type)
    
    walk(data)
    return regions


def find_matches_in_regions(regions: List[TextRegion], rules: Sequence[ReplacementRule]) -> List[TextMatch]:
    """Find all text matches in regions."""
    matches = []
    
    for region in regions:
        for rule in rules:
            start = 0
            while True:
                idx = region.text.find(rule.target, start)
                if idx == -1:
                    break
                matches.append(TextMatch(
                    region=region,
                    target=rule.target,
                    start_idx=idx,
                    end_idx=idx + len(rule.target),
                    rule=rule
                ))
                start = idx + 1
    
    return matches


def redact_image(image: Image.Image, rules: Sequence[ReplacementRule], pipeline=None) -> Tuple[Image.Image, int]:
    """
    Redact text in an image using OCR.
    Returns (redacted_image, replacement_count).
    """
    import paddle
    from paddleocr import PPStructureV3
    
    # Save image to temp file for processing
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        image.save(tmp.name, "PNG")
        tmp_path = Path(tmp.name)
    
    try:
        # Initialize pipeline if not provided
        if pipeline is None:
            device = "gpu" if paddle.is_compiled_with_cuda() else "cpu"
            eprint(f"    Initializing PPStructureV3 on {device}...")
            pipeline = PPStructureV3(device=device, use_doc_unwarping=False)
        
        # Run OCR
        predictions = pipeline.predict(input=str(tmp_path))
        
        # Collect text regions
        regions = []
        for pred in predictions:
            pred_dir = Path(tempfile.mkdtemp(prefix="ppstructure_"))
            pred.save_to_json(str(pred_dir))
            
            import json
            for json_file in sorted(pred_dir.glob("*_res.json")):
                with json_file.open("r", encoding="utf-8") as f:
                    data = json.load(f)
                    regions.extend(extract_regions_from_json(data))
        
        # Filter to OCR regions with text
        regions = [r for r in regions if r.region_type == "ocr" and r.text]
        
        if not regions:
            return image, 0
        
        # Find matches
        matches = find_matches_in_regions(regions, rules)
        
        if not matches:
            return image, 0
        
        # Process matches
        redacted_image = image.convert("RGB")
        count = 0
        
        for match in matches:
            region = match.region
            rule = match.rule
            
            font_size = estimate_font_size(region)
            font = resolve_font(font_size)
            
            sub_polygon = calculate_sub_region(region, match.start_idx, match.end_idx, font, ImageDraw.Draw(redacted_image))
            
            if rule.replacement is None:
                eprint(f"      Covering: '{match.target}' in '{region.text}'")
                draw = ImageDraw.Draw(redacted_image)
                cover_sub_region(draw, sub_polygon)
            else:
                eprint(f"      Replacing: '{match.target}' -> '{rule.replacement}' in '{region.text}'")
                redacted_image = replace_text_in_sub_region(redacted_image, region, sub_polygon, rule.replacement, font)
            
            count += 1
        
        return redacted_image, count
        
    finally:
        # Cleanup temp file
        if tmp_path.exists():
            tmp_path.unlink()


def redact_presentation(ppt_path: Path, rules: Sequence[ReplacementRule], output_path: Path) -> int:
    """
    Redact a PowerPoint presentation.
    Returns total replacement count.
    """
    eprint(f"Loading presentation: {ppt_path}")
    prs = Presentation(str(ppt_path))
    total_count = 0
    
    # Initialize OCR pipeline once for all images
    pipeline = None
    
    # Process each slide
    for slide_idx, slide in enumerate(prs.slides, 1):
        eprint(f"\nProcessing slide {slide_idx}...")
        
        for shape in slide.shapes:
            # Redact text in shape
            count = redact_shape(shape, rules)
            total_count += count
            
            # Process images in shape
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    image_bytes = image.blob
                    
                    eprint(f"  Found image in slide {slide_idx}")
                    pil_image = Image.open(io.BytesIO(image_bytes))
                    
                    redacted_image, img_count = redact_image(pil_image, rules, pipeline)
                    
                    if img_count > 0:
                        # Save redacted image back
                        output_buffer = io.BytesIO()
                        redacted_image.save(output_buffer, format="PNG")
                        output_buffer.seek(0)
                        
                        # Replace image in shape
                        # Note: python-pptx doesn't support direct image replacement,
                        # we need to use a workaround
                        shape.image.blob = output_buffer.read()
                        total_count += img_count
                        eprint(f"  Redacted {img_count} text regions in image")
                    else:
                        eprint("  No matching text found in image")
                        
                except Exception as e:
                    eprint(f"  Warning: Failed to process image: {e}")
    
    # Save presentation
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    eprint(f"\nSaved redacted presentation: {output_path}")
    
    return total_count


def convert_ppt_to_pptx(ppt_path: Path) -> Optional[Path]:
    """
    Convert .ppt to .pptx format using LibreOffice if available.
    Returns path to converted file or None if conversion failed.
    """
    import subprocess
    
    # Check if LibreOffice is available
    libreoffice_paths = [
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "libreoffice",
        "soffice",
    ]
    
    soffice = None
    for path in libreoffice_paths:
        try:
            result = subprocess.run([path, "--version"], capture_output=True, text=True)
            if result.returncode == 0:
                soffice = path
                break
        except FileNotFoundError:
            continue
    
    if not soffice:
        eprint("Warning: LibreOffice not found. Cannot convert .ppt to .pptx")
        eprint("Please install LibreOffice or convert the file manually.")
        return None
    
    eprint(f"Converting {ppt_path} to .pptx format...")
    
    output_dir = ppt_path.parent
    try:
        result = subprocess.run(
            [soffice, "--headless", "--convert-to", "pptx", "--outdir", str(output_dir), str(ppt_path)],
            capture_output=True,
            text=True,
            timeout=60
        )
        
        if result.returncode == 0:
            pptx_path = output_dir / f"{ppt_path.stem}.pptx"
            if pptx_path.exists():
                eprint(f"Converted to: {pptx_path}")
                return pptx_path
    except Exception as e:
        eprint(f"Conversion failed: {e}")
    
    return None


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Presentation Redaction Tool",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Redact a PowerPoint presentation
  %(prog)s presentation.pptx rules.csv output.pptx

Rules CSV format:
  target_text,new_text      # Replace with new text
  target_text,              # Empty = cover with █ characters (or solid color for images)

Features:
  - Replace text in slides, tables, text boxes, shapes
  - Redact text in embedded images using OCR
  - Support for .ppt format (requires LibreOffice for conversion)
"""
    )
    
    parser.add_argument("presentation", help="Input presentation path (.pptx or .ppt)")
    parser.add_argument("rules_csv", help="Rules CSV file with replacement rules")
    parser.add_argument("output", help="Output presentation path (.pptx)")
    
    args = parser.parse_args()
    
    ppt_path = Path(args.presentation).expanduser().resolve()
    rules_path = Path(args.rules_csv).expanduser().resolve()
    output_path = Path(args.output).expanduser().resolve()
    
    if not ppt_path.exists():
        eprint(f"Presentation not found: {ppt_path}")
        return 1
    
    # Handle .ppt format
    original_ppt_path = None
    if ppt_path.suffix.lower() == ".ppt":
        eprint("Detected .ppt format, attempting conversion...")
        converted_path = convert_ppt_to_pptx(ppt_path)
        if converted_path:
            original_ppt_path = ppt_path
            ppt_path = converted_path
        else:
            eprint("Error: Cannot process .ppt file without conversion.")
            return 1
    elif ppt_path.suffix.lower() != ".pptx":
        eprint(f"Unsupported file format: {ppt_path.suffix}")
        eprint("Only .pptx and .ppt files are supported")
        return 1
    
    # Load rules
    rules = load_rules(rules_path)
    eprint(f"Loaded {len(rules)} replacement rules")
    
    # Process presentation
    total_count = redact_presentation(ppt_path, rules, output_path)
    
    # Cleanup converted file if it was created
    if original_ppt_path and ppt_path.exists() and ppt_path != original_ppt_path:
        try:
            ppt_path.unlink()
            eprint(f"Cleaned up temporary file: {ppt_path}")
        except Exception:
            pass
    
    print(f"\nTotal replacements: {total_count}")
    print(str(output_path))
    
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
