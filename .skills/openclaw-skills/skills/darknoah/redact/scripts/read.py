#!/usr/bin/env python3
# /// script
# requires-python = ">=3.10"
# dependencies = [
#   "pillow",
#   "pymupdf",
#   "python-docx",
#   "python-pptx",
#   "paddleocr",
#   "paddlepaddle",
# ]
# ///
# -*- coding: utf-8 -*-
"""
Document Reader Tool - Extract text from various document formats with OCR support

Usage:
  read.py <file_path>              # Read file and output text content
  read.py <file_path> --info       # Output with page info and structure details

Supported formats:
  - Images: png, jpg, jpeg, bmp, tiff, webp, etc.
  - Documents: docx, doc
  - Presentations: pptx, ppt
  - PDF: pdf

Features:
  - Extract text from documents
  - OCR text from images embedded in documents
  - Page-by-page output with --info flag
"""

import argparse
import io
import json
import os
import platform
import subprocess
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Any, List, Optional, Tuple
import re

os.environ.setdefault("DISABLE_MODEL_SOURCE_CHECK", "True")
os.environ.setdefault("FLAGS_use_mkldnn", "0")

# Use a fixed temp directory to avoid repeated authorization prompts
_TEMP_DIR = Path.home() / ".cache" / "redact_temp"
_TEMP_DIR.mkdir(parents=True, exist_ok=True)

from PIL import Image

# ============== Constants ==============

TEXT_KEYS = ("rec_text", "text", "ocr_text", "transcription", "label", "content", "markdown", "block_content")
BOX_KEYS = ("text_region", "text_box", "dt_polys", "poly", "bbox", "box", "points", "position", "region", "block_bbox", "coordinate")

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".webp", ".gif"}

# ============== Data Classes ==============

@dataclass
class TextRegion:
    """Represents a detected text region from OCR."""
    text: str
    polygon: List[Tuple[float, float]]
    region_type: str = "ocr"
    confidence: float = 1.0


@dataclass
class PageContent:
    """Represents content of a page."""
    page_num: int
    text_content: str = ""
    images: List[Tuple[str, List[TextRegion], List[str]]] = None  # List of (image_id, ocr_regions)
    
    def __post_init__(self):
        if self.images is None:
            self.images = []


# ============== Utility Functions ==============

def eprint(*args: Any) -> None:
    """Print to stderr."""
    print(*args, file=sys.stderr)


def get_pipeline():
    """Initialize PPStructureV3 pipeline."""
    import paddle
    from paddleocr import PPStructureV3
    
    device = "gpu" if paddle.is_compiled_with_cuda() else "cpu"
    eprint(f"Initializing PPStructureV3 on {device}...")
    return PPStructureV3(device=device, use_doc_unwarping=False)


def run_pipeline(pipeline, image_path: Path, remove_images: bool = True) -> Tuple[List[Any], List[str]]:
    """Run PPStructureV3 and export results as JSON."""
    results: List[Any] = []
    markdown_texts: List[str] = []

    eprint(f"Processing: {image_path}")
    predictions = pipeline.predict(input=str(image_path))
    for idx, pred in enumerate(predictions):
        pred_dir = _TEMP_DIR / f"ppstructure_res_{idx}"
        pred_dir.mkdir(parents=True, exist_ok=True)
        pred.save_to_json(str(pred_dir))
        markdown_texts.append(pred.markdown['markdown_texts'] if pred.markdown and 'markdown_texts' in pred.markdown else "")

        for markdown_text in markdown_texts:
            markdown_text = re.sub(
                r'<div\s+style="text-align:\s*center;">\s*<img\s+[^>]*?/>\s*</div>',
                '',
                markdown_text
            )



        for json_file in sorted(pred_dir.glob("*_res.json")):
            with json_file.open("r", encoding="utf-8") as f:
                results.append(json.load(f))
        # Clean up after reading
        for f in pred_dir.iterdir():
            f.unlink()
        pred_dir.rmdir()

    return results, markdown_texts


def normalize_polygon(value: Any) -> Optional[List[Tuple[float, float]]]:
    """Convert various box formats to polygon (list of 4 points)."""
    if value is None:
        return None
    
    if isinstance(value, dict):
        for key in BOX_KEYS:
            if key in value:
                result = normalize_polygon(value[key])
                if result:
                    return result
        return None
    
    if not isinstance(value, (list, tuple)):
        return None
    
    # [x1, y1, x2, y2] format
    if len(value) == 4 and all(isinstance(v, (int, float)) for v in value):
        x1, y1, x2, y2 = value
        return [(float(x1), float(y1)), (float(x2), float(y1)), 
                (float(x2), float(y2)), (float(x1), float(y2))]
    
    # [[x, y], [x, y], ...] format
    if value and all(
        isinstance(p, (list, tuple)) and len(p) >= 2 and 
        all(isinstance(v, (int, float)) for v in p[:2])
        for p in value
    ):
        return [(float(p[0]), float(p[1])) for p in value]
    
    return None


def extract_text_from_value(value: Any) -> Optional[str]:
    """Extract text string from various formats."""
    if isinstance(value, str):
        text = value.strip()
        return text or None
    
    if isinstance(value, dict):
        for key in TEXT_KEYS:
            if key in value:
                result = extract_text_from_value(value[key])
                if result:
                    return result
        for nested in value.values():
            result = extract_text_from_value(nested)
            if result:
                return result
        return None
    
    if isinstance(value, list):
        parts: List[str] = []
        for item in value:
            result = extract_text_from_value(item)
            if result:
                parts.append(result)
        return " ".join(parts) if parts else None
    
    return None


def collect_regions(data: Any) -> List[TextRegion]:
    """Extract all text regions from PPStructureV3 result."""
    regions: List[TextRegion] = []
    seen: set = set()
    
    def add_region(text: Optional[str], polygon: Optional[List[Tuple[float, float]]], 
                   region_type: str = "ocr", confidence: float = 1.0) -> None:
        if not text or not polygon:
            return
        key = (text, tuple(polygon))
        if key in seen:
            return
        seen.add(key)
        regions.append(TextRegion(text=text, polygon=polygon, region_type=region_type, confidence=confidence))
    
    def walk(node: Any, default_type: str = "unknown") -> None:
        if isinstance(node, dict):
            # Handle parsing_res_list (block results)
            if "block_label" in node and "block_content" in node:
                text = node.get("block_content")
                bbox = node.get("block_bbox")
                label = node.get("block_label", default_type)
                # Skip table HTML content
                if label != "table" or not (text and text.startswith("<")):
                    add_region(extract_text_from_value(text), normalize_polygon(bbox), label)
            
            # Handle layout_det_res boxes
            if "label" in node and "coordinate" in node:
                text = node.get("label")
                coord = node.get("coordinate")
                score = node.get("score", 1.0)
                add_region(text, normalize_polygon(coord), text, score)
            
            # Handle overall_ocr_res / OCR results
            dt_polys = node.get("dt_polys")
            rec_texts = node.get("rec_texts")
            if dt_polys and rec_texts:
                for poly, text in zip(dt_polys, rec_texts):
                    add_region(extract_text_from_value(text), normalize_polygon(poly), "ocr")
            
            # Handle single OCR result
            if "rec_text" in node or "text" in node:
                text = extract_text_from_value(node)
                box = None
                for key in BOX_KEYS:
                    if key in node:
                        box = normalize_polygon(node[key])
                        if box:
                            break
                rtype = node.get("type", node.get("block_label", "ocr"))
                add_region(text, box, rtype)
            
            # Handle table cells
            if "cells" in node:
                for cell in node["cells"]:
                    walk(cell, "table_cell")
            
            # Recurse into nested structures
            for key, value in node.items():
                if key in ("dt_polys", "rec_texts", "boxes"):
                    continue
                if isinstance(value, (dict, list)):
                    walk(value, default_type)
        
        elif isinstance(node, list):
            for item in node:
                walk(item, default_type)
    
    walk(data)
    return regions


def ocr_image(pipeline, image: Image.Image) -> Tuple[List[TextRegion], List[str]]:
    """Run OCR on a PIL Image and return text regions."""
    # Save image to temp file in fixed directory
    tmp_path = _TEMP_DIR / "ocr_temp_image.png"
    image.save(tmp_path, "PNG")

    try:
        results, markdown_texts = run_pipeline(pipeline, tmp_path)
        regions = collect_regions(results)
        # Return all regions with text (no filtering by region_type)
        return [r for r in regions if r.text], markdown_texts
    finally:
        if tmp_path.exists():
            tmp_path.unlink()


def ocr_image_file(pipeline, image_path: Path) -> Tuple[List[TextRegion], List[str]]:
    """Run OCR on an image file and return text regions."""
    results, markdown_texts = run_pipeline(pipeline, image_path)
    regions = collect_regions(results)
    # Return all regions with text (no filtering by region_type)
    return [r for r in regions if r.text], markdown_texts


def sort_regions(regions: List[TextRegion]) -> List[TextRegion]:
    """Sort OCR regions top-to-bottom then left-to-right."""
    def sort_key(r: TextRegion):
        if r.polygon:
            xs = [p[0] for p in r.polygon]
            ys = [p[1] for p in r.polygon]
            return (min(ys), min(xs))
        return (0, 0)

    return sorted(regions, key=sort_key)


def page_content_to_dict(page: PageContent) -> dict:
    """Convert page content to requested JSON structure."""
    content = []

    if page.text_content:
        content.append({
            "type": "text",
            "text": page.text_content,
        })

    for _img_id, regions, markdown_texts in page.images:
        sorted_regions = sort_regions(regions)
        image_text = "\n".join(r.text for r in sorted_regions)
        content.append({
            "type": "image",
            "markdown": "\n".join(markdown_texts),
        })

    return {
        "page_index": page.page_num,
        "content": content,
    }


def image_content_to_dict(regions: List[TextRegion]) -> dict:
    """Convert image OCR result to requested JSON structure."""
    sorted_regions = sort_regions(regions)
    return {
        "type": "image",
        "markdown": "\n".join(r.text for r in sorted_regions),
    }


# ============== Image Reader ==============

def read_image_regions(file_path: Path) -> List[TextRegion]:
    """Read OCR regions from an image file."""
    pipeline = get_pipeline()
    return ocr_image_file(pipeline, file_path)


def read_image(file_path: Path, show_info: bool = False) -> str:
    """Read text from an image file using OCR."""
    regions = read_image_regions(file_path)
    if not regions:
        return ""
    sorted_regions = sort_regions(regions)
    return "\n".join(r.text for r in sorted_regions)


# ============== PDF Reader ==============

def read_pdf_pages(file_path: Path) -> List[PageContent]:
    """Read page content from a PDF file."""
    import fitz  # PyMuPDF
    
    doc = fitz.open(str(file_path))
    pipeline = None  # Lazy load
    pages_content: List[PageContent] = []
    
    for page_num in range(len(doc)):
        page = doc[page_num]
        eprint(f"Processing page {page_num + 1}/{len(doc)}...")
        
        page_content = PageContent(page_num=page_num + 1)
        
        # Extract text from page
        text = page.get_text().strip()
        
        # Treat a page as image-based only when no native text can be extracted.
        # This avoids misclassifying Office-exported PDF pages with short text.
        is_image_page = not text
        
        if is_image_page:
            eprint(f"  Page {page_num + 1} detected as image-based, running OCR...")
            if pipeline is None:
                pipeline = get_pipeline()
            
            # Render page to image
            mat = fitz.Matrix(150 / 72, 150 / 72)
            pix = page.get_pixmap(matrix=mat)
            img_data = pix.tobytes("png")
            image = Image.open(io.BytesIO(img_data)).convert("RGB")
            
            # OCR the page
            regions, markdown_texts = ocr_image(pipeline, image)
            
            # Image-based page: store OCR result as image content
            page_content.images.append(("page_image", regions, markdown_texts))
            page_content.text_content = ""
        else:
            # Text-based page
            page_content.text_content = text
            
            # Also check for images in the page
            images = page.get_images()
            if images:
                if pipeline is None:
                    pipeline = get_pipeline()
                
                for img_idx, img_info in enumerate(images):
                    try:
                        xref = img_info[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        image = Image.open(io.BytesIO(image_bytes)).convert("RGB")
                        
                        eprint(f"  OCR image {img_idx + 1} on page {page_num + 1}...")
                        regions, markdown_texts = ocr_image(pipeline, image)
                        
                        if regions:
                            page_content.images.append((f"image_{img_idx + 1}", regions, markdown_texts))
                    except Exception as e:
                        eprint(f"  Warning: Failed to process image: {e}")
        
        pages_content.append(page_content)
    
    doc.close()
    return pages_content


def read_pdf(file_path: Path, show_info: bool = False) -> str:
    """Read text from a PDF file."""
    pages_content = read_pdf_pages(file_path)
    return format_output(pages_content, show_info)


# ============== Office Conversion ==============

def find_soffice() -> Optional[str]:
    """Find LibreOffice/soffice executable."""
    libreoffice_paths = [
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "libreoffice",
        "soffice",
    ]

    for path in libreoffice_paths:
        try:
            result = subprocess.run([path, "--version"], capture_output=True, text=True)
            if result.returncode == 0:
                return path
        except FileNotFoundError:
            continue
    return None


def find_word_automation() -> bool:
    """Check whether Word.Application automation is available."""
    system = platform.system()

    if system == "Windows":
        try:
            result = subprocess.run(
                [
                    "powershell",
                    "-NoProfile",
                    "-Command",
                    "$word = New-Object -ComObject Word.Application; $word.Quit(); Write-Output 'ok'",
                ],
                capture_output=True,
                text=True,
                timeout=30,
            )
            return result.returncode == 0 and "ok" in result.stdout
        except Exception:
            return False

    if system == "Darwin":
        script = 'tell application id "com.microsoft.Word" to get name'
        try:
            result = subprocess.run(
                ["osascript", "-e", script],
                capture_output=True,
                text=True,
                timeout=30,
            )
            return result.returncode == 0
        except Exception:
            return False

    return False


def convert_with_word(file_path: Path, pdf_path: Path) -> Optional[Path]:
    """Convert Office document to PDF using Word automation."""
    system = platform.system()

    try:
        if system == "Windows":
            source = str(file_path)
            target = str(pdf_path)
            command = (
                "$word = $null; "
                "$doc = $null; "
                "try { "
                "  $word = New-Object -ComObject Word.Application; "
                "  $word.Visible = $false; "
                "  $word.DisplayAlerts = 0; "
                "  $doc = $word.Documents.Open('" + source.replace("'", "''") + "', $false, $true); "
                "  $doc.ExportAsFixedFormat('" + target.replace("'", "''") + "', 17); "
                "} finally { "
                "  if ($doc -ne $null) { $doc.Close(0) } "
                "  if ($word -ne $null) { $word.Quit() } "
                "}"
            )
            result = subprocess.run(
                ["powershell", "-NoProfile", "-Command", command],
                capture_output=True,
                text=True,
                timeout=180,
            )
            if result.returncode == 0 and pdf_path.exists():
                return pdf_path
            if result.stderr.strip():
                eprint(result.stderr.strip())
            return None

        if system == "Darwin":
            source = file_path.as_posix().replace('"', '\\"')
            target = pdf_path.as_posix().replace('"', '\\"')
            script = f'''
set sourceFile to POSIX file "{source}"
set targetFile to POSIX file "{target}"
tell application id "com.microsoft.Word"
    launch
    try
        set visible to false
    end try
    try
        open sourceFile
        set docRef to active document
        save as docRef file name targetFile file format format PDF
        close docRef saving no
    on error errMsg number errNum
        try
            close active document saving no
        end try
        try
            quit saving no
        end try
        error errMsg number errNum
    end try
    try
        quit saving no
    end try
end tell
'''
            result = subprocess.run(
                ["osascript", "-e", script],
                capture_output=True,
                text=True,
                timeout=180,
            )
            if result.returncode == 0 and pdf_path.exists():
                return pdf_path
            if result.stderr.strip():
                eprint(result.stderr.strip())
            return None
    except Exception as e:
        eprint(f"Word conversion failed: {e}")

    return None


def convert_with_libreoffice(file_path: Path, pdf_path: Path) -> Optional[Path]:
    """Convert Office document to PDF using LibreOffice."""
    soffice = find_soffice()
    if not soffice:
        return None

    try:
        result = subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(pdf_path.parent), str(file_path)],
            capture_output=True,
            text=True,
            timeout=120,
        )

        if result.returncode == 0 and pdf_path.exists():
            return pdf_path

        if result.stderr.strip():
            eprint(result.stderr.strip())
    except Exception as e:
        eprint(f"LibreOffice conversion failed: {e}")

    return None


def convert_office_to_pdf(file_path: Path) -> Optional[Path]:
    """Convert Office document to PDF using platform-preferred tool."""
    system = platform.system()
    # Use fixed temp directory instead of creating new one
    pdf_path = _TEMP_DIR / f"{file_path.stem}.pdf"

    word_available = find_word_automation()
    soffice_available = find_soffice() is not None

    if system in ("Darwin", "Windows"):
        if word_available:
            eprint(f"Converting {file_path.name} to PDF with Word.Application...")
            result = convert_with_word(file_path, pdf_path)
            if result:
                eprint(f"Converted to: {result}")
                return result
            if soffice_available:
                eprint("Word.Application conversion failed, falling back to LibreOffice...")
                result = convert_with_libreoffice(file_path, pdf_path)
                if result:
                    eprint(f"Converted to: {result}")
                    return result
        elif soffice_available:
            eprint("Word.Application not available, falling back to LibreOffice...")
            result = convert_with_libreoffice(file_path, pdf_path)
            if result:
                eprint(f"Converted to: {result}")
                return result
        else:
            eprint("Warning: Preferred tool is Word.Application on macOS/Windows.")
            eprint("No Word.Application or LibreOffice installation was detected.")
            return None

    if system == "Linux":
        if soffice_available:
            eprint(f"Converting {file_path.name} to PDF with LibreOffice...")
            result = convert_with_libreoffice(file_path, pdf_path)
            if result:
                eprint(f"Converted to: {result}")
                return result
            if word_available:
                eprint("LibreOffice conversion failed, falling back to Word.Application...")
                result = convert_with_word(file_path, pdf_path)
                if result:
                    eprint(f"Converted to: {result}")
                    return result
        elif word_available:
            eprint("LibreOffice not available, falling back to Word.Application...")
            result = convert_with_word(file_path, pdf_path)
            if result:
                eprint(f"Converted to: {result}")
                return result
        else:
            eprint("Warning: Preferred tool is LibreOffice on Linux.")
            eprint("No LibreOffice or Word.Application installation was detected.")
            return None

    eprint("Warning: No supported Office-to-PDF converter is available.")
    return None


# ============== Word Document Reader ==============


def read_docx_pages(file_path: Path) -> List[PageContent]:
    """Read page content from a Word document (.docx) via PDF conversion for real pagination."""
    pdf_path = convert_office_to_pdf(file_path)
    if not pdf_path:
        return []

    try:
        return read_pdf_pages(pdf_path)
    finally:
        try:
            pdf_path.unlink()
        except Exception:
            pass


def read_docx(file_path: Path, show_info: bool = False) -> str:
    """Read text from a Word document (.docx)."""
    return format_output(read_docx_pages(file_path), show_info)


def read_doc_pages(file_path: Path) -> List[PageContent]:
    """Read page content from a Word document (.doc) via PDF conversion for real pagination."""
    pdf_path = convert_office_to_pdf(file_path)
    if not pdf_path:
        return []

    try:
        return read_pdf_pages(pdf_path)
    finally:
        try:
            pdf_path.unlink()
        except Exception:
            pass


def read_doc(file_path: Path, show_info: bool = False) -> str:
    """Read text from a Word document (.doc)."""
    return format_output(read_doc_pages(file_path), show_info)


# ============== PowerPoint Reader ==============

def read_pptx_pages(file_path: Path) -> List[PageContent]:
    """Read page content from a PowerPoint presentation (.pptx)."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    
    prs = Presentation(str(file_path))
    pipeline = None  # Lazy load
    pages_content: List[PageContent] = []
    
    for slide_idx, slide in enumerate(prs.slides, 1):
        eprint(f"Processing slide {slide_idx}/{len(prs.slides)}...")
        
        page_content = PageContent(page_num=slide_idx)
        slide_text = []
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = []
                    for run in paragraph.runs:
                        if run.text.strip():
                            para_text.append(run.text.strip())
                    if para_text:
                        slide_text.append(" ".join(para_text))
            
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text_frame:
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if run.text.strip():
                                        row_text.append(run.text.strip())
                    if row_text:
                        slide_text.append(" | ".join(row_text))
            
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    image_bytes = image.blob
                    
                    if pipeline is None:
                        pipeline = get_pipeline()
                    
                    eprint(f"  OCR image on slide {slide_idx}...")
                    pil_image = Image.open(io.BytesIO(image_bytes)).convert("RGB")
                    regions, markdown_texts = ocr_image(pipeline, pil_image)
                    
                    if regions:
                        page_content.images.append(("slide_image", regions, markdown_texts))
                except Exception as e:
                    eprint(f"  Warning: Failed to process image: {e}")
        
        page_content.text_content = "\n".join(slide_text)
        pages_content.append(page_content)
    
    return pages_content


def read_pptx(file_path: Path, show_info: bool = False) -> str:
    """Read text from a PowerPoint presentation (.pptx)."""
    return format_output(read_pptx_pages(file_path), show_info)


def read_ppt(file_path: Path, show_info: bool = False) -> str:
    """Read text from a PowerPoint presentation (.ppt) by converting to .pptx."""
    # Try to convert .ppt to .pptx using LibreOffice
    libreoffice_paths = [
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "libreoffice",
        "soffice",
    ]
    
    soffice = None
    for path in libreoffice_paths:
        try:
            result = subprocess.run([path, "--version"], capture_output=True, text=True)
            if result.returncode == 0:
                soffice = path
                break
        except FileNotFoundError:
            continue
    
    if not soffice:
        eprint("Warning: LibreOffice not found. Cannot convert .ppt to .pptx")
        eprint("Please install LibreOffice or convert the file manually.")
        return ""
    
    eprint(f"Converting {file_path} to .pptx format...")
    
    output_dir = file_path.parent
    try:
        result = subprocess.run(
            [soffice, "--headless", "--convert-to", "pptx", "--outdir", str(output_dir), str(file_path)],
            capture_output=True,
            text=True,
            timeout=60
        )
        
        if result.returncode == 0:
            pptx_path = output_dir / f"{file_path.stem}.pptx"
            if pptx_path.exists():
                eprint(f"Converted to: {pptx_path}")
                content = read_pptx(pptx_path, show_info)
                # Clean up converted file
                try:
                    pptx_path.unlink()
                except Exception:
                    pass
                return content
    except Exception as e:
        eprint(f"Conversion failed: {e}")
    
    return ""


# ============== Output Formatting ==============

def format_output(pages_content: List[PageContent], show_info: bool) -> str:
    """Format the output based on whether --info flag is set."""





    if not show_info:
        # Simple output: just concatenate all text
        all_text = []
        for page in pages_content:
            if page.text_content:
                all_text.append(page.text_content)
            for img_id, regions, markdown_texts in page.images:
                if regions:
                    # Sort by position
                    def sort_key(r: TextRegion):
                        if r.polygon:
                            xs = [p[0] for p in r.polygon]
                            ys = [p[1] for p in r.polygon]
                            return (min(ys), min(xs))
                        return (0, 0)
                    regions.sort(key=sort_key)
                    all_text.extend(markdown_texts)
        return "\n".join(all_text)
    
    # Structured output with page info
    output_lines = []
    
    for page in pages_content:
        output_lines.append(f'<page index="{page.page_num}">')
        
        if page.text_content:
            output_lines.append(page.text_content)
        
        for img_id, regions, markdown_texts in page.images:
            if regions:
                output_lines.append("<image format='markdown' type='ocr'>")
                # Sort by position
                def sort_key(r: TextRegion):
                    if r.polygon:
                        xs = [p[0] for p in r.polygon]
                        ys = [p[1] for p in r.polygon]
                        return (min(ys), min(xs))
                    return (0, 0)
                regions.sort(key=sort_key)
                # for region in regions:
                #     output_lines.append(region.text)
                output_lines.append("\n".join(markdown_texts))
                output_lines.append("</image>")
        
        output_lines.append(f'</page>')
    
    return "\n".join(output_lines)


def format_json_output(data: Any) -> str:
    """Format output as JSON string."""
    return json.dumps(data, ensure_ascii=False, indent=2)


def save_output(content: str, save_path: Path) -> None:
    """Save content to file."""
    save_path.parent.mkdir(parents=True, exist_ok=True)
    with save_path.open("w", encoding="utf-8") as f:
        f.write(content)


# ============== Main ==============

def get_file_type(file_path: Path) -> str:
    """Determine file type based on extension."""
    ext = file_path.suffix.lower()
    
    if ext in IMAGE_EXTENSIONS:
        return "image"
    elif ext == ".pdf":
        return "pdf"
    elif ext == ".docx":
        return "docx"
    elif ext == ".doc":
        return "doc"
    elif ext == ".pptx":
        return "pptx"
    elif ext == ".ppt":
        return "ppt"
    else:
        return "unknown"


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Document Reader Tool - Extract text from various document formats with OCR support",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Supported formats:
  - Images: png, jpg, jpeg, bmp, tiff, webp, gif
  - Documents: docx, doc
  - Presentations: pptx, ppt
  - PDF: pdf

Examples:
  # Read text from an image
  %(prog)s image.png

  # Read text from a PDF
  %(prog)s document.pdf

  # Read text from a Word document with page info
  %(prog)s document.docx --info

  # Read text from a PowerPoint presentation
  %(prog)s presentation.pptx --info
"""
    )
    
    parser.add_argument("file_path", help="Input file path (image, pdf, docx, doc, pptx, ppt)")
    parser.add_argument("--info", action="store_true", 
                        help="Output page info and structure details (distinguish OCR text from images)")
    parser.add_argument("--mode", choices=["text", "json"], default="text",
                        help="Output mode: text (default) or json")
    parser.add_argument("--save", help="Optional path to save the output")
    
    args = parser.parse_args()
    
    file_path = Path(args.file_path).expanduser().resolve()
    
    if not file_path.exists():
        eprint(f"File not found: {file_path}")
        return 1
    
    file_type = get_file_type(file_path)
    
    if file_type == "unknown":
        eprint(f"Unsupported file format: {file_path.suffix}")
        eprint(f"Supported formats: image ({', '.join(IMAGE_EXTENSIONS)}), pdf, docx, doc, pptx, ppt")
        return 1
    
    eprint(f"File type: {file_type}")

    # Read file based on type
    if args.mode == "json":
        if file_type == "image":
            data = image_content_to_dict(read_image_regions(file_path))
        elif file_type == "pdf":
            data = {"type": "pdf", "pages": [page_content_to_dict(p) for p in read_pdf_pages(file_path)]}
        elif file_type == "docx":
            data = {"type": "docx", "pages": [page_content_to_dict(p) for p in read_docx_pages(file_path)]}
        elif file_type == "doc":
            data = {"type": "doc", "pages": [page_content_to_dict(p) for p in read_doc_pages(file_path)]}
        elif file_type == "pptx":
            data = {"type": "pptx", "pages": [page_content_to_dict(p) for p in read_pptx_pages(file_path)]}
        elif file_type == "ppt":
            content = read_ppt(file_path, args.info)
            data = {"type": "ppt", "content": content}
        else:
            eprint(f"Unsupported file type: {file_type}")
            return 1

        output = format_json_output(data)
        if args.save:
            save_path = Path(args.save).expanduser().resolve()
            save_output(output, save_path)
        print(output)
        return 0

    if file_type == "image":
        content = read_image(file_path, args.info)
    elif file_type == "pdf":
        content = read_pdf(file_path, args.info)
    elif file_type == "docx":
        content = read_docx(file_path, args.info)
    elif file_type == "doc":
        content = read_doc(file_path, args.info)
    elif file_type == "pptx":
        content = format_output(read_pptx_pages(file_path), args.info)
    elif file_type == "ppt":
        content = read_ppt(file_path, args.info)
    else:
        eprint(f"Unsupported file type: {file_type}")
        return 1
    
    if content:
        print(content)
    
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
