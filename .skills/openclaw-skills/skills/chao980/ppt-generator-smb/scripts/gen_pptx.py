#!/usr/bin/env python3
"""PPTX 生成器 - 生成 PowerPoint 文件"""

import argparse
import json
import os

def install_pptx():
    try:
        from pptx import Presentation
        return True
    except ImportError:
        import subprocess
        subprocess.check_call(['pip', 'install', 'python-pptx', '-q'])
        return True

def create_pptx(title, subtitle, slides_data, output):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 颜色
    BG = RGBColor(0x0A, 0x0A, 0x0A)
    ORANGE = RGBColor(0xFF, 0x6B, 0x35)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY = RGBColor(0xAA, 0xAA, 0xAA)
    
    def set_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG
    
    def add_text(slide, text, left, top, width, height, size=18, color=WHITE, bold=False):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = PP_ALIGN.CENTER
    
    def add_card(slide, left, top, width, height, title, content):
        shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x18, 0x18, 0x18)
        shape.line.fill.background()
        add_text(slide, title, left+0.3, top+0.3, width-0.6, 0.5, 18, ORANGE, True)
        add_text(slide, content, left+0.3, top+0.9, width-0.6, height-1.2, 14, GRAY)
    
    # Slide 1: 封面
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, "🦞", 0, 1.5, 13.333, 1.5, 72)
    add_text(slide, title, 0, 3.5, 13.333, 1, 42, WHITE, True)
    if subtitle:
        add_text(slide, subtitle, 0, 4.8, 13.333, 0.5, 18, GRAY)
    
    # 其他幻灯片
    for slide_data in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide)
        
        stype = slide_data.get("type", "content")
        stitle = slide_data.get("title", "")
        
        add_text(slide, stitle, 0, 0.5, 13.333, 0.8, 36, WHITE, True)
        
        if stype == "content":
            items = slide_data.get("items", [])
            if isinstance(items, list):
                content = "\n".join(f"• {item}" for item in items)
            else:
                content = items
            add_text(slide, content, 1, 2, 11, 4, 18, GRAY)
        
        elif stype == "cards":
            cards = slide_data.get("cards", [])
            card_width = min(3.8, 11 / len(cards)) if cards else 3.8
            start_x = (13.333 - card_width * len(cards) - 0.5 * (len(cards) - 1)) / 2
            for i, card in enumerate(cards):
                add_card(slide, start_x + i * (card_width + 0.5), 2, card_width, 4.5,
                        f"{card.get('icon', '')} {card['title']}", card['content'])
        
        elif stype == "chart":
            data = slide_data.get("data", [])
            for i, item in enumerate(data):
                x = 2.5 + i * 2.5
                h = item.get("height", 100) / 100 * 3
                shape = slide.shapes.add_shape(1, Inches(x), Inches(5.5 - h), Inches(1.2), Inches(h))
                shape.fill.solid()
                shape.fill.fore_color.rgb = ORANGE
                shape.line.fill.background()
                add_text(slide, item["value"], x-0.2, 4.8-h, 1.6, 0.4, 14, WHITE, True)
                add_text(slide, item["label"], x, 5.6, 1.2, 0.4, 14, GRAY)
        
        elif stype == "swot":
            strengths = slide_data.get("strengths", [])
            weaknesses = slide_data.get("weaknesses", [])
            s_content = "\n".join(f"• {s}" for s in strengths)
            w_content = "\n".join(f"• {w}" for w in weaknesses)
            add_card(slide, 0.5, 1.8, 5.8, 5, "💪 Strengths", s_content)
            add_card(slide, 7, 1.8, 5.8, 5, "⚠️ Weaknesses", w_content)
    
    os.makedirs(os.path.dirname(os.path.abspath(output)), exist_ok=True)
    prs.save(output)
    print(f"Done: {output}")

def main():
    parser = argparse.ArgumentParser(description='PPTX 生成器')
    parser.add_argument('--title', required=True, help='演示文稿标题')
    parser.add_argument('--subtitle', default='', help='副标题')
    parser.add_argument('--output', default='output.pptx', help='输出文件路径')
    parser.add_argument('--data', help='JSON 数据文件')
    args = parser.parse_args()
    
    install_pptx()
    
    if args.data and os.path.exists(args.data):
        with open(args.data, 'r', encoding='utf-8') as f:
            data = json.load(f)
    else:
        data = {"slides": []}
    
    create_pptx(args.title, args.subtitle, data.get("slides", []), args.output)

if __name__ == '__main__':
    main()
