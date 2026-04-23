#!/usr/bin/env python3
"""Create a sample PowerPoint presentation with multiple slide types.

Usage: python create_presentation.py [output_path]
Default output: presentation.pptx
"""

import sys

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


def create_title_slide(prs):
    """Create the title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Sample Presentation"
    slide.placeholders[1].text = "Generated automatically with python-pptx"


def create_bullet_slide(prs):
    """Create a bullet-point slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Project Overview"

    tf = slide.placeholders[1].text_frame
    tf.text = "First point: project background"

    for text, level in [
        ("Second point: core features", 0),
        ("Sub-point: feature details", 1),
        ("Third point: timeline", 0),
    ]:
        p = tf.add_paragraph()
        p.text = text
        p.level = level


def create_table_slide(prs):
    """Create a table slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Data Summary"

    data = [
        ["Metric", "Q1", "Q2", "Q3", "Q4"],
        ["Revenue (10k CNY)", "120", "150", "180", "200"],
        ["Profit (10k CNY)", "30", "45", "55", "70"],
        ["Growth Rate", "15%", "25%", "20%", "11%"],
    ]

    rows, cols = len(data), len(data[0])
    table = slide.shapes.add_table(
        rows, cols, Inches(0.5), Inches(1.8), Inches(9), Inches(3)
    ).table

    for col_idx in range(cols):
        table.columns[col_idx].width = Inches(9 / cols)

    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            if row_idx == 0:
                p.font.bold = True
                p.font.size = Pt(12)


def create_chart_slide(prs):
    """Create a chart slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Sales Trend"

    chart_data = CategoryChartData()
    chart_data.categories = ["Q1", "Q2", "Q3", "Q4"]
    chart_data.add_series("Revenue", (120, 150, 180, 200))
    chart_data.add_series("Profit", (30, 45, 55, 70))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(1.8), Inches(9), Inches(5),
        chart_data,
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False


def create_mixed_slide(prs):
    """Create a slide with text boxes and shapes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Text box
    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Custom Content Slide"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
    p.alignment = PP_ALIGN.CENTER

    # Description
    tx_box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    tf2 = tx_box2.text_frame
    tf2.word_wrap = True
    tf2.text = "This slide shows how text boxes and shapes can be combined."
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Shapes
    from pptx.enum.shapes import MSO_SHAPE

    colors = [
        (RGBColor(0x4F, 0x81, 0xBD), "Step 1"),
        (RGBColor(0xC0, 0x50, 0x4D), "Step 2"),
        (RGBColor(0x9B, 0xBB, 0x59), "Step 3"),
        (RGBColor(0x80, 0x64, 0xA2), "Step 4"),
    ]

    left = Inches(0.8)
    for color, text in colors:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(3), Inches(2), Inches(1.2)
        )
        shape.text = text
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color

        for p in shape.text_frame.paragraphs:
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.size = Pt(14)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
        shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        left += Inches(2.3)


def main():
    output_path = sys.argv[1] if len(sys.argv) > 1 else "presentation.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs)
    create_bullet_slide(prs)
    create_table_slide(prs)
    create_chart_slide(prs)
    create_mixed_slide(prs)

    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")


if __name__ == "__main__":
    main()
