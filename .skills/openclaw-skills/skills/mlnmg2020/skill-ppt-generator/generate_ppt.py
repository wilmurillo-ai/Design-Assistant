import argparse
from pptx import Presentation
from pptx.util import Inches
import re

def parse_markdown(md_text):
    """将 Markdown 文本解析为幻灯片列表"""
    slides_content = md_text.split('---')
    slides = []
    
    for slide_text in slides_content:
        slide_text = slide_text.strip()
        if not slide_text:
            continue
            
        # 初始化
        title = ""
        bullet_points = []
        paragraphs = []
        
        lines = slide_text.split('\n')
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # 检查标题 (支持 #, ##, ###)
            if line.startswith('#'):
                # 移除 # 符号并清理空格
                title = line.lstrip('#').strip()
            # 检查项目符号
            elif line.startswith('- ') or line.startswith('* '):
                bullet_points.append(line[2:].strip())
            # 其他为段落
            else:
                paragraphs.append(line)
                
        slides.append({
            'title': title,
            'bullets': bullet_points,
            'paragraphs': paragraphs
        })
    
    return slides

def create_presentation(slides_data, output_path):
    """根据解析的数据创建 PPT"""
    prs = Presentation()
    
    # 使用 "标题和内容" 版式
    content_layout = prs.slide_layouts[1]
    
    for slide_data in slides_data:
        slide = prs.slides.add_slide(content_layout)
        
        # 标题
        if slide_data['title']:
            title = slide.shapes.title
            title.text = slide_data['title']
        
        # 内容
        if slide_data['bullets'] or slide_data['paragraphs']:
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.clear()
            
            # 添加项目符号
            for bullet in slide_data['bullets']:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
                
            # 添加段落
            for para in slide_data['paragraphs']:
                p = tf.add_paragraph()
                p.text = para
                p.level = 0
                
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

def main():
    parser = argparse.ArgumentParser(description='Convert Markdown to PowerPoint')
    parser.add_argument('--input', required=True, help='Input Markdown file')
    parser.add_argument('--output', required=True, help='Output PPTX file')
    args = parser.parse_args()

    try:
        with open(args.input, 'r', encoding='utf-8') as f:
            md_text = f.read()
            
        slides_data = parse_markdown(md_text)
        create_presentation(slides_data, args.output)
        
    except Exception as e:
        print(f"Error: {str(e)}")
        exit(1)

if __name__ == "__main__":
    main()