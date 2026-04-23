#!/usr/bin/env python3
"""
analyze_template.py — Analyze a .pptx template to extract design system info.

v4 changes vs v3:
  - Primary parser: python-pptx (reliable, handles all standard OOXML)
  - Handles sysClr (system colors), prstClr (preset colors), schemeClr (theme references)
    in addition to srgbClr — the regex-only approach missed all of these
  - Fallback: zipfile + regex if python-pptx unavailable
  - Outputs placeholder_map: per-layout list of placeholder types and names,
    so template_builder.py knows exactly what slots are available
  - Outputs suggested_palette: best-guess primary/secondary/accent from theme colors
"""

import os
import sys
import json
import zipfile
import argparse
import subprocess
import re
from pathlib import Path


# ── Preset color name → hex (common PowerPoint preset colors) ────────────────
PRESET_COLORS = {
    "black": "000000", "white": "FFFFFF", "red": "FF0000", "green": "00FF00",
    "blue": "0000FF", "yellow": "FFFF00", "cyan": "00FFFF", "magenta": "FF00FF",
    "darkRed": "8B0000", "darkGreen": "006400", "darkBlue": "00008B",
    "darkYellow": "808000", "gray": "808080", "lightGray": "D3D3D3",
    "orange": "FFA500", "pink": "FFC0CB", "purple": "800080",
    "navy": "000080", "teal": "008080", "maroon": "800000",
}

# System color approximations
SYSTEM_COLORS = {
    "windowText": "000000", "window": "FFFFFF", "btnFace": "F0F0F0",
    "btnText": "000000", "highlight": "0078D4", "highlightText": "FFFFFF",
    "grayText": "808080",
}


def resolve_color_element(el) -> str | None:
    """
    Resolve an lxml color element to a hex string.
    Handles: srgbClr, sysClr, prstClr, schemeClr (returns scheme name for schemeClr).
    """
    if el is None:
        return None
    tag = el.tag.split("}")[-1] if "}" in el.tag else el.tag

    if tag == "srgbClr":
        val = el.get("val", "")
        return val.upper() if len(val) == 6 else None

    if tag == "sysClr":
        # lastClr attribute holds the actual resolved color
        last = el.get("lastClr", "")
        if last and len(last) == 6:
            return last.upper()
        name = el.get("val", "")
        return SYSTEM_COLORS.get(name, "808080").upper()

    if tag == "prstClr":
        name = el.get("val", "")
        return PRESET_COLORS.get(name, None)

    if tag == "schemeClr":
        # Return the scheme name as a placeholder — actual value depends on theme
        return f"[scheme:{el.get('val', 'unknown')}]"

    return None


def extract_theme_colors_pptx(pptx_path: str) -> tuple[list, dict]:
    """
    Use python-pptx's lxml tree to extract theme colors robustly.
    Returns (colors_list, suggested_palette).
    """
    try:
        from pptx import Presentation
        from pptx.oxml.ns import qn

        prs = Presentation(pptx_path)
        # Theme is in slide_master[0].theme_color_map or directly in theme XML
        theme_el = prs.slide_master.element.find(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}theme"
        )

        # Access theme from slide master's theme part
        try:
            theme_part = prs.slide_masters[0].theme
            theme_el   = theme_part._element
        except Exception:
            theme_el = None

        if theme_el is None:
            # fallback to zipfile regex
            return _extract_colors_regex_fallback(pptx_path)

        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"

        role_tags = [
            ("dark1",  f"{{{ns}}}dk1"),
            ("light1", f"{{{ns}}}lt1"),
            ("dark2",  f"{{{ns}}}dk2"),
            ("light2", f"{{{ns}}}lt2"),
        ]

        colors = []
        raw = {}

        for role, tag in role_tags:
            el = theme_el.find(f".//{tag}")
            if el is not None:
                for child in el:
                    resolved = resolve_color_element(child)
                    if resolved:
                        colors.append({"role": role, "hex": resolved})
                        raw[role] = resolved
                        break

        # Accents
        for i in range(1, 7):
            el = theme_el.find(f".//{{{ns}}}accent{i}")
            if el is not None:
                for child in el:
                    resolved = resolve_color_element(child)
                    if resolved and not resolved.startswith("["):
                        colors.append({"role": f"accent{i}", "hex": resolved})
                        raw[f"accent{i}"] = resolved
                        break

        suggested = _suggest_palette(raw)
        return colors, suggested

    except ImportError:
        return _extract_colors_regex_fallback(pptx_path)
    except Exception as e:
        fallback, suggested = _extract_colors_regex_fallback(pptx_path)
        fallback.append({"note": f"python-pptx error: {e}, used regex fallback"})
        return fallback, suggested


def _extract_colors_regex_fallback(pptx_path: str) -> tuple[list, dict]:
    """Regex-based fallback for when python-pptx is unavailable or fails."""
    colors = []
    raw = {}
    try:
        with zipfile.ZipFile(pptx_path, "r") as z:
            theme_files = [n for n in z.namelist() if "theme/theme" in n and n.endswith(".xml")]
            if not theme_files:
                return colors, {}
            content = z.read(theme_files[0]).decode("utf-8", errors="ignore")

            patterns = [
                ("dark1",  r'<a:dk1>.*?<a:srgbClr val="([0-9A-Fa-f]{6})"'),
                ("light1", r'<a:lt1>.*?<a:srgbClr val="([0-9A-Fa-f]{6})"'),
                ("dark2",  r'<a:dk2>.*?<a:srgbClr val="([0-9A-Fa-f]{6})"'),
                ("light2", r'<a:lt2>.*?<a:srgbClr val="([0-9A-Fa-f]{6})"'),
            ]
            for role, pat in patterns:
                m = re.search(pat, content, re.DOTALL)
                if m:
                    colors.append({"role": role, "hex": m.group(1).upper()})
                    raw[role] = m.group(1).upper()

            accents = re.findall(r'<a:accent\d>.*?<a:srgbClr val="([0-9A-Fa-f]{6})"',
                                  content, re.DOTALL)
            for i, acc in enumerate(accents[:6], 1):
                colors.append({"role": f"accent{i}", "hex": acc.upper()})
                raw[f"accent{i}"] = acc.upper()

            # Also try sysClr lastClr
            sys_patterns = [
                ("dark1",  r'<a:dk1>.*?<a:sysClr[^>]*lastClr="([0-9A-Fa-f]{6})"'),
                ("light1", r'<a:lt1>.*?<a:sysClr[^>]*lastClr="([0-9A-Fa-f]{6})"'),
            ]
            for role, pat in sys_patterns:
                if role not in raw:
                    m = re.search(pat, content, re.DOTALL)
                    if m:
                        colors.append({"role": role, "hex": m.group(1).upper()})
                        raw[role] = m.group(1).upper()

    except Exception as e:
        colors.append({"error": str(e)})

    return colors, _suggest_palette(raw)


def _suggest_palette(raw: dict) -> dict:
    """Derive a suggested primary/secondary/accent from theme color roles."""
    # Heuristic: dark1 or accent1 → primary, light1/light2 → secondary, white/light → accent
    primary   = raw.get("dark2") or raw.get("accent1") or raw.get("dark1") or "1E2761"
    secondary = raw.get("accent2") or raw.get("light2") or raw.get("accent3") or "CADCFC"
    accent    = raw.get("light1") or "FFFFFF"

    # Clean any [scheme:...] placeholders
    def clean(h):
        return h if h and not h.startswith("[") else "808080"

    return {
        "primary":   clean(primary),
        "secondary": clean(secondary),
        "accent":    clean(accent),
    }


# ── Fonts ─────────────────────────────────────────────────────────────────────

def extract_theme_fonts(pptx_path: str) -> dict:
    fonts = {}
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        theme_part = prs.slide_masters[0].theme
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        el = theme_part._element

        major = el.find(f".//{{{ns}}}majorFont/{{{ns}}}latin")
        minor = el.find(f".//{{{ns}}}minorFont/{{{ns}}}latin")
        if major is not None: fonts["heading"] = major.get("typeface", "Calibri")
        if minor is not None: fonts["body"]    = minor.get("typeface", "Calibri")
        return fonts
    except ImportError:
        pass
    except Exception:
        pass

    # Regex fallback
    try:
        with zipfile.ZipFile(pptx_path, "r") as z:
            theme_files = [n for n in z.namelist() if "theme/theme" in n and n.endswith(".xml")]
            if not theme_files:
                return fonts
            content = z.read(theme_files[0]).decode("utf-8", errors="ignore")
            major = re.search(r'<a:majorFont>.*?<a:latin typeface="([^"]+)"', content, re.DOTALL)
            minor = re.search(r'<a:minorFont>.*?<a:latin typeface="([^"]+)"', content, re.DOTALL)
            if major: fonts["heading"] = major.group(1)
            if minor: fonts["body"]    = minor.group(1)
    except Exception as e:
        fonts["error"] = str(e)
    return fonts


# ── Slide layouts with placeholder map ───────────────────────────────────────

def extract_slide_layouts(pptx_path: str) -> list:
    """
    Extract layout names AND placeholder inventory.
    This tells template_builder.py exactly what slots exist in each layout.
    """
    layouts = []
    try:
        from pptx import Presentation
        from pptx.util import Emu

        prs = Presentation(pptx_path)
        for layout in prs.slide_layouts:
            placeholders = []
            for ph in layout.placeholders:
                placeholders.append({
                    "idx":  ph.placeholder_format.idx,
                    "type": str(ph.placeholder_format.type).split(".")[-1],
                    "name": ph.name,
                })
            layouts.append({
                "name":         layout.name,
                "placeholders": placeholders,
                "ph_count":     len(placeholders),
            })
        return layouts
    except ImportError:
        pass
    except Exception:
        pass

    # Regex fallback
    try:
        with zipfile.ZipFile(pptx_path, "r") as z:
            layout_files = sorted([n for n in z.namelist()
                                    if "slideLayouts/slideLayout" in n and n.endswith(".xml")])
            for lf in layout_files:
                content = z.read(lf).decode("utf-8", errors="ignore")
                name_m = re.search(r'<p:cSld name="([^"]*)"', content)
                type_m = re.search(r'type="([^"]*)"', content)
                layouts.append({
                    "name": name_m.group(1) if name_m else "Unknown",
                    "type": type_m.group(1) if type_m else "",
                    "placeholders": [],
                    "ph_count": len(re.findall(r'<p:sp>', content)),
                })
    except Exception as e:
        layouts.append({"error": str(e)})
    return layouts


# ── Dimensions & slide count ──────────────────────────────────────────────────

def get_slide_dimensions(pptx_path: str) -> dict:
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        w = prs.slide_width
        h = prs.slide_height
        return {
            "width_inches":  round(w.inches, 3),
            "height_inches": round(h.inches, 3),
            "cx_emu": int(w),
            "cy_emu": int(h),
            "aspect": f"{round(w / h * 9, 1)}:9",
        }
    except ImportError:
        pass
    except Exception:
        pass

    # zipfile fallback
    try:
        with zipfile.ZipFile(pptx_path, "r") as z:
            if "ppt/presentation.xml" in z.namelist():
                content = z.read("ppt/presentation.xml").decode("utf-8", errors="ignore")
                m = re.search(r'<p:sldSz cx="(\d+)" cy="(\d+)"', content)
                if m:
                    cx, cy = int(m.group(1)), int(m.group(2))
                    return {
                        "width_inches":  round(cx / 914400, 3),
                        "height_inches": round(cy / 914400, 3),
                        "cx_emu": cx, "cy_emu": cy,
                        "aspect": f"{round(cx/cy * 9, 1)}:9",
                    }
    except Exception as e:
        return {"error": str(e)}
    return {}


def get_slide_count(pptx_path: str) -> int:
    try:
        from pptx import Presentation
        return len(Presentation(pptx_path).slides)
    except Exception:
        pass
    try:
        with zipfile.ZipFile(pptx_path, "r") as z:
            return len([n for n in z.namelist()
                        if "ppt/slides/slide" in n and n.endswith(".xml")])
    except:
        return 0


def get_content_overview(pptx_path: str) -> str:
    try:
        result = subprocess.run(
            ["python3", "-m", "markitdown", pptx_path],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0:
            text = result.stdout
            return text[:3000] + "\n... [truncated]" if len(text) > 3000 else text
    except Exception:
        pass
    return ""


# ── Entry point ───────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="Analyze a PPTX template")
    parser.add_argument("template", help="Path to .pptx template file")
    args = parser.parse_args()

    if not os.path.exists(args.template):
        print(json.dumps({"error": f"File not found: {args.template}"}))
        sys.exit(1)

    theme_colors, suggested_palette = extract_theme_colors_pptx(args.template)

    result = {
        "template_file":     os.path.basename(args.template),
        "slide_count":       get_slide_count(args.template),
        "dimensions":        get_slide_dimensions(args.template),
        "theme_colors":      theme_colors,
        "suggested_palette": suggested_palette,   # ready-to-use primary/secondary/accent hex
        "theme_fonts":       extract_theme_fonts(args.template),
        "slide_layouts":     extract_slide_layouts(args.template),
        "content_overview":  get_content_overview(args.template),
        "usage_notes": (
            "Pass suggested_palette.primary/secondary/accent directly to the LLM prompt. "
            "slide_layouts[].placeholders shows exactly what content slots exist per layout — "
            "use this when building the slide plan so content maps cleanly. "
            "Use template_builder.py (not XML editing) to fill content into the template."
        )
    }
    print(json.dumps(result, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
