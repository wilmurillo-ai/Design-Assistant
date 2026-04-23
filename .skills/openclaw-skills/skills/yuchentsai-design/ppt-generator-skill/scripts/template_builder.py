#!/usr/bin/env python3
"""
template_builder.py — Fill a .pptx template with slide plan content.

Replaces the broken "from-template" mode that previously just printed
instructions for manually editing XML.

Strategy:
  1. Read template .pptx into python-pptx
  2. Identify which slide layouts exist (by name/index)
  3. For each slide in the plan:
       a. Find the best matching layout in the template
       b. Add a new slide from that layout
       c. Fill all text placeholders with content from the plan
       d. Preserve all visual formatting (colors, fonts, images, shapes)
  4. Write output .pptx — template structure is NEVER directly mutated,
     output is always a fresh presentation built from template's slide masters

This approach is safe: python-pptx adds slides via the layout system,
so the template file itself is never corrupted regardless of input.

Usage:
  python3 scripts/template_builder.py \
    --template my_template.pptx \
    --plan slide_plan.json \
    --output output.pptx
"""

import os
import sys
import json
import argparse
import re
from pathlib import Path


# ── python-pptx availability check ───────────────────────────────────────────

def require_pptx():
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.enum.text import PP_ALIGN
        return True
    except ImportError:
        print(json.dumps({
            "status": "error",
            "message": "python-pptx not installed. Run: pip install python-pptx --break-system-packages"
        }))
        sys.exit(1)


# ── Layout matching ───────────────────────────────────────────────────────────

# Maps our logical layout names → keywords to search in template layout names
LAYOUT_KEYWORDS = {
    "title":      ["title slide", "title", "封面", "cover", "首页"],
    "section":    ["section", "divider", "章节", "separator", "blank", "空白"],
    "content":    ["content", "text", "内容", "body", "bullet", "正文"],
    "two-column": ["two col", "2 col", "comparison", "对比", "双栏", "columns"],
    "quote":      ["quote", "引用", "statement", "testimonial"],
    "stat":       ["blank", "空白", "content", "number"],  # fallback: use blank for stats
}

# Priority fallback order when no keyword match found
FALLBACK_ORDER = ["content", "title", "blank", ""]


def find_best_layout(slide_layouts, logical_layout: str):
    """
    Given a list of python-pptx slide layout objects, find the best match
    for our logical layout name. Returns the layout object.
    """
    keywords = LAYOUT_KEYWORDS.get(logical_layout, ["content"])

    # Build a lookup of layout names (lowercase)
    layout_map = {}
    for layout in slide_layouts:
        name_lower = (layout.name or "").lower()
        layout_map[name_lower] = layout

    # Try keyword matching (partial)
    for kw in keywords:
        kw_lower = kw.lower()
        for name, layout in layout_map.items():
            if kw_lower in name:
                return layout

    # Fallback: return first layout with "content" in name
    for name, layout in layout_map.items():
        if "content" in name:
            return layout

    # Last resort: first layout
    return slide_layouts[0] if slide_layouts else None


# ── Placeholder filling ───────────────────────────────────────────────────────

def set_placeholder_text(placeholder, text: str, bold: bool = False):
    """Safely set text on a placeholder, preserving its existing formatting."""
    try:
        tf = placeholder.text_frame
        tf.clear()
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = text
        if bold and run.font:
            run.font.bold = True
    except Exception:
        try:
            placeholder.text = text
        except Exception:
            pass  # Some placeholders are non-text (image, chart) — skip safely


def fill_slide_from_plan(prs_slide, slide_data: dict, palette: dict):
    """
    Fill placeholders on a slide using slide_data from the plan.
    Tries to match placeholders by type (title, body, subtitle).
    Does NOT touch any shapes that aren't placeholders → preserves template visuals.
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    import pptx.util as util

    title_text    = slide_data.get("title", "")
    subtitle_text = slide_data.get("subtitle", "")
    body_items    = slide_data.get("body", [])
    notes_text    = slide_data.get("notes", "")

    # Categorize placeholders by type idx
    # OOXML placeholder type indices:
    #   0 = title, 1 = body/content, 2 = subtitle, 3+ = various
    title_ph   = None
    body_ph    = None
    sub_ph     = None
    other_phs  = []

    for ph in prs_slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            title_ph = ph
        elif idx == 1:
            body_ph = ph
        elif idx == 2:
            sub_ph = ph
        else:
            other_phs.append(ph)

    # Fill title
    if title_ph and title_text:
        set_placeholder_text(title_ph, title_text, bold=True)

    # Fill subtitle (used on title slides)
    if sub_ph and subtitle_text:
        set_placeholder_text(sub_ph, subtitle_text)
    elif sub_ph and body_items:
        # If no subtitle but there's a subtitle placeholder, use first body item
        set_placeholder_text(sub_ph, body_items[0] if body_items else "")

    # Fill body / content placeholder
    if body_ph and body_items:
        try:
            tf = body_ph.text_frame
            tf.clear()
            for i, item in enumerate(body_items):
                if i == 0:
                    para = tf.paragraphs[0]
                else:
                    para = tf.add_paragraph()
                run = para.add_run()
                run.text = item
                # Preserve existing paragraph level/indent
        except Exception:
            try:
                body_ph.text = "\n".join(body_items)
            except Exception:
                pass
    elif body_ph and subtitle_text and not sub_ph:
        # No subtitle placeholder — use body for subtitle
        set_placeholder_text(body_ph, subtitle_text)

    # Speaker notes
    if notes_text:
        try:
            notes_slide = prs_slide.notes_slide
            tf = notes_slide.notes_text_frame
            tf.clear()
            para = tf.paragraphs[0]
            run = para.add_run()
            run.text = notes_text
        except Exception:
            pass


# ── Main builder ──────────────────────────────────────────────────────────────

def build_from_template(template_path: str, plan: dict, output_path: str) -> dict:
    """
    Core function: load template, add slides from plan, write output.
    Returns result dict.
    """
    require_pptx()
    from pptx import Presentation

    if not os.path.exists(template_path):
        return {"status": "error", "message": f"Template not found: {template_path}"}

    try:
        prs = Presentation(template_path)
    except Exception as e:
        return {"status": "error", "message": f"Failed to open template: {e}"}

    slide_layouts = list(prs.slide_layouts)
    if not slide_layouts:
        return {"status": "error", "message": "Template has no slide layouts"}

    slides_data  = plan.get("slides", [])
    palette      = plan.get("palette", {})

    # Record original slide count BEFORE adding new slides
    original_count = len(prs.slides)

    # Track which layouts were used and any issues
    layout_log = []
    errors     = []

    for i, slide_data in enumerate(slides_data):
        logical_layout = slide_data.get("layout", "content")
        layout_obj = find_best_layout(slide_layouts, logical_layout)

        if layout_obj is None:
            errors.append(f"slide {i+1}: no matching layout found, skipping")
            continue

        try:
            new_slide = prs.slides.add_slide(layout_obj)
            fill_slide_from_plan(new_slide, slide_data, palette)
            layout_log.append({
                "slide_index": i + 1,
                "logical_layout": logical_layout,
                "template_layout_used": layout_obj.name,
            })
        except Exception as e:
            errors.append(f"slide {i+1} ({logical_layout}): {e}")

    # Remove the original template slides (they come before our new ones)
    # original_count was captured before any add_slide calls — always accurate
    slides_xml_list = prs.slides._sldIdLst

    # Remove original slides from the front
    original_ids = list(slides_xml_list)[:original_count]
    for sld_id in original_ids:
        slides_xml_list.remove(sld_id)

    try:
        prs.save(output_path)
    except Exception as e:
        return {"status": "error", "message": f"Failed to save output: {e}"}

    return {
        "status": "ok",
        "output": output_path,
        "slides_generated": len(slides_data) - len(errors),
        "layout_mapping": layout_log,
        "errors": errors,
    }


# ── Entry point ───────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="Fill a PPTX template with slide plan content")
    parser.add_argument("--template", required=True, help="Path to .pptx template file")
    parser.add_argument("--plan",     required=True, help="Path to validated slide plan JSON")
    parser.add_argument("--output",   required=True, help="Output .pptx path")
    args = parser.parse_args()

    with open(args.plan, "r", encoding="utf-8") as f:
        plan = json.load(f)

    result = build_from_template(
        template_path=os.path.abspath(args.template),
        plan=plan,
        output_path=os.path.abspath(args.output),
    )
    print(json.dumps(result, ensure_ascii=False, indent=2))
    if result["status"] != "ok":
        sys.exit(1)


if __name__ == "__main__":
    main()
