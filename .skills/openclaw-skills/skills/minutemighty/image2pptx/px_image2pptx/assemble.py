"""PPTX assembly — place text boxes and background onto editable slides.

Pure python-pptx assembly: no ML models, no LLM calls. Takes OCR regions,
background image, and optional tight mask → produces editable .pptx.
"""

from __future__ import annotations

from pathlib import Path

import numpy as np
from PIL import Image, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor


# ── Coordinate mapping ────────────────────────────────────────


def px_to_emu(px: float, px_per_inch: float) -> int:
    """Convert image pixels to EMU (914400 per inch)."""
    return int(px / px_per_inch * 914400)


class SlideMapper:
    """Maps image pixel coordinates to slide EMU coordinates."""

    def __init__(self, img_w: int, img_h: int, slide_w_inches: float | None = None):
        self.img_w = img_w
        self.img_h = img_h
        aspect = img_w / img_h

        if slide_w_inches and slide_w_inches > 0:
            self.slide_w = slide_w_inches
            self.slide_h = slide_w_inches / aspect
        elif aspect > 1.5:
            self.slide_w, self.slide_h = 13.333, 7.5
        elif aspect > 1.2:
            self.slide_w, self.slide_h = 10.0, 7.5
        else:
            self.slide_w = 10.0
            self.slide_h = 10.0 / aspect

        self.ppi = img_w / self.slide_w

    def to_emu(self, px: float) -> int:
        return px_to_emu(px, self.ppi)

    def bbox_to_emu(self, x1, y1, x2, y2):
        return (self.to_emu(x1), self.to_emu(y1),
                self.to_emu(x2 - x1), self.to_emu(y2 - y1))


# ── Font measurement ──────────────────────────────────────────


def _load_reference_font():
    """Load a system sans-serif font for text width measurement."""
    candidates = [
        "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/System/Library/Fonts/Helvetica.ttc",
        "/Library/Fonts/Arial.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    for path in candidates:
        try:
            return ImageFont.truetype(path, 72), True
        except Exception:
            continue
    return None, False


_REF_FONT, _HAS_FONT = _load_reference_font()


def _is_cjk(ch: str) -> bool:
    cp = ord(ch)
    return (0x4E00 <= cp <= 0x9FFF or 0x3400 <= cp <= 0x4DBF or
            0x3000 <= cp <= 0x303F or 0xFF00 <= cp <= 0xFFEF or
            0xF900 <= cp <= 0xFAFF or 0x2E80 <= cp <= 0x2EFF or
            0x31C0 <= cp <= 0x31EF)


def estimate_text_width_pt(text: str, font_pt: float) -> float:
    """Measure rendered text width in points.

    Hybrid: PIL font metrics for Latin, 1.0x em for CJK.
    """
    if _HAS_FONT:
        width = 0.0
        latin_run: list[str] = []

        def flush():
            nonlocal width
            if latin_run:
                ref_w = _REF_FONT.getlength("".join(latin_run))
                width += ref_w * (font_pt / 72.0)
                latin_run.clear()

        for ch in text:
            if _is_cjk(ch):
                flush()
                width += font_pt * 1.0
            else:
                latin_run.append(ch)
        flush()
        return width

    # Fallback: heuristic
    width = 0.0
    for ch in text:
        if _is_cjk(ch):
            width += font_pt * 1.0
        elif ch == " ":
            width += font_pt * 0.25
        else:
            width += font_pt * 0.50
    return width


def autoscale_font(
    text: str,
    bbox_w_px: float,
    bbox_h_px: float,
    ppi: float,
    min_pt: int = 8,
    max_pt: int = 72,
) -> int:
    """Auto-scale font size to fill 90-94% of bbox width."""
    line_h_pt = (bbox_h_px / ppi) * 72
    pt = max(min_pt, min(max_pt, round(line_h_pt)))

    bbox_w_pt = (bbox_w_px / ppi) * 72
    lines = text.split("\n")
    longest = max(lines, key=len) if lines else text

    # Shrink to fit
    for _ in range(40):
        if estimate_text_width_pt(longest, pt) <= bbox_w_pt * 0.94 or pt <= min_pt:
            break
        pt = max(min_pt, pt - 1)

    # Grow to fill
    for _ in range(40):
        if estimate_text_width_pt(longest, pt) >= bbox_w_pt * 0.90 or pt >= max_pt:
            break
        if estimate_text_width_pt(longest, pt + 1) > bbox_w_pt * 0.94:
            break
        pt = min(max_pt, pt + 1)

    return pt


# ── Text grouping ─────────────────────────────────────────────


def group_text_lines(
    regions: list[dict],
    y_threshold: float = 0.6,
    x_gap_factor: float = 3.0,
) -> list[list[dict]]:
    """Merge word-level OCR regions into line-level groups.

    Two-pass: group by vertical proximity, then split by horizontal gaps
    to prevent merging left/right columns.
    """
    if not regions:
        return []

    for r in regions:
        b = r["bbox"]
        r["_cy"] = (b["y1"] + b["y2"]) / 2
        r["_h"] = b["y2"] - b["y1"]

    sorted_regions = sorted(regions, key=lambda r: r["_cy"])

    # Pass 1: vertical grouping
    y_lines: list[list[dict]] = []
    current = [sorted_regions[0]]
    for r in sorted_regions[1:]:
        line_cy = sum(rr["_cy"] for rr in current) / len(current)
        line_h = max(rr["_h"] for rr in current)
        if abs(r["_cy"] - line_cy) < line_h * y_threshold:
            current.append(r)
        else:
            y_lines.append(current)
            current = [r]
    y_lines.append(current)

    # Pass 2: split by horizontal gaps
    lines: list[list[dict]] = []
    for y_line in y_lines:
        y_line.sort(key=lambda r: r["bbox"]["x1"])
        if len(y_line) <= 1:
            lines.append(y_line)
            continue
        heights = sorted(rr["_h"] for rr in y_line)
        median_h = heights[len(heights) // 2]
        max_gap = median_h * x_gap_factor

        segment = [y_line[0]]
        for r in y_line[1:]:
            gap = r["bbox"]["x1"] - segment[-1]["bbox"]["x2"]
            if gap > max_gap:
                lines.append(segment)
                segment = [r]
            else:
                segment.append(r)
        lines.append(segment)

    for r in regions:
        r.pop("_cy", None)
        r.pop("_h", None)

    return lines


def group_bbox(group: list[dict]) -> tuple[int, int, int, int]:
    x1 = min(r["bbox"]["x1"] for r in group)
    y1 = min(r["bbox"]["y1"] for r in group)
    x2 = max(r["bbox"]["x2"] for r in group)
    y2 = max(r["bbox"]["y2"] for r in group)
    return x1, y1, x2, y2


def group_to_text(group: list[dict]) -> str:
    """Convert a group of OCR regions to display text."""
    if not group:
        return ""
    for r in group:
        b = r["bbox"]
        r["_cy"] = (b["y1"] + b["y2"]) / 2
        r["_h"] = b["y2"] - b["y1"]

    sorted_r = sorted(group, key=lambda r: r["_cy"])
    lines: list[list[dict]] = []
    current = [sorted_r[0]]
    for r in sorted_r[1:]:
        line_cy = sum(rr["_cy"] for rr in current) / len(current)
        line_h = max(rr["_h"] for rr in current)
        if abs(r["_cy"] - line_cy) < line_h * 0.6:
            current.append(r)
        else:
            lines.append(current)
            current = [r]
    lines.append(current)

    text_lines = []
    for line in lines:
        line.sort(key=lambda r: r["bbox"]["x1"])
        text_lines.append(" ".join(r["text"] for r in line))

    for r in group:
        r.pop("_cy", None)
        r.pop("_h", None)

    return "\n".join(text_lines)


# ── Text color detection ──────────────────────────────────────


def _local_bg_color(crop: np.ndarray, border: int = 2) -> np.ndarray:
    h, w = crop.shape[:2]
    if h < border * 2 + 1 or w < border * 2 + 1:
        return np.median(crop.reshape(-1, 3), axis=0)
    pixels = np.concatenate([
        crop[:border, :].reshape(-1, 3),
        crop[-border:, :].reshape(-1, 3),
        crop[border:-border, :border].reshape(-1, 3),
        crop[border:-border, -border:].reshape(-1, 3),
    ])
    return np.median(pixels, axis=0)


def detect_text_color(
    img_rgb: np.ndarray,
    tight_mask: np.ndarray,
    x1: int, y1: int, x2: int, y2: int,
    default: tuple[int, int, int] = (0x33, 0x33, 0x33),
    min_contrast: float = 40,
) -> tuple[int, int, int]:
    """Detect text color from original image using tight mask.

    Strategy 1: median of tight-mask ink pixels (dark text on light bg).
    Strategy 2: if color ≈ background, sample pixels most different from bg
    (handles white text on dark bg where textmask misses the text).
    """
    h, w = img_rgb.shape[:2]
    bx1, by1 = max(0, int(x1)), max(0, int(y1))
    bx2, by2 = min(w, int(x2)), min(h, int(y2))
    if bx2 <= bx1 or by2 <= by1:
        return default

    crop = img_rgb[by1:by2, bx1:bx2]
    mask_crop = tight_mask[by1:by2, bx1:bx2]
    bg = _local_bg_color(crop)

    # Strategy 1: tight mask ink pixels
    ink_pixels = crop[mask_crop > 128]
    if len(ink_pixels) >= 3:
        median = np.median(ink_pixels, axis=0)
        dist = float(((median - bg.astype(float)) ** 2).sum() ** 0.5)
        if dist >= min_contrast:
            c = median.astype(int)
            return (int(c[0]), int(c[1]), int(c[2]))

    # Strategy 2: pixels most different from background
    flat = crop.reshape(-1, 3).astype(float)
    dists = np.sqrt(((flat - bg.astype(float)) ** 2).sum(axis=1))
    threshold = np.percentile(dists, 80)
    far_pixels = flat[dists >= max(threshold, min_contrast * 0.5)]

    if len(far_pixels) >= 3:
        median = np.median(far_pixels, axis=0).astype(int)
        return (int(median[0]), int(median[1]), int(median[2]))

    return default


# ── Background detection ──────────────────────────────────────


def detect_bg_color(image_path: str, border_px: int = 20) -> tuple[int, ...] | None:
    """Sample border pixels to detect solid background color.

    Returns (r, g, b) if low variance, else None.
    """
    img = np.array(Image.open(image_path).convert("RGB"))
    h, w = img.shape[:2]
    border = np.concatenate([
        img[:border_px, :].reshape(-1, 3),
        img[-border_px:, :].reshape(-1, 3),
        img[border_px:-border_px, :border_px].reshape(-1, 3),
        img[border_px:-border_px, -border_px:].reshape(-1, 3),
    ])
    std = border.std(axis=0).mean()
    if std < 25:
        median = np.median(border, axis=0).astype(int)
        return tuple(median)
    return None


# ── PPTX assembly ─────────────────────────────────────────────


def assemble_pptx(
    image_path: str,
    ocr_regions: list[dict],
    output_path: str,
    background_path: str | None = None,
    tight_mask: np.ndarray | None = None,
    min_font: int = 8,
    max_font: int = 72,
    slide_w_inches: float | None = None,
) -> dict:
    """Assemble an editable PPTX from OCR regions and background.

    Args:
        image_path: Original input image.
        ocr_regions: List of OCR region dicts with bbox and text.
        output_path: Where to save the .pptx file.
        background_path: Inpainted background image (or None for solid bg).
        tight_mask: Pre-dilation text mask for color detection (H, W), uint8.
        min_font: Minimum font size in points.
        max_font: Maximum font size in points.
        slide_w_inches: Override slide width (auto-detected from aspect ratio).

    Returns:
        Report dict with assembly statistics.
    """
    img = Image.open(image_path)
    img_w, img_h = img.size
    mapper = SlideMapper(img_w, img_h, slide_w_inches)

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(mapper.slide_w)
    prs.slide_height = Inches(mapper.slide_h)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Background
    if background_path and Path(background_path).exists():
        slide.shapes.add_picture(
            background_path, Emu(0), Emu(0),
            Inches(mapper.slide_w), Inches(mapper.slide_h),
        )
        bg_mode = "inpainted"
    else:
        bg_color = detect_bg_color(image_path)
        if bg_color:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*bg_color)
            bg_mode = f"solid rgb{bg_color}"
        else:
            slide.shapes.add_picture(
                image_path, Emu(0), Emu(0),
                Inches(mapper.slide_w), Inches(mapper.slide_h),
            )
            bg_mode = "original"

    # Load image array for color detection
    img_rgb = None
    if tight_mask is not None:
        img_rgb = np.array(Image.open(image_path).convert("RGB"))

    # Group OCR regions into lines
    text_groups = group_text_lines(ocr_regions)

    # Add text boxes
    count = 0
    for group in text_groups:
        x1, y1, x2, y2 = group_bbox(group)
        text = group_to_text(group)
        if not text.strip():
            continue

        left, top, width, height = mapper.bbox_to_emu(x1, y1, x2, y2)
        pad = mapper.to_emu(2)
        left = max(0, left - pad)
        top = max(0, top - pad)
        width += pad * 2
        height += pad * 2

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)

        # Font size
        bbox_w = x2 - x1
        region_heights = [r["bbox"]["y2"] - r["bbox"]["y1"] for r in group]
        line_h = sorted(region_heights)[len(region_heights) // 2]
        font_size = autoscale_font(text, bbox_w, line_h, mapper.ppi, min_font, max_font)

        # Font color
        if img_rgb is not None and tight_mask is not None:
            r, g, b = detect_text_color(img_rgb, tight_mask, x1, y1, x2, y2)
        else:
            r, g, b = 0x33, 0x33, 0x33
        color = RGBColor(r, g, b)

        lines = text.split("\n")
        p = tf.paragraphs[0]
        p.text = lines[0]
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        for line in lines[1:]:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(font_size)
            p.font.color.rgb = color

        count += 1

    prs.save(output_path)

    return {
        "image_size": {"width": img_w, "height": img_h},
        "slide_size": {
            "width_inches": round(mapper.slide_w, 2),
            "height_inches": round(mapper.slide_h, 2),
        },
        "ppi": round(mapper.ppi, 1),
        "background": bg_mode,
        "text_boxes": count,
        "ocr_regions": len(ocr_regions),
    }
