#!/usr/bin/env python3
import argparse
import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT_DIR = os.path.expanduser("~/Documents")
MAX_PAGES = 10


def create_presentation(title: str, sections: str, output_path: str = None):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    if output_path is None:
        date_str = datetime.now().strftime("%Y-%m-%d")
        output_path = os.path.join(OUTPUT_DIR, f"工作汇报_{date_str}.pptx")

    blank_layout = prs.slide_layouts[6]

    # Slide 1: Title
    slide1 = prs.slides.add_slide(blank_layout)
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "SimSun"
    p.font.size = Pt(48)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Content slides
    section_list = sections.split("|")

    for section_str in section_list[: MAX_PAGES - 2]:
        if ":" not in section_str:
            continue

        slide = prs.slides.add_slide(blank_layout)

        colon_idx = section_str.index(":")
        section_title = section_str[:colon_idx].strip()
        section_content = section_str[colon_idx + 1 :].strip()

        # Section title - 36pt
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = section_title
        p.font.name = "SimSun"
        p.font.size = Pt(36)
        p.font.bold = True

        # Section content - 18pt
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.2))
        tf = content_box.text_frame
        tf.word_wrap = True

        content_parts = section_content.split("\\n")
        for i, part in enumerate(content_parts):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = part.strip()
            p.font.name = "SimSun"
            p.font.size = Pt(18)
            p.space_after = Pt(10)

    # Last slide: Thank you
    slide_end = prs.slides.add_slide(blank_layout)
    thank_box = slide_end.shapes.add_textbox(Inches(0), Inches(2.8), Inches(13.333), Inches(2))
    tf = thank_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "谢谢观看"
    p.font.name = "SimSun"
    p.font.size = Pt(48)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    prs.save(output_path)

    page_count = len(prs.slides)
    print(f"✓ 已生成PPT: {output_path}")
    print(f"  标题: {title}")
    print(f"  页数: {page_count} 页")

    return output_path


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="生成工作汇报PPT")
    parser.add_argument("--title", required=True, help="汇报标题")
    parser.add_argument("--sections", required=True, help="内容板块，格式: 板块名:内容|板块名:内容")
    parser.add_argument("--output", help="输出文件路径")

    args = parser.parse_args()

    create_presentation(args.title, args.sections, args.output)
