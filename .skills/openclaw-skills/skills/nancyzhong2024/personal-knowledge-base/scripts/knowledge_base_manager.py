#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
个人知识库向量数据库管理器
基于FAISS实现，提供知识库的创建、查询、删除等功能

支持的操作：
1. 保存、分割、向量化存储文件
2. 检索并回答问题
3. 列举知识库
4. 列举知识库下的文件
5. 删除知识库下的某个文件
6. 创建知识向量库
7. 更新知识库下的某个文件
8. 检索知识库中文件的某个文本块
"""

import os
import logging
import pickle
import numpy as np
from typing import List, Dict, Any, Optional
import hashlib
import shutil

# 导入第三方库
try:
    import faiss
    from langchain_text_splitters import RecursiveCharacterTextSplitter
    from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader, UnstructuredMarkdownLoader, TextLoader
    from langchain_core.documents import Document
    from langchain_core.documents import Document
except ImportError as e:
    print(f"缺少必要的依赖包，请安装: pip install faiss-cpu langchain langchain-community langchain-text-splitters pypdf python-docx python-pptx markdown")
    raise e

# 尝试导入python-pptx用于加载PPT文件
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class CustomPPTXLoader:
    """自定义PPT加载器，使用python-pptx"""

    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> List[Document]:
        """加载PPT文件并返回文档列表"""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx未安装，请运行: pip install python-pptx")

        documents = []
        prs = Presentation(self.file_path)

        # 遍历所有幻灯片
        for slide_num, slide in enumerate(prs.slides, start=1):
            text_content = []

            # 获取幻灯片中的所有文本
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_content.append(shape.text)

                # 获取表格中的文本
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text:
                                text_content.append(cell.text)

            if text_content:
                full_text = "\n".join(text_content)
                documents.append(Document(
                    page_content=full_text,
                    metadata={"source": f"slide_{slide_num}", "slide_number": slide_num}
                ))

        return documents

# 设置日志
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)


# ==================== 配置参数 ====================
# 这些参数可以通过调用时传入覆盖

# 文本处理常量
DEFAULT_CHUNK_SIZE = 1000  # 文本块大小
DEFAULT_CHUNK_OVERLAP_RATIO = 0.1  # 文本块重叠比例（重叠大小 = 文本块大小 * 0.1）
DEFAULT_TOP_K = 3  # 默认检索的文本块数量
DEFAULT_EMBEDDING_VECTOR_SIZE = 2048  # 嵌入向量大小 (智谱embedding-3模型的向量维度)
DEFAULT_SIMILARITY_THRESHOLD = 0.5  # 相似度过滤阈值

# 智谱AI相关常量
ZHIPU_EMBEDDING_BATCH_SIZE = 10  # 智谱嵌入模型单次批处理的最大文本数

# 默认值（当配置文件不存在或读取失败时使用）
_DEFAULT_CONFIG = {
    "workspace_root": "d:/Nancy/MyWork/WorkBuddyWorkSpace/MyKnowledgeBase",
    "embedding_model": "embedding-3",
    "DEFAULT_CHUNK_SIZE": 1000,
    "DEFAULT_TOP_K": 3,
    "DEFAULT_SIMILARITY_THRESHOLD": 0.5
}

# 全局配置缓存
_config_cache = None


def get_config_file_path():
    """获取配置文件路径"""
    # 技能脚本所在目录
    script_dir = os.path.dirname(os.path.abspath(__file__))
    # 技能根目录
    skill_dir = os.path.dirname(script_dir)
    config_path = os.path.join(skill_dir, "config.txt")
    return config_path


def load_config() -> Dict[str, Any]:
    """
    从配置文件加载配置

    Returns:
        配置字典
    """
    global _config_cache

    if _config_cache is not None:
        return _config_cache

    config_path = get_config_file_path()
    config = _DEFAULT_CONFIG.copy()

    try:
        if os.path.exists(config_path):
            logger.info(f"从配置文件加载配置: {config_path}")
            with open(config_path, 'r', encoding='utf-8') as f:
                for line in f:
                    line = line.strip()
                    # 跳过空行和注释行
                    if not line or line.startswith('#'):
                        continue

                    # 解析 key=value 格式
                    if '=' in line:
                        key, value = line.split('=', 1)
                        key = key.strip()
                        value = value.strip()

                        # 根据键名转换类型
                        if key in ["DEFAULT_CHUNK_SIZE", "DEFAULT_TOP_K"]:
                            try:
                                config[key] = int(value)
                            except ValueError:
                                logger.warning(f"配置项 {key} 值无效: {value}")
                        elif key == "DEFAULT_SIMILARITY_THRESHOLD":
                            try:
                                config[key] = float(value)
                            except ValueError:
                                logger.warning(f"配置项 {key} 值无效: {value}")
                        elif key == "workspace_root":
                            # 处理Windows路径
                            config[key] = value.replace("\\", "/")
                        elif key == "embedding_model":
                            config[key] = value
                        else:
                            config[key] = value

            logger.info(f"配置加载成功: {config}")
        else:
            logger.warning(f"配置文件不存在，使用默认配置: {config_path}")
    except Exception as e:
        logger.error(f"加载配置文件失败: {str(e)}，使用默认配置")

    _config_cache = config
    return config


def reload_config():
    """重新加载配置（清除缓存）"""
    global _config_cache
    _config_cache = None
    return load_config()


class ZhipuLLM:
    """智谱AI大语言模型包装类"""

    def __init__(self, api_key: str = None, model: str = None):
        """
        初始化智谱LLM

        Args:
            api_key: 智谱API密钥
            model: 模型名称（从配置读取，默认glm-4-flash）
        """
        # 加载配置
        config = load_config()

        # 优先使用传入的api_key，其次从环境变量读取，最后尝试配置文件
        if api_key:
            self.api_key = api_key
        else:
            self.api_key = os.getenv("ZHIPUAI_API_KEY") or config.get("ZHIPUAI_API_KEY")

        # 使用传入的model或从配置读取
        self.model = model or config.get("ZHIPUAI_LLM_MODEL", "glm-4-flash")
        self.client = None

        if not self.api_key:
            logger.warning("未找到智谱API密钥，无法使用LLM生成答案")
        else:
            try:
                from zhipuai import ZhipuAI
                self.client = ZhipuAI(api_key=self.api_key)
                logger.info("智谱LLM初始化成功")
            except ImportError as e:
                logger.warning(f"智谱AI SDK初始化失败: {str(e)}")
            except Exception as e:
                logger.warning(f"智谱LLM初始化失败: {str(e)}")

    def generate(self, prompt: str) -> str:
        """
        使用LLM生成答案

        Args:
            prompt: 提示词

        Returns:
            生成的答案
        """
        if not self.client:
            return None

        try:
            response = self.client.chat.completions.create(
                model=self.model,
                messages=[
                    {"role": "system", "content": "你是一个专业的问答助手。你需要根据提供的上下文信息，结合你的知识，规范、准确地回答用户的问题。回答要条理清晰、可读性强。"},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.7,
                max_tokens=2048
            )
            return response.choices[0].message.content
        except Exception as e:
            logger.error(f"LLM生成答案失败: {str(e)}")
            return None


class ZhipuEmbeddings:
    """智谱AI嵌入模型包装类"""

    def __init__(self, api_key: str = None, embedding_model: str = "embedding-3"):
        """
        初始化智谱嵌入模型

        Args:
            api_key: 智谱API密钥，如果不提供则从配置文件或环境变量获取
            embedding_model: 嵌入模型名称
        """
        # 优先使用传入的api_key，其次从环境变量读取，最后尝试配置文件
        if api_key:
            self.api_key = api_key
        else:
            # 优先从环境变量读取
            self.api_key = os.getenv("ZHIPUAI_API_KEY")
            # 如果环境变量没有，则尝试从配置文件读取（向后兼容）
            if not self.api_key:
                config = load_config()
                self.api_key = config.get("ZHIPUAI_API_KEY")

        self.embedding_model = embedding_model

        if not self.api_key:
            logger.warning("未找到智谱API密钥，将使用模拟嵌入向量（仅用于测试）")
            self.client = None
        else:
            try:
                from zhipuai import ZhipuAI
                self.client = ZhipuAI(api_key=self.api_key)
                logger.info("智谱AI嵌入模型初始化成功")
            except ImportError as e:
                logger.warning(f"请安装智谱AI SDK: pip install zhipuai ({str(e)})，将使用模拟嵌入向量")
                self.client = None
            except Exception as e:
                logger.warning(f"智谱AI SDK初始化失败: {str(e)}，将使用模拟嵌入向量")
                self.client = None

    def embed_documents(self, texts: List[str], vector_size: int = None) -> List[List[float]]:
        """
        批量嵌入文档（分批处理以避免API限制）

        Args:
            texts: 文档文本列表
            vector_size: 向量维度（默认为2048）

        Returns:
            嵌入向量列表
        """
        # 使用硬编码默认值
        if vector_size is None:
            vector_size = 2048
        if not texts:
            return []

        if self.client is None:
            # 使用模拟嵌入向量
            return [[hash(text) % 1000 / 1000.0 for _ in range(vector_size)] for text in texts]

        all_embeddings = []

        try:
            # 分批处理，避免超过API限制
            for i in range(0, len(texts), ZHIPU_EMBEDDING_BATCH_SIZE):
                batch_texts = texts[i:i + ZHIPU_EMBEDDING_BATCH_SIZE]

                # 批量调用智谱AI的嵌入服务
                response = self.client.embeddings.create(
                    model=self.embedding_model,
                    input=batch_texts
                )

                # 提取当前批次的嵌入向量
                for item in response.data:
                    all_embeddings.append(item.embedding)

                logger.info(f"已处理 {min(i + ZHIPU_EMBEDDING_BATCH_SIZE, len(texts))}/{len(texts)} 个文本块")

            return all_embeddings
        except Exception as e:
            logger.error(f"批量嵌入文档时发生错误: {str(e)}")
            # 如果批量处理失败，回退到逐个处理
            logger.info("批量处理失败，回退到逐个处理模式...")
            embeddings = []
            for text in texts:
                embedding = self.embed_query(text, vector_size)
                embeddings.append(embedding)
            return embeddings

    def embed_query(self, text: str, vector_size: int = None) -> List[float]:
        """
        嵌入单个查询文本

        Args:
            text: 查询文本
            vector_size: 向量维度（默认为2048）

        Returns:
            嵌入向量
        """
        # 使用硬编码默认值
        if vector_size is None:
            vector_size = 2048
        if self.client is None:
            # 使用模拟嵌入向量
            return [hash(text) % 1000 / 1000.0 for _ in range(vector_size)]

        try:
            response = self.client.embeddings.create(
                model=self.embedding_model,
                input=text
            )
            return response.data[0].embedding
        except Exception as e:
            logger.error(f"嵌入文本时发生错误: {str(e)}")
            # 返回模拟向量作为默认值
            return [hash(text) % 1000 / 1000.0 for _ in range(vector_size)]


class KnowledgeBaseManager:
    """个人知识库向量数据库管理器"""

    # 知识库元数据文件名
    KB_METADATA_FILE = "kb_metadata.pkl"

    def __init__(
        self,
        knowledge_base_name: str,
        workspace_root: str = None,
        embedding_model: str = "embedding-3",
        chunk_size: int = None,
        vector_size: int = None,
        similarity_threshold: float = None,
        top_k: int = None,
        check_exists: bool = True
    ):
        """
        初始化知识库管理器

        Args:
            knowledge_base_name: 知识库名称
            workspace_root: 工作区根目录，默认为当前工作区
            embedding_model: 嵌入模型
            chunk_size: 文本块大小（如果未指定，从知识库元数据加载）
            vector_size: 嵌入向量维度（默认2048）
            similarity_threshold: 相似度阈值（默认0.5）
            top_k: 默认检索的文本块数量（默认3）
            check_exists: 是否检查知识库是否存在
        """
        self.knowledge_base_name = knowledge_base_name
        # 加载配置
        config = load_config()
        # workspace_root 配置的是 MyKnowledgeBase 的路径
        kb_root = config.get("workspace_root", "d:/Nancy/MyWork/WorkBuddyWorkSpace/MyKnowledgeBase")
        self.workspace_root = workspace_root or kb_root

        # 知识库目录 (workspace_root 已指向 MyKnowledgeBase)
        self.knowledge_base_dir = os.path.join(self.workspace_root, knowledge_base_name)
        self.sourcefiles_dir = os.path.join(self.knowledge_base_dir, "sourcefiles")
        self.vectordb_dir = os.path.join(self.knowledge_base_dir, "vectordb")

        # 向量数据库文件
        self.index_file = os.path.join(self.vectordb_dir, "faiss_index.bin")
        self.metadata_file = os.path.join(self.vectordb_dir, "faiss_metadata.pkl")
        self.kb_metadata_file = os.path.join(self.vectordb_dir, self.KB_METADATA_FILE)

        # 加载或初始化知识库元数据
        self.kb_metadata = self._load_kb_metadata()

        # 如果指定了chunk_size，使用传入的值；否则从知识库元数据加载
        if chunk_size is not None:
            self.chunk_size = chunk_size
        else:
            self.chunk_size = self.kb_metadata.get("chunk_size", DEFAULT_CHUNK_SIZE)

        # 计算重叠大小：chunk_size * 0.1
        self.chunk_overlap = int(self.chunk_size * DEFAULT_CHUNK_OVERLAP_RATIO)

        # 其他参数使用传入值或默认值
        self.vector_size = vector_size if vector_size is not None else 2048
        self.similarity_threshold = similarity_threshold if similarity_threshold is not None else config.get("DEFAULT_SIMILARITY_THRESHOLD", 0.5)
        self.top_k = top_k if top_k is not None else config.get("DEFAULT_TOP_K", 3)

        # 检查知识库是否存在
        if check_exists and not os.path.exists(self.knowledge_base_dir):
            raise ValueError(f"知识库 '{knowledge_base_name}' 不存在，请先创建")

        # 初始化嵌入模型
        self.embeddings = ZhipuEmbeddings()

        # 初始化FAISS索引和元数据
        self.index = None
        self.metadata = []  # 存储文档元数据
        self._init_faiss_index()

    def _load_kb_metadata(self) -> Dict[str, Any]:
        """加载知识库元数据"""
        if os.path.exists(self.kb_metadata_file):
            try:
                with open(self.kb_metadata_file, 'rb') as f:
                    return pickle.load(f)
            except Exception as e:
                logger.warning(f"加载知识库元数据失败: {str(e)}")
        return {"chunk_size": DEFAULT_CHUNK_SIZE}

    def _save_kb_metadata(self):
        """保存知识库元数据"""
        try:
            os.makedirs(self.vectordb_dir, exist_ok=True)
            with open(self.kb_metadata_file, 'wb') as f:
                pickle.dump(self.kb_metadata, f)
            logger.info(f"知识库元数据已保存")
        except Exception as e:
            logger.error(f"保存知识库元数据失败: {str(e)}")

    @classmethod
    def create_knowledge_base(
        cls,
        knowledge_base_name: str,
        workspace_root: str = None,
        chunk_size: int = None
    ) -> Dict[str, Any]:
        """
        创建新的知识库

        Args:
            knowledge_base_name: 知识库名称
            workspace_root: 知识库根目录
            chunk_size: 文本块大小（如果未指定，使用config.txt中的DEFAULT_CHUNK_SIZE）

        Returns:
            操作结果字典
        """
        # 加载配置获取默认chunk_size
        config = load_config()
        if chunk_size is None:
            chunk_size = config.get("DEFAULT_CHUNK_SIZE", DEFAULT_CHUNK_SIZE)

        # 确定workspace_root
        kb_root = config.get("workspace_root", "d:/Nancy/MyWork/WorkBuddyWorkSpace/MyKnowledgeBase")
        workspace_root = workspace_root or kb_root

        # 创建知识库目录
        knowledge_base_dir = os.path.join(workspace_root, knowledge_base_name)
        sourcefiles_dir = os.path.join(knowledge_base_dir, "sourcefiles")
        vectordb_dir = os.path.join(knowledge_base_dir, "vectordb")

        # 检查知识库是否已存在
        if os.path.exists(knowledge_base_dir):
            return {
                "success": False,
                "message": f"知识库 '{knowledge_base_name}' 已存在"
            }

        try:
            # 创建目录结构
            os.makedirs(sourcefiles_dir, exist_ok=True)
            os.makedirs(vectordb_dir, exist_ok=True)

            # 保存知识库元数据
            kb_metadata = {
                "chunk_size": chunk_size,
                "created_at": str(os.path.getctime(knowledge_base_dir)) if os.path.exists(knowledge_base_dir) else None
            }
            kb_metadata_file = os.path.join(vectordb_dir, cls.KB_METADATA_FILE)
            with open(kb_metadata_file, 'wb') as f:
                pickle.dump(kb_metadata, f)

            # 创建空的FAISS索引
            import faiss
            import numpy as np
            index = faiss.IndexFlatL2(2048)
            index_file = os.path.join(vectordb_dir, "faiss_index.bin")
            with open(index_file, 'wb') as f:
                pickle.dump([], f)  # 空向量列表

            # 创建空的元数据文件
            metadata_file = os.path.join(vectordb_dir, "faiss_metadata.pkl")
            with open(metadata_file, 'wb') as f:
                pickle.dump([], f)

            logger.info(f"成功创建知识库 '{knowledge_base_name}'，文本块大小: {chunk_size}")

            return {
                "success": True,
                "message": f"成功创建知识库 '{knowledge_base_name}'，文本块大小: {chunk_size}",
                "knowledge_base_name": knowledge_base_name,
                "chunk_size": chunk_size
            }

        except Exception as e:
            logger.error(f"创建知识库失败: {str(e)}")
            return {
                "success": False,
                "message": f"创建知识库失败: {str(e)}"
            }

    def _init_faiss_index(self):
        """初始化FAISS索引"""
        try:
            # 确保目录存在
            os.makedirs(self.vectordb_dir, exist_ok=True)

            # 尝试加载现有的索引数据
            if os.path.exists(self.index_file) and os.path.exists(self.metadata_file):
                logger.info(f"加载现有的FAISS索引: {self.index_file}")

                # 加载向量和元数据
                with open(self.index_file, 'rb') as f:
                    vectors_data = pickle.load(f)

                # 加载元数据
                with open(self.metadata_file, 'rb') as f:
                    self.metadata = pickle.load(f)

                # 重建FAISS索引
                if len(vectors_data) > 0:
                    vectors = np.array(vectors_data).astype('float32')
                    dimension = vectors.shape[1] if len(vectors.shape) > 1 else self.vector_size
                    self.index = faiss.IndexFlatL2(dimension)
                    self.index.add(vectors)
                else:
                    self.index = faiss.IndexFlatL2(self.vector_size)
                    self.metadata = []

                logger.info(f"已加载索引，包含 {len(self.metadata)} 个向量")
            else:
                # 创建新的FAISS索引
                self.index = faiss.IndexFlatL2(self.vector_size)
                self.metadata = []
                logger.info("创建新的FAISS索引")

                # 立即保存一个空的索引
                self.save_index()

        except Exception as e:
            logger.error(f"初始化FAISS索引失败: {str(e)}")
            self.index = faiss.IndexFlatL2(self.vector_size)
            self.metadata = []

    def save_index(self):
        """保存FAISS索引和元数据"""
        try:
            # 确保目录存在
            os.makedirs(self.vectordb_dir, exist_ok=True)

            # 从FAISS索引中提取向量
            if self.index.ntotal > 0:
                vectors = self.index.reconstruct_n(0, self.index.ntotal)
                vectors_list = vectors.tolist()
            else:
                vectors_list = []

            # 保存向量数据
            with open(self.index_file, 'wb') as f:
                pickle.dump(vectors_list, f)

            # 保存元数据
            with open(self.metadata_file, 'wb') as f:
                pickle.dump(self.metadata, f)

            logger.info(f"FAISS索引已保存到 {self.index_file}")
        except Exception as e:
            logger.error(f"保存FAISS索引失败: {str(e)}")

    def _calculate_file_hash(self, file_path: str) -> str:
        """计算文件哈希值"""
        hash_sha256 = hashlib.sha256()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_sha256.update(chunk)
        return hash_sha256.hexdigest()

    def is_file_already_vectorized(self, file_path: str) -> bool:
        """检查文件是否已经向量化存储在向量库中"""
        try:
            file_hash = self._calculate_file_hash(file_path)

            for meta in self.metadata:
                if meta.get("file_hash") == file_hash:
                    return True

            return False
        except Exception as e:
            logger.error(f"检查文件是否已向量化时发生错误: {str(e)}")
            return False

    def extract_and_chunk_file(self, file_path: str) -> List[Document]:
        """
        提取文件内容并进行分块

        Args:
            file_path: 文件路径

        Returns:
            分块后的文档列表
        """
        file_ext = os.path.splitext(file_path)[1].lower()
        file_name = os.path.basename(file_path)

        try:
            if file_ext == '.pdf':
                loader = PyPDFLoader(file_path)
                documents = loader.load()
            elif file_ext in ['.doc', '.docx']:
                loader = Docx2txtLoader(file_path)
                documents = loader.load()
            elif file_ext in ['.md', '.markdown']:
                loader = UnstructuredMarkdownLoader(file_path)
                documents = loader.load()
            elif file_ext in ['.ppt', '.pptx']:
                loader = CustomPPTXLoader(file_path)
                documents = loader.load()
            elif file_ext == '.txt':
                loader = TextLoader(file_path, encoding='utf-8')
                documents = loader.load()
            else:
                raise ValueError(f"不支持的文件格式: {file_ext}")

            # 创建文本分块器
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=self.chunk_size,
                chunk_overlap=self.chunk_overlap,
                separators=["\n\n", "\n", "### "]
            )

            # 对文档进行分块
            chunks = text_splitter.split_documents(documents)

            # 添加文件名作为元数据
            file_hash = self._calculate_file_hash(file_path)
            for chunk in chunks:
                chunk.metadata["source_file"] = file_name
                chunk.metadata["file_hash"] = file_hash

            logger.info(f"文件 '{file_name}' 已提取并分块，共生成 {len(chunks)} 个文本块")
            return chunks

        except Exception as e:
            logger.error(f"提取和分块文件失败: {str(e)}")
            raise e

    def add_file_to_knowledge_base(self, file_path: str) -> Dict[str, Any]:
        """
        将文件添加到知识库

        Args:
            file_path: 文件路径

        Returns:
            操作结果字典
        """
        try:
            # 检查文件是否存在
            if not os.path.exists(file_path):
                return {"success": False, "message": f"文件不存在: {file_path}"}

            # 确保sourcefiles目录存在
            os.makedirs(self.sourcefiles_dir, exist_ok=True)

            # 复制文件到sourcefiles目录
            file_name = os.path.basename(file_path)
            dest_path = os.path.join(self.sourcefiles_dir, file_name)

            if os.path.exists(dest_path):
                return {"success": False, "message": "文件已经存在", "file_name": file_name}

            shutil.copy2(file_path, dest_path)

            # 检查文件是否已经向量化
            if self.is_file_already_vectorized(dest_path):
                return {"success": True, "message": "文件已存在，向量库中已有效数据，跳过向量化", "file_name": file_name}

            # 提取并分块文件内容
            chunks = self.extract_and_chunk_file(dest_path)

            if not chunks:
                return {"success": False, "message": "无法提取文件内容", "file_name": file_name}

            # 提取文本内容用于嵌入
            texts = [chunk.page_content for chunk in chunks]

            # 生成嵌入向量
            logger.info(f"正在生成 {len(texts)} 个文本块的嵌入向量...")
            embeddings = self.embeddings.embed_documents(texts, self.vector_size)

            # 将嵌入向量添加到FAISS索引
            embeddings_array = np.array(embeddings).astype('float32')
            self.index.add(embeddings_array)

            # 更新元数据，保存原始块索引（从1开始）和导入时间
            import datetime
            import_time = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            for i, chunk in enumerate(chunks):
                self.metadata.append({
                    "content": chunk.page_content,
                    "source_file": chunk.metadata.get("source_file", "unknown"),
                    "file_hash": chunk.metadata.get("file_hash", ""),
                    "chunk_index": i + 1,  # 原始块索引，从1开始
                    "import_time": import_time  # 导入时间
                })

            logger.info(f"成功将文件 '{file_name}' 的 {len(chunks)} 个文本块添加到知识库")

            # 保存更新后的索引
            self.save_index()

            return {
                "success": True,
                "message": f"成功添加文件 '{file_name}'，共 {len(chunks)} 个文本块",
                "file_name": file_name,
                "chunk_count": len(chunks)
            }

        except Exception as e:
            logger.error(f"添加文件到知识库失败: {str(e)}")
            return {"success": False, "message": f"添加文件失败: {str(e)}"}

    def query_similar_texts(self, query: str, top_k: int = None, similarity_threshold: float = None) -> List[Dict[str, Any]]:
        """
        查询语义相似的文本块

        Args:
            query: 查询文本
            top_k: 返回的最相似文本块数量
            similarity_threshold: 相似度过滤阈值

        Returns:
            相似文本块列表
        """
        top_k = top_k or self.top_k
        similarity_threshold = similarity_threshold or self.similarity_threshold

        try:
            if self.index.ntotal == 0:
                return []

            # 生成查询向量
            query_embedding = self.embeddings.embed_query(query, self.vector_size)
            query_array = np.array([query_embedding]).astype('float32')

            # 执行相似性搜索
            distances, indices = self.index.search(query_array, min(top_k * 2, self.index.ntotal))

            # 格式化结果并过滤低于阈值的相似度
            formatted_results = []

            for i, (dist, idx) in enumerate(zip(distances[0], indices[0])):
                if idx < len(self.metadata):
                    # 将L2距离转换为相似度分数
                    similarity_score = 1 / (1 + dist)
                    similarity_score = float(similarity_score)

                    # 只保留相似度高于阈值的结果
                    if similarity_score >= similarity_threshold:
                        meta = self.metadata[idx]
                        formatted_results.append({
                            "rank": len(formatted_results) + 1,
                            "content": meta["content"],
                            "source_file": meta["source_file"],
                            "chunk_index": meta.get("chunk_index", idx + 1),  # 原始块索引
                            "relevance_score": similarity_score
                        })

            # 确保只返回top_k个结果
            formatted_results = sorted(formatted_results, key=lambda x: x["relevance_score"], reverse=True)[:top_k]

            logger.info(f"查询完成，找到 {len(formatted_results)} 个相关文本块")
            return formatted_results

        except Exception as e:
            logger.error(f"查询相似文本失败: {str(e)}")
            return []

    def answer_question(self, query: str, top_k: int = None, similarity_threshold: float = None, use_llm: bool = True) -> Dict[str, Any]:
        """
        使用RAG技术回答问题

        Args:
            query: 用户问题
            top_k: 检索的文本块数量
            similarity_threshold: 相似度过滤阈值
            use_llm: 是否使用LLM生成答案（默认True）

        Returns:
            包含答案和相关信息的字典
        """
        top_k = top_k or self.top_k
        similarity_threshold = similarity_threshold or self.similarity_threshold

        try:
            # 从向量数据库检索相关文本块
            similar_texts = self.query_similar_texts(query, top_k, similarity_threshold)

            if not similar_texts:
                return {
                    "answer": "抱歉，未在知识库中找到相关信息。",
                    "generated_answer": None,
                    "related_texts": [],
                    "source_files": []
                }

            # 构建上下文
            context_parts = []
            source_files = set()
            for text_block in similar_texts:
                context_parts.append(f"来源: {text_block['source_file']}\n内容: {text_block['content']}")
                source_files.add(text_block['source_file'])

            context = "\n\n".join(context_parts)

            # 如果启用LLM，生成答案
            generated_answer = None
            if use_llm:
                llm = ZhipuLLM()
                if llm.client:
                    # 构建提示词
                    prompt = f"""请根据以下上下文信息回答问题。如果上下文信息足够回答问题，请给出规范、准确、可读性强的回答。如果上下文信息不足以回答问题，请明确说明。

上下文信息：
{context}

问题：{query}

请给出回答："""

                    generated_answer = llm.generate(prompt)
                    if generated_answer:
                        logger.info("LLM答案生成成功")

            # 构建文本块元数据列表（使用原始块索引）
            chunk_metadata_list = []
            for text_block in similar_texts:
                chunk_index = text_block.get("chunk_index", text_block["rank"])
                # 转换为Python原生int类型
                chunk_index = int(chunk_index) if chunk_index is not None else None
                chunk_metadata_list.append({
                    "chunk_index": chunk_index,
                    "source_file": text_block["source_file"],
                    "relevance_score": float(text_block["relevance_score"])
                })

            # 生成带元数据的答案
            answer_with_metadata = generated_answer if generated_answer else context
            if chunk_metadata_list:
                metadata_text = "\n\n--- 检索到的文本块信息 ---\n"
                for meta in chunk_metadata_list:
                    metadata_text += f"【文本块 {meta['chunk_index']}】文件: {meta['source_file']}, 相似度: {meta['relevance_score']:.4f}\n"
                answer_with_metadata = answer_with_metadata + metadata_text

            # 返回结果
            return {
                "answer": answer_with_metadata,
                "generated_answer": generated_answer,
                "related_texts": similar_texts,
                "source_files": list(source_files),
                "chunk_metadata": chunk_metadata_list
            }

        except Exception as e:
            logger.error(f"回答问题过程失败: {str(e)}")
            return {
                "answer": f"抱歉，在处理您的问题时遇到了错误: {str(e)}",
                "generated_answer": None,
                "related_texts": [],
                "source_files": [],
                "chunk_metadata": []
            }

    def list_files_in_knowledge_base(self) -> List[Dict[str, Any]]:
        """
        列出知识库中已向量化的文件

        Returns:
            文件信息列表（包含文件名、文本块数、导入日期时间）
        """
        try:
            if not self.metadata:
                return []

            # 统计每个文件的文本块数量和导入时间
            file_stats = {}
            for meta in self.metadata:
                filename = meta.get("source_file", "unknown")
                if filename not in file_stats:
                    file_stats[filename] = {
                        "filename": filename,
                        "chunk_count": 0,
                        "file_hash": meta.get("file_hash", ""),
                        "import_time": meta.get("import_time", "")
                    }
                file_stats[filename]["chunk_count"] += 1

            return list(file_stats.values())

        except Exception as e:
            logger.error(f"列出知识库文件失败: {str(e)}")
            return []

    def remove_file_from_knowledge_base(self, filename: str) -> Dict[str, Any]:
        """
        从知识库中删除指定文件

        Args:
            filename: 要删除的文件名

        Returns:
            操作结果字典
        """
        try:
            # 找到需要移除的向量索引
            indices_to_remove = []
            for i, meta in enumerate(self.metadata):
                if meta.get("source_file") == filename:
                    indices_to_remove.append(i)

            if not indices_to_remove:
                return {"success": False, "message": f"知识库中不存在文件: {filename}"}

            # 从后往前删除，避免索引变化影响
            indices_to_remove.reverse()
            for idx in indices_to_remove:
                del self.metadata[idx]

            # 重建索引
            if len(self.metadata) > 0:
                remaining_texts = [meta["content"] for meta in self.metadata]
                remaining_embeddings = self.embeddings.embed_documents(remaining_texts, self.vector_size)

                new_index = faiss.IndexFlatL2(self.vector_size)
                remaining_embeddings_array = np.array(remaining_embeddings).astype('float32')
                new_index.add(remaining_embeddings_array)

                self.index = new_index
            else:
                self.index = faiss.IndexFlatL2(self.vector_size)
                self.metadata = []

            # 保存更新后的索引
            self.save_index()

            # 删除sourcefiles目录中的文件
            source_file_path = os.path.join(self.sourcefiles_dir, filename)
            if os.path.exists(source_file_path):
                os.remove(source_file_path)

            logger.info(f"成功从知识库中移除文件: {filename}，共删除 {len(indices_to_remove)} 个文本块")

            return {
                "success": True,
                "message": f"成功删除文件 '{filename}'，共删除 {len(indices_to_remove)} 个文本块",
                "deleted_chunk_count": len(indices_to_remove)
            }

        except Exception as e:
            logger.error(f"从知识库中删除文件失败: {str(e)}")
            return {"success": False, "message": f"删除文件失败: {str(e)}"}

    def update_file_in_knowledge_base(self, file_path: str) -> Dict[str, Any]:
        """
        更新知识库中的文件（先删除旧文件，再添加新文件）

        Args:
            file_path: 文件路径

        Returns:
            操作结果字典
        """
        try:
            # 检查文件是否存在
            if not os.path.exists(file_path):
                return {"success": False, "message": f"文件不存在: {file_path}"}

            # 获取文件名
            file_name = os.path.basename(file_path)

            # 检查知识库中是否存在该文件
            existing_files = self.list_files_in_knowledge_base()
            file_exists = any(f["filename"] == file_name for f in existing_files)

            if not file_exists:
                return {"success": False, "message": f"知识库中不存在文件: {file_name}"}

            # 删除旧文件
            logger.info(f"正在删除旧文件: {file_name}")
            delete_result = self.remove_file_from_knowledge_base(file_name)
            if not delete_result.get("success"):
                return delete_result

            # 添加新文件
            logger.info(f"正在添加新文件: {file_name}")
            add_result = self.add_file_to_knowledge_base(file_path)

            if add_result.get("success"):
                return {
                    "success": True,
                    "message": f"成功更新文件 '{file_name}'",
                    "file_name": file_name,
                    "chunk_count": add_result.get("chunk_count", 0)
                }
            else:
                return add_result

        except Exception as e:
            logger.error(f"更新文件失败: {str(e)}")
            return {"success": False, "message": f"更新文件失败: {str(e)}"}

    def get_chunk_by_index(self, filename: str, chunk_index: int) -> Dict[str, Any]:
        """
        获取指定文件的第n个文本块内容

        Args:
            filename: 文件名
            chunk_index: 文本块索引（从1开始）

        Returns:
            文本块内容字典
        """
        try:
            # 找到该文件对应的所有文本块，使用原始块索引
            file_chunks = []
            for i, meta in enumerate(self.metadata):
                if meta.get("source_file") == filename:
                    # 优先使用存储的原始chunk_index，否则使用顺序编号
                    original_index = meta.get("chunk_index", len(file_chunks) + 1)
                    file_chunks.append({
                        "original_index": original_index,
                        "content": meta["content"]
                    })

            # 按原始块索引排序
            file_chunks = sorted(file_chunks, key=lambda x: x["original_index"])

            if not file_chunks:
                return {
                    "success": False,
                    "message": f"知识库中不存在文件: {filename}",
                    "filename": filename,
                    "chunk_index": chunk_index,
                    "content": None
                }

            # 检查索引是否有效
            if chunk_index < 1 or chunk_index > len(file_chunks):
                return {
                    "success": False,
                    "message": f"文件 '{filename}' 只有 {len(file_chunks)} 个文本块，索引 {chunk_index} 超出范围",
                    "filename": filename,
                    "chunk_index": chunk_index,
                    "total_chunks": len(file_chunks),
                    "content": None
                }

            # 返回指定索引的文本块
            chunk = file_chunks[chunk_index - 1]  # 转换为0-based索引

            return {
                "success": True,
                "message": f"成功获取文件 '{filename}' 的第 {chunk_index} 个文本块",
                "filename": filename,
                "chunk_index": chunk["original_index"],  # 使用原始块索引
                "total_chunks": len(file_chunks),
                "content": chunk["content"]
            }

        except Exception as e:
            logger.error(f"获取文本块失败: {str(e)}")
            return {
                "success": False,
                "message": f"获取文本块失败: {str(e)}",
                "filename": filename,
                "chunk_index": chunk_index,
                "content": None
            }


class KnowledgeBaseCatalog:
    """知识库目录管理器，用于列举知识库"""

    def __init__(self, workspace_root: str = None):
        """
        初始化知识库目录管理器

        Args:
            workspace_root: 工作区根目录（从config.txt配置读取）
        """
        # 加载配置
        config = load_config()
        # 使用传入的参数，如果未传入则从配置读取
        # workspace_root 配置的是 MyKnowledgeBase的路径
        kb_root = config.get("workspace_root", "d:/Nancy/MyWork/WorkBuddyWorkSpace/MyKnowledgeBase")
        self.workspace_root = workspace_root or kb_root
        self.knowledge_dbs_dir = self.workspace_root

    def list_knowledge_bases(self) -> List[str]:
        """
        列出所有知识库

        Returns:
            知识库名称列表
        """
        try:
            if not os.path.exists(self.knowledge_dbs_dir):
                return []

            knowledge_bases = []
            for item in os.listdir(self.knowledge_dbs_dir):
                item_path = os.path.join(self.knowledge_dbs_dir, item)
                if os.path.isdir(item_path):
                    knowledge_bases.append(item)

            return sorted(knowledge_bases)

        except Exception as e:
            logger.error(f"列出知识库失败: {str(e)}")
            return []

    def list_knowledge_bases_detailed(self) -> List[Dict[str, Any]]:
        """
        列出所有知识库的详细信息

        Returns:
            知识库详细信息列表（包含文件数、文本块大小、简单描述）
        """
        try:
            if not os.path.exists(self.knowledge_dbs_dir):
                return []

            knowledge_bases = []
            for item in os.listdir(self.knowledge_dbs_dir):
                item_path = os.path.join(self.knowledge_dbs_dir, item)
                if os.path.isdir(item_path):
                    kb_info = self._get_knowledge_base_info(item)
                    if kb_info:
                        knowledge_bases.append(kb_info)

            return sorted(knowledge_bases, key=lambda x: x["name"])

        except Exception as e:
            logger.error(f"列出知识库详细信息失败: {str(e)}")
            return []

    def _get_knowledge_base_info(self, knowledge_base_name: str) -> Dict[str, Any]:
        """获取单个知识库的详细信息"""
        try:
            knowledge_base_dir = os.path.join(self.knowledge_dbs_dir, knowledge_base_name)
            vectordb_dir = os.path.join(knowledge_base_dir, "vectordb")
            sourcefiles_dir = os.path.join(knowledge_base_dir, "sourcefiles")

            # 获取文本块大小元数据
            kb_metadata_file = os.path.join(vectordb_dir, KnowledgeBaseManager.KB_METADATA_FILE)
            chunk_size = DEFAULT_CHUNK_SIZE
            if os.path.exists(kb_metadata_file):
                try:
                    with open(kb_metadata_file, 'rb') as f:
                        kb_metadata = pickle.load(f)
                        chunk_size = kb_metadata.get("chunk_size", DEFAULT_CHUNK_SIZE)
                except Exception:
                    pass

            # 获取文件数量
            file_count = 0
            files_description = ""
            if os.path.exists(sourcefiles_dir):
                files = os.listdir(sourcefiles_dir)
                file_count = len([f for f in files if os.path.isfile(os.path.join(sourcefiles_dir, f))])
                # 生成简单描述（基于文件名）
                if files:
                    file_names = [f for f in files if os.path.isfile(os.path.join(sourcefiles_dir, f))]
                    if len(file_names) <= 3:
                        files_description = "、".join(file_names)
                    else:
                        files_description = f"{file_names[0]}、{file_names[1]}等{len(file_names)}个文件"

            return {
                "name": knowledge_base_name,
                "file_count": file_count,
                "chunk_size": chunk_size,
                "description": files_description if files_description else "空知识库"
            }

        except Exception as e:
            logger.error(f"获取知识库信息失败: {str(e)}")
            return None

    def knowledge_base_exists(self, knowledge_base_name: str) -> bool:
        """检查知识库是否存在"""
        knowledge_base_dir = os.path.join(self.knowledge_dbs_dir, knowledge_base_name)
        return os.path.isdir(knowledge_base_dir)


# ==================== 便捷函数 ====================

def create_knowledge_base_manager(
    knowledge_base_name: str,
    workspace_root: str = None,
    embedding_model: str = None,
    chunk_size: int = None,
    vector_size: int = None,
    similarity_threshold: float = None,
    top_k: int = None,
    check_exists: bool = True
) -> KnowledgeBaseManager:
    """
    创建知识库管理器的便捷函数

    配置参数优先从config.txt配置文件读取，如果配置文件不存在或未设置则使用默认值。

    Args:
        knowledge_base_name: 知识库名称
        workspace_root: 工作区根目录（从配置读取）
        embedding_model: 嵌入模型类型（从配置读取）
        chunk_size: 文本块大小（如果未指定，从知识库元数据加载）
        vector_size: 嵌入向量维度（默认2048）
        similarity_threshold: 相似度阈值（从配置读取）
        top_k: 默认检索的文本块数量（从配置读取）
        check_exists: 是否检查知识库是否存在（默认True）

    Returns:
        KnowledgeBaseManager实例
    """
    # 加载配置
    config = load_config()

    # 使用传入的参数，如果未传入则从配置读取
    workspace_root = workspace_root or config.get("workspace_root", _DEFAULT_CONFIG["workspace_root"])
    embedding_model = embedding_model or config.get("embedding_model", "embedding-3")
    # chunk_size不在这里设置，由KnowledgeBaseManager从知识库元数据加载
    vector_size = vector_size if vector_size is not None else 2048
    similarity_threshold = similarity_threshold if similarity_threshold is not None else config.get("DEFAULT_SIMILARITY_THRESHOLD", 0.5)
    top_k = top_k if top_k is not None else config.get("DEFAULT_TOP_K", 3)

    return KnowledgeBaseManager(
        knowledge_base_name=knowledge_base_name,
        workspace_root=workspace_root,
        embedding_model=embedding_model,
        chunk_size=chunk_size,
        vector_size=vector_size,
        similarity_threshold=similarity_threshold,
        top_k=top_k,
        check_exists=check_exists
    )


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="个人知识库向量数据库管理工具")
    parser.add_argument("--action", choices=["list-kb", "add", "query", "list-files", "remove"],
                       required=True, help="操作类型")
    parser.add_argument("--kb-name", help="知识库名称")
    parser.add_argument("--file-path", help="文件路径（添加操作时使用）")
    parser.add_argument("--query-text", help="查询文本（查询操作时使用）")
    parser.add_argument("--filename", help="文件名（删除操作时使用）")

    args = parser.parse_args()

    if args.action == "list-kb":
        catalog = KnowledgeBaseCatalog()
        kb_list = catalog.list_knowledge_bases()
        print("知识库列表:", kb_list if kb_list else "暂无知识库")

    elif args.action == "add":
        if not args.kb_name or not args.file_path:
            print("错误: 添加操作需要指定 --kb-name 和 --file-path 参数")
        else:
            manager = create_knowledge_base_manager(args.kb_name)
            result = manager.add_file_to_knowledge_base(args.file_path)
            print(result)

    elif args.action == "query":
        if not args.kb_name or not args.query_text:
            print("错误: 查询操作需要指定 --kb-name 和 --query-text 参数")
        else:
            manager = create_knowledge_base_manager(args.kb_name)
            result = manager.answer_question(args.query_text)
            print(f"答案: {result['answer']}")
            print(f"相关文件: {result['source_files']}")

    elif args.action == "list-files":
        if not args.kb_name:
            print("错误: 列出文件操作需要指定 --kb-name 参数")
        else:
            manager = create_knowledge_base_manager(args.kb_name)
            files = manager.list_files_in_knowledge_base()
            print("知识库文件列表:")
            for f in files:
                print(f"  - {f['filename']}: {f['chunk_count']} 个文本块")

    elif args.action == "remove":
        if not args.kb_name or not args.filename:
            print("错误: 删除操作需要指定 --kb-name 和 --filename 参数")
        else:
            manager = create_knowledge_base_manager(args.kb_name)
            result = manager.remove_file_from_knowledge_base(args.filename)
            print(result)
