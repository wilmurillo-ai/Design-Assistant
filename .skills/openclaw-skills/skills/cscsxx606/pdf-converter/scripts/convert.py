#!/usr/bin/env python3
"""
PDF 转 PPTX/Word 转换器

支持将 PDF 文件转换为 PowerPoint (PPTX) 或 Word (DOCX) 格式。
自动处理文件大小限制，支持批量转换。

使用方法:
    python3 convert.py <pdf 文件> [--format pptx|docx] [--max-size 30] [--dpi 150]
"""

import argparse
import glob
import os
import sys
import zipfile
import shutil
from pathlib import Path

try:
    from pdf2image import convert_from_path
    from pptx import Presentation
    from pptx.util import Inches
    from pdf2docx import Converter
    from PIL import Image
except ImportError as e:
    print(f"❌ 缺少依赖库：{e}")
    print("\n请安装依赖:")
    print("  pip3 install --break-system-packages pdf2image python-pptx pdf2docx Pillow")
    print("\n还需要安装 poppler:")
    print("  macOS: brew install poppler")
    print("  Ubuntu: sudo apt-get install poppler-utils")
    sys.exit(1)


class PDFConverter:
    """PDF 转换器类"""
    
    def __init__(self, dpi=150, quality=85, max_size_mb=30, quiet=False):
        self.dpi = dpi
        self.quality = quality
        self.max_size_mb = max_size_mb
        self.quiet = quiet
    
    def log(self, message):
        """日志输出"""
        if not self.quiet:
            print(message)
    
    def convert_to_pptx(self, pdf_path, output_path=None):
        """转换 PDF 为 PPTX"""
        if not os.path.exists(pdf_path):
            print(f"❌ 文件不存在：{pdf_path}")
            return False
        
        if not output_path:
            output_path = os.path.splitext(pdf_path)[0] + ".pptx"
        
        self.log(f"📄 源文件：{os.path.basename(pdf_path)}")
        self.log(f"📊 源大小：{os.path.getsize(pdf_path)/(1024*1024):.2f} MB")
        self.log(f"🔄 开始转换 (PPTX)...")
        
        try:
            # 将 PDF 转为图片
            self.log("📸 渲染 PDF 页面为图片...")
            images = convert_from_path(pdf_path, dpi=self.dpi)
            self.log(f"✓ 共 {len(images)} 页")
            
            # 创建 PPT
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            prs.slides._sldIdLst.clear()
            
            self.log("📊 创建幻灯片...")
            for i, img in enumerate(images):
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                
                # 保存图片
                img_path = f"/tmp/slide_{i}.jpg"
                img.save(img_path, "JPEG", quality=self.quality, optimize=True)
                
                # 插入图片到幻灯片
                slide.shapes.add_picture(img_path, 0, 0, 
                                        width=prs.slide_width, 
                                        height=prs.slide_height)
                os.remove(img_path)
                
                if (i+1) % 20 == 0:
                    self.log(f"  已处理 {i+1}/{len(images)} 页")
            
            # 保存 PPTX
            prs.save(output_path)
            
            # 检查结果
            if os.path.exists(output_path):
                size_mb = os.path.getsize(output_path)/(1024*1024)
                self.log(f"\n✅ 转换成功！")
                self.log(f"📊 输出大小：{size_mb:.2f} MB")
                
                if size_mb > self.max_size_mb:
                    self.log(f"⚠️  大小限制：≤{self.max_size_mb} MB - ❌ 超出")
                    self.log(f"💡 提示：尝试降低 --dpi 或 --quality 参数")
                else:
                    self.log(f"✓ 大小限制：≤{self.max_size_mb} MB - ✅ 符合")
                
                self.log(f"📍 保存位置：{output_path}")
                return True
            else:
                self.log("❌ 转换失败：输出文件不存在")
                return False
                
        except Exception as e:
            self.log(f"❌ 转换失败：{e}")
            return False
    
    def convert_to_docx(self, pdf_path, output_path=None):
        """转换 PDF 为 DOCX"""
        if not os.path.exists(pdf_path):
            print(f"❌ 文件不存在：{pdf_path}")
            return False
        
        if not output_path:
            output_path = os.path.splitext(pdf_path)[0] + ".docx"
        
        self.log(f"📄 源文件：{os.path.basename(pdf_path)}")
        self.log(f"📊 源大小：{os.path.getsize(pdf_path)/(1024*1024):.2f} MB")
        self.log(f"🔄 开始转换 (DOCX)...")
        
        try:
            cv = Converter(pdf_path)
            cv.convert(output_path)
            cv.close()
            
            if os.path.exists(output_path):
                size_mb = os.path.getsize(output_path)/(1024*1024)
                self.log(f"\n✅ 转换成功！")
                self.log(f"📊 输出大小：{size_mb:.2f} MB")
                
                if size_mb > self.max_size_mb:
                    self.log(f"⚠️  大小限制：≤{self.max_size_mb} MB - ❌ 超出")
                else:
                    self.log(f"✓ 大小限制：≤{self.max_size_mb} MB - ✅ 符合")
                
                self.log(f"📍 保存位置：{output_path}")
                return True
            else:
                self.log("❌ 转换失败：输出文件不存在")
                return False
                
        except Exception as e:
            self.log(f"❌ 转换失败：{e}")
            return False
    
    def convert(self, pdf_path, output_format='pptx', output_path=None):
        """转换 PDF 文件"""
        if output_format.lower() == 'pptx':
            return self.convert_to_pptx(pdf_path, output_path)
        elif output_format.lower() == 'docx':
            return self.convert_to_docx(pdf_path, output_path)
        else:
            self.log(f"❌ 不支持的格式：{output_format}")
            return False
    
    def batch_convert(self, pdf_dir, output_format='pptx'):
        """批量转换目录中的所有 PDF"""
        pdf_files = glob.glob(os.path.join(pdf_dir, "*.pdf"))
        
        if not pdf_files:
            self.log(f"❌ 目录中没有 PDF 文件：{pdf_dir}")
            return False
        
        self.log(f"📁 找到 {len(pdf_files)} 个 PDF 文件")
        self.log(f"🔄 开始批量转换 ({output_format.upper()})...\n")
        
        success_count = 0
        for i, pdf_file in enumerate(pdf_files, 1):
            self.log(f"[{i}/{len(pdf_files)}] ", end="")
            if self.convert(pdf_file, output_format):
                success_count += 1
            self.log("")
        
        self.log(f"\n✅ 批量转换完成！")
        self.log(f"📊 总计：{len(pdf_files)} 个文件")
        self.log(f"✅ 成功：{success_count} 个")
        self.log(f"❌ 失败：{len(pdf_files) - success_count} 个")
        
        return success_count == len(pdf_files)


def main():
    parser = argparse.ArgumentParser(
        description='PDF 转 PPTX/Word 转换器',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例:
  python3 convert.py file.pdf                          # 转换为 PPTX
  python3 convert.py file.pdf --format docx            # 转换为 Word
  python3 convert.py file.pdf --max-size 20            # 限制大小 20MB
  python3 convert.py file.pdf --dpi 100 --quality 70   # 降低质量减小大小
  python3 convert.py /path/to/pdfs/ --batch            # 批量转换
        """
    )
    
    parser.add_argument('input', help='PDF 文件路径或目录（批量模式）')
    parser.add_argument('--format', choices=['pptx', 'docx'], default='pptx',
                       help='输出格式 (默认：pptx)')
    parser.add_argument('--max-size', type=float, default=30,
                       help='最大文件大小 MB (默认：30)')
    parser.add_argument('--output', '-o', help='输出文件路径')
    parser.add_argument('--dpi', type=int, default=150,
                       help='渲染 DPI (默认：150)')
    parser.add_argument('--quality', type=int, default=85,
                       help='JPEG 图片质量 1-100 (默认：85)')
    parser.add_argument('--batch', action='store_true',
                       help='批量转换目录中的所有 PDF')
    parser.add_argument('--quiet', '-q', action='store_true',
                       help='静默模式，仅输出结果')
    
    args = parser.parse_args()
    
    # 创建转换器
    converter = PDFConverter(
        dpi=args.dpi,
        quality=args.quality,
        max_size_mb=args.max_size,
        quiet=args.quiet
    )
    
    # 检查输入
    if not os.path.exists(args.input):
        print(f"❌ 文件/目录不存在：{args.input}")
        sys.exit(1)
    
    # 批量模式或单文件模式
    if args.batch or os.path.isdir(args.input):
        success = converter.batch_convert(args.input, args.format)
    else:
        success = converter.convert(args.input, args.format, args.output)
    
    sys.exit(0 if success else 1)


if __name__ == "__main__":
    main()
