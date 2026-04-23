"""
多格式文档提取器 - 支持 Excel、PDF、PPTX、DOCX 等格式
"""
from typing import Dict, Any
import os
import pandas as pd
import pdfplumber
from pptx import Presentation
from docx import Document


class DocumentExtractor:
    def __init__(self, parser):
        self.parser = parser

    def extract(self, file_path: str) -> Dict[str, Any]:
        """根据文件类型调用对应提取方法"""
        file_type = self._detect_file_type(file_path, "")

        extractors = {
            "excel": self.extract_excel,
            "pdf": self.extract_pdf,
            "pptx": self.extract_pptx,
            "docx": self.extract_docx,
            "image": self.extract_image
        }

        extractor = extractors.get(file_type)
        if extractor:
            content = extractor(file_path)
            return {
                "content": content,
                "file_type": file_type,
                "file_name": os.path.basename(file_path)
            }

        return {
            "content": f"不支持的文件类型: {file_type}",
            "file_type": file_type,
            "file_name": os.path.basename(file_path)
        }

    def extract_excel(self, file_path: str) -> str:
        """提取 Excel 内容为 Markdown"""
        try:
            excel_file = pd.ExcelFile(file_path)
            result_parts = []

            for sheet_name in excel_file.sheet_names:
                df = excel_file.parse(sheet_name)
                df = self._normalize_columns(df)
                table_md = self._dataframe_to_markdown(df)
                result_parts.append(f"## {sheet_name}\n\n{table_md}\n")

            return "\n".join(result_parts) if result_parts else "Excel 文件为空或无法读取"
        except Exception as e:
            return f"提取 Excel 文件失败: {str(e)}"

    def extract_pdf(self, file_path: str) -> str:
        """提取 PDF 内容为 Markdown"""
        try:
            result_parts = []
            with pdfplumber.open(file_path) as pdf:
                for i, page in enumerate(pdf.pages, 1):
                    text = page.extract_text()
                    if text:
                        result_parts.append(f"## 第 {i} 页\n\n{text}\n")

            return "\n".join(result_parts) if result_parts else "PDF 文件为空或无法读取"
        except Exception as e:
            return f"提取 PDF 文件失败: {str(e)}"

    def extract_pptx(self, file_path: str) -> str:
        """提取 PPTX 内容为 Markdown"""
        try:
            prs = Presentation(file_path)
            result_parts = []

            for i, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())

                if slide_texts:
                    result_parts.append(f"## 幻灯片 {i}\n\n" + "\n".join(slide_texts) + "\n")

            return "\n".join(result_parts) if result_parts else "PPTX 文件为空或无法读取"
        except Exception as e:
            return f"提取 PPTX 文件失败: {str(e)}"

    def extract_docx(self, file_path: str) -> str:
        """提取 DOCX 内容为 Markdown"""
        try:
            doc = Document(file_path)
            result_parts = []

            for para in doc.paragraphs:
                if para.text.strip():
                    result_parts.append(para.text.strip())

            for table in doc.tables:
                table_text = self._extract_table_to_markdown(table)
                if table_text:
                    result_parts.append(table_text)

            return "\n".join(result_parts) if result_parts else "DOCX 文件为空或无法读取"
        except Exception as e:
            return f"提取 DOCX 文件失败: {str(e)}"

    def extract_image(self, file_path: str) -> str:
        """提取图片信息（基础实现）"""
        try:
            file_size = os.path.getsize(file_path)
            file_name = os.path.basename(file_path)
            return f"图片文件: {file_name}\n文件大小: {file_size} 字节\n\n注意: 图片内容提取需要 OCR 支持，请使用专业的 OCR 工具处理图片内容。"
        except Exception as e:
            return f"提取图片信息失败: {str(e)}"

    def _normalize_columns(self, df: pd.DataFrame) -> pd.DataFrame:
        """标准化列名"""
        df.columns = [str(col).strip() for col in df.columns]
        return df

    def _dataframe_to_markdown(self, df: pd.DataFrame) -> str:
        """DataFrame 转 Markdown"""
        if df.empty:
            return "数据为空"

        headers = df.columns.tolist()
        header_row = "| " + " | ".join(headers) + " |"
        separator_row = "| " + " | ".join(["---"] * len(headers)) + " |"

        data_rows = []
        for _, row in df.iterrows():
            row_values = [str(val) for val in row.values]
            data_rows.append("| " + " | ".join(row_values) + " |")

        return "\n".join([header_row, separator_row] + data_rows)

    def _detect_file_type(self, filename: str, content: str) -> str:
        """检测文件类型"""
        ext = os.path.splitext(filename)[1].lower()

        type_mapping = {
            ".xlsx": "excel",
            ".xls": "excel",
            ".pdf": "pdf",
            ".pptx": "pptx",
            ".docx": "docx",
            ".doc": "docx",
            ".png": "image",
            ".jpg": "image",
            ".jpeg": "image",
            ".gif": "image",
            ".bmp": "image"
        }

        return type_mapping.get(ext, "unknown")

    def _extract_table_to_markdown(self, table) -> str:
        """提取 DOCX 表格为 Markdown"""
        try:
            rows = table.rows
            if not rows:
                return ""

            result_parts = []
            for i, row in enumerate(rows):
                cells = [cell.text.strip() for cell in row.cells]
                row_str = "| " + " | ".join(cells) + " |"
                result_parts.append(row_str)

                if i == 0:
                    separator = "| " + " | ".join(["---"] * len(cells)) + " |"
                    result_parts.append(separator)

            return "\n".join(result_parts)
        except Exception:
            return ""
