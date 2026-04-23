"""
Skill 工具函数模块
"""
import os
import sys
import json
import re
from typing import Dict, List, Any, Optional
from pathlib import Path
from datetime import datetime

sys.path.insert(0, str(Path(__file__).parent.parent))
from src.config import SkillConfig

import pandas as pd

COLUMN_PATTERNS = {
    "保单号": ["保单号", "单号", "保单流水号"],
    "保险公司": ["保险公司", "公司"],
    "投保人": ["投保人"],
    "被保险人": ["被保险人", "被保人"],
    "受益人": ["受益人", "身故保险金受益人"],
    "产品名称": ["险种名称", "险种", "产品名称", "产品", "险种名称（主险）"],
    "保额": ["保障额度", "保额", "基本保额", "保险金额", "保障"],
    "缴费年限": ["交费期间", "交费年期", "缴费期间", "缴费年期", "年期"],
    "保费": ["险种保费", "期交保费", "年交保费", "保费", "保费（元）"],
    "生效日期": ["生效日", "生效日期", "保单生效日", "承保日期"],
    "保险期间": ["保险期间", "保障期间"],
    "销售渠道": ["销售渠道", "渠道"],
    "保单状态": ["保单状态", "状态"],
}


def _normalize_column_name(col_name: str) -> str:
    """标准化列名"""
    if not col_name or pd.isna(col_name):
        return ""
    col_name = str(col_name).strip()
    col_name = col_name.replace("\n", " ").replace("\r", " ")
    col_name = " ".join(col_name.split())

    for normalized, patterns in COLUMN_PATTERNS.items():
        for pattern in patterns:
            if pattern in col_name:
                return normalized
    return col_name


def _detect_company(text: str) -> str:
    """识别保险公司"""
    if not text:
        return ""
    text = str(text).replace("\n", " ")

    companies = {
        "中美联泰大都会人寿": ["大都会", "都会", "中美联泰"],
        "中国人寿": ["中国人寿", "国寿"],
        "平安": ["平安"],
        "太平洋": ["太平洋", "太保"],
        "新华": ["新华"],
        "华夏": ["华夏"],
        "泰康": ["泰康"],
        "信泰": ["信泰"],
        "横琴": ["横琴"],
        "爱心": ["爱心"],
    }

    for company, keywords in companies.items():
        for kw in keywords:
            if kw in text:
                return company
    return ""


def _extract_amount(text: str) -> str:
    """提取金额"""
    if not text:
        return ""
    text = str(text).replace("\n", " ")

    match = re.search(r"(\d+(?:\.\d+)?)\s*[万MWmw]", text)
    if match:
        return f"{match.group(1)}万"

    match = re.search(r"^[\d,]+(?:\.\d+)?$", text.strip())
    if match:
        num = match.group().replace(",", "")
        if float(num) > 0:
            return num

    return text.strip()


def _parse_value(value: Any) -> str:
    """解析单元格值"""
    if pd.isna(value):
        return ""
    if isinstance(value, datetime):
        return value.strftime("%Y-%m-%d")
    return str(value).replace("\n", " ").strip()


def _find_header_row(df) -> int:
    """查找列标题行"""
    for idx in range(min(5, len(df))):
        row = df.iloc[idx]
        row_str = " ".join([str(v) for v in row if pd.notna(v)])

        keywords = ["投保人", "被保险人", "险种名称", "险种", "保额", "保障额度", "保费", "生效", "保险公司"]
        match_count = sum(1 for kw in keywords if kw in row_str)

        if match_count >= 2:
            return idx

    return 0


def extract_single_document(file_path: str, config: SkillConfig) -> Dict:
    """
    提取单个文档并转换为结构化 JSON

    Args:
        file_path: 文档文件路径
        config: 技能配置

    Returns:
        结构化 JSON 数据
    """
    if not os.path.exists(file_path):
        return {"error": f"文件不存在: {file_path}"}

    ext = Path(file_path).suffix.lower()
    content = ""

    try:
        if ext in [".xlsx", ".xls"]:
            content = _extract_excel(file_path)
        elif ext in [".pdf"]:
            content = _extract_pdf(file_path)
        elif ext in [".docx", ".doc"]:
            content = _extract_docx(file_path)
        elif ext in [".pptx"]:
            content = _extract_pptx(file_path)
        elif ext in [".txt", ".md"]:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
        else:
            return {"error": f"不支持的文件格式: {ext}"}

        markdown = _convert_to_markdown(content, ext)
        json_data = _convert_markdown_to_json(markdown, config)

        return {
            "file_path": file_path,
            "markdown": markdown,
            "structured_data": json_data,
            "status": "success"
        }

    except Exception as e:
        return {
            "file_path": file_path,
            "error": str(e),
            "status": "error"
        }


def extract_customer_folder(folder_path: str, config: SkillConfig) -> List[Dict]:
    """
    批量提取客户文件夹中的所有文档

    Args:
        folder_path: 客户文件夹路径
        config: 技能配置

    Returns:
        该客户的 Markdown 档案列表
    """
    if not os.path.exists(folder_path):
        return [{"error": f"文件夹不存在: {folder_path}"}]

    if not os.path.isdir(folder_path):
        return [{"error": f"路径不是文件夹: {folder_path}"}]

    supported_extensions = [".xlsx", ".xls", ".pdf", ".docx", ".doc", ".pptx", ".txt", ".md"]
    results = []

    for root, dirs, files in os.walk(folder_path):
        for file in files:
            ext = Path(file).suffix.lower()
            if ext in supported_extensions:
                file_path = os.path.join(root, file)
                result = extract_single_document(file_path, config)
                result["relative_path"] = os.path.relpath(file_path, folder_path)
                results.append(result)

    return results


def merge_customer_profiles(name: str, profiles: List[Dict]) -> str:
    """
    合并多个客户档案为单一 Markdown 文件

    Args:
        name: 客户姓名
        profiles: 多个档案字典列表

    Returns:
        合并后的 Markdown 档案字符串
    """
    merged = f"# 客户档案汇总 - {name}\n\n"
    merged += f"**合并时间**: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n"
    merged += "---\n\n"

    for i, profile in enumerate(profiles, 1):
        merged += f"## 档案 {i}\n\n"

        if "error" in profile:
            merged += f"**错误**: {profile['error']}\n\n"
            continue

        if "file_path" in profile:
            merged += f"**来源文件**: {profile['file_path']}\n\n"

        if "markdown" in profile:
            merged += profile["markdown"]
        elif "structured_data" in profile:
            merged += "```json\n"
            merged += json.dumps(profile["structured_data"], ensure_ascii=False, indent=2)
            merged += "\n```\n"

        merged += "\n---\n\n"

    return merged


def query_customer_profile(profiles_dir: str, query: str, config: SkillConfig) -> str:
    """
    在客户档案目录中进行语义搜索

    Args:
        profiles_dir: 档案目录路径
        query: 查询意图
        config: 技能配置

    Returns:
        相关客户信息
    """
    if not os.path.exists(profiles_dir):
        return f"错误: 目录不存在 - {profiles_dir}"

    markdown_files = []
    for root, dirs, files in os.walk(profiles_dir):
        for file in files:
            if file.endswith(".md"):
                file_path = os.path.join(root, file)
                with open(file_path, "r", encoding="utf-8") as f:
                    content = f.read()
                    markdown_files.append({
                        "path": file_path,
                        "content": content
                    })

    if not markdown_files:
        return "未找到任何 Markdown 档案文件"

    relevant_content = _semantic_search(markdown_files, query, config)
    summary = _llm_summarize(query, relevant_content, config)

    return summary


def update_customer_profile(profile_path: str, updates: Dict) -> str:
    """
    更新客户档案内容

    Args:
        profile_path: 档案文件路径
        updates: 更新内容字典

    Returns:
        更新后的档案内容
    """
    if not os.path.exists(profile_path):
        return f"错误: 档案文件不存在 - {profile_path}"

    with open(profile_path, "r", encoding="utf-8") as f:
        content = f.read()

    if profile_path.endswith(".json"):
        data = json.loads(content) if content.strip() else {}
        data.update(updates)
        result = json.dumps(data, ensure_ascii=False, indent=2)
        with open(profile_path, "w", encoding="utf-8") as f:
            f.write(result)
        return result
    elif profile_path.endswith(".md"):
        lines = content.split("\n")
        for key, value in updates.items():
            header = f"**{key}**:"
            found = False
            for i, line in enumerate(lines):
                if line.startswith(header):
                    lines[i] = f"{header} {value}"
                    found = True
                    break
            if not found:
                lines.append(f"\n{header} {value}")

        result = "\n".join(lines)
        with open(profile_path, "w", encoding="utf-8") as f:
            f.write(result)
        return result
    else:
        return "错误: 不支持的文件格式"


def generate_customer_report(customer_name: str, profiles_dir: str, config: SkillConfig) -> str:
    """
    生成客户保障分析报告

    Args:
        customer_name: 客户姓名
        profiles_dir: 档案目录路径
        config: 技能配置

    Returns:
        保障分析报告字符串
    """
    customer_folders = []
    for item in os.listdir(profiles_dir):
        item_path = os.path.join(profiles_dir, item)
        if os.path.isdir(item_path) and customer_name in item:
            customer_folders.append(item_path)

    if not customer_folders:
        return f"未找到客户 {customer_name} 的档案目录"

    customer_folder = customer_folders[0]
    profiles = extract_customer_folder(customer_folder, config)

    all_content = []
    for profile in profiles:
        if "error" not in profile and "markdown" in profile:
            all_content.append(profile["markdown"])

    combined_content = "\n\n".join(all_content)
    report = _analyze_insurance_gaps(combined_content, config)
    recommendations = _recommend_products(combined_content, config)

    full_report = f"# 保障分析报告 - {customer_name}\n\n"
    full_report += f"**生成时间**: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n"
    full_report += "---\n\n"
    full_report += "## 现有保障分析\n\n"
    full_report += report + "\n\n"
    full_report += "## 保障缺口\n\n"
    full_report += _identify_gaps(combined_content, config) + "\n\n"
    full_report += "## 推荐产品\n\n"
    full_report += recommendations + "\n\n"

    return full_report


def _extract_excel(file_path: str) -> str:
    """从 Excel 文件提取文本内容"""
    try:
        excel_file = pd.ExcelFile(file_path)
        results = []

        for sheet_name in excel_file.sheet_names:
            try:
                df_raw = pd.read_excel(file_path, sheet_name=sheet_name, header=None, engine="openpyxl")

                if df_raw.empty:
                    continue

                header_row = _find_header_row(df_raw)
                df = pd.read_excel(file_path, sheet_name=sheet_name, header=header_row, engine="openpyxl")

                original_columns = list(df.columns)
                normalized_columns = [_normalize_column_name(col) for col in df.columns]
                df.columns = normalized_columns

                sheet_content = [f"### 工作表: {sheet_name}", ""]

                for idx, row in df.iterrows():
                    row_data = {}
                    for col_idx, (orig_col, norm_col) in enumerate(zip(original_columns, normalized_columns)):
                        if col_idx >= len(row):
                            continue
                        value = row.iloc[col_idx]
                        if norm_col and not pd.isna(value) and str(value).strip():
                            row_data[norm_col] = _parse_value(value)

                    if row_data:
                        for key, val in row_data.items():
                            sheet_content.append(f"- **{key}**: {val}")
                        sheet_content.append("")

                if len(sheet_content) > 2:
                    results.append("\n".join(sheet_content))

            except Exception as e:
                results.append(f"### 工作表: {sheet_name}\n\n错误: {str(e)}")

        return "\n---\n".join(results) if results else ""

    except ImportError:
        return "错误: 未安装 pandas/openpyxl，请运行 pip install pandas openpyxl"
    except Exception as e:
        return f"错误: {str(e)}"


def _extract_pdf(file_path: str) -> str:
    """从 PDF 文件提取文本"""
    text_parts = []

    try:
        import fitz
        doc = fitz.open(file_path)
        for page_num, page in enumerate(doc, 1):
            text = page.get_text()
            if text.strip():
                text_parts.append(f"=== 第{page_num}页 ===\n{text}")
        doc.close()
        return "\n\n".join(text_parts) if text_parts else ""

    except ImportError:
        pass
    except Exception:
        pass

    try:
        import pdfplumber
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text:
                    text_parts.append(f"=== 第{page_num}页 ===\n{text}")
        return "\n\n".join(text_parts) if text_parts else ""

    except ImportError:
        return "错误: 未安装 PDF 处理库，请运行 pip install pdfplumber 或 pip install PyMuPDF"
    except Exception as e:
        return f"错误: {str(e)}"


def _extract_docx(file_path: str) -> str:
    """从 DOCX 文件提取文本"""
    try:
        from docx import Document

        doc = Document(file_path)
        parts = []

        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)

        for table in doc.tables:
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                if any(row_data):
                    parts.append(" | ".join(row_data))

        return "\n".join(parts)

    except ImportError:
        return "错误: 未安装 python-docx，请运行 pip install python-docx"
    except Exception as e:
        return f"错误: {str(e)}"


def _extract_pptx(file_path: str) -> str:
    """从 PPTX 文件提取文本"""
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

            if slide_text:
                parts.append(f"=== 第{slide_num}页 ===\n" + "\n".join(slide_text))

        return "\n\n".join(parts) if parts else ""

    except ImportError:
        return "错误: 未安装 python-pptx，请运行 pip install python-pptx"
    except Exception as e:
        return f"错误: {str(e)}"


def _convert_to_markdown(content: str, ext: str) -> str:
    """转换为 Markdown 格式"""
    if ext in [".md", ".txt"]:
        return content
    return content


def _convert_markdown_to_json(markdown: str, config: SkillConfig) -> Dict:
    """使用 LLM 将 Markdown 转换为结构化 JSON"""
    if not markdown or len(markdown.strip()) < 10:
        return {
            "raw_text": markdown[:500] if markdown else "",
            "tables": [],
            "lists": [],
            "headings": []
        }

    try:
        from src.llm_client import LLMClient
        client = LLMClient(config)

        prompt = f"""从以下客户档案文本中提取结构化信息，返回 JSON 格式：

文本内容：
{markdown[:3000]}

请提取：姓名、手机号、险种、保额、保费、生效日期、保险公司等关键信息。
如果无法提取某项，设为空字符串。

JSON格式：
"""

        result = client.extract_info(prompt)
        return result

    except Exception:
        return {
            "raw_text": markdown[:500],
            "tables": [],
            "lists": [],
            "headings": []
        }


def _semantic_search(documents: List[Dict], query: str, config: SkillConfig) -> str:
    """改进的语义搜索"""
    if not documents:
        return "未找到相关内容"

    query_keywords = query.lower().split()
    scored_docs = []

    for doc in documents:
        content_lower = doc["content"].lower()
        score = 0

        for keyword in query_keywords:
            if keyword in content_lower:
                score += 1

        if score > 0:
            scored_docs.append((score, doc["content"][:2000]))

    scored_docs.sort(key=lambda x: x[0], reverse=True)
    relevant = [content for score, content in scored_docs[:3]]

    return "\n\n".join(relevant) if relevant else "未找到相关内容"


def _llm_summarize(query: str, content: str, config: SkillConfig) -> str:
    """使用 LLM 进行语义总结"""
    if not content or content == "未找到相关内容":
        return "未找到与查询相关的信息。"

    try:
        from src.llm_client import LLMClient
        client = LLMClient(config)

        prompt = f"""基于以下客户档案内容，回答用户问题。

用户问题：{query}

档案内容：
{content[:3000]}

请用自然语言回答用户的问题。"""

        response = client.chat([
            {"role": "user", "content": prompt}
        ])
        return response

    except Exception as e:
        return f"根据查询「{query}」的相关信息：\n\n{content[:500]}..."


def _analyze_insurance_gaps(content: str, config: SkillConfig) -> str:
    """使用 LLM 分析保障缺口"""
    if not content:
        return "未找到客户保障信息"

    try:
        from src.llm_client import LLMClient
        client = LLMClient(config)

        prompt = f"""分析以下客户档案的保障缺口。

档案内容：
{content[:3000]}

请分析：
1. 现有保障有哪些（重疾险、医疗险、寿险等）
2. 各险种的保额是多少
3. 保障缺口在哪里
4. 建议补充什么保障

请用简洁的 Markdown 格式回答。"""

        response = client.chat([
            {"role": "user", "content": prompt}
        ])
        return response

    except Exception:
        return "保障分析需要 LLM 接口支持"


def _recommend_products(content: str, config: SkillConfig) -> str:
    """使用 LLM 推荐产品"""
    if not content:
        return "未找到客户保障信息"

    try:
        from src.llm_client import LLMClient
        client = LLMClient(config)

        prompt = f"""基于以下客户档案，推荐适合的保险产品。

档案内容：
{content[:3000]}

请根据客户现有保障情况和保障缺口，推荐：
1. 适合的产品类型
2. 建议的保额
3. 建议的缴费年限

请用简洁的 Markdown 格式回答。"""

        response = client.chat([
            {"role": "user", "content": prompt}
        ])
        return response

    except Exception:
        return "产品推荐需要 LLM 接口支持"


def _identify_gaps(content: str, config: SkillConfig) -> str:
    """识别保障缺口"""
    if not content:
        return "未找到客户保障信息"

    try:
        from src.llm_client import LLMClient
        client = LLMClient(config)

        prompt = f"""从以下客户档案中识别保障缺口。

档案内容：
{content[:3000]}

请识别：
1. 哪些险种缺失
2. 保额是否充足
3. 是否有待确认的信息

请用简洁的 Markdown 格式回答。"""

        response = client.chat([
            {"role": "user", "content": prompt}
        ])
        return response

    except Exception:
        return "保障缺口识别需要 LLM 接口支持"


def calculate_renewal_dates(policies: List[Dict]) -> List[Dict]:
    """计算保单续费日期

    Args:
        policies: 保单列表

    Returns:
        包含续费日期的保单列表
    """
    today = datetime.now()
    results = []

    for policy in policies:
        effective_date_str = policy.get("生效日期", "")
        if not effective_date_str:
            continue

        try:
            if isinstance(effective_date_str, str):
                effective_date = datetime.strptime(effective_date_str, "%Y-%m-%d")
            else:
                continue

            next_renewal = effective_date.replace(year=today.year)
            if next_renewal < today:
                next_renewal = next_renewal.replace(year=today.year + 1)

            days_until = (next_renewal - today).days

            policy_copy = policy.copy()
            policy_copy["下次续费日期"] = next_renewal.strftime("%Y-%m-%d")
            policy_copy["距续费天数"] = days_until

            if days_until <= 30:
                policy_copy["续费提醒"] = "紧急"
            elif days_until <= 90:
                policy_copy["续费提醒"] = "关注"
            else:
                policy_copy["续费提醒"] = "正常"

            results.append(policy_copy)

        except Exception:
            continue

    results.sort(key=lambda x: x.get("距续费天数", 999))
    return results


def generate_reminder_list(profiles_dir: str, days_ahead: int = 30) -> str:
    """生成跟进提醒清单

    Args:
        profiles_dir: 档案目录路径
        days_ahead: 提前多少天提醒

    Returns:
        提醒清单 Markdown 格式
    """
    markdown_files = []
    for root, dirs, files in os.walk(profiles_dir):
        for file in files:
            if file.endswith(".md"):
                file_path = os.path.join(root, file)
                try:
                    with open(file_path, "r", encoding="utf-8") as f:
                        content = f.read()
                        markdown_files.append({
                            "path": file_path,
                            "name": file.replace(".md", ""),
                            "content": content
                        })
                except Exception:
                    continue

    today = datetime.now()
    reminders = []

    for mf in markdown_files:
        content = mf["content"]

        policy_pattern = r"\*\*生效日期\*\*:\s*(\d{4}-\d{2}-\d{2})"
        matches = re.findall(policy_pattern, content)

        for date_str in matches:
            try:
                effective_date = datetime.strptime(date_str, "%Y-%m-%d")
                next_renewal = effective_date.replace(year=today.year)
                if next_renewal < today:
                    next_renewal = next_renewal.replace(year=today.year + 1)

                days_until = (next_renewal - today).days

                if 0 <= days_until <= days_ahead:
                    reminders.append({
                        "客户": mf["name"],
                        "续费日期": date_str,
                        "距续费": f"{days_until}天",
                        "紧急程度": "紧急" if days_until <= 7 else ("关注" if days_until <= 30 else "正常")
                    })
            except Exception:
                continue

    reminders.sort(key=lambda x: x["距续费"])

    if not reminders:
        return f"**未来{days_ahead}天内没有续费提醒**"

    result = f"## 续费提醒清单（未来{days_ahead}天）\n\n"
    result += f"**生成时间**: {today.strftime('%Y-%m-%d %H:%M:%S')}\n\n"
    result += "| 客户 | 续费日期 | 距续费 | 紧急程度 |\n"
    result += "|------|----------|--------|----------|\n"

    for r in reminders:
        result += f"| {r['客户']} | {r['续费日期']} | {r['距续费']} | {r['紧急程度']} |\n"

    return result
