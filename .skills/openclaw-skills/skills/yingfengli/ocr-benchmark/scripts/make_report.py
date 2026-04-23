#!/usr/bin/env python3
"""
Generate OCR Benchmark PPT report from JSON results.

Usage:
  python3 make_report.py --results-dir ./ocr-results --images img1.jpg img2.jpg --output report.pptx
  python3 make_report.py --results-dir ./ocr-results --images img1.jpg --scores scores.json
"""

import argparse, json, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

DARK_BG = (30, 30, 40)
TITLE_COLOR = (255, 200, 60)
WHITE = (255, 255, 255)
LIGHT_GRAY = (200, 200, 210)

MODEL_COLORS = {
    'opus': (255, 150, 50),
    'sonnet': (100, 180, 255),
    'haiku': (120, 220, 160),
    'gemini3pro': (220, 100, 255),
    'gemini3flash': (255, 100, 180),
    'paddleocr': (255, 220, 100),
}
DEFAULT_COLOR = (180, 180, 180)


def set_slide_bg(slide, r, g, b):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def add_textbox(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(*color)
    p.font.bold = bold
    p.font.name = 'Arial'
    p.alignment = alignment
    return tf


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_color)
    shape.line.fill.background()
    return shape


def acc_color(v):
    if v >= 90: return (100, 255, 100)
    if v >= 75: return (200, 255, 100)
    if v >= 50: return (255, 200, 100)
    return (255, 100, 100)


def load_results(results_dir, images, models):
    """Load all JSON result files."""
    data = {}
    for img_name in images:
        data[img_name] = {}
        for m in models:
            path = os.path.join(results_dir, f"{img_name}.{m}.json")
            if os.path.exists(path):
                with open(path, 'r', encoding='utf-8') as f:
                    data[img_name][m] = json.load(f)
    return data


def generate_report(results_dir, image_paths, models, scores_path, output_path):
    images = [os.path.splitext(os.path.basename(p))[0] for p in image_paths]
    data = load_results(results_dir, images, models)

    # Load scores if available
    scores = {}
    if scores_path and os.path.exists(scores_path):
        with open(scores_path, 'r', encoding='utf-8') as f:
            scores = json.load(f)

    # Sort models by overall score
    def sort_key(m):
        if m in scores:
            return -scores[m].get('overall_pct', 0)
        return 0
    ranked = sorted(models, key=sort_key)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, *DARK_BG)
    add_textbox(slide, 1, 1.5, 11.33, 1.2, 'OCR Benchmark Report', font_size=48, color=TITLE_COLOR, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1, 3.0, 11.33, 0.8, f'{len(models)} Models × {len(images)} Images', font_size=30, color=WHITE, alignment=PP_ALIGN.CENTER)
    labels = [data.get(images[0], {}).get(m, {}).get('model', m) for m in ranked]
    add_textbox(slide, 1, 4.2, 11.33, 0.5, ' | '.join(labels), font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # ── Slide 2: Ranking ──
    if scores:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide, *DARK_BG)
        add_textbox(slide, 0.4, 0.2, 12, 0.7, '🏆 Final Ranking', font_size=36, color=TITLE_COLOR, bold=True)

        y = 1.2
        for i, m in enumerate(ranked):
            s = scores.get(m, {})
            mc = MODEL_COLORS.get(m, DEFAULT_COLOR)
            label = data.get(images[0], {}).get(m, {}).get('model', m)
            pct = s.get('overall_pct', 0)
            medal = ['🥇','🥈','🥉'][i] if i < 3 else f'{i+1}.'

            add_rect(slide, 0.5, y, 0.15, 0.55, mc)
            add_textbox(slide, 0.8, y, 3, 0.55, f'{medal} {label}', font_size=18, color=mc, bold=True)
            # Accuracy bar
            bar_w = pct / 100 * 7
            add_rect(slide, 4.2, y + 0.1, bar_w, 0.35, mc)
            add_textbox(slide, 4.2 + bar_w + 0.1, y, 1.5, 0.55, f'{pct:.1f}%', font_size=18, color=acc_color(pct), bold=True)
            y += 0.7

    # ── Per-image slides ──
    for idx, img_name in enumerate(images):
        img_path = image_paths[idx] if idx < len(image_paths) else None

        # 3 models per slide (split into pages)
        for page_start in range(0, len(ranked), 3):
            page_models = ranked[page_start:page_start+3]
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_slide_bg(slide, *DARK_BG)
            page_label = 'Top' if page_start == 0 else f'#{page_start+1}-{page_start+len(page_models)}'
            add_textbox(slide, 0.3, 0.15, 12, 0.5, f'{img_name} [{page_label}]', font_size=26, color=TITLE_COLOR, bold=True)

            if img_path and os.path.exists(img_path):
                try:
                    pic = slide.shapes.add_picture(img_path, Inches(0.3), Inches(0.8), height=Inches(5.8))
                    if pic.width > Inches(3.2):
                        ratio = Inches(3.2) / pic.width
                        pic.width = Inches(3.2)
                        pic.height = int(pic.height * ratio)
                except:
                    pass

            col_width = 3.1
            for i, m in enumerate(page_models):
                x = 3.8 + i * col_width
                mc = MODEL_COLORS.get(m, DEFAULT_COLOR)
                d = data.get(img_name, {}).get(m, {})
                pct = scores.get(m, {}).get('images', {}).get(img_name, {}).get('pct', '?')

                add_rect(slide, x, 0.8, col_width - 0.15, 0.5, mc)
                add_textbox(slide, x, 0.82, col_width - 0.15, 0.25, d.get('model', m), font_size=13, color=DARK_BG, bold=True, alignment=PP_ALIGN.CENTER)
                add_textbox(slide, x, 1.05, col_width - 0.15, 0.25, f'Accuracy: {pct}%', font_size=11, color=DARK_BG, alignment=PP_ALIGN.CENTER)

                latency = d.get('latency_seconds', '?')
                add_textbox(slide, x, 1.35, col_width - 0.15, 0.3, f'⏱ {latency}s', font_size=10, color=(160, 160, 175))

                texts = d.get('text_extracted', ['(no data)'])
                ocr_text = '\n'.join(texts[:20])
                if len(texts) > 20:
                    ocr_text += f'\n... +{len(texts)-20} more'
                add_textbox(slide, x, 1.65, col_width - 0.15, 5.5, ocr_text, font_size=9, color=WHITE)

    prs.save(output_path)
    sz = os.path.getsize(output_path)
    print(f"✅ Report saved: {output_path} ({sz/1024:.1f} KiB, {len(prs.slides)} slides)")


def main():
    parser = argparse.ArgumentParser(description='Generate OCR Benchmark PPT Report')
    parser.add_argument('--results-dir', required=True, help='Directory with JSON results')
    parser.add_argument('--images', nargs='+', required=True, help='Original image paths')
    parser.add_argument('--models', nargs='+', default=list(MODEL_COLORS.keys()), help='Model keys to include')
    parser.add_argument('--scores', help='Path to scores.json')
    parser.add_argument('--output', default='OCR_Benchmark_Report.pptx', help='Output PPTX path')
    args = parser.parse_args()

    generate_report(args.results_dir, args.images, args.models, args.scores, args.output)


if __name__ == '__main__':
    main()
