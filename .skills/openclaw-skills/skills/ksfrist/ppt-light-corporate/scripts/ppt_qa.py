#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""ppt_qa.py — Automated QA verification for generated PPTs.

Two core checks:
1. Loop Verification: boundary + text extraction + data cross-check
2. Layout Diversity: detect repeated adjacent layouts

Usage:
    from ppt_qa import PPTQualityChecker
    qc = PPTQualityChecker("output.pptx")
    report = qc.full_report(source_text="原始素材文本...")
    print(report)
"""
import re
from collections import Counter
from pptx import Presentation
from pptx.util import Emu


class PPTQualityChecker:
    """Automated QA for generated PPTs."""

    CANVAS_W = 10.0
    CANVAS_H = 5.625
    SAFE_MARGIN = 0.3  # slightly relaxed for checking (0.4 for generation)

    # Layout fingerprints: shape count + approximate positions
    # Used to detect when adjacent slides use identical layouts
    DENSITY_HIGH = 300  # chars per slide = too dense
    DENSITY_LOW = 20    # chars per slide = too empty (excl cover/ending)

    def __init__(self, pptx_path):
        self.path = pptx_path
        self.prs = Presentation(pptx_path)
        self.slides = list(self.prs.slides)
        self.n_slides = len(self.slides)

    # ── 1. Boundary Check ─────────────────────────────────

    def check_boundaries(self):
        """Check all shapes stay within canvas bounds."""
        violations = []
        for idx, slide in enumerate(self.slides):
            for shape in slide.shapes:
                l = Emu(shape.left).inches
                t = Emu(shape.top).inches
                w = Emu(shape.width).inches
                h = Emu(shape.height).inches
                r = l + w
                b = t + h
                name = shape.text_frame.text[:40] if shape.has_text_frame else shape.name
                if l < -0.01:
                    violations.append(f"S{idx+1}: \"{name}\" left={l:.2f}\" (超出左边界)")
                if t < -0.01:
                    violations.append(f"S{idx+1}: \"{name}\" top={t:.2f}\" (超出上边界)")
                if r > self.CANVAS_W + 0.01:
                    violations.append(f"S{idx+1}: \"{name}\" right={r:.2f}\" (超出右边界 {self.CANVAS_W}\")")
                if b > self.CANVAS_H + 0.01:
                    violations.append(f"S{idx+1}: \"{name}\" bottom={b:.2f}\" (超出下边界 {self.CANVAS_H}\")")
        return violations

    # ── 2. Text Extraction ────────────────────────────────

    def extract_all_text(self):
        """Extract text from every slide, return list of (slide_num, text)."""
        result = []
        for idx, slide in enumerate(self.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        texts.append(t)
            result.append((idx + 1, '\n'.join(texts)))
        return result

    # ── 3. Data Cross-Check ───────────────────────────────

    def extract_numbers_dates(self):
        """Extract all numbers and date-like strings from PPT."""
        findings = []
        for idx, slide in enumerate(self.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    # Numbers (with Chinese units)
                    nums = re.findall(r'[\d,]+\.?\d*\s*[万亿千百%+件条个项人次年月日天]+', text)
                    # Pure percentages
                    pcts = re.findall(r'\d+\.?\d*%', text)
                    # Year patterns
                    years = re.findall(r'20[0-9]{2}(?:\.[0-9]{1,2})?(?:-[0-9]{1,2})?', text)
                    # Date patterns
                    dates = re.findall(r'\d{1,2}月\d{1,2}日', text)
                    for n in set(nums + pcts + years + dates):
                        findings.append((idx + 1, n))
        # Deduplicate
        return sorted(set(findings))

    def cross_check_content(self, source_text):
        """Check if key data from source appears in PPT."""
        if not source_text:
            return []
        # Extract numbers from source
        src_nums = set(re.findall(r'[\d,]+\.?\d*\s*[万亿千百%+件条个项人次年月日天]+', source_text))
        src_nums.update(re.findall(r'\d+\.?\d*%', source_text))
        # Extract numbers from PPT
        ppt_data = self.extract_numbers_dates()
        ppt_nums = set(n for _, n in ppt_data)
        # Find missing
        missing = []
        for num in src_nums:
            # Fuzzy match: check if any PPT number contains same digits
            num_digits = re.sub(r'[^\d.]', '', num)
            found = False
            for ppt_num in ppt_nums:
                if num_digits in re.sub(r'[^\d.]', '', ppt_num):
                    found = True
                    break
            if not found and len(num_digits) >= 2:  # ignore single digits
                missing.append(num)
        return sorted(missing)

    # ── 4. Text Truncation Risk ───────────────────────────

    def check_text_overflow_risk(self):
        """Estimate if text might overflow its text box, including title checks."""
        risks = []
        for idx, slide in enumerate(self.slides):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                w_inches = Emu(shape.width).inches
                h_inches = Emu(shape.height).inches
                t_inches = Emu(shape.top).inches
                is_title = (t_inches < 0.6 and h_inches < 0.6)
                if is_title:
                    max_chars = int(w_inches * 3.5)
                    if len(text) > max_chars:
                        risks.append(f"S{idx+1}: 标题\"{text[:20]}...\" ({len(text)}字>框容{max_chars}字) — 会换行溢出")
                else:
                    chars_per_line = max(1, int(w_inches * 7))
                    lines_needed = len(text) / chars_per_line
                    height_needed = lines_needed * 0.22
                    if height_needed > h_inches * 1.5 and len(text) > 20:
                        risks.append(f"S{idx+1}: \"{text[:30]}...\" ({len(text)}字 in {w_inches:.1f}\"×{h_inches:.1f}\" box)")
        return risks

    # ── 5. Information Density ────────────────────────────

    def check_density(self):
        """Check per-slide character count for density balance."""
        issues = []
        slide_texts = self.extract_all_text()
        for slide_num, text in slide_texts:
            char_count = len(text.replace('\n', '').replace(' ', ''))
            if slide_num == 1 or slide_num == self.n_slides:
                continue  # skip cover/ending
            if char_count > self.DENSITY_HIGH:
                issues.append(f"S{slide_num}: {char_count}字 (过密，建议<{self.DENSITY_HIGH}字或拆页)")
            elif char_count < self.DENSITY_LOW:
                issues.append(f"S{slide_num}: {char_count}字 (过空，建议补充内容)")
        return issues

    # ── 5a. Element Overlap Detection ───────────────────────

    def check_element_overlap(self):
        """Check if text elements overlap with other shapes."""
        issues = []
        for idx, slide in enumerate(self.slides):
            text_shapes = []
            all_shapes = []
            for shape in slide.shapes:
                t = Emu(shape.top).inches
                b = t + Emu(shape.height).inches
                l = Emu(shape.left).inches
                r = l + Emu(shape.width).inches
                has_text = shape.has_text_frame and shape.text_frame.text.strip()
                entry = {"t": t, "b": b, "l": l, "r": r, "text": has_text, "txt": shape.text_frame.text.strip()[:20] if has_text else ""}
                all_shapes.append(entry)
                if has_text:
                    text_shapes.append(entry)
            # Check title/subtitle vs content shapes
            for ts in text_shapes:
                if ts["t"] > 1.0:  # not a title
                    continue
                for other in all_shapes:
                    if other is ts:
                        continue
                    if other["t"] < 1.0:  # both in title zone, skip
                        continue
                    # Check vertical overlap
                    v_overlap = min(ts["b"], other["b"]) - max(ts["t"], other["t"])
                    h_overlap = min(ts["r"], other["r"]) - max(ts["l"], other["l"])
                    if v_overlap > 0.05 and h_overlap > 0.1:
                        issues.append(
                            f"S{idx+1}: \"{ts['txt']}\" (y={ts['t']:.2f}~{ts['b']:.2f}) "
                            f"与内容元素(y={other['t']:.2f}~{other['b']:.2f})重叠{v_overlap:.2f}\"")
                        break  # one overlap per title is enough
        return issues

    # ── 5b. Vertical Balance ────────────────────────────────

    def check_vertical_balance(self):
        """Check if content is vertically centered in the page."""
        issues = []
        for idx, slide in enumerate(self.slides):
            if idx == 0 or idx == self.n_slides - 1:
                continue  # skip cover/ending
            tops = []
            bottoms = []
            for shape in slide.shapes:
                t = Emu(shape.top).inches
                b = t + Emu(shape.height).inches
                if t < 0.8:  # skip title area
                    continue
                tops.append(t)
                bottoms.append(b)
            if not tops:
                continue
            content_top = min(tops)
            content_bottom = max(bottoms)
            # Available area: 1.0" to 5.2" (below title, above margin)
            area_top = 1.0
            area_bottom = self.CANVAS_H - 0.4
            space_above = content_top - area_top
            space_below = area_bottom - content_bottom
            # If content uses less than 60% of area AND is heavily skewed
            content_h = content_bottom - content_top
            area_h = area_bottom - area_top
            if content_h < area_h * 0.6:
                ratio = space_below / max(space_above, 0.01) if space_above > 0 else 999
                if ratio > 3.0:
                    issues.append(
                        f"S{idx+1}: 内容偏上(上留白{space_above:.1f}\"/下留白{space_below:.1f}\")"
                        f" — 建议垂直居中")
                elif ratio < 0.33:
                    issues.append(
                        f"S{idx+1}: 内容偏下(上留白{space_above:.1f}\"/下留白{space_below:.1f}\")"
                        f" — 建议垂直居中")
        return issues

    # ── 6. Layout Diversity ───────────────────────────────

    def _slide_fingerprint(self, slide):
        """Generate a layout fingerprint for a slide.
        Based on: shape count, approximate grid positions, shape types.
        """
        shapes_info = []
        for shape in slide.shapes:
            # Quantize position to 0.5" grid
            l = round(Emu(shape.left).inches * 2) / 2
            t = round(Emu(shape.top).inches * 2) / 2
            w = round(Emu(shape.width).inches * 2) / 2
            h = round(Emu(shape.height).inches * 2) / 2
            st = str(shape.shape_type)
            shapes_info.append(f"{st}@{l},{t},{w},{h}")
        shapes_info.sort()
        return '|'.join(shapes_info)

    def check_layout_diversity(self):
        """Detect adjacent slides using identical/similar layouts."""
        issues = []
        fingerprints = []
        for slide in self.slides:
            fingerprints.append(self._slide_fingerprint(slide))

        # Check adjacent pairs
        for i in range(len(fingerprints) - 1):
            if i == 0:  # skip cover
                continue
            fp1 = fingerprints[i]
            fp2 = fingerprints[i + 1]
            if fp1 == fp2:
                issues.append(f"S{i+1} ↔ S{i+2}: 布局完全相同，建议换一种版式")
            elif self._fingerprint_similarity(fp1, fp2) > 0.8:
                issues.append(f"S{i+1} ↔ S{i+2}: 布局高度相似({self._fingerprint_similarity(fp1, fp2):.0%})，建议增加变化")

        # Check overall: any layout used 3+ times
        fp_counter = Counter(fingerprints[1:])  # exclude cover
        for fp, count in fp_counter.items():
            if count >= 3:
                # Find which slides
                slide_nums = [i+1 for i, f in enumerate(fingerprints) if f == fp]
                issues.append(f"同一布局出现{count}次 (S{','.join(map(str,slide_nums))})，建议增加版式多样性")

        return issues

    def _fingerprint_similarity(self, fp1, fp2):
        """Compute Jaccard similarity between two fingerprints."""
        set1 = set(fp1.split('|'))
        set2 = set(fp2.split('|'))
        if not set1 and not set2:
            return 1.0
        intersection = len(set1 & set2)
        union = len(set1 | set2)
        return intersection / union if union > 0 else 0

    # ── Full Report ───────────────────────────────────────

    def full_report(self, source_text=None):
        """Run all checks and return formatted report string."""
        lines = [f"📋 审校报告 ({self.path})", f"总页数: {self.n_slides}", ""]

        # 1. Boundaries
        violations = self.check_boundaries()
        if violations:
            lines.append(f"❌ 边界溢出: {len(violations)} 个问题")
            for v in violations:
                lines.append(f"   - {v}")
        else:
            lines.append("✅ 边界溢出: 0 violations")

        # 2. Text overflow risk
        risks = self.check_text_overflow_risk()
        if risks:
            lines.append(f"⚠️ 文字截断风险: {len(risks)} 个")
            for r in risks:
                lines.append(f"   - {r}")
        else:
            lines.append("✅ 文字截断: 无风险")

        # 3. Density
        density_issues = self.check_density()
        if density_issues:
            lines.append(f"⚠️ 信息密度: {len(density_issues)} 个问题")
            for d in density_issues:
                lines.append(f"   - {d}")
        else:
            lines.append("✅ 信息密度: 均衡")

        # 4. Layout diversity
        # 4a. Element overlap
        overlap_issues = self.check_element_overlap()
        if overlap_issues:
            lines.append(f"⚠️ 元素重叠: {len(overlap_issues)} 个问题")
            for o in overlap_issues:
                lines.append(f"   - {o}")
        else:
            lines.append("✅ 元素重叠: 无")

        # 4b. Vertical balance
        vb_issues = self.check_vertical_balance()
        if vb_issues:
            lines.append(f"⚠️ 垂直平衡: {len(vb_issues)} 个问题")
            for v in vb_issues:
                lines.append(f"   - {v}")
        else:
            lines.append("✅ 垂直平衡: OK")

        layout_issues = self.check_layout_diversity()
        if layout_issues:
            lines.append(f"⚠️ 布局多样性: {len(layout_issues)} 个问题")
            for l in layout_issues:
                lines.append(f"   - {l}")
        else:
            lines.append("✅ 布局多样性: OK")

        # 5. Data cross-check
        data = self.extract_numbers_dates()
        if data:
            lines.append(f"📊 数据核对 ({len(data)} 项数据，请人工确认):")
            for slide_num, num in data:
                lines.append(f"   - S{slide_num}: {num}")
        else:
            lines.append("ℹ️ 数据核对: 未发现数字/日期数据")

        # 6. Content completeness (if source provided)
        if source_text:
            missing = self.cross_check_content(source_text)
            if missing:
                lines.append(f"⚠️ 内容完整性: {len(missing)} 项素材数据未在PPT中找到")
                for m in missing:
                    lines.append(f"   - 缺失: {m}")
            else:
                lines.append("✅ 内容完整性: 素材关键数据已全部覆盖")

        return '\n'.join(lines)


# CLI entry
if __name__ == '__main__':
    import sys
    if len(sys.argv) < 2:
        print("Usage: python ppt_qa.py <pptx_path> [source_text_file]")
        sys.exit(1)
    pptx_path = sys.argv[1]
    source_text = None
    if len(sys.argv) >= 3:
        with open(sys.argv[2], 'r', encoding='utf-8') as f:
            source_text = f.read()
    qc = PPTQualityChecker(pptx_path)
    print(qc.full_report(source_text))
