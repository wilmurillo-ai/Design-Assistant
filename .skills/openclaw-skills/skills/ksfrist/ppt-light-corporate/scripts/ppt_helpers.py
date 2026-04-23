#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""ppt_helpers.py — Reusable helpers for the light corporate template.

Usage:
    from ppt_helpers import TemplateBuilder
    tb = TemplateBuilder()
    s = tb.add_content_slide("Page Title")
    tb.text(s, 0.5, 1.2, 4.2, 0.3, "Hello", sz=13, color=tb.RED3, bold=True)
    tb.save("output.pptx")
"""
import sys, os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml


class TemplateBuilder:
    """Build PPTs from the light corporate template."""

    # ── Colors ───────────────────────────────────────────
    RED    = RGBColor(0xC0, 0x00, 0x00)
    RED2   = RGBColor(0xEA, 0x00, 0x00)
    RED3   = RGBColor(0xD8, 0x22, 0x2A)
    DARK   = RGBColor(0x34, 0x34, 0x34)
    MID    = RGBColor(0x4B, 0x49, 0x48)
    LIGHT  = RGBColor(0xA0, 0xA0, 0xA0)
    BG     = RGBColor(0xF1, 0xF1, 0xF1)
    BORDER = RGBColor(0xE7, 0xE7, 0xE7)
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    ACCENT_GREY = RGBColor(0x6B, 0x6B, 0x6B)  # warm grey for alternating accents
    CORAL  = RGBColor(0xE0, 0x6C, 0x5C)         # soft coral red
    DARK_RED = RGBColor(0x8B, 0x1A, 0x1A)        # dark red
    LIGHT_RED = RGBColor(0xF2, 0xD5, 0xD5)       # light red tint
    LIGHT_RED2 = RGBColor(0xF5, 0xE0, 0xE0)      # lighter red tint
    FONT   = "微软雅黑"

    def __init__(self, template_path=None):
        if template_path is None:
            skill_dir = Path(__file__).parent.parent
            template_path = skill_dir / "assets" / "template-light.pptx"
        self.prs = Presentation(str(template_path))
        self._clear_slides()
        self._icon_lib = None  # lazy-loaded

    @property
    def icons(self):
        """Lazy-load the icon library."""
        if self._icon_lib is None:
            from icon_library import IconLibrary
            self._icon_lib = IconLibrary()
        return self._icon_lib

    def place_icon(self, slide, query, x, y, size=0.5, color=None, category=None):
        """Search for an icon and place it on the slide.
        color: hex string e.g. 'FFFFFF' for white on dark backgrounds.
        Returns icon entry if found, None otherwise.
        """
        results = self.icons.search(query, category=category, limit=1)
        if results:
            self.icons.copy_icon(slide, results[0], x, y, size, color=color)
            return results[0]
        return None

    def _clear_slides(self):
        """Remove all demo slides from template."""
        ns = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'
        while len(self.prs.slides._sldIdLst) > 0:
            el = self.prs.slides._sldIdLst[0]
            rId = el.get(f'{ns}id') or el.attrib.get('r:id')
            if rId:
                try: self.prs.part.drop_rel(rId)
                except: pass
            self.prs.slides._sldIdLst.remove(el)

    # ── Slide creation ───────────────────────────────────
    def add_content_slide(self, title=None):
        """Add a content slide (Layout 3: 绝密-仅标题)."""
        s = self.prs.slides.add_slide(self.prs.slide_layouts[3])
        if title:
            self._page_title(s, title)
        return s

    def add_cover_slide(self):
        """Add a cover/section slide (Layout 6: 节标题)."""
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def add_two_column_slide(self, title=None):
        """Add a two-column slide (Layout 1: 绝密-两栏内容)."""
        s = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        if title:
            self._page_title(s, title)
        return s

    def add_comparison_slide(self, title=None):
        """Add a comparison slide (Layout 2: 绝密-比较)."""
        s = self.prs.slides.add_slide(self.prs.slide_layouts[2])
        if title:
            self._page_title(s, title)
        return s

    # ── Page title ───────────────────────────────────────
    def _page_title(self, slide, title):
        # Auto-adjust: wider box + shrink font for long titles
        title_w = 7.8  # nearly full width (was 4.2")
        title_h = 0.4
        # Estimate chars that fit: ~3.5 CJK chars per inch at 20pt
        max_chars_20pt = int(title_w * 3.5)
        if len(title) > max_chars_20pt:
            sz = 16  # shrink for very long titles
        elif len(title) > int(title_w * 2.8):
            sz = 18  # slightly shrink
        else:
            sz = 20  # default
        box = self.text(slide, 1.1, 0.4, title_w, title_h, title,
                        sz=sz, color=self.RED, bold=True)
        # Zero insets to match template
        self._zero_insets(box)
        self._set_line_spacing(box.text_frame.paragraphs[0], 120)
        return box

    # ── Shape helpers ────────────────────────────────────
    def rect(self, slide, l, t, w, h, color, border=None):
        sh = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
        sh.fill.solid(); sh.fill.fore_color.rgb = color
        if border:
            sh.line.color.rgb = border; sh.line.width = Pt(0.5)
        else:
            sh.line.fill.background()
        return sh

    def rounded_rect(self, slide, l, t, w, h, color, border=None):
        sh = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
        sh.fill.solid(); sh.fill.fore_color.rgb = color
        if border:
            sh.line.color.rgb = border; sh.line.width = Pt(0.5)
        else:
            sh.line.fill.background()
        return sh

    def card(self, slide, l, t, w, h):
        """Standard card: BG fill + BORDER outline."""
        return self.rounded_rect(slide, l, t, w, h, self.BG, border=self.BORDER)

    def divider(self, slide, l, t, w, color=None):
        """Thin horizontal line."""
        color = color or self.BORDER
        return self.rect(slide, l, t, w, 0.005, color)

    # ── Text helpers ─────────────────────────────────────
    def text(self, slide, l, t, w, h, text, sz=11, color=None,
             bold=False, align=PP_ALIGN.LEFT, shrink=False):
        color = color or self.DARK
        box = slide.shapes.add_textbox(
            Inches(l), Inches(t), Inches(w), Inches(h))
        tf = box.text_frame; tf.word_wrap = True
        # Auto-shrink text to fit box when enabled
        if shrink:
            from pptx.enum.text import MSO_AUTO_SIZE
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(sz); p.font.color.rgb = color
        p.font.bold = bold
        # Force BOTH latin and East Asian font on ALL runs
        from lxml.etree import SubElement
        from pptx.oxml.ns import qn
        for run in p.runs:
            rPr = run._r.get_or_add_rPr()
            # Latin font
            latin = rPr.find(qn('a:latin'))
            if latin is None:
                latin = SubElement(rPr, qn('a:latin'))
            latin.set('typeface', self.FONT)
            # East Asian font
            ea = rPr.find(qn('a:ea'))
            if ea is None:
                ea = SubElement(rPr, qn('a:ea'))
            ea.set('typeface', self.FONT)
        p.alignment = align
        self._set_line_spacing(p, 120)
        return box

    def multi_text(self, slide, l, t, w, h, lines, sz=10, color=None, sp=4):
        """Multi-line text box.
        lines: list of str or (str, dict) where dict overrides: sz, color, bold, align, sa
        """
        color = color or self.DARK
        box = slide.shapes.add_textbox(
            Inches(l), Inches(t), Inches(w), Inches(h))
        tf = box.text_frame; tf.word_wrap = True
        for i, item in enumerate(lines):
            txt, ov = (item, {}) if isinstance(item, str) else (item[0], item[1])
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = txt
            p.font.size = Pt(ov.get("sz", sz))
            p.font.color.rgb = ov.get("color", color)
            p.font.bold = ov.get("bold", False)
            p.font.name = ov.get("font", self.FONT)
            # Force BOTH latin and East Asian font on ALL runs
            from lxml.etree import SubElement
            from pptx.oxml.ns import qn
            for run in p.runs:
                rPr = run._r.get_or_add_rPr()
                latin = rPr.find(qn('a:latin'))
                if latin is None:
                    latin = SubElement(rPr, qn('a:latin'))
                latin.set('typeface', ov.get("font", self.FONT))
                ea = rPr.find(qn('a:ea'))
                if ea is None:
                    ea = SubElement(rPr, qn('a:ea'))
                ea.set('typeface', ov.get("font", self.FONT))
            p.alignment = ov.get("align", PP_ALIGN.LEFT)
            p.space_after = Pt(ov.get("sa", sp))
            self._set_line_spacing(p, 120)
        return box

    def section_header(self, slide, l, t, text, w=4.2):
        """Red section sub-header."""
        return self.text(slide, l, t, w, 0.3, text,
                         sz=13, color=self.RED3, bold=True)

    # ── Table helpers ────────────────────────────────────
    def table_header(self, slide, x0, y0, headers, col_widths, row_height=0.36):
        """Red header row for a comparison table."""
        x = x0
        for j, h in enumerate(headers):
            self.rect(slide, x, y0, col_widths[j], row_height, self.RED)
            self.text(slide, x, y0, col_widths[j], row_height,
                      h, sz=10, color=self.WHITE, bold=True, align=PP_ALIGN.CENTER)
            x += col_widths[j]

    def table_row(self, slide, x0, y, cells, col_widths, row_height=0.36,
                  bg=None, is_label_col=True):
        """Single data row. First column treated as label (RED, bold, left-align)."""
        bg = bg or self.WHITE
        x = x0
        for j, cell in enumerate(cells):
            self.rect(slide, x, y, col_widths[j], row_height, bg, border=self.BORDER)
            self.text(slide, x + 0.08, y, col_widths[j] - 0.16, row_height,
                      cell, sz=10,
                      color=self.RED if (j == 0 and is_label_col) else self.DARK,
                      bold=(j == 0 and is_label_col),
                      align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)
            x += col_widths[j]

    # ── Two-Column Layout helpers (learned from 两栏模板) ──

    def two_col_numbered_cards(self, slide, items, card_top=1.6):
        """Side-by-side cards with numbered circles.
        items: list of 2 dicts: {num, title, body}
        Learned from: 两栏模板 Slide 1
        """
        card_w, card_h = 4.15, 3.2
        positions = [(0.6, items[0]), (5.25, items[1])]
        for x, item in positions:
            # Card background
            self.card(slide, x, card_top, card_w, card_h)
            # Number circle
            c = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x + 0.15), Inches(card_top - 0.25),
                Inches(0.65), Inches(0.65))
            c.fill.solid(); c.fill.fore_color.rgb = self.RED
            c.line.fill.background()
            self.text(slide, x + 0.15, card_top - 0.22, 0.65, 0.65,
                      item['num'], sz=20, color=self.WHITE, bold=True,
                      align=PP_ALIGN.CENTER)
            # Title
            self.text(slide, x + 0.25, card_top + 0.55, card_w - 0.5, 0.35,
                      item['title'], sz=16, color=self.RED, bold=True)
            # Divider line
            self.divider(slide, x + 0.2, card_top + 1.0, card_w - 0.4)
            # Body
            self.text(slide, x + 0.25, card_top + 1.15, card_w - 0.5, 1.8,
                      item['body'], sz=11, color=self.DARK,
                      shrink=True)

    def two_row_numbered_bars(self, slide, items, bar_top=1.3):
        """Stacked horizontal bars with left-side numbers.
        items: list of 2 dicts: {num, title, body}
        Learned from: 两栏模板 Slide 4
        """
        bar_w, bar_h = 8.5, 1.75
        bar_x = 0.75
        gap = 0.3
        for i, item in enumerate(items[:2]):
            y = bar_top + i * (bar_h + gap)
            self.card(slide, bar_x, y, bar_w, bar_h)
            self.text(slide, bar_x + 0.2, y + 0.35, 0.8, 0.7,
                      item['num'], sz=28, color=self.RED, bold=True,
                      align=PP_ALIGN.CENTER)
            self.text(slide, bar_x + 1.2, y + 0.2, bar_w - 1.6, 0.35,
                      item['title'], sz=16, color=self.DARK, bold=True)
            self.text(slide, bar_x + 1.2, y + 0.65, bar_w - 1.6, 1.0,
                      item['body'], sz=11, color=self.MID,
                      shrink=True)

    def two_row_label_tabs(self, slide, items, row_top=1.3):
        """Full-width rows with colored label tabs on the left.
        items: list of 2 dicts: {label, body}
        Learned from: 两栏模板 Slide 5
        """
        row_w, row_h = 8.5, 1.75
        row_x = 0.75
        tab_w, tab_h = 2.6, 0.85
        gap = 0.3
        for i, item in enumerate(items[:2]):
            y = row_top + i * (row_h + gap)
            self.card(slide, row_x, y, row_w, row_h)
            # Red label tab
            self.rounded_rect(slide, row_x - 0.15,
                              y + (row_h - tab_h) / 2,
                              tab_w, tab_h, self.RED)
            self.text(slide, row_x - 0.15,
                      y + (row_h - tab_h) / 2 + 0.15,
                      tab_w, tab_h - 0.3,
                      item['label'], sz=14, color=self.WHITE, bold=True,
                      align=PP_ALIGN.CENTER)
            # Body text
            self.text(slide, row_x + tab_w + 0.1, y + 0.25,
                      row_w - tab_w - 0.3, row_h - 0.5,
                      item['body'], sz=11, color=self.DARK,
                      shrink=True)

    def two_col_circle_cards(self, slide, items, card_top=2.0):
        """Left-right cards with circle avatars/icons floating above.
        items: list of 2 dicts: {icon_text, title, body}
        Learned from: 两栏模板 Slide 3
        """
        card_w, card_h = 4.15, 2.8
        positions = [(0.6, items[0]), (5.25, items[1])]
        circle_size = 1.2
        for x, item in positions:
            self.card(slide, x, card_top, card_w, card_h)
            cx = x + (card_w - circle_size) / 2
            cy = card_top - circle_size * 0.5
            c = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(cx), Inches(cy),
                Inches(circle_size), Inches(circle_size))
            c.fill.solid(); c.fill.fore_color.rgb = self.RED3
            c.line.fill.background()
            self.text(slide, cx, cy + 0.25, circle_size, circle_size - 0.5,
                      item['icon_text'], sz=11, color=self.WHITE, bold=True,
                      align=PP_ALIGN.CENTER)
            self.text(slide, x + 0.25, card_top + 0.75, card_w - 0.5, 0.35,
                      item['title'], sz=16, color=self.RED, bold=True)
            self.text(slide, x + 0.25, card_top + 1.2, card_w - 0.5, 1.4,
                      item['body'], sz=11, color=self.DARK,
                      shrink=True)

    def left_text_right_stacked_cards(self, slide, left_title, left_body,
                                       cards, left_x=0.6, right_x=5.0):
        """Left narrative text + right stacked mini cards.
        cards: list of 2-3 dicts: {num, title, body}
        Left text uses larger font to fill the space proportionally.
        """
        n = min(len(cards), 3)
        # Left text — large title + generous body
        self.text(slide, left_x, 1.2, 4.0, 0.4,
                  left_title, sz=18, color=self.RED3, bold=True)
        self.text(slide, left_x, 1.75, 4.0, 3.0,
                  left_body, sz=13, color=self.DARK,
                  shrink=True)

        # Right stacked cards — adaptive height
        total_h = 3.6  # available height
        gap = 0.2
        mini_w = 4.6
        mini_h = (total_h - (n - 1) * gap) / n

        for i, card in enumerate(cards[:n]):
            y = 1.2 + i * (mini_h + gap)
            self.card(slide, right_x, y, mini_w, mini_h)
            # Number
            self.text(slide, right_x + 0.15, y + 0.1, 0.6, 0.45,
                      card['num'], sz=20, color=self.RED, bold=True,
                      align=PP_ALIGN.CENTER)
            # Title
            self.text(slide, right_x + 0.8, y + 0.1, mini_w - 1.0, 0.28,
                      card['title'], sz=12, color=self.DARK, bold=True)
            # Body
            body_top = y + 0.42
            body_h = mini_h - 0.52
            self.text(slide, right_x + 0.8, body_top, mini_w - 1.0, body_h,
                      card['body'], sz=10, color=self.MID,
                      shrink=True)

    def three_col_vertical_cards(self, slide, items, card_top=1.2):
        """Three tall vertical cards side by side.
        items: list of 3 dicts: {title, body}
        Learned from: 三栏模板 Slide 4
        """
        card_w, card_h = 2.85, 3.7
        gap = 0.35
        start_x = (10.0 - (3 * card_w + 2 * gap)) / 2  # centered
        for i, item in enumerate(items[:3]):
            x = start_x + i * (card_w + gap)
            self.card(slide, x, card_top, card_w, card_h)
            # Title
            self.text(slide, x + 0.15, card_top + 0.15, card_w - 0.3, 0.55,
                      item['title'], sz=14, color=self.RED, bold=True)
            # Divider
            self.divider(slide, x + 0.15, card_top + 0.8, card_w - 0.3)
            # Body
            self.text(slide, x + 0.15, card_top + 0.95, card_w - 0.3, card_h - 1.2,
                      item['body'], sz=10, color=self.DARK,
                      shrink=True)

    def three_row_numbered_bars(self, slide, items, bar_top=1.1):
        """Three stacked full-width bars with numbered circles.
        items: list of 3 dicts: {num, title, body}
        Learned from: 三栏模板 Slide 1
        """
        bar_w, bar_h = 8.8, 1.15
        bar_x = 0.6
        gap = 0.2
        for i, item in enumerate(items[:3]):
            y = bar_top + i * (bar_h + gap)
            self.card(slide, bar_x, y, bar_w, bar_h)
            # Number circle
            c = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(bar_x + 0.15), Inches(y + (bar_h - 0.5) / 2),
                Inches(0.5), Inches(0.5))
            c.fill.solid(); c.fill.fore_color.rgb = self.RED
            c.line.fill.background()
            self.text(slide, bar_x + 0.15, y + (bar_h - 0.5) / 2 + 0.02,
                      0.5, 0.5, item['num'],
                      sz=18, color=self.WHITE, bold=True, align=PP_ALIGN.CENTER)
            # Title
            self.text(slide, bar_x + 0.85, y + 0.1, bar_w - 1.1, 0.3,
                      item['title'], sz=14, color=self.DARK, bold=True)
            # Body
            self.text(slide, bar_x + 0.85, y + 0.45, bar_w - 1.1, 0.65,
                      item['body'], sz=10, color=self.MID,
                      shrink=True)

    def three_row_label_bars(self, slide, items, row_top=1.1):
        """Three stacked rows with left colored label + right body text.
        items: list of 3 dicts: {label, body}
        Learned from: 三栏模板 Slide 6
        """
        label_w, label_h = 1.2, 1.15
        body_w = 7.6
        row_x = 0.5
        gap = 0.15
        for i, item in enumerate(items[:3]):
            y = row_top + i * (label_h + gap)
            # Label shape
            self.rounded_rect(slide, row_x, y, label_w, label_h, self.RED)
            self.text(slide, row_x, y + (label_h - 0.35) / 2, label_w, 0.35,
                      item['label'], sz=14, color=self.WHITE, bold=True,
                      align=PP_ALIGN.CENTER)
            # Body area
            self.card(slide, row_x + label_w + 0.1, y, body_w, label_h)
            self.text(slide, row_x + label_w + 0.25, y + 0.15,
                      body_w - 0.3, label_h - 0.3,
                      item['body'], sz=11, color=self.DARK,
                      shrink=True)

    def three_col_icon_cards(self, slide, items, card_top=1.2):
        """Three columns with icon circle + title + body.
        items: list of 3 dicts: {icon_text, title, body}
        Learned from: 三栏模板 Slide 7 (adapted to H3C style)
        """
        card_w, card_h = 2.85, 3.7
        gap = 0.35
        start_x = (10.0 - (3 * card_w + 2 * gap)) / 2  # centered
        for i, item in enumerate(items[:3]):
            x = start_x + i * (card_w + gap)
            self.card(slide, x, card_top, card_w, card_h)
            # Icon circle at top center
            cx = x + (card_w - 0.7) / 2
            c = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(cx), Inches(card_top + 0.2),
                Inches(0.7), Inches(0.7))
            c.fill.solid(); c.fill.fore_color.rgb = self.RED
            c.line.fill.background()
            self.text(slide, cx, card_top + 0.3, 0.7, 0.5,
                      item['icon_text'], sz=11, color=self.WHITE, bold=True,
                      align=PP_ALIGN.CENTER)
            # Title below icon
            self.text(slide, x + 0.15, card_top + 1.1, card_w - 0.3, 0.35,
                      item['title'], sz=13, color=self.RED, bold=True,
                      align=PP_ALIGN.CENTER)
            # Body
            self.text(slide, x + 0.15, card_top + 1.55, card_w - 0.3, card_h - 1.8,
                      item['body'], sz=10, color=self.DARK,
                      shrink=True)

    # ── Four-Column Layout helpers (learned from 四栏模板) ──

    def four_col_tall_cards(self, slide, items, card_top=1.2):
        """Four tall vertical cards with icon circle + title + divider + body.
        items: list of 4 dicts: {title, body}
        Learned from: 四栏模板 Slide 3
        """
        card_w, card_h = 2.1, 3.3
        gap = 0.35
        start_x = (10.0 - (4 * card_w + 3 * gap)) / 2  # centered
        for i, item in enumerate(items[:4]):
            x = start_x + i * (card_w + gap)
            self.card(slide, x, card_top, card_w, card_h)
            # Icon circle at top center
            cx = x + (card_w - 0.6) / 2
            c = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(cx), Inches(card_top + 0.15),
                Inches(0.6), Inches(0.6))
            c.fill.solid(); c.fill.fore_color.rgb = self.RED
            c.line.fill.background()
            # Title — auto-size height based on line count
            title_text = item['title']
            title_lines = title_text.count('\n') + 1
            title_h = 0.25 * title_lines + 0.1  # ~0.35 for 1 line, ~0.6 for 2
            title_y = card_top + 0.9
            self.text(slide, x + 0.1, title_y, card_w - 0.2, title_h,
                      title_text, sz=12, color=self.RED, bold=True,
                      align=PP_ALIGN.CENTER)
            # Divider — below title with gap
            div_y = title_y + title_h + 0.08
            self.divider(slide, x + 0.1, div_y, card_w - 0.2)
            # Body — below divider
            body_y = div_y + 0.12
            body_h = card_top + card_h - body_y - 0.15
            self.text(slide, x + 0.1, body_y, card_w - 0.2, body_h,
                      item['body'], sz=10, color=self.DARK,
                      shrink=True)

    def four_row_badge_bars(self, slide, items, row_top=1.1):
        """Four rows with left badge/label + right body text.
        items: list of 4 dicts: {badge, title, body}
        Learned from: 四栏模板 Slide 4
        """
        badge_w, badge_h = 1.5, 0.5
        body_w = 7.0
        row_x = 0.6
        row_h = 0.9
        gap = 0.15
        for i, item in enumerate(items[:4]):
            y = row_top + i * (row_h + gap)
            # Badge
            self.rounded_rect(slide, row_x, y + (row_h - badge_h) / 2,
                              badge_w, badge_h, self.RED)
            self.text(slide, row_x, y + (row_h - badge_h) / 2 + 0.05,
                      badge_w, badge_h - 0.1,
                      item['badge'], sz=11, color=self.WHITE, bold=True,
                      align=PP_ALIGN.CENTER)
            # Title + body
            self.text(slide, row_x + badge_w + 0.3, y + 0.05,
                      body_w, 0.3,
                      item['title'], sz=12, color=self.DARK, bold=True)
            self.text(slide, row_x + badge_w + 0.3, y + 0.4,
                      body_w, row_h - 0.45,
                      item['body'], sz=10, color=self.MID,
                      shrink=True)

    def four_col_ascending_steps(self, slide, items, base_top=1.2):
        """Four columns in ascending step pattern (left to right, bottom to top).
        items: list of 4 dicts: {title, body}
        Learned from: 四栏模板 Slide 5
        """
        col_w = 2.1
        start_x = 0.5
        gap = 0.35
        step_y = 0.6  # each column is higher than the previous
        for i, item in enumerate(items[:4]):
            x = start_x + i * (col_w + gap)
            y = base_top + (3 - i) * step_y  # ascending: rightmost is highest
            card_h = 1.2 + i * 0.3  # taller cards as they go up
            self.card(slide, x, y, col_w, card_h)
            self.text(slide, x + 0.1, y + 0.1, col_w - 0.2, 0.3,
                      item['title'], sz=12, color=self.RED, bold=True)
            self.text(slide, x + 0.1, y + 0.5, col_w - 0.2, card_h - 0.6,
                      item['body'], sz=10, color=self.DARK,
                      shrink=True)
            # Dot connector at bottom center
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x + col_w / 2 - 0.06), Inches(y + card_h + 0.05),
                Inches(0.12), Inches(0.12))
            dot.fill.solid(); dot.fill.fore_color.rgb = self.RED
            dot.line.fill.background()

    # Canvas bounds constants (10" x 5.625")
    CANVAS_W = 10.0
    CANVAS_H = 5.625
    SAFE_MARGIN = 0.4  # minimum margin from any edge

    # ── Timeline / Process Flow helpers (learned from 时间线流程图模板) ──

    def _clamp(self, val, lo, hi):
        """Clamp a value to [lo, hi] to prevent canvas overflow."""
        return max(lo, min(val, hi))

    def timeline_horizontal_ascending(self, slide, items):
        """Horizontal timeline with nodes, direction arrow, adaptive spacing.
        items: list of 2-6 dicts: {year, title, body}
        """
        n = min(len(items), 6)
        # ── Adaptive layout ──
        left_pad = 1.2          # breathing room on left
        right_pad = 0.9         # room for arrow on right
        arrow_w = 0.25
        spine_x_start = left_pad
        spine_x_end = self.CANVAS_W - right_pad
        spine_w = spine_x_end - spine_x_start

        top_limit = 1.0
        bottom_limit = self.CANVAS_H - self.SAFE_MARGIN
        spine_y = (top_limit + bottom_limit) / 2

        node_inset = 0.15
        first_cx = spine_x_start + node_inset
        last_cx = spine_x_end - node_inset - arrow_w
        node_span = last_cx - first_cx
        node_spacing = node_span / max(n - 1, 1) if n > 1 else 0
        content_w = min(node_spacing * 0.8, 1.6) if n > 1 else 1.6

        # ── Red spine ──
        self.rect(slide, spine_x_start, spine_y, spine_w, 0.025, self.RED)

        # ── Direction arrow at right end ──
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            Inches(spine_x_end - arrow_w), Inches(spine_y - 0.12),
            Inches(arrow_w), Inches(arrow_w * 0.9))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = self.RED
        arrow.line.fill.background()
        arrow.rotation = 90.0

        for i, item in enumerate(items[:n]):
            cx = first_cx + i * node_spacing
            dot_sz = 0.18
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(cx - dot_sz / 2), Inches(spine_y - dot_sz / 2 + 0.012),
                Inches(dot_sz), Inches(dot_sz))
            dot.fill.solid(); dot.fill.fore_color.rgb = self.RED
            dot.line.fill.background()

            content_x = self._clamp(cx - content_w / 2, self.SAFE_MARGIN,
                                    self.CANVAS_W - content_w - self.SAFE_MARGIN)

            if i % 2 == 0:  # ABOVE
                stub_len = 0.2
                badge_h = 0.26; badge_w = 0.75
                stub_bottom = spine_y - dot_sz / 2
                stub_top = stub_bottom - stub_len
                badge_top = stub_top - badge_h - 0.02
                title_top = badge_top - 0.28
                body_top = title_top - 0.35
                if body_top < top_limit:
                    shift = top_limit - body_top
                    body_top += shift; title_top += shift
                    badge_top += shift; stub_top += shift
                self.rect(slide, cx - 0.007, stub_top, 0.014, stub_len, self.BORDER)
                self.rounded_rect(slide, cx - badge_w / 2, badge_top, badge_w, badge_h, self.RED)
                self.text(slide, cx - badge_w / 2, badge_top + 0.02, badge_w, badge_h - 0.04,
                          item['year'], sz=10, color=self.WHITE, bold=True,
                          align=PP_ALIGN.CENTER)
                self.text(slide, content_x, title_top, content_w, 0.24,
                          item['title'], sz=10, color=self.DARK, bold=True,
                          align=PP_ALIGN.CENTER)
                self.text(slide, content_x, body_top, content_w, 0.32,
                          item['body'], sz=10, color=self.MID,
                          align=PP_ALIGN.CENTER, shrink=True)
            else:  # BELOW
                stub_len = 0.2
                badge_h = 0.26; badge_w = 0.75
                stub_top = spine_y + dot_sz / 2
                badge_top = stub_top + stub_len + 0.02
                title_top = badge_top + badge_h + 0.05
                body_top = title_top + 0.28
                body_bottom = body_top + 0.35
                if body_bottom > bottom_limit:
                    shift = body_bottom - bottom_limit
                    body_top -= shift; title_top -= shift; badge_top -= shift
                self.rect(slide, cx - 0.007, stub_top, 0.014, stub_len, self.BORDER)
                self.rounded_rect(slide, cx - badge_w / 2, badge_top, badge_w, badge_h, self.RED)
                self.text(slide, cx - badge_w / 2, badge_top + 0.02, badge_w, badge_h - 0.04,
                          item['year'], sz=10, color=self.WHITE, bold=True,
                          align=PP_ALIGN.CENTER)
                self.text(slide, content_x, title_top, content_w, 0.24,
                          item['title'], sz=10, color=self.DARK, bold=True,
                          align=PP_ALIGN.CENTER)
                self.text(slide, content_x, body_top, content_w, 0.32,
                          item['body'], sz=10, color=self.MID,
                          align=PP_ALIGN.CENTER, shrink=True)

    def process_flow_arrows(self, slide, items):
        """Horizontal process flow with numbered step blocks.
        items: list of 3 dicts: {num, title, body}
        Learned from: 时间线流程图 Slide 2 (optimized for 10" x 5.625")
        """
        n = min(len(items), 5)
        block_w = 7.5 / n
        gap = 0.3
        actual_w = (7.5 - gap * (n - 1)) / n
        start_x = 0.7
        step_y = 2.65  # vertical center
        for i, item in enumerate(items[:n]):
            x = start_x + i * (actual_w + gap)
            # Step number block
            self.rounded_rect(slide, x, step_y, actual_w, 0.45, self.RED)
            self.text(slide, x, step_y + 0.03, actual_w, 0.4,
                      item['num'], sz=14, color=self.WHITE, bold=True,
                      align=PP_ALIGN.CENTER)
            # Arrow between blocks
            if i < n - 1:
                arrow_x = x + actual_w
                self.text(slide, arrow_x + 0.02, step_y + 0.02, gap - 0.04, 0.4,
                          '▶', sz=14, color=self.RED, align=PP_ALIGN.CENTER)
            # Title above
            self.text(slide, x, 1.15, actual_w, 0.3,
                      item['title'], sz=12, color=self.RED, bold=True)
            # Body above step block
            self.text(slide, x, 1.5, actual_w, 1.0,
                      item['body'], sz=10, color=self.DARK,
                      shrink=True)
            # Vertical connector line from content to step block
            self.rect(slide, x + actual_w / 2 - 0.008, 2.5, 0.015, 0.15, self.BORDER)
            # Small label below step block
            self.text(slide, x, step_y + 0.55, actual_w, 0.2,
                      f'STEP {item["num"]}', sz=10, color=self.LIGHT,
                      align=PP_ALIGN.CENTER)

    def process_three_stage(self, slide, items):
        """Three-stage process with large numbers and content cards.
        items: list of 3 dicts: {num, title, body}
        Learned from: 时间线流程图 Slide 3 (optimized for 10" x 5.625")
        """
        n = min(len(items), 4)
        col_w = (8.5 - 0.3 * (n - 1)) / n
        gap = 0.3
        start_x = 0.6
        for i, item in enumerate(items[:n]):
            x = start_x + i * (col_w + gap)
            # Content card (top portion)
            self.card(slide, x, 1.1, col_w, 2.5)
            # Large number at top of card
            self.text(slide, x + 0.1, 1.2, col_w - 0.2, 0.45,
                      item['num'], sz=28, color=self.RED, bold=True)
            # Thin red accent bar
            self.rect(slide, x + 0.1, 1.72, col_w * 0.35, 0.03, self.RED)
            # Title
            self.text(slide, x + 0.1, 1.85, col_w - 0.2, 0.3,
                      item['title'], sz=12, color=self.DARK, bold=True)
            # Body
            self.text(slide, x + 0.1, 2.2, col_w - 0.2, 1.3,
                      item['body'], sz=10, color=self.MID,
                      shrink=True)
            # Bottom accent strip
            self.rect(slide, x, 3.55, col_w, 0.06, self.RED)
            # Arrow to next
            if i < n - 1:
                ax = x + col_w + 0.02
                self.text(slide, ax, 2.2, gap - 0.04, 0.3,
                          '▶', sz=14, color=self.RED, align=PP_ALIGN.CENTER)

    # ── Timeline helper (original) ─────────────────────
    def timeline_item(self, slide, x, y, date, title, desc=""):
        """Single milestone in a vertical timeline."""
        self.rect(slide, x, y, 0.8, 0.28, self.RED)
        self.text(slide, x, y, 0.8, 0.28, date,
                  sz=10, color=self.WHITE, bold=True, align=PP_ALIGN.CENTER)
        self.text(slide, x + 1.0, y, 2.8, 0.22, title,
                  sz=11, color=self.DARK, bold=True)
        if desc:
            self.text(slide, x + 1.0, y + 0.25, 2.8, 0.2, desc,
                      sz=10, color=self.LIGHT)

    # ── Internal utilities ───────────────────────────────
    def _zero_insets(self, textbox):
        bp = textbox.text_frame._txBody.find(qn('a:bodyPr'))
        if bp is None: return
        for attr in ('lIns', 'tIns', 'rIns', 'bIns'):
            bp.set(attr, '0')

    def _set_line_spacing(self, paragraph, pct):
        """Set line spacing as percentage (e.g. 120 = 1.2x)."""
        pPr = paragraph._p.get_or_add_pPr()
        existing = pPr.find(qn('a:lnSpc'))
        if existing is not None:
            pPr.remove(existing)
        spc = parse_xml(
            f'<a:lnSpc xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:spcPct val="{pct * 1000}"/></a:lnSpc>')
        pPr.append(spc)

    # ── Red Base Four-Column Cards (learned from 市场需求 reference) ──

    # Accent color for alternating layouts
    BLUE = RGBColor(0x2B, 0x7A, 0xB8)
    WARM_GREY = RGBColor(0x6B, 0x6B, 0x6B)  # elegant dark grey for alternating

    def timeline_icon_flow(self, slide, items, alt_color=None):
        """Horizontal timeline with large circle icons, chevron arrows, year labels & accent bars.
        items: list of 3-5 dicts: {year, icon_text, title, body}
        alt_color: secondary color for alternating (default: WARM_GREY)
        Layout: large colored circles (red/grey alternating) linked by grey chevron arrows.
        Below each: grey bar with year, colored accent line, title + body text.
        Vertically centered in content area.
        """
        n = min(len(items), 5)
        alt = alt_color or self.WARM_GREY
        colors = [self.RED, alt, self.RED, alt, self.RED]
        margin_x = 0.5
        usable_w = self.CANVAS_W - 2 * margin_x
        col_w = usable_w / n

        # Vertical centering: content block ~3.4" tall, area 1.0"~5.2"
        content_h = 3.4
        area_top = 1.0
        area_bottom = self.CANVAS_H - self.SAFE_MARGIN
        area_h = area_bottom - area_top
        base_y = area_top + (area_h - content_h) / 2  # vertically center

        for i, item in enumerate(items[:n]):
            cx = margin_x + i * col_w + col_w / 2  # center X
            col_left = margin_x + i * col_w + 0.1
            col_inner_w = col_w - 0.2
            color = colors[i]

            # ── Large circle icon
            circ_d = 0.9
            circ_x = cx - circ_d / 2
            circ_y = base_y
            circ = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(circ_x), Inches(circ_y),
                Inches(circ_d), Inches(circ_d))
            circ.fill.solid()
            circ.fill.fore_color.rgb = color
            circ.line.fill.background()
            # Icon text
            self.text(slide, circ_x, circ_y + 0.15, circ_d, circ_d - 0.3,
                      item['icon_text'], sz=14, color=self.WHITE, bold=True,
                      align=PP_ALIGN.CENTER)

            # ── Chevron arrow between circles
            if i < n - 1:
                arrow_x = cx + circ_d / 2 + 0.05
                arrow_w = col_w - circ_d - 0.1
                if arrow_w > 0.1:
                    self.text(slide, arrow_x, circ_y + 0.2, arrow_w, circ_d - 0.4,
                              '\u276f', sz=16, color=self.BORDER, bold=True,
                              align=PP_ALIGN.CENTER)

            # ── Grey bar with year
            bar_y = circ_y + circ_d + 0.2
            bar_h = 0.3
            self.rect(slide, col_left, bar_y, col_inner_w, bar_h, self.BG)
            self.text(slide, col_left + 0.08, bar_y + 0.02, col_inner_w - 0.16, bar_h - 0.04,
                      item['year'], sz=10, color=self.DARK, bold=True)

            # ── Colored accent line
            accent_y = bar_y + bar_h
            self.rect(slide, col_left, accent_y, col_inner_w, 0.04, color)

            # ── Title
            title_y = accent_y + 0.12
            self.text(slide, col_left, title_y, col_inner_w, 0.25,
                      item.get('title', ''), sz=10, color=self.DARK, bold=True)

            # ── Body
            body_y = title_y + 0.3
            body_h = self._clamp(self.CANVAS_H - self.SAFE_MARGIN - body_y, 0.3, 2.0)
            self.text(slide, col_left, body_y, col_inner_w, body_h,
                      item['body'], sz=10, color=self.MID,
                      shrink=True)

    def four_col_red_base_cards(self, slide, items, subtitle=None):
        """Four white cards on a red background base with ring icons & big numbers.
        items: list of 4 dicts: {num, icon_text, title, body}
        subtitle: optional subtitle text below page title
        Layout: upper half white, lower half red. White cards span both zones.
        Each card: ring icon at top, red title, body text, faded number at bottom-right.
        """
        n = min(len(items), 4)

        # Subtitle (integrated, so we can calculate spacing)
        sub_bottom = 0.95  # default: just below title
        if subtitle:
            sub_y = 0.85
            sub_h = 0.3
            self.text(slide, 1.1, sub_y, 7.8, sub_h, subtitle,
                      sz=10, color=self.MID)
            sub_bottom = sub_y + sub_h + 0.1  # 0.1" gap

        # Card top: ensure at least 0.15" below subtitle
        card_top = max(sub_bottom + 0.15, 1.35)

        # Red base rectangle (starts at 55% of card height)
        card_h = self.CANVAS_H - self.SAFE_MARGIN - card_top
        card_h = min(card_h, 3.6)  # cap height
        base_top = card_top + card_h * 0.4
        base_h = self.CANVAS_H - base_top
        self.rect(slide, 0, base_top, self.CANVAS_W, base_h, self.RED)

        # Card dimensions
        card_w = 1.95
        gap = 0.3
        total_w = n * card_w + (n - 1) * gap
        start_x = (self.CANVAS_W - total_w) / 2
        # card_top and card_h already calculated above

        for i, item in enumerate(items[:n]):
            x = start_x + i * (card_w + gap)

            # White card (no border, slight visual lift from red bg)
            self.card(slide, x, card_top, card_w, card_h)

            # Ring icon at top of card
            ring_d = 0.85
            ring_x = x + (card_w - ring_d) / 2
            ring_y = card_top + 0.2
            # Outer ring
            outer = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(ring_x), Inches(ring_y),
                Inches(ring_d), Inches(ring_d))
            outer.fill.solid()
            outer.fill.fore_color.rgb = self.BG
            outer.line.color.rgb = self.BORDER
            outer.line.width = Pt(1)
            # Icon text inside ring
            self.text(slide, ring_x, ring_y + 0.12, ring_d, ring_d - 0.24,
                      item['icon_text'], sz=14, color=self.RED, bold=True,
                      align=PP_ALIGN.CENTER)

            # Title below ring
            title_y = ring_y + ring_d + 0.15
            self.text(slide, x + 0.1, title_y, card_w - 0.2, 0.28,
                      item['title'], sz=11, color=self.RED, bold=True,
                      align=PP_ALIGN.CENTER)

            # Body text
            body_y = title_y + 0.35
            self.text(slide, x + 0.12, body_y, card_w - 0.24, 1.4,
                      item['body'], sz=10, color=self.DARK,
                      shrink=True)

            # Big faded number at bottom-right of card
            num_y = card_top + card_h - 0.55
            self.text(slide, x + card_w - 0.75, num_y, 0.65, 0.45,
                      item['num'], sz=24, color=self.BORDER, bold=True,
                      align=PP_ALIGN.RIGHT)

    # ── Ring Badge Three-Column (learned from 推荐+社交+搜索 reference) ──

    def three_col_numbered_card_blocks(self, slide, items, alt_color=None):
        """Three light-grey cards with half-embedded numbered badges, colored titles & dividers.
        items: list of 3 dicts: {num, title, body}
        alt_color: second accent color (default: ACCENT_GREY). Pattern: RED, alt, RED.
        Layout: badge top-half outside card, left-aligned title + divider + body inside.
        Vertically centered.
        """
        n = min(len(items), 3)
        alt = alt_color or self.ACCENT_GREY
        colors = [self.RED, alt, self.RED]

        card_w = 2.6
        card_h = 2.8
        gap = 0.4
        badge_h = 0.4
        badge_w = 0.55
        embed = badge_h / 2  # half outside, half inside

        total_w = n * card_w + (n - 1) * gap
        start_x = (self.CANVAS_W - total_w) / 2

        # Vertical center (account for badge overhang)
        content_h = card_h + embed
        area_top = 1.0
        area_bottom = self.CANVAS_H - self.SAFE_MARGIN
        area_h = area_bottom - area_top
        base_y = area_top + (area_h - content_h) / 2
        card_top = base_y + embed  # card starts below badge overhang

        for i, item in enumerate(items[:n]):
            x = start_x + i * (card_w + gap)
            color = colors[i]
            pad = 0.2  # internal padding

            # ── Light grey card
            card = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(card_top),
                Inches(card_w), Inches(card_h))
            card.fill.solid()
            card.fill.fore_color.rgb = self.BG  # light grey
            card.line.fill.background()

            # ── Half-embedded badge (straddles card top edge)
            badge_x = x + pad
            badge_y = card_top - embed
            badge = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(badge_x), Inches(badge_y),
                Inches(badge_w), Inches(badge_h))
            badge.fill.solid()
            badge.fill.fore_color.rgb = color
            badge.line.fill.background()
            # Badge number
            self.text(slide, badge_x, badge_y + 0.02, badge_w, badge_h - 0.04,
                      item['num'], sz=14, color=self.WHITE, bold=True,
                      align=PP_ALIGN.CENTER)

            # ── Title (colored, left-aligned)
            title_y = card_top + embed + 0.15
            self.text(slide, x + pad, title_y, card_w - 2 * pad, 0.3,
                      item['title'], sz=13, color=color, bold=True)

            # ── Divider line (colored)
            div_y = title_y + 0.38
            self.rect(slide, x + pad, div_y, card_w - 2 * pad, 0.02, color)

            # ── Body text
            body_y = div_y + 0.15
            body_h = card_top + card_h - body_y - 0.15
            body_h = max(body_h, 0.3)
            self.text(slide, x + pad, body_y, card_w - 2 * pad, body_h,
                      item['body'], sz=10, color=self.DARK,
                      shrink=True)

    # ── Triple Overlapping Circles (learned from 三圆交叠 reference) ──

    def triple_overlap_circles(self, slide, items):
        """Three overlapping circles in inverted triangle with text labels outside.
        items: list of 3 dicts: {icon_text, title, body}
        Layout: circles[0]=top-left, circles[1]=top-right, circles[2]=bottom-center.
        Text: left-side right-aligned, right-side left-aligned, bottom centered.
        Colors: RED, CORAL, DARK_RED (warm red family).
        """
        import math
        colors = [self.RED, self.CORAL, self.DARK_RED]
        circ_d = 1.6
        overlap = circ_d * 0.3  # 30% overlap

        # Triangle geometry
        cx_page = self.CANVAS_W / 2
        cy_page = 2.8  # vertical center of circles group
        h_offset = (circ_d - overlap) / 2 + 0.1
        v_offset = (circ_d - overlap) / 2

        positions = [
            (cx_page - h_offset, cy_page - v_offset),  # top-left
            (cx_page + h_offset, cy_page - v_offset),  # top-right
            (cx_page, cy_page + v_offset),              # bottom-center
        ]

        # Draw circles (back to front: bottom, then top-left, then top-right)
        draw_order = [2, 0, 1]
        for di in draw_order:
            pcx, pcy = positions[di]
            color = colors[di]
            circ = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(pcx - circ_d / 2), Inches(pcy - circ_d / 2),
                Inches(circ_d), Inches(circ_d))
            circ.fill.solid()
            circ.fill.fore_color.rgb = color
            circ.line.fill.background()
            # Icon text inside
            self.text(slide,
                      pcx - circ_d / 2, pcy - 0.2,
                      circ_d, 0.4,
                      items[di]['icon_text'],
                      sz=16, color=self.WHITE, bold=True,
                      align=PP_ALIGN.CENTER)

        # Text labels outside circles
        txt_w = 2.8
        txt_h_title = 0.25
        txt_h_body = 0.8

        # [0] Top-left: text on LEFT side, right-aligned
        lx = positions[0][0] - circ_d / 2 - txt_w - 0.15
        ly = positions[0][1] - 0.5
        self.text(slide, lx, ly, txt_w, txt_h_title,
                  items[0]['title'], sz=11, color=colors[0], bold=True,
                  align=PP_ALIGN.RIGHT)
        self.text(slide, lx, ly + 0.3, txt_w, txt_h_body,
                  items[0]['body'], sz=10, color=self.MID,
                  align=PP_ALIGN.RIGHT, shrink=True)

        # [1] Top-right: text on RIGHT side, left-aligned
        rx = positions[1][0] + circ_d / 2 + 0.15
        ry = positions[1][1] - 0.5
        self.text(slide, rx, ry, txt_w, txt_h_title,
                  items[1]['title'], sz=11, color=colors[1], bold=True)
        self.text(slide, rx, ry + 0.3, txt_w, txt_h_body,
                  items[1]['body'], sz=10, color=self.MID,
                  shrink=True)

        # [2] Bottom: text BELOW, centered
        bx = positions[2][0] - txt_w / 2
        by = positions[2][1] + circ_d / 2 + 0.1
        self.text(slide, bx, by, txt_w, txt_h_title,
                  items[2]['title'], sz=11, color=colors[2], bold=True,
                  align=PP_ALIGN.CENTER)
        self.text(slide, bx, by + 0.3, txt_w, txt_h_body,
                  items[2]['body'], sz=10, color=self.MID,
                  align=PP_ALIGN.CENTER, shrink=True)

    # ── Diamond Quadrant Layout (learned from 菱形四象限 reference) ──

    def four_col_numbered_dividers(self, slide, items, alt_color=None):
        """Four equal columns with numbered badges, colored titles, vertical dividers.
        items: list of 4 dicts: {num, title, body}
        alt_color: second accent (default: ACCENT_GREY). Pattern: RED, alt, RED, alt.
        Divider lines between columns match alternating colors.
        Vertically centered.
        """
        n = min(len(items), 4)
        alt = alt_color or self.ACCENT_GREY
        colors = [self.RED, alt, self.RED, alt]
        div_colors = [alt, self.RED, alt]  # dividers between columns

        margin_x = 0.8
        usable_w = self.CANVAS_W - 2 * margin_x
        col_w = usable_w / n

        # Vertical centering
        content_h = 3.0  # badge(0.4) + gap(0.25) + title(0.3) + gap(0.15) + body(1.5) + padding
        area_top = 1.0
        area_bottom = self.CANVAS_H - self.SAFE_MARGIN
        area_h = area_bottom - area_top
        base_y = area_top + (area_h - content_h) / 2

        badge_size = 0.4

        for i, item in enumerate(items[:n]):
            cx = margin_x + i * col_w + col_w / 2
            col_left = margin_x + i * col_w
            color = colors[i]

            # ── Number badge (small colored square)
            badge_x = cx - badge_size / 2
            badge_y = base_y
            badge = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(badge_x), Inches(badge_y),
                Inches(badge_size), Inches(badge_size))
            badge.fill.solid()
            badge.fill.fore_color.rgb = color
            badge.line.fill.background()
            self.text(slide, badge_x, badge_y + 0.02, badge_size, badge_size - 0.04,
                      item['num'], sz=14, color=self.WHITE, bold=True,
                      align=PP_ALIGN.CENTER)

            # ── Title (colored, centered)
            title_y = badge_y + badge_size + 0.2
            self.text(slide, col_left + 0.1, title_y, col_w - 0.2, 0.3,
                      item['title'], sz=11, color=color, bold=True,
                      align=PP_ALIGN.CENTER)

            # ── Body (dark, centered)
            body_y = title_y + 0.4
            body_h = content_h - (body_y - base_y) - 0.1
            body_h = max(body_h, 0.5)
            self.text(slide, col_left + 0.15, body_y, col_w - 0.3, body_h,
                      item['body'], sz=10, color=self.DARK,
                      align=PP_ALIGN.CENTER, shrink=True)

            # ── Vertical divider line (between columns, not after last)
            if i < n - 1:
                div_x = margin_x + (i + 1) * col_w
                div_top = title_y
                div_bottom = body_y + body_h
                div_color = div_colors[i] if i < len(div_colors) else alt
                self.rect(slide, div_x - 0.005, div_top, 0.01, div_bottom - div_top, div_color)

    def diamond_quadrants(self, slide, center_text, items):
        """Central diamond with 4 surrounding text quadrants.
        center_text: text inside the diamond
        items: list of 4 dicts: {title, body} for TL, TR, BL, BR
        Text alignment: TL/BL right-aligned (toward center), TR/BR left-aligned.
        """
        from pptx.enum.shapes import MSO_SHAPE

        # Page center for content area
        cx = self.CANVAS_W / 2
        cy = 3.0  # slightly below vertical center to account for title

        # ── Diamond (rotated square)
        dia_size = 1.5  # side length before rotation
        dia = slide.shapes.add_shape(
            MSO_SHAPE.DIAMOND,
            Inches(cx - dia_size / 2), Inches(cy - dia_size / 2),
            Inches(dia_size), Inches(dia_size))
        dia.fill.solid()
        dia.fill.fore_color.rgb = self.WHITE
        dia.line.color.rgb = self.RED
        dia.line.width = Pt(2.5)
        # Center text
        self.text(slide, cx - dia_size / 2 + 0.15, cy - 0.15,
                  dia_size - 0.3, 0.3,
                  center_text, sz=12, color=self.RED, bold=True,
                  align=PP_ALIGN.CENTER)

        # ── Four quadrant text blocks
        txt_w = 2.8
        txt_h = 1.2
        gap = 0.25  # gap from diamond edge
        d_half = dia_size / 2 + gap

        # Quadrant positions: (x, y, align)
        quads = [
            # TL: right-aligned, above-left of diamond
            (cx - d_half - txt_w, cy - d_half - txt_h + 0.4, PP_ALIGN.RIGHT),
            # TR: left-aligned, above-right
            (cx + d_half, cy - d_half - txt_h + 0.4, PP_ALIGN.LEFT),
            # BL: right-aligned, below-left
            (cx - d_half - txt_w, cy + d_half - 0.4, PP_ALIGN.RIGHT),
            # BR: left-aligned, below-right
            (cx + d_half, cy + d_half - 0.4, PP_ALIGN.LEFT),
        ]

        for i, item in enumerate(items[:4]):
            qx, qy, align = quads[i]
            # Light grey background block
            bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(qx), Inches(qy),
                Inches(txt_w), Inches(txt_h))
            bg.fill.solid()
            bg.fill.fore_color.rgb = self.BG
            bg.line.fill.background()

            # Title
            self.text(slide, qx + 0.1, qy + 0.1, txt_w - 0.2, 0.25,
                      item.get('title', ''), sz=11, color=self.RED, bold=True,
                      align=align)
            # Body
            self.text(slide, qx + 0.1, qy + 0.4, txt_w - 0.2, txt_h - 0.5,
                      item['body'], sz=10, color=self.DARK,
                      align=align, shrink=True)

    # ── Parallelogram Four-Column (learned from 斜四栏 reference) ──

    LIGHT_GREY2 = RGBColor(0xEE, 0xEE, 0xEE)  # very light grey

    def four_panel_on_platform(self, slide, platform_text, items, alt_bg=None):
        """Four panels arranged above an elliptical platform base.
        platform_text: main title on the elliptical base
        items: list of 4 dicts: {title, body}
        Layout: bottom ellipse as "base", four tall cards above in slight fan.
        Colors: alternating light-red / light-grey.
        """
        n = min(len(items), 4)
        alt = alt_bg or self.LIGHT_GREY2
        bg_colors = [self.LIGHT_RED, alt, self.LIGHT_RED, alt]
        title_colors = [self.RED, self.ACCENT_GREY, self.RED, self.ACCENT_GREY]

        # ── Vertical layout: panels + platform need to be centered
        panel_h = 2.5
        plat_h = 1.0
        gap_panel_plat = 0.1  # overlap amount
        total_content_h = panel_h + plat_h - gap_panel_plat

        area_top = 1.0  # below title
        area_bottom = self.CANVAS_H - self.SAFE_MARGIN
        area_h = area_bottom - area_top
        base_y = area_top + (area_h - total_content_h) / 2

        panel_top = base_y
        panel_bottom = panel_top + panel_h

        # ── Elliptical platform
        plat_w = 6.0
        plat_x = (self.CANVAS_W - plat_w) / 2
        plat_y = panel_bottom - gap_panel_plat
        plat = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(plat_x), Inches(plat_y),
            Inches(plat_w), Inches(plat_h))
        plat.fill.solid()
        plat.fill.fore_color.rgb = self.LIGHT_RED2
        plat.line.fill.background()
        # Platform text
        self.text(slide, plat_x + 0.5, plat_y + 0.2, plat_w - 1.0, 0.5,
                  platform_text, sz=16, color=self.RED, bold=True,
                  align=PP_ALIGN.CENTER)

        # ── Four panels above platform
        panel_w = 1.8
        gap = 0.25
        total_panels_w = n * panel_w + (n - 1) * gap
        start_x = (self.CANVAS_W - total_panels_w) / 2

        for i, item in enumerate(items[:n]):
            x = start_x + i * (panel_w + gap)
            bg_color = bg_colors[i]
            t_color = title_colors[i]

            # Panel card (rounded rectangle)
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(panel_top),
                Inches(panel_w), Inches(panel_h))
            card.fill.solid()
            card.fill.fore_color.rgb = bg_color
            card.line.fill.background()

            # Title inside panel (top area)
            self.text(slide, x + 0.12, panel_top + 0.2, panel_w - 0.24, 0.3,
                      item['title'], sz=11, color=t_color, bold=True,
                      align=PP_ALIGN.CENTER)

            # Body inside panel
            body_y = panel_top + 0.6
            body_h = panel_h - 0.8
            self.text(slide, x + 0.12, body_y, panel_w - 0.24, body_h,
                      item['body'], sz=10, color=self.DARK,
                      align=PP_ALIGN.CENTER, shrink=True)

    def four_col_parallelogram(self, slide, items, alt_bg=None):
        """Four slanted parallelogram columns with circle icons on top.
        items: list of 4 dicts: {icon_text, title, body}
        Alternating light-red / light-grey background.
        Circle icon at top in deeper color, title + body below.
        """
        n = min(len(items), 4)
        alt = alt_bg or self.LIGHT_GREY2
        bg_colors = [self.LIGHT_RED, alt, self.LIGHT_RED, alt]
        icon_colors = [self.RED, self.ACCENT_GREY, self.RED, self.ACCENT_GREY]

        margin_x = 0.6
        usable_w = self.CANVAS_W - 2 * margin_x
        gap = 0.25
        col_w = (usable_w - (n - 1) * gap) / n
        col_h = 3.2

        # Vertical centering
        area_top = 1.0
        area_bottom = self.CANVAS_H - self.SAFE_MARGIN
        area_h = area_bottom - area_top
        circ_overhang = 0.25  # circle sticks out above parallelogram
        total_h = col_h + circ_overhang
        base_y = area_top + (area_h - total_h) / 2 + circ_overhang

        for i, item in enumerate(items[:n]):
            x = margin_x + i * (col_w + gap)
            cx = x + col_w / 2
            bg_color = bg_colors[i]
            ic_color = icon_colors[i]

            # ── Parallelogram (using PARALLELOGRAM shape)
            para = slide.shapes.add_shape(
                MSO_SHAPE.PARALLELOGRAM,
                Inches(x), Inches(base_y),
                Inches(col_w), Inches(col_h))
            para.fill.solid()
            para.fill.fore_color.rgb = bg_color
            para.line.fill.background()
            # Adjust the parallelogram skew via adjustments
            if para.adjustments:
                try:
                    para.adjustments[0] = 0.15  # gentle skew
                except:
                    pass

            # ── Circle icon at top (overlapping parallelogram top edge)
            circ_d = 0.7
            circ_x = cx - circ_d / 2
            circ_y = base_y - circ_overhang
            circ = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(circ_x), Inches(circ_y),
                Inches(circ_d), Inches(circ_d))
            circ.fill.solid()
            circ.fill.fore_color.rgb = ic_color
            circ.line.fill.background()
            self.text(slide, circ_x, circ_y + 0.1, circ_d, circ_d - 0.2,
                      item['icon_text'], sz=14, color=self.WHITE, bold=True,
                      align=PP_ALIGN.CENTER)

            # ── Title
            title_y = base_y + 0.55
            self.text(slide, x + 0.15, title_y, col_w - 0.3, 0.3,
                      item['title'], sz=11, color=ic_color, bold=True,
                      align=PP_ALIGN.CENTER)

            # ── Body
            body_y = title_y + 0.4
            body_h = col_h - (body_y - base_y) - 0.2
            body_h = max(body_h, 0.5)
            self.text(slide, x + 0.15, body_y, col_w - 0.3, body_h,
                      item['body'], sz=10, color=self.DARK,
                      align=PP_ALIGN.CENTER, shrink=True)

    def timeline_zigzag(self, slide, items):
        """Horizontal timeline with numbered circles and text blocks alternating above/below.
        items: list of 3-5 dicts: {num, title, body}
        Odd items (01/03) have text above, even items (02/04) below.
        Connected by a thin horizontal line through circle centers.
        Learned from: user-shared reference image 2026-04-14.
        """
        n = min(len(items), 5)
        margin_x = 0.6
        usable_w = self.CANVAS_W - 2 * margin_x
        col_w = usable_w / n

        # Vertical layout constants
        circ_d = 0.65
        line_y = self.CANVAS_H / 2 + 0.15  # horizontal line Y (center of circles)
        circ_top = line_y - circ_d / 2
        text_w = col_w - 0.3
        text_h_above = 1.3  # space for text block above
        text_h_below = 1.3  # space for text block below

        # ── Horizontal connecting line (behind circles)
        line_x1 = margin_x + col_w / 2
        line_x2 = margin_x + (n - 1) * col_w + col_w / 2
        line_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(line_x1), Inches(line_y - 0.015),
            Inches(line_x2 - line_x1), Inches(0.03))
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = self.BORDER
        line_shape.line.fill.background()

        for i, item in enumerate(items[:n]):
            cx = margin_x + i * col_w + col_w / 2
            col_left = cx - text_w / 2
            num = item.get('num', f'{i+1:02d}')
            is_above = (i % 2 == 0)  # 0,2,4 = above; 1,3 = below
            color = self.RED if (i % 2 == 0) else self.DARK

            # ── Numbered circle
            circ = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(cx - circ_d / 2), Inches(circ_top),
                Inches(circ_d), Inches(circ_d))
            circ.fill.solid()
            circ.fill.fore_color.rgb = color
            circ.line.fill.background()
            self.text(slide, cx - circ_d / 2, circ_top + 0.08, circ_d, circ_d - 0.16,
                      num, sz=14, color=self.WHITE, bold=True,
                      align=PP_ALIGN.CENTER)

            # ── Text block (title + body) — directly adjacent to circle
            gap = 0.15  # small gap between circle and text
            if is_above:
                title_y = circ_top - gap - text_h_above
                body_y = title_y + 0.35
                body_h = text_h_above - 0.4
            else:
                title_y = circ_top + circ_d + gap
                body_y = title_y + 0.35
                body_h = text_h_below - 0.4

            # Title
            self.text(slide, col_left, title_y, text_w, 0.3,
                      item.get('title', ''), sz=12, color=color, bold=True,
                      align=PP_ALIGN.CENTER)
            # Body
            self.text(slide, col_left, body_y, text_w, body_h,
                      item['body'], sz=10, color=self.MID,
                      align=PP_ALIGN.CENTER, shrink=True)

    # ── Save ─────────────────────────────────────────────
    def save(self, path):
        self.prs.save(path)
        return path
