#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""ppt_scorer.py — Per-page scoring & deck-level evaluation system.

Three layers:
1. Page Score: 6 dimensions, 100 points per page
2. Deck Score: 4 dimensions, 100 points total
3. Learning Loop: feedback recording + best practice capture

Usage:
    from ppt_scorer import PPTScorer
    scorer = PPTScorer("output.pptx")
    report = scorer.full_score_report()
    print(report)
"""
import re
import math
from collections import Counter
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor


class PPTScorer:
    """Per-page and deck-level PPT scoring system."""

    CANVAS_W = 10.0
    CANVAS_H = 5.625
    SAFE_MARGIN = 0.4
    CONTENT_TOP = 1.0  # below title
    CONTENT_BOTTOM = 5.225  # CANVAS_H - SAFE_MARGIN

    # Template colors (for color harmony check)
    TEMPLATE_COLORS = {
        (0xC0, 0x00, 0x00),  # RED
        (0xBB, 0x00, 0x00),  # RED2
        (0x33, 0x33, 0x33),  # DARK
        (0x66, 0x66, 0x66),  # MID
        (0x99, 0x99, 0x99),  # LIGHT
        (0xF5, 0xF5, 0xF5),  # BG
        (0xFF, 0xFF, 0xFF),  # WHITE
        (0xE0, 0xE0, 0xE0),  # BORDER
        (0x6B, 0x6B, 0x6B),  # ACCENT_GREY
        (0x8B, 0x1A, 0x1A),  # DARK_RED
        (0xE0, 0x6C, 0x5C),  # CORAL
        (0xFC, 0xE4, 0xE4),  # LIGHT_RED
        (0xEE, 0xEE, 0xEE),  # LIGHT_GREY2
        (0xF5, 0xD5, 0xD5),  # LIGHT_RED2
        (0x2B, 0x7A, 0xB8),  # BLUE (legacy)
        (0x90, 0x00, 0x00),  # darker red borders
        (0xA0, 0x00, 0x00),  # darker red
        (0xE8, 0x50, 0x50),  # lighter red (ring)
        (0xE0, 0xC0, 0xC0),  # light pink halo
        (0x00, 0x00, 0x00),  # pure black (rare but valid)
    }

    def __init__(self, pptx_path):
        self.path = pptx_path
        self.prs = Presentation(pptx_path)
        self.slides = list(self.prs.slides)
        self.n_slides = len(self.slides)

    # ══════════════════════════════════════════════════════
    # LAYER 1: Per-Page Scoring (6 dimensions, 100 points)
    # ══════════════════════════════════════════════════════

    def score_page(self, slide_idx):
        """Score a single page across 6 dimensions. Returns dict with scores + suggestions."""
        slide = self.slides[slide_idx]
        is_cover = (slide_idx == 0 or slide_idx == self.n_slides - 1)

        shapes_data = self._extract_shapes(slide)
        text_shapes = [s for s in shapes_data if s['text']]
        all_text = ' '.join(s['text'] for s in text_shapes)
        char_count = len(all_text.replace(' ', '').replace('\n', ''))

        scores = {}
        suggestions = []

        # ── D1: Structure (20 pts) ──
        scores['结构合理'] = self._score_structure(shapes_data, is_cover, suggestions)

        # ── D2: Whitespace (20 pts) ──
        scores['留白呼吸'] = self._score_whitespace(shapes_data, is_cover, suggestions)

        # ── D3: Information Load (20 pts) ──
        scores['信息承载'] = self._score_info_load(char_count, is_cover, suggestions)

        # ── D4: Visual Hierarchy (20 pts) ──
        scores['视觉层次'] = self._score_visual_hierarchy(text_shapes, suggestions)

        # ── D5: Color Harmony (10 pts) ──
        scores['配色协调'] = self._score_color_harmony(shapes_data, suggestions)

        # ── D6: Layout Fit (10 pts) ──
        scores['布局适配'] = self._score_layout_fit(shapes_data, char_count, suggestions)

        total = sum(scores.values())
        return {
            'slide': slide_idx + 1,
            'total': total,
            'scores': scores,
            'suggestions': suggestions,
            'is_cover': is_cover,
            'char_count': char_count,
        }

    def _extract_shapes(self, slide):
        """Extract shape metadata from a slide."""
        shapes = []
        for shape in slide.shapes:
            t = Emu(shape.top).inches
            l = Emu(shape.left).inches
            w = Emu(shape.width).inches
            h = Emu(shape.height).inches
            text = ''
            font_sizes = []
            colors = []
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            font_sizes.append(Emu(run.font.size).pt if hasattr(Emu(run.font.size), 'pt') else run.font.size / 12700)
                        if run.font.color and run.font.color.type is not None:
                            try:
                                c = run.font.color.rgb
                                if c:
                                    colors.append((int(str(c)[0:2], 16),
                                                  int(str(c)[2:4], 16),
                                                  int(str(c)[4:6], 16)))
                            except (AttributeError, ValueError):
                                pass
            # Shape fill color
            fill_color = None
            try:
                if shape.fill and shape.fill.fore_color and shape.fill.fore_color.rgb:
                    fc = shape.fill.fore_color.rgb
                    fill_color = (fc[0] if isinstance(fc[0], int) else int(str(fc)[0:2], 16),
                                 fc[1] if isinstance(fc[1], int) else int(str(fc)[2:4], 16),
                                 fc[2] if isinstance(fc[2], int) else int(str(fc)[4:6], 16))
            except:
                pass

            shapes.append({
                'top': t, 'left': l, 'width': w, 'height': h,
                'bottom': t + h, 'right': l + w,
                'text': text, 'font_sizes': font_sizes,
                'colors': colors, 'fill_color': fill_color,
            })
        return shapes

    # ── Scoring Functions ─────────────────────────────────

    def _score_structure(self, shapes, is_cover, suggestions):
        """D1: Structure (20 pts) — vertical centering, spacing uniformity."""
        if is_cover:
            return 18  # covers are manually placed, light check

        score = 20
        content_shapes = [s for s in shapes if s['top'] >= 0.8]
        if not content_shapes:
            return 15

        # Vertical centering check
        content_top = min(s['top'] for s in content_shapes)
        content_bottom = max(s['bottom'] for s in content_shapes)
        content_h = content_bottom - content_top
        area_h = self.CONTENT_BOTTOM - self.CONTENT_TOP

        if content_h < area_h * 0.55:
            space_above = content_top - self.CONTENT_TOP
            space_below = self.CONTENT_BOTTOM - content_bottom
            if space_above > 0:
                ratio = space_below / space_above
            else:
                ratio = 10
            if ratio > 4 or ratio < 0.25:
                score -= 8
                suggestions.append('垂直严重偏移，建议居中放置内容')
            elif ratio > 2.5 or ratio < 0.4:
                score -= 4
                suggestions.append('垂直略偏，可微调居中')
            elif ratio > 1.8 or ratio < 0.55:
                score -= 2
                suggestions.append('垂直可再优化')

        # Element spacing uniformity (check if gaps are roughly even)
        if len(content_shapes) >= 3:
            tops = sorted(set(round(s['top'], 1) for s in content_shapes))
            if len(tops) >= 3:
                gaps = [tops[i+1] - tops[i] for i in range(len(tops)-1)]
                if gaps:
                    avg_gap = sum(gaps) / len(gaps)
                    max_deviation = max(abs(g - avg_gap) for g in gaps)
                    if max_deviation > avg_gap * 0.8 and avg_gap > 0.3:
                        score -= 3
                        suggestions.append('元素间距不均匀')

        return max(score, 0)

    def _score_whitespace(self, shapes, is_cover, suggestions):
        """D2: Whitespace (20 pts) — content fills 40-75% of available area."""
        if is_cover:
            return 18

        score = 20
        content_shapes = [s for s in shapes if s['top'] >= 0.8]
        if not content_shapes:
            return 10

        # Calculate bounding box area of content
        content_top = min(s['top'] for s in content_shapes)
        content_bottom = max(s['bottom'] for s in content_shapes)
        content_left = min(s['left'] for s in content_shapes)
        content_right = max(s['right'] for s in content_shapes)
        content_area = (content_bottom - content_top) * (content_right - content_left)

        avail_area = (self.CONTENT_BOTTOM - self.CONTENT_TOP) * (self.CANVAS_W - 2 * self.SAFE_MARGIN)
        fill_ratio = content_area / avail_area if avail_area > 0 else 0

        if fill_ratio < 0.25:
            score -= 10
            suggestions.append(f'内容过少(填充{fill_ratio:.0%})，页面空旷')
        elif fill_ratio < 0.40:
            score -= 5
            suggestions.append(f'留白偏多(填充{fill_ratio:.0%})，可增加内容或放大元素')
        elif fill_ratio > 0.85:
            score -= 8
            suggestions.append(f'内容过密(填充{fill_ratio:.0%})，缺少呼吸空间')
        elif fill_ratio > 0.75:
            score -= 3
            suggestions.append(f'略显拥挤(填充{fill_ratio:.0%})，可精简或拆页')

        return max(score, 0)

    def _score_info_load(self, char_count, is_cover, suggestions):
        """D3: Information Load (20 pts) — 60-220 chars ideal for content slides."""
        if is_cover:
            return 18 if char_count < 80 else 15

        score = 20
        if char_count < 30:
            score -= 10
            suggestions.append(f'文字过少({char_count}字)，内容薄弱')
        elif char_count < 60:
            score -= 5
            suggestions.append(f'文字偏少({char_count}字)，可补充要点')
        elif char_count > 350:
            score -= 10
            suggestions.append(f'文字过多({char_count}字)，建议拆页或精简')
        elif char_count > 250:
            score -= 5
            suggestions.append(f'文字偏多({char_count}字)，可精简')
        elif char_count > 220:
            score -= 2
            suggestions.append(f'文字量偏高({char_count}字)')

        return max(score, 0)

    def _score_visual_hierarchy(self, text_shapes, suggestions):
        """D4: Visual Hierarchy (20 pts) — font size gradient, color layers."""
        score = 20

        # Collect all font sizes
        all_sizes = []
        for s in text_shapes:
            all_sizes.extend(s['font_sizes'])

        if not all_sizes:
            return 12

        unique_sizes = sorted(set(round(s) for s in all_sizes), reverse=True)

        # Font size gradient: should have at least 2 distinct sizes
        if len(unique_sizes) < 2:
            score -= 8
            suggestions.append('字号单一，缺少层次(标题/正文应有大小区分)')
        elif len(unique_sizes) < 3:
            score -= 3
            suggestions.append('字号层次可再丰富(考虑标题/副标题/正文三级)')

        # Size ratio: largest should be at least 1.5x smallest
        if len(unique_sizes) >= 2:
            ratio = unique_sizes[0] / unique_sizes[-1] if unique_sizes[-1] > 0 else 1
            if ratio < 1.3:
                score -= 4
                suggestions.append('最大字号与最小字号差异不够，层次感弱')

        # Color diversity: should have at least 2 distinct text colors
        all_colors = []
        for s in text_shapes:
            all_colors.extend(s['colors'])
        unique_colors = set(all_colors)
        if len(unique_colors) < 2:
            score -= 4
            suggestions.append('文字颜色单一，缺少视觉重点')

        return max(score, 0)

    def _score_color_harmony(self, shapes, suggestions):
        """D5: Color Harmony (10 pts) — colors within template palette."""
        score = 10

        all_fill_colors = []
        for s in shapes:
            if s['fill_color']:
                all_fill_colors.append(s['fill_color'])

        if not all_fill_colors:
            return 8  # no fills to check

        off_palette = 0
        for c in all_fill_colors:
            # Check if color is close to any template color (tolerance 30)
            min_dist = min(
                math.sqrt((c[0]-tc[0])**2 + (c[1]-tc[1])**2 + (c[2]-tc[2])**2)
                for tc in self.TEMPLATE_COLORS
            )
            if min_dist > 50:  # significantly off-palette
                off_palette += 1

        if off_palette > 3:
            score -= 6
            suggestions.append(f'{off_palette}个元素使用了模板外颜色')
        elif off_palette > 1:
            score -= 3
            suggestions.append(f'{off_palette}个元素颜色偏离模板色系')
        elif off_palette == 1:
            score -= 1

        return max(score, 0)

    def _score_layout_fit(self, shapes, char_count, suggestions):
        """D6: Layout Fit (10 pts) — layout type matches content volume."""
        score = 10
        n_shapes = len(shapes)

        # Heuristic: too many shapes for little text = over-decorated
        if n_shapes > 25 and char_count < 80:
            score -= 4
            suggestions.append('元素多但文字少，可能过度装饰')

        # Too few shapes for lots of text = wall of text
        if n_shapes < 5 and char_count > 200:
            score -= 4
            suggestions.append('文字多但元素少，缺少结构化展示')

        return max(score, 0)

    # ══════════════════════════════════════════════════════
    # LAYER 2: Deck Score (4 dimensions, 100 points)
    # ══════════════════════════════════════════════════════

    def score_deck(self, page_scores=None):
        """Score the entire deck. Returns dict."""
        if page_scores is None:
            page_scores = [self.score_page(i) for i in range(self.n_slides)]

        scores = {}
        suggestions = []

        # ── Narrative Flow (25 pts): titles tell a story ──
        titles = []
        for ps in page_scores:
            slide = self.slides[ps['slide'] - 1]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = Emu(shape.top).inches
                    h = Emu(shape.height).inches
                    if t < 0.6 and h < 0.6:
                        titles.append(shape.text_frame.text.strip())
                        break
        score_narrative = 25
        if len(titles) < 3:
            score_narrative -= 5
        # Check for duplicate titles
        title_counter = Counter(titles)
        dupes = sum(1 for t, c in title_counter.items() if c > 1 and t)
        if dupes > 0:
            score_narrative -= 5 * dupes
            suggestions.append(f'{dupes}个重复标题，叙事线不清晰')
        scores['叙事线'] = max(score_narrative, 0)

        # ── Rhythm (25 pts): density variation ──
        densities = [ps['char_count'] for ps in page_scores if not ps['is_cover']]
        score_rhythm = 25
        if densities:
            avg_d = sum(densities) / len(densities)
            std_d = math.sqrt(sum((d - avg_d)**2 for d in densities) / len(densities)) if len(densities) > 1 else 0
            cv = std_d / avg_d if avg_d > 0 else 0
            if cv < 0.1 and len(densities) > 3:
                score_rhythm -= 8
                suggestions.append('各页信息量过于均匀，缺少节奏感(建议有疏有密)')
            elif cv < 0.15 and len(densities) > 3:
                score_rhythm -= 4
                suggestions.append('信息节奏可再丰富')
        scores['节奏感'] = max(score_rhythm, 0)

        # ── Diversity (25 pts): layout variety ──
        from ppt_qa import PPTQualityChecker
        qc = PPTQualityChecker(self.path)
        layout_issues = qc.check_layout_diversity()
        score_diversity = 25
        score_diversity -= 5 * len(layout_issues)
        if layout_issues:
            for li in layout_issues:
                suggestions.append(li)
        scores['多样性'] = max(score_diversity, 0)

        # ── Average Page Score (25 pts) ──
        content_scores = [ps['total'] for ps in page_scores if not ps['is_cover']]
        avg_score = sum(content_scores) / len(content_scores) if content_scores else 50
        scores['页均分'] = min(25, int(avg_score / 4))  # scale 100 -> 25

        total = sum(scores.values())
        return {
            'total': total,
            'scores': scores,
            'suggestions': suggestions,
        }

    # ══════════════════════════════════════════════════════
    # LAYER 3: Learning Loop (feedback + best practice)
    # ══════════════════════════════════════════════════════

    @staticmethod
    def record_feedback(memory_dir, layout_name, issue, fix):
        """Record feedback to memory/ppt-feedback.md."""
        import os
        from datetime import datetime
        feedback_path = os.path.join(memory_dir, 'ppt-feedback.md')

        # Create file if needed
        if not os.path.exists(feedback_path):
            with open(feedback_path, 'w', encoding='utf-8') as f:
                f.write('# PPT Feedback Log\n\n')
                f.write('| 日期 | 布局 | 问题 | 修复 |\n')
                f.write('|------|------|------|------|\n')

        date_str = datetime.now().strftime('%Y-%m-%d')
        line = f'| {date_str} | {layout_name} | {issue} | {fix} |\n'
        with open(feedback_path, 'a', encoding='utf-8') as f:
            f.write(line)

    @staticmethod
    def record_best_practice(memory_dir, slide_num, score, layout_name, char_count, notes=''):
        """Record high-scoring pages as best practices."""
        import os
        from datetime import datetime
        bp_path = os.path.join(memory_dir, 'ppt-best-practices.md')

        if not os.path.exists(bp_path):
            with open(bp_path, 'w', encoding='utf-8') as f:
                f.write('# PPT Best Practices\n\n')
                f.write('| 日期 | 页面 | 分数 | 布局 | 字数 | 备注 |\n')
                f.write('|------|------|------|------|------|------|\n')

        date_str = datetime.now().strftime('%Y-%m-%d')
        line = f'| {date_str} | S{slide_num} | {score}/100 | {layout_name} | {char_count} | {notes} |\n'
        with open(bp_path, 'a', encoding='utf-8') as f:
            f.write(line)

    # ══════════════════════════════════════════════════════
    # Full Score Report
    # ══════════════════════════════════════════════════════

    def full_score_report(self):
        """Generate complete scoring report: per-page + deck."""
        lines = []
        lines.append(f'🏆 PPT 评分报告 ({self.path})')
        lines.append(f'总页数: {self.n_slides}')
        lines.append('')

        # Per-page scores
        page_scores = []
        for i in range(self.n_slides):
            ps = self.score_page(i)
            page_scores.append(ps)

            slide_label = '封面' if i == 0 else ('结尾' if i == self.n_slides - 1 else f'内容')
            grade = self._grade(ps['total'])
            lines.append(f'📊 S{ps["slide"]} [{slide_label}] {ps["total"]}/100 {grade}')

            for dim, val in ps['scores'].items():
                max_val = 20 if dim in ('结构合理', '留白呼吸', '信息承载', '视觉层次') else 10
                bar = '█' * int(val / max_val * 5) + '░' * (5 - int(val / max_val * 5))
                lines.append(f'  {dim}: {val}/{max_val} {bar}')

            if ps['suggestions']:
                lines.append(f'  💡 建议:')
                for s in ps['suggestions']:
                    lines.append(f'     - {s}')
            lines.append('')

        # Deck score
        deck = self.score_deck(page_scores)
        grade = self._grade(deck['total'])
        lines.append(f'═══════════════════════════════')
        lines.append(f'🎯 整体评分: {deck["total"]}/100 {grade}')
        for dim, val in deck['scores'].items():
            bar = '█' * int(val / 25 * 5) + '░' * (5 - int(val / 25 * 5))
            lines.append(f'  {dim}: {val}/25 {bar}')
        if deck['suggestions']:
            lines.append(f'  💡 整体建议:')
            for s in deck['suggestions']:
                lines.append(f'     - {s}')

        # Summary
        lines.append('')
        avg_page = sum(ps['total'] for ps in page_scores) / len(page_scores) if page_scores else 0
        best_page = max(page_scores, key=lambda x: x['total']) if page_scores else None
        worst_page = min((ps for ps in page_scores if not ps['is_cover']), key=lambda x: x['total'], default=None)

        lines.append(f'📈 页均分: {avg_page:.0f}/100')
        if best_page:
            lines.append(f'🌟 最佳页: S{best_page["slide"]} ({best_page["total"]}/100)')
        if worst_page:
            lines.append(f'🔧 待优化: S{worst_page["slide"]} ({worst_page["total"]}/100)')

        return '\n'.join(lines)

    @staticmethod
    def _grade(score):
        """Convert score to letter grade."""
        if score >= 90:
            return '⭐ A'
        elif score >= 80:
            return '✅ B'
        elif score >= 70:
            return '🔶 C'
        elif score >= 60:
            return '⚠️ D'
        else:
            return '❌ F'


# CLI entry
if __name__ == '__main__':
    import sys
    if len(sys.argv) < 2:
        print("Usage: python ppt_scorer.py <pptx_path>")
        sys.exit(1)
    scorer = PPTScorer(sys.argv[1])
    print(scorer.full_score_report())
