#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""icon_library.py — Icon search, copy, and recolor engine.

Usage:
    from icon_library import IconLibrary
    lib = IconLibrary()
    results = lib.search('服务器', limit=3)
    lib.copy_icon(slide, results[0], x=1.0, y=1.0, size=0.6, color='FFFFFF')
"""
import json, copy, re, os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.oxml.ns import qn
from lxml import etree


class IconLibrary:
    """Search, copy, and recolor icons from the icon library PPTX."""

    def __init__(self):
        base = Path(__file__).parent.parent / 'assets'
        self._pptx_path = base / 'icons-library.pptx'
        self._index_path = base / 'icon-index.json'
        
        with open(self._index_path, 'r', encoding='utf-8') as f:
            self._index = json.load(f)
        
        self._prs = Presentation(str(self._pptx_path))
        self._slide = self._prs.slides[0]
        self._shapes = list(self._slide.shapes)

    def search(self, query, category=None, limit=5):
        """Search icons by keyword (Chinese or English) and optional category.
        Returns list of icon entries sorted by relevance.
        """
        query = query.lower().strip()
        results = []
        
        for entry in self._index:
            score = 0
            name = entry['name'].lower()
            tags_cn = ' '.join(entry.get('tags_cn', []))
            tags_en = ' '.join(entry.get('tags_en', [])).lower()
            cat = entry.get('category', '')
            
            # Category filter
            if category and cat != category:
                continue
            
            # Exact name match
            if query in name:
                score += 10
            
            # Tag match (Chinese)
            if query in tags_cn:
                score += 8
            
            # Tag match (English)
            if query in tags_en:
                score += 8
            
            # Category match
            if query in cat:
                score += 5
            
            # Partial match
            for word in query.split():
                if word in name:
                    score += 3
                if word in tags_cn or word in tags_en:
                    score += 2
            
            if score > 0:
                results.append((score, entry))
        
        results.sort(key=lambda x: -x[0])
        return [r[1] for r in results[:limit]]

    def copy_icon(self, dst_slide, icon_entry, x, y, size=0.5, color=None):
        """Copy an icon to a destination slide, centered in a size x size box at (x, y).
        
        Args:
            dst_slide: target slide object
            icon_entry: dict from search results
            x, y: top-left of the bounding box in inches
            size: bounding box dimension in inches (icon is centered within)
            color: optional hex color string (e.g. 'FFFFFF' for white).
                   If set, all fills in the icon are replaced with this color.
        
        Returns: the new XML element
        """
        idx = icon_entry['shape_idx']
        src_shape = self._shapes[idx]
        
        # Deep copy XML
        new_el = copy.deepcopy(src_shape._element)
        
        # Scale to fit within size x size box, preserving aspect ratio
        orig_w = src_shape.width
        orig_h = src_shape.height
        target = Inches(size)
        scale = target / max(orig_w, orig_h)
        new_w = int(orig_w * scale)
        new_h = int(orig_h * scale)
        
        # Center within the bounding box
        offset_x = Inches(x) + (target - new_w) // 2
        offset_y = Inches(y) + (target - new_h) // 2
        
        # Position and resize
        is_group = src_shape.shape_type == 6
        if is_group:
            pr = new_el.find(qn('p:grpSpPr'))
        else:
            pr = new_el.find(qn('p:spPr'))
        
        if pr is not None:
            xfrm = pr.find(qn('a:xfrm'))
            if xfrm is not None:
                off = xfrm.find(qn('a:off'))
                ext = xfrm.find(qn('a:ext'))
                if off is not None:
                    off.set('x', str(offset_x))
                    off.set('y', str(offset_y))
                if ext is not None:
                    ext.set('cx', str(new_w))
                    ext.set('cy', str(new_h))
        
        # Recolor if requested
        if color:
            self._recolor(new_el, color)
        
        # Add to slide
        sp_tree = dst_slide._element.find(qn('p:cSld')).find(qn('p:spTree'))
        sp_tree.append(new_el)
        return new_el

    def _recolor(self, el, hex_color):
        """Replace all solidFill colors in an element tree."""
        for fill in el.iter(qn('a:solidFill')):
            for child in list(fill):
                fill.remove(child)
            srgb = etree.SubElement(fill, qn('a:srgbClr'))
            srgb.set('val', hex_color)

    def list_categories(self):
        """List all icon categories with counts."""
        from collections import Counter
        cats = Counter(e.get('category', '未分类') for e in self._index)
        return dict(cats.most_common())

    def __len__(self):
        return len(self._index)
