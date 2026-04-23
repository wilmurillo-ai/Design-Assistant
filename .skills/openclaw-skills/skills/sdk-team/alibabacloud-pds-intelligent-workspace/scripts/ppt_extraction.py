#!/usr/bin/env python3
"""
PDS 视频精读 PPT 提取脚本

功能:
- 从视频精读结果中提取 PPT 图片
- 生成 PPTX 文件
- 支持添加备注信息 (页码、时间戳等)
"""

import json
import requests
import argparse
from pathlib import Path
from io import BytesIO

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    print("❌ 缺少依赖库 python-pptx")
    print("请运行：pip install -r scripts/requirements.txt")
    exit(1)


def ms_to_timestamp(ms):
    """毫秒转时间戳格式"""
    seconds = ms // 1000
    hours = seconds // 3600
    minutes = (seconds % 3600) // 60
    secs = seconds % 60
    return f"{hours:02d}:{minutes:02d}:{secs:02d}"


def download_image(url):
    """下载图片到内存"""
    response = requests.get(url, timeout=30)
    response.raise_for_status()
    return BytesIO(response.content)


def create_pptx_from_video_analysis(result_json, output_path="output.pptx", keep_aspect_ratio=False):
    """
    从视频分析结果创建 PPTX 文件
    
    参数:
        result_json: 精读 API 返回的完整结果 (dict 或 JSON 文件路径)
        output_path: 输出 PPTX 文件路径
        keep_aspect_ratio: 是否保持图片宽高比 (默认 False，填充整个幻灯片)
    
    返回:
        bool: 成功返回 True，失败返回 False
    """
    # 1. 加载结果数据
    if isinstance(result_json, str):
        with open(result_json, 'r', encoding='utf-8') as f:
            result = json.load(f)
    else:
        result = result_json

    # 2. 检查 ppt_details 是否存在
    if "ppt_details" not in result or not result["ppt_details"]:
        print("❌ 该视频中未检测到 PPT 内容")
        return False

    # 3. 下载 ppt_details JSON 文件
    ppt_details_url = result["ppt_details"][0]
    print(f"📥 下载 PPT 详情：{ppt_details_url}")
    ppt_data = requests.get(ppt_details_url, timeout=30).json()

    # 4. 按 PPTShotIndex 排序
    ppt_data.sort(key=lambda x: x["PPTShotIndex"])
    print(f"📊 检测到 {len(ppt_data)} 页 PPT")

    # 5. 获取 images 映射
    images_map = result.get("images", {})

    # 6. 创建 PPTX 文件
    prs = Presentation()

    # 设置幻灯片尺寸 (16:9 宽屏)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # 7. 逐页添加 PPT 图片
    for i, ppt_page in enumerate(ppt_data, 1):
        image_path = ppt_page["ImagePath"]
        ppt_index = ppt_page["PPTShotIndex"]
        start_time = ms_to_timestamp(ppt_page["StartTime"])

        print(f"  - 处理第 {i}/{len(ppt_data)} 页 (索引：{ppt_index}, 时间：{start_time})")

        # 获取图片 URL（兼容大小写）
        if image_path not in images_map:
            print(f"    ⚠️  警告：图片路径 {image_path} 未在 images 中找到，跳过")
            continue
        
        image_info = images_map[image_path]
        image_url = image_info.get("Url") or image_info.get("url")

        # 下载图片
        try:
            image_stream = download_image(image_url)
        except Exception as e:
            print(f"    ❌ 下载图片失败：{e}")
            continue

        # 添加空白幻灯片
        blank_slide_layout = prs.slide_layouts[6]  # 6 表示空白布局
        slide = prs.slides.add_slide(blank_slide_layout)

        # 插入图片
        if keep_aspect_ratio:
            # 保持宽高比插入
            add_picture_with_aspect_ratio(
                slide, 
                image_stream, 
                prs.slide_width, 
                prs.slide_height
            )
        else:
            # 填充整个幻灯片
            left = Inches(0)
            top = Inches(0)
            width = prs.slide_width
            height = prs.slide_height
            
            slide.shapes.add_picture(
                image_stream,
                left, top,
                width=width,
                height=height
            )

        # 添加备注信息
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = (
            f"页码：{i}\n"
            f"索引：{ppt_index}\n"
            f"出现时间：{start_time}\n"
            f"图片路径：{image_path}"
        )

    # 8. 保存 PPTX 文件
    prs.save(output_path)
    print(f"✅ PPTX 文件已生成：{output_path}")
    print(f"   总页数：{len(prs.slides)}")

    return True


def add_picture_with_aspect_ratio(slide, image_stream, slide_width, slide_height):
    """
    插入图片并保持宽高比
    
    参数:
        slide: PPTX 幻灯片对象
        image_stream: 图片流 (BytesIO)
        slide_width: 幻灯片宽度
        slide_height: 幻灯片高度
    """
    from PIL import Image
    
    # 获取图片尺寸
    img = Image.open(image_stream)
    img_width, img_height = img.size
    img_aspect = img_width / img_height

    slide_aspect = slide_width / slide_height

    if img_aspect > slide_aspect:
        # 图片更宽，以宽度为基准
        width = slide_width
        height = slide_width / img_aspect
        left = Inches(0)
        top = (slide_height - height) / 2
    else:
        # 图片更高，以高度为基准
        height = slide_height
        width = slide_height * img_aspect
        top = Inches(0)
        left = (slide_width - width) / 2

    # 重置流位置
    image_stream.seek(0)

    slide.shapes.add_picture(image_stream, left, top, width=width, height=height)


def validate_pptx(pptx_path, expected_slide_count):
    """
    验证生成的 PPTX 文件
    
    参数:
        pptx_path: PPTX 文件路径
        expected_slide_count: 期望的幻灯片数量
    
    返回:
        bool: 验证是否通过
    """
    try:
        prs = Presentation(pptx_path)
        actual_count = len(prs.slides)

        print(f"\n📊 PPTX 验证结果:")
        print(f"   文件路径：{pptx_path}")
        print(f"   期望页数：{expected_slide_count}")
        print(f"   实际页数：{actual_count}")

        if actual_count == expected_slide_count:
            print("   ✅ 页数匹配")
        else:
            print("   ⚠️  页数不匹配")

        # 检查每页是否包含图片
        missing_images = []
        for i, slide in enumerate(prs.slides, 1):
            has_picture = any(
                shape.shape_type == 13  # 13 表示图片
                for shape in slide.shapes
            )
            if not has_picture:
                missing_images.append(i)
        
        if missing_images:
            print(f"   ⚠️  以下页面未包含图片：{missing_images}")
        else:
            print("   ✅ 所有页面都包含图片")

        return actual_count == expected_slide_count and not missing_images

    except Exception as e:
        print(f"❌ 验证失败：{e}")
        return False


def main():
    """主函数"""
    parser = argparse.ArgumentParser(
        description='PDS 视频精读 PPT 提取工具',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例用法:
  python ppt_extraction.py video_analysis_result.json
  python ppt_extraction.py video_analysis_result.json -o extracted.pptx
  python ppt_extraction.py video_analysis_result.json --keep-aspect-ratio
  python ppt_extraction.py video_analysis_result.json --validate
        """
    )
    
    parser.add_argument(
        'input_file',
        help='视频精读 API 返回的 JSON 结果文件路径'
    )
    
    parser.add_argument(
        '-o', '--output',
        default='extracted_ppt.pptx',
        help='输出 PPTX 文件路径 (默认：extracted_ppt.pptx)'
    )
    
    parser.add_argument(
        '--keep-aspect-ratio',
        action='store_true',
        help='保持图片宽高比 (默认填充整个幻灯片)'
    )
    
    parser.add_argument(
        '--validate',
        action='store_true',
        help='生成后验证 PPTX 文件'
    )
    
    args = parser.parse_args()
    
    # 检查输入文件是否存在
    input_path = Path(args.input_file)
    if not input_path.exists():
        print(f"❌ 文件不存在：{args.input_file}")
        return 1
    
    try:
        success = create_pptx_from_video_analysis(
            args.input_file,
            args.output,
            args.keep_aspect_ratio
        )
        
        if success and args.validate:
            # 加载 ppt_details 获取期望页数
            with open(args.input_file, 'r', encoding='utf-8') as f:
                result = json.load(f)
            
            if "ppt_details" in result and result["ppt_details"]:
                ppt_details_url = result["ppt_details"][0]
                ppt_data = requests.get(ppt_details_url, timeout=30).json()
                expected_count = len(ppt_data)
                validate_pptx(args.output, expected_count)
        
        return 0 if success else 1
        
    except Exception as e:
        print(f"❌ 处理失败：{e}")
        import traceback
        traceback.print_exc()
        return 1


if __name__ == "__main__":
    exit(main())
