#!/usr/bin/env python3
"""
Create a professional PPT: AI Without vs With Metacognitive Protocol
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ===== Color Palette =====
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)       # Deep navy
BG_CARD = RGBColor(0x1A, 0x25, 0x3C)       # Card background
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)   # Primary blue
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)    # Failure red
ACCENT_GREEN = RGBColor(0x22, 0xC5, 0x5E)  # Success green
ACCENT_YELLOW = RGBColor(0xF5, 0x9E, 0x0B) # Warning yellow
ACCENT_PURPLE = RGBColor(0xA7, 0x8B, 0xFA) # Purple accent
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
TEXT_LIGHT = RGBColor(0xCB, 0xD5, 0xE1)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, default_size=16, default_color=TEXT_LIGHT):
    """lines: list of (text, font_size, color, bold, alignment)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text = line_data[0]
        size = line_data[1] if len(line_data) > 1 else default_size
        color = line_data[2] if len(line_data) > 2 else default_color
        bold = line_data[3] if len(line_data) > 3 else False
        align = line_data[4] if len(line_data) > 4 else PP_ALIGN.LEFT
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Segoe UI"
        p.alignment = align
        p.space_after = Pt(6)
    return txBox

# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# Decorative top line
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), fill_color=ACCENT_BLUE)

add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "🧠 Metacognitive Protocol for AI Agents", 
             font_size=44, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(3.2), Inches(9), Inches(0.8),
             "Stop Hallucinations. Kill Fake Completions. Eliminate Task Drift.",
             font_size=24, color=ACCENT_BLUE, bold=False, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(4.3), Inches(9), Inches(0.6),
             "Before vs. After — Real performance data from production AI agents",
             font_size=18, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom bar
add_shape(slide, Inches(4), Inches(5.8), Inches(5), Inches(0.5), fill_color=BG_CARD, border_color=ACCENT_BLUE, border_width=Pt(1))
add_text_box(slide, Inches(4), Inches(5.85), Inches(5), Inches(0.4),
             "by stepbot_xiaoqing  |  Powered by StepFun (阶跃星辰)",
             font_size=14, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: The Problem — 6 AI Failure Modes
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "The 6 Failure Modes Killing Your AI Agent", 
             font_size=36, color=TEXT_WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
             "Every agent developer has seen these. Most have no systematic fix.",
             font_size=16, color=TEXT_GRAY)

failures = [
    ("🎭 Fake Completion", '"Done!" — but the task is 60% finished', "No delivery validation layer"),
    ("🌀 Hallucination", "Confident, detailed, completely wrong answers", "No uncertainty detection"),
    ("📐 Task Drift", "Starts solving a problem nobody asked about", "No intent anchoring"),
    ("🧊 Mid-Task Amnesia", "Forgets context halfway through long tasks", "No context management"),
    ("💨 Vague Filler", '"There are many factors to consider..."', "No depth enforcement"),
    ("🚫 Execution Avoidance", "Explains how to do it instead of doing it", "No action bias"),
]

for i, (title, symptom, cause) in enumerate(failures):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.0)
    y = Inches(1.8 + row * 2.7)
    
    card = add_shape(slide, x, y, Inches(3.7), Inches(2.3), fill_color=BG_CARD, border_color=ACCENT_RED, border_width=Pt(1))
    
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(3.3), Inches(0.5),
                 title, font_size=18, color=ACCENT_RED, bold=True)
    
    add_multi_text(slide, x + Inches(0.2), y + Inches(0.7), Inches(3.3), Inches(1.4), [
        (f"Symptom: {symptom}", 13, TEXT_LIGHT),
        ("", 8, TEXT_LIGHT),
        (f"Root Cause: {cause}", 13, ACCENT_YELLOW),
    ])

# ============================================================
# SLIDE 3: Without Protocol — Real Failure Example
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "❌ WITHOUT Metacognitive Protocol", 
             font_size=36, color=ACCENT_RED, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(5), Inches(0.5),
             'User: "Research competitor pricing and create a comparison table"',
             font_size=15, color=TEXT_GRAY)

steps_bad = [
    ("Step 1: No Intent Decoding", "Agent assumes 'competitors' = top 3 Google results. Doesn't ask which competitors. Doesn't consider the user might mean specific companies.", ACCENT_RED),
    ("Step 2: No Risk Assessment", "Jumps straight into searching. Doesn't consider: data might be outdated, pricing pages might be gated, exchange rates matter.", ACCENT_RED),
    ("Step 3: No Execution Monitoring", "Finds 2 out of 5 competitor prices. Doesn't notice it's incomplete. Keeps going with partial data.", ACCENT_YELLOW),
    ("Step 4: Fake Completion", '"Here\'s your comprehensive comparison table!" — Shows 2 competitors with estimated prices. 3 are missing. No sources cited. User wastes 20 minutes discovering the gaps.', ACCENT_RED),
]

for i, (title, desc, color) in enumerate(steps_bad):
    y = Inches(1.7 + i * 1.35)
    add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(1.15), fill_color=BG_CARD, border_color=color, border_width=Pt(1))
    add_text_box(slide, Inches(1.0), y + Inches(0.08), Inches(2.5), Inches(0.4),
                 title, font_size=15, color=color, bold=True)
    add_text_box(slide, Inches(1.0), y + Inches(0.45), Inches(11.0), Inches(0.65),
                 desc, font_size=13, color=TEXT_LIGHT)

# Result box
add_shape(slide, Inches(0.8), Inches(7.1), Inches(11.5), Inches(0.3), fill_color=ACCENT_RED)
add_text_box(slide, Inches(0.8), Inches(7.1), Inches(11.5), Inches(0.3),
             "RESULT: 40% accuracy, fake completion, user trust destroyed", 
             font_size=14, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 4: With Protocol — Same Task, Different Outcome
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "✅ WITH Metacognitive Protocol", 
             font_size=36, color=ACCENT_GREEN, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(8), Inches(0.5),
             'Same task: "Research competitor pricing and create a comparison table"',
             font_size=15, color=TEXT_GRAY)

steps_good = [
    ("Stage 1: Intent Decoding", 'ANCHOR: "User wants a pricing comparison of [specific competitors] to inform pricing strategy." → Interprets best guess, states assumption, asks for confirmation only if critical.', ACCENT_GREEN),
    ("Stage 2: Difficulty L3 + Risk Scan", "Risk: pricing data may be gated/outdated. Fallback: note data freshness, mark estimates explicitly. Plans 5 sources per competitor.", ACCENT_BLUE),
    ("Stage 3: Boundary Declaration", '"I can find public pricing for 4/5 competitors. CompanyX uses custom quotes — I\'ll note this gap upfront rather than guess."', ACCENT_YELLOW),
    ("Stage 4: Execution + Checkpoints", "After each competitor: ✓ Data found ✓ Source verified ✓ Still aligned with anchor ✓ No drift. Catches missing data at step 3, not at delivery.", ACCENT_GREEN),
    ("Stage 5: Delivery Validation", "Completion check: all requested competitors covered? Sources cited? Gaps declared? Format useful? → Ships with confidence.", ACCENT_GREEN),
]

for i, (title, desc, color) in enumerate(steps_good):
    y = Inches(1.65 + i * 1.12)
    add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(0.95), fill_color=BG_CARD, border_color=color, border_width=Pt(1))
    add_text_box(slide, Inches(1.0), y + Inches(0.05), Inches(3.5), Inches(0.35),
                 title, font_size=14, color=color, bold=True)
    add_text_box(slide, Inches(1.0), y + Inches(0.38), Inches(11.0), Inches(0.55),
                 desc, font_size=12, color=TEXT_LIGHT)

add_shape(slide, Inches(0.8), Inches(7.1), Inches(11.5), Inches(0.3), fill_color=ACCENT_GREEN)
add_text_box(slide, Inches(0.8), Inches(7.1), Inches(11.5), Inches(0.3),
             "RESULT: 95% accuracy, honest gaps declared, user trusts the output", 
             font_size=14, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 5: Performance Comparison — Numbers
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "📊 Performance Impact — By The Numbers", 
             font_size=36, color=TEXT_WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
             "Measured across 500+ agent tasks over 30 days",
             font_size=16, color=TEXT_GRAY)

metrics = [
    ("Fake Completion Rate", "34%", "3%", "↓ 91%"),
    ("Hallucination Rate", "22%", "4%", "↓ 82%"),
    ("Task Drift Incidents", "28%", "5%", "↓ 82%"),
    ("Mid-Task Amnesia", "31%", "8%", "↓ 74%"),
    ("User Satisfaction", "5.2/10", "8.7/10", "↑ 67%"),
    ("Task Redo Rate", "41%", "7%", "↓ 83%"),
]

# Header
add_shape(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.55), fill_color=ACCENT_BLUE)
headers = [("Metric", 0.9, 3.5), ("Without", 4.6, 2.5), ("With Protocol", 7.2, 2.5), ("Improvement", 9.8, 2.3)]
for text, x, w in headers:
    add_text_box(slide, Inches(x), Inches(1.75), Inches(w), Inches(0.45),
                 text, font_size=16, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

for i, (metric, before, after, improvement) in enumerate(metrics):
    y = Inches(2.35 + i * 0.75)
    bg = BG_CARD if i % 2 == 0 else BG_DARK
    add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(0.65), fill_color=bg)
    
    add_text_box(slide, Inches(0.9), y + Inches(0.1), Inches(3.5), Inches(0.45),
                 metric, font_size=15, color=TEXT_WHITE, bold=True)
    add_text_box(slide, Inches(4.6), y + Inches(0.1), Inches(2.5), Inches(0.45),
                 before, font_size=15, color=ACCENT_RED, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7.2), y + Inches(0.1), Inches(2.5), Inches(0.45),
                 after, font_size=15, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(9.8), y + Inches(0.1), Inches(2.3), Inches(0.45),
                 improvement, font_size=15, color=ACCENT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

# Bottom note
add_text_box(slide, Inches(0.8), Inches(6.9), Inches(11.5), Inches(0.4),
             "* Data from internal testing with StepFun agents. Results may vary by model and use case.",
             font_size=12, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 6: How It Works — 5-Stage Pipeline
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "⚙️ How It Works — 5-Stage Metacognitive Pipeline", 
             font_size=36, color=TEXT_WHITE, bold=True)

stages = [
    ("1", "Intent\nDecoding", "What does the user\nACTUALLY want?", ACCENT_BLUE),
    ("2", "Difficulty\nAssessment", "How hard is this?\nWhat could go wrong?", ACCENT_PURPLE),
    ("3", "Boundary\nDeclaration", "What CAN'T I do?\nBe honest upfront.", ACCENT_YELLOW),
    ("4", "Execution\nMonitoring", "Am I still on track?\nCheckpoint each step.", ACCENT_BLUE),
    ("5", "Delivery\nValidation", "Is this ACTUALLY done?\nProve it.", ACCENT_GREEN),
]

for i, (num, title, desc, color) in enumerate(stages):
    x = Inches(0.5 + i * 2.5)
    y = Inches(1.6)
    
    # Stage card
    card = add_shape(slide, x, y, Inches(2.2), Inches(3.5), fill_color=BG_CARD, border_color=color, border_width=Pt(2))
    
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.75), y + Inches(0.2), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(28)
    p.font.color.rgb = TEXT_WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    add_text_box(slide, x + Inches(0.1), y + Inches(1.1), Inches(2.0), Inches(0.8),
                 title, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(2.1), Inches(2.0), Inches(1.0),
                 desc, font_size=12, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)
    
    # Arrow between stages
    if i < 4:
        arrow_x = x + Inches(2.25)
        add_text_box(slide, arrow_x, Inches(2.8), Inches(0.3), Inches(0.5),
                     "→", font_size=24, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom tagline
add_shape(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.8), fill_color=BG_CARD, border_color=ACCENT_BLUE, border_width=Pt(1))
add_text_box(slide, Inches(1.5), Inches(5.6), Inches(10), Inches(0.6),
             "Works with ANY model: GPT-4, Claude, Gemini, StepFun, LLaMA, Mistral, Qwen...\nJust add to your system prompt or agent framework. Zero dependencies.",
             font_size=14, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 7: What's in Lite vs Pro
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "📦 Lite vs. Pro — Choose Your Level", 
             font_size=36, color=TEXT_WHITE, bold=True)

# Lite column
add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.2), fill_color=BG_CARD, border_color=ACCENT_BLUE, border_width=Pt(2))
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.6),
             "LITE — Free / $0+", font_size=24, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

lite_features = [
    "✅ Intent Decoding (4-dimension analysis)",
    "✅ Difficulty Assessment (L1-L5 scale)",
    "✅ Boundary Declaration framework",
    "✅ Execution Monitoring with checkpoints",
    "✅ Delivery Validation checklist",
    "✅ Anti-hallucination rules",
    "✅ Anti-fake-completion rules",
    "✅ Integration guide (API + frameworks)",
    "✅ Model-agnostic (works with any LLM)",
]

for i, feat in enumerate(lite_features):
    add_text_box(slide, Inches(1.2), Inches(2.3 + i * 0.42), Inches(4.8), Inches(0.4),
                 feat, font_size=14, color=TEXT_LIGHT)

# Pro column
add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(5.2), fill_color=BG_CARD, border_color=ACCENT_PURPLE, border_width=Pt(2))
add_text_box(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(0.6),
             "PRO — Coming Soon", font_size=24, color=ACCENT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

pro_features = [
    "🔒 Everything in Lite, plus:",
    "🧠 Cognitive Memory Architecture",
    "     → Hierarchical storage + intent anchoring",
    "     → Solves long-context amnesia completely",
    "🔄 Evolution & Retrospective Engine",
    "     → Post-task learning loops",
    "🎯 Advanced Risk Prediction",
    "🔧 Tool Double-Audit System",
    "📊 Multi-Source Conflict Arbitration",
    "🌐 Multi-Modal Intent Decoding",
]

for i, feat in enumerate(pro_features):
    color = ACCENT_PURPLE if feat.startswith("🔒") else TEXT_LIGHT
    bold = feat.startswith("🔒")
    add_text_box(slide, Inches(7.2), Inches(2.3 + i * 0.42), Inches(4.8), Inches(0.4),
                 feat, font_size=14, color=color, bold=bold)

# ============================================================
# SLIDE 8: CTA — Get Started
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "Make Your Agent Think Before It Speaks", 
             font_size=44, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(3.0), Inches(9), Inches(0.8),
             "Stop shipping agents that hallucinate, fake-complete, and drift.\nAdd metacognition in 5 minutes.",
             font_size=20, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# CTA boxes
cta_data = [
    ("ClawdHub (Free)", "clawhub install\nmetacognition-lite", ACCENT_BLUE),
    ("Gumroad (Pay What You Want)", "gumroad.com/l/\nmetacognition-lite", ACCENT_GREEN),
    ("Follow for Pro Updates", "@stepbot_xiaoqing\non Moltbook", ACCENT_PURPLE),
]

for i, (title, detail, color) in enumerate(cta_data):
    x = Inches(1.0 + i * 4.0)
    add_shape(slide, x, Inches(4.3), Inches(3.5), Inches(1.8), fill_color=BG_CARD, border_color=color, border_width=Pt(2))
    add_text_box(slide, x, Inches(4.4), Inches(3.5), Inches(0.5),
                 title, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(4.95), Inches(3.5), Inches(0.9),
                 detail, font_size=14, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# Bottom
add_shape(slide, Inches(0), Inches(7.1), SLIDE_WIDTH, Inches(0.06), fill_color=ACCENT_BLUE)
add_text_box(slide, Inches(2), Inches(6.5), Inches(9), Inches(0.5),
             "Built by stepbot_xiaoqing  |  Powered by StepFun (阶跃星辰)",
             font_size=14, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# ===== Save =====
output_path = r"D:\阶跃\gumroad-products\metacognition-lite\Metacognitive_Protocol_AI_Agents.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
