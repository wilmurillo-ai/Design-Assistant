"""
PowerPoint PPT生成器实现
基于python-pptx库
"""

from .base_generator import BasePPTGenerator, PPTStructure, SlideContent, PPTSoftware, PPTTemplate
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from typing import Any, Optional
import os
from datetime import datetime


class PPTXGenerator(BasePPTGenerator):
    """PowerPoint PPT生成器"""
    
    def __init__(self, template: PPTTemplate = PPTTemplate.BUSINESS):
        super().__init__(software_type=PPTSoftware.POWERPOINT, template=template)
        self.prs = Presentation()
        self._setup_styles()
    
    def _setup_styles(self):
        """设置样式"""
        self.styles = {
            "colors": self._get_template_colors(),
            "fonts": self._get_font_sizes(),
            "layouts": {
                "title": 0,      # 标题幻灯片
                "content": 1,    # 标题和内容
                "two_content": 3, # 两栏内容
                "comparison": 4,  # 比较
                "blank": 6        # 空白
            }
        }
    
    def _get_template_colors(self) -> dict:
        """根据模板获取颜色方案"""
        color_schemes = {
            PPTTemplate.BUSINESS: {
                "primary": RGBColor(0, 122, 255),    # 商务蓝
                "secondary": RGBColor(142, 142, 147), # 科技灰
                "accent": RGBColor(255, 149, 0),     # 活力橙
                "background": RGBColor(242, 242, 247), # 浅灰白
                "text_dark": RGBColor(28, 28, 30)    # 深灰黑
            },
            PPTTemplate.EDUCATION: {
                "primary": RGBColor(52, 199, 89),    # 教育绿
                "secondary": RGBColor(175, 82, 222), # 紫色
                "accent": RGBColor(255, 59, 48),     # 红色
                "background": RGBColor(255, 255, 255), # 白色
                "text_dark": RGBColor(0, 0, 0)       # 黑色
            },
            PPTTemplate.TECHNICAL: {
                "primary": RGBColor(88, 86, 214),    # 技术紫
                "secondary": RGBColor(255, 45, 85),  # 粉色
                "accent": RGBColor(255, 204, 0),     # 黄色
                "background": RGBColor(29, 29, 31),  # 深色背景
                "text_dark": RGBColor(255, 255, 255) # 白色文字
            },
            PPTTemplate.CREATIVE: {
                "primary": RGBColor(255, 69, 58),    # 创意红
                "secondary": RGBColor(48, 176, 199), # 青色
                "accent": RGBColor(255, 214, 10),    # 亮黄
                "background": RGBColor(255, 255, 255), # 白色
                "text_dark": RGBColor(0, 0, 0)       # 黑色
            },
            PPTTemplate.MINIMALIST: {
                "primary": RGBColor(0, 0, 0),        # 黑色
                "secondary": RGBColor(150, 150, 150), # 灰色
                "accent": RGBColor(0, 0, 0),         # 黑色
                "background": RGBColor(255, 255, 255), # 白色
                "text_dark": RGBColor(0, 0, 0)       # 黑色
            }
        }
        
        return color_schemes.get(self.template, color_schemes[PPTTemplate.BUSINESS])
    
    def _get_font_sizes(self) -> dict:
        """获取字体大小配置"""
        return {
            "title_size": Pt(44),
            "subtitle_size": Pt(28),
            "heading_size": Pt(32),
            "body_size": Pt(18),
            "caption_size": Pt(14)
        }
    
    def create_cover(self, structure: PPTStructure) -> Any:
        """
        创建封面页
        
        Args:
            structure: PPT结构
            
        Returns:
            PowerPoint幻灯片对象
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.styles["layouts"]["title"]])
        
        # 设置标题
        title_shape = slide.shapes.title
        title_shape.text = structure.title
        title_shape.text_frame.paragraphs[0].font.size = self.styles["fonts"]["title_size"]
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.styles["colors"]["primary"]
        
        # 设置副标题
        if structure.subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = structure.subtitle
            subtitle_shape.text_frame.paragraphs[0].font.size = self.styles["fonts"]["subtitle_size"]
            subtitle_shape.text_frame.paragraphs[0].font.color.rgb = self.styles["colors"]["secondary"]
        
        # 添加作者和日期
        footer_text = ""
        if structure.author:
            footer_text += f"作者: {structure.author}"
        if structure.date:
            if footer_text:
                footer_text += " | "
            footer_text += structure.date
        
        if footer_text:
            left = Inches(1)
            top = Inches(6.5)
            width = Inches(8)
            height = Inches(0.5)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            p = text_frame.add_paragraph()
            p.text = footer_text
            p.font.size = self.styles["fonts"]["caption_size"]
            p.font.color.rgb = self.styles["colors"]["secondary"]
            p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def create_content_slide(self, content: SlideContent) -> Any:
        """
        创建内容页
        
        Args:
            content: 幻灯片内容
            
        Returns:
            PowerPoint幻灯片对象
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.styles["layouts"]["content"]])
        
        # 设置标题
        title_shape = slide.shapes.title
        title_shape.text = content.title
        title_shape.text_frame.paragraphs[0].font.size = self.styles["fonts"]["heading_size"]
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.styles["colors"]["primary"]
        
        # 设置内容
        if content.content:
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()  # 清除默认文本
            
            for i, item in enumerate(content.content):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = self.styles["fonts"]["body_size"]
                p.font.color.rgb = self.styles["colors"]["text_dark"]
                
                if content.bullet_points and i == 0:
                    p.level = 0
                elif content.bullet_points:
                    p.level = 1
        
        return slide
    
    def create_two_column_slide(self, content: SlideContent) -> Any:
        """
        创建两栏内容页
        
        Args:
            content: 幻灯片内容
            
        Returns:
            PowerPoint幻灯片对象
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.styles["layouts"]["two_content"]])
        
        # 设置标题
        title_shape = slide.shapes.title
        title_shape.text = content.title
        title_shape.text_frame.paragraphs[0].font.size = self.styles["fonts"]["heading_size"]
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # 设置左侧内容
        if content.left_content:
            left_shape = slide.placeholders[1]
            text_frame = left_shape.text_frame
            text_frame.clear()
            
            for item in content.left_content:
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = self.styles["fonts"]["body_size"]
        
        # 设置右侧内容
        if content.right_content:
            right_shape = slide.placeholders[2]
            text_frame = right_shape.text_frame
            text_frame.clear()
            
            for item in content.right_content:
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = self.styles["fonts"]["body_size"]
        
        return slide
    
    def create_comparison_slide(self, content: SlideContent) -> Any:
        """
        创建比较页
        
        Args:
            content: 幻灯片内容
            
        Returns:
            PowerPoint幻灯片对象
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.styles["layouts"]["comparison"]])
        
        # 设置主标题
        title_shape = slide.shapes.title
        title_shape.text = content.title
        title_shape.text_frame.paragraphs[0].font.size = self.styles["fonts"]["heading_size"]
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # 设置左侧标题和内容
        if content.left_content:
            left_title_shape = slide.placeholders[1]
            if content.left_content:
                left_title_shape.text = content.title  # 使用主标题作为左侧标题
            
            left_content_shape = slide.placeholders[2]
            text_frame = left_content_shape.text_frame
            text_frame.clear()
            
            for item in content.left_content:
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = self.styles["fonts"]["body_size"]
        
        # 设置右侧标题和内容
        if content.right_content:
            right_title_shape = slide.placeholders[3]
            if content.right_content:
                right_title_shape.text = content.title  # 使用主标题作为右侧标题
            
            right_content_shape = slide.placeholders[4]
            text_frame = right_content_shape.text_frame
            text_frame.clear()
            
            for item in content.right_content:
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = self.styles["fonts"]["body_size"]
        
        return slide
    
    def create_summary_slide(self, content: SlideContent) -> Any:
        """
        创建总结页
        
        Args:
            content: 幻灯片内容
            
        Returns:
            PowerPoint幻灯片对象
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.styles["layouts"]["content"]])
        
        # 设置标题
        title_shape = slide.shapes.title
        title_shape.text = content.title
        title_shape.text_frame.paragraphs[0].font.size = self.styles["fonts"]["heading_size"]
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # 添加关键点
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        if content.key_points:
            for point in content.key_points:
                p = text_frame.add_paragraph()
                p.text = f"• {point}"
                p.font.size = self.styles["fonts"]["body_size"]
                p.font.bold = True
                p.font.color.rgb = self.styles["colors"]["primary"]
                p.space_after = Pt(10)
        
        # 添加结论
        if content.conclusion:
            p = text_frame.add_paragraph()
            p.text = content.conclusion
            p.font.size = self.styles["fonts"]["body_size"]
            p.font.italic = True
            p.space_before = Pt(20)
        
        return slide
    
    def save_ppt(self, file_path: str) -> str:
        """
        保存PPT文件
        
        Args:
            file_path: 保存路径
            
        Returns:
            保存的文件路径
        """
        # 确保目录存在
        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        
        # 保存文件
        self.prs.save(file_path)
        
        return file_path