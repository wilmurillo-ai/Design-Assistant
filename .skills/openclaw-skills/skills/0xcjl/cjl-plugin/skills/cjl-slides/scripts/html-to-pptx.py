#!/usr/bin/env python3
"""
html-to-pptx.py — 将 HTML 幻灯片转换为 PPTX 格式

保真原则：颜色（十六进制原样）、字体（精确字族）、字号（精确值）、布局（坐标还原）
用法：python html-to-pptx.py <input.html> [output.pptx]
"""

import sys
import re
import html
from pathlib import Path
from html.parser import HTMLParser
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# 默认 PPTX 尺寸：16:9 (914.4 x 514.35 cm -> 9144000 x 5143500 EMU)
SLIDE_WIDTH = Emu(9144000)
SLIDE_HEIGHT = Emu(5143500)


class StyleResolver:
    """从 HTML 解析计算后的样式"""

    def __init__(self):
        self.styles = {}
        self.compute_order = []

    def parse_color(self, color_str):
        """解析颜色字符串为 RGBColor"""
        if not color_str or color_str in ('inherit', 'transparent', 'none', ''):
            return None
        color_str = color_str.strip()
        if color_str.startswith('#'):
            if len(color_str) == 7:
                r = int(color_str[1:3], 16)
                g = int(color_str[3:5], 16)
                b = int(color_str[5:7], 16)
                return RGBColor(r, g, b)
            elif len(color_str) == 4:
                r = int(color_str[1] * 2, 16)
                g = int(color_str[2] * 2, 16)
                b = int(color_str[3] * 2, 16)
                return RGBColor(r, g, b)
        elif color_str.startswith('rgb'):
            match = re.findall(r'\d+', color_str)
            if len(match) >= 3:
                return RGBColor(int(match[0]), int(match[1]), int(match[2]))
        return None

    def parse_size(self, size_str, fallback=Pt(18)):
        """解析尺寸字符串为 Pt"""
        if not size_str or size_str in ('inherit', 'auto', ''):
            return fallback
        size_str = size_str.strip()
        match = re.match(r'([\d.]+)(px|em|rem|vw|vh|pt|cm)?', size_str)
        if not match:
            return fallback
        value = float(match.group(1))
        unit = match.group(2) or 'px'
        if unit == 'pt':
            return Pt(value)
        elif unit == 'px':
            return Pt(value * 0.75)
        elif unit == 'em' or unit == 'rem':
            return Pt(value * 16 * 0.75)
        elif unit == 'vw':
            return Pt(value * 9.144)
        elif unit == 'vh':
            return Pt(value * 5.1435)
        return fallback


class HtmlToPptx:
    """HTML 幻灯片 → PPTX 转换器"""

    def __init__(self, html_content: str):
        self.html = html_content
        self.resolver = StyleResolver()
        self.presentation = Presentation()
        self.presentation.slide_width = SLIDE_WIDTH
        self.presentation.slide_height = SLIDE_HEIGHT
        self._parse_html_tree()

    def _parse_html_tree(self):
        """用 lxml 解析 HTML"""
        try:
            root = etree.HTML(self.html)
        except Exception:
            root = etree.fromstring(f'<html><body>{self.html}</body></html>')
        self.root = root

    def _get_element_style(self, elem) -> dict:
        """从 style 属性解析内联样式"""
        style_attr = elem.get('style', '')
        style = {}
        for item in style_attr.split(';'):
            if ':' in item:
                key, value = item.split(':', 1)
                style[key.strip().lower()] = value.strip()
        return style

    def _elem_text(self, elem) -> str:
        """提取元素的所有文本内容（递归）"""
        if elem is None:
            return ''
        if isinstance(elem, str):
            return elem
        text = ''
        if elem.text:
            text += elem.text
        for child in elem:
            text += self._elem_text(child)
            if child.tail:
                text += child.tail
        return text

    def _tag_name(self, elem) -> str:
        if elem is None:
            return ''
        tag = elem.tag
        if isinstance(tag, str):
            return tag.lower()
        return ''

    def convert(self, output_path: str):
        """执行转换"""
        slides = self.root.xpath('//div[contains(concat(" ", normalize-space(@class), " "), " slide ")]')
        if not slides:
            slides = self.root.xpath('//section')
        if not slides:
            slides = self.root.xpath('//div[@id]')

        if not slides:
            blank_slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])
            body = blank_slide.shapes.placeholders[0]
            tf = body.text_frame
            tf.text = self._elem_text(self.root)
        else:
            for slide_elem in slides:
                self._convert_slide(slide_elem)

        self.presentation.save(output_path)
        return output_path

    def _convert_slide(self, slide_elem):
        """转换单张幻灯片"""
        blank_slide_layout = self.presentation.slide_layouts[6]
        slide = self.presentation.slides.add_slide(blank_slide_layout)

        bg_shape = slide.shapes.add_shape(
            1, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
        )
        bg_shape.line.fill.background()

        bg_style = self._get_element_style(slide_elem)
        if 'background-color' in bg_style:
            bg_color = self.resolver.parse_color(bg_style['background-color'])
            if bg_color:
                bg_shape.fill.solid()
                bg_shape.fill.fore_color.rgb = bg_color
        elif 'background' in bg_style:
            bg_color = self.resolver.parse_color(bg_style['background'])
            if bg_color:
                bg_shape.fill.solid()
                bg_shape.fill.fore_color.rgb = bg_color

        self._add_children_to_slide(slide, slide_elem, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)

    def _add_children_to_slide(self, slide, parent, x, y, width, height):
        """递归将子元素添加到幻灯片"""
        children = list(parent)

        for child in children:
            tag = self._tag_name(child)

            if tag in ('style', 'script', 'meta', 'link'):
                continue

            style = self._get_element_style(child)

            left = x + self._parse_emu(style.get('left', '0'), width)
            top = y + self._parse_emu(style.get('top', '0'), height)
            child_width = self._parse_emu(style.get('width', f'{width}'), width)
            child_height = self._parse_emu(style.get('height', '0'), height)

            if tag == 'div' or tag == 'section':
                self._add_children_to_slide(slide, child, left, top, child_width, child_height)
            elif tag in ('h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'span'):
                self._add_text_to_slide(slide, child, left, top, child_width, child_height, style)
            elif tag == 'img':
                self._add_image_to_slide(slide, child, left, top, child_width, child_height)
            elif tag == 'svg':
                pass

    def _parse_emu(self, value, parent_value):
        """解析为 EMU 单位"""
        if not value or value in ('auto', 'inherit', ''):
            return 0
        value = value.strip()
        match = re.match(r'([\d.]+)(px|em|rem|vw|vh|pt|cm|%)?', value)
        if not match:
            return 0
        num = float(match.group(1))
        unit = match.group(2) or 'px'
        if unit == 'px':
            return int(num * 914400 / 96)
        elif unit == 'pt':
            return int(num * 12700)
        elif unit == 'em' or unit == 'rem':
            return int(num * 16 * 914400 / 96)
        elif unit == 'vw':
            return int(num / 100 * 9144000)
        elif unit == 'vh':
            return int(num / 100 * 5143500)
        elif unit == 'cm':
            return int(num * 360000)
        elif unit == '%':
            return int(num / 100 * parent_value)
        return 0

    def _add_text_to_slide(self, slide, elem, x, y, width, height, style):
        """添加文本框"""
        if width < Emu(10000):
            width = Emu(914400)

        text = self._elem_text(elem)
        if not text.strip():
            return

        font_size = self.resolver.parse_size(style.get('font-size', '24pt'), Pt(24))
        font_bold = 'bold' in style.get('font-weight', '').lower()
        font_italic = 'italic' in style.get('font-style', '').lower()
        font_color = self.resolver.parse_color(style.get('color', '#000000'))
        text_align = style.get('text-align', 'left')
        font_family = style.get('font-family', 'Arial')
        line_height = style.get('line-height', '1.2')

        try:
            txBox = slide.shapes.add_textbox(x, y, width, max(height, Emu(500000)))
            tf = txBox.text_frame
            tf.word_wrap = True

            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text

            run.font.size = font_size
            run.font.bold = font_bold
            run.font.italic = font_italic
            if font_color:
                run.font.color.rgb = font_color
            run.font.name = self._clean_font_name(font_family)

            if text_align == 'center':
                p.alignment = PP_ALIGN.CENTER
            elif text_align == 'right':
                p.alignment = PP_ALIGN.RIGHT

            from pptx.util import Pt as PtUnit
            lh = float(line_height) if line_height else 1.2
            tf.paragraphs[0].line_spacing = PtUnit(font_size.pt * lh)
        except Exception:
            pass

    def _add_image_to_slide(self, slide, elem, x, y, width, height):
        """添加图片占位符（PPTX 不直接支持 SVG 内联）"""
        src = elem.get('src', '')
        alt = elem.get('alt', 'image')
        if not src:
            return

        try:
            from pptx.util import Inches
            pic = slide.shapes.add_picture(src, x, y, width, height)
        except Exception:
            shape = slide.shapes.add_shape(
                1, x, y, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(200, 200, 200)
            tf = shape.text_frame
            tf.text = f'[{alt}]'
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _clean_font_name(self, font_family):
        """清理字体名称"""
        fonts = font_family.split(',')
        primary = fonts[0].strip().strip('"\'')
        replacements = {
            'playfair display': 'Playfair Display',
            'fraunces': 'Fraunces',
            'dm serif display': 'DM Serif Display',
            'cormorant garamond': 'Cormorant Garamond',
            'syne': 'Syne',
            'bebas neue': 'Bebas Neue',
            'dm sans': 'DM Sans',
            'outfit': 'Outfit',
            'figtree': 'Figtree',
            'epilogue': 'Epilogue',
            'noto serif sc': 'Noto Serif SC',
            'noto sans sc': 'Noto Sans SC',
            'arial': 'Arial',
            'helvetica': 'Helvetica',
        }
        lower = primary.lower()
        for key, val in replacements.items():
            if key in lower:
                return val
        return primary


def main():
    if len(sys.argv) < 2:
        print("用法: python html-to-pptx.py <input.html> [output.pptx]")
        sys.exit(1)

    input_html = sys.argv[1]
    output_pptx = sys.argv[2] if len(sys.argv) > 2 else str(Path(input_html).with_suffix('.pptx'))

    html_path = Path(input_html)
    if not html_path.exists():
        print(f"文件不存在: {input_html}")
        sys.exit(1)

    html_content = html_path.read_text(encoding='utf-8')

    try:
        from pptx import Presentation
    except ImportError:
        print("需要 python-pptx 库: pip install python-pptx")
        sys.exit(1)

    converter = HtmlToPptx(html_content)
    result_path = converter.convert(output_pptx)
    print(f"PPTX 已保存: {result_path}")


if __name__ == '__main__':
    main()
