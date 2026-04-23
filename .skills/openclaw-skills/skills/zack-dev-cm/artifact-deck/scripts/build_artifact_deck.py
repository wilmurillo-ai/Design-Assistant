#!/usr/bin/env python3
"""Build a PPTX deck from an Artifact Deck manifest."""

from __future__ import annotations

import argparse
from pathlib import Path
from typing import Any, Dict, List

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from deck_common import ensure_existing_path, load_json, slide_titles, write_json


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(prs: Presentation, title: str, bullets: List[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    text_frame = slide.shapes.placeholders[1].text_frame
    text_frame.clear()
    for index, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(20)


def add_image_slide(prs: Presentation, title: str, image_path: Path, caption: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_picture(str(image_path), Inches(0.6), Inches(1.4), width=Inches(12.1))
    if caption:
        box = slide.shapes.add_textbox(Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.5))
        text_frame = box.text_frame
        text_frame.text = caption
        text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        text_frame.paragraphs[0].font.size = Pt(14)


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--manifest", required=True, help="Artifact Deck manifest JSON.")
    parser.add_argument("--deck-out", required=True, help="Output PPTX path.")
    parser.add_argument("--out", required=True, help="Output build JSON path.")
    args = parser.parse_args()

    manifest_path = ensure_existing_path(args.manifest, "--manifest")
    manifest = load_json(manifest_path)
    deck_out = Path(args.deck_out).expanduser().resolve()
    out_path = Path(args.out).expanduser().resolve()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title = str(manifest.get("title") or "").strip()
    subtitle = str(manifest.get("subtitle") or "").strip()
    add_title_slide(prs, title, subtitle)

    slides = manifest.get("slides") or []
    for slide in slides:
        slide_title = str(slide.get("title") or "").strip()
        bullets = [str(item).strip() for item in (slide.get("bullets") or []) if str(item).strip()]
        if slide_title and bullets:
            add_bullet_slide(prs, slide_title, bullets)

    images = manifest.get("images") or []
    for image in images:
        image_title = str(image.get("title") or "").strip()
        image_path = ensure_existing_path(str(image.get("path") or ""), "--image path")
        caption = str(image.get("caption") or "").strip()
        if image_title:
            add_image_slide(prs, image_title, image_path, caption)

    deck_out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(deck_out))

    payload: Dict[str, Any] = {
        "manifest": str(manifest_path),
        "deck_path": str(deck_out),
        "slide_count": len(prs.slides),
        "image_slide_count": len(images),
        "slide_titles": slide_titles(manifest),
    }
    write_json(out_path, payload)
    print(out_path)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
