#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
知识库文档导入脚本
将 Word/Excel/PPT/PDF/MD 转换为 Markdown 格式
图片上传到图床，生成外部 URL 链接

依赖：
    pip3 install python-docx python-pptx openpyxl pdfplumber

配置：
    复制 config.py.example 为 config.py，并填入实际值
"""

import os
import sys
import argparse
import time
import base64
import uuid
from pathlib import Path
from urllib.parse import urljoin
from urllib.request import urlopen, Request
from urllib.error import URLError, HTTPError

# 导入配置
try:
    from config import DUFs_CONFIG, DEFAULT_KNOWLEDGE_BASE, UPLOAD_PATH
except ImportError:
    # 使用默认值（仅供测试）
    DUFs_CONFIG = {
        "server_url": os.getenv("DUFS_SERVER_URL", "http://localhost:5000"),
        "timeout": 30,
        "retry_times": 3,
        "retry_delay": 2,
    }
    DEFAULT_KNOWLEDGE_BASE = os.getenv(
        "KNOWLEDGE_BASE_PATH",
        os.path.expanduser("~/Obsidian/Obsidian")
    )
    UPLOAD_PATH = "/Picture"


# ============================================================
# 分类规则配置
# ============================================================

# 分类关键词映射（按优先级排序）
CATEGORY_KEYWORDS = {
    "申报方案": ["申报", "投标", "建设方案申请", "申请"],
    "解决方案": ["方案", "解决", "产品", "介绍"],
    "技术文档": ["技术", "文档", "手册", "使用", "部署"],
    "正式文件": ["正式", "原始", "V1", "V2", "V3"],
    "世校赛": ["世校赛", "技能大赛", "大赛", "总决赛", "赛道", "承办校"],
    "职业教育": ["职业", "教育", "教学", "人才培养", "专业建设"],
}


def auto_detect_category(filename: str) -> str:
    """
    根据文件名自动检测分类
    返回: 分类名称，如果未匹配则返回空字符串
    """
    # 移除扩展名
    name = os.path.splitext(filename)[0]
    
    for category, keywords in CATEGORY_KEYWORDS.items():
        for keyword in keywords:
            if keyword in name:
                return category
    
    return ""


def get_output_subdir(filename: str, category: str = None) -> str:
    """
    获取输出子目录
    1. 如果指定了 category，直接使用
    2. 否则根据文件名自动检测
    3. 检测不到则返回空字符串（直接保存到根目录）
    """
    if category:
        return category
    
    detected = auto_detect_category(filename)
    return detected

# 文档解析库
try:
    from docx import Document as DocxDocument
except Exception as e:
    print(f"警告: python-docx 导入失败 - {e}")
    DocxDocument = None

try:
    from pptx import Presentation
except Exception as e:
    print(f"警告: python-pptx 导入失败 - {e}")
    Presentation = None

try:
    import openpyxl
except Exception as e:
    print(f"警告: openpyxl 导入失败 - {e}")
    openpyxl = None

try:
    import pdfplumber
except Exception as e:
    print(f"警告: pdfplumber 导入失败 - {e}")
    pdfplumber = None


# ============================================================
# 图床上传
# ============================================================

def upload_to_image_host(image_path: str) -> str:
    """
    上传图片到图床
    返回: 成功返回图片URL，失败返回 None
    """
    if not os.path.exists(image_path):
        print(f"  警告: 图片文件不存在 - {image_path}")
        return None

    server_url = DUFs_CONFIG["server_url"]

    # 生成唯一文件名
    ext = os.path.splitext(image_path)[1].lower()
    if not ext or ext == '.jpeg':
        ext = '.jpg'
    filename = f"{uuid.uuid4().hex}{ext}"

    # 上传路径
    url = f"{server_url}{UPLOAD_PATH}/{filename}"

    for attempt in range(DUFs_CONFIG["retry_times"]):
        try:
            with open(image_path, 'rb') as f:
                image_data = f.read()

            # Dufs 使用 PUT 请求上传
            req = Request(
                url,
                data=image_data,
                method='PUT',
                headers={'Content-Type': 'application/octet-stream'},
            )

            with urlopen(req, timeout=DUFs_CONFIG["timeout"]) as response:
                if response.status in (200, 201):
                    return url

            print(f"  警告: 上传响应状态异常 - {response.status}")

        except HTTPError as e:
            print(f"  上传失败 (HTTP {e.code}), 重试 {attempt + 1}/{DUFs_CONFIG['retry_times']}...")
        except URLError as e:
            print(f"  上传失败 (网络错误: {e.reason}), 重试 {attempt + 1}/{DUFs_CONFIG['retry_times']}...")
        except Exception as e:
            print(f"  上传失败 ({type(e).__name__}: {e}), 重试 {attempt + 1}/{DUFs_CONFIG['retry_times']}...")

        if attempt < DUFs_CONFIG["retry_times"] - 1:
            time.sleep(DUFs_CONFIG["retry_delay"])

    print(f"  错误: 图片上传失败 - {image_path}")
    return None


def upload_images_to_host(image_paths: list) -> list:
    """
    批量上传图片到图床
    返回: 列表，每个元素是 (原路径, 图片URL) 元组
    """
    results = []

    for i, img_path in enumerate(image_paths, 1):
        print(f"  正在上传第 {i}/{len(image_paths)} 张图片到图床...")

        url = upload_to_image_host(img_path)
        if url:
            results.append((img_path, url))
            print(f"    ✅ {url}")
        else:
            # 上传失败，降级为 base64
            print(f"    ⚠️ 上传失败，使用 Base64 内嵌...")
            img_b64 = image_to_base64(img_path)
            results.append((img_path, img_b64))

    return results


# ============================================================
# 文档解析
# ============================================================

def extract_images_from_docx(doc, output_dir: str) -> list:
    """从 Word 文档提取图片"""
    if not os.path.exists(output_dir):
        os.makedirs(output_dir, exist_ok=True)

    image_paths = []

    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            try:
                image = rel.target_part.blob
                ext = rel.target_part.content_type.split('/')[-1]
                if ext == 'jpeg':
                    ext = 'jpg'

                timestamp = int(time.time() * 1000)
                img_name = f"word_{timestamp}_{len(image_paths) + 1}.{ext}"
                img_path = os.path.join(output_dir, img_name)

                with open(img_path, 'wb') as f:
                    f.write(image)

                image_paths.append(img_path)
            except Exception as e:
                print(f"  警告: 提取图片失败 - {e}")

    return image_paths


def extract_images_from_pptx(prs, output_dir: str) -> list:
    """从 PPT 提取图片"""
    if not os.path.exists(output_dir):
        os.makedirs(output_dir, exist_ok=True)

    image_paths = []

    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                try:
                    image = shape.image
                    ext = image.ext.lower()
                    if ext == 'jpeg':
                        ext = 'jpg'

                    timestamp = int(time.time() * 1000)
                    img_name = f"ppt_{timestamp}_{slide_num}_{len(image_paths) + 1}.{ext}"
                    img_path = os.path.join(output_dir, img_name)

                    with open(img_path, 'wb') as f:
                        f.write(image.blob)

                    image_paths.append(img_path)
                except Exception as e:
                    print(f"  警告: 提取图片失败 - {e}")

    return image_paths


def extract_text_from_docx(file_path: str) -> str:
    """从 Word 文档提取文本"""
    if DocxDocument is None:
        return "[错误] python-docx 未安装，请运行: pip3 install python-docx"

    try:
        doc = DocxDocument(file_path)
        content = []

        filename = os.path.basename(file_path)
        content.append(f"# {os.path.splitext(filename)[0]}\n")
        content.append(f"**源文件:** {filename}\n")
        content.append("---\n")

        # 提取段落文本
        for para in doc.paragraphs:
            if para.text.strip():
                content.append(para.text)

        # 处理表格
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join([cell.text.strip() for cell in row.cells])
                content.append(f"| {row_text} |")

        # 提取图片
        assets_dir = os.path.join("/tmp", f"docx_assets_{int(time.time())}")
        os.makedirs(assets_dir, exist_ok=True)

        print(f"  正在提取图片...")
        image_paths = extract_images_from_docx(doc, assets_dir)

        if image_paths:
            content.append("\n\n## 图片内容\n")

            uploaded = upload_images_to_host(image_paths)

            for i, (img_path, img_url) in enumerate(uploaded, 1):
                content.append(f"\n### 图片 {i}: {os.path.basename(img_path)}\n")
                content.append(f"![{os.path.basename(img_path)}]({img_url})\n")

        return "\n\n".join(content)
    except Exception as e:
        return f"[错误] 读取 Word 文档失败: {e}"


def extract_text_from_pptx(file_path: str, extract_images: bool = False) -> str:
    """从 PPT 提取文本
    
    Args:
        file_path: PPT 文件路径
        extract_images: 是否提取图片（默认 False，PPT 图片较多，节省时间和空间）
    """
    if Presentation is None:
        return "[错误] python-pptx 未安装，请运行: pip3 install python-pptx"

    try:
        prs = Presentation(file_path)
        content = []

        filename = os.path.basename(file_path)
        content.append(f"# {os.path.splitext(filename)[0]}\n")
        content.append(f"**源文件:** {filename}\n")
        content.append(f"**总页数:** {len(prs.slides)}\n")
        content.append("---\n")

        for i, slide in enumerate(prs.slides, 1):
            content.append(f"\n## 第 {i} 页\n")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content.append(shape.text)

            content.append("\n---")

        # PPT 默认不提取图片（图片较多，会导致上传慢且占用图床空间）
        if extract_images:
            assets_dir = os.path.join("/tmp", f"pptx_assets_{int(time.time())}")
            os.makedirs(assets_dir, exist_ok=True)

            print(f"  正在提取图片...")
            image_paths = extract_images_from_pptx(prs, assets_dir)

            if image_paths:
                content.append("\n\n## 图片内容\n")

                uploaded = upload_images_to_host(image_paths)

                for i, (img_path, img_url) in enumerate(uploaded, 1):
                    content.append(f"\n### 图片 {i}: {os.path.basename(img_path)}\n")
                    content.append(f"![{os.path.basename(img_path)}]({img_url})\n")

        return "\n".join(content)
    except Exception as e:
        return f"[错误] 读取 PPT 失败: {e}"


def extract_text_from_xlsx(file_path: str) -> str:
    """从 Excel 提取文本"""
    if openpyxl is None:
        return "[错误] openpyxl 未安装，请运行: pip3 install openpyxl"

    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        content = []

        filename = os.path.basename(file_path)
        content.append(f"# {os.path.splitext(filename)[0]}\n")
        content.append(f"**源文件:** {filename}\n")
        content.append("---\n")

        for sheet_name in wb.sheetnames:
            content.append(f"\n## 工作表: {sheet_name}\n")

            ws = wb[sheet_name]

            rows = []
            for row in ws.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    rows.append(row)

            if rows:
                for i, row in enumerate(rows):
                    row_str = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    content.append(f"| {row_str} |")

                    if i == 0:
                        content.append("|" + "|".join(["---"] * len(row)) + "|")

        return "\n".join(content)
    except Exception as e:
        return f"[错误] 读取 Excel 失败: {e}"


def extract_images_from_pdf(file_path: str, output_dir: str) -> list:
    """从 PDF 提取嵌入的图片（使用 PyMuPDF）"""
    try:
        import fitz
    except ImportError:
        print("  警告: PyMuPDF 未安装，跳过 PDF 图片提取")
        print("  安装命令: pip3 install PyMuPDF --break-system-packages")
        return []
    
    if not os.path.exists(output_dir):
        os.makedirs(output_dir, exist_ok=True)
    
    image_paths = []
    
    try:
        doc = fitz.open(file_path)
        print(f"  正在提取 PDF 嵌入图片...")
        
        img_count = 0
        for page_num in range(len(doc)):
            page = doc[page_num]
            images = page.get_images()
            
            for img_index, img in enumerate(images):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                
                img_count += 1
                img_name = f"pdf_img_p{page_num+1}_{img_index+1}.{image_ext}"
                img_path = os.path.join(output_dir, img_name)
                
                with open(img_path, "wb") as f:
                    f.write(image_bytes)
                
                image_paths.append(img_path)
        
        doc.close()
        print(f"  PDF 图片提取完成，共 {img_count} 张")
        
    except Exception as e:
        print(f"  错误: PDF 图片提取失败 - {e}")
    
    return image_paths


def extract_text_from_pdf(file_path: str) -> str:
    """从 PDF 提取文本和图片"""
    if pdfplumber is None:
        return "[错误] pdfplumber 未安装，请运行: pip3 install pdfplumber"

    try:
        content = []

        filename = os.path.basename(file_path)
        content.append(f"# {os.path.splitext(filename)[0]}\n")
        content.append(f"**源文件:** {filename}\n")
        content.append("---\n")

        with pdfplumber.open(file_path) as pdf:
            content.append(f"**总页数:** {len(pdf.pages)}\n")

            for i, page in enumerate(pdf.pages, 1):
                content.append(f"\n## 第 {i} 页\n")

                text = page.extract_text()
                if text:
                    content.append(text)

                # 提取表格
                tables = page.extract_tables()
                if tables:
                    content.append("\n### 表格\n")
                    for table in tables:
                        for row in table:
                            row_str = " | ".join([str(cell) if cell else "" for cell in row])
                            content.append(f"| {row_str} |")

        # 提取图片
        assets_dir = os.path.join("/tmp", f"pdf_assets_{int(time.time())}")
        os.makedirs(assets_dir, exist_ok=True)

        print(f"  正在提取图片...")
        image_paths = extract_images_from_pdf(file_path, assets_dir)

        if image_paths:
            content.append("\n\n## PDF 嵌入图片\n")
            content.append("(以下为 PDF 中提取的原始图片)\n")

            uploaded = upload_images_to_host(image_paths)

            for i, (img_path, img_url) in enumerate(uploaded, 1):
                content.append(f"\n### 图片 {i}: {os.path.basename(img_path)}\n")
                content.append(f"![图片 {i}]({img_url})\n")

        return "\n".join(content)
    except Exception as e:
        return f"[错误] 读取 PDF 失败: {e}"


def extract_text_from_md(file_path: str) -> str:
    """读取 Markdown 文件"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        return f"[错误] 读取 Markdown 失败: {e}"


def image_to_base64(image_path: str) -> str:
    """将图片转换为 base64 格式"""
    try:
        with open(image_path, "rb") as f:
            img_data = f.read()

        ext = os.path.splitext(image_path)[1].lower()
        mime_types = {
            '.jpg': 'image/jpeg',
            '.jpeg': 'image/jpeg',
            '.png': 'image/png',
            '.gif': 'image/gif',
            '.webp': 'image/webp',
            '.bmp': 'image/bmp',
        }
        mime_type = mime_types.get(ext, 'image/png')

        img_b64 = base64.b64encode(img_data).decode("utf-8")
        return f"data:{mime_type};base64,{img_b64}"
    except Exception as e:
        return f"[图片读取错误: {e}]"


def get_file_type(file_path: str) -> str:
    """获取文件类型"""
    ext = os.path.splitext(file_path)[1].lower()

    type_map = {
        '.docx': 'docx',
        '.doc': 'doc',
        '.xlsx': 'xlsx',
        '.xls': 'xls',
        '.pptx': 'pptx',
        '.ppt': 'ppt',
        '.pdf': 'pdf',
        '.md': 'md',
        '.markdown': 'md',
    }

    return type_map.get(ext, 'unknown')


def convert_file(file_path: str, output_dir: str = None, category: str = None) -> str:
    """转换单个文件
    
    Args:
        file_path: 要转换的文件路径
        output_dir: 输出目录（默认为知识库根目录）
        category: 指定分类目录（如"解决方案/UBT"），优先级高于自动检测
    """
    if not os.path.exists(file_path):
        return f"[错误] 文件不存在: {file_path}"

    file_type = get_file_type(file_path)

    if file_type == 'unknown':
        return f"[错误] 不支持的格式: {file_path}"

    if output_dir is None:
        output_dir = DEFAULT_KNOWLEDGE_BASE

    filename = os.path.basename(file_path)
    
    # 获取分类子目录
    subdir = get_output_subdir(filename, category)
    
    # 构建完整输出路径
    if subdir:
        # 如果 category 包含子目录（如"解决方案/UBT"），直接使用
        if "/" in subdir:
            target_dir = os.path.join(output_dir, subdir)
        else:
            # 否则在对应一级分类下创建子目录（以文件名命名）
            file_stem = os.path.splitext(filename)[0]
            target_dir = os.path.join(output_dir, subdir, file_stem)
    else:
        target_dir = output_dir

    os.makedirs(target_dir, exist_ok=True)

    print(f"正在处理: {filename}")
    print(f"  分类: {subdir if subdir else '（根目录）'}")

    if file_type == 'docx':
        content = extract_text_from_docx(file_path)
    elif file_type == 'xlsx':
        content = extract_text_from_xlsx(file_path)
    elif file_type == 'pptx':
        content = extract_text_from_pptx(file_path)
    elif file_type == 'pdf':
        content = extract_text_from_pdf(file_path)
    elif file_type == 'md':
        content = extract_text_from_md(file_path)
    else:
        return f"[错误] 暂不支持此格式: {file_type}"

    # 保存文件
    output_filename = f"{os.path.splitext(filename)[0]}.md"
    output_file = os.path.join(target_dir, output_filename)

    try:
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(content)
        return f"[成功] 已保存到: {output_file}"
    except Exception as e:
        return f"[错误] 保存文件失败: {e}"


def batch_convert(source_dir: str, output_dir: str = None, category: str = None) -> str:
    """批量转换目录下的所有支持的文件
    
    Args:
        source_dir: 源目录
        output_dir: 输出目录
        category: 指定分类目录
    """
    if not os.path.isdir(source_dir):
        return f"[错误] 目录不存在: {source_dir}"

    if output_dir is None:
        output_dir = DEFAULT_KNOWLEDGE_BASE

    os.makedirs(output_dir, exist_ok=True)

    supported_exts = ['.docx', '.xlsx', '.pptx', '.pdf', '.md', '.markdown']
    results = []

    for root, dirs, files in os.walk(source_dir):
        for file in files:
            ext = os.path.splitext(file)[1].lower()
            if ext in supported_exts:
                file_path = os.path.join(root, file)

                rel_path = os.path.relpath(file_path, source_dir)
                rel_dir = os.path.dirname(rel_path)

                if rel_dir:
                    target_dir = os.path.join(output_dir, rel_dir)
                else:
                    target_dir = output_dir

                result = convert_file(file_path, target_dir, category)
                results.append(result)
                results.append("")

    return "\n".join(results)


def main():
    parser = argparse.ArgumentParser(
        description='知识库文档导入工具（图片上传到图床，支持自动分类）',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
分类说明:
  自动分类规则（根据文件名关键词）:
    - 申报方案: 申报、投标、建设方案申请
    - 解决方案: 方案、解决、产品、介绍
    - 技术文档: 技术、文档、手册、使用、部署
    - 正式文件: 正式、原始、V1/V2/V3版本
    - 世校赛: 世校赛、技能大赛、总决赛、赛道、承办校
    - 职业教育: 职业、教育、教学、人才培养、专业建设

  也支持手动指定分类: --category 申报方案/无人机

示例:
  python3 import_doc.py /path/to/document.pdf
  python3 import_doc.py /path/to/document.pdf --category 解决方案/UBT
  python3 import_doc.py /path/to/document.docx /path/to/output
  python3 import_doc.py --batch /path/to/folder
  python3 import_doc.py --batch /path/to/folder --category 申报方案

依赖:
  pip3 install python-docx python-pptx openpyxl pdfplumber PyMuPDF
        """
    )
    parser.add_argument('file', nargs='?', help='要转换的文件路径')
    parser.add_argument('--batch', action='store_true', help='批量转换模式')
    parser.add_argument('output', nargs='?', help='输出目录（可选）')
    parser.add_argument('--category', '-c', help='指定分类目录（如：解决方案/UBT、申报方案/无人机）')

    args = parser.parse_args()

    if args.batch:
        if not args.file:
            print("[错误] 批量模式需要指定源目录")
            sys.exit(1)
        result = batch_convert(args.file, args.output, args.category)
        print(result)
    else:
        if not args.file:
            print("用法:")
            print("  python3 import_doc.py <文件> [输出目录] [--category 分类]")
            print("  python3 import_doc.py --batch <源目录> [输出目录] [--category 分类]")
            print("\n详细帮助: python3 import_doc.py --help")
            sys.exit(1)
        result = convert_file(args.file, args.output, args.category)
        print(result)


if __name__ == '__main__':
    main()
