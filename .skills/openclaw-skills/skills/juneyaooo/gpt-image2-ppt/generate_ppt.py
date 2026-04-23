#!/usr/bin/env python3
"""
PPT Generator - Generate PPT slide images using OpenAI gpt-image-2 (Images API).

This script generates PPT slide images based on a slide plan and style template,
then creates an HTML viewer for playback.
"""

import argparse
import json
import os
import re
import sys
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, Optional

from dotenv import load_dotenv


# =============================================================================
# Constants
# =============================================================================

DEFAULT_TEMPLATE_PATH = "templates/viewer.html"
OUTPUT_BASE_DIR = "outputs"

SCRIPT_DIR = Path(__file__).parent
CWD = Path.cwd()


# =============================================================================
# Environment Configuration
# =============================================================================

def find_and_load_env() -> bool:
    """Load .env from explicit, scoped locations only.

    Search order (first match wins, no parent-directory walking):
    1. $GPT_IMAGE2_PPT_ENV (explicit override)
    2. <script_dir>/.env  -- the skill's own .env
    3. ~/.claude/skills/gpt-image2-ppt-skills/.env  -- default Claude Code install
    4. ~/skills/gpt-image2-ppt/.env  -- default OpenClaw install

    Intentionally does NOT walk parent directories of the script or cwd, to avoid
    accidentally loading unrelated secrets from a surrounding project's .env.
    """
    env_locations = []
    explicit = os.getenv("GPT_IMAGE2_PPT_ENV")
    if explicit:
        env_locations.append(Path(explicit))
    env_locations.extend([
        SCRIPT_DIR / ".env",
        Path.home() / ".claude" / "skills" / "gpt-image2-ppt-skills" / ".env",
        Path.home() / "skills" / "gpt-image2-ppt" / ".env",
    ])

    for env_path in env_locations:
        if env_path.exists():
            load_dotenv(env_path, override=True)
            print(f"Loaded environment from: {env_path}")
            return True

    # 故意不再 load_dotenv() 兜底----避免 dotenv 的隐式 cwd/父目录搜索
    # 误吃无关项目的 .env。如果 scoped 位置都没 .env，就只用进程 env vars。
    print("Warning: No .env file found in scoped locations; using process env vars only.")
    return False


# =============================================================================
# Style Template
# =============================================================================

def load_style_template(style_path: str) -> str:
    """Extract the '## 基础提示词模板' section from a style markdown file."""
    with open(style_path, "r", encoding="utf-8") as f:
        content = f.read()

    base_prompt_marker = "## 基础提示词模板"
    start_idx = content.find(base_prompt_marker)

    if start_idx == -1:
        print("Warning: '## 基础提示词模板' section not found, using fallback extraction")
        start_idx = content.find("## ")
        end_idx = content.find("## ", start_idx + 3)
        if start_idx == -1 or end_idx == -1:
            return content
        return content[start_idx + 3:end_idx].strip()

    section_start = start_idx + len(base_prompt_marker)
    next_section_idx = content.find("\n## ", section_start)

    if next_section_idx == -1:
        extracted = content[section_start:]
    else:
        extracted = content[section_start:next_section_idx]

    return extracted.strip()


# =============================================================================
# Prompt Generation
# =============================================================================

LANGUAGE_FONT_RULE = """

【强制语言与字体要求】
1. 幻灯片上所有文字必须使用简体中文，严禁出现任何英文单词或句子（产品名称等专有名词可保留英文，其余一律用中文）。
2. 中文字体使用思源黑体（Source Han Sans）或苹方（PingFang SC），字形清晰、笔画规整，严禁使用草书、艺术字或变形字体。
3. 标题字体粗体，正文字体常规，字号对比清晰，确保在演示场景下可读性极高。
"""


def generate_prompt(
    style_template: str,
    page_type: str,
    content_text: str,
    slide_number: int,
    total_slides: int,
) -> str:
    """Generate a complete prompt for a single slide.

    Per-page composition rules live inside each style's `## 基础提示词模板`
    (cover / content / data sub-blocks); this function只在尾部追加一个中性的
    page_type 提示，让模型按本风格自己的规范处理。
    """
    is_cover = page_type == "cover" or slide_number == 1
    is_data = page_type == "data" or slide_number == total_slides
    if is_cover:
        label = "封面页（cover）"
        hint = "标题/副标题处理为视觉焦点，按本风格的封面构图规范处理。"
    elif is_data:
        label = "数据页（data）"
        hint = "突出关键数字、对比或结论；按本风格的数据/总结构图规范处理。"
    else:
        label = "内容页（content）"
        hint = "把要点按本风格的内容构图规范结构化呈现，注意层级、对齐、留白。"

    return (
        style_template
        + "\n\n---\n\n"
        + f"现在请生成本组中的【{label}】，{hint}\n"
        + "本页要呈现的内容如下（请按本风格美学重新设计版式，不要原样照搬文本节奏）：\n\n"
        + content_text
        + LANGUAGE_FONT_RULE
    )


# =============================================================================
# Image Generation
# =============================================================================

def generate_slide(
    prompt: str,
    slide_number: int,
    output_dir: str,
    reference_image_path: Optional[str] = None,
) -> Optional[str]:
    """Generate a single PPT slide image using gpt-image-2."""
    sys.path.insert(0, str(SCRIPT_DIR))
    from image_generator import GptImage2Generator

    print(f"  Generating slide {slide_number} via gpt-image-2 ...")

    try:
        generator = GptImage2Generator(aspect_ratio="16:9")
        image_path = os.path.join(output_dir, "images", f"slide-{slide_number:02d}.png")

        scene_data = {
            "index": slide_number,
            "image_prompt": prompt,
        }
        generator.generate_scene_image(
            scene_data=scene_data,
            output_path=image_path,
            reference_image_path=reference_image_path,
        )
        print(f"  Slide {slide_number} saved: {image_path}")
        return image_path

    except Exception as e:
        print(f"  Slide {slide_number} failed: {e}")
        return None


# =============================================================================
# Output Generation
# =============================================================================

def generate_viewer_html(
    output_dir: str,
    slide_count: int,
    template_path: str,
) -> str:
    """Generate HTML viewer for slides playback."""
    if not os.path.isabs(template_path):
        template_path = str(SCRIPT_DIR / template_path)

    with open(template_path, "r", encoding="utf-8") as f:
        html_template = f.read()

    slides_list = [f"'images/slide-{i:02d}.png'" for i in range(1, slide_count + 1)]

    html_content = html_template.replace(
        "/* IMAGE_LIST_PLACEHOLDER */",
        ",\n            ".join(slides_list),
    )

    html_path = os.path.join(output_dir, "index.html")
    with open(html_path, "w", encoding="utf-8") as f:
        f.write(html_content)

    print(f"  Viewer HTML generated: {html_path}")
    return html_path


def save_prompts(output_dir: str, prompts_data: Dict[str, Any]) -> str:
    """Save all prompts to JSON file."""
    prompts_path = os.path.join(output_dir, "prompts.json")
    with open(prompts_path, "w", encoding="utf-8") as f:
        json.dump(prompts_data, f, ensure_ascii=False, indent=2)
    print(f"  Prompts saved: {prompts_path}")
    return prompts_path


def generate_pptx(
    output_dir: str,
    slide_count: int,
    title: str = "Untitled",
) -> Optional[str]:
    """把 images/slide-XX.png 打包成 16:9 .pptx，每页填满。

    需要 python-pptx；如果没装就跳过并提示。
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Emu
    except ImportError:
        print("(!)  跳过 .pptx 生成（缺 python-pptx，pip install python-pptx 后重试）")
        return None

    prs = Presentation()
    # 标准 16:9 PPT 尺寸：13.333 x 7.5 英寸（1280x720pt）
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # 完全空白布局

    img_dir = os.path.join(output_dir, "images")
    added = 0
    for i in range(1, slide_count + 1):
        img_path = os.path.join(img_dir, f"slide-{i:02d}.png")
        if not os.path.exists(img_path):
            print(f"  跳过 slide-{i:02d}.png（文件不存在）")
            continue
        slide = prs.slides.add_slide(blank)
        # 图片填满整页（如果原图比例不是 16:9，python-pptx 默认按指定 width/height 拉伸）
        slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        added += 1

    if added == 0:
        print("(!)  没有可用图片，未生成 .pptx")
        return None

    # 文件名用 plan title（去除非法字符）
    safe_title = re.sub(r"[^\w\u4e00-\u9fff\-]+", "_", title)[:60] or "deck"
    pptx_path = os.path.join(output_dir, f"{safe_title}.pptx")
    prs.save(pptx_path)
    print(f"  📑 PPTX generated: {pptx_path}  ({added} slides)")
    return pptx_path


# =============================================================================
# Main Entry Point
# =============================================================================

def create_argument_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="PPT Generator - Generate PPT images using OpenAI gpt-image-2",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Example usage:
  python generate_ppt.py --plan slides_plan.json --style styles/gradient-glass.md
  python generate_ppt.py --plan slides_plan.json --style styles/clean-tech-blue.md --slides 1,3,5

Environment variables (set in .env file):
  OPENAI_BASE_URL:        Images API base URL (default: https://api.openai.com)
  OPENAI_API_KEY:         API key (required)
  GPT_IMAGE_MODEL_NAME:   Model name (default: gpt-image-2)
  GPT_IMAGE_QUALITY:      low / medium / high / auto (default: high)
""",
    )

    parser.add_argument("--plan", required=True, help="Path to slides plan JSON file")
    parser.add_argument("--style", help="Path to style template file (与 --template-pptx 二选一)")
    parser.add_argument("--output", help="Output directory path (default: outputs/TIMESTAMP)")
    parser.add_argument(
        "--template",
        default=DEFAULT_TEMPLATE_PATH,
        help=f"HTML template path (default: {DEFAULT_TEMPLATE_PATH})",
    )
    parser.add_argument(
        "--template-pptx",
        help="用户的 .pptx 模板路径，启用「仿模板」模式",
    )
    parser.add_argument(
        "--template-images",
        help="模板每页 PNG 所在目录（强烈建议传，没有则只读 .pptx XML，不能跑 vision）",
    )
    parser.add_argument(
        "--template-strict",
        action="store_true",
        help="高保真模式：把模板对应页作为 image reference 传给 gpt-image-2 出新图",
    )
    parser.add_argument(
        "--rebuild-template-cache",
        action="store_true",
        help="无视模板缓存重新跑 vision",
    )
    parser.add_argument(
        "--slides",
        help="Only generate specific slides, e.g. '1,3,5'",
    )
    parser.add_argument(
        "--concurrency",
        type=int,
        default=int(os.getenv("GPT_IMAGE_CONCURRENCY", "10")),
        help="并发请求数（默认 10，可用 GPT_IMAGE_CONCURRENCY 环境变量覆盖）",
    )
    parser.add_argument(
        "--no-pptx",
        action="store_true",
        help="不生成 .pptx 文件（默认会自动打包成 16:9 PPTX）",
    )

    return parser


def main() -> None:
    find_and_load_env()

    parser = create_argument_parser()
    args = parser.parse_args()

    # 校验：style 与 template-pptx 至少有一个
    use_template = bool(args.template_pptx or args.template_images)
    if not use_template and not args.style:
        parser.error("必须传 --style 或 --template-pptx / --template-images 至少其一")

    style_template = ""
    if args.style:
        style_path = args.style
        if not os.path.isabs(style_path):
            candidate = SCRIPT_DIR / style_path
            if candidate.exists():
                style_path = str(candidate)
        style_template = load_style_template(style_path)
    else:
        style_path = "(template-derived)"

    # 模板模式：跑 vision 拿 TemplateProfile（带缓存）
    template_profile: Optional[Dict[str, Any]] = None
    if use_template:
        sys.path.insert(0, str(SCRIPT_DIR))
        # 只给了 .pptx 没给 PNG -> 自动渲染到 <cwd>/template_renders/<stem>/
        if args.template_pptx and not args.template_images:
            from render_template import render_pptx_to_pngs
            print(f"🖨️  --template-images 未指定，自动渲染 {args.template_pptx}")
            args.template_images = str(render_pptx_to_pngs(args.template_pptx))
        from template_analyzer import analyze_template
        template_profile = analyze_template(
            pptx_path=args.template_pptx,
            images_dir=args.template_images,
            rebuild=args.rebuild_template_cache,
        )
        if not template_profile.get("layouts"):
            print("(!)  模板分析未产出 layouts（缺 --template-images？），将回退到自由风格 prompt")
            template_profile = None

    with open(args.plan, "r", encoding="utf-8") as f:
        slides_plan = json.load(f)

    if args.output:
        output_dir = args.output
    else:
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        # 默认输出到调用者当前工作目录，而不是 skill 安装目录
        output_dir = str(CWD / OUTPUT_BASE_DIR / timestamp)

    os.makedirs(os.path.join(output_dir, "images"), exist_ok=True)

    slides = slides_plan["slides"]
    total_slides = len(slides)

    if args.slides:
        target_nums = set(int(x.strip()) for x in args.slides.split(","))
        slides = [s for s in slides if s.get("slide_number") in target_nums]

    print("=" * 60)
    print("PPT Generator (gpt-image-2) Started")
    print("=" * 60)
    print(f"Style: {style_path}")
    if template_profile:
        print(f"Template: {template_profile.get('source')} (hash={template_profile.get('source_hash')}, "
              f"{len(template_profile.get('layouts', []))} layouts)")
        print(f"Strict mode: {args.template_strict}")
    print(f"Slides: {len(slides)} / {total_slides}")
    print(f"Output: {output_dir}")
    print(f"Concurrency: {args.concurrency}")
    print("=" * 60)
    print()

    prompts_data: Dict[str, Any] = {
        "metadata": {
            "title": slides_plan.get("title", "Untitled Presentation"),
            "total_slides": total_slides,
            "model": os.getenv("GPT_IMAGE_MODEL_NAME", "gpt-image-2"),
            "style": style_path,
            "template": template_profile.get("source") if template_profile else None,
            "template_strict": args.template_strict if template_profile else False,
            "generated_at": datetime.now().isoformat(),
        },
        "slides": [],
    }

    if template_profile:
        from template_analyzer import (
            match_layout,
            coerce_fields,
            render_prompt_from_template,
            check_layout_reuse,
        )
        # layout 复用检测：在派发任务前打出建议，让用户决定是否中断
        reuse_warnings = check_layout_reuse(slides, template_profile)
        if reuse_warnings:
            print()
            print("=" * 60)
            print("📐 Layout 复用检测（建议尽量做到 1 page : 1 layout）")
            print("=" * 60)
            for w in reuse_warnings:
                print(w)
            print("=" * 60)
            print()

    # 收集所有待跑任务（跳过已存在的）
    pending_tasks = []
    for slide_info in slides:
        slide_number = slide_info["slide_number"]
        page_type = slide_info.get("page_type", "content")
        content_text = slide_info.get("content", "")

        existing = os.path.join(output_dir, "images", f"slide-{slide_number:02d}.png")
        if os.path.exists(existing):
            print(f"Slide {slide_number}: already exists, skipping.")
            prompts_data["slides"].append({
                "slide_number": slide_number,
                "page_type": page_type,
                "content": content_text,
                "prompt": "(skipped - already exists)",
                "image_path": existing,
            })
            continue

        reference_image = None
        if template_profile:
            layout = match_layout(slide_info, template_profile)
            if layout is None:
                # 模板未匹配 -> 回退到 style_template
                prompt = generate_prompt(
                    style_template, page_type, content_text, slide_number, total_slides
                )
                matched_layout_id = None
            else:
                fields = coerce_fields(slide_info, layout)
                prompt = render_prompt_from_template(
                    profile=template_profile,
                    layout=layout,
                    fields=fields,
                    language_rule=LANGUAGE_FONT_RULE.strip(),
                )
                matched_layout_id = layout.get("id")
                if args.template_strict:
                    reference_image = layout.get("reference_image")
        else:
            prompt = generate_prompt(
                style_template, page_type, content_text, slide_number, total_slides
            )
            matched_layout_id = None

        pending_tasks.append({
            "slide_number": slide_number,
            "page_type": page_type,
            "content": content_text,
            "prompt": prompt,
            "reference_image": reference_image,
            "layout_id": matched_layout_id,
        })

    if pending_tasks:
        from concurrent.futures import ThreadPoolExecutor, as_completed

        worker_count = max(1, min(args.concurrency, len(pending_tasks)))
        print(f"📦 派发 {len(pending_tasks)} 个任务到 {worker_count} 个并发 worker...\n")

        results: Dict[int, Optional[str]] = {}

        def _run(task):
            n = task["slide_number"]
            print(f">️  [slide {n}] start ({task['page_type']}{' / ref' if task.get('reference_image') else ''})")
            try:
                path = generate_slide(
                    task["prompt"], n, output_dir,
                    reference_image_path=task.get("reference_image"),
                )
                print(f"[OK] [slide {n}] done")
                return n, path
            except Exception as e:
                print(f"[X] [slide {n}] failed: {e}")
                return n, None

        with ThreadPoolExecutor(max_workers=worker_count) as executor:
            futures = [executor.submit(_run, t) for t in pending_tasks]
            for fut in as_completed(futures):
                n, path = fut.result()
                results[n] = path

        # 按原顺序写回 prompts_data
        for task in pending_tasks:
            n = task["slide_number"]
            prompts_data["slides"].append({
                "slide_number": n,
                "page_type": task["page_type"],
                "content": task["content"],
                "layout_id": task.get("layout_id"),
                "reference_image": task.get("reference_image"),
                "prompt": task["prompt"],
                "image_path": results.get(n),
            })

    # 按 slide_number 排序，保证 prompts.json 与播放顺序一致
    prompts_data["slides"].sort(key=lambda s: s["slide_number"])
    print()

    save_prompts(output_dir, prompts_data)
    generate_viewer_html(output_dir, total_slides, args.template)

    pptx_path = None
    if not args.no_pptx:
        pptx_path = generate_pptx(
            output_dir,
            total_slides,
            title=slides_plan.get("title", "Untitled"),
        )

    print()
    print("=" * 60)
    print("Generation Complete!")
    print("=" * 60)
    print(f"Output directory: {output_dir}")
    print(f"Viewer HTML: {os.path.join(output_dir, 'index.html')}")
    if pptx_path:
        print(f"PPTX file:   {pptx_path}")
    print()
    print("Open viewer in browser:")
    print(f"  open {os.path.join(output_dir, 'index.html')}")
    print()


if __name__ == "__main__":
    main()
