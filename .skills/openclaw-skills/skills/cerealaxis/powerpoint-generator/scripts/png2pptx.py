#!/usr/bin/env python3
"""PNG to PPTX -- 将 PNG 图片打包为 PPTX 演示文稿

每张 PNG 作为独立幻灯片，图片填满整个幻灯片背景。
主要用于跨平台 100% 视觉还原（文字不可编辑）。

用法:
  python png2pptx.py <png_dir_or_file> -o output.pptx
  python png2pptx.py ppt-output/png/ -o ppt-output/presentation-png.pptx
"""

import argparse
import base64
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN


# -------------------------------------------------------------------
# 配置
# -------------------------------------------------------------------
SLIDE_W = 12192000  # 10 inches in EMUs
SLIDE_H = 6858000   # 7.5 inches in EMUs (4:3) 或 6858 (16:9)


# -------------------------------------------------------------------
# 核心转换逻辑
# -------------------------------------------------------------------
def png_to_slide(prs: Presentation, png_path: Path, slide_index: int) -> bool:
    """将一张 PNG 转为幻灯片，PNG 作为背景填满整页。"""
    try:
        from PIL import Image
    except ImportError:
        print("ERROR: Pillow not installed. Run: pip install Pillow", file=sys.stderr)
        return False

    try:
        # 读取 PNG 尺寸
        with Image.open(png_path) as img:
            img_w, img_h = img.size

        # 添加空白幻灯片（使用与图片相同比例）
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 填充背景图片
        left = Emu(0)
        top = Emu(0)
        pic_width = Emu(img_w * 9525)  # 像素转 EMU
        pic_height = Emu(img_h * 9525)

        # 如果图片比例与幻灯片不同，以宽度或高度为准进行裁切
        slide_width_emu = prs.slide_width
        slide_height_emu = prs.slide_height

        # 计算缩放比例（填满）
        scale_w = slide_width_emu / pic_width
        scale_h = slide_height_emu / pic_height
        scale = max(scale_w, scale_h)  # 取较大值确保填满

        final_w = int(pic_width * scale)
        final_h = int(pic_height * scale)

        # 居中裁切
        left = Emu((slide_width_emu - final_w) // 2)
        top = Emu((slide_height_emu - final_h) // 2)

        pic = slide.shapes.add_picture(
            str(png_path),
            left, top,
            width=Emu(final_w),
            height=Emu(final_h)
        )

        # 将图片移到最底层
        spTree = slide.shapes._spTree
        pic_elem = pic._element
        spTree.remove(pic_elem)
        spTree.insert(2, pic_elem)

        return True

    except Exception as e:
        print(f"  [!] {png_path.name}: {e}", file=sys.stderr)
        return False


def convert_dir(png_dir: Path, output_path: Path) -> int:
    """将目录下所有 PNG 转为 PPTX。"""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 清空默认空白幻灯片
    if prs.slides:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    png_files = sorted(png_dir.glob("*.png"))
    if not png_files:
        print(f"WARNING: No PNG files found in {png_dir}", file=sys.stderr)
        return 0

    print(f"Converting {len(png_files)} PNG files to PPTX...")
    success = 0
    for i, png_file in enumerate(png_files, 1):
        if png_to_slide(prs, png_file, i):
            success += 1

    if success > 0:
        prs.save(str(output_path))
        print(f"Saved: {output_path} ({success} slides)")
    else:
        print("ERROR: No slides were created", file=sys.stderr)

    return success


def convert_file(png_file: Path, output_path: Path) -> bool:
    """将单个 PNG 转为 PPTX（单页）。"""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    if prs.slides:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    result = png_to_slide(prs, png_file, 1)
    if result:
        prs.save(str(output_path))
        print(f"Saved: {output_path}")
    return result


def main():
    parser = argparse.ArgumentParser(description="PNG to PPTX converter")
    parser.add_argument("png_source", type=Path, help="PNG file or directory containing PNG files")
    parser.add_argument("-o", "--output", type=Path, required=True, help="Output PPTX file path")
    parser.add_argument("--width", type=int, default=1280, help="Slide width in px (default: 1280)")
    parser.add_argument("--height", type=int, default=720, help="Slide height in px (default: 720)")

    args = parser.parse_args()

    global SLIDE_W, SLIDE_H
    # 将像素转为 EMU（1px = 9525 EMU）
    SLIDE_W = Emu(args.width * 9525)
    SLIDE_H = Emu(args.height * 9525)

    if not args.png_source.exists():
        print(f"ERROR: Source not found: {args.png_source}", file=sys.stderr)
        sys.exit(1)

    # 确保输出目录存在
    args.output.parent.mkdir(parents=True, exist_ok=True)

    if args.png_source.is_file():
        success = convert_file(args.png_source, args.output)
    else:
        success = convert_dir(args.png_source, args.output)

    sys.exit(0 if success > 0 else 1)


if __name__ == "__main__":
    main()
