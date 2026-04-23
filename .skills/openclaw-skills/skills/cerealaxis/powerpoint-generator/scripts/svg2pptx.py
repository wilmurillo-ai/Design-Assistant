#!/usr/bin/env python3
"""SVG to PPTX -- 将 SVG 元素解析为原生 OOXML 形状

支持: rect, text+tspan, circle, ellipse, line, path, image(data URI + file)
      linearGradient, radialGradient, transform(translate/scale/matrix)
      group opacity 传递, 首屏 rect 自动设为幻灯片背景

用法:
  python svg2pptx.py <svg_dir_or_file> -o output.pptx
"""

import argparse
import base64
import io
import math
import re
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

# -------------------------------------------------------------------
# 常量
# -------------------------------------------------------------------
SVG_NS = 'http://www.w3.org/2000/svg'
XLINK_NS = 'http://www.w3.org/1999/xlink'
NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}
EMU_PX = 9525
SLIDE_W = 12192000
SLIDE_H = 6858000

# CSS 完整命名颜色表（常用子集）
CSS_COLORS = {
    'aliceblue': 'f0f8ff', 'antiquewhite': 'faebd7', 'aqua': '00ffff',
    'aquamarine': '7fffd4', 'azure': 'f0ffff', 'beige': 'f5f5dc',
    'bisque': 'ffe4c4', 'black': '000000', 'blanchedalmond': 'ffebcd',
    'blue': '0000ff', 'blueviolet': '8a2be2', 'brown': 'a52a2a',
    'burlywood': 'deb887', 'cadetblue': '5f9ea0', 'chartreuse': '7fff00',
    'chocolate': 'd2691e', 'coral': 'ff7f50', 'cornflowerblue': '6495ed',
    'cornsilk': 'fff8dc', 'crimson': 'dc143c', 'cyan': '00ffff',
    'darkblue': '00008b', 'darkcyan': '008b8b', 'darkgoldenrod': 'b8860b',
    'darkgray': 'a9a9a9', 'darkgreen': '006400', 'darkgrey': 'a9a9a9',
    'darkkhaki': 'bdb76b', 'darkmagenta': '8b008b', 'darkolivegreen': '556b2f',
    'darkorange': 'ff8c00', 'darkorchid': '9932cc', 'darkred': '8b0000',
    'darksalmon': 'e9967a', 'darkseagreen': '8fbc8f', 'darkslateblue': '483d8b',
    'darkslategray': '2f4f4f', 'darkturquoise': '00ced1', 'darkviolet': '9400d3',
    'deeppink': 'ff1493', 'deepskyblue': '00bfff', 'dimgray': '696969',
    'dodgerblue': '1e90ff', 'firebrick': 'b22222', 'floralwhite': 'fffaf0',
    'forestgreen': '228b22', 'fuchsia': 'ff00ff', 'gainsboro': 'dcdcdc',
    'ghostwhite': 'f8f8ff', 'gold': 'ffd700', 'goldenrod': 'daa520',
    'gray': '808080', 'green': '008000', 'greenyellow': 'adff2f',
    'grey': '808080', 'honeydew': 'f0fff0', 'hotpink': 'ff69b4',
    'indianred': 'cd5c5c', 'indigo': '4b0082', 'ivory': 'fffff0',
    'khaki': 'f0e68c', 'lavender': 'e6e6fa', 'lawngreen': '7cfc00',
    'lemonchiffon': 'fffacd', 'lightblue': 'add8e6', 'lightcoral': 'f08080',
    'lightcyan': 'e0ffff', 'lightgoldenrodyellow': 'fafad2', 'lightgray': 'd3d3d3',
    'lightgreen': '90ee90', 'lightpink': 'ffb6c1', 'lightsalmon': 'ffa07a',
    'lightseagreen': '20b2aa', 'lightskyblue': '87cefa', 'lightslategray': '778899',
    'lightsteelblue': 'b0c4de', 'lightyellow': 'ffffe0', 'lime': '00ff00',
    'limegreen': '32cd32', 'linen': 'faf0e6', 'magenta': 'ff00ff',
    'maroon': '800000', 'mediumaquamarine': '66cdaa', 'mediumblue': '0000cd',
    'mediumorchid': 'ba55d3', 'mediumpurple': '9370db', 'mediumseagreen': '3cb371',
    'mediumslateblue': '7b68ee', 'mediumspringgreen': '00fa9a',
    'mediumturquoise': '48d1cc', 'mediumvioletred': 'c71585', 'midnightblue': '191970',
    'mintcream': 'f5fffa', 'mistyrose': 'ffe4e1', 'moccasin': 'ffe4b5',
    'navajowhite': 'ffdead', 'navy': '000080', 'oldlace': 'fdf5e6',
    'olive': '808000', 'olivedrab': '6b8e23', 'orange': 'ffa500',
    'orangered': 'ff4500', 'orchid': 'da70d6', 'palegoldenrod': 'eee8aa',
    'palegreen': '98fb98', 'paleturquoise': 'afeeee', 'palevioletred': 'db7093',
    'papayawhip': 'ffefd5', 'peachpuff': 'ffdab9', 'peru': 'cd853f',
    'pink': 'ffc0cb', 'plum': 'dda0dd', 'powderblue': 'b0e0e6',
    'purple': '800080', 'rebeccapurple': '663399', 'red': 'ff0000',
    'rosybrown': 'bc8f8f', 'royalblue': '4169e1', 'saddlebrown': '8b4513',
    'salmon': 'fa8072', 'sandybrown': 'f4a460', 'seagreen': '2e8b57',
    'seashell': 'fff5ee', 'sienna': 'a0522d', 'silver': 'c0c0c0',
    'skyblue': '87ceeb', 'slateblue': '6a5acd', 'slategray': '708090',
    'snow': 'fffafa', 'springgreen': '00ff7f', 'steelblue': '4682b4',
    'tan': 'd2b48c', 'teal': '008080', 'thistle': 'd8bfd8',
    'tomato': 'ff6347', 'turquoise': '40e0d0', 'violet': 'ee82ee',
    'wheat': 'f5deb3', 'white': 'ffffff', 'whitesmoke': 'f5f5f5',
    'yellow': 'ffff00', 'yellowgreen': '9acd32',
}

# 字体回退链
FONT_FALLBACK = {
    'PingFang SC': 'Microsoft YaHei',
    'SF Pro Display': 'Arial',
    'Helvetica Neue': 'Arial',
    'Helvetica': 'Arial',
    'system-ui': 'Microsoft YaHei',
    'sans-serif': 'Microsoft YaHei',
}


def px(v):
    return int(float(v) * EMU_PX)

def font_sz(svg_px):
    return max(100, int(float(svg_px) * 75))

def strip_unit(v):
    return re.sub(r'[a-z%]+', '', str(v))

def resolve_font(ff_str):
    """解析 font-family 字符串，返回 PPT 可用字体。"""
    ff_str = ff_str.replace('&quot;', '').replace('"', '').replace("'", '')
    fonts = [f.strip() for f in ff_str.split(',') if f.strip()]
    for f in fonts:
        if f in FONT_FALLBACK:
            return FONT_FALLBACK[f]
        if f and f not in ('sans-serif', 'serif', 'monospace', 'system-ui'):
            return f
    return 'Microsoft YaHei'


# -------------------------------------------------------------------
# 颜色解析（完整 CSS 命名颜色）
# -------------------------------------------------------------------
def parse_color(s):
    if not s or s.strip() == 'none':
        return None
    s = s.strip()
    if s.startswith('url('):
        m = re.search(r'#([\w-]+)', s)
        return ('grad', m.group(1)) if m else None
    m = re.match(r'rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*(?:,\s*([\d.]+))?\s*\)', s)
    if m:
        r, g, b = int(m.group(1)), int(m.group(2)), int(m.group(3))
        a = float(m.group(4)) if m.group(4) else 1.0
        return (f'{r:02x}{g:02x}{b:02x}', int(a * 100000))
    if s.startswith('#'):
        h = s[1:]
        if len(h) == 3:
            h = h[0]*2 + h[1]*2 + h[2]*2
        return (h.lower().ljust(6, '0')[:6], 100000)
    c = CSS_COLORS.get(s.lower())
    return (c, 100000) if c else None


# -------------------------------------------------------------------
# OOXML 元素构造
# -------------------------------------------------------------------
def _el(tag, attrib=None, text=None, children=None):
    pre, local = tag.split(':') if ':' in tag else ('a', tag)
    el = etree.Element(f'{{{NS[pre]}}}{local}')
    if attrib:
        for k, v in attrib.items():
            el.set(k, str(v))
    if text is not None:
        el.text = str(text)
    for c in (children or []):
        if c is not None:
            el.append(c)
    return el

def _srgb(hex6, alpha=100000):
    el = _el('a:srgbClr', {'val': hex6})
    if alpha < 100000:
        el.append(_el('a:alpha', {'val': str(alpha)}))
    return el

def make_fill(fill_str, grads, opacity=1.0):
    c = parse_color(fill_str)
    if c is None:
        return _el('a:noFill')
    if c[0] == 'grad':
        gdef = grads.get(c[1])
        return _make_grad(gdef) if gdef else _el('a:noFill')
    hex6, alpha = c
    alpha = int(alpha * opacity)
    return _el('a:solidFill', children=[_srgb(hex6, alpha)])

def _make_grad(gdef):
    gs_lst = _el('a:gsLst')
    for stop in gdef['stops']:
        pos = int(stop['offset'] * 1000)
        sc = parse_color(stop['color_str'])
        if not sc or sc[0] == 'grad':
            continue
        hex6, alpha = sc
        alpha = int(alpha * stop.get('opacity', 1.0))
        gs_lst.append(_el('a:gs', {'pos': str(pos)}, children=[_srgb(hex6, alpha)]))

    if gdef.get('type') == 'radial':
        # 径向渐变
        path = _el('a:path', {'path': 'circle'}, children=[
            _el('a:fillToRect', {'l': '50000', 't': '50000', 'r': '50000', 'b': '50000'})
        ])
        return _el('a:gradFill', {'rotWithShape': '1'}, children=[gs_lst, path])
    else:
        # 线性渐变
        dx = gdef.get('x2', 1) - gdef.get('x1', 0)
        dy = gdef.get('y2', 1) - gdef.get('y1', 0)
        ang = int(math.degrees(math.atan2(dy, dx)) * 60000)
        if ang < 0:
            ang += 21600000
        lin = _el('a:lin', {'ang': str(ang), 'scaled': '0'})
        return _el('a:gradFill', children=[gs_lst, lin])

def make_line(stroke_str, stroke_w=1):
    c = parse_color(stroke_str)
    if not c or c[0] == 'grad':
        return None
    hex6, alpha = c
    w = max(1, int(float(strip_unit(stroke_w)) * 12700))
    return _el('a:ln', {'w': str(w)},
               children=[_el('a:solidFill', children=[_srgb(hex6, alpha)])])

def make_shape(sid, name, x, y, cx, cy, preset='rect',
               fill_el=None, line_el=None, rx=0, geom_el=None):
    sp = _el('p:sp')
    sp.append(_el('p:nvSpPr', children=[
        _el('p:cNvPr', {'id': str(sid), 'name': name}),
        _el('p:cNvSpPr'), _el('p:nvPr'),
    ]))
    sp_pr = _el('p:spPr')
    sp_pr.append(_el('a:xfrm', children=[
        _el('a:off', {'x': str(max(0, int(x))), 'y': str(max(0, int(y)))}),
        _el('a:ext', {'cx': str(max(0, int(cx))), 'cy': str(max(0, int(cy)))}),
    ]))
    if geom_el is not None:
        sp_pr.append(geom_el)
    else:
        geom = _el('a:prstGeom', {'prst': preset})
        av = _el('a:avLst')
        if preset == 'roundRect' and rx > 0:
            shorter = max(min(cx, cy), 1)
            adj = min(50000, int(rx / (shorter / 2) * 50000))
            av.append(_el('a:gd', {'name': 'adj', 'fmla': f'val {adj}'}))
        geom.append(av)
        sp_pr.append(geom)
    sp_pr.append(fill_el if fill_el is not None else _el('a:noFill'))
    if line_el is not None:
        sp_pr.append(line_el)
    sp.append(sp_pr)
    return sp

def make_textbox(sid, name, x, y, cx, cy, paragraphs):
    """paragraphs = [[{text,sz,bold,hex,alpha,font}, ...], ...]"""
    sp = _el('p:sp')
    sp.append(_el('p:nvSpPr', children=[
        _el('p:cNvPr', {'id': str(sid), 'name': name}),
        _el('p:cNvSpPr', {'txBox': '1'}), _el('p:nvPr'),
    ]))
    sp.append(_el('p:spPr', children=[
        _el('a:xfrm', children=[
            _el('a:off', {'x': str(max(0, int(x))), 'y': str(max(0, int(y)))}),
            _el('a:ext', {'cx': str(max(0, int(cx))), 'cy': str(max(0, int(cy)))}),
        ]),
        _el('a:prstGeom', {'prst': 'rect'}, children=[_el('a:avLst')]),
        _el('a:noFill'), _el('a:ln', children=[_el('a:noFill')]),
    ]))
    tx = _el('p:txBody', children=[
        _el('a:bodyPr', {'wrap': 'none', 'lIns': '0', 'tIns': '0',
                         'rIns': '0', 'bIns': '0', 'anchor': 't'}),
        _el('a:lstStyle'),
    ])
    for runs in paragraphs:
        p_el = _el('a:p')
        for run in runs:
            rpr_a = {'lang': 'zh-CN', 'dirty': '0'}
            if run.get('sz'):
                rpr_a['sz'] = str(run['sz'])
            if run.get('bold'):
                rpr_a['b'] = '1'
            rpr = _el('a:rPr', rpr_a)
            rpr.append(_el('a:solidFill', children=[
                _srgb(run.get('hex', '000000'), run.get('alpha', 100000))
            ]))
            font = run.get('font', 'Microsoft YaHei')
            rpr.append(_el('a:latin', {'typeface': font}))
            rpr.append(_el('a:ea', {'typeface': font}))
            p_el.append(_el('a:r', children=[rpr, _el('a:t', text=run.get('text', ''))]))
        tx.append(p_el)
    sp.append(tx)
    return sp


# -------------------------------------------------------------------
# SVG Path 解析器 -> OOXML custGeom
# -------------------------------------------------------------------
_PATH_RE = re.compile(r'([mMzZlLhHvVcCsSqQtTaA])|([+-]?(?:\d+\.?\d*|\.\d+)(?:[eE][+-]?\d+)?)')

def parse_path_to_custgeom(d_str, bbox):
    """SVG path d -> OOXML a:custGeom 元素。bbox=(x,y,w,h) 用于坐标偏移。"""
    bx, by, bw, bh = bbox
    scale = 100000  # OOXML 路径坐标空间

    def coord(v, is_x=True):
        base = bw if is_x else bh
        offset = bx if is_x else by
        if base <= 0:
            return 0
        return int((float(v) - offset) / base * scale)

    tokens = _PATH_RE.findall(d_str)
    items = []
    for cmd_match, num_match in tokens:
        if cmd_match:
            items.append(cmd_match)
        elif num_match:
            items.append(float(num_match))

    path_el = _el('a:path', {'w': str(scale), 'h': str(scale)})
    i = 0
    cx_p, cy_p = 0, 0  # current point (absolute)
    cmd = None
    rel = False

    while i < len(items):
        if isinstance(items[i], str):
            cmd = items[i].lower()
            rel = items[i].islower()
            i += 1
            if cmd == 'z':
                path_el.append(_el('a:close'))
                continue

        if cmd is None:
            i += 1
            continue

        try:
            if cmd == 'm':
                x, y = float(items[i]), float(items[i+1])
                if rel:
                    x += cx_p; y += cy_p
                cx_p, cy_p = x, y
                path_el.append(_el('a:moveTo', children=[
                    _el('a:pt', {'x': str(coord(x, True)), 'y': str(coord(y, False))})
                ]))
                i += 2
                cmd = 'l'  # implicit lineTo after moveTo

            elif cmd == 'l':
                x, y = float(items[i]), float(items[i+1])
                if rel:
                    x += cx_p; y += cy_p
                cx_p, cy_p = x, y
                path_el.append(_el('a:lnTo', children=[
                    _el('a:pt', {'x': str(coord(x, True)), 'y': str(coord(y, False))})
                ]))
                i += 2

            elif cmd == 'h':
                x = float(items[i])
                if rel:
                    x += cx_p
                cx_p = x
                path_el.append(_el('a:lnTo', children=[
                    _el('a:pt', {'x': str(coord(cx_p, True)), 'y': str(coord(cy_p, False))})
                ]))
                i += 1

            elif cmd == 'v':
                y = float(items[i])
                if rel:
                    y += cy_p
                cy_p = y
                path_el.append(_el('a:lnTo', children=[
                    _el('a:pt', {'x': str(coord(cx_p, True)), 'y': str(coord(cy_p, False))})
                ]))
                i += 1

            elif cmd == 'c':
                x1, y1 = float(items[i]), float(items[i+1])
                x2, y2 = float(items[i+2]), float(items[i+3])
                x, y = float(items[i+4]), float(items[i+5])
                if rel:
                    x1 += cx_p; y1 += cy_p
                    x2 += cx_p; y2 += cy_p
                    x += cx_p; y += cy_p
                cx_p, cy_p = x, y
                path_el.append(_el('a:cubicBezTo', children=[
                    _el('a:pt', {'x': str(coord(x1, True)), 'y': str(coord(y1, False))}),
                    _el('a:pt', {'x': str(coord(x2, True)), 'y': str(coord(y2, False))}),
                    _el('a:pt', {'x': str(coord(x, True)), 'y': str(coord(y, False))}),
                ]))
                i += 6

            elif cmd in ('s', 'q', 't', 'a'):
                # 简化处理：跳过复杂曲线
                skip = {'s': 4, 'q': 4, 't': 2, 'a': 7}.get(cmd, 2)
                i += skip
            else:
                i += 1
        except (IndexError, ValueError):
            i += 1

    cust_geom = _el('a:custGeom', children=[
        _el('a:avLst'), _el('a:gdLst'), _el('a:ahLst'), _el('a:cxnLst'),
        _el('a:rect', {'l': 'l', 't': 't', 'r': 'r', 'b': 'b'}),
        _el('a:pathLst', children=[path_el]),
    ])
    return cust_geom


# -------------------------------------------------------------------
# SVG -> PPTX 转换器
# -------------------------------------------------------------------
class SvgConverter:
    def __init__(self, on_progress=None):
        self.sid = 100
        self.grads = {}
        self.bg_set = False  # 是否已设置幻灯片背景
        self.on_progress = on_progress  # 进度回调 (i, total, filename)
        self.stats = {'shapes': 0, 'skipped': 0, 'errors': 0}

    def _id(self):
        self.sid += 1
        return self.sid

    def convert(self, svg_path, slide):
        self.bg_set = False
        self.stats = {'shapes': 0, 'skipped': 0, 'errors': 0}
        tree = etree.parse(str(svg_path))
        root = tree.getroot()
        self._parse_grads(root)
        sp_tree = None
        for d in slide._element.iter():
            if d.tag.endswith('}spTree'):
                sp_tree = d
                break
        if sp_tree is None:
            return
        self._walk(root, sp_tree, 0, 0, 1.0, 1.0, slide)

    def _parse_grads(self, root):
        self.grads = {}
        pct = lambda v: float(v.rstrip('%')) / 100 if '%' in str(v) else float(v)
        for g in root.iter(f'{{{SVG_NS}}}linearGradient'):
            gid = g.get('id')
            if not gid:
                continue
            stops = []
            for s in g.findall(f'{{{SVG_NS}}}stop'):
                off = s.get('offset', '0%')
                off = float(off.rstrip('%')) if '%' in off else float(off) * 100
                stops.append({'offset': off, 'color_str': s.get('stop-color', '#000'),
                              'opacity': float(s.get('stop-opacity', '1'))})
            self.grads[gid] = {
                'type': 'linear', 'stops': stops,
                'x1': pct(g.get('x1', '0%')), 'y1': pct(g.get('y1', '0%')),
                'x2': pct(g.get('x2', '100%')), 'y2': pct(g.get('y2', '100%')),
            }
        for g in root.iter(f'{{{SVG_NS}}}radialGradient'):
            gid = g.get('id')
            if not gid:
                continue
            stops = []
            for s in g.findall(f'{{{SVG_NS}}}stop'):
                off = s.get('offset', '0%')
                off = float(off.rstrip('%')) if '%' in off else float(off) * 100
                stops.append({'offset': off, 'color_str': s.get('stop-color', '#000'),
                              'opacity': float(s.get('stop-opacity', '1'))})
            self.grads[gid] = {'type': 'radial', 'stops': stops}

    def _tag(self, el):
        t = el.tag
        return t.split('}')[1] if isinstance(t, str) and '}' in t else (t if isinstance(t, str) else '')

    def _parse_transform(self, el):
        """解析 transform -> (dx, dy, sx, sy)。"""
        t = el.get('transform', '')
        dx, dy, sx, sy = 0.0, 0.0, 1.0, 1.0
        # translate
        m = re.search(r'translate\(\s*([\d.\-]+)[,\s]+([\d.\-]+)', t)
        if m:
            dx, dy = float(m.group(1)), float(m.group(2))
        # scale
        m = re.search(r'scale\(\s*([\d.\-]+)(?:[,\s]+([\d.\-]+))?\s*\)', t)
        if m:
            sx = float(m.group(1))
            sy = float(m.group(2)) if m.group(2) else sx
        # matrix(a,b,c,d,e,f) -> e=translateX, f=translateY
        m = re.search(r'matrix\(\s*([\d.\-]+)[,\s]+([\d.\-]+)[,\s]+([\d.\-]+)[,\s]+([\d.\-]+)[,\s]+([\d.\-]+)[,\s]+([\d.\-]+)', t)
        if m:
            dx = float(m.group(5))
            dy = float(m.group(6))
            sx = float(m.group(1))
            sy = float(m.group(4))
        return dx, dy, sx, sy

    def _walk(self, el, sp, ox, oy, group_opacity, scale, slide):
        tag = self._tag(el)
        try:
            if tag == 'rect':
                self._rect(el, sp, ox, oy, group_opacity, scale, slide)
            elif tag == 'text':
                self._text(el, sp, ox, oy, group_opacity, scale)
            elif tag == 'circle':
                self._circle(el, sp, ox, oy, group_opacity, scale)
            elif tag == 'ellipse':
                self._ellipse(el, sp, ox, oy, group_opacity, scale)
            elif tag == 'line':
                self._line(el, sp, ox, oy, scale)
            elif tag == 'path':
                self._path(el, sp, ox, oy, group_opacity, scale)
            elif tag == 'image':
                self._image(el, sp, ox, oy, group_opacity, scale, slide)
            elif tag == 'g':
                dx, dy, sx, sy = self._parse_transform(el)
                el_opacity = float(el.get('opacity', '1'))
                child_opacity = group_opacity * el_opacity
                # scale 累积：父级scale * 当前g的scale
                child_scale = scale * sx  # 假设sx==sy（等比缩放）
                new_ox = ox + dx * scale
                new_oy = oy + dy * scale
                for c in el:
                    self._walk(c, sp, new_ox, new_oy,
                               child_opacity, child_scale, slide)
            elif tag in ('defs', 'style', 'linearGradient', 'radialGradient',
                         'stop', 'pattern', 'clipPath', 'filter', 'mask'):
                pass
            else:
                for c in el:
                    self._walk(c, sp, ox, oy, group_opacity, scale, slide)
        except Exception as e:
            self.stats['errors'] += 1
            print(f"    Warning: {tag} element failed: {e}", file=sys.stderr)

    def _rect(self, el, sp, ox, oy, opacity, scale, slide):
        x = (float(el.get('x', 0)) * scale) + ox
        y = (float(el.get('y', 0)) * scale) + oy
        w = float(el.get('width', 0)) * scale
        h = float(el.get('height', 0)) * scale
        if w <= 0 or h <= 0:
            return

        # 过滤面积 < 4px 的纯装饰元素
        if w < 4 and h < 4:
            self.stats['skipped'] += 1
            return

        fill_s = el.get('fill', '')
        stroke_s = el.get('stroke', '')
        c = parse_color(fill_s)

        # 跳过全透明无边框矩形
        if c and c[0] != 'grad' and c[1] == 0 and not stroke_s:
            return

        el_opacity = float(el.get('opacity', '1')) * opacity

        # 首个全屏 rect -> 幻灯片背景
        if not self.bg_set and w >= 1270 and h >= 710:
            self.bg_set = True
            bg = slide._element.find(f'.//{{{NS["p"]}}}bg')
            if bg is None:
                cSld = slide._element.find(f'{{{NS["p"]}}}cSld')
                if cSld is not None:
                    bg_el = _el('p:bg', children=[
                        _el('p:bgPr', children=[
                            make_fill(fill_s, self.grads, el_opacity),
                            _el('a:effectLst'),
                        ])
                    ])
                    cSld.insert(0, bg_el)
            return  # 不再作为形状添加

        r = max(float(el.get('rx', 0)), float(el.get('ry', 0)))
        preset = 'roundRect' if r > 0 else 'rect'
        fill_el = make_fill(fill_s, self.grads, el_opacity)
        line_el = make_line(stroke_s, el.get('stroke-width', '1')) if stroke_s else None
        shape = make_shape(self._id(), f'R{self.sid}',
                           px(x), px(y), px(w), px(h),
                           preset=preset, fill_el=fill_el, line_el=line_el, rx=px(r))
        sp.append(shape)
        self.stats['shapes'] += 1

    def _text(self, el, sp, ox, oy, opacity, scale):
        """每个 tspan 保持独立文本框，保留精确 x/y 坐标。"""
        fill_s = el.get('fill', el.get('color', ''))
        fsz = el.get('font-size', '14px').replace('px', '')
        fw = el.get('font-weight', '')
        ff = el.get('font-family', '')
        baseline = el.get('dominant-baseline', '')
        anchor = el.get('text-anchor', 'start')

        tspans = list(el.findall(f'{{{SVG_NS}}}tspan'))

        if tspans:
            for ts in tspans:
                txt = ts.text
                if not txt or not txt.strip():
                    continue
                x = float(ts.get('x', 0)) * scale + ox
                y = float(ts.get('y', 0)) * scale + oy
                tlen = float(ts.get('textLength', 0))
                ts_fsz = ts.get('font-size', fsz).replace('px', '')
                ts_fw = ts.get('font-weight', fw)
                ts_fill = ts.get('fill', fill_s)
                ts_ff = ts.get('font-family', ff)
                fh = float(ts_fsz)
                # baseline偏移: text-after-edge -> y是底部减全高; auto -> y是baseline减85%
                if 'after-edge' in baseline:
                    y -= fh
                else:
                    y -= fh * 0.85
                c = parse_color(ts_fill)
                hex6 = c[0] if c and c[0] != 'grad' else '000000'
                alpha = c[1] if c and c[0] != 'grad' else 100000
                alpha = int(alpha * opacity)
                cx_v = px(tlen) if tlen > 0 else px(len(txt) * float(ts_fsz) * 0.7)
                cy_v = px(fh * 1.5)
                # text-anchor 偏移: middle -> x减半宽, end -> x减全宽
                if anchor == 'middle':
                    x -= cx_v / EMU_PX / 2
                elif anchor == 'end':
                    x -= cx_v / EMU_PX
                run = {
                    'text': txt.strip(), 'sz': font_sz(ts_fsz),
                    'bold': ts_fw in ('bold', '700', '800', '900'),
                    'hex': hex6, 'alpha': alpha,
                    'font': resolve_font(ts_ff),
                }
                shape = make_textbox(self._id(), f'T{self.sid}',
                                     px(x), px(y), cx_v, cy_v, [[run]])
                sp.append(shape)
                self.stats['shapes'] += 1

        elif el.text and el.text.strip():
            x = float(el.get('x', 0)) * scale + ox
            y = float(el.get('y', 0)) * scale + oy
            fh = float(fsz)
            # baseline偏移
            if 'after-edge' in baseline:
                y -= fh
            else:
                y -= fh * 0.85
            c = parse_color(fill_s)
            hex6 = c[0] if c and c[0] != 'grad' else '000000'
            alpha = c[1] if c and c[0] != 'grad' else 100000
            alpha = int(alpha * opacity)
            txt = el.text.strip()
            txt_w = len(txt) * float(fsz) * 0.7
            # text-anchor 偏移
            if anchor == 'middle':
                x -= txt_w / 2
            elif anchor == 'end':
                x -= txt_w
            run = {
                'text': txt, 'sz': font_sz(fsz),
                'bold': fw in ('bold', '700', '800', '900'),
                'hex': hex6, 'alpha': alpha, 'font': resolve_font(ff),
            }
            shape = make_textbox(self._id(), f'T{self.sid}',
                                 px(x), px(y),
                                 px(len(txt) * float(fsz) * 0.7),
                                 px(fh * 1.5), [[run]])
            sp.append(shape)
            self.stats['shapes'] += 1

    def _circle(self, el, sp, ox, oy, opacity, scale):
        cx_v = float(el.get('cx', 0)) * scale + ox
        cy_v = float(el.get('cy', 0)) * scale + oy
        r = float(el.get('r', 0)) * scale
        if r <= 0 or r < 2:
            self.stats['skipped'] += 1
            return

        el_opacity = float(el.get('opacity', '1')) * opacity
        fill_s = el.get('fill', '')
        stroke_s = el.get('stroke', '')
        stroke_w_s = el.get('stroke-width', '1')
        dasharray = el.get('stroke-dasharray', '')

        # 环形图特殊处理：fill=none + stroke + dasharray -> OOXML arc + 粗描边
        if (fill_s == 'none' or not fill_s) and stroke_s and dasharray:
            sw = float(strip_unit(stroke_w_s))
            # 解析 dasharray (格式: "188.1 188.5" 或 "113.097px, 150.796px")
            dash_parts = [float(strip_unit(p.strip())) for p in dasharray.replace(',', ' ').split() if p.strip()]
            if len(dash_parts) >= 2:
                circumference = 2 * math.pi * r
                arc_len = dash_parts[0]
                angle_pct = min(arc_len / circumference, 1.0)

                # 检查 rotate transform
                transform = el.get('transform', '')
                start_angle = 0
                rot_m = re.search(r'rotate\(\s*([\d.\-]+)', transform)
                if rot_m:
                    start_angle = float(rot_m.group(1))

                # SVG -> PowerPoint 角度转换
                # SVG rotate(-90) = 从 12 点钟方向开始
                # PowerPoint arc: adj1=startAngle, adj2=endAngle (从3点钟顺时针, 60000单位/度)
                ppt_start = (start_angle + 90) % 360
                sweep = angle_pct * 360
                ppt_end = (ppt_start + sweep) % 360

                adj1 = int(ppt_start * 60000)
                adj2 = int(ppt_end * 60000)

                # 用 arc 预设 (只画弧线轮廓) + 粗描边 = 环形弧
                geom = _el('a:prstGeom', {'prst': 'arc'})
                av = _el('a:avLst')
                av.append(_el('a:gd', {'name': 'adj1', 'fmla': f'val {adj1}'}))
                av.append(_el('a:gd', {'name': 'adj2', 'fmla': f'val {adj2}'}))
                geom.append(av)

                # 描边颜色 = SVG 的 stroke 颜色（支持渐变引用）
                stroke_color = parse_color(stroke_s)
                ln_children = []
                if stroke_color and stroke_color[0] == 'grad':
                    # stroke 引用渐变 -> 提取渐变的第一个 stop 颜色作为实色
                    gdef = self.grads.get(stroke_color[1])
                    if gdef and gdef.get('stops'):
                        first_stop = gdef['stops'][0]
                        sc = parse_color(first_stop['color_str'])
                        if sc and sc[0] != 'grad':
                            ln_children.append(_el('a:solidFill', children=[
                                _srgb(sc[0], int(sc[1] * el_opacity))
                            ]))
                    # 也尝试用渐变填充（OOXML线条支持渐变）
                    if not ln_children and gdef:
                        grad_fill = _make_grad(gdef)
                        if grad_fill is not None:
                            ln_children.append(grad_fill)
                elif stroke_color and stroke_color[0] != 'grad':
                    ln_children.append(_el('a:solidFill', children=[
                        _srgb(stroke_color[0], int(stroke_color[1] * el_opacity))
                    ]))
                ln_children.append(_el('a:round'))
                line_el = _el('a:ln', {'w': str(int(sw * 12700))}, children=ln_children)

                shape = _el('p:sp')
                shape.append(_el('p:nvSpPr', children=[
                    _el('p:cNvPr', {'id': str(self._id()), 'name': f'Arc{self.sid}'}),
                    _el('p:cNvSpPr'), _el('p:nvPr'),
                ]))
                sp_pr = _el('p:spPr')
                sp_pr.append(_el('a:xfrm', children=[
                    _el('a:off', {'x': str(max(0, px(cx_v - r))),
                                  'y': str(max(0, px(cy_v - r)))}),
                    _el('a:ext', {'cx': str(px(2 * r)),
                                  'cy': str(px(2 * r))}),
                ]))
                sp_pr.append(geom)
                sp_pr.append(_el('a:noFill'))
                sp_pr.append(line_el)
                shape.append(sp_pr)
                sp.append(shape)
                self.stats['shapes'] += 1
                return

        # fill=none + stroke (无dasharray) -> 空心圆 + 粗描边
        if (fill_s == 'none' or not fill_s) and stroke_s and stroke_s != 'none':
            sw = float(strip_unit(stroke_w_s))
            stroke_color = parse_color(stroke_s)
            ln_children = []
            if stroke_color and stroke_color[0] != 'grad':
                ln_children.append(_el('a:solidFill', children=[
                    _srgb(stroke_color[0], int(stroke_color[1] * el_opacity))
                ]))
            ln_children.append(_el('a:round'))
            line_el = _el('a:ln', {'w': str(int(sw * 12700))}, children=ln_children)

            sp.append(make_shape(self._id(), f'C{self.sid}',
                                 px(cx_v - r), px(cy_v - r), px(2*r), px(2*r),
                                 preset='ellipse',
                                 fill_el=_el('a:noFill'),
                                 line_el=line_el))
            self.stats['shapes'] += 1
            return

        # 普通圆形
        fill_el = make_fill(fill_s, self.grads, el_opacity)
        line_el = make_line(stroke_s, stroke_w_s) if stroke_s and stroke_s != 'none' else None
        sp.append(make_shape(self._id(), f'C{self.sid}',
                             px(cx_v - r), px(cy_v - r), px(2*r), px(2*r),
                             preset='ellipse', fill_el=fill_el, line_el=line_el))
        self.stats['shapes'] += 1

    def _ellipse(self, el, sp, ox, oy, opacity, scale):
        cx_v = float(el.get('cx', 0)) * scale + ox
        cy_v = float(el.get('cy', 0)) * scale + oy
        rx = float(el.get('rx', 0)) * scale
        ry = float(el.get('ry', 0)) * scale
        if rx <= 0 or ry <= 0:
            return
        el_opacity = float(el.get('opacity', '1')) * opacity
        fill_el = make_fill(el.get('fill', ''), self.grads, el_opacity)
        sp.append(make_shape(self._id(), f'E{self.sid}',
                             px(cx_v - rx), px(cy_v - ry), px(2*rx), px(2*ry),
                             preset='ellipse', fill_el=fill_el))
        self.stats['shapes'] += 1

    def _line(self, el, sp, ox, oy, scale):
        x1 = float(el.get('x1', 0)) * scale + ox
        y1 = float(el.get('y1', 0)) * scale + oy
        x2 = float(el.get('x2', 0)) * scale + ox
        y2 = float(el.get('y2', 0)) * scale + oy
        line_el = make_line(el.get('stroke', '#000'), el.get('stroke-width', '1'))
        if line_el is None:
            return
        mx, my = min(x1, x2), min(y1, y2)
        w, h = abs(x2 - x1) or 1, abs(y2 - y1) or 1
        shape = make_shape(self._id(), f'L{self.sid}',
                           px(mx), px(my), px(w), px(h),
                           preset='line', fill_el=_el('a:noFill'), line_el=line_el)
        xfrm = shape.find(f'.//{{{NS["a"]}}}xfrm')
        if x1 > x2:
            xfrm.set('flipH', '1')
        if y1 > y2:
            xfrm.set('flipV', '1')
        sp.append(shape)
        self.stats['shapes'] += 1

    def _path(self, el, sp, ox, oy, opacity, scale):
        """SVG <path> -> OOXML custGeom 形状。"""
        d = el.get('d', '')
        if not d or 'nan' in d:
            return
        # 计算 bounding box（简化：从 path 数据提取所有数字坐标）
        nums = re.findall(r'[+-]?(?:\d+\.?\d*|\.\d+)', d)
        if len(nums) < 4:
            return
        coords = [float(n) for n in nums]
        xs = coords[0::2]
        ys = coords[1::2] if len(coords) > 1 else [0]
        bx, by = min(xs), min(ys)
        bw = max(xs) - bx or 1
        bh = max(ys) - by or 1

        # 过滤极小路径
        if bw < 4 and bh < 4:
            self.stats['skipped'] += 1
            return

        geom_el = parse_path_to_custgeom(d, (bx, by, bw, bh))
        el_opacity = float(el.get('opacity', '1')) * opacity
        fill_el = make_fill(el.get('fill', ''), self.grads, el_opacity)
        line_el = make_line(el.get('stroke', ''), el.get('stroke-width', '1')) if el.get('stroke') else None

        shape = make_shape(self._id(), f'P{self.sid}',
                           px((bx + ox) * scale) if scale != 1.0 else px(bx + ox),
                           px((by + oy) * scale) if scale != 1.0 else px(by + oy),
                           px(bw * scale), px(bh * scale),
                           fill_el=fill_el, line_el=line_el, geom_el=geom_el)
        sp.append(shape)
        self.stats['shapes'] += 1

    def _image(self, el, sp, ox, oy, opacity, scale, slide):
        href = el.get(f'{{{XLINK_NS}}}href') or el.get('href', '')
        x = float(el.get('x', 0)) * scale + ox
        y = float(el.get('y', 0)) * scale + oy
        w = float(el.get('width', 0)) * scale
        h = float(el.get('height', 0)) * scale
        el_opacity = float(el.get('opacity', '1')) * opacity
        if not href or w <= 0 or h <= 0:
            return

        img_source = None
        if href.startswith('data:'):
            m = re.match(r'data:image/\w+;base64,(.*)', href, re.DOTALL)
            if m:
                img_source = io.BytesIO(base64.b64decode(m.group(1)))
        elif href.startswith('file://'):
            p = Path(href.replace('file://', ''))
            if p.exists():
                img_source = str(p)
        elif not href.startswith('http'):
            p = Path(href)
            if p.exists():
                img_source = str(p)

        if img_source is None:
            return

        # 获取图片原始尺寸以计算宽高比
        try:
            from PIL import Image as PILImage
            if isinstance(img_source, io.BytesIO):
                img_source.seek(0)
                pil_img = PILImage.open(img_source)
                img_w, img_h = pil_img.size
                # 不 close -- PIL close 会关掉底层 BytesIO
                del pil_img
                img_source.seek(0)
            else:
                with PILImage.open(img_source) as pil_img:
                    img_w, img_h = pil_img.size
        except ImportError:
            # 没有 PIL，退回直接拉伸
            pic = slide.shapes.add_picture(img_source,
                                           Emu(px(x)), Emu(px(y)),
                                           Emu(px(w)), Emu(px(h)))
            self.stats['shapes'] += 1
            return

        # object-fit: cover -- 按比例放大到覆盖容器，然后裁剪
        container_w = px(w)
        container_h = px(h)
        img_ratio = img_w / img_h
        container_ratio = container_w / container_h

        if img_ratio > container_ratio:
            # 图片更宽 -> 按高度填满，裁剪左右
            scale_h = container_h
            scale_w = int(scale_h * img_ratio)
        else:
            # 图片更高 -> 按宽度填满，裁剪上下
            scale_w = container_w
            scale_h = int(scale_w / img_ratio)

        # 放置缩放后的图片（居中裁剪）
        offset_x = (scale_w - container_w) / 2
        offset_y = (scale_h - container_h) / 2

        pic = slide.shapes.add_picture(img_source,
                                       Emu(px(x)), Emu(px(y)),
                                       Emu(scale_w), Emu(scale_h))

        # 用 crop 实现裁剪（值为比例 0.0-1.0）
        if scale_w > 0 and scale_h > 0:
            crop_lr = offset_x / scale_w  # 左右各裁多少比例
            crop_tb = offset_y / scale_h  # 上下各裁多少比例
            pic.crop_left = crop_lr
            pic.crop_right = crop_lr
            pic.crop_top = crop_tb
            pic.crop_bottom = crop_tb

        # 应用透明度（通过 OOXML alphaModFix）
        if el_opacity < 0.99:
            from pptx.oxml.ns import qn
            sp_pr = pic._element.find(qn('p:spPr'))
            if sp_pr is None:
                sp_pr = pic._element.find(qn('pic:spPr'))
            # 在 blipFill 的 blip 上设置 alphaModFix
            blip = pic._element.find('.//' + qn('a:blip'))
            if blip is not None:
                alpha_val = int(el_opacity * 100000)
                alpha_el = _el('a:alphaModFix', {'amt': str(alpha_val)})
                blip.append(alpha_el)

        self.stats['shapes'] += 1


# -------------------------------------------------------------------
# 主流程
# -------------------------------------------------------------------
def convert(svg_input, output_path, on_progress=None):
    svg_input = Path(svg_input)
    if svg_input.is_file():
        svg_files = [svg_input]
    elif svg_input.is_dir():
        svg_files = sorted(svg_input.glob('*.svg'))
    else:
        print(f"Error: {svg_input} not found", file=sys.stderr)
        sys.exit(1)

    if not svg_files:
        print("Error: No SVG files found", file=sys.stderr)
        sys.exit(1)

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    blank = prs.slide_layouts[6]
    converter = SvgConverter(on_progress=on_progress)
    total = len(svg_files)

    for i, svg_file in enumerate(svg_files):
        slide = prs.slides.add_slide(blank)
        converter.convert(svg_file, slide)
        s = converter.stats
        print(f"  [{i+1}/{total}] {svg_file.name} "
              f"({s['shapes']} shapes, {s['skipped']} skipped, {s['errors']} errors)")
        if on_progress:
            on_progress(i + 1, total, svg_file.name)

    prs.save(str(output_path))
    print(f"Saved: {output_path} ({total} slides)")


def main():
    parser = argparse.ArgumentParser(description="SVG to PPTX (native shapes)")
    parser.add_argument('svg', help='SVG file or directory')
    parser.add_argument('-o', '--output', default='presentation.pptx')
    parser.add_argument('--html-dir', default=None,
                        help='HTML source directory (for future notes extraction)')
    args = parser.parse_args()
    convert(args.svg, args.output)


if __name__ == '__main__':
    main()
