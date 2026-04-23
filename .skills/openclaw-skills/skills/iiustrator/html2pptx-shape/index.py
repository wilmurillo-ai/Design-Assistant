#!/usr/bin/env python3
"""
html2pptx-shape — HTML 转 PPTX 形状转换器

将 HTML 幻灯片转换为完全可编辑的 PPTX，保留 CSS 样式、布局，
映射为 PPTX 原生形状（矩形/文本框/图片等）。
"""

import sys
import os
import re
from pathlib import Path
from bs4 import BeautifulSoup
import cssutils
import premailer
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from PIL import Image
import io
import base64
import requests



class CSSStyleParser:
    """解析 CSS 样式并转换为 PPTX 可用的属性"""
    
    def __init__(self):
        self.css_vars = {}
    
    def parse_css_vars(self, css_text):
        """提取 CSS 变量定义"""
        # 匹配 :root { --var-name: value; } 或任何选择器中的变量定义
        # 包括 .tpl-xhs-white-editorial { --xw-bg: #fff; }
        for block_match in re.finditer(r'[^{]+\{([^}]+)\}', css_text):
            block_content = block_match.group(1)
            for var_match in re.finditer(r'--([\w-]+)\s*:\s*([^;]+);', block_content):
                var_name = f'--{var_match.group(1)}'
                var_value = var_match.group(2).strip()
                self.css_vars[var_name] = var_value
    
    def resolve_var(self, value):
        """解析 CSS 变量引用"""
        if not value:
            return value
        var_match = re.match(r'var\((--[\w-]+)\)', value)
        if var_match:
            var_name = var_match.group(1)
            return self.css_vars.get(var_name, value)
        return value
    
    def parse_color(self, color_str):
        """解析颜色字符串为 RGB 元组"""
        if not color_str:
            return None
        
        color_str = self.resolve_var(color_str).strip()
        
        # 十六进制颜色 #RRGGBB 或 #RGB
        if color_str.startswith('#'):
            hex_color = color_str[1:]
            if len(hex_color) == 3:
                hex_color = ''.join([c*2 for c in hex_color])
            if len(hex_color) == 6:
                return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        
        # rgb(r, g, b)
        rgb_match = re.match(r'rgb\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)', color_str)
        if rgb_match:
            return tuple(int(rgb_match.group(i)) for i in (1, 2, 3))
        
        # rgba(r, g, b, a)
        rgba_match = re.match(r'rgba\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*,\s*([\d.]+)\s*\)', color_str)
        if rgba_match:
            return tuple(int(rgba_match.group(i)) for i in (1, 2, 3))
        
        # 颜色名称映射
        color_names = {
            'black': (0, 0, 0),
            'white': (255, 255, 255),
            'red': (255, 0, 0),
            'green': (0, 128, 0),
            'blue': (0, 0, 255),
            'yellow': (255, 255, 0),
            'gray': (128, 128, 128),
            'grey': (128, 128, 128),
        }
        return color_names.get(color_str.lower())
    
    def parse_gradient(self, gradient_str):
        """解析线性渐变"""
        if not gradient_str or 'gradient' not in gradient_str:
            return None
        
        # 简化处理：提取渐变色标
        colors = re.findall(r'#[0-9a-fA-F]{3,6}|rgb\([^)]+\)', gradient_str)
        if len(colors) >= 2:
            return [self.parse_color(c) for c in colors[:2]]
        return None
    
    def parse_font_size(self, size_str):
        """解析字体大小为 Points"""
        if not size_str:
            return Pt(12)
        
        size_str = self.resolve_var(size_str).strip()
        
        # px to pt (approximate: 1px = 0.75pt)
        px_match = re.match(r'(\d+(?:\.\d+)?)\s*px', size_str)
        if px_match:
            return Pt(float(px_match.group(1)) * 0.75)
        
        # pt
        pt_match = re.match(r'(\d+(?:\.\d+)?)\s*pt', size_str)
        if pt_match:
            return Pt(float(pt_match.group(1)))
        
        # em (assume base 16px)
        em_match = re.match(r'(\d+(?:\.\d+)?)\s*em', size_str)
        if em_match:
            return Pt(float(em_match.group(1)) * 12)
        
        return Pt(12)
    
    def parse_size(self, size_str, base_dpi=96):
        """解析尺寸为 Inches"""
        if not size_str:
            return Inches(0)
        
        size_str = self.resolve_var(size_str).strip()
        
        # px
        px_match = re.match(r'(\d+(?:\.\d+)?)\s*px', size_str)
        if px_match:
            return Inches(float(px_match.group(1)) / base_dpi)
        
        # em
        em_match = re.match(r'(\d+(?:\.\d+)?)\s*em', size_str)
        if em_match:
            return Inches(float(em_match.group(1)) * 0.625)
        
        # % (viewport based)
        pct_match = re.match(r'(\d+(?:\.\d+)?)\s*%', size_str)
        if pct_match:
            return Inches(float(pct_match.group(1)) / 100 * 10)
        
        return Inches(0)
    
    def parse_border_radius(self, radius_str):
        """解析圆角半径"""
        if not radius_str:
            return 0
        
        radius_str = self.resolve_var(radius_str).strip()
        px_match = re.match(r'(\d+(?:\.\d+)?)\s*px', radius_str)
        if px_match:
            return float(px_match.group(1))
        return 0
    
    def parse_text_align(self, align_str):
        """解析文本对齐方式"""
        if not align_str:
            return PP_ALIGN.LEFT
        
        align_str = align_str.strip().lower()
        mapping = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY,
        }
        return mapping.get(align_str, PP_ALIGN.LEFT)


class HTMLToPPTXConverter:
    """HTML 到 PPTX 转换器"""
    
    def __init__(self):
        self.css_parser = CSSStyleParser()
        self.slide_width = Inches(13.333)  # 16:9
        self.slide_height = Inches(7.5)
    
    def embed_external_css(self, html_content, base_path):
        """将外部 CSS 文件嵌入到 HTML 内部"""
        # 查找所有 <link rel="stylesheet" href="..."> 标签
        link_pattern = r'<link\s+rel=["\']stylesheet["\']\s+href=["\']([^"\']+)["\'][^>]*>'
        
        matches = list(re.finditer(link_pattern, html_content, re.IGNORECASE))
        
        if not matches:
            return html_content
        
        print(f"Embedding {len(matches)} external CSS file(s)...")
        
        # 从后往前替换，避免位置偏移
        for match in reversed(matches):
            css_href = match.group(1)
            
            # 跳过远程 CSS
            if css_href.startswith(('http://', 'https://')):
                print(f"  ⚠️  Skipping remote CSS: {css_href}")
                continue
            
            css_path = Path(base_path) / css_href
            
            if not css_path.exists():
                print(f"  ⚠️  CSS file not found: {css_path}")
                continue
            
            # 读取 CSS 内容
            with open(css_path, 'r', encoding='utf-8') as f:
                css_content = f.read()
            
            # 创建 <style> 标签
            style_tag = f'<style>\n{css_content}\n</style>'
            
            # 替换 <link> 标签
            html_content = html_content[:match.start()] + style_tag + html_content[match.end():]
            print(f"  ✅ Embedded: {css_path.name}")
        
        return html_content
    
    def inline_css_manual(self, soup, css_text):
        """手动内联 CSS - 使用 cssutils 解析并应用样式到元素"""
        cssutils.log.setLevel('CRITICAL')
        
        # 解析 CSS
        sheet = cssutils.parseString(css_text)
        
        # 第一遍：收集所有 CSS 变量定义
        css_vars = {}
        for rule in sheet:
            if hasattr(rule, 'style'):
                for prop in rule.style:
                    if prop.name.startswith('--'):
                        css_vars[prop.name] = prop.value
        
        # 解析 CSS 变量到实际值（处理变量引用变量）
        def resolve_value(value, depth=0):
            if depth > 10 or not value:
                return value
            var_match = re.search(r'var\((--[\w-]+)\)', value)
            if var_match:
                var_name = var_match.group(1)
                resolved = css_vars.get(var_name, value)
                # 递归解析嵌套变量
                return resolve_value(resolved, depth + 1)
            return value
        
        # 第二遍：收集要应用的规则（解析变量）
        rules_to_apply = []
        for rule in sheet:
            if hasattr(rule, 'selectorText') and hasattr(rule, 'style'):
                selector = rule.selectorText
                style_dict = {}
                for prop in rule.style:
                    if not prop.name.startswith('--'):  # 跳过变量定义
                        resolved_value = resolve_value(prop.value)
                        style_dict[prop.name] = resolved_value
                priority = getattr(rule.style, 'priority', 0) or 0
                rules_to_apply.append((selector, style_dict, priority))
        
        # 解析 CSS 变量到实际值
        def resolve_value(value, depth=0):
            if depth > 10 or not value:
                return value
            var_match = re.match(r'var\((--([\w-]+))\)', value.strip())
            if var_match:
                var_name = var_match.group(1)
                resolved = css_vars.get(var_name, value)
                return resolve_value(resolved, depth + 1)
            return value
        
        # 应用 CSS 规则到匹配的元素
        for selector, style_dict, priority in rules_to_apply:
            try:
                elements = soup.select(selector)
                for element in elements:
                    if not hasattr(element, 'name') or element.name is None:
                        continue
                    
                    # 获取或创建 style 属性
                    existing_style = element.get('style', '')
                    existing_dict = self.parse_style_attribute(existing_style)
                    
                    # 应用新样式（已解析 CSS 变量）
                    for prop_name, prop_value in style_dict.items():
                        resolved_value = resolve_value(prop_value)
                        existing_dict[prop_name] = resolved_value
                    
                    # 转换回 style 字符串
                    style_str = '; '.join(f'{k}: {v}' for k, v in existing_dict.items())
                    element['style'] = style_str
            except Exception as e:
                print(f"  ⚠️  Selector error: {selector} - {e}")
        
        return soup
    
    def inline_css(self, html_content, base_path=None):
        """内联 CSS 样式"""
        # 先嵌入外部 CSS 文件
        html_content = self.embed_external_css(html_content, base_path)
        
        # 解析 HTML
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # 提取所有 <style> 标签的 CSS
        all_css = ""
        for style_tag in soup.find_all('style'):
            if style_tag.string:
                all_css += style_tag.string + "\n"
        
        # 手动内联 CSS
        soup = self.inline_css_manual(soup, all_css)
        
        return str(soup)
    
    def parse_style_attribute(self, style_str):
        """解析 style 属性字符串为字典"""
        styles = {}
        if not style_str:
            return styles
        
        for declaration in style_str.split(';'):
            if ':' in declaration:
                key, value = declaration.split(':', 1)
                styles[key.strip().lower()] = value.strip()
        
        return styles
    
    def create_slide(self, prs, slide_element):
        """从 HTML slide 元素创建 PPTX 幻灯片"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # 解析 slide 的 style 属性
        style_str = slide_element.get('style', '')
        styles = self.parse_style_attribute(style_str)
        
        # 应用背景
        bg_color = styles.get('background-color') or styles.get('background')
        if bg_color and 'gradient' not in str(bg_color):
            rgb = self.css_parser.parse_color(bg_color)
            if rgb:
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(*rgb)
        
        # 遍历子元素
        for child in slide_element.children:
            self.process_element(slide, child, Inches(0), Inches(0))
        
        return slide
    
    def process_element(self, slide, element, offset_x, offset_y, parent_shape=None):
        """处理 HTML 元素并转换为 PPTX 形状"""
        # 跳过文本节点
        if not hasattr(element, 'name') or element.name is None:
            return
        
        tag = element.name
        
        # 跳过 script, style 等非内容元素
        if tag in ['script', 'style', 'meta', 'link', 'head']:
            return
        
        # 获取样式
        style_str = element.get('style', '')
        styles = self.parse_style_attribute(style_str)
        
        # 计算位置
        left = offset_x + self.css_parser.parse_size(styles.get('left', '0'))
        top = offset_y + self.css_parser.parse_size(styles.get('top', '0'))
        
        # 如果没有显式位置，使用相对定位
        if styles.get('position') not in ['absolute', 'fixed']:
            # 简化处理：使用默认位置
            pass
        
        width = self.css_parser.parse_size(styles.get('width', '100%'))
        height = self.css_parser.parse_size(styles.get('height', 'auto'))
        
        # 根据标签类型创建不同的形状
        if tag == 'img':
            self.create_image_shape(slide, element, left, top, width, height)
        elif tag in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'span', 'div', 'li']:
            text = element.get_text(strip=True)
            if text:
                self.create_textbox(slide, element, text, left, top, width, height, parent_shape)
        elif tag == 'svg':
            self.create_svg_shape(slide, element, left, top, width, height)
        
        # 递归处理子元素
        for child in element.children:
            if child.name and child.name not in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'span', 'img', 'svg']:
                self.process_element(slide, child, left, top, None)
    
    def create_textbox(self, slide, element, text, left, top, width, height, parent_shape=None):
        """创建文本框"""
        style_str = element.get('style', '')
        styles = self.parse_style_attribute(style_str)
        
        # 创建文本框
        if width == Inches(0):
            width = Inches(2)
        if height == Inches(0):
            height = Inches(0.5)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        p.text = text
        
        # 应用字体样式
        font = p.font
        font_size = self.css_parser.parse_font_size(styles.get('font-size'))
        font.size = font_size
        
        color = styles.get('color')
        if color:
            rgb = self.css_parser.parse_color(color)
            if rgb:
                font.color.rgb = RGBColor(*rgb)
        
        # 字重
        font_weight = styles.get('font-weight', 'normal')
        if font_weight in ['bold', '700', '800', '900']:
            font.bold = True
        
        # 斜体
        font_style = styles.get('font-style', 'normal')
        if font_style == 'italic':
            font.italic = True
        
        # 对齐
        text_align = self.css_parser.parse_text_align(styles.get('text-align'))
        p.alignment = text_align
        
        # 背景填充
        bg_color = styles.get('background-color')
        if bg_color:
            rgb = self.css_parser.parse_color(bg_color)
            if rgb:
                fill = textbox.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(*rgb)
        
        # 边框
        border = styles.get('border')
        if border and border != 'none':
            # 简化处理：只处理 solid 边框
            line = textbox.line
            line.width = Pt(1)
            border_color = styles.get('border-color')
            if border_color:
                rgb = self.css_parser.parse_color(border_color)
                if rgb:
                    line.color.rgb = RGBColor(*rgb)
        
        # 圆角 - PPTX 中文本框不支持圆角，跳过
    
    def create_image_shape(self, slide, element, left, top, width, height):
        """创建图片形状"""
        src = element.get('src', '')
        
        if not src:
            return
        
        # 处理图片源
        try:
            if src.startswith('data:image'):
                # Base64 图片
                match = re.match(r'data:image/(\w+);base64,(.+)', src)
                if match:
                    img_data = base64.b64decode(match.group(2))
                    img = io.BytesIO(img_data)
                    slide.shapes.add_picture(img, left, top, width=width if width != Inches(0) else None, height=height if height != Inches(0) else None)
            elif src.startswith(('http://', 'https://')):
                # 网络图片
                response = requests.get(src)
                img = io.BytesIO(response.content)
                slide.shapes.add_picture(img, left, top, width=width if width != Inches(0) else None, height=height if height != Inches(0) else None)
            elif os.path.exists(src):
                # 本地文件
                slide.shapes.add_picture(src, left, top, width=width if width != Inches(0) else None, height=height if height != Inches(0) else None)
        except Exception as e:
            print(f"Image loading warning: {e}")
    
    def create_svg_shape(self, slide, element, left, top, width, height):
        """创建 SVG 形状（简化为占位符）"""
        # SVG 转换复杂，暂时创建占位符文本框
        textbox = slide.shapes.add_textbox(left, top, width or Inches(2), height or Inches(2))
        tf = textbox.text_frame
        tf.text = "[SVG Placeholder]"
    
    def convert(self, html_path, output_path=None):
        """执行转换"""
        html_path = Path(html_path)
        
        if not html_path.exists():
            raise FileNotFoundError(f"HTML file not found: {html_path}")
        
        # 读取 HTML
        with open(html_path, 'r', encoding='utf-8') as f:
            html_content = f.read()
        
        # 内联 CSS
        base_path = str(html_path.parent)
        inlined_html = self.inline_css(html_content, base_path)
        
        # 解析 HTML
        soup = BeautifulSoup(inlined_html, 'html.parser')
        
        # 提取 CSS 变量
        style_tags = soup.find_all('style')
        for style_tag in style_tags:
            if style_tag.string:
                self.css_parser.parse_css_vars(style_tag.string)
        
        # 创建 PPTX
        prs = Presentation()
        prs.slide_width = self.slide_width
        prs.slide_height = self.slide_height
        
        # 查找所有 slide
        slides = soup.find_all('section', class_='slide')
        
        if not slides:
            # 如果没有找到 slide 结构，将整个 HTML 作为单页
            body = soup.find('body')
            if body:
                slides = [body]
        
        # 转换每个 slide
        for i, slide_element in enumerate(slides):
            print(f"Converting slide {i+1}/{len(slides)}...")
            self.create_slide(prs, slide_element)
        
        # 保存 PPTX
        if not output_path:
            output_path = html_path.with_suffix('')
            output_path = output_path.with_name(f"{output_path.name}_converted.pptx")
        else:
            output_path = Path(output_path)
        
        prs.save(str(output_path))
        print(f"Saved: {output_path}")
        
        return {
            'output_file': str(output_path),
            'slides_count': len(slides)
        }


def run(args=None):
    """主入口函数"""
    if args is None:
        args = sys.argv[1:]
    
    if len(args) < 1:
        print("Usage: python3 index.py <html-file> [output-file]")
        sys.exit(1)
    
    html_path = args[0]
    output_path = args[1] if len(args) > 1 else None
    
    converter = HTMLToPPTXConverter()
    result = converter.convert(html_path, output_path)
    
    return result


if __name__ == '__main__':
    run()
