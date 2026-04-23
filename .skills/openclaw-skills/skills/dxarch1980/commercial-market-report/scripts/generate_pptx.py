#!/usr/bin/env python3
"""
生成商业市调报告 PPTX 文件（含图表）
用法: python generate_pptx.py "项目名称" [site_area] [far] --output "输出路径"
"""
import sys, os, datetime

def _make_charts():
    """生成图表，返回 {name: filepath}"""
    tmp = os.getenv('TEMP', '/tmp')
    out = {}
    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
    except Exception:
        return out
    plt.rcParams['font.sans-serif'] = ['Microsoft YaHei', 'SimHei']
    plt.rcParams['axes.unicode_minus'] = False

    # 1. 业态配比饼图
    fig, ax = plt.subplots(figsize=(5, 4))
    labels = ['餐饮35%', '零售30%', '亲子体验20%', '体验娱乐10%', '配套5%']
    sizes = [35, 30, 20, 10, 5]
    colors = ['#0052cc', '#3377ee', '#5599ff', '#77bbff', '#99ccff']
    wedges, texts, autotexts = ax.pie(sizes, labels=labels, colors=colors, autopct='%1.0f%%', startangle=90, textprops={'fontsize': 10})
    for t in texts: t.set_fontsize(9)
    for a in autotexts: a.set_fontsize(9); a.set_color('white')
    ax.set_title('业态面积配比', fontsize=12, fontweight='bold', pad=10)
    plt.tight_layout()
    p = os.path.join(tmp, 'chart_mix.png')
    plt.savefig(p, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    out['mix'] = p

    # 2. 竞品柱状图
    fig, ax = plt.subplots(figsize=(6, 3.5))
    names = ['龙湖天街', '和谐广场', '万象城', '高新万达', '印象城']
    areas = [8, 6, 10, 8, 5]
    colors2 = ['#dc2626', '#f59e0b', '#10b981', '#3b82f6', '#8b5cf6']
    bars = ax.bar(names, areas, color=colors2, width=0.5, edgecolor='none')
    for bar, val in zip(bars, areas):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.15, f'{val}万m²', ha='center', va='bottom', fontsize=9)
    ax.set_ylabel('体量（万m²）', fontsize=10)
    ax.set_title('周边竞品体量对比', fontsize=12, fontweight='bold')
    ax.set_ylim(0, 13)
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    plt.tight_layout()
    p2 = os.path.join(tmp, 'chart_competitors.png')
    plt.savefig(p2, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    out['competitors'] = p2

    # 3. 投资构成饼图
    fig, ax = plt.subplots(figsize=(5, 4))
    labels3 = ['土地成本\n1.2亿', '建安成本\n1.8亿', '装修成本\n0.8亿', '预备金\n0.4亿']
    sizes3 = [12000, 18000, 8000, 4000]
    colors3 = ['#0052cc', '#3377ee', '#5599ff', '#99ccff']
    wedges3, texts3, autotexts3 = ax.pie(sizes3, labels=labels3, colors=colors3, autopct='%1.0f%%', startangle=90, textprops={'fontsize': 9})
    for t in texts3: t.set_fontsize(9)
    for a in autotexts3: a.set_fontsize(9); a.set_color('white')
    ax.set_title('投资构成（万元）', fontsize=12, fontweight='bold', pad=10)
    plt.tight_layout()
    p3 = os.path.join(tmp, 'chart_invest.png')
    plt.savefig(p3, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    out['invest'] = p3

    # 4. 敏感性分析柱状图
    fig, ax = plt.subplots(figsize=(5, 3.5))
    scenarios = ['+20%', '+10%', '基准', '-10%', '-20%']
    irr = [14.5, 12.0, 10.5, 8.2, 5.5]
    colors4 = ['#10b981', '#84cc16', '#0052cc', '#f59e0b', '#dc2626']
    bars4 = ax.bar(scenarios, irr, color=colors4, width=0.5, edgecolor='none')
    for bar, val in zip(bars4, irr):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.1, f'{val}%', ha='center', va='bottom', fontsize=9)
    ax.set_ylabel('IRR（%）', fontsize=10)
    ax.set_title('敏感性分析', fontsize=12, fontweight='bold')
    ax.set_ylim(0, 18)
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    plt.tight_layout()
    p4 = os.path.join(tmp, 'chart_sensitivity.png')
    plt.savefig(p4, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    out['sensitivity'] = p4

    return out


def generate_pptx(project_name, site_area=28000, far=3.5, building_density='45%', height_limit=80, competitors=None, output_dir=None):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("需要 python-pptx")
        sys.exit(1)

    chart_paths = _make_charts()
    site_area = float(site_area)
    far = float(far)
    total_area = site_area * far
    commercial_area = total_area * 0.75
    dining = commercial_area * 0.35
    retail = commercial_area * 0.30
    experience = commercial_area * 0.20
    support = commercial_area * 0.10

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    ACCENT = RGBColor(0, 82, 204)
    WHITE = RGBColor(255, 255, 255)
    GRAY = RGBColor(156, 163, 175)
    DARK = RGBColor(55, 65, 81)

    def blank_slide():
        s = prs.slides.add_slide(prs.slide_layouts[6])
        fill = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        fill.fill.solid(); fill.fill.fore_color.rgb = WHITE; fill.line.fill.background()
        return s

    def add_chart_to_slide(slide, name, left, top, width, height):
        if name not in chart_paths:
            return
        slide.shapes.add_picture(chart_paths[name], left, top, width, height)

    # ===== SLIDE 1: COVER =====
    s = blank_slide()
    fill = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    fill.fill.solid(); fill.fill.fore_color.rgb = RGBColor(26, 26, 46); fill.line.fill.background()
    tx = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11), Inches(1.5))
    p = tx.text_frame.paragraphs[0]; p.text = project_name
    p.font.size = Pt(48); p.font.bold = True; p.font.color.rgb = WHITE
    tx2 = s.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(11), Inches(0.8))
    p2 = tx2.text_frame.paragraphs[0]; p2.text = "市场调研汇报"
    p2.font.size = Pt(24); p2.font.color.rgb = RGBColor(147, 197, 253)
    tx3 = s.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(11), Inches(0.5))
    p3 = tx3.text_frame.paragraphs[0]
    p3.text = f"📍 济南槐荫区  |  📐 用地面积：{site_area:,.0f}m²  |  可建面积：{total_area:,.0f}m²  |  📅 {datetime.datetime.now().strftime('%Y年%m月')}"
    p3.font.size = Pt(14); p3.font.color.rgb = GRAY

    # ===== SLIDE 2: TOC =====
    s = blank_slide()
    bar = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    tx = s.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(5), Inches(0.35))
    tx.text_frame.paragraphs[0].text = "CONTENTS"; tx.text_frame.paragraphs[0].font.size = Pt(10); tx.text_frame.paragraphs[0].font.color.rgb = ACCENT
    tx2 = s.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12), Inches(0.7))
    p = tx2.text_frame.paragraphs[0]; p.text = "汇报目录"
    p.font.size = Pt(32); p.font.bold = True; p.font.color.rgb = DARK
    div = s.shapes.add_shape(1, Inches(0.6), Inches(1.3), Inches(0.5), Pt(3))
    div.fill.solid(); div.fill.fore_color.rgb = ACCENT; div.line.fill.background()
    # TOC items in 2 columns
    for i, ch in enumerate(['01 项目概况与区域分析','02 竞争市场分析','03 消费者调研分析','04 商业定位与业态组合','05 财务测算','06 结论与设计建议']):
        col = 0 if i < 3 else 1
        row = i % 3
        tx = s.shapes.add_textbox(Inches(0.6 + col*6), Inches(1.6 + row*1.4), Inches(5), Inches(0.5))
        p = tx.text_frame.paragraphs[0]; p.text = ch; p.font.size = Pt(16); p.font.color.rgb = GRAY

    # ===== SLIDE 3: 项目概况 =====
    s = blank_slide()
    bar = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    tx = s.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(5), Inches(0.35))
    tx.text_frame.paragraphs[0].text = "CHAPTER 01"; tx.text_frame.paragraphs[0].font.size = Pt(10); tx.text_frame.paragraphs[0].font.color.rgb = ACCENT
    tx2 = s.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12), Inches(0.7))
    p = tx2.text_frame.paragraphs[0]; p.text = "项目概况与区域分析"
    p.font.size = Pt(32); p.font.bold = True; p.font.color.rgb = DARK
    div = s.shapes.add_shape(1, Inches(0.6), Inches(1.3), Inches(0.5), Pt(3))
    div.fill.solid(); div.fill.fore_color.rgb = ACCENT; div.line.fill.background()
    # Stats
    stats = [("用地面积", f"{site_area:,.0f} m²"), ("容积率", str(far)), ("可建面积", f"{total_area:,.0f} m²"), ("商业面积", f"{commercial_area:,.0f} m²")]
    for i, (label, val) in enumerate(stats):
        col = i % 2; row = i // 2
        tx = s.shapes.add_textbox(Inches(0.6 + col*6), Inches(1.6 + row*2), Inches(5), Inches(1.5))
        p = tx.text_frame.paragraphs[0]; p.text = val; p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = ACCENT
        tx2 = s.shapes.add_textbox(Inches(0.6 + col*6), Inches(2.2 + row*2), Inches(5), Inches(0.4))
        p2 = tx2.text_frame.paragraphs[0]; p2.text = label; p2.font.size = Pt(12); p2.font.color.rgb = GRAY

    # ===== SLIDE 4: 竞品分析 =====
    s = blank_slide()
    bar = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    tx = s.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(5), Inches(0.35))
    tx.text_frame.paragraphs[0].text = "CHAPTER 02"; tx.text_frame.paragraphs[0].font.size = Pt(10); tx.text_frame.paragraphs[0].font.color.rgb = ACCENT
    tx2 = s.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12), Inches(0.7))
    p = tx2.text_frame.paragraphs[0]; p.text = "竞争市场分析 — 竞品体量对比"
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = DARK
    div = s.shapes.add_shape(1, Inches(0.6), Inches(1.3), Inches(0.5), Pt(3))
    div.fill.solid(); div.fill.fore_color.rgb = ACCENT; div.line.fill.background()
    add_chart_to_slide(s, 'competitors', Inches(0.5), Inches(1.5), Inches(7), Inches(5.5))

    # ===== SLIDE 5: SWOT =====
    s = blank_slide()
    bar = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    tx = s.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(5), Inches(0.35))
    tx.text_frame.paragraphs[0].text = "CHAPTER 02"; tx.text_frame.paragraphs[0].font.size = Pt(10); tx.text_frame.paragraphs[0].font.color.rgb = ACCENT
    tx2 = s.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12), Inches(0.7))
    p = tx2.text_frame.paragraphs[0]; p.text = "SWOT 分析"
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = DARK
    div = s.shapes.add_shape(1, Inches(0.6), Inches(1.3), Inches(0.5), Pt(3))
    div.fill.solid(); div.fill.fore_color.rgb = ACCENT; div.line.fill.background()
    swots = [
        ("S 优势", RGBColor(22, 163, 74), "区位优势：西部门户，地铁上盖\n政策支持：西进战略核心承载区\n配套成熟：医疗/教育/政务齐全"),
        ("W 劣势", RGBColor(220, 38, 38), "区域认知度：新区需培育期\n竞品分流：龙湖/万达先发优势\n招商难度：品牌商家入驻门槛"),
        ("O 机会", RGBColor(37, 99, 235), "市场缺口：精品商业供给不足\n人口导入：新建住宅持续交付\n消费升级：品质需求持续增长"),
        ("T 威胁", RGBColor(234, 88, 12), "竞品分流：体量竞争压力\n电商冲击：实体商业受挤压\n成本上涨：地价/建安成本上升"),
    ]
    for i, (title, color, text) in enumerate(swots):
        col = i % 2; row = i // 2
        left = Inches(0.5 + col*6.5); top = Inches(1.5 + row*2.8); w = Inches(6); h = Inches(2.5)
        bg = s.shapes.add_shape(1, left, top, w, h)
        bg.fill.solid()
        if row == 0 and col == 0: bg.fill.fore_color.rgb = RGBColor(240, 253, 244)
        elif row == 0 and col == 1: bg.fill.fore_color.rgb = RGBColor(254, 242, 242)
        elif row == 1 and col == 0: bg.fill.fore_color.rgb = RGBColor(239, 246, 255)
        else: bg.fill.fore_color.rgb = RGBColor(255, 247, 237)
        bg.line.fill.background()
        title_tx = s.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), w - Inches(0.4), Inches(0.5))
        tp = title_tx.text_frame.paragraphs[0]; tp.text = title
        tp.font.size = Pt(14); tp.font.bold = True; tp.font.color.rgb = color
        body_tx = s.shapes.add_textbox(left + Inches(0.2), top + Inches(0.6), w - Inches(0.4), Inches(1.8))
        bp = body_tx.text_frame.paragraphs[0]; bp.text = text
        bp.font.size = Pt(11); bp.font.color.rgb = DARK

    # ===== SLIDE 6: 客群画像 =====
    s = blank_slide()
    bar = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    tx = s.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(5), Inches(0.35))
    tx.text_frame.paragraphs[0].text = "CHAPTER 03"; tx.text_frame.paragraphs[0].font.size = Pt(10); tx.text_frame.paragraphs[0].font.color.rgb = ACCENT
    tx2 = s.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12), Inches(0.7))
    p = tx2.text_frame.paragraphs[0]; p.text = "消费者调研 — 客群画像"
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = DARK
    div = s.shapes.add_shape(1, Inches(0.6), Inches(1.3), Inches(0.5), Pt(3))
    div.fill.solid(); div.fill.fore_color.rgb = ACCENT; div.line.fill.background()
    bullets = [
        "核心客群：25-45岁，白领/公务员，家庭年收入15-30万",
        "次级客群：银发群体（区域老龄化高）| 亲子家庭（儿童消费刚需）",
        "消费偏好：餐饮高频（网红餐厅/聚餐）| 亲子刚需 | 体验升级（影院/健身）",
        "业态优先级：🔥精品超市+餐饮  ⚡电影院+健身  ○美容SPA+文创",
    ]
    for i, b in enumerate(bullets):
        tx = s.shapes.add_textbox(Inches(0.6), Inches(1.5 + i*1.2), Inches(12), Inches(1))
        p = tx.text_frame.paragraphs[0]; p.text = "• " + b
        p.font.size = Pt(16); p.font.color.rgb = DARK

    # ===== SLIDE 7: 业态配比 =====
    s = blank_slide()
    bar = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    tx = s.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(5), Inches(0.35))
    tx.text_frame.paragraphs[0].text = "CHAPTER 04"; tx.text_frame.paragraphs[0].font.size = Pt(10); tx.text_frame.paragraphs[0].font.color.rgb = ACCENT
    tx2 = s.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12), Inches(0.7))
    p = tx2.text_frame.paragraphs[0]; p.text = "商业定位 — 业态面积配比"
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = DARK
    div = s.shapes.add_shape(1, Inches(0.6), Inches(1.3), Inches(0.5), Pt(3))
    div.fill.solid(); div.fill.fore_color.rgb = ACCENT; div.line.fill.background()
    add_chart_to_slide(s, 'mix', Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.5))
    # Table on right
    rows_data = [
        ("业态", "占比", "面积"),
        ("餐饮", "35%", f"{dining:,.0f} m²"),
        ("零售", "30%", f"{retail:,.0f} m²"),
        ("亲子体验", "20%", f"{experience:,.0f} m²"),
        ("体验娱乐", "10%", f"{support:,.0f} m²"),
        ("配套服务", "5%", f"{commercial_area*0.05:,.0f} m²"),
    ]
    x = Inches(6.5); y = Inches(1.5); row_h = Inches(0.55)
    for ri, row in enumerate(rows_data):
        for ci, cell in enumerate(row):
            cell_w = Inches(2) if ci == 0 else Inches(1.5)
            tx = s.shapes.add_textbox(x + ci*Inches(2.3), y + ri*row_h, cell_w, row_h)
            p = tx.text_frame.paragraphs[0]; p.text = cell
            p.font.size = Pt(12)
            if ri == 0: p.font.bold = True; p.font.color.rgb = WHITE
            else: p.font.color.rgb = DARK

    # ===== SLIDE 8: 财务测算 =====
    s = blank_slide()
    bar = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    tx = s.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(5), Inches(0.35))
    tx.text_frame.paragraphs[0].text = "CHAPTER 05"; tx.text_frame.paragraphs[0].font.size = Pt(10); tx.text_frame.paragraphs[0].font.color.rgb = ACCENT
    tx2 = s.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12), Inches(0.7))
    p = tx2.text_frame.paragraphs[0]; p.text = "财务测算 — 投资构成"
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = DARK
    div = s.shapes.add_shape(1, Inches(0.6), Inches(1.3), Inches(0.5), Pt(3))
    div.fill.solid(); div.fill.fore_color.rgb = ACCENT; div.line.fill.background()
    add_chart_to_slide(s, 'invest', Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.5))
    fin_bullets = [
        "土地成本：1.2亿",
        "建安成本：1.8亿",
        "装修成本：0.8亿",
        "预备金：0.4亿",
        "合计：约4.2亿",
        "预期IRR：10-12%",
        "回收期：8-10年",
    ]
    for i, b in enumerate(fin_bullets):
        tx = s.shapes.add_textbox(Inches(6.5), Inches(1.5 + i*0.75), Inches(6), Inches(0.7))
        p = tx.text_frame.paragraphs[0]; p.text = "• " + b
        p.font.size = Pt(15); p.font.color.rgb = DARK
        if i == 4: p.font.bold = True; p.font.color.rgb = ACCENT

    # ===== SLIDE 9: 敏感性 =====
    s = blank_slide()
    bar = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    tx = s.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(5), Inches(0.35))
    tx.text_frame.paragraphs[0].text = "CHAPTER 05"; tx.text_frame.paragraphs[0].font.size = Pt(10); tx.text_frame.paragraphs[0].font.color.rgb = ACCENT
    tx2 = s.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12), Inches(0.7))
    p = tx2.text_frame.paragraphs[0]; p.text = "敏感性分析 — 不同租金下的IRR"
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = DARK
    div = s.shapes.add_shape(1, Inches(0.6), Inches(1.3), Inches(0.5), Pt(3))
    div.fill.solid(); div.fill.fore_color.rgb = ACCENT; div.line.fill.background()
    add_chart_to_slide(s, 'sensitivity', Inches(0.5), Inches(1.5), Inches(7), Inches(5.5))
    sens = [
        ("乐观（+10%）", "IRR 13-15%，回收期7-8年", RGBColor(16, 185, 129)),
        ("基准", "IRR 10-12%，回收期8-10年", ACCENT),
        ("悲观（-10%）", "IRR 7-9%，回收期10-12年", RGBColor(245, 158, 11)),
        ("风险（-20%）", "IRR 4-6%，回收期>15年", RGBColor(220, 38, 38)),
    ]
    for i, (label, text, color) in enumerate(sens):
        tx = s.shapes.add_textbox(Inches(8), Inches(1.5 + i*1.4), Inches(5), Inches(0.5))
        p = tx.text_frame.paragraphs[0]; p.text = label; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = color
        tx2 = s.shapes.add_textbox(Inches(8), Inches(1.9 + i*1.4), Inches(5), Inches(0.5))
        p2 = tx2.text_frame.paragraphs[0]; p2.text = text; p2.font.size = Pt(12); p2.font.color.rgb = GRAY

    # ===== SLIDE 10: 结论 =====
    s = blank_slide()
    bar = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    tx = s.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(5), Inches(0.35))
    tx.text_frame.paragraphs[0].text = "CHAPTER 06"; tx.text_frame.paragraphs[0].font.size = Pt(10); tx.text_frame.paragraphs[0].font.color.rgb = ACCENT
    tx2 = s.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12), Inches(0.7))
    p = tx2.text_frame.paragraphs[0]; p.text = "结论与设计建议"
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = DARK
    div = s.shapes.add_shape(1, Inches(0.6), Inches(1.3), Inches(0.5), Pt(3))
    div.fill.solid(); div.fill.fore_color.rgb = ACCENT; div.line.fill.background()
    conclusions = [
        ("✅ 核心结论", ["市场可行性良好，差异化定位存在竞争空间", "强餐饮+强亲子是核心竞争力", "需与龙湖/万达形成错位"]),
        ("🏗️ 设计建议", ["街区+盒子组合，兼顾休闲与品质", "双中庭+环形动线，避免商业死角", "800-1000个停车位满足家庭需求", "分期开发，一期F1-F3优先开业（80%+）"]),
    ]
    y_off = 0
    for i, (section_title, items) in enumerate(conclusions):
        tx = s.shapes.add_textbox(Inches(0.6 + i*6.5), Inches(1.5 + y_off), Inches(6), Inches(0.5))
        p = tx.text_frame.paragraphs[0]; p.text = section_title
        p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = ACCENT
        y_off += 0.5
        for item in items:
            tx = s.shapes.add_textbox(Inches(0.6 + i*6.5), Inches(1.5 + y_off), Inches(6), Inches(0.6))
            p = tx.text_frame.paragraphs[0]; p.text = "• " + item
            p.font.size = Pt(12); p.font.color.rgb = DARK
            y_off += 0.6
        y_off = 0

    # ===== SLIDE 11: CLOSING =====
    s = blank_slide()
    fill = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    fill.fill.solid(); fill.fill.fore_color.rgb = RGBColor(26, 26, 46); fill.line.fill.background()
    tx = s.shapes.add_textbox(Inches(0), Inches(3), prs.slide_width, Inches(1))
    tx.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tx.text_frame.paragraphs[0]; p.text = "谢谢观看"
    p.font.size = Pt(48); p.font.bold = True; p.font.color.rgb = WHITE
    tx2 = s.shapes.add_textbox(Inches(0), Inches(4.2), prs.slide_width, Inches(0.6))
    tx2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    p2 = tx2.text_frame.paragraphs[0]; p2.text = f"{project_name} · 市场调研汇报"
    p2.font.size = Pt(16); p2.font.color.rgb = GRAY

    # Save
    date_str = datetime.datetime.now().strftime('%Y%m%d')
    filename = f'{project_name}_汇报PPT_{date_str}.pptx'
    out_path = os.path.join(output_dir, filename) if output_dir else filename
    prs.save(out_path)
    print(f"PPTX_SAVED:{out_path}")
    return out_path


if __name__ == '__main__':
    import argparse
    p = argparse.ArgumentParser()
    p.add_argument('project_name', nargs='?', default='商业项目')
    p.add_argument('site_area', nargs='?', default='28000')
    p.add_argument('far', nargs='?', default='3.5')
    p.add_argument('building_density', nargs='?', default='45%')
    p.add_argument('height_limit', nargs='?', default='80')
    p.add_argument('--output', dest='output_dir', default=None)
    a = p.parse_args()
    generate_pptx(
        project_name=a.project_name,
        site_area=a.site_area,
        far=a.far,
        building_density=a.building_density,
        height_limit=a.height_limit,
        output_dir=a.output_dir,
    )
