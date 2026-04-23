"""
PPTX converter — rich visual slides from HTML content.

Workflow:
1) Parse HTML into structured sections (headings, paragraphs, tables, stats, cards).
2) Ask LLM to produce a typed slide plan (title_slide / content / stats / table / two_column / highlight).
3) Render each slide with colored backgrounds, rounded-rect cards, shadows, and proper typography.
"""

from __future__ import annotations

import json
import os
import re
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from bs4 import BeautifulSoup, Tag
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}

PALETTE = {
    "primary": "#2563EB",
    "secondary": "#1E40AF",
    "accent": "#10B981",
    "warning": "#F59E0B",
    "danger": "#EF4444",
    "dark": "#1E293B",
    "light": "#F1F5F9",
    "white": "#FFFFFF",
    "muted": "#64748B",
}

SLIDE_BG_COLORS = [
    "#FFFFFF",
    "#F8FAFC",
    "#F1F5F9",
    "#FFFBEB",
    "#F0FDF4",
    "#EFF6FF",
    "#FDF2F8",
]


def _hex(color: str) -> RGBColor:
    h = color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _add_shadow(shape, blur: int = 50800, dist: int = 38100, alpha: int = 25):
    sp = shape._element
    spPr = sp.find(".//p:spPr", NSMAP)
    if spPr is None:
        return
    effectLst = spPr.find("a:effectLst", NSMAP)
    if effectLst is None:
        effectLst = etree.SubElement(spPr, "{%s}effectLst" % NSMAP["a"])
    outerShdw = etree.SubElement(effectLst, "{%s}outerShdw" % NSMAP["a"])
    outerShdw.set("blurRad", str(blur))
    outerShdw.set("dist", str(dist))
    outerShdw.set("dir", "2700000")
    outerShdw.set("algn", "tl")
    outerShdw.set("rotWithShape", "0")
    clr = etree.SubElement(outerShdw, "{%s}srgbClr" % NSMAP["a"])
    clr.set("val", "000000")
    a = etree.SubElement(clr, "{%s}alpha" % NSMAP["a"])
    a.set("val", str(alpha * 1000))


# ---------------------------------------------------------------------------
# Drawing helpers
# ---------------------------------------------------------------------------

class SlideBuilder:
    """Low-level drawing primitives for a single slide."""

    W, H = 16.0, 9.0

    def __init__(self, prs: Presentation):
        self.prs = prs

    def add_blank(self) -> Any:
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    @staticmethod
    def rect(slide, x, y, w, h, fill="#FFFFFF", border=None, border_w=1.0, shadow=False):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex(fill)
        if border:
            shape.line.color.rgb = _hex(border)
            shape.line.width = Pt(border_w)
        else:
            shape.line.fill.background()
        if shadow:
            _add_shadow(shape)
        return shape

    @staticmethod
    def rrect(slide, x, y, w, h, fill="#FFFFFF", border=None, border_w=1.0, radius=0.15, shadow=False):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex(fill)
        if border:
            shape.line.color.rgb = _hex(border)
            shape.line.width = Pt(border_w)
        else:
            shape.line.fill.background()
        sp = shape._element
        prstGeom = sp.find(".//a:prstGeom", NSMAP)
        if prstGeom is not None:
            avLst = prstGeom.find("a:avLst", NSMAP)
            if avLst is None:
                avLst = etree.SubElement(prstGeom, "{%s}avLst" % NSMAP["a"])
            for child in list(avLst):
                avLst.remove(child)
            gd = etree.SubElement(avLst, "{%s}gd" % NSMAP["a"])
            gd.set("name", "adj")
            gd.set("fmla", f"val {int(radius * 100000 / min(w, h))}")
        if shadow:
            _add_shadow(shape)
        return shape

    @staticmethod
    def text(slide, x, y, w, h, txt, *, font="Arial", size=24, color="#000000",
             bold=False, italic=False, align="left", valign="top", wrap=True):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        tf.auto_size = MSO_AUTO_SIZE.NONE
        anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
        tf.vertical_anchor = anchor.get(valign, MSO_ANCHOR.TOP)
        p = tf.paragraphs[0]
        p.text = txt
        p.font.name = font
        p.font.size = Pt(size)
        p.font.color.rgb = _hex(color)
        p.font.bold = bold
        p.font.italic = italic
        al = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
        p.alignment = al.get(align, PP_ALIGN.LEFT)
        return tb

    @staticmethod
    def multiline(slide, x, y, w, h, lines: List[Dict], default_font="Arial"):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = ln.get("text", "")
            p.font.name = ln.get("font", default_font)
            p.font.size = Pt(ln.get("size", 20))
            p.font.color.rgb = _hex(ln.get("color", "#000000"))
            p.font.bold = ln.get("bold", False)
            p.font.italic = ln.get("italic", False)
            al = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
            p.alignment = al.get(ln.get("align", "left"), PP_ALIGN.LEFT)
            p.space_after = Pt(ln.get("space_after", 4))
        return tb

    @staticmethod
    def bg(slide, fill: str):
        """Full-slide background rectangle pushed to back (after required spTree headers)."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            Inches(SlideBuilder.W), Inches(SlideBuilder.H),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex(fill)
        shape.line.fill.background()
        sp = shape._element
        tree = sp.getparent()
        tree.remove(sp)
        # OOXML requires nvGrpSpPr[0] + grpSpPr[1] first; shapes start at index 2
        tree.insert(2, sp)
        return shape


# ---------------------------------------------------------------------------
# HTML extraction
# ---------------------------------------------------------------------------

def _extract_rich_payload(html_path: Path) -> Dict[str, Any]:
    html = html_path.read_text(encoding="utf-8")
    soup = BeautifulSoup(html, "html.parser")

    title = (soup.title.get_text(strip=True) if soup.title else "") or html_path.stem

    headings = []
    for node in soup.find_all(["h1", "h2", "h3"]):
        t = node.get_text(" ", strip=True)
        if t:
            headings.append(t)

    paragraphs = []
    for node in soup.find_all("p"):
        t = node.get_text(" ", strip=True)
        if t and len(t) > 10:
            paragraphs.append(t)

    # Extract tables
    tables: List[Dict] = []
    for tbl in soup.find_all("table"):
        headers = [th.get_text(" ", strip=True) for th in tbl.find_all("th")]
        rows = []
        for tr in tbl.find_all("tr"):
            cells = [td.get_text(" ", strip=True) for td in tr.find_all("td")]
            if cells:
                rows.append(cells)
        if headers or rows:
            tables.append({"headers": headers[:8], "rows": [r[:8] for r in rows[:10]]})

    # Extract stat items
    stats: List[Dict] = []
    for item in soup.find_all(class_=lambda c: c and "stat-item" in c):
        num_el = item.find(class_=lambda c: c and "stat-number" in c)
        label_el = item.find(class_=lambda c: c and "stat-label" in c)
        if num_el and label_el:
            stats.append({
                "value": num_el.get_text(strip=True),
                "label": label_el.get_text(strip=True),
            })

    # Extract cards
    cards: List[Dict] = []
    for card in soup.find_all(class_=lambda c: c and "policy-card" in c):
        card_title = ""
        h4 = card.find("h4")
        if h4:
            card_title = h4.get_text(" ", strip=True)
        highlight = ""
        hl = card.find(class_="highlight")
        if hl:
            highlight = hl.get_text(" ", strip=True)
        items = [li.get_text(" ", strip=True) for li in card.find_all("li")]
        cards.append({"title": card_title, "highlight": highlight, "items": items[:6]})

    # Extract highlight boxes
    highlights: List[Dict] = []
    for box in soup.find_all(class_=lambda c: c and "highlight-box" in c):
        h_title = ""
        h4 = box.find("h4")
        if h4:
            h_title = h4.get_text(" ", strip=True)
        body_parts = []
        for el in box.find_all(["p", "li"]):
            t = el.get_text(" ", strip=True)
            if t and t != h_title:
                body_parts.append(t)
        highlights.append({"title": h_title, "body": " ".join(body_parts[:4])})

    # Extract references (text + URL pairs)
    references: List[Dict] = []
    for li in soup.find_all(class_=lambda c: c and "ref-item" in c):
        text = li.get_text(" ", strip=True)
        link = li.find("a")
        url = (link.get("href") or "").strip() if link else ""
        if text:
            references.append({"text": text, "url": url})
    # Fallback: look for any section/div with "references" or "参考" in heading
    if not references:
        for section in soup.find_all(["section", "div"]):
            heading = section.find(["h2", "h3"])
            if heading and any(kw in heading.get_text() for kw in ("参考", "引用", "reference", "Reference", "来源", "Source")):
                for li in section.find_all("li"):
                    text = li.get_text(" ", strip=True)
                    link = li.find("a")
                    url = (link.get("href") or "").strip() if link else ""
                    if text:
                        references.append({"text": text, "url": url})
                break

    return {
        "title": title,
        "headings": headings[:20],
        "paragraphs": paragraphs[:60],
        "tables": tables[:5],
        "stats": stats[:10],
        "cards": cards[:10],
        "highlights": highlights[:5],
        "references": references[:20],
    }


# ---------------------------------------------------------------------------
# PPTXConverter
# ---------------------------------------------------------------------------

class PPTXConverter:
    """Generate rich PPTX from HTML using LLM-driven slide planning."""

    SLIDE_TYPES = ["title_slide", "content", "stats", "table", "two_column", "highlight", "references"]

    def __init__(self, model: Optional[str] = None):
        self.model = model or os.getenv("TEXT_MODEL") or "gpt-4.1-2025-04-14"
        self.api_key = os.getenv("OPENAI_API_KEY")
        self.base_url = os.getenv("OPENAI_BASE_URL")

    # ------------------------------------------------------------------
    # LLM slide planning
    # ------------------------------------------------------------------

    def _plan_slides_with_llm(self, payload: Dict[str, Any], max_slides: int) -> Optional[List[Dict[str, Any]]]:
        if not self.api_key:
            print("[PPTXConverter] No OPENAI_API_KEY found, using fallback plan.")
            return None
        try:
            from openai import OpenAI
        except Exception:
            print("[PPTXConverter] openai package not installed, using fallback plan.")
            return None

        print(f"[PPTXConverter] Calling LLM ({self.model}) via {self.base_url or 'https://api.openai.com'} ...")
        client = OpenAI(api_key=self.api_key, base_url=self.base_url) if self.base_url else OpenAI(api_key=self.api_key)

        system_prompt = (
            "You are a presentation designer. Convert structured HTML content into a slide plan. "
            "Return JSON only. Keep slide text concise and presentation-ready."
        )

        user_prompt = {
            "max_slides": max_slides,
            "instructions": (
                "Create a slide deck from the payload. Use varied slide types. "
                "Each slide should have: type, title, subtitle (optional), accent_color (#hex), and items. "
                "Slide types: title_slide (items=[]), content (items=[bullet strings]), "
                'stats (items=[{"value":"...","label":"..."}]), '
                'table (items={"headers":[...],"rows":[[...],...]}), '
                'two_column (items=[{"title":"...","points":["..."]}]), '
                'highlight (items={"heading":"...","body":"..."}), '
                'references (items=[{"text":"...","url":"..."}]). '
                "If the payload contains references, always include a references slide as the last slide."
            ),
            "payload": payload,
        }

        try:
            resp = client.chat.completions.create(
                model=self.model,
                temperature=0.3,
                max_tokens=4000,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": json.dumps(user_prompt, ensure_ascii=False)},
                ],
            )
        except Exception as e:
            print(f"[PPTXConverter] LLM call failed: {e}")
            return None

        text = (resp.choices[0].message.content or "").strip() if resp and resp.choices else ""
        if not text:
            print("[PPTXConverter] LLM returned empty response, using fallback.")
            return None

        import re
        # qwen3 models wrap reasoning in <think>...</think> — extract content AFTER it
        if "<think>" in text:
            # Try content after </think> first
            after_think = re.sub(r"^.*</think>", "", text, count=1, flags=re.DOTALL).strip()
            if after_think and "{" in after_think:
                text = after_think
            else:
                # JSON might be inside the think block; strip the tags only
                text = re.sub(r"</?think>", "", text).strip()

        # Strip markdown code fences
        if text.startswith("```"):
            lines = text.split("\n")
            if lines[0].strip().startswith("```"):
                lines = lines[1:]
            if lines and lines[-1].strip() == "```":
                lines = lines[:-1]
            text = "\n".join(lines).strip()

        # Find the JSON object boundary
        start = text.find("{")
        end = text.rfind("}")
        if start != -1 and end != -1 and end > start:
            text = text[start:end + 1]

        try:
            data = json.loads(text)
        except json.JSONDecodeError as e:
            print(f"[PPTXConverter] Failed to parse LLM JSON: {e}")
            print(f"[PPTXConverter] Raw text (first 800 chars): {text[:800]}")
            return None

        slides = data.get("slides") if isinstance(data, dict) else None
        if not isinstance(slides, list):
            print(f"[PPTXConverter] LLM response missing 'slides' key. data type={type(data).__name__}, keys={list(data.keys()) if isinstance(data, dict) else 'N/A'}")
            return None

        sanitized: List[Dict[str, Any]] = []
        for item in slides[:max_slides]:
            if not isinstance(item, dict):
                continue
            sanitized.append(self._normalize_slide(item))

        if sanitized:
            types_summary = ", ".join(s["type"] for s in sanitized)
            print(f"[PPTXConverter] LLM planned {len(sanitized)} slides: [{types_summary}]")
        else:
            print("[PPTXConverter] Sanitization produced 0 slides, using fallback.")
        return sanitized or None

    @staticmethod
    def _normalize_slide(raw: Dict[str, Any]) -> Dict[str, Any]:
        """Map LLM's free-form slide to our typed schema."""
        raw_type = str(raw.get("type", "content")).lower().replace("-", "_").replace(" ", "_")

        TYPE_MAP = {
            "title": "title_slide", "title_slide": "title_slide", "cover": "title_slide",
            "thank_you": "title_slide", "thankyou": "title_slide", "ending": "title_slide",
            "content": "content", "text": "content", "bullet": "content", "bullets": "content",
            "summary": "content", "overview": "content", "reference": "content", "quote": "highlight",
            "stats": "stats", "stat": "stats", "metric": "stats", "metrics": "stats",
            "infographic": "stats", "kpi": "stats",
            "table": "table", "data_table": "table", "comparison": "table",
            "two_column": "two_column", "columns": "two_column", "card_grid": "two_column",
            "cards": "two_column", "grid": "two_column",
            "highlight": "highlight", "callout": "highlight", "key_point": "highlight",
            "takeaway": "highlight",
            "references": "references", "reference": "references", "sources": "references",
            "bibliography": "references", "citations": "references",
        }
        stype = TYPE_MAP.get(raw_type, "content")

        title = str(raw.get("title", "") or "").strip()
        subtitle = str(raw.get("subtitle", "") or "").strip()
        accent = str(raw.get("accent_color", "") or raw.get("color", "") or PALETTE["primary"]).strip()

        # Extract items from whichever field the LLM used
        items = raw.get("items")
        if items is None:
            # Try common alternative field names
            for key in ("body", "bullets", "points", "content", "cards", "stats",
                        "quote", "data", "list", "sections"):
                if key in raw and raw[key]:
                    items = raw[key]
                    break

        # table: LLM may put headers/rows at top level
        if stype == "table" and not isinstance(items, dict):
            if "headers" in raw or "rows" in raw:
                items = {"headers": raw.get("headers", []), "rows": raw.get("rows", [])}

        # stats: if items is a list of dicts with value/label, keep; otherwise try "stats" field
        if stype == "stats" and isinstance(items, list) and items:
            if isinstance(items[0], dict) and ("value" in items[0] or "label" in items[0]):
                pass  # already correct
            elif isinstance(items[0], str):
                items = [{"value": s, "label": ""} for s in items]

        # two_column: if items is a list of dicts with title/points, keep
        if stype == "two_column" and isinstance(items, list) and items:
            normalized_cols = []
            for col in items:
                if isinstance(col, dict):
                    col_title = str(col.get("title", "") or col.get("heading", "") or "").strip()
                    col_points = col.get("points", col.get("items", col.get("bullets", [])))
                    if isinstance(col_points, str):
                        col_points = [col_points]
                    normalized_cols.append({"title": col_title, "points": col_points or []})
                elif isinstance(col, str):
                    normalized_cols.append({"title": col, "points": []})
            items = normalized_cols

        # highlight: ensure items is a dict with heading/body
        if stype == "highlight":
            if isinstance(items, str):
                items = {"heading": title, "body": items}
            elif isinstance(items, dict):
                items = {
                    "heading": str(items.get("heading", "") or items.get("title", "") or "").strip(),
                    "body": str(items.get("body", "") or items.get("text", "") or items.get("quote", "") or "").strip(),
                }
            elif isinstance(items, list):
                items = {"heading": title, "body": " ".join(str(x) for x in items)}
            else:
                # Check for quote/author pattern
                if "quote" in raw:
                    q = str(raw["quote"])
                    author = str(raw.get("author", ""))
                    items = {"heading": author, "body": q}
                else:
                    items = {"heading": title, "body": subtitle}

        # references: ensure items is a list of {text, url}
        if stype == "references":
            if isinstance(items, list):
                normalized_refs = []
                for ref in items:
                    if isinstance(ref, dict):
                        normalized_refs.append({
                            "text": str(ref.get("text", "") or "").strip(),
                            "url": str(ref.get("url", "") or ref.get("link", "") or ref.get("href", "") or "").strip(),
                        })
                    elif isinstance(ref, str):
                        normalized_refs.append({"text": ref, "url": ""})
                items = normalized_refs
            elif isinstance(items, dict):
                items = [{"text": str(v), "url": ""} for v in items.values() if v]
            else:
                items = []

        # content: ensure items is a list of strings
        if stype == "content":
            if isinstance(items, str):
                items = [items]
            elif isinstance(items, list):
                flat = []
                for x in items:
                    if isinstance(x, str):
                        flat.append(x)
                    elif isinstance(x, dict):
                        flat.append(str(x.get("text", "") or x.get("title", "") or x.get("point", "") or str(x)))
                items = flat
            elif isinstance(items, dict):
                items = [str(v) for v in items.values() if v]
            else:
                items = []
            # Also pull in note/highlight fields as extra bullets
            for extra_key in ("note", "highlight", "summary"):
                extra = raw.get(extra_key)
                if extra and isinstance(extra, str):
                    items.append(extra)

        if items is None:
            items = []

        return {
            "type": stype,
            "title": title,
            "subtitle": subtitle,
            "accent_color": accent,
            "items": items,
        }

    # ------------------------------------------------------------------
    # Fallback plan (no LLM)
    # ------------------------------------------------------------------

    def _fallback_plan(self, payload: Dict[str, Any], max_slides: int) -> List[Dict[str, Any]]:
        slides: List[Dict[str, Any]] = []

        slides.append({
            "type": "title_slide",
            "title": payload["title"],
            "subtitle": "",
            "accent_color": PALETTE["primary"],
            "items": [],
        })

        if payload.get("stats"):
            slides.append({
                "type": "stats",
                "title": "关键指标",
                "subtitle": "",
                "accent_color": PALETTE["accent"],
                "items": payload["stats"][:4],
            })

        if payload.get("cards"):
            cols = []
            for c in payload["cards"][:3]:
                cols.append({"title": c["title"], "points": c["items"][:4]})
            slides.append({
                "type": "two_column",
                "title": payload["headings"][1] if len(payload["headings"]) > 1 else "内容概览",
                "subtitle": "",
                "accent_color": PALETTE["secondary"],
                "items": cols,
            })

        if payload.get("tables"):
            tbl = payload["tables"][0]
            slides.append({
                "type": "table",
                "title": payload["headings"][2] if len(payload["headings"]) > 2 else "详细数据",
                "subtitle": "",
                "accent_color": PALETTE["danger"],
                "items": tbl,
            })

        para_idx = 0
        paras = payload.get("paragraphs", [])
        for i, heading in enumerate(payload.get("headings", [])[3:]):
            if len(slides) >= max_slides - 1:
                break
            bullets = []
            for _ in range(4):
                if para_idx < len(paras):
                    bullets.append(paras[para_idx][:160])
                    para_idx += 1
            slides.append({
                "type": "content",
                "title": heading[:80],
                "subtitle": "",
                "accent_color": PALETTE["primary"],
                "items": bullets,
            })

        if payload.get("highlights"):
            hl = payload["highlights"][0]
            slides.append({
                "type": "highlight",
                "title": hl.get("title", "要点"),
                "subtitle": "",
                "accent_color": PALETTE["warning"],
                "items": hl,
            })

        if payload.get("references"):
            slides.append({
                "type": "references",
                "title": "参考文献与来源",
                "subtitle": "",
                "accent_color": PALETTE["muted"],
                "items": payload["references"][:14],
            })

        return slides[:max_slides]

    # ------------------------------------------------------------------
    # Rendering
    # ------------------------------------------------------------------

    def _render_title_slide(self, sb: SlideBuilder, plan: Dict):
        slide = sb.add_blank()
        accent = plan.get("accent_color", PALETTE["primary"])
        sb.bg(slide, PALETTE["dark"])
        sb.rect(slide, 0, 0, 16, 0.15, fill=accent)
        sb.rrect(slide, -0.5, 7.5, 4, 2.5, fill=accent, radius=0.4, shadow=True)
        sb.rrect(slide, 13.5, -0.5, 3.5, 2.5, fill=accent, radius=0.4, shadow=True)
        sb.text(slide, 1, 2.8, 14, 1.5, plan["title"],
                font="Helvetica Neue", size=48, color="#FFFFFF", bold=True, align="center")
        if plan.get("subtitle"):
            sb.text(slide, 1, 4.8, 14, 1, plan["subtitle"],
                    font="Arial", size=26, color="#94A3B8", italic=True, align="center")

    def _render_content_slide(self, sb: SlideBuilder, plan: Dict, idx: int):
        slide = sb.add_blank()
        bg_color = SLIDE_BG_COLORS[idx % len(SLIDE_BG_COLORS)]
        accent = plan.get("accent_color", PALETTE["primary"])
        sb.bg(slide, bg_color)
        sb.rect(slide, 0, 0, 16, 1.3, fill=accent)
        sb.text(slide, 0.6, 0.25, 14, 0.9, plan["title"],
                font="Helvetica Neue", size=36, color="#FFFFFF", bold=True)

        items = plan.get("items", [])
        if not isinstance(items, list):
            items = []

        y = 1.8
        for bullet in items[:8]:
            if not isinstance(bullet, str):
                continue
            sb.rrect(slide, 0.6, y, 14.8, 0.75, fill="#FFFFFF", border=accent, border_w=0.5, radius=0.1, shadow=True)
            sb.text(slide, 0.9, y + 0.1, 14.2, 0.55, f"•  {bullet}",
                    font="Arial", size=18, color=PALETTE["dark"])
            y += 0.88

    def _render_stats_slide(self, sb: SlideBuilder, plan: Dict, idx: int):
        slide = sb.add_blank()
        accent = plan.get("accent_color", PALETTE["accent"])
        sb.bg(slide, "#F8FAFC")
        sb.rect(slide, 0, 0, 16, 1.3, fill=PALETTE["dark"])
        sb.text(slide, 0.6, 0.25, 14, 0.9, plan["title"],
                font="Helvetica Neue", size=36, color="#FFFFFF", bold=True)

        items = plan.get("items", [])
        if not isinstance(items, list):
            items = []
        n = min(len(items), 4)
        if n == 0:
            return

        card_w = min(3.5, (14.8 - (n - 1) * 0.4) / n)
        total_w = n * card_w + (n - 1) * 0.4
        start_x = (16 - total_w) / 2

        colors = [accent, PALETTE["primary"], PALETTE["warning"], PALETTE["danger"]]

        for i, stat in enumerate(items[:n]):
            if not isinstance(stat, dict):
                continue
            x = start_x + i * (card_w + 0.4)
            c = colors[i % len(colors)]
            sb.rrect(slide, x, 2.2, card_w, 4.5, fill=c, radius=0.25, shadow=True)
            value = str(stat.get("value", ""))
            label = str(stat.get("label", ""))
            sb.text(slide, x, 3.0, card_w, 1.5, value,
                    font="Helvetica Neue", size=36, color="#FFFFFF", bold=True, align="center", valign="middle")
            sb.text(slide, x, 5.0, card_w, 1.0, label,
                    font="Arial", size=18, color="#E2E8F0", align="center", valign="top")

    def _render_table_slide(self, sb: SlideBuilder, plan: Dict, idx: int):
        slide = sb.add_blank()
        accent = plan.get("accent_color", PALETTE["danger"])
        sb.bg(slide, "#FFFFFF")
        sb.rect(slide, 0, 0, 16, 1.3, fill=accent)
        sb.text(slide, 0.6, 0.25, 14, 0.9, plan["title"],
                font="Helvetica Neue", size=36, color="#FFFFFF", bold=True)

        tbl_data = plan.get("items", {})
        if not isinstance(tbl_data, dict):
            return
        headers = tbl_data.get("headers", [])
        rows = tbl_data.get("rows", [])
        if not headers and not rows:
            return

        n_cols = max(len(headers), max((len(r) for r in rows), default=0))
        if n_cols == 0:
            return
        n_rows = min(len(rows), 6)

        tbl_left, tbl_top = 0.6, 1.8
        tbl_w, row_h = 14.8, 0.8
        col_w = tbl_w / n_cols

        if headers:
            for ci, hdr in enumerate(headers[:n_cols]):
                x = tbl_left + ci * col_w
                sb.rect(slide, x, tbl_top, col_w, row_h, fill=PALETTE["secondary"], border="#FFFFFF", border_w=0.5)
                sb.text(slide, x + 0.15, tbl_top + 0.1, col_w - 0.3, row_h - 0.2, str(hdr),
                        font="Arial", size=16, color="#FFFFFF", bold=True, valign="middle")
            tbl_top += row_h

        for ri, row in enumerate(rows[:n_rows]):
            fill = "#F8FAFC" if ri % 2 == 0 else "#FFFFFF"
            for ci in range(n_cols):
                x = tbl_left + ci * col_w
                cell_text = str(row[ci]) if ci < len(row) else ""
                sb.rect(slide, x, tbl_top, col_w, row_h, fill=fill, border="#E2E8F0", border_w=0.5)
                sb.text(slide, x + 0.15, tbl_top + 0.05, col_w - 0.3, row_h - 0.1, cell_text,
                        font="Arial", size=14, color=PALETTE["dark"], valign="middle")
            tbl_top += row_h

    def _render_two_column_slide(self, sb: SlideBuilder, plan: Dict, idx: int):
        slide = sb.add_blank()
        accent = plan.get("accent_color", PALETTE["secondary"])
        sb.bg(slide, "#F8FAFC")
        sb.rect(slide, 0, 0, 16, 1.3, fill=accent)
        sb.text(slide, 0.6, 0.25, 14, 0.9, plan["title"],
                font="Helvetica Neue", size=36, color="#FFFFFF", bold=True)

        items = plan.get("items", [])
        if not isinstance(items, list):
            items = []
        n = min(len(items), 3)
        if n == 0:
            return

        card_w = (14.8 - (n - 1) * 0.4) / n
        accent_variants = [accent, PALETTE["accent"], PALETTE["warning"]]

        for i, col in enumerate(items[:n]):
            if not isinstance(col, dict):
                continue
            x = 0.6 + i * (card_w + 0.4)
            c = accent_variants[i % len(accent_variants)]

            sb.rrect(slide, x, 1.7, card_w, 6.8, fill="#FFFFFF", border=c, border_w=2, radius=0.2, shadow=True)
            sb.rect(slide, x, 1.7, card_w, 0.12, fill=c)

            col_title = str(col.get("title", ""))
            sb.text(slide, x + 0.3, 2.0, card_w - 0.6, 0.7, col_title,
                    font="Helvetica Neue", size=20, color=c, bold=True)

            points = col.get("points", [])
            if not isinstance(points, list):
                points = []
            lines = []
            for pt in points[:6]:
                lines.append({"text": f"•  {str(pt)[:120]}", "size": 14, "color": PALETTE["dark"], "space_after": 6})
            if lines:
                sb.multiline(slide, x + 0.3, 2.8, card_w - 0.6, 5.2, lines)

    def _render_highlight_slide(self, sb: SlideBuilder, plan: Dict, idx: int):
        slide = sb.add_blank()
        accent = plan.get("accent_color", PALETTE["warning"])
        sb.bg(slide, PALETTE["dark"])

        sb.rrect(slide, 1.5, 1.5, 13, 6, fill="#FFFFFF", radius=0.3, shadow=True)
        sb.rect(slide, 1.5, 1.5, 0.2, 6, fill=accent)

        heading = plan.get("title", "")
        body = ""
        items = plan.get("items", {})
        if isinstance(items, dict):
            heading = items.get("heading", heading) or heading
            body = items.get("body", "")
        elif isinstance(items, str):
            body = items

        sb.text(slide, 2.5, 2.2, 11, 1.2, heading,
                font="Helvetica Neue", size=36, color=accent, bold=True)
        sb.text(slide, 2.5, 3.8, 11, 3.2, body,
                font="Arial", size=22, color=PALETTE["dark"])

    def _render_references_slide(self, sb: SlideBuilder, plan: Dict, idx: int):
        slide = sb.add_blank()
        accent = plan.get("accent_color", PALETTE["muted"])
        sb.bg(slide, "#F8FAFC")
        sb.rect(slide, 0, 0, 16, 1.3, fill=PALETTE["secondary"])
        sb.text(slide, 0.6, 0.25, 14, 0.9, plan.get("title") or "参考文献与来源",
                font="Helvetica Neue", size=36, color="#FFFFFF", bold=True)

        items = plan.get("items", [])
        if not isinstance(items, list):
            items = []

        y = 1.7
        col_w = 7.2
        for i, ref in enumerate(items[:14]):
            if not isinstance(ref, dict):
                continue
            text = str(ref.get("text", "")).strip()
            url = str(ref.get("url", "")).strip()
            if not text:
                continue
            text = re.sub(r"^\[\d+\]\s*", "", text)

            col = 0 if i < 7 else 1
            row = i if i < 7 else i - 7
            x = 0.6 + col * (col_w + 0.4)
            ry = y + row * 0.95

            sb.rrect(slide, x, ry, col_w, 0.8, fill="#FFFFFF", border="#E2E8F0", border_w=0.5, radius=0.08, shadow=True)

            ref_num = f"[{i + 1}]  "
            display_text = ref_num + text[:80]

            tb = sb.text(slide, x + 0.2, ry + 0.1, col_w - 0.4, 0.6, display_text,
                         font="Arial", size=12, color=PALETTE["dark"])

            if url:
                tf = tb.text_frame
                p = tf.paragraphs[0]
                p.clear()
                run_num = p.add_run()
                run_num.text = ref_num
                run_num.font.size = Pt(12)
                run_num.font.color.rgb = _hex(PALETTE["primary"])
                run_num.font.bold = True

                run_link = p.add_run()
                run_link.text = text[:70]
                run_link.font.size = Pt(11)
                run_link.font.color.rgb = _hex(PALETTE["primary"])
                run_link.hyperlink.address = url

    def _render_slide(self, sb: SlideBuilder, plan: Dict, idx: int):
        stype = plan.get("type", "content")
        dispatch = {
            "title_slide": self._render_title_slide,
            "content": self._render_content_slide,
            "stats": self._render_stats_slide,
            "table": self._render_table_slide,
            "two_column": self._render_two_column_slide,
            "highlight": self._render_highlight_slide,
            "references": self._render_references_slide,
        }
        renderer = dispatch.get(stype, self._render_content_slide)
        if stype == "title_slide":
            renderer(sb, plan)
        else:
            renderer(sb, plan, idx)

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def generate_pptx(self, html_path: Path, output_path: Path, max_slides: int = 10) -> str:
        payload = _extract_rich_payload(html_path)
        plan = self._plan_slides_with_llm(payload, max_slides=max_slides) or self._fallback_plan(payload, max_slides=max_slides)

        # Ensure references page exists if HTML has references
        has_ref_slide = any(s.get("type") == "references" for s in plan)
        if not has_ref_slide and payload.get("references"):
            plan.append({
                "type": "references",
                "title": "参考文献与来源",
                "subtitle": "",
                "accent_color": PALETTE["muted"],
                "items": payload["references"][:14],
            })

        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)

        sb = SlideBuilder(prs)
        for i, slide_plan in enumerate(plan):
            self._render_slide(sb, slide_plan, i)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        return str(output_path)
