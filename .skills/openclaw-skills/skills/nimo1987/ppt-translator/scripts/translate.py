#!/usr/bin/env python3
"""
PPT Translator - Agent Loop Core
用法:
  python3 translate.py --input file.pptx --output out.pptx --translations '{"中文": "English"}' [--max-iter 5]

检测溢出：本脚本只做渲染+写回，溢出检测由外部调用（Agent Vision）完成。
完整 loop 由 Agent 驱动：写PPT → 渲染 → Vision检测 → 缩字号 → 重复。
"""

import argparse
import json
import subprocess
import sys
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt


def get_para_text(para):
    return "".join(run.text for run in para.runs).strip()


def write_translated_pptx(input_path: str, output_path: str,
                           translations: dict, scale: float = 1.0):
    """
    Write translated PPTX with optional font scale.
    translations: {original_text: translated_text}
    scale: font size multiplier (1.0 = no change, 0.85 = shrink 15%)
    """
    prs = Presentation(input_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            tf.word_wrap = True

            for para in tf.paragraphs:
                orig = get_para_text(para)
                if not orig or orig not in translations:
                    continue
                translated = translations[orig]
                if not para.runs:
                    continue

                # Get original font size
                orig_pt = next((r.font.size.pt for r in para.runs if r.font.size), 12.0)
                new_pt = max(6.0, round(orig_pt * scale, 1))

                # Apply: first run gets full text, others cleared
                para.runs[0].text = translated
                para.runs[0].font.size = Pt(new_pt)
                for run in para.runs[1:]:
                    run.text = ""

    prs.save(output_path)
    return output_path


def render_to_png(pptx_path: str, output_dir: str) -> str:
    """Render PPTX to PNG using LibreOffice. Returns PNG path."""
    result = subprocess.run(
        ['libreoffice', '--headless', '--convert-to', 'png',
         '--outdir', output_dir, pptx_path],
        capture_output=True, text=True, timeout=60
    )
    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice failed: {result.stderr}")

    stem = Path(pptx_path).stem
    png_path = Path(output_dir) / f"{stem}.png"
    if not png_path.exists():
        pngs = sorted(Path(output_dir).glob(f"{stem}*.png"))
        if not pngs:
            raise RuntimeError(f"No PNG generated in {output_dir}")
        png_path = pngs[0]
    return str(png_path)


def extract_texts(input_path: str) -> list:
    """Extract all non-empty paragraph texts from PPTX."""
    prs = Presentation(input_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = get_para_text(para)
                    if t:
                        texts.append(t)
    return list(dict.fromkeys(texts))  # deduplicated


def main():
    parser = argparse.ArgumentParser(description="PPT Translator Core")
    parser.add_argument("--input", required=True, help="Input PPTX path")
    parser.add_argument("--output", required=False, help="Output PPTX path")
    parser.add_argument("--translations", required=False,
                        help="JSON string: {original: translated}")
    parser.add_argument("--scale", type=float, default=1.0,
                        help="Font scale factor (default: 1.0)")
    parser.add_argument("--render", action="store_true",
                        help="Also render to PNG after writing")
    parser.add_argument("--extract", action="store_true",
                        help="Just extract texts and print as JSON, then exit")
    args = parser.parse_args()

    if args.extract:
        texts = extract_texts(args.input)
        print(json.dumps(texts, ensure_ascii=False, indent=2))
        return

    if not args.translations:
        print("Error: --translations required unless --extract", file=sys.stderr)
        sys.exit(1)

    translations = json.loads(args.translations)

    output_path = write_translated_pptx(
        args.input, args.output, translations, args.scale
    )
    print(json.dumps({"status": "ok", "output": output_path, "scale": args.scale}))

    if args.render:
        output_dir = str(Path(args.output).parent)
        png_path = render_to_png(output_path, output_dir)
        print(json.dumps({"status": "rendered", "png": png_path}))


if __name__ == "__main__":
    main()
