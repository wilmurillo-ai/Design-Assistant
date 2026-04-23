#!/usr/bin/env python3
"""
标准化PPT报告生成脚本 - v4.2 优化版

严格按照以下结构生成PPT：
1. 封面
2. 目录
3. 用户画像总结
4. 地域深耕方向（如有）
5-N. 各维度详细分析（左右分栏+图形元素）
倒数第2页. 综合策略建议
最后一页. 深度分析方向

🎨 设计优化 v4.2：
- 文字大小优化（13-16pt）
- 左右分栏布局（充分利用空间）
- 图形元素：柱状图、色块、icon、标签
- 数据突出：加粗、高亮背景、颜色强调
- 零空白设计（内容铺满整页）

使用方式：
  python generate_ppt_report.py path/to/report.md
"""

import re
import sys
from pathlib import Path
from typing import Dict, List, Optional, Any, Tuple

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("❌ 缺少 python-pptx，请运行：pip install python-pptx")
    sys.exit(1)


class PPTStandardGenerator:
    """标准化PPT生成器 v4.2 - 优化版"""
    
    COLORS = {
        'primary': RGBColor(0, 112, 192),        # 深蓝 - 主色
        'primary_dark': RGBColor(6, 90, 130),   # 更深的蓝
        'accent': RGBColor(255, 107, 107),      # 红色 - 强调
        'accent_orange': RGBColor(255, 152, 0), # 橙色 - 辅助强调
        'text_dark': RGBColor(47, 60, 126),     # 深蓝灰 - 正文
        'text_light': RGBColor(100, 100, 100),  # 浅灰 - 辅助
        'white': RGBColor(255, 255, 255),       # 白色
        'bg_light': RGBColor(245, 245, 245),    # 浅灰背景
        'bg_lighter': RGBColor(230, 240, 250),  # 更浅的蓝背景
        'bg_dark': RGBColor(6, 90, 130),        # 深蓝背景
        'highlight_bg': RGBColor(255, 243, 224), # 浅橙背景（突出）
        'chart_blue': RGBColor(0, 176, 240),    # 图表蓝
        'chart_orange': RGBColor(255, 192, 0),  # 图表橙
        'chart_green': RGBColor(112, 173, 71),  # 图表绿
    }
    
    def __init__(self, output_dir: str = "reports"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(exist_ok=True)
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        self.slide_count = 0
        self.parsed_data = {}
    
    def _add_slide(self):
        """添加空白幻灯片"""
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])
    
    def _add_shape_rectangle(self, slide, x, y, width, height, color):
        """添加矩形"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color
        return shape
    
    def _add_title(self, slide, text: str, x=0.5, y=0.4, width=9, size=44, color=None):
        """添加标题"""
        if color is None:
            color = self.COLORS['primary']
        
        textbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(1.2))
        frame = textbox.text_frame
        frame.word_wrap = True
        p = frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.LEFT
        return textbox
    
    def _add_subtitle(self, slide, text: str, x=0.5, y=1.3, width=9, size=24):
        """添加副标题"""
        textbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(0.6))
        frame = textbox.text_frame
        frame.word_wrap = True
        p = frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = self.COLORS['accent_orange']
        p.alignment = PP_ALIGN.LEFT
        return textbox
    
    def _add_text(self, slide, text: str, x=0.5, y=1.5, width=9, size=12, color=None):
        """添加文本"""
        if color is None:
            color = self.COLORS['text_dark']
        
        textbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(5.5))
        frame = textbox.text_frame
        frame.word_wrap = True
        p = frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        return textbox
    
    def _add_divider(self, slide, y=1.35):
        """添加分隔线"""
        self._add_shape_rectangle(slide, 0.5, y, 9, 0.02, self.COLORS['accent'])
    
    def _add_bullet_list(self, slide, items: List[str], x=0.5, y=1.5, width=9, size=12, bold=False, color=None):
        """添加项目列表"""
        if color is None:
            color = self.COLORS['text_dark']
        
        textbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(5.5))
        frame = textbox.text_frame
        frame.word_wrap = True
        
        for i, item in enumerate(items):
            if i == 0:
                p = frame.paragraphs[0]
            else:
                p = frame.add_paragraph()
            
            p.text = "• " + item  # 添加项目符号
            p.font.size = Pt(size)
            p.font.bold = bold
            p.font.color.rgb = color
            p.space_before = Pt(4)
            p.space_after = Pt(4)
    
    def _add_shape_circle(self, slide, x, y, diameter, color):
        """添加圆形 - 用于icon"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x), Inches(y), Inches(diameter), Inches(diameter)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color
        return shape
    
    def _add_colored_tag(self, slide, text: str, x: float, y: float, width=1.5, bg_color=None, text_color=None):
        """添加彩色标签（用于突出数据）"""
        if bg_color is None:
            bg_color = self.COLORS['accent_orange']
        if text_color is None:
            text_color = self.COLORS['white']
        
        # 背景矩形
        self._add_shape_rectangle(slide, x, y, width, 0.4, bg_color)
        
        # 标签文字
        textbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(0.4))
        frame = textbox.text_frame
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER
    
    def _add_chart_bars(self, slide, x: float, y: float, items: List[Tuple[str, float]], max_width=2.5, bar_height=0.2):
        """添加简单柱状图
        
        Args:
            slide: 幻灯片对象
            x, y: 图表左上角坐标
            items: [(标签, 值), ...] 的列表
            max_width: 最大柱子宽度（英寸）
            bar_height: 每根柱子的高度（英寸）
        """
        if not items:
            return
        
        max_value = max(v for _, v in items)
        if max_value == 0:
            return
        
        for i, (label, value) in enumerate(items[:5]):  # 最多显示5条
            bar_y = y + i * (bar_height + 0.08)
            
            # 标签
            label_box = slide.shapes.add_textbox(Inches(x), Inches(bar_y), Inches(0.8), Inches(bar_height))
            frame = label_box.text_frame
            frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            frame.word_wrap = False
            p = frame.paragraphs[0]
            p.text = label[:6]
            p.font.size = Pt(8)
            p.font.color.rgb = self.COLORS['text_dark']
            
            # 柱子宽度计算
            bar_width = (value / max_value) * max_width
            
            # 绘制柱子
            self._add_shape_rectangle(slide, x + 0.85, bar_y + 0.02, bar_width, bar_height - 0.04, self.COLORS['chart_blue'])
            
            # 数值标签
            value_text = f"{value:.1f}"
            value_box = slide.shapes.add_textbox(Inches(x + 0.85 + bar_width + 0.05), Inches(bar_y), Inches(0.4), Inches(bar_height))
            frame = value_box.text_frame
            frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = frame.paragraphs[0]
            p.text = value_text
            p.font.size = Pt(7)
            p.font.bold = True
            p.font.color.rgb = self.COLORS['accent']

    def _add_two_column_layout(self, slide, left_content: str, right_content: str, title: str = "", left_title: str = "核心特征", right_title: str = "数据分析"):
        """添加左右分栏布局"""
        # 页面标题
        if title:
            self._add_title(slide, title)
            self._add_divider(slide)
            y_start = 1.5
        else:
            y_start = 0.5
        
        # 左列
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_start), Inches(4.5), Inches(5.5))
        frame = left_box.text_frame
        frame.word_wrap = True
        
        p = frame.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['primary']
        
        p = frame.add_paragraph()
        p.text = left_content
        p.font.size = Pt(12)
        p.font.color.rgb = self.COLORS['text_dark']
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        
        # 右列
        right_box = slide.shapes.add_textbox(Inches(5.2), Inches(y_start), Inches(4.3), Inches(5.5))
        frame = right_box.text_frame
        frame.word_wrap = True
        
        p = frame.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['primary']
        
        p = frame.add_paragraph()
        p.text = right_content
        p.font.size = Pt(12)
        p.font.color.rgb = self.COLORS['text_dark']
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    def _parse_markdown(self, filepath: str) -> Dict[str, Any]:
        """完整解析Markdown文件"""
        with open(filepath, 'r', encoding='utf-8') as f:
            content = f.read()
        
        result = {
            'title': '',
            'persona_name': '',
            'dimension_count': 0,
            'user_profile': {},
            'geography': {},
            'dimensions': [],
            'strategy': {},
            'deep_analysis': {}
        }
        
        # 提取标题和人群名称
        title_match = re.search(r'^# (.+?) 人群洞察分析报告', content, re.MULTILINE)
        if title_match:
            result['persona_name'] = title_match.group(1)
            result['title'] = title_match.group(1)
        
        # 提取用户画像部分
        profile_section = re.search(r'## 一、用户画像\n(.*?)(?=## 二、)', content, re.DOTALL)
        if profile_section:
            profile_text = profile_section.group(1)
            
            # 提取一句话核心特征（支持多种格式）
            # 格式1: ### 一句话核心特征\n
            # 格式2: 一句话核心特征：xxx
            feature_match = re.search(r'[#]*\s*一句话核心特征[：:](.*?)(?=\n[a-z]|提炼|###|\Z)', profile_text, re.DOTALL)
            if feature_match:
                result['user_profile']['feature'] = feature_match.group(1).strip()
            
            # 提取核心定位标签（支持多种格式）
            # 格式1: ### 核心定位标签\n
            # 格式2: 提炼核心定位：xxx
            position_match = re.search(r'[#]*\s*(?:提炼)?核心定位(?:标签)?[：:](.*?)(?=\n###|\Z)', profile_text, re.DOTALL)
            if position_match:
                result['user_profile']['position'] = position_match.group(1).strip()
        
        # 提取地域深耕方向
        geo_section = re.search(r'### 地域深耕方向\n(.*?)(?=\n### 【|\n---|\Z)', content, re.DOTALL)
        if geo_section:
            geo_text = geo_section.group(1)
            result['geography']['text'] = geo_text.strip()
            
            # 提取核心城市
            core_match = re.search(r'\*\*核心城市\*\*.*?\n(.+?)(?=\n\*\*|\Z)', geo_text, re.DOTALL)
            if core_match:
                result['geography']['core'] = core_match.group(1).strip()
            
            # 提取潜力城市
            potential_match = re.search(r'\*\*潜力(.+?)\*\*.*?\n(.+?)(?=\n\*\*|\Z)', geo_text, re.DOTALL)
            if potential_match:
                result['geography']['potential'] = potential_match.group(2).strip()
        
        # 提取各维度分析（支持两种格式）
        # 格式1: ### 【维度名】维度
        # 格式2: ### 维度名维度
        dimensions = []
        
        # 先尝试格式1
        dims_1 = re.findall(r'### 【(.+?)】维度\n(.*?)(?=###|\Z)', content, re.DOTALL)
        if dims_1:
            dimensions = dims_1
        
        # 如果格式1没有找到，尝试格式2
        if not dimensions:
            dims_2 = re.findall(r'### (.+?维度)\n(.*?)(?=###|\Z)', content, re.DOTALL)
            # 清理格式2的维度名（去掉"维度"后缀后半部分）
            dimensions = [(name.replace('维度', ''), content) for name, content in dims_2]
        
        for dim_name, dim_content in dimensions:
            dim_data = {
                'name': dim_name,
                'features': [],
                'tgi_avg': 0,
                'tgi_max': 0,
                'reason': '',
                'distribution': ''
            }
            
            # 提取特征列表（支持两种格式）
            # 格式1: 表格格式 (| 排名 | 特征 |...)
            table_match = re.search(r'\| 排名 \| 特征.*?\n\|.*?\n(.*?)(?=\n\n|\*\*|###|\Z)', dim_content, re.DOTALL)
            if table_match:
                lines = table_match.group(1).strip().split('\n')
                for line in lines:
                    if '|' in line and '排名' not in line:
                        parts = [p.strip() for p in line.split('|')[1:-1]]
                        if len(parts) >= 3:
                            feature = {
                                'rank': parts[0],
                                'name': parts[1],
                                'tgi': parts[2] if len(parts) > 2 else '',
                                'proportion': parts[3] if len(parts) > 3 else ''
                            }
                            dim_data['features'].append(feature)
            
            # 格式2: 列表格式 (1. 特征名（TGI xxx）)
            if not dim_data['features']:
                list_matches = re.findall(r'\s+(\d+)\.\s+(.+?)\（TGI\s+([\d.]+)\）', dim_content)
                for rank, name, tgi in list_matches[:10]:  # 最多取前10个
                    feature = {
                        'rank': rank,
                        'name': name.strip(),
                        'tgi': tgi.strip(),
                        'proportion': ''
                    }
                    dim_data['features'].append(feature)
            
            # 提取统计数据（最高TGI）
            # 从特征列表中自动提取最高TGI
            if dim_data['features']:
                tgi_values = [float(f['tgi']) for f in dim_data['features'] if f['tgi']]
                if tgi_values:
                    dim_data['tgi_max'] = max(tgi_values)
                    dim_data['tgi_avg'] = sum(tgi_values) / len(tgi_values)
            
            # 也可以从"平均TGI"文本中提取
            stats_match = re.search(r'平均TGI[：:]\s*([\d.]+)', dim_content)
            if stats_match:
                dim_data['tgi_avg'] = float(stats_match.group(1))
            
            # 提取判断理由（支持两种格式）
            reason_match = re.search(r'\*\*判断理由[：:]*\*\*\n(.*?)(?=\n\n|\*\*|###|\Z)', dim_content, re.DOTALL)
            if reason_match:
                dim_data['reason'] = reason_match.group(1).strip()
            
            # 提取数据分布（支持多种格式）
            distribution_match = re.search(r'\*\*数据分布[情况]*说明[：:]*\*\*\n(.*?)(?=\n\n|\*\*|###|\Z)', dim_content, re.DOTALL)
            if distribution_match:
                dim_data['distribution'] = distribution_match.group(1).strip()
            
            result['dimensions'].append(dim_data)
        
        result['dimension_count'] = len(result['dimensions'])
        
        # 提取综合策略建议（支持两种格式）
        strategy_section = re.search(r'## 三、综合策略建议\n(.*?)(?=## 四、|## 五、|\Z)', content, re.DOTALL)
        if strategy_section:
            strategy_text = strategy_section.group(1)
            
            # 格式1: ### 【产品定位】
            # 格式2: 【产品定位】
            
            # 产品定位
            pos_match = re.search(r'[#]*\s*【产品定位[】:：]*(.*?)(?=【|##|\Z)', strategy_text, re.DOTALL)
            if pos_match:
                result['strategy']['position'] = pos_match.group(1).strip()
            
            # 渠道布局
            chan_match = re.search(r'【渠道布局[^】]*[】:：]*(.*?)(?=【|##|\Z)', strategy_text, re.DOTALL)
            if chan_match:
                result['strategy']['channel'] = chan_match.group(1).strip()
            
            # 内容营销
            cont_match = re.search(r'【内容营销[^】]*[】:：]*(.*?)(?=【|##|\Z)', strategy_text, re.DOTALL)
            if cont_match:
                result['strategy']['content'] = cont_match.group(1).strip()
        
        # 提取深度分析建议
        analysis_section = re.search(r'## 四、深度分析建议\n(.*?)(?=\n---|\Z)', content, re.DOTALL)
        if analysis_section:
            analysis_text = analysis_section.group(1)
            
            # 单维度分析：提取从"### 单维度"到"### 交叉"之间的所有内容
            single_match = re.search(r'###\s*单维度[^\n]*\n(.*?)(?=###|\Z)', analysis_text, re.DOTALL)
            if single_match:
                single_text = single_match.group(1).strip()
                # 整理：保留所有列表项，格式化为易读的文本
                # 将 "- **标签**：内容" 转换为 "标签：内容"
                single_text = re.sub(r'^-\s+', '', single_text, flags=re.MULTILINE)
                result['deep_analysis']['single'] = single_text
            
            # 交叉维度分析：提取从"### 交叉"到"### 优先级"之间的所有内容
            cross_match = re.search(r'###\s*交叉[^\n]*\n(.*?)(?=###|\Z)', analysis_text, re.DOTALL)
            if cross_match:
                cross_text = cross_match.group(1).strip()
                # 整理：保留所有列表项
                cross_text = re.sub(r'^-\s+', '', cross_text, flags=re.MULTILINE)
                result['deep_analysis']['cross'] = cross_text
        
        return result
    
    def _create_cover_slide(self, persona_name: str, dimension_count: int):
        """P1: 封面页"""
        slide = self._add_slide()
        self.slide_count += 1
        
        # 深蓝背景
        self._add_shape_rectangle(slide, 0, 0, 10, 7.5, self.COLORS['bg_dark'])
        
        # 主标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        frame = title_box.text_frame
        frame.word_wrap = True
        p = frame.paragraphs[0]
        p.text = "DMP用户画像洞察分析报告"
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        
        # 副标题
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
        frame = subtitle_box.text_frame
        p = frame.paragraphs[0]
        p.text = f"基于{dimension_count}个维度的深度分析"
        p.font.size = Pt(28)
        p.font.color.rgb = self.COLORS['accent_orange']
        p.alignment = PP_ALIGN.CENTER
        
        # 人群名称
        name_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(0.8))
        frame = name_box.text_frame
        p = frame.paragraphs[0]
        p.text = persona_name
        p.font.size = Pt(32)
        p.font.color.rgb = self.COLORS['white']
        p.alignment = PP_ALIGN.CENTER
    
    def _create_contents_slide(self, dimension_count: int, has_geography: bool):
        """P2: 目录页 - 优化版"""
        slide = self._add_slide()
        self.slide_count += 1
        
        # 添加侧边装饰条
        self._add_shape_rectangle(slide, 0, 0, 0.08, 7.5, self.COLORS['primary'])
        
        self._add_title(slide, "内容导航")
        self._add_divider(slide)
        
        y = 1.6
        items = [
            ("用户画像总结", "核心特征 × 定位标签"),
        ]
        
        if has_geography:
            items.append(("地域深耕方向", "优先布局 × 潜力城市"))
            items.append((f"各维度详细分析", f"{dimension_count}个维度深度剖析"))
            items.append(("综合策略建议", "产品 × 渠道 × 内容"))
            items.append(("深度分析方向", "单维度 × 交叉维度"))
        else:
            items.append((f"各维度详细分析", f"{dimension_count}个维度深度剖析"))
            items.append(("综合策略建议", "产品 × 渠道 × 内容"))
            items.append(("深度分析方向", "单维度 × 交叉维度"))
        
        colors_cycle = [self.COLORS['accent'], self.COLORS['accent_orange'], self.COLORS['primary'], 
                        self.COLORS['chart_green'], self.COLORS['chart_blue']]
        
        for idx, (title, subtitle) in enumerate(items):
            color = colors_cycle[idx % len(colors_cycle)]
            
            # 背景色块
            self._add_shape_rectangle(slide, 0.5, y - 0.05, 9, 0.9, self.COLORS['bg_light'])
            
            # 序号 circle
            self._add_shape_circle(slide, 0.5, y, 0.4, color)
            
            # 标题
            title_box = slide.shapes.add_textbox(Inches(1), Inches(y + 0.05), Inches(4), Inches(0.3))
            frame = title_box.text_frame
            frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = frame.paragraphs[0]
            p.text = f"{idx + 1}. {title}"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = color
            
            # 副标题
            subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(y + 0.35), Inches(4), Inches(0.3))
            frame = subtitle_box.text_frame
            p = frame.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(11)
            p.font.color.rgb = self.COLORS['text_light']
            
            y += 1.0
    
    def _create_user_profile_slide(self, profile: Dict):
        """P3: 用户画像总结 - 优化版（完整显示所有内容）"""
        slide = self._add_slide()
        self.slide_count += 1
        
        # 添加背景色块（左侧装饰）
        self._add_shape_rectangle(slide, 0, 0, 0.08, 7.5, self.COLORS['primary'])
        
        self._add_title(slide, "用户画像总结")
        self._add_divider(slide)
        
        y = 1.5
        
        # 一句话核心特征 - 突出显示
        if 'feature' in profile and profile['feature']:
            feature_text = profile['feature']
            # 背景色块
            self._add_shape_rectangle(slide, 0.5, y - 0.1, 9, 2.0, self.COLORS['bg_lighter'])
            
            # 标题
            label_box = slide.shapes.add_textbox(Inches(0.5), Inches(y - 0.05), Inches(1.5), Inches(0.3))
            frame = label_box.text_frame
            p = frame.paragraphs[0]
            p.text = "核心特征"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = self.COLORS['primary']
            
            # 内容（完整显示，不截断）
            feature_box = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.25), Inches(8.6), Inches(1.65))
            frame = feature_box.text_frame
            frame.word_wrap = True
            p = frame.paragraphs[0]
            p.text = feature_text
            p.font.size = Pt(11)
            p.font.color.rgb = self.COLORS['text_dark']
            
            y += 2.2
        
        # 核心定位标签 - 高亮显示
        if 'position' in profile and profile['position']:
            # 背景色块
            self._add_shape_rectangle(slide, 0.5, y - 0.1, 9, 1.8, self.COLORS['highlight_bg'])
            
            # 标题
            label_box = slide.shapes.add_textbox(Inches(0.5), Inches(y - 0.05), Inches(1.5), Inches(0.3))
            frame = label_box.text_frame
            p = frame.paragraphs[0]
            p.text = "定位标签"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = self.COLORS['accent_orange']
            
            # 内容（完整显示）
            position_text = profile['position']
            position_box = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.25), Inches(8.6), Inches(1.45))
            frame = position_box.text_frame
            frame.word_wrap = True
            p = frame.paragraphs[0]
            p.text = position_text
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = self.COLORS['accent']
    
    def _create_geography_slide(self, geography: Dict):
        """P4: 地域深耕方向 - 优化版"""
        if not geography:
            return
        
        slide = self._add_slide()
        self.slide_count += 1
        
        # 添加侧边装饰条
        self._add_shape_rectangle(slide, 0, 0, 0.08, 7.5, self.COLORS['primary'])
        
        self._add_title(slide, "地域深耕方向")
        self._add_divider(slide)
        
        y = 1.5
        
        # 核心布局城市
        if 'core' in geography:
            # 背景色块
            self._add_shape_rectangle(slide, 0.5, y - 0.1, 9, 2, self.COLORS['bg_lighter'])
            
            # 标题
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(2), Inches(0.35))
            frame = title_box.text_frame
            p = frame.paragraphs[0]
            p.text = "🎯 优先布局城市"
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = self.COLORS['primary']
            
            # 内容
            core_text = geography['core'][:180]
            core_box = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.4), Inches(8.8), Inches(1.5))
            frame = core_box.text_frame
            frame.word_wrap = True
            p = frame.paragraphs[0]
            p.text = core_text
            p.font.size = Pt(12)
            p.font.color.rgb = self.COLORS['text_dark']
            
            y += 2.2
        
        # 潜力市场
        if 'potential' in geography:
            # 背景色块
            self._add_shape_rectangle(slide, 0.5, y - 0.1, 9, 2, self.COLORS['highlight_bg'])
            
            # 标题
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(2), Inches(0.35))
            frame = title_box.text_frame
            p = frame.paragraphs[0]
            p.text = "📈 潜力城市"
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = self.COLORS['accent_orange']
            
            # 内容
            potential_text = geography['potential'][:180]
            potential_box = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.4), Inches(8.8), Inches(1.5))
            frame = potential_box.text_frame
            frame.word_wrap = True
            p = frame.paragraphs[0]
            p.text = potential_text
            p.font.size = Pt(12)
            p.font.color.rgb = self.COLORS['text_dark']
    
    def _create_dimension_slides(self, dimensions: List[Dict]):
        """P5-N: 各维度详细分析 - 优化版（左右分栏+图形元素）"""
        for dim in dimensions:
            slide = self._add_slide()
            self.slide_count += 1
            
            # 添加侧边装饰条
            self._add_shape_rectangle(slide, 0, 0, 0.08, 7.5, self.COLORS['primary'])
            
            self._add_title(slide, f"{dim['name']}维度分析", x=0.5)
            self._add_divider(slide)
            
            y_start = 1.5
            
            # 左列：核心特征列表
            if dim['features']:
                # 先添加背景色块（确保在文字下面）
                self._add_shape_rectangle(slide, 0.5, y_start + 0.35, 4.8, 4.7, self.COLORS['bg_light'])
                
                left_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_start), Inches(4.8), Inches(5.5))
                frame = left_box.text_frame
                frame.word_wrap = True
                
                p = frame.paragraphs[0]
                p.text = "Top 特征"
                p.font.size = Pt(13)
                p.font.bold = True
                p.font.color.rgb = self.COLORS['primary']
                p.space_after = Pt(8)
                
                for i, feature in enumerate(dim['features'][:8]):
                    if i == 0:
                        p = frame.add_paragraph()
                    else:
                        p = frame.add_paragraph()
                    
                    feature_name = feature['name'][:15]
                    tgi_val = feature['tgi'] if feature['tgi'] else '-'
                    
                    p.text = f"{i+1}. {feature_name}"
                    p.font.size = Pt(12)
                    p.font.bold = True if i < 3 else False  # 前3个加粗
                    p.font.color.rgb = self.COLORS['accent'] if i < 3 else self.COLORS['text_dark']
                    p.space_before = Pt(4)
                    p.space_after = Pt(4)
                    
                    # TGI值
                    p = frame.add_paragraph()
                    p.text = f"    TGI: {tgi_val}"
                    p.font.size = Pt(11)
                    p.font.color.rgb = self.COLORS['text_light']
                    p.space_after = Pt(6)
            
            # 右列：数据分析
            # 先添加背景色块（确保在文字下面）
            self._add_shape_rectangle(slide, 5.3, y_start + 0.35, 4.2, 4.7, self.COLORS['bg_lighter'])
            
            right_box = slide.shapes.add_textbox(Inches(5.3), Inches(y_start), Inches(4.2), Inches(5.5))
            frame = right_box.text_frame
            frame.word_wrap = True
            
            p = frame.paragraphs[0]
            p.text = "分析洞察"
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = self.COLORS['primary']
            p.space_after = Pt(8)
            
            # TGI统计
            if dim['tgi_max'] > 0:
                p = frame.add_paragraph()
                p.text = f"最高TGI"
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = self.COLORS['primary']
                p.space_after = Pt(2)
                
                p = frame.add_paragraph()
                p.text = f"{dim['tgi_max']:.2f}"
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = self.COLORS['accent']
                p.space_after = Pt(8)
            
            # 判断理由
            if dim['reason']:
                p = frame.add_paragraph()
                p.text = "判断理由"
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = self.COLORS['primary']
                p.space_after = Pt(4)
                
                # 完整显示判断理由（不截断）
                reason_text = dim['reason']
                p = frame.add_paragraph()
                p.text = reason_text
                p.font.size = Pt(9)  # 略微减小字号以容纳更多文字
                p.font.color.rgb = self.COLORS['text_dark']
                p.space_after = Pt(6)
            
            # 数据分布图表
            # 在文本框之外绘制柱状图
            if dim['features']:
                # 准备图表数据（Top 5）
                chart_items = []
                for feature in dim['features'][:5]:
                    try:
                        tgi_val = float(feature['tgi']) if feature['tgi'] else 0
                        chart_items.append((feature['name'][:6], tgi_val))
                    except:
                        pass
                
                # 绘制柱状图（右下角位置）
                if chart_items:
                    self._add_chart_bars(slide, 5.4, 4.5, chart_items, max_width=2.5, bar_height=0.2)
    
    def _create_strategy_slide(self, strategy: Dict):
        """倒数第2页: 综合策略建议 - 优化版"""
        slide = self._add_slide()
        self.slide_count += 1
        
        # 添加侧边装饰条（绿色）
        self._add_shape_rectangle(slide, 0, 0, 0.08, 7.5, RGBColor(112, 173, 71))
        
        self._add_title(slide, "综合策略建议")
        self._add_divider(slide)
        
        y = 1.5
        section_height = 1.8  # 增加高度以容纳更多内容
        
        strategies = [
            ('产品定位', strategy.get('position', ''), self.COLORS['accent']),
            ('渠道布局', strategy.get('channel', ''), self.COLORS['chart_orange']),
            ('内容营销', strategy.get('content', ''), self.COLORS['primary']),
        ]
        
        for idx, (title, content, color) in enumerate(strategies):
            if not content:
                continue
            
            # 背景色块（轻色）
            self._add_shape_rectangle(slide, 0.5, y - 0.05, 9, section_height, self.COLORS['bg_light'])
            
            # 标题 + 颜色标记
            self._add_shape_rectangle(slide, 0.5, y, 0.05, 0.35, color)  # 彩色条
            
            # 标题文本
            title_box = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(2), Inches(0.35))
            frame = title_box.text_frame
            frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = frame.paragraphs[0]
            p.text = title
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = color
            
            # 内容（完整显示，不截断）
            content_box = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.35), Inches(8.8), Inches(1.4))
            frame = content_box.text_frame
            frame.word_wrap = True
            p = frame.paragraphs[0]
            p.text = content  # 完整显示，不截断
            p.font.size = Pt(10)
            p.font.color.rgb = self.COLORS['text_dark']
            
            y += section_height + 0.1
    
    def _create_deep_analysis_combined_slide(self, analysis: Dict):
        """最后一页: 深度分析建议 - 一页显示单维度和交叉维度分析（通过调整框大小）"""
        slide = self._add_slide()
        self.slide_count += 1
        
        # 添加侧边装饰条
        self._add_shape_rectangle(slide, 0, 0, 0.08, 7.5, self.COLORS['primary'])
        
        self._add_title(slide, "深度分析建议")
        self._add_divider(slide)
        
        y = 1.5
        
        # 单维度分析（缩小背景框，使用更小字号）
        if 'single' in analysis and analysis['single']:
            # 背景色块（高度缩小到 3.0）
            self._add_shape_rectangle(slide, 0.5, y - 0.1, 9, 3.0, self.COLORS['bg_lighter'])
            
            # 标题
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(3), Inches(0.3))
            frame = title_box.text_frame
            p = frame.paragraphs[0]
            p.text = "📊 单维度深度分析"
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = self.COLORS['primary']
            
            # 内容（高度缩小到 2.6）
            single_text = analysis['single']
            content_box = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.35), Inches(8.6), Inches(2.6))
            frame = content_box.text_frame
            frame.word_wrap = True
            p = frame.paragraphs[0]
            p.text = single_text
            p.font.size = Pt(8)  # 字号改为 8pt 以容纳更多内容
            p.font.color.rgb = self.COLORS['text_dark']
            
            y += 3.2
        
        # 交叉维度分析（缩小背景框，使用更小字号）
        if 'cross' in analysis and analysis['cross']:
            # 背景色块（高度缩小到 3.0）
            self._add_shape_rectangle(slide, 0.5, y - 0.1, 9, 3.0, self.COLORS['highlight_bg'])
            
            # 标题
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(3), Inches(0.3))
            frame = title_box.text_frame
            p = frame.paragraphs[0]
            p.text = "🔗 交叉维度分析"
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = self.COLORS['accent_orange']
            
            # 内容（高度缩小到 2.6）
            cross_text = analysis['cross']
            content_box = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.35), Inches(8.6), Inches(2.6))
            frame = content_box.text_frame
            frame.word_wrap = True
            p = frame.paragraphs[0]
            p.text = cross_text
            p.font.size = Pt(8)  # 字号改为 8pt 以容纳更多内容
            p.font.color.rgb = self.COLORS['text_dark']
    
    def _validate_content(self) -> bool:
        """验证解析的内容是否完整"""
        print("\n📋 **内容验证**")
        issues = []
        
        # 检查用户画像
        profile = self.parsed_data.get('user_profile', {})
        if not profile.get('feature'):
            issues.append("⚠️  用户画像缺失：核心特征")
        else:
            print(f"  ✓ 核心特征: {profile['feature'][:50]}...")
        
        if not profile.get('position'):
            issues.append("⚠️  用户画像缺失：定位标签")
        else:
            print(f"  ✓ 定位标签: {profile['position'][:50]}...")
        
        # 检查维度数据
        dimensions = self.parsed_data.get('dimensions', [])
        print(f"  ✓ 维度数: {len(dimensions)}")
        for dim in dimensions:
            if not dim.get('features'):
                issues.append(f"⚠️  维度 '{dim['name']}' 缺失特征数据")
            if not dim.get('reason'):
                issues.append(f"⚠️  维度 '{dim['name']}' 缺失判断理由")
            if not dim.get('distribution'):
                issues.append(f"⚠️  维度 '{dim['name']}' 缺失数据分布")
        
        # 检查策略
        strategy = self.parsed_data.get('strategy', {})
        if not strategy.get('position'):
            issues.append("⚠️  策略缺失：产品定位")
        else:
            print(f"  ✓ 产品定位: {strategy['position'][:40]}...")
        
        if not strategy.get('channel'):
            issues.append("⚠️  策略缺失：渠道布局")
        else:
            print(f"  ✓ 渠道布局: {strategy['channel'][:40]}...")
        
        if not strategy.get('content'):
            issues.append("⚠️  策略缺失：内容营销")
        else:
            print(f"  ✓ 内容营销: {strategy['content'][:40]}...")
        
        # 检查深度分析
        analysis = self.parsed_data.get('deep_analysis', {})
        if not analysis.get('single'):
            issues.append("⚠️  分析缺失：单维度分析")
        else:
            print(f"  ✓ 单维度分析: {analysis['single'][:40]}...")
        
        if not analysis.get('cross'):
            issues.append("⚠️  分析缺失：交叉维度分析")
        else:
            print(f"  ✓ 交叉维度分析: {analysis['cross'][:40]}...")
        
        if issues:
            print("\n⚠️  **发现问题：**")
            for issue in issues:
                print(f"  {issue}")
            return False
        
        print("\n✅ **所有内容验证通过！**")
        return True
    
    def generate(self, report_path: str) -> str:
        """完整生成PPT"""
        print(f"\n📊 解析报告: {report_path}")
        self.parsed_data = self._parse_markdown(report_path)
        
        persona_name = self.parsed_data.get('persona_name', '未知人群')
        dimension_count = self.parsed_data.get('dimension_count', 0)
        has_geography = bool(self.parsed_data.get('geography'))
        
        print(f"✓ 人群名称: {persona_name}")
        print(f"✓ 维度数: {dimension_count}")
        print(f"✓ 地域数据: {'有' if has_geography else '无'}")
        
        # 内容验证
        if not self._validate_content():
            print("\n❌ 内容验证失败，请检查 Markdown 报告格式！")
            return ""
        
        # 生成各页
        print(f"\n生成PPT页面：")
        print(f"  P1: 封面")
        self._create_cover_slide(persona_name, dimension_count)
        
        print(f"  P2: 目录")
        self._create_contents_slide(dimension_count, has_geography)
        
        print(f"  P3: 用户画像总结")
        if self.parsed_data.get('user_profile'):
            self._create_user_profile_slide(self.parsed_data['user_profile'])
        
        if has_geography:
            print(f"  P4: 地域深耕方向")
            self._create_geography_slide(self.parsed_data['geography'])
        
        print(f"  P{4 if has_geography else 4}-P{3 + dimension_count if has_geography else 2 + dimension_count}: 维度分析({dimension_count}个维度)")
        if self.parsed_data.get('dimensions'):
            self._create_dimension_slides(self.parsed_data['dimensions'])
        
        print(f"  P{4 + dimension_count if has_geography else 3 + dimension_count}: 综合策略建议")
        if self.parsed_data.get('strategy'):
            self._create_strategy_slide(self.parsed_data['strategy'])
        
        print(f"  P{5 + dimension_count if has_geography else 4 + dimension_count}: 深度分析建议")
        if self.parsed_data.get('deep_analysis'):
            self._create_deep_analysis_combined_slide(self.parsed_data['deep_analysis'])
        
        # 保存
        output_file = self.output_dir / f"{Path(report_path).stem}_标准格式_v4.1.pptx"
        self.prs.save(output_file)
        
        print(f"\n✅ PPT生成完成！")
        print(f"📄 文件: {output_file}")
        print(f"📊 页数: {self.slide_count}页")
        print(f"🎨 设计: 标准格式（深蓝+红橙配色）")
        
        return str(output_file)


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("用法: python generate_ppt_standard.py <report_markdown_path> [--output <dir>]")
        sys.exit(1)
    
    report_path = sys.argv[1]
    output_dir = 'reports'
    
    if len(sys.argv) > 3 and sys.argv[2] == '--output':
        output_dir = sys.argv[3]
    
    generator = PPTStandardGenerator(output_dir)
    generator.generate(report_path)
