#!/usr/bin/env python3
import argparse
import json
import os
import zipfile
import tarfile
import sys
from typing import Dict, List, Optional, Tuple

# 尝试导入各种文档处理库
try:
    import textract
    TEXTRACT_AVAILABLE = True
except ImportError:
    TEXTRACT_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

class DocumentReader:
    def __init__(self):
        pass
    
    def get_file_extension(self, filename: str) -> str:
        """获取文件扩展名（小写）"""
        return os.path.splitext(filename)[1].lower().lstrip('.')
    
    def read_xlsx(self, filepath: str) -> str:
        """读取Excel文件"""
        if not XLSX_AVAILABLE:
            raise RuntimeError("openpyxl not installed. Please install with: pip install openpyxl")
        
        wb = load_workbook(filepath, read_only=True, data_only=True)
        output = []
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            output.append(f"=== Sheet: {sheet_name} ===")
            
            # 获取最大行列
            max_row = ws.max_row
            max_col = ws.max_column
            
            if max_row == 0:
                output.append("(Empty sheet)")
                continue
            
            # 输出为类似CSV的格式
            for row in ws.iter_rows(min_row=1, max_row=max_row, max_col=max_col):
                row_values = []
                for cell in row:
                    if cell.value is not None:
                        val = str(cell.value).replace('\n', ' ').replace('\r', '')
                        row_values.append(val)
                    else:
                        row_values.append('')
                output.append('\t'.join(row_values))
            output.append('')
        
        return '\n'.join(output)
    
    def read_docx(self, filepath: str) -> str:
        """读取Word文档"""
        if not DOCX_AVAILABLE:
            raise RuntimeError("python-docx not installed. Please install with: pip install python-docx")
        
        doc = DocxDocument(filepath)
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)
        return '\n\n'.join(paragraphs)
    
    def read_pptx(self, filepath: str) -> str:
        """读取PowerPoint文档"""
        if not PPTX_AVAILABLE:
            raise RuntimeError("python-pptx not installed. Please install with: pip install python-pptx")
        
        prs = Presentation(filepath)
        output = []
        
        for i, slide in enumerate(prs.slides, 1):
            output.append(f"=== Slide {i} ===")
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if slide_text:
                output.append('\n'.join(slide_text))
            output.append('')
        
        return '\n'.join(output)
    
    def read_text(self, filepath: str) -> str:
        """读取纯文本文件"""
        with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
            return f.read()
    
    def read_document(self, filepath: str) -> Dict:
        """读取文档，根据格式自动选择方法"""
        ext = self.get_file_extension(filepath)
        
        result = {
            'success': True,
            'filename': os.path.basename(filepath),
            'extension': ext,
            'content': '',
            'error': None
        }
        
        try:
            if not os.path.exists(filepath):
                raise FileNotFoundError(f"File not found: {filepath}")
            
            if ext == 'xlsx':
                result['content'] = self.read_xlsx(filepath)
            elif ext == 'docx':
                result['content'] = self.read_docx(filepath)
            elif ext == 'pptx':
                result['content'] = self.read_pptx(filepath)
            elif ext in ['txt', 'md', 'markdown', 'json', 'xml', 'html', 'htm', 'css', 'js', 'py', 'sh', 'bat']:
                result['content'] = self.read_text(filepath)
            elif TEXTRACT_AVAILABLE:
                # 对于其他格式，使用textract
                result['content'] = textract.process(filepath).decode('utf-8', errors='replace')
            else:
                # textract不可用时，尝试直接作为文本读取
                result['content'] = self.read_text(filepath)
        except Exception as e:
            result['success'] = False
            result['error'] = str(e)
        
        return result
    
    def list_zip(self, filepath: str) -> List[str]:
        """列出ZIP包内所有文件"""
        files = []
        with zipfile.ZipFile(filepath, 'r') as zf:
            for info in zf.infolist():
                if not info.is_dir():
                    files.append(info.filename)
        return files
    
    def list_tar(self, filepath: str) -> List[str]:
        """列出tar包内所有文件"""
        files = []
        mode = 'r'
        if filepath.endswith('.gz') or filepath.endswith('.tgz'):
            mode = 'r:gz'
        elif filepath.endswith('.bz2'):
            mode = 'r:bz2'
        
        with tarfile.open(filepath, mode) as tf:
            for info in tf:
                if info.isfile():
                    files.append(info.name)
        return files
    
    def list_rar(self, filepath: str) -> List[str]:
        """列出RAR包内所有文件"""
        files = []
        try:
            import rarfile
            with rarfile.RarFile(filepath, 'r') as rf:
                for info in rf.infolist():
                    if not info.is_dir():
                        files.append(info.filename)
        except ImportError:
            raise RuntimeError("rarfile not installed. Please install with: pip install rarfile")
        return files
    
    def list_7z(self, filepath: str) -> List[str]:
        """列出7z包内所有文件"""
        files = []
        try:
            import py7zr
            with py7zr.SevenZipFile(filepath, 'r') as zf:
                for info in zf.list():
                    if not info.is_directory:
                        files.append(info.filename)
        except ImportError:
            raise RuntimeError("py7zr not installed. Please install with: pip install py7zr")
        return files
    
    def read_from_zip(self, zip_path: str, inner_path: str) -> Dict:
        """从ZIP包读取指定文件"""
        try:
            with zipfile.ZipFile(zip_path, 'r') as zf:
                # 检查文件是否存在
                found = None
                for info in zf.infolist():
                    if info.filename == inner_path and not info.is_dir():
                        found = info
                        break
                    # 尝试不区分大小写匹配
                    if info.filename.lower() == inner_path.lower() and not info.is_dir():
                        found = info
                        break
                
                if found is None:
                    return {
                        'success': False,
                        'error': f"File '{inner_path}' not found in ZIP archive"
                    }
                
                # 读取文件内容
                with zf.open(found) as f:
                    content_bytes = f.read()
                
                # 保存到临时文件处理（因为有些库需要文件路径）
                temp_path = f"/tmp/doc_reader_{os.path.basename(inner_path)}"
                with open(temp_path, 'wb') as tf:
                    tf.write(content_bytes)
                
                try:
                    result = self.read_document(temp_path)
                    result['archive_path'] = zip_path
                    result['inner_path'] = inner_path
                    return result
                finally:
                    if os.path.exists(temp_path):
                        os.unlink(temp_path)
        
        except Exception as e:
            return {
                'success': False,
                'error': str(e)
            }
    
    def read_from_tar(self, tar_path: str, inner_path: str) -> Dict:
        """从tar包读取指定文件"""
        try:
            mode = 'r'
            if tar_path.endswith('.gz') or tar_path.endswith('.tgz'):
                mode = 'r:gz'
            elif tar_path.endswith('.bz2'):
                mode = 'r:bz2'
            
            with tarfile.open(tar_path, mode) as tf:
                # 查找文件
                found = None
                for info in tf:
                    if info.name == inner_path and info.isfile():
                        found = info
                        break
                if found is None:
                    return {
                        'success': False,
                        'error': f"File '{inner_path}' not found in TAR archive"
                    }
                
                # 读取内容
                f = tf.extractfile(found)
                if f is None:
                    return {
                        'success': False,
                        'error': f"Cannot extract file '{inner_path}'"
                    }
                content_bytes = f.read()
                
                # 保存到临时文件
                temp_path = f"/tmp/doc_reader_{os.path.basename(inner_path)}"
                with open(temp_path, 'wb') as tf_tmp:
                    tf_tmp.write(content_bytes)
                
                try:
                    result = self.read_document(temp_path)
                    result['archive_path'] = tar_path
                    result['inner_path'] = inner_path
                    return result
                finally:
                    if os.path.exists(temp_path):
                        os.unlink(temp_path)
        
        except Exception as e:
            return {
                'success': False,
                'error': str(e)
            }
    
    def read_from_rar(self, rar_path: str, inner_path: str) -> Dict:
        """从RAR包读取指定文件"""
        try:
            import rarfile
            with rarfile.RarFile(rar_path, 'r') as rf:
                # 查找文件
                found = None
                for info in rf.infolist():
                    if info.filename == inner_path and not info.isdir():
                        found = info
                        break
                if found is None:
                    # 尝试不区分大小写匹配
                    for info in rf.infolist():
                        if info.filename.lower() == inner_path.lower() and not info.isdir():
                            found = info
                            break
                if found is None:
                    return {
                        'success': False,
                        'error': f"File '{inner_path}' not found in RAR archive"
                    }
                
                # 读取内容
                content_bytes = rf.read(inner_path)
                
                # 保存到临时文件
                temp_path = f"/tmp/doc_reader_{os.path.basename(inner_path)}"
                with open(temp_path, 'wb') as tf_tmp:
                    tf_tmp.write(content_bytes)
                
                try:
                    result = self.read_document(temp_path)
                    result['archive_path'] = rar_path
                    result['inner_path'] = inner_path
                    return result
                finally:
                    if os.path.exists(temp_path):
                        os.unlink(temp_path)
        
        except Exception as e:
            return {
                'success': False,
                'error': str(e)
            }
    
    def read_from_7z(self, sz_path: str, inner_path: str) -> Dict:
        """从7z包读取指定文件"""
        try:
            import py7zr
            with py7zr.SevenZipFile(sz_path, 'r') as zf:
                # 查找文件
                found = None
                for info in zf.list():
                    if info.filename == inner_path and not info.is_directory:
                        found = info
                        break
                if found is None:
                    # 尝试不区分大小写匹配
                    for info in zf.list():
                        if info.filename.lower() == inner_path.lower() and not info.is_directory:
                            found = info
                            break
                if found is None:
                    return {
                        'success': False,
                        'error': f"File '{inner_path}' not found in 7z archive"
                    }
                
                # 读取内容
                extracted = zf.read([inner_path])
                content_bytes = extracted[inner_path].read()
                extracted[inner_path].close()
                
                # 保存到临时文件
                temp_path = f"/tmp/doc_reader_{os.path.basename(inner_path)}"
                with open(temp_path, 'wb') as tf_tmp:
                    tf_tmp.write(content_bytes)
                
                try:
                    result = self.read_document(temp_path)
                    result['archive_path'] = sz_path
                    result['inner_path'] = inner_path
                    return result
                finally:
                    if os.path.exists(temp_path):
                        os.unlink(temp_path)
        
        except Exception as e:
            return {
                'success': False,
                'error': str(e)
            }

def main():
    parser = argparse.ArgumentParser(description='Universal Document Reader')
    group = parser.add_mutually_exclusive_group(required=True)
    group.add_argument('--file', help='Path to document file to read')
    group.add_argument('--list', help='List contents of archive file')
    
    parser.add_argument('--inner-path', help='Path to file inside archive')
    parser.add_argument('--format', choices=['text', 'json'], default='text', help='Output format')
    args = parser.parse_args()
    
    reader = DocumentReader()
    
    if args.list:
        # 列出压缩包内容
        ext = reader.get_file_extension(args.list)
        if ext == 'zip':
            files = reader.list_zip(args.list)
        elif ext in ['tar', 'gz', 'tgz', 'bz2']:
            files = reader.list_tar(args.list)
        elif ext == 'rar':
            files = reader.list_rar(args.list)
        elif ext == '7z':
            files = reader.list_7z(args.list)
        else:
            print(f"Error: Unsupported archive format: {ext}")
            sys.exit(1)
        
        if args.format == 'json':
            print(json.dumps({
                'archive': args.list,
                'files': files,
                'count': len(files)
            }, indent=2, ensure_ascii=False))
        else:
            print(f"Archive: {args.list}")
            print(f"Found {len(files)} file(s):")
            print()
            for f in files:
                print(f"  {f}")
        return
    
    if args.file:
        # 读取文件
        ext = reader.get_file_extension(args.file)
        is_archive = ext in ['zip', 'tar', 'gz', 'tgz', 'bz2', 'rar', '7z']
        
        if is_archive:
            if not args.inner_path:
                print(f"Error: {args.file} is an archive. Please use --inner-path to specify which file to read, or --list to see contents.")
                sys.exit(1)
            
            if ext == 'zip':
                result = reader.read_from_zip(args.file, args.inner_path)
            elif ext in ['tar', 'gz', 'tgz', 'bz2']:
                result = reader.read_from_tar(args.file, args.inner_path)
            elif ext == 'rar':
                result = reader.read_from_rar(args.file, args.inner_path)
            elif ext == '7z':
                result = reader.read_from_7z(args.file, args.inner_path)
            else:
                result = {'success': False, 'error': f"Unsupported archive format: {ext}"}
        else:
            result = reader.read_document(args.file)
        
        if args.format == 'json':
            print(json.dumps(result, indent=2, ensure_ascii=False))
        else:
            if not result['success']:
                print(f"Error reading {result['filename']}: {result['error']}")
                sys.exit(1)
            
            print(f"=== {result['filename']} ===")
            print()
            print(result['content'])
            
            if len(result['content']) > 10000:
                print()
                print(f"[...] Content truncated. Total length: {len(result['content'])} characters")
        return

if __name__ == '__main__':
    main()
