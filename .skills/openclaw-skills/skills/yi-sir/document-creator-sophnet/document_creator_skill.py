#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import os
import sys
import json
import tempfile
import requests
from pathlib import Path
from datetime import datetime

# OpenClaw config path
OPENCLAW_CONFIG_PATH = Path.home() / '.openclaw' / 'openclaw.json'


def get_api_key() -> str:
    """Get API key"""
    soph_api_key = os.environ.get("SOPH_API_KEY")
    if soph_api_key:
        return soph_api_key
    
    if OPENCLAW_CONFIG_PATH.exists():
        try:
            with open(OPENCLAW_CONFIG_PATH, 'r', encoding='utf-8') as f:
                config = json.load(f)
            
            api_key = config.get('models', {}).get('providers', {}).get('sophnet', {}).get('apiKey')
            if api_key:
                return api_key
        except (json.JSONDecodeError, KeyError, IOError):
            pass
    
    raise ValueError("API key not found. Please set SOPH_API_KEY environment variable or ensure OpenClaw config contains sophnet API key")


def upload_file(file_path: str) -> str:
    """Upload file to OSS"""
    if not os.path.isfile(file_path):
        raise ValueError(f"File '{file_path}' does not exist")
    
    api_key = get_api_key()
    upload_url = "https://www.sophnet.com/api/open-apis/projects/upload"
    
    file_name = os.path.basename(file_path)
    with open(file_path, 'rb') as f:
        headers = {"Authorization": f"Bearer {api_key}"}
        files = {'file': (file_name, f, 'application/octet-stream')}
        response = requests.post(upload_url, headers=headers, files=files, timeout=30)
        
    response.raise_for_status()
    json_data = response.json()
    
    result = json_data.get("result")
    if not result or not isinstance(result, dict):
        raise ValueError("Response format error")
    
    signed_url = result.get("signedUrl")
    if not signed_url:
        raise ValueError("No valid URL returned")
    
    return signed_url


def create_docx_document(params: dict) -> dict:
    """Create Word document"""
    try:
        from docx import Document
        from docx.shared import Pt
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        return {
            'success': False,
            'error': 'python-docx library not installed, please run: pip install python-docx'
        }
    
    title = params.get('title', 'Untitled Document')
    content = params.get('content', '')
    author = params.get('author', 'OpenClaw Assistant')
    
    # Create document
    doc = Document()
    
    # Set document style
    style = doc.styles['Normal']
    style.font.name = 'Arial'
    style.font.size = Pt(12)
    
    # Title
    title_para = doc.add_heading(title, 0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # Author and time
    doc.add_paragraph(f"Author: {author}")
    doc.add_paragraph(f"Created: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    
    # Add content
    if content:
        # Simple Markdown parsing
        lines = content.split('\n')
        for line in lines:
            line = line.strip()
            if line.startswith('# '):
                doc.add_heading(line[2:], level=1)
            elif line.startswith('## '):
                doc.add_heading(line[3:], level=2)
            elif line.startswith('### '):
                doc.add_heading(line[4:], level=3)
            elif line.startswith('- '):
                doc.add_paragraph(line[2:], style='List Bullet')
            elif line.startswith('* '):
                doc.add_paragraph(line[2:], style='List Bullet')
            elif line:
                doc.add_paragraph(line)
    
    # Generate filename
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    filename = f"{title.replace(' ', '_')}_{timestamp}.docx"
    
    # Save document
    doc.save(filename)
    
    return {
        'success': True,
        'filename': filename,
        'type': 'docx'
    }


def create_pptx_presentation(params: dict) -> dict:
    """Create PowerPoint presentation"""
    try:
        from pptx import Presentation
    except ImportError:
        return {
            'success': False,
            'error': 'python-pptx library not installed, please run: pip install python-pptx'
        }
    
    title = params.get('title', 'Untitled Presentation')
    slides = params.get('slides', 5)
    
    # Create presentation
    prs = Presentation()
    
    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = f"Created: {datetime.now().strftime('%Y-%m-%d')}"
    
    # Add content slides
    for i in range(1, slides):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = f"Slide {i}"
        content_shape.text = f"This is content for slide {i}.\n\n• Point 1\n• Point 2\n• Point 3"
    
    # Generate filename
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    filename = f"{title.replace(' ', '_')}_{timestamp}.pptx"
    
    # Save presentation
    prs.save(filename)
    
    return {
        'success': True,
        'filename': filename,
        'type': 'pptx'
    }


def handle_document_creation(params: dict) -> dict:
    """Handle document creation request"""
    try:
        doc_type = params.get('type', 'docx')
        
        # Create document
        if doc_type == 'docx':
            result = create_docx_document(params)
        elif doc_type == 'pptx':
            result = create_pptx_presentation(params)
        else:
            return {
                'success': False,
                'error': f'Unsupported document type: {doc_type}'
            }
        
        if not result['success']:
            return result
        
        # Upload file
        upload = params.get('upload', True)
        if upload:
            try:
                url = upload_file(result['filename'])
                result['url'] = url
                result['message'] = f'{doc_type.upper()} document created and uploaded successfully'
            except Exception as e:
                result['success'] = False
                result['error'] = f'Upload failed: {str(e)}'
        else:
            result['message'] = f'{doc_type.upper()} document created successfully'
        
        return result
        
    except Exception as e:
        return {
            'success': False,
            'error': str(e)
        }


def skill_main():
    """Skill main function"""
    if len(sys.argv) > 1:
        # Command line mode
        import argparse
        
        parser = argparse.ArgumentParser(description='Document creation skill')
        parser.add_argument('type', choices=['docx', 'pptx'], help='Document type')
        parser.add_argument('--title', required=True, help='Document title')
        parser.add_argument('--content', help='Document content')
        parser.add_argument('--slides', type=int, default=5, help='Number of slides')
        parser.add_argument('--author', default='OpenClaw Assistant', help='Author name')
        parser.add_argument('--upload', type=bool, default=True, help='Upload file')
        
        args = parser.parse_args()
        
        params = {
            'type': args.type,
            'title': args.title,
            'content': args.content,
            'slides': args.slides,
            'author': args.author,
            'upload': args.upload
        }
        
        result = handle_document_creation(params)
        
        if result['success']:
            print(f"Success: {result.get('message', 'Document created successfully')}")
            if 'url' in result:
                print(f"URL: {result['url']}")
            if 'filename' in result:
                print(f"File: {result['filename']}")
        else:
            print(f"Error: {result.get('error', 'Unknown error')}")
            sys.exit(1)
    else:
        # OpenClaw Skill mode
        print("Usage: python document_creator_skill.py <docx|pptx> --title Title [--content Content] [--slides 5] [--author Author]")


if __name__ == "__main__":
    skill_main()