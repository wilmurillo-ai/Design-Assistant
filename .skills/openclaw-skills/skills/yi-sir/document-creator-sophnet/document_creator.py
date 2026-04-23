#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import os
import sys
import json
import argparse
import tempfile
import requests
from pathlib import Path
from datetime import datetime

# OpenClaw config path
OPENCLAW_CONFIG_PATH = Path.home() / '.openclaw' / 'openclaw.json'


def get_api_key() -> str:
    """Get API key, priority: environment variable > config file"""
    # 1. Priority: environment variable
    soph_api_key = os.environ.get("SOPH_API_KEY")
    if soph_api_key:
        return soph_api_key
    
    # 2. From OpenClaw config file
    if OPENCLAW_CONFIG_PATH.exists():
        try:
            with open(OPENCLAW_CONFIG_PATH, 'r', encoding='utf-8') as f:
                config = json.load(f)
            
            api_key = config.get('models', {}).get('providers', {}).get('sophnet', {}).get('apiKey')
            if api_key:
                return api_key
        except (json.JSONDecodeError, KeyError, IOError):
            pass
    
    raise ValueError("API key not found. Please set SOPH_API_KEY environment variable or ensure OpenClaw config contains sophnet API key")


def upload_file(file_path: str) -> str:
    """Upload file to OSS and return signed URL"""
    if not os.path.isfile(file_path):
        raise ValueError(f"File '{file_path}' does not exist")
    
    api_key = get_api_key()
    upload_url = "https://www.sophnet.com/api/open-apis/projects/upload"
    
    file_name = os.path.basename(file_path)
    with open(file_path, 'rb') as f:
        headers = {"Authorization": f"Bearer {api_key}"}
        files = {'file': (file_name, f, 'application/octet-stream')}
        response = requests.post(upload_url, headers=headers, files=files, timeout=30)
        
    response.raise_for_status()
    json_data = response.json()
    
    result = json_data.get("result")
    if not result or not isinstance(result, dict):
        raise ValueError("Response format error")
    
    signed_url = result.get("signedUrl")
    if not signed_url:
        raise ValueError("No valid URL returned")
    
    return signed_url


def create_docx_document(title: str, content: str = None, markdown_file: str = None, author: str = "OpenClaw Assistant") -> str:
    """Create Word document"""
    from docx import Document
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    
    # Create document
    doc = Document()
    
    # Set document style
    style = doc.styles['Normal']
    style.font.name = 'Arial'
    style.font.size = Pt(12)
    
    # Title
    title_para = doc.add_heading(title, 0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # Author and time
    doc.add_paragraph(f"Author: {author}")
    doc.add_paragraph(f"Created: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    
    # Add content
    if markdown_file and os.path.exists(markdown_file):
        # Read content from Markdown file
        with open(markdown_file, 'r', encoding='utf-8') as f:
            content = f.read()
    
    if content:
        # Simple Markdown parsing (supports headers and lists)
        lines = content.split('\n')
        for line in lines:
            line = line.strip()
            if line.startswith('# '):
                doc.add_heading(line[2:], level=1)
            elif line.startswith('## '):
                doc.add_heading(line[3:], level=2)
            elif line.startswith('### '):
                doc.add_heading(line[4:], level=3)
            elif line.startswith('- '):
                doc.add_paragraph(line[2:], style='List Bullet')
            elif line.startswith('* '):
                doc.add_paragraph(line[2:], style='List Bullet')
            elif line:
                doc.add_paragraph(line)
    
    # Generate filename
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    filename = f"{title.replace(' ', '_')}_{timestamp}.docx"
    
    # Save document
    doc.save(filename)
    return filename


def create_pptx_presentation(title: str, slides: int = 5, template: str = "minimal", 
                           markdown_file: str = None, json_file: str = None) -> str:
    """Create PowerPoint presentation"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        # Fallback if python-pptx not available
        return create_simple_pptx(title, slides)
    
    # Create presentation
    prs = Presentation()
    
    # Title slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = f"Created: {datetime.now().strftime('%Y-%m-%d')}"
    
    # Add content slides
    for i in range(1, slides):
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = f"Slide {i}"
        content_shape.text = f"This is content for slide {i}.\n\n• Point 1\n• Point 2\n• Point 3"
    
    # Generate filename
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    filename = f"{title.replace(' ', '_')}_{timestamp}.pptx"
    
    # Save presentation
    prs.save(filename)
    return filename


def create_simple_pptx(title: str, slides: int = 5) -> str:
    """Simple PPTX creation (when python-pptx not available)"""
    # Create a simple text file as alternative
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    filename = f"{title.replace(' ', '_')}_{timestamp}.txt"
    
    content = f"{title}\n"
    content += "=" * len(title) + "\n\n"
    content += f"Created: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
    content += f"Number of slides: {slides}\n\n"
    
    for i in range(1, slides + 1):
        content += f"Slide {i}:\n"
        content += f"- Content point 1\n"
        content += f"- Content point 2\n"
        content += f"- Content point 3\n\n"
    
    with open(filename, 'w', encoding='utf-8') as f:
        f.write(content)
    
    return filename


def main():
    """Main function"""
    parser = argparse.ArgumentParser(description='Document creation and upload tool')
    parser.add_argument('type', choices=['docx', 'pptx'], help='Document type')
    parser.add_argument('--title', required=True, help='Document title')
    parser.add_argument('--content', help='Document content')
    parser.add_argument('--markdown', help='Markdown file path')
    parser.add_argument('--json', help='JSON structure file')
    parser.add_argument('--slides', type=int, default=5, help='Number of slides (pptx only)')
    parser.add_argument('--template', default='minimal', help='Template name (pptx only)')
    parser.add_argument('--author', default='OpenClaw Assistant', help='Author name (docx only)')
    parser.add_argument('--output', help='Output filename')
    parser.add_argument('--upload', type=bool, default=True, help='Auto upload')
    
    args = parser.parse_args()
    
    try:
        # Create document
        if args.type == 'docx':
            filename = create_docx_document(
                title=args.title,
                content=args.content,
                markdown_file=args.markdown,
                author=args.author
            )
        else:  # pptx
            filename = create_pptx_presentation(
                title=args.title,
                slides=args.slides,
                template=args.template,
                markdown_file=args.markdown,
                json_file=args.json
            )
        
        # Rename file if output filename specified
        if args.output:
            os.rename(filename, args.output)
            filename = args.output
        
        # Upload file
        if args.upload:
            url = upload_file(filename)
            result = {
                'success': True,
                'filename': filename,
                'url': url,
                'message': f'{args.type.upper()} document created and uploaded successfully'
            }
        else:
            result = {
                'success': True,
                'filename': filename,
                'message': f'{args.type.upper()} document created successfully'
            }
        
        print(json.dumps(result, ensure_ascii=False, indent=2))
        
    except Exception as e:
        result = {
            'success': False,
            'error': str(e),
            'message': f'{args.type.upper()} document creation failed'
        }
        print(json.dumps(result, ensure_ascii=False, indent=2))
        sys.exit(1)


if __name__ == '__main__':
    main()