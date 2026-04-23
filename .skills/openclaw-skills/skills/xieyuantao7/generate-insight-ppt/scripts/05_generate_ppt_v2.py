#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
OpenHarness 技术洞察报告 - PPT 生成脚本 (优化版 v2)
字体: Microsoft YaHei (微软雅黑) | 加粗: 小标题、条目开头
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# 配色
C = {
    'primary': RGBColor(192, 0, 0),       # 深红色
    'text': RGBColor(0, 0, 0),             # 黑色
    'muted': RGBColor(102, 102, 102),      # 灰色
    'lightGray': RGBColor(245, 245, 245),  # 淡灰色
    'effectBg': RGBColor(150, 220, 223),   # 淡蓝色
    'gold': RGBColor(212, 175, 55),        # 金色
    'white': RGBColor(255, 255, 255),
    'lightRed': RGBColor(255, 204, 204),
}

FONT = 'Microsoft YaHei'

def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_text(slide, text, left, top, width, height, font_size=12, bold=False, color=None, align='left'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.name = FONT
    run.font.bold = bold
    run.font.color.rgb = color or C['text']
    if align == 'center':
        p.alignment = PP_ALIGN.CENTER
    return txBox

def add_bold_text(slide, title, desc, left, top, width, height=0.35, font_size=11):
    """添加加粗标题 + 普通描述"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = FONT
    run.font.bold = True
    run.font.size = Pt(font_size)
    run.font.color.rgb = C['text']
    run2 = p.add_run()
    run2.text = desc
    run2.font.name = FONT
    run2.font.bold = False
    run2.font.size = Pt(font_size)
    run2.font.color.rgb = C['text']
    return txBox

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_title(slide, text):
    return add_text(slide, text, 0.5, 0.3, 9, 0.6, font_size=28, bold=True, color=C['primary'])

def create_part1():
    """Part 1: 封面、目录、执行摘要"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # 1. 封面 - 白色背景，深红+黑色
    s = new_slide(prs)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C['white']
    # 左上角深红色装饰块
    add_rect(s, 0, 0, 0.15, 5.625, C['primary'])
    add_rect(s, 0, 0, 10, 0.08, C['primary'])
    add_rect(s, 0, 5.545, 10, 0.08, C['primary'])
    # 主标题 - 深红色
    add_text(s, 'OpenHarness 技术洞察报告', 0.5, 1.8, 9, 1.0, font_size=44, bold=True, color=C['primary'], align='center')
    # 副标题 - 黑色
    add_text(s, '企业级 DevOps 平台开源实践', 0.5, 2.9, 9, 0.6, font_size=22, color=C['text'], align='center')
    # 标签 - 黑色
    add_text(s, 'DevOps | CI/CD | GitOps | Kubernetes', 0.5, 3.6, 9, 0.4, font_size=14, color=C['text'], align='center')
    # 日期 - 灰色
    add_text(s, '2026年4月', 0.5, 5.0, 9, 0.3, font_size=12, color=C['muted'], align='center')
    
    # 2. 目录
    s = new_slide(prs)
    add_title(s, '目 录')
    items = ['执行摘要', '技术架构总览', 'Delegate 代理模式', 'CI 持续集成引擎', 
             'CD 持续交付流水线', 'RBAC 权限治理', 'Policy as Code', 
             '持续验证 (CV)', '竞品对比分析', '总结与建议']
    for i, item in enumerate(items):
        y = 1.1 + i * 0.43
        add_rect(s, 0.5, y, 0.4, 0.38, C['primary'])
        add_text(s, f'{(i+1):02d}', 0.5, y, 0.4, 0.38, font_size=12, bold=True, color=C['white'], align='center')
        add_text(s, item, 1.0, y, 8, 0.38, font_size=14)
    
    # 3. 执行摘要
    s = new_slide(prs)
    add_title(s, '执行摘要')
    stats = [('GitHub Stars', '34.4k'), ('Forks', '2.9k'), ('支持平台', '10+'), ('CI/CD 集成', '100+')]
    for i, (label, value) in enumerate(stats):
        x = 0.5 + i * 2.3
        add_rect(s, x, 1.05, 2.1, 0.95, C['lightGray'])
        add_text(s, value, x, 1.1, 2.1, 0.5, font_size=26, bold=True, color=C['primary'], align='center')
        add_text(s, label, x, 1.6, 2.1, 0.3, font_size=10, color=C['muted'], align='center')
    add_rect(s, 0.5, 2.15, 1.5, 0.35, C['primary'])
    add_text(s, '核心发现', 0.5, 2.15, 1.5, 0.35, font_size=13, bold=True, color=C['white'], align='center')
    findings = [
        ('34.4k GitHub Stars', ' - 全球最大规模企业级 DevOps 开源平台'),
        ('Go 语言微服务架构', ' - 高性能、易部署的现代化架构'),
        ('GitOps 原生设计', ' - Git 作为单一真相来源的声明式交付'),
        ('Delegate 代理模式', ' - SaaS 平台与客户内网的安全连接'),
        ('完整能力栈', ' - 从源码管理到持续交付的全链路覆盖'),
    ]
    for i, (title, desc) in enumerate(findings):
        y = 2.65 + i * 0.52
        add_rect(s, 0.5, y, 0.05, 0.4, C['primary'])
        add_bold_text(s, title, desc, 0.65, y, 8.7, 0.4)
    
    prs.save('D:/techinsight/reports/insight_openharness/output/part1_v2.pptx')
    print('Part 1 done')

def create_part2():
    """Part 2: 技术架构总览、Delegate 代理模式"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # 技术架构总览
    s = new_slide(prs)
    add_title(s, '技术架构总览')
    add_rect(s, 0.5, 1.05, 9, 0.35, C['effectBg'])
    add_text(s, 'Harness Platform (SaaS)', 0.65, 1.1, 8.7, 0.3, font_size=13, bold=True, color=C['primary'])
    
    # Account 层
    add_rect(s, 0.5, 1.55, 9, 0.6, C['lightGray'])
    add_text(s, 'Account', 0.6, 1.6, 1.5, 0.5, font_size=12, bold=True, color=C['primary'])
    add_text(s, '全局配置 | SSO | 审计日志 | 计费', 2.1, 1.6, 7.2, 0.5, font_size=11, color=C['text'])
    
    # Organization 层
    add_rect(s, 0.5, 2.2, 9, 0.6, C['lightGray'])
    add_text(s, 'Organization', 0.6, 2.25, 1.5, 0.5, font_size=12, bold=True, color=C['primary'])
    add_text(s, '业务隔离 | 共享密钥 | 连接器', 2.1, 2.25, 7.2, 0.5, font_size=11, color=C['text'])
    
    # Project 层
    add_rect(s, 0.5, 2.85, 9, 0.6, C['lightGray'])
    add_text(s, 'Project', 0.6, 2.9, 1.5, 0.5, font_size=12, bold=True, color=C['primary'])
    add_text(s, 'CI/CD | Feature Flags | Chaos Engineering | Governance', 2.1, 2.9, 7.2, 0.5, font_size=11, color=C['text'])
    
    # 通信层
    add_rect(s, 0.5, 3.5, 9, 0.6, C['effectBg'])
    add_text(s, '通信层', 0.6, 3.55, 1.5, 0.5, font_size=12, bold=True, color=C['primary'])
    add_text(s, 'WebSocket TLS | gRPC | Token 认证', 2.1, 3.55, 7.2, 0.5, font_size=11, color=C['text'])
    
    # Delegate 层
    add_rect(s, 0.5, 4.15, 9, 0.6, C['lightGray'])
    add_text(s, 'Delegate', 0.6, 4.2, 1.5, 0.5, font_size=12, bold=True, color=C['primary'])
    add_text(s, 'K8s/Docker/Helm Delegate | Task级最小权限 | 高可用部署', 2.1, 4.2, 7.2, 0.5, font_size=11, color=C['text'])
    
    # Delegate 代理模式
    s = new_slide(prs)
    add_title(s, 'Delegate 代理模式')
    add_rect(s, 0.5, 0.95, 9, 0.35, C['effectBg'])
    add_text(s, '混合云安全集成架构', 0.65, 0.98, 8.7, 0.3, font_size=14, color=C['muted'])
    add_rect(s, 0.5, 1.4, 4, 2.8, C['lightGray'])
    add_text(s, '架构原理', 0.65, 1.5, 3.7, 0.3, font_size=12, bold=True, color=C['primary'])
    add_rect(s, 1.3, 1.9, 1.4, 0.5, C['primary'])
    add_text(s, 'SaaS Platform', 1.3, 1.9, 1.4, 0.5, font_size=10, bold=True, color=C['white'], align='center')
    add_rect(s, 1.95, 2.5, 0.1, 0.3, C['muted'])
    add_text(s, 'WebSocket TLS', 1.3, 2.85, 1.4, 0.25, font_size=8, color=C['muted'], align='center')
    add_rect(s, 1.3, 3.15, 1.4, 0.5, C['effectBg'])
    add_text(s, 'Delegate', 1.3, 3.15, 1.4, 0.5, font_size=11, bold=True, color=C['primary'], align='center')
    add_text(s, 'Customer Network', 1.3, 3.65, 1.4, 0.25, font_size=8, color=C['muted'], align='center')
    add_text(s, '技术细节', 4.7, 1.4, 4.8, 0.3, font_size=12, bold=True, color=C['primary'])
    add_bold_text(s, '双向 TLS 认证:', ' TLS 1.3 加密，证书自动轮换，端到端安全。', 4.7, 1.75, 4.8, 0.35)
    add_bold_text(s, 'Task级最小权限:', ' 每个任务仅授予所需最小权限，防止横向移动。', 4.7, 2.2, 4.8, 0.35)
    add_bold_text(s, '出站连接架构:', ' 仅需建立出站 WebSocket，无需防火墙入站端口。', 4.7, 2.65, 4.8, 0.35)
    add_bold_text(s, '高可用部署:', ' 支持 3+ Delegate 实例集群，K8s 支持 HPA 自动扩缩容。', 4.7, 3.1, 4.8, 0.35)
    add_rect(s, 0.5, 4.35, 9, 0.9, C['effectBg'])
    add_text(s, '效果: 部署时间 < 5分钟 | 资源占用 ~100MB | 支持 Linux/macOS/Windows', 0.65, 4.5, 8.7, 0.6, font_size=11)
    
    prs.save('D:/techinsight/reports/insight_openharness/output/part2_v2.pptx')
    print('Part 2 done')

def create_part3():
    """Part 3: CI 持续集成引擎"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    s = new_slide(prs)
    add_title(s, 'CI 持续集成引擎')
    add_rect(s, 0.5, 0.95, 9, 0.35, C['effectBg'])
    add_text(s, '云原生容器化 CI 执行引擎，支持多语言并行构建', 0.65, 0.98, 8.7, 0.3, font_size=14, color=C['muted'])
    
    add_rect(s, 0.5, 1.4, 4.3, 1.8, C['lightGray'])
    add_text(s, '核心特性', 0.65, 1.5, 4, 0.3, font_size=12, bold=True, color=C['primary'])
    features = [
        ('Docker 容器化:', ' 每个步骤在独立容器中执行，环境隔离'),
        ('智能缓存:', ' 支持依赖缓存、层缓存、多级缓存策略'),
        ('DAG 并行调度:', ' 基于有向无环图的并行执行优化'),
        ('资源弹性:', ' 自动扩缩容构建资源，按需使用'),
    ]
    for i, (title, desc) in enumerate(features):
        add_bold_text(s, title, desc, 0.65, 1.85 + i * 0.38, 4, 0.35)
    
    add_rect(s, 5, 1.4, 4.5, 1.8, C['lightGray'])
    add_text(s, '支持语言', 5.15, 1.5, 4.2, 0.3, font_size=12, bold=True, color=C['primary'])
    langs = ['Go', 'Java', 'Node.js', 'Python', 'Ruby', 'C++', '.NET', 'Rust']
    for i, lang in enumerate(langs):
        col = i % 2
        row = i // 2
        x = 5.15 + col * 2.1
        y = 1.85 + row * 0.32
        add_text(s, lang, x, y, 2, 0.3, font_size=10, bold=True, color=C['text'])
    
    add_rect(s, 0.5, 3.35, 9, 1.0, C['lightGray'])
    add_text(s, '技术亮点', 0.65, 3.4, 8.7, 0.3, font_size=12, bold=True, color=C['primary'])
    highlights = [
        '步进级日志实时流式输出，支持构建历史回放',
        '内置 100+ 步骤模板，支持自定义步骤开发',
        '与 GitHub/GitLab/Bitbucket 深度集成',
        '支持 Pull Request 验证和分支策略',
    ]
    for i, h in enumerate(highlights):
        add_text(s, f'- {h}', 0.65, 3.75 + i * 0.28, 8.7, 0.28, font_size=10)
    
    add_rect(s, 0.5, 4.5, 9, 0.75, C['effectBg'])
    add_text(s, '效果: 构建时间缩短 60% | 缓存命中率 > 80% | 并行度提升 5x', 0.65, 4.6, 8.7, 0.5, font_size=11)
    
    prs.save('D:/techinsight/reports/insight_openharness/output/part3_v2.pptx')
    print('Part 3 done')

def create_part4():
    """Part 4: CD 持续交付流水线、RBAC 权限治理"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # CD 持续交付流水线
    s = new_slide(prs)
    add_title(s, 'CD 持续交付流水线')
    add_rect(s, 0.5, 0.95, 9, 0.35, C['effectBg'])
    add_text(s, '声明式 YAML 流水线定义，支持复杂部署策略', 0.65, 0.98, 8.7, 0.3, font_size=14, color=C['muted'])
    
    add_rect(s, 0.5, 1.4, 4.3, 1.6, C['lightGray'])
    add_text(s, '部署策略', 0.65, 1.5, 4, 0.3, font_size=12, bold=True, color=C['primary'])
    strategies = [
        ('Rolling Update:', ' 滚动更新，逐步替换旧版本实例'),
        ('Blue-Green:', ' 双环境切换，零 downtime 部署'),
        ('Canary:', ' 金丝雀发布，小流量验证后全量'),
        ('Immutable:', ' 不可变基础设施，完整重建部署'),
    ]
    for i, (title, desc) in enumerate(strategies):
        add_bold_text(s, title, desc, 0.65, 1.85 + i * 0.36, 4, 0.35)
    
    add_rect(s, 5, 1.4, 4.5, 1.6, C['lightGray'])
    add_text(s, '核心能力', 5.15, 1.5, 4.2, 0.3, font_size=12, bold=True, color=C['primary'])
    caps = ['GitOps 工作流', '环境模板化', '配置即代码', '回滚自动化']
    for i, cap in enumerate(caps):
        add_text(s, f'- {cap}', 5.15, 1.85 + i * 0.36, 4.2, 0.35, font_size=11)
    
    add_rect(s, 0.5, 3.15, 9, 0.8, C['effectBg'])
    add_text(s, '工作流: 代码提交 -> CI 构建 -> 安全扫描 -> 制品入库 -> CD 部署 -> 持续验证', 0.65, 3.25, 8.7, 0.5, font_size=11)
    
    # RBAC 权限治理
    s = new_slide(prs)
    add_title(s, 'RBAC 权限治理')
    add_rect(s, 0.5, 0.95, 9, 0.35, C['effectBg'])
    add_text(s, 'Account-Organization-Project 三层权限模型', 0.65, 0.98, 8.7, 0.3, font_size=14, color=C['muted'])
    
    add_rect(s, 0.5, 1.4, 9, 1.5, C['lightGray'])
    add_text(s, '权限模型', 0.65, 1.5, 8.7, 0.3, font_size=12, bold=True, color=C['primary'])
    elements = [
        ('主体 (Principal):', ' User | Service Account | API Key | SSO Group'),
        ('资源组 (Resource Group):', ' Pipeline | Environment | Connector | Secret'),
        ('角色 (Role):', ' Admin | Executor | Viewer | Custom'),
        ('权限 (Permission):', ' pipeline_execute | env_deploy | secret_manage'),
    ]
    for i, (title, desc) in enumerate(elements):
        add_bold_text(s, title, desc, 0.65, 1.85 + i * 0.35, 8.7, 0.35)
    
    add_rect(s, 0.5, 3.05, 4.3, 1.3, C['lightGray'])
    add_text(s, '预置角色', 0.65, 3.1, 4, 0.3, font_size=12, bold=True, color=C['primary'])
    roles = ['Account Admin', 'Organization Admin', 'Project Admin', 'Pipeline Executor', 'Environment Viewer']
    for i, role in enumerate(roles):
        add_text(s, f'- {role}', 0.65, 3.45 + i * 0.28, 4, 0.28, font_size=10)
    
    add_rect(s, 5, 3.05, 4.5, 1.3, C['lightGray'])
    add_text(s, '安全特性', 5.15, 3.1, 4.2, 0.3, font_size=12, bold=True, color=C['primary'])
    sec = ['最小权限原则', '权限继承与覆盖', '审计日志记录', 'API 访问控制', 'OAuth 2.0 集成']
    for i, s_item in enumerate(sec):
        add_text(s, f'- {s_item}', 5.15, 3.45 + i * 0.28, 4.2, 0.28, font_size=10)
    
    add_rect(s, 0.5, 4.5, 9, 0.75, C['effectBg'])
    add_text(s, '效果: 细粒度权限控制 | 合规审计追溯 | SSO 单点登录 | 多因素认证支持', 0.65, 4.6, 8.7, 0.5, font_size=11)
    
    prs.save('D:/techinsight/reports/insight_openharness/output/part4_v2.pptx')
    print('Part 4 done')

def create_part5():
    """Part 5: Policy as Code、持续验证 (CV)"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Policy as Code
    s = new_slide(prs)
    add_title(s, 'Policy as Code')
    add_rect(s, 0.5, 0.95, 9, 0.35, C['effectBg'])
    add_text(s, '基于 OPA Rego 语言的声明式策略引擎', 0.65, 0.98, 8.7, 0.3, font_size=14, color=C['muted'])
    
    add_rect(s, 0.5, 1.4, 4.3, 1.8, C['lightGray'])
    add_text(s, 'OPA 策略引擎', 0.65, 1.5, 4, 0.3, font_size=12, bold=True, color=C['primary'])
    opa = [
        ('Rego 语言:', ' 声明式策略定义语言，简洁表达'),
        ('Bundle 部署:', ' 策略版本化管理与分发'),
        ('内置函数库:', ' 丰富的字符串、正则、时间处理'),
        ('测试框架:', ' 支持策略单元测试与覆盖率'),
    ]
    for i, (title, desc) in enumerate(opa):
        add_bold_text(s, title, desc, 0.65, 1.85 + i * 0.38, 4, 0.35)
    
    add_rect(s, 5, 1.4, 4.5, 1.8, C['lightGray'])
    add_text(s, '应用场景', 5.15, 1.5, 4.2, 0.3, font_size=12, bold=True, color=C['primary'])
    scenarios = [
        '准入控制 (Admission Control)',
        '流水线策略校验',
        'Secret 访问审计',
        '资源配额限制',
        '合规性检查',
    ]
    for i, sc in enumerate(scenarios):
        add_text(s, f'- {sc}', 5.15, 1.85 + i * 0.32, 4.2, 0.32, font_size=11)
    
    add_rect(s, 0.5, 3.35, 9, 0.8, C['lightGray'])
    add_text(s, '策略示例: 检查 Kubernetes 资源标签合规性，强制环境隔离，拒绝未加密 Secret 部署', 0.65, 3.45, 8.7, 0.6, font_size=11)
    
    # 持续验证 CV
    s = new_slide(prs)
    add_title(s, '持续验证 (CV)')
    add_rect(s, 0.5, 0.95, 9, 0.35, C['effectBg'])
    add_text(s, '实时 APM/日志/指标集成，AI 驱动的异常检测与自动回滚', 0.65, 0.98, 8.7, 0.3, font_size=14, color=C['muted'])
    
    add_rect(s, 0.5, 1.4, 4.3, 1.6, C['lightGray'])
    add_text(s, '监控集成', 0.65, 1.5, 4, 0.3, font_size=12, bold=True, color=C['primary'])
    monitors = [
        ('APM 集成:', ' New Relic | Datadog | AppDynamics | Dynatrace'),
        ('日志集成:', ' Splunk | ELK | CloudWatch | Sumo Logic'),
        ('指标集成:', ' Prometheus | CloudWatch | Stackdriver'),
        ('自定义:', ' REST API | Prometheus Metrics'),
    ]
    for i, (title, desc) in enumerate(monitors):
        add_bold_text(s, title, desc, 0.65, 1.85 + i * 0.35, 4, 0.35)
    
    add_rect(s, 5, 1.4, 4.5, 1.6, C['lightGray'])
    add_text(s, '异常检测', 5.15, 1.5, 4.2, 0.3, font_size=12, bold=True, color=C['primary'])
    anomalies = [
        '时序数据分析',
        '基线对比分析',
        '多指标联合分析',
        '自动阈值调整',
    ]
    for i, an in enumerate(anomalies):
        add_text(s, f'- {an}', 5.15, 1.85 + i * 0.35, 4.2, 0.35, font_size=11)
    
    add_rect(s, 0.5, 3.15, 9, 0.8, C['effectBg'])
    add_text(s, '验证流程: 部署 -> 数据采集 -> 异常检测 -> 风险评估 -> 自动回滚 (可选)', 0.65, 3.25, 8.7, 0.5, font_size=11)
    
    add_rect(s, 0.5, 4.1, 9, 0.8, C['lightGray'])
    add_text(s, '关键指标', 0.65, 4.15, 8.7, 0.3, font_size=12, bold=True, color=C['primary'])
    metrics = ['Error Rate | Latency | Throughput | Custom Metrics | Business KPIs']
    add_text(s, metrics[0], 0.65, 4.5, 8.7, 0.3, font_size=11)
    
    add_rect(s, 0.5, 5.0, 9, 0.5, C['effectBg'])
    add_text(s, '效果: MTTR 缩短 70% | 主动发现问题 | 零手动回滚 | SLA 保障', 0.65, 5.05, 8.7, 0.4, font_size=11)
    
    prs.save('D:/techinsight/reports/insight_openharness/output/part5_v2.pptx')
    print('Part 5 done')

def create_part6():
    """Part 6: 竞品对比分析、总结与建议、感谢页"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # 竞品对比
    s = new_slide(prs)
    add_title(s, '竞品对比分析')
    add_rect(s, 0.5, 0.95, 9, 0.35, C['effectBg'])
    add_text(s, 'Jenkins X | Argo CD | Tekton | GitLab CI vs OpenHarness', 0.65, 0.98, 8.7, 0.3, font_size=14, color=C['muted'])
    
    # 对比表格
    headers = ['维度', 'OpenHarness', 'Jenkins X', 'Argo CD', 'GitLab CI']
    data = [
        ['部署模式', 'SaaS + 自托管', '自托管', '自托管', 'SaaS + 自托管'],
        ['K8s 原生', 'Yes', 'Yes', 'Yes', 'Limited'],
        ['GitOps', 'Native', 'Argo', 'Native', 'Limited'],
        ['RBAC', 'Enterprise', 'Basic', 'Basic', 'Basic'],
        ['CI/CD 一体化', 'Yes', 'Yes', 'No', 'Yes'],
        ['可观测性集成', '100+', 'Limited', 'Limited', 'Limited'],
        ['学习曲线', 'Medium', 'Steep', 'Medium', 'Low'],
    ]
    
    col_widths = [1.5, 2.5, 1.8, 1.8, 1.8]
    for i, header in enumerate(headers):
        x = 0.5 + sum(col_widths[:i])
        add_rect(s, x, 1.4, col_widths[i], 0.4, C['primary'])
        add_text(s, header, x, 1.4, col_widths[i], 0.4, font_size=10, bold=True, color=C['white'], align='center')
    
    for row_i, row in enumerate(data):
        y = 1.8 + row_i * 0.42
        bg = C['lightGray'] if row_i % 2 == 0 else C['white']
        for col_i, cell in enumerate(row):
            x = 0.5 + sum(col_widths[:col_i])
            add_rect(s, x, y, col_widths[col_i], 0.4, bg)
            bold = col_i == 0 or (col_i == 1 and row_i < len(data))
            add_text(s, cell, x, y, col_widths[col_i], 0.4, font_size=9, bold=bold, align='center')
    
    add_rect(s, 0.5, 4.85, 9, 0.4, C['effectBg'])
    add_text(s, '结论: OpenHarness 在企业级特性、可观测性集成、GitOps 支持方面具有明显优势', 0.65, 4.9, 8.7, 0.3, font_size=11)
    
    # 总结与建议
    s = new_slide(prs)
    add_title(s, '总结与建议')
    add_rect(s, 0.5, 0.95, 9, 0.35, C['effectBg'])
    add_text(s, '基于技术洞察的战略建议', 0.65, 0.98, 8.7, 0.3, font_size=14, color=C['muted'])
    
    add_rect(s, 0.5, 1.4, 4.3, 1.8, C['lightGray'])
    add_text(s, '技术优势', 0.65, 1.5, 4, 0.3, font_size=12, bold=True, color=C['primary'])
    advantages = [
        '34.4k Stars 验证社区认可度',
        'Go 语言微服务架构高性能',
        'Delegate 模式安全灵活',
        'GitOps 原生支持现代化交付',
        '100+ 可观测性集成开箱即用',
    ]
    for i, adv in enumerate(advantages):
        add_text(s, f'- {adv}', 0.65, 1.85 + i * 0.3, 4, 0.3, font_size=10)
    
    add_rect(s, 5, 1.4, 4.5, 1.8, C['lightGray'])
    add_text(s, '应用建议', 5.15, 1.5, 4.2, 0.3, font_size=12, bold=True, color=C['primary'])
    suggestions = [
        '适合: 大规模微服务 DevOps 转型',
        '适合: 混合云/多集群部署场景',
        '适合: 需要强合规审计的企业',
        '推荐: 从 Delegate 试点开始',
        '推荐: 利用 Harness OPA 强化治理',
    ]
    for i, sug in enumerate(suggestions):
        add_text(s, f'- {sug}', 5.15, 1.85 + i * 0.3, 4.2, 0.3, font_size=10)
    
    add_rect(s, 0.5, 3.35, 9, 1.5, C['lightGray'])
    add_text(s, '关键观察', 0.65, 3.4, 8.7, 0.3, font_size=12, bold=True, color=C['primary'])
    observations = [
        '1. OpenHarness 的 Delegate 代理模式解决了 SaaS 安全集成难题',
        '2. Policy as Code 能力使其成为合规敏感行业的优选',
        '3. 持续验证 (CV) 功能弥合了部署与运营的 gap',
        '4. 34.4k Stars 表明开源社区的强劲信心',
    ]
    for i, obs in enumerate(observations):
        add_text(s, obs, 0.65, 3.75 + i * 0.3, 8.7, 0.3, font_size=10)
    
    add_rect(s, 0.5, 5.0, 9, 0.5, C['effectBg'])
    add_text(s, '行动建议: POC 试点 -> Delegate 部署 -> CI 集成 -> CD 扩展 -> CV 启用', 0.65, 5.05, 8.7, 0.4, font_size=11)
    
    # 感谢页 - 白色背景，深红+黑色
    s = new_slide(prs)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C['white']
    # 左上角深红色装饰块
    add_rect(s, 0, 0, 0.15, 5.625, C['primary'])
    add_rect(s, 0, 0, 10, 0.08, C['primary'])
    add_rect(s, 0, 5.545, 10, 0.08, C['primary'])
    # 主标题 - 深红色
    add_text(s, '感谢阅读', 0.5, 1.8, 9, 1.0, font_size=48, bold=True, color=C['primary'], align='center')
    # 副标题 - 黑色
    add_text(s, 'OpenHarness 技术洞察报告', 0.5, 2.9, 9, 0.6, font_size=22, color=C['text'], align='center')
    # 标签 - 黑色
    add_text(s, 'DevOps | CI/CD | GitOps | Kubernetes', 0.5, 3.6, 9, 0.4, font_size=14, color=C['text'], align='center')
    # 日期 - 灰色
    add_text(s, '2026年4月', 0.5, 5.0, 9, 0.3, font_size=12, color=C['muted'], align='center')
    
    prs.save('D:/techinsight/reports/insight_openharness/output/part6_v2.pptx')
    print('Part 6 done')

if __name__ == '__main__':
    import os
    os.makedirs('D:/techinsight/reports/insight_openharness/output', exist_ok=True)
    create_part1()
    create_part2()
    create_part3()
    create_part4()
    create_part5()
    create_part6()
    print('All done!')
