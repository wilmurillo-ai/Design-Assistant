# -*- coding: utf-8 -*-
"""
ppt_lib.py - 暗色科技风 PPT 生成工具库
用法: from ppt_lib import *
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ======================== 调色板 ========================
DARK_BG    = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT     = RGBColor(0x00, 0xD2, 0xFF)
ACCENT2    = RGBColor(0x7C, 0x3A, 0xED)
ACCENT3    = RGBColor(0x10, 0xB9, 0x81)
ACCENT4    = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT5    = RGBColor(0xEF, 0x44, 0x44)
ACCENT6    = RGBColor(0x8B, 0x5C, 0xF6)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY   = RGBColor(0x99, 0x99, 0x99)
CARD_BG    = RGBColor(0x25, 0x25, 0x40)
CARD_BG2   = RGBColor(0x2D, 0x2D, 0x4A)

PALETTE = [ACCENT, ACCENT2, ACCENT3, ACCENT4, ACCENT5, ACCENT6]


def create_presentation(wide=True):
    """创建16:9演示文稿"""
    prs = Presentation()
    if wide:
        prs.slide_width  = Inches(13.333)
        prs.slide_height = Inches(7.5)
    else:
        prs.slide_width  = Inches(10)
        prs.slide_height = Inches(7.5)
    return prs


def save_ppt(prs, path):
    """保存PPT"""
    prs.save(path)
    print(f'Done: {path}')


# ======================== 基础元素 ========================

def set_bg(slide, color=DARK_BG):
    """设置幻灯片背景色"""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, corner=None):
    """添加矩形/圆角矩形"""
    if corner:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = corner
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text,
             size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
             font='Microsoft YaHei'):
    """添加文本框"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tb


def add_line(slide, left, top, width, color=ACCENT, thickness=Pt(3)):
    """添加强调色装饰线"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, thickness)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_badge(slide, left, top, text, color=ACCENT, width=Inches(1.5)):
    """添加标签/徽章"""
    shape = add_rect(slide, left, top, width, Inches(0.38), color, corner=0.15)
    add_text(slide, left, top + Pt(2), width, Inches(0.34),
             text, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    return shape


# ======================== 复合组件 ========================

def add_card(slide, left, top, width, height, title, lines,
             icon='', color=ACCENT, title_size=16, line_size=13):
    """创建带标题+内容的卡片"""
    add_rect(slide, left, top, width, height, CARD_BG, corner=0.03)
    add_rect(slide, left, top, width, Pt(4), color)
    x_icon = left + Inches(0.25) if icon else left + Inches(0.2)
    add_text(slide, x_icon, top + Inches(0.15), Inches(0.5), Inches(0.45),
             icon, size=20, color=color, bold=True)
    title_x = left + Inches(0.25) if not icon else left + Inches(0.6)
    add_text(slide, title_x, top + Inches(0.15), width - Inches(0.8), Inches(0.45),
             title, size=title_size, color=WHITE, bold=True)
    y = top + Inches(0.65)
    for line in lines:
        add_text(slide, left + Inches(0.35), y, width - Inches(0.6), Inches(0.30),
                 line, size=line_size, color=LIGHT_GRAY)
        y += Inches(0.30)


def add_flow_arrow(slide, left, top):
    """添加流程箭头"""
    add_text(slide, left, top, Inches(0.4), Inches(0.5),
             '→', size=28, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ======================== 页面模板 ========================

def make_cover_slide(prs, title, subtitle='', date_text=''):
    """封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), prs.slide_height, ACCENT)
    add_text(slide, Inches(1.2), Inches(1.8), Inches(10), Inches(1.2),
             title, size=44, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(1.2), Inches(3.0), Inches(10), Inches(0.7),
                 subtitle, size=22, color=ACCENT)
        add_line(slide, Inches(1.2), Inches(3.8), Inches(3))
    if date_text:
        add_text(slide, Inches(1.2), Inches(4.3), Inches(5), Inches(0.5),
                 date_text, size=16, color=MID_GRAY)
    return slide


def make_section_slide(prs, number, title, bullets, accent_color=ACCENT):
    """章节页: 编号 + 标题 + 要点列表"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    num_text = f'{number:02d}  {title}' if isinstance(number, int) else f'{number}  {title}'
    add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             num_text, size=30, color=accent_color, bold=True)
    add_line(slide, Inches(0.8), Inches(1.0), Inches(2), accent_color)

    # 要点
    add_rect(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.8 + len(bullets) * 0.45),
             CARD_BG, corner=0.02)
    for i, bullet in enumerate(bullets):
        add_text(slide, Inches(1.3), Inches(1.55 + i * 0.45), Inches(10.5), Inches(0.4),
                 f'▸  {bullet}', size=16, color=LIGHT_GRAY)
    return slide


def make_content_slide(prs, title, cards_data, accent_color=ACCENT):
    """内容页: 标题 + 卡片列表 (自动换行, 每行2个)
    cards_data: [(card_title, [lines], color), ...]
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             title, size=30, color=accent_color, bold=True)
    add_line(slide, Inches(0.8), Inches(1.0), Inches(2), accent_color)

    for i, (ct, cl, cc) in enumerate(cards_data):
        col = i % 2
        row = i // 2
        left = Inches(0.6 + col * 6.3)
        top = Inches(1.4 + row * 2.9)
        add_card(slide, left, top, Inches(5.9), Inches(2.5), ct, cl, color=cc)
    return slide


def make_cards_slide(prs, title, cards_data, cols=3, card_w=Inches(3.7),
                     card_h=Inches(2.0), accent_color=ACCENT):
    """卡片网格页
    cards_data: [(card_title, [lines], color), ...]
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             title, size=30, color=accent_color, bold=True)
    add_line(slide, Inches(0.8), Inches(1.0), Inches(2), accent_color)

    gap = Inches(0.15)
    start_x = Inches(0.6)
    for i, (ct, cl, cc) in enumerate(cards_data):
        col = i % cols
        row = i // cols
        left = start_x + col * (card_w + gap)
        top = Inches(1.3) + row * (card_h + gap)
        add_card(slide, left, top, card_w, card_h, ct, cl, color=cc)
    return slide


def make_compare_slide(prs, title, left_title, left_lines, right_title, right_lines,
                       left_color=ACCENT, right_color=ACCENT2):
    """对比页: 左右两栏"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             title, size=30, color=ACCENT, bold=True)
    add_line(slide, Inches(0.8), Inches(1.0), Inches(2))

    # 左栏
    add_card(slide, Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.5),
             left_title, left_lines, color=left_color, line_size=14)
    # 右栏
    add_card(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.5),
             right_title, right_lines, color=right_color, line_size=14)
    return slide


def make_summary_slide(prs, title, values, end_text='谢谢！'):
    """总结页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), prs.slide_height, ACCENT)
    add_text(slide, Inches(1.2), Inches(1.5), Inches(10), Inches(0.8),
             title, size=38, color=WHITE, bold=True)
    add_line(slide, Inches(1.2), Inches(2.3), Inches(2.5))

    for i, (icon, val_title, desc) in enumerate(values):
        y = Inches(2.8 + i * 0.85)
        add_text(slide, Inches(1.5), y, Inches(0.5), Inches(0.5),
                 icon, size=24, color=ACCENT3)
        add_text(slide, Inches(2.1), y, Inches(4), Inches(0.4),
                 val_title, size=20, color=WHITE, bold=True)
        add_text(slide, Inches(2.1), y + Inches(0.38), Inches(8), Inches(0.35),
                 desc, size=15, color=LIGHT_GRAY)

    if end_text:
        add_text(slide, Inches(1.2), Inches(6.3), Inches(10), Inches(0.5),
                 end_text, size=28, color=ACCENT, bold=True)
    return slide


def make_table_slide(prs, title, headers, rows, accent_color=ACCENT):
    """表格页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             title, size=30, color=accent_color, bold=True)
    add_line(slide, Inches(0.8), Inches(1.0), Inches(2), accent_color)

    col_count = len(headers)
    col_w = Inches(11.5) / col_count
    row_h = Inches(0.45)
    table_top = Inches(1.5)

    # 表头
    for j, h in enumerate(headers):
        x = Inches(0.8) + j * col_w
        add_rect(slide, x, table_top, col_w, row_h, accent_color)
        add_text(slide, x, table_top + Pt(3), col_w, row_h,
                 h, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # 数据行
    for i, row in enumerate(rows):
        y = table_top + (i + 1) * row_h
        bg = CARD_BG if i % 2 == 0 else CARD_BG2
        for j, cell in enumerate(row):
            x = Inches(0.8) + j * col_w
            add_rect(slide, x, y, col_w, row_h, bg)
            add_text(slide, x, y + Pt(2), col_w, row_h,
                     str(cell), size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    return slide
