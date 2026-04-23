#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
文件创建脚本 - 支持 docx, xlsx, pptx, pdf

用法:
    python create_doc.py --type docx --output "文件.docx" --content "内容"
    python create_doc.py --type xlsx --output "表格.xlsx" --data data.csv
"""

import argparse
import sys
from pathlib import Path

def create_docx(output: str, content: str, template: str = None):
    """创建 Word 文档"""
    try:
        from docx import Document
        from docx.shared import Pt, Cm
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        
        if template and Path(template).exists():
            doc = Document(template)
        else:
            doc = Document()
        
        # 添加标题
        title = content.split('\n')[0] if '\n' in content else "文档标题"
        heading = doc.add_heading(title, level=1)
        heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # 添加正文
        body = '\n'.join(content.split('\n')[1:]) if '\n' in content else content
        for paragraph in body.split('\n\n'):
            if paragraph.strip():
                p = doc.add_paragraph(paragraph.strip())
                p.paragraph_format.line_spacing = Cm(0.5)
        
        # 保存
        doc.save(output)
        print(f"[OK] 已创建 Word 文档：{output}")
        return True
    except ImportError:
        print("错误：需要安装 python-docx 库")
        print("运行：pip install python-docx")
        return False
    except Exception as e:
        print(f"错误：{e}")
        return False


def create_xlsx(output: str, data_file: str = None):
    """创建 Excel 表格"""
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment, PatternFill
        
        wb = Workbook()
        ws = wb.active
        ws.title = "数据表"
        
        # 设置表头样式
        header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        header_font = Font(bold=True, color="FFFFFF")
        
        if data_file and Path(data_file).exists():
            # 从 CSV 导入数据
            import csv
            with open(data_file, 'r', encoding='utf-8') as f:
                reader = csv.reader(f)
                for row_idx, row in enumerate(reader, 1):
                    for col_idx, value in enumerate(row, 1):
                        cell = ws.cell(row=row_idx, column=col_idx, value=value)
                        if row_idx == 1:
                            cell.fill = header_fill
                            cell.font = header_font
                            cell.alignment = Alignment(horizontal='center')
        else:
            # 创建示例表格
            headers = ["序号", "项目名称", "负责人", "状态", "备注"]
            for col, header in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col, value=header)
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = Alignment(horizontal='center')
            
            # 示例数据
            sample_data = [
                [1, "项目 A", "张三", "进行中", ""],
                [2, "项目 B", "李四", "已完成", "提前交付"],
                [3, "项目 C", "王五", "未开始", ""],
            ]
            for row_idx, row_data in enumerate(sample_data, 2):
                for col_idx, value in enumerate(row_data, 1):
                    ws.cell(row=row_idx, column=col_idx, value=value)
        
        # 自动调整列宽
        for column in ws.columns:
            max_length = 0
            column_letter = column[0].column_letter
            for cell in column:
                try:
                    if len(str(cell.value)) > max_length:
                        max_length = len(str(cell.value))
                except:
                    pass
            adjusted_width = min(max_length + 2, 50)
            ws.column_dimensions[column_letter].width = adjusted_width
        
        wb.save(output)
        print(f"[OK] 已创建 Excel 表格：{output}")
        return True
    except ImportError:
        print("错误：需要安装 openpyxl 库")
        print("运行：pip install openpyxl")
        return False
    except Exception as e:
        print(f"错误：{e}")
        return False


def create_pptx(output: str, template: str = None, content: str = None):
    """创建 PPT 演示文稿"""
    try:
        from pptx import Presentation
        from pptx.util import Cm, Pt
        from pptx.enum.text import PP_ALIGN
        
        if template and Path(template).exists():
            prs = Presentation(template)
        else:
            prs = Presentation()
        
        # 如果没有内容，创建示例幻灯片
        if not content:
            content = "标题\n\n第一点\n第二点\n第三点"
        
        slides_data = content.split('---')
        
        for i, slide_content in enumerate(slides_data):
            if i == 0:
                # 标题页
                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                lines = slide_content.strip().split('\n')
                slide.shapes.title.text = lines[0] if lines else "标题"
                if len(lines) > 1:
                    slide.placeholders[1].text = '\n'.join(lines[1:])
            else:
                # 内容页
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                lines = slide_content.strip().split('\n')
                slide.shapes.title.text = lines[0] if lines else "章节标题"
                if len(lines) > 1:
                    body = slide.placeholders[1]
                    tf = body.text_frame
                    tf.text = '\n'.join(lines[1:])
        
        prs.save(output)
        print(f"[OK] 已创建 PPT 演示文稿：{output}")
        return True
    except ImportError:
        print("错误：需要安装 python-pptx 库")
        print("运行：pip install python-pptx")
        return False
    except Exception as e:
        print(f"错误：{e}")
        return False


def create_pdf(output: str, content: str):
    """创建 PDF 文件"""
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfgen import canvas
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        from reportlab.lib.units import cm
        
        # 尝试注册中文字体
        font_paths = [
            r"C:\Windows\Fonts\simhei.ttf",  # 黑体
            r"C:\Windows\Fonts\simsun.ttc",  # 宋体
            r"C:\Windows\Fonts\msyh.ttc",    # 微软雅黑
        ]
        
        c = canvas.Canvas(output, pagesize=A4)
        width, height = A4
        
        # 注册第一个可用的中文字体
        font_registered = False
        for font_path in font_paths:
            if Path(font_path).exists():
                try:
                    pdfmetrics.registerFont(TTFont('Chinese', font_path))
                    font_registered = True
                    break
                except:
                    pass
        
        if not font_registered:
            print("警告：未找到中文字体，PDF 可能无法正确显示中文")
        
        # 设置字体
        if font_registered:
            c.setFont('Chinese', 12)
        else:
            c.setFont('Helvetica', 12)
        
        # 添加内容
        y_position = height - 2*cm
        lines = content.split('\n')
        
        for line in lines:
            if y_position < 2*cm:
                c.showPage()
                if font_registered:
                    c.setFont('Chinese', 12)
                else:
                    c.setFont('Helvetica', 12)
                y_position = height - 2*cm
            
            # 处理中文编码
            try:
                c.drawString(2*cm, y_position, line)
            except:
                # 如果中文显示失败，尝试编码转换
                try:
                    c.drawString(2*cm, y_position, line.encode('utf-8').decode('latin-1'))
                except:
                    c.drawString(2*cm, y_position, "[中文内容]")
            
            y_position -= 0.5*cm
        
        c.save()
        print(f"[OK] 已创建 PDF 文件：{output}")
        return True
    except ImportError:
        print("错误：需要安装 reportlab 库")
        print("运行：pip install reportlab")
        return False
    except Exception as e:
        print(f"错误：{e}")
        return False


def main():
    parser = argparse.ArgumentParser(description='文件创建工具 - 支持 docx, xlsx, pptx, pdf')
    parser.add_argument('--type', required=True, choices=['docx', 'xlsx', 'pptx', 'pdf'],
                        help='文件类型')
    parser.add_argument('--output', required=True, help='输出文件名')
    parser.add_argument('--content', help='文件内容')
    parser.add_argument('--data', help='数据文件 (CSV 格式，用于 xlsx)')
    parser.add_argument('--template', help='模板文件路径')
    
    args = parser.parse_args()
    
    # 如果没有指定内容，从 stdin 读取
    content = args.content
    if not content and not args.data:
        print("请输入内容 (以 EOF 结束):")
        content = sys.stdin.read()
    
    # 根据类型创建文件
    if args.type == 'docx':
        create_docx(args.output, content or "", args.template)
    elif args.type == 'xlsx':
        create_xlsx(args.output, args.data)
    elif args.type == 'pptx':
        create_pptx(args.output, args.template, content)
    elif args.type == 'pdf':
        create_pdf(args.output, content or "")


if __name__ == '__main__':
    main()
