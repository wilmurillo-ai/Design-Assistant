#!/usr/bin/env python3
import argparse
import csv
import json
from pathlib import Path

from pptx import Presentation


def normalize_shape_text(text: str) -> str:
    lines = [line.strip() for line in text.splitlines()]
    lines = [line for line in lines if line]
    return " | ".join(lines)


def extract_outline(pptx_path: Path):
    prs = Presentation(str(pptx_path))
    slides = []
    for index, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                normalized = normalize_shape_text(shape.text)
                if normalized:
                    texts.append(normalized)
        slides.append(
            {
                "slide": index,
                "text": " || ".join(texts),
            }
        )
    return slides


def write_csv(slides, output: Path):
    with output.open("w", encoding="utf-8", newline="") as fh:
        writer = csv.DictWriter(fh, fieldnames=["slide", "text"])
        writer.writeheader()
        writer.writerows(slides)


def write_json(slides, output: Path):
    with output.open("w", encoding="utf-8") as fh:
        json.dump(slides, fh, ensure_ascii=False, indent=2)


def main():
    parser = argparse.ArgumentParser(description="Extract slide text outline from a PPTX file.")
    parser.add_argument("--pptx", required=True, help="Path to the source PPTX file.")
    parser.add_argument("--out", required=True, help="Path to the output CSV or JSON file.")
    args = parser.parse_args()

    pptx_path = Path(args.pptx).expanduser().resolve()
    output = Path(args.out).expanduser().resolve()

    if not pptx_path.exists():
        raise SystemExit(f"PPTX not found: {pptx_path}")

    slides = extract_outline(pptx_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    if output.suffix.lower() == ".json":
        write_json(slides, output)
    else:
        write_csv(slides, output)

    print(f"OUTLINE_FILE={output}")
    print(f"SLIDE_COUNT={len(slides)}")


if __name__ == "__main__":
    main()
