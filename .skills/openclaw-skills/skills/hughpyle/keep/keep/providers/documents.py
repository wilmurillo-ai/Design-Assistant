"""
Document providers for fetching content from various URI schemes.
"""

from pathlib import Path

from .base import Document, DocumentProvider, get_registry


def extract_html_text(html_content: str) -> str:
    """
    Extract readable text from HTML, removing scripts and styles.

    Used by both FileDocumentProvider and HttpDocumentProvider to ensure
    consistent content regularization for embedding and summarization.

    Args:
        html_content: Raw HTML string

    Returns:
        Extracted text with whitespace normalized

    Raises:
        ImportError: If beautifulsoup4 is not installed
    """
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(html_content, "html.parser")

    # Remove script and style elements
    for script in soup(["script", "style"]):
        script.decompose()

    # Get text
    text = soup.get_text()

    # Clean up whitespace
    lines = (line.strip() for line in text.splitlines())
    chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
    return '\n'.join(chunk for chunk in chunks if chunk)


class FileDocumentProvider:
    """
    Fetches documents from the local filesystem.

    Supports file:// URIs and attempts to detect content type from extension.
    Performs text extraction for PDF and HTML files.
    """

    # Default max file size: 100MB
    MAX_FILE_SIZE = 100_000_000

    def __init__(self, max_size: int | None = None):
        self.max_size = max_size or self.MAX_FILE_SIZE

    EXTENSION_TYPES = {
        ".md": "text/markdown",
        ".markdown": "text/markdown",
        ".txt": "text/plain",
        ".py": "text/x-python",
        ".js": "text/javascript",
        ".ts": "text/typescript",
        ".json": "application/json",
        ".yaml": "text/yaml",
        ".yml": "text/yaml",
        ".html": "text/html",
        ".htm": "text/html",
        ".css": "text/css",
        ".xml": "application/xml",
        ".rst": "text/x-rst",
        ".pdf": "application/pdf",
        # Office documents
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        # Audio
        ".mp3": "audio/mpeg",
        ".flac": "audio/flac",
        ".aiff": "audio/aiff",
        ".aif": "audio/aiff",
        ".ogg": "audio/ogg",
        ".wav": "audio/wav",
        ".m4a": "audio/mp4",
        ".alac": "audio/mp4",
        ".wma": "audio/x-ms-wma",
        # Images
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".png": "image/png",
        ".tiff": "image/tiff",
        ".tif": "image/tiff",
        ".webp": "image/webp",
    }
    
    def supports(self, uri: str) -> bool:
        """Check if this is a file:// URI or bare path."""
        return uri.startswith("file://") or uri.startswith("/")
    
    def fetch(self, uri: str) -> Document:
        """Read file content from the filesystem with text extraction for PDF/HTML."""
        # Normalize to path
        if uri.startswith("file://"):
            path_str = uri.removeprefix("file://")
        else:
            path_str = uri

        path = Path(path_str).resolve()

        if not path.exists():
            raise IOError(f"File not found: {path}")

        if not path.is_file():
            raise IOError(f"Not a file: {path}")

        # Reject paths outside user's home directory as a safety boundary
        home = Path.home().resolve()
        if not path.is_relative_to(home):
            raise IOError(f"Path traversal blocked: {path} is outside home directory")

        # Check file size before processing
        file_size = path.stat().st_size
        if file_size > self.max_size:
            raise IOError(
                f"File too large: {file_size:,} bytes "
                f"(limit: {self.max_size:,} bytes). "
                f"Configure max_file_size in store config to increase."
            )

        # Detect content type
        suffix = path.suffix.lower()
        content_type = self.EXTENSION_TYPES.get(suffix, "text/plain")

        # Extract text based on file type
        extracted_tags: dict[str, str] | None = None
        if suffix == ".pdf":
            content = self._extract_pdf_text(path)
        elif suffix in (".html", ".htm"):
            content = self._extract_html_text(path)
        elif suffix in (".docx",):
            content, extracted_tags = self._extract_docx(path)
        elif suffix in (".pptx",):
            content, extracted_tags = self._extract_pptx(path)
        elif content_type and content_type.startswith("audio/"):
            content, extracted_tags = self._extract_audio_metadata(path)
        elif content_type and content_type.startswith("image/"):
            content, extracted_tags = self._extract_image_metadata(path)
        else:
            # Read as plain text
            try:
                content = path.read_text(encoding="utf-8")
            except UnicodeDecodeError:
                raise IOError(f"Cannot read file as text: {path}")

        # Gather metadata
        stat = path.stat()
        metadata = {
            "size": stat.st_size,
            "modified": stat.st_mtime,
            "name": path.name,
        }

        return Document(
            uri=f"file://{path.resolve()}",  # Normalize to absolute
            content=content,
            content_type=content_type,
            metadata=metadata,
            tags=extracted_tags,
        )

    def _extract_docx(self, path: Path) -> tuple[str, dict[str, str]]:
        """Extract text and metadata from DOCX file."""
        try:
            from docx import Document as DocxDocument
        except ImportError:
            raise IOError(
                f"DOCX support requires 'python-docx' library. "
                f"Install with: pip install python-docx\n"
                f"Cannot read DOCX: {path}"
            )

        try:
            doc = DocxDocument(path)
            parts = []

            # Extract paragraph text
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    parts.append(text)

            # Extract table text
            for table in doc.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))

            if not parts:
                raise IOError(f"No text extracted from DOCX: {path}")

            content = "\n\n".join(parts)

            # Extract metadata tags
            tags: dict[str, str] = {}
            props = doc.core_properties
            if props.author:
                tags["author"] = props.author
            if props.title:
                tags["title"] = props.title

            return content, tags or None
        except ImportError:
            raise
        except IOError:
            raise
        except Exception as e:
            raise IOError(f"Failed to extract text from DOCX {path}: {e}")

    def _extract_pptx(self, path: Path) -> tuple[str, dict[str, str]]:
        """Extract text and metadata from PPTX file."""
        try:
            from pptx import Presentation
        except ImportError:
            raise IOError(
                f"PPTX support requires 'python-pptx' library. "
                f"Install with: pip install python-pptx\n"
                f"Cannot read PPTX: {path}"
            )

        try:
            prs = Presentation(path)
            parts = []

            for i, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                slide_texts.append(text)
                if slide_texts:
                    parts.append(f"Slide {i}:\n" + "\n".join(slide_texts))

                # Extract notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        parts.append(f"Notes:\n{notes}")

            if not parts:
                raise IOError(f"No text extracted from PPTX: {path}")

            content = "\n\n".join(parts)

            # Extract metadata tags
            tags: dict[str, str] = {}
            props = prs.core_properties
            if props.author:
                tags["author"] = props.author
            if props.title:
                tags["title"] = props.title

            return content, tags or None
        except ImportError:
            raise
        except IOError:
            raise
        except Exception as e:
            raise IOError(f"Failed to extract text from PPTX {path}: {e}")

    def _extract_audio_metadata(self, path: Path) -> tuple[str, dict[str, str]]:
        """Extract metadata from audio file as structured text."""
        try:
            from tinytag import TinyTag
        except ImportError:
            raise IOError(
                f"Audio metadata support requires 'tinytag' library. "
                f"Install with: pip install tinytag\n"
                f"Cannot read audio metadata: {path}"
            )

        try:
            tag = TinyTag.get(path)
            lines = []
            tags: dict[str, str] = {}

            if tag.title:
                lines.append(f"Title: {tag.title}")
                tags["title"] = tag.title
            if tag.artist:
                lines.append(f"Artist: {tag.artist}")
                tags["artist"] = tag.artist
            if tag.album:
                lines.append(f"Album: {tag.album}")
                tags["album"] = tag.album
            if tag.albumartist:
                lines.append(f"Album Artist: {tag.albumartist}")
            if tag.track:
                lines.append(f"Track: {tag.track}")
            if tag.genre:
                lines.append(f"Genre: {tag.genre}")
                tags["genre"] = tag.genre
            if tag.year:
                lines.append(f"Year: {tag.year}")
                tags["year"] = str(tag.year)
            if tag.duration:
                mins, secs = divmod(int(tag.duration), 60)
                lines.append(f"Duration: {mins}:{secs:02d}")
            if tag.bitrate:
                lines.append(f"Bitrate: {int(tag.bitrate)} kbps")
            if tag.samplerate:
                lines.append(f"Sample Rate: {tag.samplerate} Hz")
            if tag.comment:
                lines.append(f"Comment: {tag.comment}")

            if not lines:
                lines.append(f"Audio file: {path.name}")

            content = "\n".join(lines)
            return content, tags or None
        except ImportError:
            raise
        except Exception as e:
            raise IOError(f"Failed to extract audio metadata from {path}: {e}")

    def _extract_image_metadata(self, path: Path) -> tuple[str, dict[str, str]]:
        """Extract EXIF metadata from image file as structured text."""
        try:
            from PIL import Image
            from PIL.ExifTags import TAGS as EXIF_TAGS
        except ImportError:
            raise IOError(
                f"Image metadata support requires 'Pillow' library. "
                f"Install with: pip install Pillow\n"
                f"Cannot read image metadata: {path}"
            )

        try:
            img = Image.open(path)
            lines = []
            tags: dict[str, str] = {}

            # Basic image info
            w, h = img.size
            lines.append(f"Dimensions: {w}x{h}")
            tags["dimensions"] = f"{w}x{h}"
            lines.append(f"Format: {img.format}")

            # EXIF data
            exif = img.getexif()
            if exif:
                # Camera model
                model = exif.get(0x0110)  # Model
                make = exif.get(0x010F)  # Make
                if model:
                    camera = f"{make} {model}".strip() if make else model
                    lines.append(f"Camera: {camera}")
                    tags["camera"] = camera

                # Date taken
                date = exif.get(0x9003) or exif.get(0x0132)  # DateTimeOriginal or DateTime
                if date:
                    lines.append(f"Date: {date}")
                    tags["date"] = date

                # Focal length
                focal = exif.get(0x920A)  # FocalLength
                if focal:
                    if hasattr(focal, 'numerator'):
                        focal_val = focal.numerator / focal.denominator
                        lines.append(f"Focal Length: {focal_val:.0f}mm")
                    else:
                        lines.append(f"Focal Length: {focal}mm")

                # ISO
                iso = exif.get(0x8827)  # ISOSpeedRatings
                if iso:
                    lines.append(f"ISO: {iso}")

                # Exposure time
                exposure = exif.get(0x829A)  # ExposureTime
                if exposure:
                    if hasattr(exposure, 'numerator'):
                        if exposure.numerator == 1:
                            lines.append(f"Exposure: 1/{exposure.denominator}s")
                        else:
                            lines.append(f"Exposure: {exposure.numerator}/{exposure.denominator}s")
                    else:
                        lines.append(f"Exposure: {exposure}s")

                # F-number
                fnumber = exif.get(0x829D)  # FNumber
                if fnumber:
                    if hasattr(fnumber, 'numerator'):
                        fval = fnumber.numerator / fnumber.denominator
                        lines.append(f"Aperture: f/{fval:.1f}")
                    else:
                        lines.append(f"Aperture: f/{fnumber}")

            img.close()

            if not lines:
                lines.append(f"Image file: {path.name}")

            content = "\n".join(lines)
            return content, tags or None
        except ImportError:
            raise
        except Exception as e:
            raise IOError(f"Failed to extract image metadata from {path}: {e}")

    def _extract_pdf_text(self, path: Path) -> str:
        """Extract text from PDF file."""
        try:
            from pypdf import PdfReader
        except ImportError:
            raise IOError(
                f"PDF support requires 'pypdf' library. "
                f"Install with: pip install pypdf\n"
                f"Cannot read PDF: {path}"
            )

        try:
            reader = PdfReader(path)
            text_parts = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)

            if not text_parts:
                raise IOError(f"No text extracted from PDF: {path}")

            return "\n\n".join(text_parts)
        except Exception as e:
            raise IOError(f"Failed to extract text from PDF {path}: {e}")

    def _extract_html_text(self, path: Path) -> str:
        """Extract text from HTML file."""
        try:
            html_content = path.read_text(encoding="utf-8")
            return extract_html_text(html_content)
        except ImportError:
            raise IOError(
                f"HTML text extraction requires 'beautifulsoup4' library. "
                f"Install with: pip install beautifulsoup4\n"
                f"Cannot extract text from HTML: {path}"
            )
        except Exception as e:
            raise IOError(f"Failed to extract text from HTML {path}: {e}")


class HttpDocumentProvider:
    """
    Fetches documents from HTTP/HTTPS URLs.
    
    Requires the `requests` library (optional dependency).
    """
    
    def __init__(self, timeout: int = 30, max_size: int = 10_000_000):
        """
        Args:
            timeout: Request timeout in seconds
            max_size: Maximum content size in bytes
        """
        self.timeout = timeout
        self.max_size = max_size
    
    def supports(self, uri: str) -> bool:
        """Check if this is an HTTP(S) URL."""
        return uri.startswith("http://") or uri.startswith("https://")
    
    @staticmethod
    def _is_private_url(uri: str) -> bool:
        """Check if URL targets a private/internal network address.

        Note: DNS resolution here is inherently TOCTOU — the hostname could
        resolve to a different address by the time requests.get() connects.
        Sufficient for CLI use; a hosted service should enforce this at the
        network layer (firewall/VPC rules) rather than relying on client checks.
        """
        from urllib.parse import urlparse
        import ipaddress
        import socket

        parsed = urlparse(uri)
        hostname = parsed.hostname
        if not hostname:
            return True

        # Block known metadata endpoints and localhost
        if hostname in ("metadata.google.internal",):
            return True

        try:
            addr = ipaddress.ip_address(hostname)
            return (addr.is_private or addr.is_loopback or addr.is_link_local
                    or addr.is_reserved or addr.is_unspecified or addr.is_multicast)
        except ValueError:
            pass  # Not an IP literal — resolve it

        try:
            for _, _, _, _, sockaddr in socket.getaddrinfo(hostname, None):
                addr = ipaddress.ip_address(sockaddr[0])
                if (addr.is_private or addr.is_loopback or addr.is_link_local
                        or addr.is_reserved or addr.is_unspecified or addr.is_multicast):
                    return True
        except socket.gaierror:
            pass  # DNS failure will be caught by requests

        return False

    _MAX_REDIRECTS = 5

    def fetch(self, uri: str) -> Document:
        """Fetch content from HTTP URL with text extraction for HTML."""
        if self._is_private_url(uri):
            raise IOError(f"Blocked request to private/internal address: {uri}")

        try:
            import requests
        except ImportError:
            raise RuntimeError("HTTP document fetching requires 'requests' library")

        from keep import __version__

        # Follow redirects manually so each hop is validated against SSRF
        target = uri
        for _ in range(self._MAX_REDIRECTS):
            resp = requests.get(
                target,
                timeout=self.timeout,
                headers={"User-Agent": f"keep/{__version__}"},
                stream=True,
                allow_redirects=False,
            )
            if resp.is_redirect:
                target = resp.headers.get("Location", "")
                if not target.startswith(("http://", "https://")):
                    raise IOError(f"Redirect to unsupported scheme: {target}")
                if self._is_private_url(target):
                    raise IOError(f"Redirect to private/internal address blocked: {target}")
                resp.close()
                continue
            break
        else:
            raise IOError(f"Too many redirects fetching {uri}")

        try:
            with resp:
                resp.raise_for_status()

                # Check declared size
                content_length = resp.headers.get("content-length")
                if content_length:
                    try:
                        if int(content_length) > self.max_size:
                            raise IOError(f"Content too large: {content_length} bytes")
                    except ValueError:
                        pass  # Malformed header — enforce via iter_content below

                # Read content in chunks with enforced size limit
                chunks: list[bytes] = []
                downloaded = 0
                for chunk in resp.iter_content(chunk_size=65536):
                    downloaded += len(chunk)
                    if downloaded > self.max_size:
                        chunks.append(chunk[:self.max_size - (downloaded - len(chunk))])
                        break
                    chunks.append(chunk)
                raw = b"".join(chunks)

                # Decode using the response's detected encoding
                encoding = resp.encoding or "utf-8"
                content = raw.decode(encoding, errors="replace")

                # Get content type
                content_type = resp.headers.get("content-type", "text/plain")
                if ";" in content_type:
                    content_type = content_type.split(";")[0].strip()

                # Extract text from HTML content
                if content_type == "text/html":
                    try:
                        content = extract_html_text(content)
                    except ImportError:
                        # Graceful fallback: use raw HTML if bs4 not installed
                        pass

                return Document(
                    uri=uri,
                    content=content,
                    content_type=content_type,
                    metadata={
                        "status_code": resp.status_code,
                        "headers": dict(resp.headers),
                    },
                )
        except requests.RequestException as e:
            raise IOError(f"Failed to fetch {uri}: {e}")


class CompositeDocumentProvider:
    """
    Combines multiple document providers, delegating to the appropriate one.
    
    This is the default provider used by Keeper.
    """
    
    def __init__(self, providers: list[DocumentProvider] | None = None):
        """
        Args:
            providers: List of providers to try. If None, uses defaults.
        """
        if providers is None:
            self._providers = [
                FileDocumentProvider(),
                HttpDocumentProvider(),
            ]
        else:
            self._providers = list(providers)
    
    def supports(self, uri: str) -> bool:
        """Check if any provider supports this URI."""
        return any(p.supports(uri) for p in self._providers)
    
    def fetch(self, uri: str) -> Document:
        """Fetch using the first provider that supports this URI."""
        for provider in self._providers:
            if provider.supports(uri):
                return provider.fetch(uri)
        
        raise ValueError(f"No provider supports URI: {uri}")
    
    def add_provider(self, provider: DocumentProvider) -> None:
        """Add a provider to the list (checked first)."""
        self._providers.insert(0, provider)


# Register providers
_registry = get_registry()
_registry.register_document("file", FileDocumentProvider)
_registry.register_document("http", HttpDocumentProvider)
_registry.register_document("composite", CompositeDocumentProvider)
