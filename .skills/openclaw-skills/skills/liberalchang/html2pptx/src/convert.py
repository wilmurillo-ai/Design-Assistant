#!/usr/bin/env python3
"""
HTML to PPTX Converter v3.1
优化布局、图标、色彩还原
"""

import sys
import re
from pathlib import Path
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


def rgb_to_color(rgb_tuple):
    if not rgb_tuple:
        return RGBColor(255, 255, 255)
    return RGBColor(*rgb_tuple)


def extract_text(element):
    if not element:
        return ""
    text = element.get_text(separator=' ', strip=True)
    text = re.sub(r'\s+', ' ', text)
    return text.strip()


def get_slide_type(slide_elem):
    classes = slide_elem.get('class', [])
    type_map = {
        'title-slide': 'title',
        'closing-slide': 'closing',
        'overview-slide': 'overview',
        'pain-points-slide': 'pain_points',
        'solution-slide': 'solution',
        'features-slide': 'features',
        'workflow-slide': 'workflow',
        'tech-slide': 'tech',
        'timeline-slide': 'timeline',
        'business-slide': 'business',
    }
    for cls, stype in type_map.items():
        if cls in classes:
            return stype
    return 'content'


def add_slide_content(slide, soup, slide_type, colors):
    """根据类型添加幻灯片内容"""

    # 设置背景
    background = slide.background
    fill = background.fill
    fill.solid()

    if slide_type in ['title', 'closing']:
        fill.fore_color.rgb = rgb_to_color(colors['bg_primary'])
    elif slide_type in ['overview', 'timeline', 'workflow']:
        fill.fore_color.rgb = rgb_to_color(colors['bg_secondary'])
    else:
        fill.fore_color.rgb = rgb_to_color(colors['bg_primary'])

    # 提取标题
    section_title = soup.find(class_='section-title')
    main_title = soup.find('h1', class_='main-title') or soup.find('h1')
    subtitle = soup.find('p', class_='subtitle')

    y_pos = Inches(0.4)

    # 添加主标题（标题页）
    if slide_type == 'title' and main_title:
        title_text = extract_text(main_title)
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5),
            Inches(12.333), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = rgb_to_color(colors['accent_cyan'])
        p.alignment = PP_ALIGN.CENTER

        # 副标题
        if subtitle:
            sub_text = extract_text(subtitle)
            sub_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.2),
                Inches(12.333), Inches(0.8)
            )
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = sub_text
            p.font.size = Pt(24)
            p.font.color.rgb = rgb_to_color(colors['text_secondary'])
            p.alignment = PP_ALIGN.CENTER
        return

    # 添加章节标题
    if section_title:
        title_text = extract_text(section_title)
        title_box = slide.shapes.add_textbox(
            Inches(0.5), y_pos,
            Inches(12.333), Inches(1)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = rgb_to_color(colors['text_primary'])
        y_pos = Inches(1.5)

    # 根据类型处理内容
    if slide_type == 'overview':
        # 2x2卡片网格
        cards = soup.find_all(class_='overview-card')[:4]
        card_width = Inches(5.8)
        card_height = Inches(2.2)
        gap = Inches(0.4)
        start_x = Inches(0.5)
        start_y = Inches(1.6)

        for idx, card in enumerate(cards):
            col = idx % 2
            row = idx // 2
            x = start_x + col * (card_width + gap)
            y = start_y + row * (card_height + gap)

            # 卡片背景
            card_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_width, card_height
            )
            card_shape.fill.solid()
            card_shape.fill.fore_color.rgb = rgb_to_color((17, 24, 39))
            card_shape.line.color.rgb = rgb_to_color(colors['accent_cyan'])

            # 提取内容
            icon = extract_text(card.find(class_='card-icon'))
            title = extract_text(card.find(class_='card-title'))
            desc = extract_text(card.find(class_='card-desc'))

            # 图标
            if icon:
                icon_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), Inches(0.8), Inches(0.6))
                tf = icon_box.text_frame
                tf.paragraphs[0].text = icon
                tf.paragraphs[0].font.size = Pt(28)

            # 标题
            if title:
                t_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.8), card_width - Inches(0.4), Inches(0.5))
                tf = t_box.text_frame
                tf.paragraphs[0].text = title
                tf.paragraphs[0].font.size = Pt(18)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = rgb_to_color(colors['text_primary'])

            # 描述
            if desc:
                d_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.3), card_width - Inches(0.4), Inches(0.8))
                tf = d_box.text_frame
                tf.word_wrap = True
                tf.paragraphs[0].text = desc
                tf.paragraphs[0].font.size = Pt(13)
                tf.paragraphs[0].font.color.rgb = rgb_to_color(colors['text_secondary'])

    elif slide_type == 'pain_points':
        # 垂直列表
        items = soup.find_all(class_='pain-item')[:6]
        for idx, item in enumerate(items):
            y = Inches(1.5) + idx * Inches(0.7)
            icon = extract_text(item.find(class_='pain-icon'))
            text = extract_text(item.find(class_='pain-text'))

            # 图标
            if icon:
                i_box = slide.shapes.add_textbox(Inches(0.5), y, Inches(0.5), Inches(0.5))
                i_box.text_frame.paragraphs[0].text = icon
                i_box.text_frame.paragraphs[0].font.size = Pt(20)

            # 文本
            if text:
                t_box = slide.shapes.add_textbox(Inches(1.1), y, Inches(11.7), Inches(0.6))
                tf = t_box.text_frame
                tf.word_wrap = True
                tf.paragraphs[0].text = text
                tf.paragraphs[0].font.size = Pt(16)
                tf.paragraphs[0].font.color.rgb = rgb_to_color(colors['text_primary'])

    elif slide_type == 'solution':
        # 三步骤横向
        steps = soup.find_all(class_='solution-step')[:3]
        step_width = Inches(3.9)
        step_height = Inches(3.5)
        gap = Inches(0.3)
        start_x = Inches(0.5)
        start_y = Inches(2.0)

        for idx, step in enumerate(steps):
            x = start_x + idx * (step_width + gap)

            # 卡片
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, step_width, step_height)
            card.fill.solid()
            card.fill.fore_color.rgb = rgb_to_color((17, 24, 39))
            card.line.color.rgb = rgb_to_color(colors['accent_cyan'])

            # 编号
            num = extract_text(step.find(class_='step-number'))
            if num:
                num_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.55), start_y + Inches(0.3), Inches(0.8), Inches(0.8))
                num_shape.fill.solid()
                num_shape.fill.fore_color.rgb = rgb_to_color(colors['accent_cyan'])
                num_shape.line.fill.background()
                tf = num_shape.text_frame
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                tf.paragraphs[0].text = num
                tf.paragraphs[0].font.size = Pt(24)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = rgb_to_color(colors['bg_primary'])

            # 标题
            title = extract_text(step.find(class_='step-title'))
            if title:
                t_box = slide.shapes.add_textbox(x + Inches(0.2), start_y + Inches(1.3), step_width - Inches(0.4), Inches(0.6))
                tf = t_box.text_frame
                tf.paragraphs[0].text = title
                tf.paragraphs[0].font.size = Pt(20)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = rgb_to_color(colors['text_primary'])
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER

            # 描述
            desc = extract_text(step.find(class_='step-desc'))
            if desc:
                d_box = slide.shapes.add_textbox(x + Inches(0.2), start_y + Inches(2.0), step_width - Inches(0.4), Inches(1.2))
                tf = d_box.text_frame
                tf.word_wrap = True
                tf.paragraphs[0].text = desc
                tf.paragraphs[0].font.size = Pt(14)
                tf.paragraphs[0].font.color.rgb = rgb_to_color(colors['text_secondary'])
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    elif slide_type == 'features':
        # 4列网格
        features = soup.find_all(class_='feature-card')[:12]
        card_width = Inches(2.9)
        card_height = Inches(1.3)
        gap = Inches(0.3)
        start_x = Inches(0.5)
        start_y = Inches(1.3)

        for idx, feature in enumerate(features):
            col = idx % 4
            row = idx // 4
            x = start_x + col * (card_width + gap)
            y = start_y + row * (card_height + gap)

            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, card_height)
            card.fill.solid()
            card.fill.fore_color.rgb = rgb_to_color((17, 24, 39))
            card.line.color.rgb = rgb_to_color(colors['accent_cyan'])

            name = extract_text(feature.find(class_='feature-name'))
            fid = extract_text(feature.find(class_='feature-id'))

            if name:
                n_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.15), card_width - Inches(0.2), Inches(0.5))
                tf = n_box.text_frame
                tf.word_wrap = True
                tf.paragraphs[0].text = name
                tf.paragraphs[0].font.size = Pt(13)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = rgb_to_color(colors['text_primary'])

            if fid:
                i_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.7), card_width - Inches(0.2), Inches(0.4))
                tf = i_box.text_frame
                tf.paragraphs[0].text = fid
                tf.paragraphs[0].font.size = Pt(11)
                tf.paragraphs[0].font.color.rgb = rgb_to_color(colors['accent_cyan'])

    elif slide_type == 'workflow':
        # 垂直步骤列表
        steps = soup.find_all(class_='workflow-step')[:9]
        for idx, step in enumerate(steps):
            y = Inches(1.2) + idx * Inches(0.55)
            icon = extract_text(step.find(class_='workflow-icon'))
            text = extract_text(step.find(class_='workflow-text'))

            # 图标背景
            icon_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y, Inches(0.5), Inches(0.45))
            icon_bg.fill.solid()
            icon_bg.fill.fore_color.rgb = rgb_to_color(colors['accent_cyan'])
            icon_bg.line.fill.background()

            if icon:
                i_box = slide.shapes.add_textbox(Inches(0.5), y, Inches(0.5), Inches(0.45))
                tf = i_box.text_frame
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                tf.paragraphs[0].text = icon
                tf.paragraphs[0].font.size = Pt(16)

            if text:
                t_box = slide.shapes.add_textbox(Inches(1.1), y, Inches(11.7), Inches(0.5))
                tf = t_box.text_frame
                tf.word_wrap = True
                tf.paragraphs[0].text = text
                tf.paragraphs[0].font.size = Pt(14)
                tf.paragraphs[0].font.color.rgb = rgb_to_color(colors['text_primary'])

    elif slide_type == 'tech':
        # 技术层列表
        layers = soup.find_all(class_='tech-layer')[:5]
        for idx, layer in enumerate(layers):
            y = Inches(1.5) + idx * Inches(0.9)

            # 层背景
            layer_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y, Inches(12.333), Inches(0.8))
            layer_bg.fill.solid()
            layer_bg.fill.fore_color.rgb = rgb_to_color((17, 24, 39))
            layer_bg.line.color.rgb = rgb_to_color(colors['accent_cyan'])

            label = extract_text(layer.find(class_='layer-label'))
            content = extract_text(layer.find(class_='layer-content'))

            if label:
                l_box = slide.shapes.add_textbox(Inches(0.7), y + Inches(0.2), Inches(1.5), Inches(0.4))
                tf = l_box.text_frame
                tf.paragraphs[0].text = label
                tf.paragraphs[0].font.size = Pt(14)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = rgb_to_color(colors['accent_cyan'])

            if content:
                c_box = slide.shapes.add_textbox(Inches(2.3), y + Inches(0.2), Inches(10.0), Inches(0.4))
                tf = c_box.text_frame
                tf.word_wrap = True
                tf.paragraphs[0].text = content
                tf.paragraphs[0].font.size = Pt(14)
                tf.paragraphs[0].font.color.rgb = rgb_to_color(colors['text_primary'])

    elif slide_type == 'timeline':
        # 时间线阶段
        phases = soup.find_all(class_='timeline-phase')[:5]
        phase_colors = [(16, 185, 129), (59, 130, 246), (245, 158, 11), (139, 92, 246), (236, 72, 153)]

        for idx, phase in enumerate(phases):
            y = Inches(1.5) + idx * Inches(0.95)
            color = phase_colors[idx % len(phase_colors)]

            # 阶段背景
            phase_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y, Inches(12.333), Inches(0.85))
            phase_bg.fill.solid()
            phase_bg.fill.fore_color.rgb = rgb_to_color((17, 24, 39))
            phase_bg.line.color.rgb = rgb_to_color(color)
            phase_bg.line.width = Pt(2)

            duration = extract_text(phase.find(class_='phase-duration'))
            title = extract_text(phase.find(class_='phase-title'))
            desc = extract_text(phase.find(class_='phase-desc'))

            if duration:
                d_box = slide.shapes.add_textbox(Inches(0.7), y + Inches(0.25), Inches(1.0), Inches(0.4))
                tf = d_box.text_frame
                tf.paragraphs[0].text = duration
                tf.paragraphs[0].font.size = Pt(16)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = rgb_to_color(color)

            if title:
                t_box = slide.shapes.add_textbox(Inches(1.8), y + Inches(0.15), Inches(3.0), Inches(0.4))
                tf = t_box.text_frame
                tf.paragraphs[0].text = title
                tf.paragraphs[0].font.size = Pt(18)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = rgb_to_color(colors['text_primary'])

            if desc:
                de_box = slide.shapes.add_textbox(Inches(1.8), y + Inches(0.5), Inches(10.5), Inches(0.4))
                tf = de_box.text_frame
                tf.word_wrap = True
                tf.paragraphs[0].text = desc
                tf.paragraphs[0].font.size = Pt(12)
                tf.paragraphs[0].font.color.rgb = rgb_to_color(colors['text_secondary'])

    elif slide_type == 'business':
        # 2x2业务卡片
        cards = soup.find_all(class_='business-card')[:4]
        card_width = Inches(5.8)
        card_height = Inches(2.5)
        gap = Inches(0.4)
        start_x = Inches(0.5)
        start_y = Inches(1.5)

        for idx, card in enumerate(cards):
            col = idx % 2
            row = idx // 2
            x = start_x + col * (card_width + gap)
            y = start_y + row * (card_height + gap)

            # 卡片背景
            card_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, card_height)
            card_shape.fill.solid()
            card_shape.fill.fore_color.rgb = rgb_to_color((17, 24, 39))
            card_shape.line.color.rgb = rgb_to_color(colors['accent_cyan'])

            title = extract_text(card.find(class_='business-title'))
            if title:
                t_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), card_width - Inches(0.4), Inches(0.5))
                tf = t_box.text_frame
                tf.paragraphs[0].text = title
                tf.paragraphs[0].font.size = Pt(20)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = rgb_to_color(colors['text_primary'])

            # 列表项
            list_items = card.find_all('li')
            list_text = '• ' + '\n• '.join([extract_text(li) for li in list_items[:3]])
            if list_text:
                l_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.7), card_width - Inches(0.4), Inches(1.6))
                tf = l_box.text_frame
                tf.word_wrap = True
                tf.paragraphs[0].text = list_text
                tf.paragraphs[0].font.size = Pt(14)
                tf.paragraphs[0].font.color.rgb = rgb_to_color(colors['text_secondary'])

    elif slide_type == 'closing':
        # 结尾页
        closing_title = soup.find('h2', class_='closing-title')
        closing_sub = soup.find('p', class_='closing-subtitle')

        if closing_title:
            t_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
            tf = t_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = extract_text(closing_title)
            tf.paragraphs[0].font.size = Pt(48)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = rgb_to_color(colors['text_primary'])
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        if closing_sub:
            s_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
            tf = s_box.text_frame
            tf.paragraphs[0].text = extract_text(closing_sub)
            tf.paragraphs[0].font.size = Pt(24)
            tf.paragraphs[0].font.color.rgb = rgb_to_color(colors['text_secondary'])
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    else:
        # 默认内容页 - 提取所有文本
        content_elems = soup.find_all(['p', 'li'])
        y = y_pos
        for elem in content_elems[:8]:
            text = extract_text(elem)
            if text and len(text) > 5:
                t_box = slide.shapes.add_textbox(Inches(0.5), y, Inches(12.333), Inches(0.5))
                tf = t_box.text_frame
                tf.word_wrap = True
                tf.paragraphs[0].text = text
                tf.paragraphs[0].font.size = Pt(16)
                tf.paragraphs[0].font.color.rgb = rgb_to_color(colors['text_primary'])
                y += Inches(0.6)


def convert_html_to_pptx(html_path, output_path=None):
    """转换HTML到PPTX"""
    html_path = Path(html_path)
    if not html_path.exists():
        print(f"错误: 文件不存在 {html_path}")
        return False

    if not output_path:
        output_path = html_path.with_suffix('.pptx')
    else:
        output_path = Path(output_path)

    print(f"正在转换: {html_path}")
    print(f"输出到: {output_path}")

    # 读取HTML
    with open(html_path, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f.read(), 'html.parser')

    # 创建演示文稿
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 查找所有slide
    slides = soup.find_all('section', class_='slide')
    if not slides:
        slides = soup.find_all('div', class_='slide')

    if not slides:
        print("错误: 未找到幻灯片元素")
        return False

    print(f"找到 {len(slides)} 页幻灯片")

    # 颜色配置
    colors = {
        'bg_primary': (10, 15, 28),
        'bg_secondary': (17, 24, 39),
        'accent_cyan': (0, 255, 204),
        'accent_magenta': (255, 0, 170),
        'text_primary': (255, 255, 255),
        'text_secondary': (156, 163, 175),
    }

    # 处理每页
    for idx, slide_elem in enumerate(slides):
        print(f"  处理第 {idx + 1} 页...")
        slide_type = get_slide_type(slide_elem)

        # 创建空白幻灯片
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # 添加内容
        add_slide_content(slide, slide_elem, slide_type, colors)

    # 保存
    prs.save(output_path)
    print(f"\n✓ 转换完成: {output_path}")
    print(f"  共 {len(slides)} 页幻灯片")
    return True


def main():
    if len(sys.argv) < 2:
        print("用法: python convert.py <input.html> [output.pptx]")
        sys.exit(1)

    html_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else None

    success = convert_html_to_pptx(html_path, output_path)
    sys.exit(0 if success else 1)


if __name__ == "__main__":
    main()
