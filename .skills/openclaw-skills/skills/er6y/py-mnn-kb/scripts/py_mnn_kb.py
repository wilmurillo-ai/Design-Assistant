"""
py_mnn_kb.py - Python Knowledge Base for MNN Embedding
Fully aligns with Android OfflineAI RAG implementation.

Three public interfaces:
  1. build_kb(files, kb_name=None)    - Build KB from file list
  2. add_note(text, kb_name=None)     - Add a text note directly
  3. query_kb(prompt, kb_name=None)   - RAG retrieval, returns combined context string

Database schema mirrors Android KnowledgeGraphDatabase exactly:
  - documents:     id, collection, content, source, metadata, embedding(BLOB), created_at
  - entities:      id, collection, entity_text, entity_type, language, frequency,
                   avg_confidence, first_seen, last_seen
  - entity_edges:  id, collection, from_entity, to_entity, weight, chunk_ids, created_at,
                   updated_at
  - chunk_entities: chunk_id, entity_text, entity_type, confidence
  - metadata:      key, value

Embedding backend:
  - MNN Python package (MNN.llm / pyMNN) only
"""

import argparse
import contextlib
import json
import os
import re
import sqlite3
import struct
import sys
import time
import math
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from typing import List, Optional, Tuple, Dict, Set


def _build_android_embedding_chat_template() -> str:
    open2 = chr(123) + chr(123)
    close2 = chr(125) + chr(125)
    sq = chr(39)
    core = (
        open2
        + " messages | map(attribute="
        + sq
        + "content"
        + sq
        + ") | join("
        + sq
        + sq
        + ") "
        + close2
    )
    return core + "<|endoftext|>"


def _build_android_embedding_runtime_config(thread_num: int = 4) -> dict:
    return {
        "backend_type": "cpu",
        "memory": "low",
        "power": "high",
        "precision": "low",
        "thread_num": thread_num
    }


def _ensure_android_embedding_model_config(model_dir: Path) -> None:
    llm_cfg_path = model_dir / "llm_config.json"
    if not llm_cfg_path.exists():
        return
    with open(llm_cfg_path, "r", encoding="utf-8") as f:
        llm_cfg = json.load(f)
    expected_template = _build_android_embedding_chat_template()
    jinja_cfg = llm_cfg.get("jinja") or {}
    current_template = jinja_cfg.get("chat_template")
    current_eos = jinja_cfg.get("eos")
    if current_template == expected_template and current_eos == "<|im_end|>":
        return
    llm_cfg["jinja"] = {
        "chat_template": expected_template,
        "eos": "<|im_end|>"
    }
    with open(llm_cfg_path, "w", encoding="utf-8") as f:
        json.dump(llm_cfg, f, ensure_ascii=False, indent=4)
    print(f"[EmbeddingClient] Updated llm_config.json in place: {llm_cfg_path}", file=sys.stderr)


# ──────────────────────────────────────────────
# Config loader
# ──────────────────────────────────────────────

_DEFAULT_CONFIG_PATH = Path(__file__).parent.parent / "config.json"


def load_config(config_path=None) -> dict:
    """Load config.json, resolve relative paths against config file location."""
    path = Path(config_path) if config_path else _DEFAULT_CONFIG_PATH
    with open(path, "r", encoding="utf-8") as f:
        cfg = json.load(f)
    base = path.parent.resolve()

    def _resolve(p):
        if not p:
            return p
        r = Path(p)
        if not r.is_absolute():
            r = (base / r).resolve()
        return str(r)

    cfg["embedding"]["model_dir"] = _resolve(cfg["embedding"]["model_dir"])
    cfg["knowledge_base"]["storage_dir"] = _resolve(cfg["knowledge_base"]["storage_dir"])
    if cfg["graph_ner"].get("custom_dict_path"):
        cfg["graph_ner"]["custom_dict_path"] = _resolve(cfg["graph_ner"]["custom_dict_path"])
    if cfg["graph_ner"].get("stopwords_path"):
        cfg["graph_ner"]["stopwords_path"] = _resolve(cfg["graph_ner"]["stopwords_path"])
    return cfg


# ──────────────────────────────────────────────
# Model auto-download (ModelScope)
# ──────────────────────────────────────────────

_MODEL_FILES = {
    "Qwen3-Embedding-0.6B-MNN-int4": {
        "config.json":          "https://www.modelscope.cn/models/MNN/Qwen3-Embedding-0.6B-MNN/resolve/master/config.json",
        "configuration.json":   "https://www.modelscope.cn/models/MNN/Qwen3-Embedding-0.6B-MNN/resolve/master/configuration.json",
        "llm_config.json":      "https://www.modelscope.cn/models/MNN/Qwen3-Embedding-0.6B-MNN/resolve/master/llm_config.json",
        "embedding.mnn":        "https://www.modelscope.cn/models/MNN/Qwen3-Embedding-0.6B-MNN/resolve/master/embedding.mnn",
        "embedding.mnn.weight": "https://www.modelscope.cn/models/MNN/Qwen3-Embedding-0.6B-MNN/resolve/master/embedding.mnn.weight",
        "embeddings_int4.bin":  "https://www.modelscope.cn/models/MNN/Qwen3-Embedding-0.6B-MNN/resolve/master/embeddings_int4.bin",
        "tokenizer.txt":        "https://www.modelscope.cn/models/MNN/Qwen3-Embedding-0.6B-MNN/resolve/master/tokenizer.txt",
    }
}

_REQUIRED_MODEL_FILES = {"config.json", "llm_config.json", "embedding.mnn", "tokenizer.txt"}


def ensure_model_downloaded(model_dir: str) -> None:
    """
    Check if the embedding model exists; if not, auto-download from ModelScope.
    Only downloads missing files. Skips download if all required files are present.
    """
    model_path = Path(model_dir)
    model_name = model_path.name

    file_map = _MODEL_FILES.get(model_name)
    if file_map is None:
        return

    missing = [f for f in _REQUIRED_MODEL_FILES if not (model_path / f).exists()]
    if not missing:
        return

    print(f"[ModelDownload] Model not found at: {model_path}", file=sys.stderr)
    print(f"[ModelDownload] Auto-downloading {model_name} from ModelScope ...", file=sys.stderr)
    model_path.mkdir(parents=True, exist_ok=True)

    import urllib.request

    total = len(file_map)
    for i, (fname, url) in enumerate(file_map.items(), 1):
        dest = model_path / fname
        if dest.exists():
            print(f"[ModelDownload] [{i}/{total}] Skip (exists): {fname}", file=sys.stderr)
            continue
        print(f"[ModelDownload] [{i}/{total}] Downloading: {fname} ...", file=sys.stderr)
        try:
            tmp = dest.with_suffix(dest.suffix + ".tmp")
            def _reporthook(block_num, block_size, total_size):
                if total_size > 0:
                    pct = min(100, block_num * block_size * 100 // total_size)
                    print(f"\r  {pct}%", end="", flush=True, file=sys.stderr)
            urllib.request.urlretrieve(url, str(tmp), reporthook=_reporthook)
            print(file=sys.stderr)
            tmp.rename(dest)
            print(f"[ModelDownload] [{i}/{total}] Done: {fname} ({dest.stat().st_size // 1024} KB)", file=sys.stderr)
        except Exception as e:
            if tmp.exists():
                tmp.unlink(missing_ok=True)
            raise RuntimeError(f"[ModelDownload] Failed to download {fname}: {e}") from e

    print(f"[ModelDownload] All files downloaded to: {model_path}", file=sys.stderr)


# ──────────────────────────────────────────────
# Embedding Client
# Uses pyMNN (MNN.llm) with Android-aligned chat_template.
# Mirrors Android EmbeddingHandler.computeEmbedding().
# Key alignment: llm_config.json jinja.chat_template is replaced with
#   the simplified template used by Android (concatenate message contents),
#   so tokenization is identical to Android KB build path.
# ──────────────────────────────────────────────

class EmbeddingClient:
    """
    Compute text embedding vectors using pyMNN (MNN.llm).
    Uses Android-aligned chat_template (cos~0.991 vs Android vectors).
    """

    def __init__(self, model_dir: str, thread_num: int = 4):
        self.model_dir = Path(model_dir)
        if not (self.model_dir / "config.json").exists():
            raise FileNotFoundError(
                f"Embedding model config not found: {self.model_dir / 'config.json'}"
            )

        self._mnn_embedding = None
        self._dim = None
        self.thread_num = max(1, int(thread_num or 1))
        _ensure_android_embedding_model_config(self.model_dir)
        self.config_json = self.model_dir / "config.json"

        runtime_cfg = json.dumps(
            _build_android_embedding_runtime_config(self.thread_num),
            ensure_ascii=False
        )
        try:
            import MNN.llm as mnn_llm  # type: ignore
            llm = mnn_llm.create(str(self.config_json), True)
            llm.load()
            llm.set_config(runtime_cfg)
            test_var = llm.txt_embedding("test")
            test_vec = self._varp_to_list(test_var)
            if not test_vec:
                raise RuntimeError("pyMNN embedding returned empty vector during warmup")
            self._mnn_embedding = llm
            self._dim = len(test_vec)
            import MNN
            mnn_ver = getattr(MNN, '__version__', '?')
            print(
                f"[EmbeddingClient] pyMNN v{mnn_ver}, model={self.model_dir.name}, "
                f"config={self.config_json}, dim={self._dim}, threads={self.thread_num}",
                file=sys.stderr
            )
        except Exception as e:
            raise RuntimeError(f"pyMNN unavailable: {e}") from e

    def compute_embedding(self, text: str) -> List[float]:
        """Compute embedding vector for text. Returns list of float32 values."""
        return self._compute_via_mnn(text)

    @staticmethod
    def _varp_to_list(var) -> List[float]:
        """
        Convert MNN VARP (shape [1, dim] or [dim]) to Python float list.
        MNN 3.x: MNN.numpy.array(var).flatten() returns a 1-D MNN array.
        """
        import MNN.numpy as mnn_np  # type: ignore
        flat = mnn_np.array(var).flatten()
        return [float(flat[i]) for i in range(flat.shape[0])]

    def _compute_via_mnn(self, text: str) -> List[float]:
        """Compute embedding via pyMNN."""
        try:
            var = self._mnn_embedding.txt_embedding(text)
            result = self._varp_to_list(var)
            if not result:
                raise RuntimeError("pyMNN embedding returned empty vector")
            return result
        except Exception as e:
            raise RuntimeError(f"pyMNN embedding failed: {e}") from e


# ──────────────────────────────────────────────
# Text splitter - mirrors Android DocumentByCharacterSplitter
# ──────────────────────────────────────────────

class TextSplitter:
    """
    Pure character sliding-window splitter.
    Mirrors Android LangChain4j DocumentByCharacterSplitter(chunkSize, chunkOverlap) exactly.
    No separator logic - just slides a window of chunkSize chars with chunkOverlap step-back.
    """

    def __init__(self, chunk_size: int = 500, chunk_overlap: int = 100, min_chunk_size: int = 10):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.min_chunk_size = min_chunk_size

    def split_text(self, text: str) -> List[str]:
        if not text or not text.strip():
            return []
        chunks = []
        start = 0
        length = len(text)
        while start < length:
            end = min(start + self.chunk_size, length)
            chunk = text[start:end]
            if len(chunk) >= self.min_chunk_size:
                chunks.append(chunk)
            if end >= length:
                break
            start += self.chunk_size - self.chunk_overlap
        return chunks


# ──────────────────────────────────────────────
# JSON dataset processor - mirrors Android JsonDatasetProcessor
# ──────────────────────────────────────────────

class JsonDatasetProcessor:
    """
    Processes JSON training datasets in Alpaca/CoT/DPO/Conversation formats.
    Mirrors Android JsonDatasetProcessor exactly.
    """

    def __init__(self, min_chunk_size: int = 10):
        self.min_chunk_size = min_chunk_size

    def is_json_content(self, text: str) -> bool:
        t = text.strip()
        return (t.startswith("{") and t.endswith("}")) or \
               (t.startswith("[") and t.endswith("]"))

    def process(self, json_content: str) -> List[str]:
        """Main entry: auto-detect format and extract text chunks."""
        text = json_content.strip()
        chunks = []
        try:
            if text.startswith("["):
                data = json.loads(text)
                chunks = self._process_array(data)
                print(f"[JsonDatasetProcessor] array items={len(data)}, extracted_units={len(chunks)}", file=sys.stderr)
            elif text.startswith("{"):
                data = json.loads(text)
                chunks = self._process_object(data)
                print(f"[JsonDatasetProcessor] object extracted_units={len(chunks)}", file=sys.stderr)
        except json.JSONDecodeError as e:
            print(f"[JsonDatasetProcessor] JSON parse error: {e}", file=sys.stderr)
        return chunks

    def _process_array(self, arr: list) -> List[str]:
        if not arr:
            return []
        first = arr[0] if isinstance(arr[0], dict) else {}
        # Detect format
        if self._is_alpaca(first):
            return self._process_alpaca(arr)
        elif self._is_cot(first):
            return self._process_cot(arr)
        elif self._is_dpo(first):
            return self._process_dpo(arr)
        elif self._is_conversation(first):
            return self._process_conversation(arr)
        else:
            return self._process_generic_array(arr)

    def _is_alpaca(self, item: dict) -> bool:
        return "instruction" in item and ("output" in item or "response" in item)

    def _is_cot(self, item: dict) -> bool:
        return "question" in item and ("answer" in item or "chain_of_thought" in item)

    def _is_dpo(self, item: dict) -> bool:
        return "chosen" in item and "rejected" in item

    def _is_conversation(self, item: dict) -> bool:
        return "conversations" in item or "messages" in item

    def _process_alpaca(self, arr: list) -> List[str]:
        chunks = []
        for item in arr:
            if not isinstance(item, dict):
                continue
            instruction = item.get("instruction", "").strip()
            input_text = item.get("input", "").strip()
            output = (item.get("output") or item.get("response") or "").strip()
            system_text = item.get("system", "").strip()
            parts = []
            if instruction:
                parts.append(f"指令: {instruction}")
            if input_text:
                parts.append(f"输入: {input_text}")
            if output:
                parts.append(f"输出: {output}")
            text = "\n\n".join(parts)
            if system_text:
                text = f"{text}\n\n系统: {system_text}" if text else f"系统: {system_text}"
            if len(text) >= self.min_chunk_size:
                chunks.append(text)
        return chunks

    def _process_cot(self, arr: list) -> List[str]:
        chunks = []
        for item in arr:
            if not isinstance(item, dict):
                continue
            question = item.get("question", "").strip()
            cot = item.get("chain_of_thought", "").strip()
            answer = item.get("answer", "").strip()
            parts = []
            if question:
                parts.append(f"Question: {question}")
            if cot:
                parts.append(f"Reasoning: {cot}")
            if answer:
                parts.append(f"Answer: {answer}")
            text = "\n".join(parts)
            if len(text) >= self.min_chunk_size:
                chunks.append(text)
        return chunks

    def _process_dpo(self, arr: list) -> List[str]:
        chunks = []
        for item in arr:
            if not isinstance(item, dict):
                continue
            prompt = item.get("prompt", "").strip()
            chosen = item.get("chosen", "")
            if isinstance(chosen, dict):
                chosen = chosen.get("content", "")
            chosen = chosen.strip()
            parts = []
            if prompt:
                parts.append(f"Prompt: {prompt}")
            if chosen:
                parts.append(f"Response: {chosen}")
            text = "\n".join(parts)
            if len(text) >= self.min_chunk_size:
                chunks.append(text)
        return chunks

    def _process_conversation(self, arr: list) -> List[str]:
        chunks = []
        for item in arr:
            if not isinstance(item, dict):
                continue
            messages = item.get("conversations") or item.get("messages") or []
            parts = []
            for msg in messages:
                if not isinstance(msg, dict):
                    continue
                role = (msg.get("from") or msg.get("role") or "").strip()
                content = (msg.get("value") or msg.get("content") or "").strip()
                if role and content:
                    parts.append(f"{role}: {content}")
            text = "\n".join(parts)
            if len(text) >= self.min_chunk_size:
                chunks.append(text)
        return chunks

    def _process_generic_array(self, arr: list) -> List[str]:
        chunks = []
        for item in arr:
            text = None
            if isinstance(item, str):
                text = item.strip()
            elif isinstance(item, dict):
                for key in ("text", "content", "data", "value"):
                    if key in item:
                        text = str(item[key]).strip()
                        break
                if not text:
                    text = " ".join(str(v) for v in item.values()).strip()
            if text and len(text) >= self.min_chunk_size:
                chunks.append(text)
        return chunks

    def _process_object(self, obj: dict) -> List[str]:
        chunks = []
        for key in ("data", "items", "records", "examples", "train"):
            if key in obj and isinstance(obj[key], list):
                chunks.extend(self._process_array(obj[key]))
                break
        if not chunks:
            for v in obj.values():
                if isinstance(v, str) and len(v) >= self.min_chunk_size:
                    chunks.append(v)
        return chunks


def _append_text_part(parts: List[str], text: str, min_chars: int = 0) -> None:
    """Append text part if non-empty. No length filter - aligns with Android XWPFWordExtractor full-text extraction."""
    if not text:
        return
    cleaned = text.strip()
    if cleaned and len(cleaned) >= min_chars:
        parts.append(cleaned)


def _collect_docx_table_parts(table, parts: List[str], min_row_chars: int = 52) -> None:
    last_row = None
    for row in table.rows:
        row_cells = []
        for cell in row.cells:
            cell_text = cell.text.strip()
            if cell_text:
                row_cells.append(cell_text)
        if row_cells:
            row_text = "\t".join(row_cells)
            if len(row_text) >= min_row_chars and row_text != last_row:
                parts.append(row_text)
            last_row = row_text


def _collect_pptx_shape_parts(shape, parts: List[str]) -> None:
    shape_text = getattr(shape, "text", None)
    if isinstance(shape_text, str) and shape_text.strip():
        parts.append(shape_text.strip())
    if getattr(shape, "has_table", False):
        try:
            for row in shape.table.rows:
                row_cells = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_cells.append(cell_text)
                if row_cells:
                    parts.append("\t".join(row_cells))
        except Exception:
            pass
    child_shapes = getattr(shape, "shapes", None)
    if child_shapes is not None:
        try:
            for child in child_shapes:
                _collect_pptx_shape_parts(child, parts)
        except Exception:
            pass


# ──────────────────────────────────────────────
# Document parsers - extract text from various file types
# ──────────────────────────────────────────────

def extract_text_from_file(file_path: str, json_processor: JsonDatasetProcessor,
                            enable_json_splitting: bool = True) -> List[Tuple[str, str, bool]]:
    """
    Extract (text_chunk, source) pairs from a file.
    Supported: txt, md, json, jsonl, pdf, docx, pptx, xlsx, csv, html, htm.
    Returns list of (text, source_filename).
    """
    path = Path(file_path)
    if not path.exists():
        print(f"[Parser] File not found: {file_path}", file=sys.stderr)
        return []

    ext = path.suffix.lower()
    source = path.name

    try:
        if ext in (".txt", ".md", ".log", ".rst", ".csv"):
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                text = f.read()
                print(f"[Parser] {source}: text chars={len(text)}", file=sys.stderr)
                return [(text, source, False)]

        elif ext in (".json", ".jsonl"):
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                raw = f.read()
            if enable_json_splitting and json_processor.is_json_content(raw):
                chunks = json_processor.process(raw)
                if chunks:
                    print(f"[Parser] {source}: json dataset split into {len(chunks)} units", file=sys.stderr)
                    return [(c, source, True) for c in chunks]
            print(f"[Parser] {source}: json raw chars={len(raw)}", file=sys.stderr)
            return [(raw, source, False)]

        elif ext == ".pdf":
            try:
                import pdfplumber
                texts = []
                with pdfplumber.open(str(path)) as pdf:
                    for page in pdf.pages:
                        t = page.extract_text()
                        if t:
                            texts.append(t)
                text = "\n".join(texts)
                print(f"[Parser] {source}: pdf pages={len(pdf.pages)}, chars={len(text)}", file=sys.stderr)
                return [(text, source, False)]
            except ImportError:
                print("[Parser] pdfplumber not installed. Install: pip install pdfplumber", file=sys.stderr)
                return []

        elif ext in (".docx",):
            try:
                import docx
                import re as _re
                doc = docx.Document(str(path))
                parts: List[str] = []
                for paragraph in doc.paragraphs:
                    _append_text_part(parts, paragraph.text)
                for table in doc.tables:
                    _collect_docx_table_parts(table, parts)
                raw = "\n".join(parts)
                # Align with Android cleanText: unify line endings and collapse excessive blank lines
                raw = raw.replace("\r\n", "\n").replace("\r", "\n")
                text = _re.sub(r"\n{3,}", "\n\n", raw).strip()
                print(f"[Parser] {source}: docx paragraphs={len(doc.paragraphs)}, text_parts={len(parts)}, chars={len(text)}", file=sys.stderr)
                return [(text, source, False)]
            except ImportError:
                print("[Parser] python-docx not installed. Install: pip install python-docx", file=sys.stderr)
                return []

        elif ext in (".pptx",):
            try:
                from pptx import Presentation
                prs = Presentation(str(path))
                parts = []
                slide_count = len(prs.slides)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        _collect_pptx_shape_parts(shape, parts)
                text = "\n".join(parts)
                print(f"[Parser] {source}: pptx slides={slide_count}, text_blocks={len(parts)}, chars={len(text)}", file=sys.stderr)
                return [(text, source, False)]
            except ImportError:
                print("[Parser] python-pptx not installed. Install: pip install python-pptx", file=sys.stderr)
                return []

        elif ext in (".xlsx", ".xls"):
            try:
                import openpyxl
                wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
                rows = []
                for ws in wb.worksheets:
                    for row in ws.iter_rows(values_only=True):
                        row_text = "\t".join(str(c) if c is not None else "" for c in row)
                        if row_text.strip():
                            rows.append(row_text)
                text = "\n".join(rows)
                print(f"[Parser] {source}: sheets={len(wb.worksheets)}, rows={len(rows)}, chars={len(text)}", file=sys.stderr)
                return [(text, source, False)]
            except ImportError:
                print("[Parser] openpyxl not installed. Install: pip install openpyxl", file=sys.stderr)
                return []

        elif ext in (".html", ".htm"):
            try:
                from bs4 import BeautifulSoup
                with open(path, "r", encoding="utf-8", errors="replace") as f:
                    soup = BeautifulSoup(f.read(), "html.parser")
                text = soup.get_text(separator="\n")
                print(f"[Parser] {source}: html chars={len(text)}", file=sys.stderr)
                return [(text, source, False)]
            except ImportError:
                # Fallback: strip HTML tags with regex
                with open(path, "r", encoding="utf-8", errors="replace") as f:
                    raw = f.read()
                text = re.sub(r"<[^>]+>", " ", raw)
                print(f"[Parser] {source}: html fallback chars={len(text)}", file=sys.stderr)
                return [(text, source, False)]

        else:
            print(f"[Parser] Unsupported file type: {ext}, treating as text", file=sys.stderr)
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                return [(f.read(), source, False)]

    except Exception as e:
        print(f"[Parser] Error reading {file_path}: {e}", file=sys.stderr)
        return []


# ──────────────────────────────────────────────
# Simple Python NER - mirrors HanLpNerHandler logic
# (HanLP is Java-only; Python side uses jieba + heuristics as equivalent)
# ──────────────────────────────────────────────

class PythonNerHandler:
    """
    Python-side NER using jieba segmentation (mirrors Android HanLpNerHandler).
    - Uses jieba for Chinese word segmentation
    - Applies custom dictionary if provided
    - Skips single-char, punctuation, particles (mirrors shouldSkip())
    - Applies stopwords filtering (mirrors GraphStopwordsMatcher)
    - Applies alias normalization (mirrors aliasToCanonical map)
    """

    # POS tags to skip - mirrors Android shouldSkip()
    SKIP_POS = {"x", "w", "ul", "uj", "ud", "ug", "p", "c", "u", "d", "e",
                "y", "o", "k", "q", "m", "r"}

    def __init__(self, custom_dict_path: str = None, stopwords_path: str = None):
        self.alias_to_canonical: Dict[str, str] = {}
        self.custom_words: Set[str] = set()
        self.stopwords: Set[str] = set()
        self.stopword_prefixes: List[str] = []
        self.stopword_regexes: list = []  # compiled re.Pattern objects
        self._jieba_available = False

        self._init_jieba()
        if custom_dict_path:
            self._load_custom_dict(custom_dict_path)
        if stopwords_path:
            self._load_stopwords(stopwords_path)

    def _init_jieba(self):
        try:
            import jieba
            import jieba.posseg as pseg
            self._jieba = jieba
            self._pseg = pseg
            self._jieba_available = True
            self._jieba.setLogLevel("WARNING")
            print("[NER] jieba available for Chinese NER", file=sys.stderr)
        except ImportError:
            print("[NER] jieba not installed. Install: pip install jieba", file=sys.stderr)
            print("[NER]   Graph RAG NER disabled (vector-only mode)", file=sys.stderr)

    def _load_custom_dict(self, path: str):
        """Load JSON custom dictionary - mirrors Android HanLpNerHandler.loadDictionary()"""
        try:
            with open(path, "r", encoding="utf-8") as f:
                raw = json.load(f)
            entries = raw if isinstance(raw, list) else raw.get("entries", [])
            count = 0
            for entry in entries:
                word = entry.get("word", "")
                if not word:
                    continue
                freq = entry.get("frequency", 10000)
                if self._jieba_available:
                    self._jieba.add_word(word, freq=freq)
                canonical_key = word.strip()
                if canonical_key and canonical_key not in self.alias_to_canonical:
                    self.alias_to_canonical[canonical_key] = canonical_key
                self.custom_words.add(word.strip())
                count += 1
                for alias in entry.get("aliases", []):
                    alias = alias.strip()
                    if not alias:
                        continue
                    if self._jieba_available:
                        self._jieba.add_word(alias, freq=freq)
                    if alias not in self.alias_to_canonical and canonical_key:
                        self.alias_to_canonical[alias] = canonical_key
                    self.custom_words.add(alias)
                    count += 1
            print(f"[NER] Loaded {count} words from custom dictionary: {Path(path).name}", file=sys.stderr)
        except Exception as e:
            print(f"[NER] Failed to load custom dict {path}: {e}", file=sys.stderr)

    def _load_stopwords(self, path: str):
        """
        Load stopwords JSON - mirrors Android GraphStopwordsMatcher.loadFromFile().
        Supports three keys: exact (Set), prefix (List), regex (List[Pattern]).
        Regex compiled patterns use fullmatch semantics (Java Pattern.matcher().matches()).
        """
        try:
            with open(path, "r", encoding="utf-8") as f:
                raw = json.load(f)
            if isinstance(raw, list):
                # Legacy plain list format
                for item in raw:
                    if isinstance(item, str):
                        self.stopwords.add(item.strip())
                    elif isinstance(item, dict):
                        w = item.get("word", "")
                        if w:
                            self.stopwords.add(w.strip())
            elif isinstance(raw, dict):
                for w in raw.get("exact", []):
                    if isinstance(w, str) and w.strip():
                        self.stopwords.add(w.strip())
                for p in raw.get("prefix", []):
                    if isinstance(p, str) and p.strip():
                        self.stopword_prefixes.append(p.strip())
                # Regex patterns compiled for fullmatch (Java .matches() is whole-string)
                for pat in raw.get("regex", []):
                    if not isinstance(pat, str) or not pat.strip():
                        continue
                    try:
                        self.stopword_regexes.append(re.compile(pat.strip()))
                    except re.error as e:
                        print(f"[NER] Invalid stopword regex '{pat}': {e}", file=sys.stderr)
            print(f"[NER] Loaded stopwords from {Path(path).name}: "
                  f"exact={len(self.stopwords)}, "
                  f"prefix={len(self.stopword_prefixes)}, "
                  f"regex={len(self.stopword_regexes)}", file=sys.stderr)
        except Exception as e:
            print(f"[NER] Failed to load stopwords {path}: {e}", file=sys.stderr)

    def should_skip(self, pos: str, word: str) -> bool:
        """Mirror Android HanLpNerHandler.shouldSkip()"""
        if pos in self.SKIP_POS:
            return True
        if len(word.strip()) <= 1:
            return True
        return False

    def matches_stopword(self, text: str) -> bool:
        """
        Check if entity text should be filtered.
        Mirrors Java GraphStopwordsMatcher.matches():
          1. exact set lookup
          2. prefix startsWith
          3. regex fullmatch (NOT search - Java .matches() = whole-string match)
        """
        if not text:
            return False
        if text in self.stopwords:
            return True
        for prefix in self.stopword_prefixes:
            if prefix and text.startswith(prefix):
                return True
        for pat in self.stopword_regexes:
            if pat.fullmatch(text):
                return True
        return False

    def normalize_text(self, text: str) -> Optional[str]:
        """Mirror Android HanLpNerHandler.normalizeTextForGraph()"""
        if not text:
            return None
        t = text.strip()
        if not t:
            return None
        canonical = self.alias_to_canonical.get(t)
        if canonical:
            return canonical
        return t

    def extract_entities(self, text: str, confidence_threshold: float = 0.7) -> List[dict]:
        """
        Extract entities from text. Returns list of {text, type, confidence}.
        Mirrors Android HanLpNerHandler.extractEntities().
        """
        if not self._jieba_available or not text.strip():
            return []
        try:
            words = self._pseg.cut(text)
            entities = []
            for word, pos in words:
                pos_str = str(pos)
                if self.should_skip(pos_str, word):
                    continue
                normalized = self.normalize_text(word)
                if normalized is None:
                    continue
                if self.matches_stopword(normalized):
                    continue
                entities.append({
                    "text": normalized,
                    "type": pos_str,
                    "confidence": 1.0
                })
            return entities
        except Exception as e:
            print(f"[NER] Extraction error: {e}", file=sys.stderr)
            return []


# ──────────────────────────────────────────────
# Knowledge Graph Database - mirrors Android KnowledgeGraphDatabase
# ──────────────────────────────────────────────

DB_VERSION = 1

CREATE_DOCUMENTS = """
CREATE TABLE IF NOT EXISTS documents (
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    collection TEXT NOT NULL,
    content TEXT NOT NULL,
    source TEXT,
    metadata TEXT,
    embedding BLOB NOT NULL,
    created_at INTEGER DEFAULT (strftime('%s', 'now'))
)"""

CREATE_ENTITIES = """
CREATE TABLE IF NOT EXISTS entities (
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    collection TEXT NOT NULL,
    entity_text TEXT NOT NULL,
    entity_type TEXT NOT NULL,
    language TEXT,
    frequency INTEGER DEFAULT 1,
    avg_confidence REAL,
    first_seen INTEGER DEFAULT (strftime('%s', 'now')),
    last_seen INTEGER DEFAULT (strftime('%s', 'now')),
    UNIQUE(collection, entity_text, entity_type)
)"""

CREATE_EDGES = """
CREATE TABLE IF NOT EXISTS entity_edges (
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    collection TEXT NOT NULL,
    from_entity TEXT NOT NULL,
    to_entity TEXT NOT NULL,
    weight REAL DEFAULT 1.0,
    chunk_ids TEXT,
    created_at INTEGER DEFAULT (strftime('%s', 'now')),
    updated_at INTEGER DEFAULT (strftime('%s', 'now')),
    UNIQUE(collection, from_entity, to_entity)
)"""

CREATE_CHUNK_ENTITIES = """
CREATE TABLE IF NOT EXISTS chunk_entities (
    chunk_id INTEGER NOT NULL,
    entity_text TEXT NOT NULL,
    entity_type TEXT NOT NULL,
    confidence REAL,
    PRIMARY KEY(chunk_id, entity_text, entity_type),
    FOREIGN KEY(chunk_id) REFERENCES documents(id) ON DELETE CASCADE
)"""

CREATE_METADATA = """
CREATE TABLE IF NOT EXISTS metadata (
    key TEXT PRIMARY KEY,
    value TEXT
)"""

INDEXES = [
    "CREATE INDEX IF NOT EXISTS idx_doc_collection ON documents(collection)",
    "CREATE INDEX IF NOT EXISTS idx_entity_text ON entities(entity_text)",
    "CREATE INDEX IF NOT EXISTS idx_entity_type ON entities(entity_type)",
    "CREATE INDEX IF NOT EXISTS idx_entity_collection ON entities(collection)",
    "CREATE INDEX IF NOT EXISTS idx_entity_freq ON entities(frequency)",
    "CREATE INDEX IF NOT EXISTS idx_edge_from ON entity_edges(from_entity)",
    "CREATE INDEX IF NOT EXISTS idx_edge_to ON entity_edges(to_entity)",
    "CREATE INDEX IF NOT EXISTS idx_edge_collection ON entity_edges(collection)",
    "CREATE INDEX IF NOT EXISTS idx_edge_weight ON entity_edges(weight)",
    "CREATE INDEX IF NOT EXISTS idx_ce_chunk ON chunk_entities(chunk_id)",
    "CREATE INDEX IF NOT EXISTS idx_ce_entity ON chunk_entities(entity_text)",
]


def _vector_to_blob(vector: List[float]) -> bytes:
    """
    Serialize float[] to bytes - mirrors Android KnowledgeGraphDatabase.vectorToBlob().
    Uses little-endian 4-byte floats, identical to Java ByteBuffer.LITTLE_ENDIAN.putFloat().
    """
    return struct.pack(f"<{len(vector)}f", *vector)


def _blob_to_vector(blob: bytes) -> List[float]:
    """
    Deserialize bytes to float[] - mirrors Android KnowledgeGraphDatabase.blobToVector().
    """
    count = len(blob) // 4
    return list(struct.unpack(f"<{count}f", blob))


def _bm25_tokenize(text: str) -> List[str]:
    """
    Simple tokenizer for BM25: split on whitespace/punctuation, extract CJK chars
    as individual tokens, keep ASCII word tokens (>=2 chars).
    Intentionally lightweight - no jieba dependency needed for BM25.
    """
    tokens: List[str] = []
    # ASCII alphanumeric tokens (min length 2)
    for m in re.finditer(r'[A-Za-z0-9][A-Za-z0-9_\-\.]*', text):
        t = m.group(0).lower().strip('._-')
        if len(t) >= 2:
            tokens.append(t)
    # CJK character bigrams
    cjk = re.findall(r'[\u4e00-\u9fff\u3400-\u4dbf\u20000-\u2a6df]', text)
    for i in range(len(cjk)):
        tokens.append(cjk[i])
        if i + 1 < len(cjk):
            tokens.append(cjk[i] + cjk[i + 1])
    return tokens


def _build_bm25_index(corpus: List[str], k1: float = 1.5, b: float = 0.75):
    """
    Build BM25 index from a list of document strings.
    Pre-computes per-document tf_map and dl for fast query-time scoring.
    Returns a dict with internal state for scoring.
    """
    N = len(corpus)
    tokenized = [_bm25_tokenize(doc) for doc in corpus]
    avg_dl = sum(len(d) for d in tokenized) / max(N, 1)
    # Pre-compute tf_map and dl per document (avoids rebuilding at query time)
    doc_tf: List[Dict[str, int]] = []
    doc_dl: List[int] = []
    df: Dict[str, int] = {}
    for doc_tokens in tokenized:
        tf_map: Dict[str, int] = {}
        for t in doc_tokens:
            tf_map[t] = tf_map.get(t, 0) + 1
        doc_tf.append(tf_map)
        doc_dl.append(len(doc_tokens))
        for term in tf_map:
            df[term] = df.get(term, 0) + 1
    return {
        "doc_tf": doc_tf,
        "doc_dl": doc_dl,
        "avg_dl": avg_dl,
        "df": df,
        "N": N,
        "k1": k1,
        "b": b,
    }


def _bm25_score(index: dict, query_tokens: List[str], doc_idx: int) -> float:
    """Compute BM25 score for a single document using pre-built index."""
    doc_tf = index["doc_tf"]
    doc_dl = index["doc_dl"]
    avg_dl = index["avg_dl"]
    df = index["df"]
    N = index["N"]
    k1 = index["k1"]
    b = index["b"]

    tf_map = doc_tf[doc_idx]
    dl = doc_dl[doc_idx]

    score = 0.0
    for term in query_tokens:
        tf = tf_map.get(term, 0)
        if tf == 0:
            continue
        n_q = df.get(term, 0)
        idf = math.log((N - n_q + 0.5) / (n_q + 0.5) + 1.0)
        numerator = tf * (k1 + 1.0)
        denominator = tf + k1 * (1.0 - b + b * dl / max(avg_dl, 1.0))
        score += idf * numerator / denominator
    return score


def _run_bm25_search(query: str, all_chunks: List[dict], top_k: int) -> List[dict]:
    """
    BM25 search over all_chunks [{id, content, source}].
    Returns top_k chunks sorted by BM25 score descending, each with 'bm25_score'.
    """
    if not all_chunks:
        return []
    corpus = [c["content"] for c in all_chunks]
    index = _build_bm25_index(corpus)
    query_tokens = _bm25_tokenize(query)
    if not query_tokens:
        return []
    scored = []
    for i, chunk in enumerate(all_chunks):
        s = _bm25_score(index, query_tokens, i)
        if s > 0:
            scored.append({**chunk, "bm25_score": s, "similarity": 0.0})
    scored.sort(key=lambda x: x["bm25_score"], reverse=True)
    return scored[:top_k]


def _rrf_fusion(vec_results: List[dict], bm25_results: List[dict],
                top_k: int, rrf_k: int = 60) -> List[dict]:
    """
    Reciprocal Rank Fusion of two ranked lists.
    Each result: RRF_score = 1/(k + rank_vec) + 1/(k + rank_bm25)
    Missing from one list gets penalty rank = max(len_list, 1) * 2 + 1,
    ensuring absent items score strictly lower than the last present item.
    Edge cases:
      - vec_results empty  → pure BM25 ordering preserved
      - bm25_results empty → pure vec ordering preserved
      - top_k <= 0         → returns empty list
      - top_k > total candidates → returns all candidates
    """
    if top_k <= 0:
        return []
    # Degenerate: one list empty → return the other as-is (up to top_k)
    if not vec_results:
        return [dict(r) for r in bm25_results[:top_k]]
    if not bm25_results:
        return [dict(r) for r in vec_results[:top_k]]

    # Build id -> chunk dict (vec takes priority for similarity field)
    id_to_chunk: Dict[int, dict] = {}
    for r in bm25_results:
        id_to_chunk[r["id"]] = dict(r)
    for r in vec_results:
        existing = id_to_chunk.get(r["id"])
        if existing:
            existing["similarity"] = r["similarity"]
        else:
            id_to_chunk[r["id"]] = dict(r)
        if "bm25_score" not in id_to_chunk[r["id"]]:
            id_to_chunk[r["id"]]["bm25_score"] = 0.0

    vec_rank: Dict[int, int] = {r["id"]: i + 1 for i, r in enumerate(vec_results)}
    bm25_rank: Dict[int, int] = {r["id"]: i + 1 for i, r in enumerate(bm25_results)}
    # Penalty rank for absent items: well beyond the end of each list
    # Using len*2+1 ensures absent items always score lower than last-rank present items
    vec_penalty = len(vec_results) * 2 + 1
    bm25_penalty = len(bm25_results) * 2 + 1

    rrf_scores: Dict[int, float] = {}
    all_ids = set(vec_rank.keys()) | set(bm25_rank.keys())
    for cid in all_ids:
        rv = vec_rank.get(cid, vec_penalty)
        rb = bm25_rank.get(cid, bm25_penalty)
        rrf_scores[cid] = 1.0 / (rrf_k + rv) + 1.0 / (rrf_k + rb)

    merged = sorted(all_ids, key=lambda cid: rrf_scores[cid], reverse=True)
    result = []
    for cid in merged[:top_k]:
        chunk = id_to_chunk[cid]
        chunk["rrf_score"] = rrf_scores[cid]
        result.append(chunk)
    return result


def _cosine_similarity(a: List[float], b: List[float]) -> float:
    """
    Cosine similarity - mirrors Android KnowledgeGraphDatabase.cosineSimilarity().
    dot(a,b) / (norm(a) * norm(b))
    Identical algorithm to Java and C++ implementations.
    """
    if not a or not b or len(a) != len(b):
        return 0.0
    dot = sum(x * y for x, y in zip(a, b))
    norm_a = math.sqrt(sum(x * x for x in a))
    norm_b = math.sqrt(sum(y * y for y in b))
    if norm_a == 0.0 or norm_b == 0.0:
        return 0.0
    return dot / (norm_a * norm_b)


def _extract_query_terms(text: str, ner_handler=None) -> List[Tuple[str, float]]:
    text = (text or "").strip()
    if not text:
        return []
    terms: List[Tuple[str, float]] = []
    seen: Set[str] = set()

    def _add_term(term: str, weight: float):
        normalized = (term or "").strip()
        if len(normalized) < 2:
            return
        key = normalized.casefold()
        if key in seen:
            return
        if ner_handler and ner_handler.matches_stopword(normalized):
            return
        seen.add(key)
        terms.append((normalized, float(weight)))

    for match in re.finditer(r"[A-Za-z0-9][A-Za-z0-9_:/\.\-]{1,}", text):
        token = match.group(0).strip("_:/.-")
        if len(token) < 2:
            continue
        weight = 2.5 if (any(ch.isdigit() for ch in token) or token.upper() == token) else 1.6
        _add_term(token, weight)

    if ner_handler:
        for ent in ner_handler.extract_entities(text):
            normalized = ner_handler.normalize_text(ent.get("text", "")) or ent.get("text", "")
            if normalized:
                _add_term(normalized, 1.0)
        alias_keys = sorted(ner_handler.alias_to_canonical.keys(), key=len, reverse=True)
        for alias in alias_keys:
            if len(alias) < 2 or alias not in text:
                continue
            normalized = ner_handler.normalize_text(alias) or alias
            weight = 2.5 if re.search(r"[A-Za-z0-9]", normalized) else 1.0
            _add_term(normalized, weight)

    filtered_terms: List[Tuple[str, float]] = []
    for term, weight in terms:
        is_nested = False
        if re.fullmatch(r"[A-Za-z0-9]+", term):
            for other_term, _ in terms:
                if other_term == term:
                    continue
                if len(other_term) > len(term) and re.fullmatch(r"[A-Za-z0-9]+", other_term):
                    if term.casefold() in other_term.casefold():
                        is_nested = True
                        break
        if not is_nested:
            filtered_terms.append((term, weight))

    return filtered_terms


def _extract_model_terms(text: str) -> List[str]:
    text = (text or "").strip()
    if not text:
        return []
    seen: Set[str] = set()
    result: List[str] = []
    for match in re.finditer(r"[A-Za-z]+\d{2,}[A-Za-z0-9]*", text):
        token = match.group(0).strip()
        key = token.casefold()
        if key in seen:
            continue
        seen.add(key)
        result.append(token)
    return result


def _score_model_term_matches(content: str, query_model_terms: List[str]) -> Tuple[int, int]:
    if not content or not query_model_terms:
        return 0, 0
    haystack = content.casefold()
    hits = 0
    for term in query_model_terms:
        if term.casefold() in haystack:
            hits += 1
    return hits, len(query_model_terms)


def _score_content_against_query_terms(content: str, query_terms: List[Tuple[str, float]]) -> Tuple[float, int, int]:
    if not content or not query_terms:
        return 0.0, 0, 0
    haystack = content.casefold()
    matched_weight = 0.0
    strong_hits = 0
    strong_total = 0
    total_weight = 0.0
    for term, weight in query_terms:
        total_weight += weight
        is_strong = weight >= 1.5
        if is_strong:
            strong_total += 1
        if term.casefold() in haystack:
            matched_weight += weight
            if is_strong:
                strong_hits += 1
    if total_weight <= 0.0:
        return 0.0, strong_hits, strong_total
    return matched_weight / total_weight, strong_hits, strong_total


def _apply_lexical_rerank(search_results: List[dict], user_query: str,
                          ner_handler, cfg_retrieval: dict) -> List[dict]:
    if not search_results:
        return search_results
    if not cfg_retrieval.get("lexical_match_enabled", True):
        return search_results
    query_terms = _extract_query_terms(user_query, ner_handler)
    query_model_terms = _extract_model_terms(user_query)
    if not query_terms:
        return search_results

    lexical_weight = float(cfg_retrieval.get("lexical_match_weight", 0.35))
    strong_penalty = float(cfg_retrieval.get("lexical_strong_term_penalty", 0.18))
    model_missing_penalty = float(cfg_retrieval.get("lexical_model_missing_penalty", 0.45))
    model_match_bonus = float(cfg_retrieval.get("lexical_model_match_bonus", 0.20))
    sim_scores = [r.get("similarity", 0.0) for r in search_results]
    sim_min = min(sim_scores)
    sim_max = max(sim_scores)
    reranked: List[dict] = []

    for rank, result in enumerate(search_results):
        lexical_score, strong_hits, strong_total = _score_content_against_query_terms(
            result.get("content", ""), query_terms
        )
        model_hits, model_total = _score_model_term_matches(result.get("content", ""), query_model_terms)
        similarity = result.get("similarity", 0.0)
        sim_norm = ((similarity - sim_min) / (sim_max - sim_min)
                    if sim_max > sim_min else (1.0 if sim_max > 0 else 0.0))
        hybrid_score = (1.0 - lexical_weight) * sim_norm + lexical_weight * lexical_score
        if strong_total > 0 and strong_hits == 0:
            hybrid_score -= strong_penalty
        if model_total > 0:
            if model_hits == 0:
                hybrid_score -= model_missing_penalty
            else:
                hybrid_score += model_match_bonus * (model_hits / model_total)
        item = dict(result)
        item["lexical_score"] = lexical_score
        item["strong_term_hits"] = strong_hits
        item["strong_term_total"] = strong_total
        item["model_term_hits"] = model_hits
        item["model_term_total"] = model_total
        item["hybrid_score"] = hybrid_score
        item["vector_rank"] = rank
        reranked.append(item)

    reranked.sort(
        key=lambda x: (x.get("hybrid_score", 0.0), x.get("similarity", 0.0)),
        reverse=True
    )
    top_preview = ", ".join(
        f"id={r['id']}:lex={r.get('lexical_score', 0.0):.3f}/hy={r.get('hybrid_score', 0.0):.3f}"
        for r in reranked[:5]
    )
    print(f"[query_kb] Lexical rerank terms={len(query_terms)} top: {top_preview}", file=sys.stderr)
    return reranked


class KnowledgeGraphDB:
    """
    Python-side knowledge graph database.
    Schema and operations mirror Android KnowledgeGraphDatabase exactly.
    """

    def __init__(self, db_path: str, collection: str):
        self.db_path = db_path
        self._provided_collection = collection  # Only for new DB creation
        os.makedirs(Path(db_path).parent, exist_ok=True)
        self.conn = sqlite3.connect(db_path, check_same_thread=False)
        self.conn.execute("PRAGMA journal_mode=WAL")
        self.conn.execute("PRAGMA foreign_keys=ON")
        self._init_schema()
        
        # Auto-detect actual collection from DB, don't force directory name
        self.collection = self._detect_collection()
        print(f"[KnowledgeGraphDB] path={db_path}, detected_collection={self.collection}", file=sys.stderr)

    def _init_schema(self):
        c = self.conn
        c.execute(CREATE_DOCUMENTS)
        c.execute(CREATE_ENTITIES)
        c.execute(CREATE_EDGES)
        c.execute(CREATE_CHUNK_ENTITIES)
        c.execute(CREATE_METADATA)
        for idx in INDEXES:
            c.execute(idx)
        # Init metadata
        c.execute("INSERT OR IGNORE INTO metadata(key,value) VALUES(?,?)",
                  ("schema_version", str(DB_VERSION)))
        c.execute("INSERT OR IGNORE INTO metadata(key,value) VALUES(?,?)",
                  ("created_at", str(int(time.time() * 1000))))
        # Set SQLite user_version for Android SQLiteOpenHelper compatibility
        c.execute(f"PRAGMA user_version = {DB_VERSION}")
        c.commit()

    def _detect_collection(self) -> str:
        """
        Auto-detect actual collection name from existing data in DB.
        If DB is empty, use the provided collection name for new data.
        This allows copying/renaming KB directories without breaking queries.
        """
        # Try to get collection from existing documents
        row = self.conn.execute(
            "SELECT DISTINCT collection FROM documents LIMIT 1"
        ).fetchone()
        if row:
            return row[0]
        
        # Try to get from entities
        row = self.conn.execute(
            "SELECT DISTINCT collection FROM entities LIMIT 1"
        ).fetchone()
        if row:
            return row[0]
        
        # Empty DB: use provided collection name
        return self._provided_collection

    def update_metadata(self, model_dir: str, embedding_dim: int,
                         hub_threshold_build: int = 0):
        """Mirror Android KnowledgeGraphDatabase.updateMetadata()"""
        ts = int(time.time())
        updates = [
            ("modeldir", model_dir),
            ("embedding_model", model_dir),
            ("embedding_dimension", str(embedding_dim)),
        ]
        if hub_threshold_build > 0:
            updates.append(("hub_threshold", str(hub_threshold_build)))
        for key, val in updates:
            self.conn.execute(
                "INSERT OR REPLACE INTO metadata(key,value) VALUES(?,?)", (key, val)
            )
        self.conn.commit()

    def get_metadata(self) -> dict:
        cursor = self.conn.execute("SELECT key, value FROM metadata")
        return {row[0]: row[1] for row in cursor.fetchall()}

    def get_chunk_count(self) -> int:
        row = self.conn.execute(
            "SELECT COUNT(*) FROM documents WHERE collection=?", (self.collection,)
        ).fetchone()
        return row[0] if row else 0

    def add_chunk(self, content: str, source: str, embedding: List[float],
                  metadata: str = "") -> int:
        """
        Insert document chunk. Returns chunk_id.
        Mirrors Android KnowledgeGraphDatabase.addChunk().
        """
        blob = _vector_to_blob(embedding)
        cursor = self.conn.execute(
            "INSERT INTO documents(collection,content,source,metadata,embedding) VALUES(?,?,?,?,?)",
            (self.collection, content, source, metadata, blob)
        )
        self.conn.commit()
        chunk_id = cursor.lastrowid
        return chunk_id

    def search_similar(self, query_vector: List[float], top_k: int) -> List[dict]:
        """
        Vector similarity search - mirrors Android KnowledgeGraphDatabase.searchSimilar().
        Full scan + cosine similarity sort.
        """
        cursor = self.conn.execute(
            "SELECT id, content, source, embedding FROM documents WHERE collection=?",
            (self.collection,)
        )
        scored = []
        for row in cursor.fetchall():
            chunk_id, content, source, blob = row
            vec = _blob_to_vector(blob)
            sim = _cosine_similarity(query_vector, vec)
            scored.append({"id": chunk_id, "content": content, "source": source,
                            "similarity": sim})
        scored.sort(key=lambda x: x["similarity"], reverse=True)
        return scored[:top_k]

    def get_all_chunks_for_bm25(self) -> List[dict]:
        """
        Return all chunks as [{id, content, source}] without embedding blob.
        Used for BM25 full-corpus scoring - avoids deserializing vectors.
        """
        cursor = self.conn.execute(
            "SELECT id, content, source FROM documents WHERE collection=?",
            (self.collection,)
        )
        return [{"id": row[0], "content": row[1], "source": row[2]}
                for row in cursor.fetchall()]

    def add_entity(self, entity_text: str, entity_type: str,
                   confidence: float = 1.0) -> int:
        """Upsert entity. Returns entity_id."""
        ts = int(time.time())
        row = self.conn.execute(
            "SELECT id, frequency, avg_confidence FROM entities "
            "WHERE collection=? AND entity_text=? AND entity_type=?",
            (self.collection, entity_text, entity_type)
        ).fetchone()
        if row:
            eid, old_freq, old_conf = row
            new_freq = old_freq + 1
            new_conf = (old_conf * old_freq + confidence) / new_freq
            self.conn.execute(
                "UPDATE entities SET frequency=?, avg_confidence=?, last_seen=? WHERE id=?",
                (new_freq, new_conf, ts, eid)
            )
            return eid
        else:
            cursor = self.conn.execute(
                "INSERT INTO entities(collection,entity_text,entity_type,language,"
                "frequency,avg_confidence,first_seen,last_seen) VALUES(?,?,?,?,?,?,?,?)",
                (self.collection, entity_text, entity_type, "zh",
                 1, confidence, ts, ts)
            )
            return cursor.lastrowid

    def link_chunk_to_entity(self, chunk_id: int, entity_text: str,
                              entity_type: str, confidence: float):
        """Mirror Android KnowledgeGraphDatabase.linkChunkToEntity()"""
        try:
            self.conn.execute(
                "INSERT OR IGNORE INTO chunk_entities(chunk_id,entity_text,entity_type,confidence) "
                "VALUES(?,?,?,?)",
                (chunk_id, entity_text, entity_type, confidence)
            )
        except Exception as e:
            print(f"[KnowledgeGraphDB] link_chunk_to_entity error: {e}", file=sys.stderr)

    def add_edge(self, from_entity_id: int, to_entity_id: int,
                  weight: float = 1.0, chunk_id: int = None):
        """
        Add/update co-occurrence edge. Mirrors Android KnowledgeGraphDatabase.addEdgeInternal().
        Consistent ordering: smaller text first (mirrors Java makeEdgeKey compareTo).
        chunk_id is appended to the chunk_ids JSON array if provided.
        """
        if from_entity_id > to_entity_id:
            from_entity_id, to_entity_id = to_entity_id, from_entity_id

        from_text = self._get_entity_text(from_entity_id)
        to_text = self._get_entity_text(to_entity_id)
        if not from_text or not to_text:
            return
        # Ensure canonical ordering by text (mirrors Java makeEdgeKey: compareTo)
        if from_text > to_text:
            from_text, to_text = to_text, from_text

        row = self.conn.execute(
            "SELECT id, weight, chunk_ids FROM entity_edges "
            "WHERE collection=? AND from_entity=? AND to_entity=?",
            (self.collection, from_text, to_text)
        ).fetchone()
        ts = int(time.time())
        if row:
            new_weight = row[1] + weight
            # Merge chunk_ids JSON array
            existing_ids = set()
            if row[2]:
                try:
                    existing_ids = set(json.loads(row[2]))
                except Exception:
                    pass
            if chunk_id is not None:
                existing_ids.add(chunk_id)
            chunk_ids_json = json.dumps(sorted(existing_ids))
            self.conn.execute(
                "UPDATE entity_edges SET weight=?, chunk_ids=?, updated_at=? WHERE id=?",
                (new_weight, chunk_ids_json, ts, row[0])
            )
        else:
            chunk_ids_json = json.dumps([chunk_id]) if chunk_id is not None else json.dumps([])
            self.conn.execute(
                "INSERT INTO entity_edges(collection,from_entity,to_entity,weight,chunk_ids,created_at,updated_at) "
                "VALUES(?,?,?,?,?,?,?)",
                (self.collection, from_text, to_text, weight, chunk_ids_json, ts, ts)
            )

    def add_entities_and_build_graph(self, chunk_id: int, entities: List[dict],
                                      ner_handler: PythonNerHandler = None):
        """
        Batch: addEntity + linkChunkToEntity + addEdge for all entity pairs.
        Mirrors Android KnowledgeGraphDatabase.addEntitiesAndBuildGraph().
        """
        if not entities:
            return 0
        entity_ids = []
        for ent in entities:
            text = ent["text"]
            etype = ent["type"]
            conf = ent.get("confidence", 1.0)
            if ner_handler:
                text = ner_handler.normalize_text(text) or text
            eid = self.add_entity(text, etype, conf)
            if eid > 0:
                self.link_chunk_to_entity(chunk_id, text, etype, conf)
                entity_ids.append(eid)

        # Build bidirectional co-occurrence edges - mirrors Java loop
        # Note: add_edge already normalizes ordering (smaller id first), so one call per pair
        edge_count = 0
        for j in range(len(entity_ids)):
            for k in range(j + 1, len(entity_ids)):
                self.add_edge(entity_ids[j], entity_ids[k], 1.0, chunk_id)
                edge_count += 1
        self.conn.commit()
        return len(entity_ids)

    def get_entities_for_chunks(self, chunk_ids: List[int]) -> Dict[int, List[str]]:
        """Mirror Android KnowledgeGraphDatabase.getEntitiesForChunks()"""
        if not chunk_ids:
            return {}
        placeholders = ",".join("?" * len(chunk_ids))
        rows = self.conn.execute(
            f"SELECT chunk_id, entity_text FROM chunk_entities WHERE chunk_id IN ({placeholders})",
            chunk_ids
        ).fetchall()
        result: Dict[int, List[str]] = {}
        for chunk_id, entity_text in rows:
            result.setdefault(chunk_id, []).append(entity_text)
        return result

    def get_connected_entities(self, seed_entities: Set[str],
                                min_weight: int, max_results: int) -> List[dict]:
        """
        1-hop graph expansion. Mirrors Android KnowledgeGraphDatabase.getConnectedEntities().
        """
        if not seed_entities:
            return []
        placeholders = ",".join("?" * len(seed_entities))
        seeds = list(seed_entities)
        query = (
            f"SELECT to_entity, weight FROM entity_edges "
            f"WHERE collection=? AND from_entity IN ({placeholders}) AND weight>=? "
            f"UNION "
            f"SELECT from_entity, weight FROM entity_edges "
            f"WHERE collection=? AND to_entity IN ({placeholders}) AND weight>=? "
            f"ORDER BY weight DESC LIMIT ?"
        )
        args = [self.collection] + seeds + [min_weight, self.collection] + seeds + [min_weight, max_results]
        rows = self.conn.execute(query, args).fetchall()
        result = []
        for entity_text, weight in rows:
            if entity_text not in seed_entities:
                result.append({"entity_text": entity_text, "weight": weight})
        return result

    def get_chunk_ids_by_entities(self, entity_texts: List[str]) -> List[int]:
        """
        Mirror Android KnowledgeGraphDatabase.getChunkIdsByEntities().
        Returns DISTINCT chunk_ids that contain any of the given entity texts.
        """
        if not entity_texts:
            return []
        placeholders = ",".join("?" * len(entity_texts))
        rows = self.conn.execute(
            f"SELECT DISTINCT chunk_id FROM chunk_entities WHERE entity_text IN ({placeholders})",
            entity_texts
        ).fetchall()
        return [row[0] for row in rows]

    def get_chunks_by_ids(self, chunk_ids: List[int]) -> List[dict]:
        """
        Mirror Android KnowledgeGraphDatabase.getChunksByIds().
        Returns document rows for the given chunk IDs (similarity=0.0 for graph-expanded chunks).
        """
        if not chunk_ids:
            return []
        placeholders = ",".join("?" * len(chunk_ids))
        rows = self.conn.execute(
            f"SELECT id, content, source FROM documents WHERE id IN ({placeholders})",
            chunk_ids
        ).fetchall()
        return [{"id": row[0], "content": row[1], "source": row[2], "similarity": 0.0}
                for row in rows]

    def get_hub_entities(self, threshold: int,
                          protected: Set[str] = None) -> Set[str]:
        """
        Query hub entities by edge degree + total weight.
        Mirrors Android KnowledgeGraphDatabase.getHubEntities(threshold, protectedEntities).

        Hub criteria (normal entity):
          degree >= threshold OR total_weight >= threshold
        Protected entities (custom dictionary words) use 5x fallback:
          degree >= threshold*5 OR total_weight >= threshold*5
        """
        if threshold <= 0:
            return set()
        # UNION ALL both edge directions to get per-entity degree and total weight
        hub_query = (
            "SELECT entity_text, COUNT(DISTINCT neighbor) AS degree, "
            "SUM(weight) AS total_weight "
            "FROM ("
            "  SELECT from_entity AS entity_text, to_entity AS neighbor, weight "
            "  FROM entity_edges WHERE collection=? "
            "  UNION ALL "
            "  SELECT to_entity AS entity_text, from_entity AS neighbor, weight "
            "  FROM entity_edges WHERE collection=? "
            ") AS edges "
            "GROUP BY entity_text"
        )
        rows = self.conn.execute(hub_query, (self.collection, self.collection)).fetchall()
        fallback_threshold = threshold * 5
        hub_entities: Set[str] = set()
        for entity_text, degree, total_weight in rows:
            is_protected = protected is not None and entity_text in protected
            if not is_protected:
                if degree >= threshold or total_weight >= threshold:
                    hub_entities.add(entity_text)
            else:
                if fallback_threshold > 0 and (
                        degree >= fallback_threshold or total_weight >= fallback_threshold):
                    hub_entities.add(entity_text)
        return hub_entities

    def apply_hub_threshold(self, threshold: int,
                             protected: Set[str] = None) -> int:
        """
        Remove hub entities from DB (entities + edges + chunk_entities).
        Mirrors Android KnowledgeGraphDatabase.applyHubThreshold().
        Used by add_note (post-insert cleanup on a small DB).
        Protected entities use 5x fallback threshold.
        Returns number of hub entities removed.
        """
        if threshold <= 0:
            return 0
        hub_entities = self.get_hub_entities(threshold, protected)
        if not hub_entities:
            return 0
        hub_count = 0
        old_isolation = self.conn.isolation_level
        try:
            self.conn.isolation_level = None  # manual transaction
            self.conn.execute("BEGIN")
            for hub in hub_entities:
                self.conn.execute(
                    "DELETE FROM entities WHERE collection=? AND entity_text=?",
                    (self.collection, hub)
                )
                self.conn.execute(
                    "DELETE FROM chunk_entities "
                    "WHERE entity_text=? AND chunk_id IN "
                    "(SELECT id FROM documents WHERE collection=?)",
                    (hub, self.collection)
                )
                self.conn.execute(
                    "DELETE FROM entity_edges "
                    "WHERE collection=? AND (from_entity=? OR to_entity=?)",
                    (self.collection, hub, hub)
                )
                hub_count += 1
            self.conn.execute("COMMIT")
        except Exception as e:
            self.conn.execute("ROLLBACK")
            print(f"[KnowledgeGraphDB] apply_hub_threshold error: {e}", file=sys.stderr)
        finally:
            self.conn.isolation_level = old_isolation
        print(f"[KnowledgeGraphDB] apply_hub_threshold: removed={hub_count} threshold={threshold}",
              file=sys.stderr)
        return hub_count

    def _get_entity_text(self, entity_id: int) -> Optional[str]:
        row = self.conn.execute(
            "SELECT entity_text FROM entities WHERE id=?", (entity_id,)
        ).fetchone()
        return row[0] if row else None

    def close(self):
        self.conn.close()


# ──────────────────────────────────────────────
# InMemoryGraphBuilder - mirrors Android TextChunkProcessor.InMemoryGraphBuilder
# Build graph entirely in memory, apply hub filter, then flush to SQLite.
# Avoids N² individual SQLite transactions during build phase.
# ──────────────────────────────────────────────

class InMemoryGraphBuilder:
    """
    Accumulates entities and edges in memory during KB construction.
    After all chunks are processed, applies hub threshold filter and then
    flushes the cleaned graph to SQLite in a single transaction.
    Mirrors Android TextChunkProcessor.InMemoryGraphBuilder exactly.
    """

    def __init__(self, protected_entity_texts: Set[str] = None,
                 hub_threshold: int = 0):
        # entity_key -> {text, type, frequency, avg_confidence, first_seen, last_seen}
        self._entity_map: Dict[tuple, dict] = {}
        # chunk_id -> list of {text, type, confidence}
        self._chunk_entity_map: Dict[int, list] = {}
        # edge_key (from_text, to_text) with from_text <= to_text -> {weight: int, chunk_ids: set}
        self._edge_map: Dict[tuple, dict] = {}

        self._protected = set(protected_entity_texts) if protected_entity_texts else set()
        self._hub_threshold = hub_threshold
        self._online_hub_threshold = hub_threshold
        self._online_protected_hub_threshold = hub_threshold * 5 if hub_threshold > 0 else 0
        # Per-entity online hub stats: text -> {degree, total_weight}
        self._online_hub_stats: Dict[str, dict] = {}
        self._runtime_hub_entities: Set[str] = set()

    def add_chunk_entities(self, chunk_id: int, entities: List[dict],
                            ner_handler=None):
        """
        Process one chunk: upsert entities, link chunk->entities,
        build bidirectional co-occurrence edges.
        Mirrors the inner loop of Java addEntitiesAndBuildGraph().
        """
        if not entities:
            return {"input_entities": 0, "effective_entities": 0, "edges_added": 0}
        ts = int(time.time())
        effective_entities = []
        edges_added = 0

        for ent in entities:
            raw_text = ent["text"]
            etype = ent["type"]
            conf = ent.get("confidence", 1.0)

            if ner_handler:
                raw_text = ner_handler.normalize_text(raw_text) or raw_text
            if not raw_text or raw_text in self._runtime_hub_entities:
                continue

            effective_entities.append({"text": raw_text, "type": etype, "confidence": conf})

        if not effective_entities:
            return {"input_entities": len(entities), "effective_entities": 0, "edges_added": 0}

        for ent in effective_entities:
            raw_text = ent["text"]
            etype = ent["type"]
            conf = ent["confidence"]

            key = (raw_text, etype)
            if key in self._entity_map:
                e = self._entity_map[key]
                old_freq = e["frequency"]
                new_freq = old_freq + 1
                e["frequency"] = new_freq
                e["avg_confidence"] = (e["avg_confidence"] * old_freq + conf) / new_freq
                e["last_seen"] = ts
            else:
                self._entity_map[key] = {
                    "text": raw_text, "type": etype, "language": "zh",
                    "frequency": 1, "avg_confidence": conf,
                    "first_seen": ts, "last_seen": ts
                }

            # Link chunk -> entity
            self._chunk_entity_map.setdefault(chunk_id, []).append(
                {"text": raw_text, "type": etype, "confidence": conf}
            )

        for j in range(len(effective_entities)):
            ft = effective_entities[j]["text"]
            if ft in self._runtime_hub_entities:
                continue
            for k in range(j + 1, len(effective_entities)):
                tt = effective_entities[k]["text"]
                if tt in self._runtime_hub_entities:
                    continue
                self._increment_online_hub_stats(ft, 1)
                self._increment_online_hub_stats(tt, 1)
                self._check_and_handle_runtime_hub(ft)
                self._check_and_handle_runtime_hub(tt)
                if ft in self._runtime_hub_entities or tt in self._runtime_hub_entities:
                    continue
                if ft > tt:
                    ft, tt = tt, ft
                edge_key = (ft, tt)
                if edge_key in self._edge_map:
                    self._edge_map[edge_key]["weight"] += 1
                    self._edge_map[edge_key]["chunk_ids"].add(chunk_id)
                else:
                    self._edge_map[edge_key] = {"weight": 1, "chunk_ids": {chunk_id}}
                edges_added += 1

        return {
            "input_entities": len(entities),
            "effective_entities": len(effective_entities),
            "edges_added": edges_added
        }

    def _increment_online_hub_stats(self, entity_text: str, weight: float):
        if self._online_hub_threshold <= 0:
            return
        s = self._online_hub_stats.get(entity_text)
        if s is None:
            s = {"degree": 0, "total_weight": 0}
            self._online_hub_stats[entity_text] = s
        s["degree"] += 1
        s["total_weight"] += weight

    def _check_and_handle_runtime_hub(self, entity_text: str):
        if self._online_hub_threshold <= 0:
            return
        if entity_text in self._runtime_hub_entities:
            return
        s = self._online_hub_stats.get(entity_text)
        if s is None:
            return

        is_protected = entity_text in self._protected
        limit = (self._online_protected_hub_threshold if is_protected
                 else self._online_hub_threshold)
        if limit > 0 and (s["degree"] >= limit or s["total_weight"] >= limit):
            self._runtime_hub_entities.add(entity_text)
            self._online_hub_stats.pop(entity_text, None)
            self._prune_entity_from_graph(entity_text)

    def _prune_entity_from_graph(self, entity_text: str):
        for key in list(self._entity_map.keys()):
            if key[0] == entity_text:
                del self._entity_map[key]
        for cid in list(self._chunk_entity_map.keys()):
            refs = self._chunk_entity_map[cid]
            self._chunk_entity_map[cid] = [r for r in refs if r["text"] != entity_text]
        for key in list(self._edge_map.keys()):
            if key[0] == entity_text or key[1] == entity_text:
                del self._edge_map[key]  # mirrors Java pruneEntityFromGraph()

    def apply_hub_filter(self, threshold: int) -> int:
        """
        Remove hub entities and their edges from in-memory structures.
        Mirrors Java InMemoryGraphBuilder.applyHubFilter().
        Uses edge degree+weight across all stored edges (not frequency field).
        Protected entities use 5x fallback threshold.
        Returns number of hub entities removed.
        """
        if threshold <= 0:
            return 0

        # Compute degree and total_weight per entity from edge map
        hub_stats: Dict[str, dict] = {}
        for (ft, tt), ev in self._edge_map.items():
            w = ev["weight"]
            for text in (ft, tt):
                s = hub_stats.get(text)
                if s is None:
                    s = {"degree": 0, "total_weight": 0}
                    hub_stats[text] = s
                s["degree"] += 1
                s["total_weight"] += w

        fallback = threshold * 5
        hub_entities: Set[str] = set()
        for text, s in hub_stats.items():
            is_protected = text in self._protected
            if not is_protected:
                if s["degree"] >= threshold or s["total_weight"] >= threshold:
                    hub_entities.add(text)
            else:
                if fallback > 0 and (s["degree"] >= fallback
                                     or s["total_weight"] >= fallback):
                    hub_entities.add(text)

        if not hub_entities:
            return 0

        print(f"[InMemoryGraphBuilder] hub_filter: threshold={threshold} candidates={len(hub_entities)}",
              file=sys.stderr)

        # Remove from entity map
        for key in list(self._entity_map.keys()):
            if key[0] in hub_entities:
                del self._entity_map[key]

        # Remove from chunk-entity map
        for cid in list(self._chunk_entity_map.keys()):
            refs = self._chunk_entity_map[cid]
            cleaned = [r for r in refs if r["text"] not in hub_entities]
            self._chunk_entity_map[cid] = cleaned

        # Remove edges involving hubs
        for key in list(self._edge_map.keys()):
            if key[0] in hub_entities or key[1] in hub_entities:
                del self._edge_map[key]

        return len(hub_entities)

    def _prune_edge_stats(self, entity_text: str):
        """Remove edges involving entity_text from _edge_map (used by online pruning)."""
        for key in list(self._edge_map.keys()):
            if key[0] == entity_text or key[1] == entity_text:
                del self._edge_map[key]

    def get_runtime_hub_snapshot(self) -> str:
        """Return comma-separated runtime hub entity texts for metadata storage."""
        return ",".join(sorted(self._runtime_hub_entities))

    def flush_to_db(self, conn: sqlite3.Connection, collection: str):
        """
        Write all in-memory entities, chunk_entities, and edges to SQLite.
        Clears existing data for this collection first (full rebuild semantics).
        Mirrors Java InMemoryGraphBuilder.flushToDatabase().
        Uses a single transaction for performance.
        """
        ts = int(time.time())
        print(f"[InMemoryGraphBuilder] Flushing to DB: entities={len(self._entity_map)} "
              f"edges={len(self._edge_map)}", file=sys.stderr)
        old_isolation = conn.isolation_level
        try:
            conn.isolation_level = None  # manual transaction mode
            conn.execute("BEGIN")
            # Clear existing graph data - mirrors Java flushToDatabase() exactly:
            # entities/edges filtered by collection; chunk_entities deleted entirely (no collection col)
            conn.execute("DELETE FROM entities WHERE collection=?", (collection,))
            conn.execute("DELETE FROM entity_edges WHERE collection=?", (collection,))
            conn.execute("DELETE FROM chunk_entities")  # Java: db.delete("chunk_entities", null, null)

            # Insert entities, collect text->id mapping
            text_to_id: Dict[str, int] = {}
            for (text, etype), e in self._entity_map.items():
                cursor = conn.execute(
                    "INSERT INTO entities(collection,entity_text,entity_type,language,"
                    "frequency,avg_confidence,first_seen,last_seen) VALUES(?,?,?,?,?,?,?,?)",
                    (collection, text, etype, e["language"],
                     e["frequency"], e["avg_confidence"],
                     e["first_seen"], e["last_seen"])
                )
                text_to_id[text] = cursor.lastrowid

            # Insert chunk_entities
            for chunk_id, refs in self._chunk_entity_map.items():
                for ref in refs:
                    conn.execute(
                        "INSERT OR IGNORE INTO chunk_entities"
                        "(chunk_id,entity_text,entity_type,confidence) VALUES(?,?,?,?)",
                        (chunk_id, ref["text"], ref["type"], ref["confidence"])
                    )

            # Insert edges - mirrors Java flushToDatabase() with chunk_ids JSON array
            for (ft, tt), ev in self._edge_map.items():
                w = ev["weight"]
                chunk_ids_json = json.dumps(sorted(ev["chunk_ids"]))
                conn.execute(
                    "INSERT OR REPLACE INTO entity_edges"
                    "(collection,from_entity,to_entity,weight,chunk_ids,created_at,updated_at) "
                    "VALUES(?,?,?,?,?,?,?)",
                    (collection, ft, tt, w, chunk_ids_json, ts, ts)
                )

            conn.execute("COMMIT")
            print(f"[InMemoryGraphBuilder] Flush complete: entities={len(text_to_id)} "
                  f"edges={len(self._edge_map)}", file=sys.stderr)
        except Exception as e:
            try:
                conn.execute("ROLLBACK")
            except Exception:
                pass
            print(f"[InMemoryGraphBuilder] Flush failed: {e}", file=sys.stderr)
            raise
        finally:
            conn.isolation_level = old_isolation


# ──────────────────────────────────────────────
# Graph RAG fusion pipeline - mirrors Android RagQueryManager.runGraphRagPipeline()
# ──────────────────────────────────────────────

def _normalize_entity_text(text: str) -> Optional[str]:
    """Mirror Android RagQueryManager.normalizeEntityTextForManager()"""
    if not text:
        return None
    t = text.strip()
    return t if t else None


def run_graph_rag_pipeline(
        user_query: str,
        search_results: List[dict],
        vector_db: KnowledgeGraphDB,
        retrieval_count: int,
        ner_handler: PythonNerHandler,
        cfg_retrieval: dict
) -> List[dict]:
    """
    Graph RAG fusion. 100% mirrors Android RagQueryManager.runGraphRagPipeline().
    Fuses vector score + graph score + entity overlap with configurable presets.
    No extra heuristics beyond the Java implementation.
    """
    if not search_results:
        return search_results

    # Config - mirrors Java ConfigManager defaults
    min_edge_weight = cfg_retrieval.get("graph_min_edge_weight", 2)
    max_expand_entities = cfg_retrieval.get("graph_max_expand_entities", 50)
    hub_threshold = cfg_retrieval.get("graph_hub_threshold_query", 300)
    confidence_threshold = cfg_retrieval.get("graph_entity_confidence_threshold", 0.7)
    max_expand_chunks = cfg_retrieval.get("graph_max_expand_chunks", 50)
    preset = cfg_retrieval.get("graph_rag_weight_preset", 1)

    # Weight presets - mirrors Java switch(preset) in runGraphRagPipeline()
    if preset == 0:   # vector first
        alpha, beta, gamma = 0.9, 0.1, 0.0
    elif preset == 2: # graph enhanced
        alpha, beta, gamma = 0.4, 0.4, 0.2
    else:             # balanced (default=1)
        alpha, beta, gamma = 0.7, 0.2, 0.1

    # Hub entities - mirrors Java vectorDb.getHubEntities(hubThreshold, protectedEntities)
    protected = ner_handler.custom_words if ner_handler else set()
    hub_entities = vector_db.get_hub_entities(hub_threshold, protected)

    def _is_stopword(entity_text: str) -> bool:
        if ner_handler:
            return ner_handler.matches_stopword(entity_text)
        return False

    # --- Seed entity collection (mirrors Java lines 2248-2304) ---
    # Part 1: query NER entities
    query_entities = ner_handler.extract_entities(user_query, confidence_threshold) if ner_handler else []
    query_entity_texts: Set[str] = set()
    seed_order: List[str] = []
    alias_normalized_count = 0
    for ent in query_entities:
        if ent.get("confidence", 1.0) < confidence_threshold:
            continue
        norm = _normalize_entity_text(ent["text"])
        if norm and ner_handler:
            canonical = ner_handler.normalize_text(norm)
            if canonical and canonical != norm:
                alias_normalized_count += 1
                norm = canonical
        if not norm:
            continue
        if norm in hub_entities or _is_stopword(norm):
            continue
        if norm not in query_entity_texts:
            query_entity_texts.add(norm)
            seed_order.append(norm)

    # Part 2: entities from vector top-5 chunks (Java: maxEntitySource = Math.min(5, searchResults.size()))
    max_entity_source = min(5, len(search_results))
    top_chunk_ids = [r["id"] for r in search_results[:max_entity_source]]
    entities_for_top = vector_db.get_entities_for_chunks(top_chunk_ids)
    for ent_list in entities_for_top.values():
        for entity_text in (ent_list or []):
            norm = _normalize_entity_text(entity_text)
            if norm and ner_handler:
                canonical = ner_handler.normalize_text(norm)
                if canonical and canonical != norm:
                    alias_normalized_count += 1
                    norm = canonical
            if not norm:
                continue
            if norm in hub_entities or _is_stopword(norm):
                continue
            if norm not in query_entity_texts:
                query_entity_texts.add(norm)
                seed_order.append(norm)

    if alias_normalized_count > 0:
        print(f"[GraphRAG] alias_normalized={alias_normalized_count}", file=sys.stderr)

    # Cap seed set at GRAPH_RAG_MAX_SEED_ENTITIES_MANAGER=50 (Java constant)
    seed_entities: Set[str] = set(seed_order[:50])

    if not seed_entities:
        print("[GraphRAG] No seed entities after filtering, fallback to vector-only", file=sys.stderr)
        return search_results[:retrieval_count]

    # After building seed set, use it as queryEntityTexts for overlap (mirrors Java line 2343-2344)
    query_entity_texts.clear()
    query_entity_texts.update(seed_entities)

    # --- Graph expansion (mirrors Java lines 2345-2416) ---
    connected = vector_db.get_connected_entities(seed_entities, min_edge_weight, max_expand_entities)
    graph_weight_map: Dict[str, int] = {}
    for ce in connected:
        norm = _normalize_entity_text(ce["entity_text"])
        if norm and ner_handler:
            canonical = ner_handler.normalize_text(norm)
            if canonical:
                norm = canonical
        if not norm or norm in hub_entities or _is_stopword(norm):
            continue
        existing = graph_weight_map.get(norm, 0)
        if ce["weight"] > existing:
            graph_weight_map[norm] = ce["weight"]

    all_entity_texts: Set[str] = set(seed_entities)
    for ce in connected:
        norm = _normalize_entity_text(ce["entity_text"])
        if norm and ner_handler:
            canonical = ner_handler.normalize_text(norm)
            if canonical:
                norm = canonical
        if norm and norm not in hub_entities and not _is_stopword(norm):
            all_entity_texts.add(norm)

    # Get graph-expanded chunk IDs via getChunkIdsByEntities
    entity_text_list = list(all_entity_texts)
    graph_chunk_ids = vector_db.get_chunk_ids_by_entities(entity_text_list)

    # Cap at max_expand_chunks (mirrors Java lines 2411-2413)
    if max_expand_chunks > 0 and len(graph_chunk_ids) > max_expand_chunks:
        graph_chunk_ids = graph_chunk_ids[:max_expand_chunks]

    graph_chunks = vector_db.get_chunks_by_ids(graph_chunk_ids)

    # --- Build candidate map (mirrors Java lines 2424-2447) ---
    # First add all vector search results (vectorRank = i)
    candidate_map: Dict[int, dict] = {}
    for i, result in enumerate(search_results):
        if result["id"] not in candidate_map:
            candidate_map[result["id"]] = {
                "result": result,
                "vector_score": result["similarity"],
                "graph_score": 0.0,
                "entity_overlap": 0,
                "final_score": result["similarity"],
                "vector_rank": i,
            }
    # Then add graph-expanded chunks not already present (vectorRank = -1)
    for result in graph_chunks:
        if result["id"] not in candidate_map:
            candidate_map[result["id"]] = {
                "result": result,
                "vector_score": result["similarity"],
                "graph_score": 0.0,
                "entity_overlap": 0,
                "final_score": result["similarity"],
                "vector_rank": -1,
            }

    # --- First pass: compute raw graph and overlap scores + collect min/max (mirrors Java lines 2484-2528) ---
    all_chunk_ids = list(candidate_map.keys())
    entities_for_all = vector_db.get_entities_for_chunks(all_chunk_ids)

    candidates = []
    vec_min, vec_max = float("inf"), float("-inf")
    graph_min, graph_max = float("inf"), float("-inf")
    overlap_min, overlap_max = float("inf"), float("-inf")

    for cid, cand in candidate_map.items():
        ents = entities_for_all.get(cid, [])
        overlap = 0
        g_score = 0.0
        for entity_text in (ents or []):
            norm = _normalize_entity_text(entity_text)
            if not norm or _is_stopword(norm) or norm in hub_entities:
                continue
            if norm in query_entity_texts:
                overlap += 1
            w = graph_weight_map.get(norm, 0)
            if w:
                g_score += w
        if g_score < 0.0:
            g_score = 0.0
        cand["entity_overlap"] = overlap
        # Adjust graph score by overlap ratio: candidates with no query-entity overlap
        # get graphScore=0, preventing hub-entity pollution from unrelated documents.
        overlap_ratio = overlap / max(1, len(query_entity_texts))
        cand["graph_score"] = math.log1p(g_score) * overlap_ratio

        v = cand["vector_score"]
        g = cand["graph_score"]
        o = cand["entity_overlap"]
        if v < vec_min: vec_min = v
        if v > vec_max: vec_max = v
        if g < graph_min: graph_min = g
        if g > graph_max: graph_max = g
        if o < overlap_min: overlap_min = o
        if o > overlap_max: overlap_max = o

        candidates.append(cand)

    # --- Second pass: normalize and compute final score (mirrors Java lines 2531-2553) ---
    for cand in candidates:
        vec_norm = ((cand["vector_score"] - vec_min) / (vec_max - vec_min)
                    if vec_max > vec_min else (1.0 if vec_max > 0 else 0.0))
        graph_norm = ((cand["graph_score"] - graph_min) / (graph_max - graph_min)
                      if graph_max > graph_min else (1.0 if graph_max > 0 else 0.0))
        overlap_norm = ((cand["entity_overlap"] - overlap_min) / (overlap_max - overlap_min)
                        if overlap_max > overlap_min else (1.0 if overlap_max > 0 else 0.0))
        # Java: c.finalScore = alpha * vecNorm + beta * graphNorm + gamma * overlapNorm
        cand["final_score"] = alpha * vec_norm + beta * graph_norm + gamma * overlap_norm

    candidates.sort(key=lambda x: x["final_score"], reverse=True)
    limit = min(retrieval_count, len(candidates))

    print(f"[GraphRAG] Fused top-{limit} candidates (preset={preset}, alpha={alpha}, beta={beta}, gamma={gamma})")
    for i, cand in enumerate(candidates[:limit]):
        print(f"  #{i+1} id={cand['result']['id']} vecRank={cand['vector_rank']} "
              f"vec={cand['vector_score']:.3f} graph={cand['graph_score']:.3f} "
              f"overlap={cand['entity_overlap']} final={cand['final_score']:.3f}")

    return [cand["result"] for cand in candidates[:limit]]


# ──────────────────────────────────────────────
# Main KnowledgeBase class - three public interfaces
# ──────────────────────────────────────────────

class KnowledgeBase:
    """
    Main entry point. Provides three interfaces matching Android RAG system.
    All behavior aligns with Android TextChunkProcessor + RagQueryManager.
    """

    def __init__(self, config_path=None):
        self.cfg = load_config(config_path)
        self._embedding_client: Optional[EmbeddingClient] = None
        self._ner_handler: Optional[PythonNerHandler] = None
        self._text_splitter: Optional[TextSplitter] = None
        self._json_processor: Optional[JsonDatasetProcessor] = None
        # BM25 index cache: kb_name -> {"index": ..., "chunks": [...], "chunk_count": N}
        # Invalidated when chunk count changes (after build/add_note)
        self._bm25_cache: Dict[str, dict] = {}
        self._init_components()

    def _init_components(self):
        emb_cfg = self.cfg["embedding"]
        chunk_cfg = self.cfg["chunking"]
        ner_cfg = self.cfg["graph_ner"]
        embedding_threads = int(emb_cfg.get("thread_num", 0) or 0)
        if embedding_threads <= 0:
            embedding_threads = min(max(1, os.cpu_count() or 4), 8)

        print("[KnowledgeBase] Initializing components...", file=sys.stderr)

        ensure_model_downloaded(emb_cfg["model_dir"])

        self._embedding_client = EmbeddingClient(
            model_dir=emb_cfg["model_dir"],
            thread_num=embedding_threads
        )

        self._text_splitter = TextSplitter(
            chunk_size=chunk_cfg["chunk_size"],
            chunk_overlap=chunk_cfg["chunk_overlap"],
            min_chunk_size=chunk_cfg["min_chunk_size"]
        )

        self._json_processor = JsonDatasetProcessor(
            min_chunk_size=chunk_cfg["min_chunk_size"]
        )

        self._ner_handler = PythonNerHandler(
            custom_dict_path=ner_cfg.get("custom_dict_path"),
            stopwords_path=ner_cfg.get("stopwords_path")
        ) if ner_cfg.get("enabled", True) else None

        print("[KnowledgeBase] Ready.", file=sys.stderr)

    def _persist_kb_metadata(self, kb_name: str, embedding_dim: int,
                              hub_threshold_build: int = 0,
                              runtime_hubs: str = "") -> None:
        storage = Path(self.cfg["knowledge_base"]["storage_dir"])
        kb_dir = storage / kb_name
        kb_dir.mkdir(parents=True, exist_ok=True)

        model_dir_cfg = Path(self.cfg["embedding"]["model_dir"])
        model_dir_name = model_dir_cfg.name
        reranker_cfg = self.cfg.get("reranker") or self.cfg.get("rerank") or {}
        reranker_model = (
            reranker_cfg.get("model_dir")
            or reranker_cfg.get("model")
            or reranker_cfg.get("reranker_model")
            or ""
        )
        if reranker_model:
            reranker_model = Path(reranker_model).name

        metadata_json = {
            "knowledgeBase": kb_name,
            "embeddingModel": model_dir_name,
            "modeldir": model_dir_name,
            "embeddingDimension": embedding_dim,
            "updated": int(time.time() * 1000)
        }
        if reranker_model:
            metadata_json["rerankerModel"] = reranker_model
        if hub_threshold_build > 0:
            metadata_json["hubThreshold"] = hub_threshold_build
        if runtime_hubs:
            metadata_json["runtimeHubEntities"] = runtime_hubs

        metadata_path = kb_dir / "metadata.json"
        with open(metadata_path, "w", encoding="utf-8") as f:
            json.dump(metadata_json, f, ensure_ascii=False, separators=(",", ":"))

    def _get_db(self, kb_name: str = None) -> KnowledgeGraphDB:
        """Open (or create) knowledge base database."""
        name = kb_name or self.cfg["knowledge_base"]["default_name"]
        storage = Path(self.cfg["knowledge_base"]["storage_dir"])
        kb_dir = storage / name
        kb_dir.mkdir(parents=True, exist_ok=True)
        db_path = str(kb_dir / "knowledge_graph.db")
        return KnowledgeGraphDB(db_path, name)

    def _get_preprocess_workers(self, file_count: int) -> int:
        """Determine number of workers for parallel file preprocessing."""
        chunk_cfg = self.cfg["chunking"]
        workers = chunk_cfg.get("preprocess_workers", 4)
        if workers <= 0:
            workers = min(max(1, os.cpu_count() or 4), 8)
        # For small file counts, reduce workers to avoid overhead
        return min(workers, max(1, file_count))

    def _format_time_brief(self, seconds: float) -> str:
        """Format seconds to brief human-readable string."""
        if seconds < 60:
            return f"{seconds:.1f}s"
        m = int(seconds // 60)
        s = seconds % 60
        return f"{m}m{s:.0f}s"

    def _estimate_eta(self, done: int, total: int, elapsed: float) -> str:
        """Estimate remaining time based on progress."""
        if done <= 0 or total <= 0:
            return "?"
        rate = done / elapsed if elapsed > 0 else 0
        if rate <= 0:
            return "?"
        remaining = (total - done) / rate
        return self._format_time_brief(remaining)

    def _build_file_plan(self, file_path: str, enable_json: bool = True) -> dict:
        """
        Parse a file, split into chunks, return a plan dict.
        Mirrors Android TextChunkProcessor file processing logic.
        Returns:
          {
            file_path, file_name,
            units: [{source, text, chars, optimized_json, chunks}],
            unit_count, chunk_count
          }
        """
        from pathlib import Path as _Path
        fpath = str(file_path)
        fname = _Path(fpath).name

        # Extract raw (text, source, is_json_optimized) units from file
        raw_units = extract_text_from_file(fpath, self._json_processor, enable_json)

        units = []
        for raw_text, source, optimized_json in raw_units:
            if not raw_text or not raw_text.strip():
                continue
            # JSON-optimized units are stored as-is (mirrors Android: JSON items skip splitTextIntoChunks)
            if optimized_json:
                chunks = [raw_text]
            else:
                chunks = self._text_splitter.split_text(raw_text)
            if not chunks:
                continue
            units.append({
                "source": source,
                "text": raw_text,
                "chars": len(raw_text),
                "optimized_json": optimized_json,
                "chunks": chunks,
            })

        total_chunks = sum(len(u["chunks"]) for u in units)
        return {
            "file_path": fpath,
            "file_name": fname,
            "units": units,
            "unit_count": len(units),
            "chunk_count": total_chunks,
        }

    def _embed_chunk_to_db(self, db: KnowledgeGraphDB, text: str, source: str,
                            metadata: str = "") -> int:
        """
        Embed one chunk and insert the document row only (no graph writes).
        Graph entities are accumulated via InMemoryGraphBuilder separately.
        Returns chunk_id.
        """
        embedding = self._embedding_client.compute_embedding(text)
        chunk_id = db.add_chunk(text, source, embedding, metadata)
        return chunk_id

    def build_kb(self, files: List[str], kb_name: str = None,
                  progress_callback=None) -> dict:
        """
        Build/append knowledge base from file list.
        Mirrors Android TextChunkProcessor.processFilesAndBuildKnowledgeBase().

        Graph construction uses InMemoryGraphBuilder:
          1. Embed all chunks, insert document rows one by one.
          2. Accumulate entities/edges entirely in RAM.
          3. Apply hub threshold filter in memory.
          4. Flush cleaned graph to SQLite in a single transaction.
        This avoids N^2 individual SQLite transactions during build.

        Supported file types: txt, md, json, jsonl, pdf, docx, pptx, xlsx, csv, html
        JSON training set formats: Alpaca, CoT, DPO, Conversation (auto-detected)

        Returns: dict with stats (chunks_added, entities_added, etc.)
        """
        if not files:
            print("[build_kb] No files provided", file=sys.stderr)
            return {"status": "error", "message": "No files provided"}

        chunk_cfg = self.cfg["chunking"]
        retrieval_cfg = self.cfg["retrieval"]
        ner_cfg = self.cfg["graph_ner"]
        enable_json = chunk_cfg.get("enable_json_dataset_opt", True)
        hub_threshold_build = retrieval_cfg.get("graph_hub_threshold_build", 1000)
        ner_enabled = ner_cfg.get("enabled", True) and self._ner_handler is not None

        db = self._get_db(kb_name)
        db_name = kb_name or self.cfg["knowledge_base"]["default_name"]

        # InMemoryGraphBuilder mirrors Java TextChunkProcessor.InMemoryGraphBuilder
        protected = self._ner_handler.custom_words if self._ner_handler else set()
        graph_builder = InMemoryGraphBuilder(
            protected_entity_texts=protected,
            hub_threshold=hub_threshold_build
        ) if ner_enabled else None

        total_chunks_added = 0
        total_files = 0
        total_errors = 0
        emb_dim = 0
        total_extracted_units = 0
        total_entities_seen = 0
        total_effective_entities = 0
        total_edges_added = 0
        start_time = time.time()
        planning_sec = 0.0
        embedding_sec = 0.0
        graph_sec = 0.0
        hub_filter_sec = 0.0
        flush_sec = 0.0
        planned_units = 0
        planned_chunks = 0
        preprocess_workers = self._get_preprocess_workers(len(files))

        print(f"[build_kb] START  files={len(files)} → KB='{db_name}'  "
              f"chunk={chunk_cfg['chunk_size']}/{chunk_cfg['chunk_overlap']} hub={hub_threshold_build}",
              file=sys.stderr)

        try:
            planning_start = time.time()
            indexed_files = [(index, str(fpath)) for index, fpath in enumerate(files, start=1)]
            planned_results = []

            if preprocess_workers <= 1:
                for file_index, fpath in indexed_files:
                    try:
                        file_plan = self._build_file_plan(fpath, enable_json)
                        planned_results.append((file_index, file_plan))
                        print(
                            f"[build_kb][1/3 extract] [{file_index}/{len(files)}] {file_plan['file_name']} "
                            f"units={file_plan['unit_count']} chunks={file_plan['chunk_count']}",
                            file=sys.stderr
                        )
                    except Exception as e:
                        print(f"[build_kb][1/3 extract] ERROR {fpath}: {e}", file=sys.stderr)
                        total_errors += 1
            else:
                with ThreadPoolExecutor(max_workers=preprocess_workers) as executor:
                    future_map = {
                        executor.submit(self._build_file_plan, fpath, enable_json): (file_index, fpath)
                        for file_index, fpath in indexed_files
                    }
                    completed_files = 0
                    for future in as_completed(future_map):
                        file_index, fpath = future_map[future]
                        completed_files += 1
                        try:
                            file_plan = future.result()
                            planned_results.append((file_index, file_plan))
                            print(
                                f"[build_kb][1/3 extract] [{completed_files}/{len(files)}] {file_plan['file_name']} "
                                f"units={file_plan['unit_count']} chunks={file_plan['chunk_count']}",
                                file=sys.stderr
                            )
                        except Exception as e:
                            print(f"[build_kb][1/3 extract] ERROR {fpath}: {e}", file=sys.stderr)
                            total_errors += 1

            planned_results.sort(key=lambda item: item[0])
            file_plans = [item[1] for item in planned_results]
            planning_sec = time.time() - planning_start
            total_extracted_units = sum(plan["unit_count"] for plan in file_plans)
            planned_units = total_extracted_units
            planned_chunks = sum(plan["chunk_count"] for plan in file_plans)

            print(
                f"[build_kb][1/3 extract] DONE  files={len(file_plans)}/{len(files)} "
                f"units={planned_units} chunks={planned_chunks} "
                f"elapsed={self._format_time_brief(planning_sec)}",
                file=sys.stderr
            )

            global_unit_index = 0
            progress_interval = max(1, planned_chunks // 100)  # Print every 1%
            last_printed_pct = -1
            
            for file_index, file_plan in enumerate(file_plans, start=1):
                try:
                    for file_unit_index, unit_plan in enumerate(file_plan["units"], start=1):
                        global_unit_index += 1
                        source = unit_plan["source"]
                        chunks = unit_plan["chunks"]

                        for i, chunk_text in enumerate(chunks):
                            embed_start = time.time()
                            chunk_id = self._embed_chunk_to_db(
                                db, chunk_text, source,
                                f'{{"chunk_index":{i}}}'
                            )
                            embedding_sec += time.time() - embed_start
                            total_chunks_added += 1

                            if graph_builder is not None:
                                graph_start = time.time()
                                entities = self._ner_handler.extract_entities(chunk_text)
                                ner_count = len(entities)
                                if entities:
                                    entities = [
                                        e for e in entities
                                        if not self._ner_handler.matches_stopword(e["text"])
                                    ]
                                    graph_stats = graph_builder.add_chunk_entities(
                                        chunk_id, entities, self._ner_handler
                                    )
                                else:
                                    graph_stats = {"input_entities": 0, "effective_entities": 0, "edges_added": 0}
                                graph_sec += time.time() - graph_start
                                total_entities_seen += ner_count
                                total_effective_entities += graph_stats["effective_entities"]
                                total_edges_added += graph_stats["edges_added"]
                                pass  # per-chunk graph debug suppressed to reduce noise

                            # Print progress every 1% or at completion
                            current_pct = int((total_chunks_added * 100.0 / planned_chunks)) if planned_chunks else 100
                            if (total_chunks_added % progress_interval == 0 and current_pct != last_printed_pct) or total_chunks_added == planned_chunks:
                                elapsed = time.time() - start_time
                                eta = self._estimate_eta(total_chunks_added, planned_chunks, elapsed)
                                print(
                                    f"[build_kb][2/3 embed]  [{total_chunks_added}/{planned_chunks}] "
                                    f"{current_pct}%  elapsed={self._format_time_brief(elapsed)} eta={eta}",
                                    file=sys.stderr
                                )
                                last_printed_pct = current_pct
                                if progress_callback:
                                    progress_callback(current_pct, total_chunks_added, planned_chunks)

                    total_files += 1

                except Exception as e:
                    print(f"[build_kb][2/3 embed]  ERROR {file_plan['file_path']}: {e}", file=sys.stderr)
                    total_errors += 1

            # Step 3: Apply hub threshold filter in memory
            hubs_removed = 0
            if graph_builder is not None and hub_threshold_build > 0:
                hub_start = time.time()
                hubs_removed = graph_builder.apply_hub_filter(hub_threshold_build)
                hub_filter_sec = time.time() - hub_start
                print(
                    f"[build_kb][3/3 graph]  hub_filter: removed={hubs_removed} "
                    f"elapsed={self._format_time_brief(hub_filter_sec)}",
                    file=sys.stderr
                )

            # Step 4: Flush cleaned graph to DB in a single transaction
            if graph_builder is not None and total_chunks_added > 0:
                flush_start = time.time()
                graph_builder.flush_to_db(db.conn, db_name)
                flush_sec = time.time() - flush_start
                print(
                    f"[build_kb][3/3 graph]  flush: entities={len(graph_builder._entity_map)} "
                    f"edges={len(graph_builder._edge_map)} elapsed={self._format_time_brief(flush_sec)}",
                    file=sys.stderr
                )

            # Update DB metadata
            row = db.conn.execute(
                "SELECT embedding FROM documents WHERE collection=? LIMIT 1",
                (db_name,)
            ).fetchone()
            emb_dim = len(_blob_to_vector(row[0])) if row else 0
            runtime_hubs = (
                graph_builder.get_runtime_hub_snapshot()
                if graph_builder else ""
            )
            db.update_metadata(
                model_dir=Path(self.cfg["embedding"]["model_dir"]).name,
                embedding_dim=emb_dim,
                hub_threshold_build=hub_threshold_build
            )
            if runtime_hubs:
                db.conn.execute(
                    "INSERT OR REPLACE INTO metadata(key,value) VALUES(?,?)",
                    ("runtime_hub_entities", runtime_hubs)
                )
                db.conn.commit()
            self._persist_kb_metadata(
                kb_name=db_name,
                embedding_dim=emb_dim,
                hub_threshold_build=hub_threshold_build,
                runtime_hubs=runtime_hubs
            )

        finally:
            db.close()

        elapsed = time.time() - start_time
        stats = {
            "status": "ok",
            "kb_name": db_name,
            "files_processed": total_files,
            "files_errors": total_errors,
            "chunks_added": total_chunks_added,
            "extracted_units": total_extracted_units,
            "planned_chunks": planned_chunks,
            "entities_seen": total_entities_seen,
            "effective_entities": total_effective_entities,
            "edges_added": total_edges_added,
            "hubs_removed": hubs_removed if graph_builder else 0,
            "embedding_dim": emb_dim,
            "planning_sec": round(planning_sec, 1),
            "embedding_sec": round(embedding_sec, 1),
            "graph_sec": round(graph_sec, 1),
            "hub_filter_sec": round(hub_filter_sec, 1),
            "flush_sec": round(flush_sec, 1),
            "elapsed_sec": round(elapsed, 1)
        }
        print(
            f"[build_kb] DONE  kb='{db_name}' files={total_files}/{total_files+total_errors} "
            f"chunks={total_chunks_added} entities={total_effective_entities} "
            f"edges={total_edges_added} hubs_removed={hubs_removed if graph_builder else 0} "
            f"dim={emb_dim} elapsed={self._format_time_brief(elapsed)}",
            file=sys.stderr
        )
        return stats

    def add_note(self, text: str, kb_name: str = None, source: str = "note", title: str = "") -> dict:
        """
        Add a free-form text note directly to the knowledge base.
        Mirrors Android KnowledgeNoteFragment behavior:
          - title + content are treated as one complete structured note
          - one note => one embedding => one document row
          - entities/edges are appended incrementally to existing graph
          - applies hub threshold (build threshold) after insert
        Returns dict with stats.
        """
        if not text or not text.strip():
            return {"status": "error", "message": "Empty text"}

        db = self._get_db(kb_name)
        db_name = kb_name or self.cfg["knowledge_base"]["default_name"]
        retrieval_cfg = self.cfg["retrieval"]
        ner_cfg = self.cfg["graph_ner"]
        hub_threshold_build = retrieval_cfg.get("graph_hub_threshold_build", 1000)
        ner_enabled = ner_cfg.get("enabled", True) and self._ner_handler is not None

        start_time = time.time()
        note_title = title.strip()
        note_text = text.strip()
        full_text = f"{note_title}\n\n{note_text}" if note_title else note_text
        note_source = note_title if note_title else source
        print(f"[add_note] kb='{db_name}' chars={len(full_text)}", file=sys.stderr)

        total_chunks_added = 0
        emb_dim = 0
        hubs_removed = 0
        try:
            chunk_id = self._embed_chunk_to_db(db, full_text, note_source, "")
            total_chunks_added = 1

            if ner_enabled:
                entities = self._ner_handler.extract_entities(full_text)
                if entities:
                    entities = [
                        e for e in entities
                        if not self._ner_handler.matches_stopword(e["text"])
                    ]
                    if entities:
                        db.add_entities_and_build_graph(
                            chunk_id, entities, self._ner_handler
                        )

            if hub_threshold_build > 0:
                protected = self._ner_handler.custom_words if self._ner_handler else set()
                hubs_removed = db.apply_hub_threshold(hub_threshold_build, protected)
                print(f"[add_note] hub_filter removed={hubs_removed}", file=sys.stderr)

            # Update metadata
            row = db.conn.execute(
                "SELECT embedding FROM documents WHERE collection=? LIMIT 1",
                (db_name,)
            ).fetchone()
            emb_dim = len(_blob_to_vector(row[0])) if row else 0
            db.update_metadata(
                model_dir=Path(self.cfg["embedding"]["model_dir"]).name,
                embedding_dim=emb_dim,
                hub_threshold_build=hub_threshold_build
            )
            self._persist_kb_metadata(
                kb_name=db_name,
                embedding_dim=emb_dim,
                hub_threshold_build=hub_threshold_build
            )

        finally:
            db.close()

        elapsed = time.time() - start_time
        stats = {
            "status": "ok",
            "kb_name": db_name,
            "chunks_added": total_chunks_added,
            "hubs_removed": hubs_removed,
            "elapsed_sec": round(elapsed, 1)
        }
        print(f"[add_note] DONE  chunks_added={total_chunks_added} elapsed={elapsed:.1f}s", file=sys.stderr)
        return stats

    # ── Interface 3: RAG Query ───────────────────────────────────────────────

    def query_kb(self, prompt: str, kb_name: str = None) -> str:
        """
        RAG retrieval. Returns combined context string (N chunks joined).
        Mirrors Android RagQueryManager.runRagRetrievalPipeline() + runGraphRagPipeline().

        Steps:
          1. Compute query embedding
          2. Vector search (topK = search_depth + graph_expand if graph enabled)
          2b. BM25 search (if bm25_enabled) + RRF fusion with vector results
          3. Graph RAG fusion (if enabled)
          4. Return top-N results as combined text string
        """
        if not prompt or not prompt.strip():
            return ""

        db = self._get_db(kb_name)
        db_name = kb_name or self.cfg["knowledge_base"]["default_name"]
        retrieval_cfg = self.cfg["retrieval"]

        search_depth = retrieval_cfg.get("search_depth", 20)
        graph_enabled = retrieval_cfg.get("graph_rag_enabled", True)
        graph_expand = retrieval_cfg.get("graph_rag_vector_expand", 20) if graph_enabled else 0
        vector_top_k = search_depth + graph_expand
        bm25_enabled = retrieval_cfg.get("bm25_enabled", True)

        print(f"[query_kb] kb='{db_name}' depth={search_depth} graph={graph_enabled} bm25={bm25_enabled}",
              file=sys.stderr)

        try:
            # Step 1: Compute query embedding
            query_vector = self._embedding_client.compute_embedding(prompt)

            # Step 2: Vector search
            vec_results = db.search_similar(query_vector, vector_top_k)
            if not vec_results:
                print("[query_kb] no results", file=sys.stderr)
                return ""

            sim_scores = [f"{r['similarity']:.3f}" for r in vec_results]
            print(f"[query_kb] vec={len(vec_results)} top_sim={', '.join(sim_scores[:3])}",
                  file=sys.stderr)

            # Step 2b: BM25 + RRF fusion
            search_results = vec_results
            if bm25_enabled:
                t0 = time.time()
                # Use cached BM25 index; invalidate if chunk count changed
                cache = self._bm25_cache.get(db_name)
                current_count = db.get_chunk_count()
                if cache is None or cache.get("chunk_count") != current_count:
                    all_chunks = db.get_all_chunks_for_bm25()
                    bm25_index = _build_bm25_index([c["content"] for c in all_chunks])
                    self._bm25_cache[db_name] = {
                        "index": bm25_index,
                        "chunks": all_chunks,
                        "chunk_count": current_count,
                    }
                    cache_status = "built"
                else:
                    all_chunks = cache["chunks"]
                    bm25_index = cache["index"]
                    cache_status = "hit"
                bm25_top_k = max(vector_top_k * 2, 60)
                # Score all docs with cached index
                query_tokens = _bm25_tokenize(prompt)
                bm25_results = []
                if query_tokens:
                    for i, chunk in enumerate(all_chunks):
                        s = _bm25_score(bm25_index, query_tokens, i)
                        if s > 0:
                            bm25_results.append({**chunk, "bm25_score": s, "similarity": 0.0})
                    bm25_results.sort(key=lambda x: x["bm25_score"], reverse=True)
                    bm25_results = bm25_results[:bm25_top_k]
                # RRF fusion: expand to larger pool for GraphRAG, then trim to vector_top_k
                search_results = _rrf_fusion(vec_results, bm25_results, vector_top_k)
                elapsed_ms = int((time.time() - t0) * 1000)
                bm25_ids_top5 = [r["id"] for r in bm25_results[:5]]
                rrf_ids_top5 = [r["id"] for r in search_results[:5]]
                print(f"[query_kb] bm25+rrf cache={cache_status} corpus={len(all_chunks)} {elapsed_ms}ms",
                      file=sys.stderr)

            # Step 3: Graph RAG fusion
            final_results = search_results
            if graph_enabled and self._ner_handler and search_results:
                final_results = run_graph_rag_pipeline(
                    user_query=prompt,
                    search_results=search_results,
                    vector_db=db,
                    retrieval_count=search_depth,
                    ner_handler=self._ner_handler,
                    cfg_retrieval=retrieval_cfg
                )
            else:
                final_results = search_results[:search_depth]

            # Step 4: Build combined context string
            # Mirrors Android RagQueryManager.buildPromptWithKnowledgeBase()
            parts = []
            for i, r in enumerate(final_results):
                parts.append(f"Document{i+1} [ID:{r['id']}]:\n{r['content']}")

            context = "\n\n".join(parts)
            print(f"[query_kb] DONE  returned={len(final_results)} docs chars={len(context)}",
                  file=sys.stderr)
            return context

        finally:
            db.close()

    def get_build_status(self, kb_name: str = None) -> dict:
        """
        Return current build status. Instant return, no KB init needed.
        Delegates to module-level _is_build_running() + _read_build_status().
        Returns dict with keys: status (building|ok|error|idle), progress, message, stats, error.
        """
        running = _is_build_running()
        st = _read_build_status()

        if running:
            return {
                "status": "building",
                "progress": st.get("progress", 0),
                "message": st.get("message", ""),
            }

        state = st.get("state", "idle")
        if state == "ok":
            return {"status": "ok", "message": st.get("message", "Build complete"), "stats": st.get("stats", {})}
        elif state == "error":
            return {"status": "error", "error": st.get("error", "Unknown error")}
        elif state == "building":
            # Process died without updating status (crash or SIGKILL)
            err = "Build process exited unexpectedly (crash or killed)"
            _write_build_status("error", progress=0, message=err, error=err)
            return {"status": "error", "error": err}
        else:
            return {"status": "idle", "message": "No build has run yet."}


# ──────────────────────────────────────────────
# Module-level convenience functions
# ──────────────────────────────────────────────

_default_kb: Optional[KnowledgeBase] = None


def _get_default_kb(config_path=None) -> KnowledgeBase:
    global _default_kb
    if _default_kb is None:
        _default_kb = KnowledgeBase(config_path)
    return _default_kb


def build_kb(files: List[str], kb_name: str = None,
             config_path: str = None) -> dict:
    """Module-level interface 1: build knowledge base from files."""
    return _get_default_kb(config_path).build_kb(files, kb_name)


def add_note(text: str, kb_name: str = None,
             config_path: str = None, title: str = "") -> dict:
    """Module-level interface 2: add text note to knowledge base."""
    return _get_default_kb(config_path).add_note(text, kb_name, title=title)


def query_kb(prompt: str, kb_name: str = None,
             config_path: str = None) -> str:
    """Module-level interface 3: RAG retrieval, returns context string."""
    return _get_default_kb(config_path).query_kb(prompt, kb_name)


# ──────────────────────────────────────────────
# CLI entry point
# Run directly: python py_mnn_kb.py build/note/query/status ...
# build runs as detached background process (pidfile + status json)
# note/query run synchronously; both are blocked while build is running
# ──────────────────────────────────────────────

# Build state files (project root, one level above scripts/)
_HERE = Path(__file__).parent.parent
_BUILD_PID_FILE = _HERE / "build.pid"
_BUILD_STATUS_FILE = _HERE / "build_status.json"


def _write_build_status(state: str, progress: int = 0, message: str = "",
                        error: str = "", stats: dict = None) -> None:
    """Write build status to build_status.json (called from build subprocess)."""
    payload = {
        "state": state,          # building | ok | error
        "progress": progress,    # 0-100
        "message": message,
        "error": error,
        "stats": stats or {},
        "updated_at": time.time(),
    }
    try:
        tmp = str(_BUILD_STATUS_FILE) + ".tmp"
        with open(tmp, "w", encoding="utf-8") as f:
            json.dump(payload, f, ensure_ascii=False)
        import os as _os
        _os.replace(tmp, str(_BUILD_STATUS_FILE))
    except Exception:
        pass


def _read_build_status() -> dict:
    """Read build_status.json; returns empty dict if missing."""
    try:
        with open(_BUILD_STATUS_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {}


def _is_build_running() -> bool:
    """Check whether the build subprocess is still alive via pidfile."""
    try:
        pid_text = _BUILD_PID_FILE.read_text(encoding="utf-8").strip()
        pid = int(pid_text)
    except Exception:
        return False
    try:
        import psutil
        return psutil.pid_exists(pid)
    except ImportError:
        # Fallback without psutil: send signal 0
        import os as _os
        try:
            _os.kill(pid, 0)
            return True
        except (OSError, ProcessLookupError):
            return False


def _building_response(progress: int) -> dict:
    """Standard 'build in progress' response payload."""
    return {
        "status": "building",
        "progress": progress,
        "message": f"Build in progress ({progress}%). Run 'status' to check again.",
    }


def _call_llm_api(prompt: str, context: str, cfg_llm: dict,
                  response_stream=None) -> str:
    """Call OpenAI-compatible LLM API for RAG generation."""
    try:
        import openai
    except ImportError:
        print("[LLM] openai package not installed. Install: pip install openai", file=sys.stderr)
        return ""

    base_url = cfg_llm.get("base_url", "http://localhost:11434/v1")
    api_key = cfg_llm.get("api_key", "none")
    model = cfg_llm.get("model", "qwen2.5:7b")
    max_tokens = cfg_llm.get("max_tokens", 2048)
    temperature = cfg_llm.get("temperature", 0.7)
    system_prompt = cfg_llm.get("system_prompt",
        "You are a helpful assistant. Answer the question based on the provided context.")

    client = openai.OpenAI(base_url=base_url, api_key=api_key)
    user_content = f"Context:\n{context}\n\nQuestion: {prompt}"

    print(f"[LLM] Calling {model} @ {base_url} ...", file=sys.stderr)
    print(f"[LLM] Context: {len(context)} chars, question: {len(prompt)} chars", file=sys.stderr)

    try:
        response = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_content}
            ],
            max_tokens=max_tokens,
            temperature=temperature,
            stream=True
        )
        if response_stream is not None:
            print("[LLM] Response:", file=sys.stderr)
            print("─" * 60, file=sys.stderr)
        full_response = ""
        for chunk in response:
            delta = chunk.choices[0].delta
            if delta.content:
                if response_stream is not None:
                    print(delta.content, end="", flush=True, file=response_stream)
                full_response += delta.content
        if response_stream is not None:
            print(file=response_stream, flush=True)
            print("─" * 60, file=sys.stderr)
        return full_response

    except Exception as e:
        print(f"[LLM] API call failed: {e}", file=sys.stderr)
        return ""


def _emit_json(payload: dict) -> None:
    json.dump(payload, sys.stdout, ensure_ascii=False)
    sys.stdout.write("\n")
    sys.stdout.flush()


_SUPPORTED_EXTS = {
    ".txt", ".md", ".log", ".rst",
    ".json", ".jsonl",
    ".pdf", ".docx", ".pptx",
    ".xlsx", ".xls", ".csv",
    ".html", ".htm"
}


def _collect_files(input_path: Path):
    """Collect supported files from a path (file or directory)."""
    if input_path.is_file():
        return [str(input_path)] if input_path.suffix.lower() in _SUPPORTED_EXTS else []
    return [str(f) for f in input_path.rglob("*")
            if f.is_file() and f.suffix.lower() in _SUPPORTED_EXTS]


def _cmd_build(args, cfg: dict):
    """
    Async build: spawn a detached subprocess to run the actual build,
    write pidfile immediately, return 'started' to the caller.
    The subprocess runs with --_bg_build flag and updates build_status.json.
    """
    import subprocess

    # If already building, return current progress
    if _is_build_running():
        st = _read_build_status()
        progress = st.get("progress", 0)
        print(f"[build] Already building ({progress}%), ignoring new request.", file=sys.stderr)
        return _building_response(progress)

    # Validate path and collect files before launching subprocess
    input_path = Path(args.dir)
    if not input_path.exists():
        raise FileNotFoundError(f"[build] Path not found: {input_path}")
    files = _collect_files(input_path)
    if not files:
        raise ValueError(f"[build] No supported files found in: {input_path}")

    kb_name = args.kb or cfg["knowledge_base"]["default_name"]
    config_path = args.config or str(_HERE / "config.json")

    print(f"[build] Launching background build: {len(files)} files → KB='{kb_name}'",
          file=sys.stderr)

    # Mark status as building before spawning (avoids race if caller polls immediately)
    _write_build_status("building", progress=0, message="Starting...")

    # Build subprocess command: same interpreter, same script, internal --_bg_build flag
    cmd = [
        sys.executable, str(Path(__file__).resolve()),
        "--config", str(config_path),
        "_bg_build",
        "--dir", str(input_path.resolve()),
        "--kb", kb_name,
    ]

    # Build log file: background process writes stderr here
    build_log = _HERE / "build.log"

    if sys.platform == "win32":
        DETACHED_PROCESS = 0x00000008
        CREATE_NEW_PROCESS_GROUP = 0x00000200
        creation_flags = DETACHED_PROCESS | CREATE_NEW_PROCESS_GROUP
        log_f = open(build_log, "w", encoding="utf-8", errors="replace")
        proc = subprocess.Popen(
            cmd,
            stdout=subprocess.DEVNULL,
            stderr=log_f,
            creationflags=creation_flags,
        )
        log_f.close()
    else:
        log_f = open(build_log, "w", encoding="utf-8", errors="replace")
        proc = subprocess.Popen(
            cmd,
            stdout=subprocess.DEVNULL,
            stderr=log_f,
            start_new_session=True,
        )
        log_f.close()

    # Write pidfile
    _BUILD_PID_FILE.write_text(str(proc.pid), encoding="utf-8")
    print(f"[build] Background process started (pid={proc.pid})", file=sys.stderr)

    return {
        "status": "building",
        "command": "build",
        "kb_name": kb_name,
        "pid": proc.pid,
        "files": len(files),
        "message": f"Build started in background (pid={proc.pid}). "
                   f"Use 'status' command to check progress.",
    }


def _run_bg_build(args, cfg: dict):
    """
    Internal: called inside the background subprocess (--_bg_build flag).
    Runs build synchronously, updates build_status.json throughout.
    """
    input_path = Path(args.dir)
    files = _collect_files(input_path)
    kb_name = args.kb or cfg["knowledge_base"]["default_name"]
    config_path = args.config or str(_HERE / "config.json")

    _write_build_status("building", progress=0,
                        message=f"Initializing... ({len(files)} files)")
    print(f"[bg_build] START  files={len(files)} kb='{kb_name}'", file=sys.stderr)

    def on_progress(pct: int, done: int, total: int):
        _write_build_status("building", progress=pct,
                            message=f"Embedding {done}/{total} chunks ({pct}%)")

    try:
        kb = KnowledgeBase(config_path)
        _write_build_status("building", progress=1, message="Model loaded, building...")

        stats = kb.build_kb(files, kb_name, progress_callback=on_progress)

        _write_build_status(
            "ok", progress=100,
            message=f"Build complete: {stats.get('chunks_added', 0)} chunks",
            stats=stats,
        )
        print(f"[bg_build] DONE  chunks={stats.get('chunks_added')} "
              f"elapsed={stats.get('elapsed_sec')}s", file=sys.stderr)

    except Exception as e:
        _write_build_status("error", progress=0,
                            message="Build failed", error=str(e))
        print(f"[bg_build] ERROR  {e}", file=sys.stderr)
        sys.exit(1)


def _cmd_status(args):
    """Return current build status (always instant, no KB init needed)."""
    running = _is_build_running()
    st = _read_build_status()

    if running:
        progress = st.get("progress", 0)
        msg = st.get("message", "")
        print(f"[status] Building... {progress}%  {msg}", file=sys.stderr)
        return {
            "status": "building",
            "progress": progress,
            "message": msg,
        }

    state = st.get("state", "idle")
    if state == "ok":
        stats = st.get("stats", {})
        msg = st.get("message", "Build complete")
        print(f"[status] {msg}", file=sys.stderr)
        return {"status": "ok", "message": msg, "stats": stats}
    elif state == "error":
        err = st.get("error", "Unknown error")
        print(f"[status] Last build failed: {err}", file=sys.stderr)
        return {"status": "error", "error": err}
    elif state == "building":
        # Process died without updating status (crash or SIGKILL)
        err = "Build process exited unexpectedly (crash or killed)"
        _write_build_status("error", progress=0, message=err, error=err)
        print(f"[status] {err}", file=sys.stderr)
        return {"status": "error", "error": err}
    else:
        print("[status] Idle (no build has run yet)", file=sys.stderr)
        return {"status": "idle", "message": "No build has run yet."}


def _cmd_note(args, kb: KnowledgeBase, cfg: dict):
    # Guard: refuse while build is running
    if _is_build_running():
        st = _read_build_status()
        progress = st.get("progress", 0)
        print(f"[note] Build in progress ({progress}%), note refused.", file=sys.stderr)
        return _building_response(progress)

    text = args.text
    if not text.strip():
        raise ValueError("[note] Empty text provided")

    kb_name = args.kb or cfg["knowledge_base"]["default_name"]
    note_title = (args.title or "").strip()
    print(f"[note] Adding note ({len(text)} chars) to KB '{kb_name}'"
          + (f" title='{note_title}'" if note_title else ""), file=sys.stderr)

    stats = kb.add_note(text, kb_name, title=note_title)
    print(f"[note] Done: chunks_added={stats.get('chunks_added')} elapsed={stats.get('elapsed_sec')}s",
          file=sys.stderr)

    return {"status": stats.get("status", "ok"), "command": "note",
            "kb_name": stats.get("kb_name"), "title": note_title, "stats": stats}


def _cmd_query(args, kb: KnowledgeBase, cfg: dict):
    # Guard: refuse while build is running
    if _is_build_running():
        st = _read_build_status()
        progress = st.get("progress", 0)
        print(f"[query] Build in progress ({progress}%), query refused.", file=sys.stderr)
        return _building_response(progress)

    prompt = args.prompt
    if not prompt.strip():
        raise ValueError("[query] Empty prompt provided")

    kb_name = args.kb or cfg["knowledge_base"]["default_name"]
    print(f"[query] KB='{kb_name}' prompt={len(prompt)}chars", file=sys.stderr)

    context = kb.query_kb(prompt, kb_name)

    if not context:
        print("[query] No context retrieved", file=sys.stderr)
        return {"status": "ok", "command": "query", "kb_name": kb_name,
                "prompt": prompt, "llm_used": False, "context": "",
                "context_chars": 0, "answer": "", "answer_chars": 0}

    doc_count = context.count("\n\nDocument") or (1 if context.strip() else 0)
    print(f"[query] Retrieved {doc_count} docs ({len(context)} chars)", file=sys.stderr)

    if args.no_llm:
        if args.output == "text":
            lines = [f"[Retrieved {doc_count} documents from KB '{kb_name}']"]
            for part in context.split("\n\nDocument")[1:]:
                header_end = part.find(":\n")
                if header_end >= 0:
                    snippet = part[header_end + 2:header_end + 200].replace("\n", " ").strip()
                    lines.append(f"  Doc{part[:header_end].strip()}: {snippet}...")
            print("\n".join(lines))
        return {"status": "ok", "command": "query", "kb_name": kb_name,
                "prompt": prompt, "llm_used": False, "context": context,
                "context_chars": len(context), "answer": "", "answer_chars": 0}

    cfg_llm = cfg.get("llm_api", {})
    if not cfg_llm.get("base_url"):
        print("[query] No LLM API configured in config.json [llm_api]", file=sys.stderr)
        if args.output == "text":
            print(context, end="" if context.endswith("\n") else "\n")
        return {"status": "ok", "command": "query", "kb_name": kb_name,
                "prompt": prompt, "llm_used": False, "context": context,
                "context_chars": len(context), "answer": "", "answer_chars": 0}

    print(f"[query] Calling LLM ({cfg_llm.get('model', '?')}) ...", file=sys.stderr)
    response_stream = sys.stdout if args.output == "text" else None
    answer = _call_llm_api(prompt, context, cfg_llm, response_stream=response_stream)

    return {"status": "ok", "command": "query", "kb_name": kb_name,
            "prompt": prompt, "llm_used": True, "context": context,
            "context_chars": len(context), "answer": answer, "answer_chars": len(answer)}


def main():
    # Force UTF-8 output on Windows
    if hasattr(sys.stdout, "reconfigure") and sys.stdout.encoding and \
            sys.stdout.encoding.lower() not in ("utf-8", "utf8"):
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    if hasattr(sys.stderr, "reconfigure") and sys.stderr.encoding and \
            sys.stderr.encoding.lower() not in ("utf-8", "utf8"):
        sys.stderr.reconfigure(encoding="utf-8", errors="replace")

    parser = argparse.ArgumentParser(
        description="py_mnn_kb — MNN Knowledge Base CLI / OpenClaw Skill",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""Commands:
  build <dir>     Index files in background; returns immediately
  status          Check build progress / last result
  note  "<text>"  Add a text note (synchronous; refused while building)
  query "<prompt>" RAG retrieval + optional LLM (synchronous; refused while building)

Examples:
  python py_mnn_kb.py build ../docs/ --kb myproject
  python py_mnn_kb.py status
  python py_mnn_kb.py note "Key finding: ..." --kb myproject --title "Meeting note"
  python py_mnn_kb.py query "What is the deployment process?" --kb myproject
  python py_mnn_kb.py query "..." --kb myproject --no-llm
  python py_mnn_kb.py --output json query "..." --kb myproject
"""
    )
    parser.add_argument("--config", default=None,
                        help="Path to config.json (default: same dir as py_mnn_kb.py)")
    parser.add_argument("--output", default="text", choices=["json", "text"],
                        help="Output format: text (default) or json")

    sub = parser.add_subparsers(dest="command", required=True)

    p_build = sub.add_parser("build", help="Build KB from files (background, non-blocking)")
    p_build.add_argument("dir", help="Directory or file path to index")
    p_build.add_argument("--kb", default=None, help="Knowledge base name")

    sub.add_parser("status", help="Check build progress or last build result")

    p_note = sub.add_parser("note", help="Add a text note to knowledge base (synchronous)")
    p_note.add_argument("text", help="Text content to add")
    p_note.add_argument("--kb", default=None, help="Knowledge base name")
    p_note.add_argument("--title", default="", help="Optional title for the note")

    p_query = sub.add_parser("query", help="RAG retrieval + LLM generation (synchronous)")
    p_query.add_argument("prompt", help="Question or query text")
    p_query.add_argument("--kb", default=None, help="Knowledge base name")
    p_query.add_argument("--no-llm", action="store_true",
                         help="Return raw retrieved context without LLM generation")

    # Internal subcommand used by background subprocess (not shown in help)
    p_bg = sub.add_parser("_bg_build")
    p_bg.add_argument("--dir", required=True)
    p_bg.add_argument("--kb", default=None)

    args = parser.parse_args()

    # Internal background build path: no pidfile check, no guard, just run and exit
    if args.command == "_bg_build":
        _here = Path(__file__).parent.parent
        config_path = args.config or str(_here / "config.json")
        args.config = config_path
        cfg = load_config(config_path)
        _run_bg_build(args, cfg)
        return

    _here = Path(__file__).parent.parent
    config_path = args.config or str(_here / "config.json")
    args.config = config_path

    if not Path(config_path).exists():
        msg = f"Config not found: {config_path}"
        if args.output == "json":
            _emit_json({"status": "error", "command": args.command, "error": msg})
        else:
            print(msg, file=sys.stderr)
        sys.exit(1)

    try:
        # status needs no KB init
        if args.command == "status":
            result = _cmd_status(args)
            if args.output == "json":
                _emit_json(result)
            return

        # build needs no KB init either (spawns subprocess)
        cfg = load_config(config_path)
        if args.command == "build":
            result = _cmd_build(args, cfg)
            if args.output == "json":
                _emit_json(result)
            return

        # note / query need KB init
        kb = KnowledgeBase(config_path)
        if args.command == "note":
            result = _cmd_note(args, kb, cfg)
        elif args.command == "query":
            result = _cmd_query(args, kb, cfg)
        else:
            raise ValueError(f"Unknown command: {args.command}")

        if args.output == "json":
            _emit_json(result)

    except Exception as e:
        if args.output == "json":
            _emit_json({"status": "error", "command": args.command, "error": str(e)})
        else:
            print(f"[error] {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
