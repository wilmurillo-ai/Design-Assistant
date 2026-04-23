"""
chart_builder.py — Native PPT chart generation from structured data.

Creates editable charts (bar, pie, line, scatter, doughnut, radar, area)
using python-pptx's chart API. Charts are fully editable in MS PPT / WPS.

Supported chart types (python-pptx natively supports 29 of 73 XL_CHART_TYPE):
  - bar/column: clustered, stacked, stacked_100
  - line: plain, with markers, stacked
  - pie: standard, exploded
  - doughnut: standard, exploded
  - area: standard, stacked, stacked_100
  - scatter: xy, with lines, smooth
  - radar: standard, filled, with markers
  - bubble: standard, 3d effect

References:
  python-pptx charts: https://python-pptx.readthedocs.io/en/stable/user/charts.html
  XL_CHART_TYPE enum: https://python-pptx.readthedocs.io/en/stable/api/enum/XlChartType.html
  ChartData API: https://python-pptx.readthedocs.io/en/stable/api/chart-data.html
"""

from lxml import etree
from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import (
    XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION,
)
from pptx.oxml.ns import qn
from pptx.util import Pt, Emu

CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "bar_stacked": XL_CHART_TYPE.BAR_STACKED,
    "bar_stacked_100": XL_CHART_TYPE.BAR_STACKED_100,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column_stacked": XL_CHART_TYPE.COLUMN_STACKED,
    "column_stacked_100": XL_CHART_TYPE.COLUMN_STACKED_100,
    "line": XL_CHART_TYPE.LINE,
    "line_markers": XL_CHART_TYPE.LINE_MARKERS,
    "line_stacked": XL_CHART_TYPE.LINE_STACKED,
    "pie": XL_CHART_TYPE.PIE,
    "pie_exploded": XL_CHART_TYPE.PIE_EXPLODED,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "doughnut_exploded": XL_CHART_TYPE.DOUGHNUT_EXPLODED,
    "area": XL_CHART_TYPE.AREA,
    "area_stacked": XL_CHART_TYPE.AREA_STACKED,
    "area_stacked_100": XL_CHART_TYPE.AREA_STACKED_100,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
    "scatter_lines": XL_CHART_TYPE.XY_SCATTER_LINES,
    "scatter_smooth": XL_CHART_TYPE.XY_SCATTER_SMOOTH,
    "radar": XL_CHART_TYPE.RADAR,
    "radar_filled": XL_CHART_TYPE.RADAR_FILLED,
    "radar_markers": XL_CHART_TYPE.RADAR_MARKERS,
    "bubble": XL_CHART_TYPE.BUBBLE,
}

PALETTE_DARK = [
    "22D3EE", "3B82F6", "8B5CF6", "EC4899", "F59E0B",
    "10B981", "EF4444", "6366F1", "14B8A6", "F97316",
]

PALETTE_LIGHT = [
    "60A5FA", "34D399", "FBBF24", "F87171", "A78BFA",
    "2DD4BF", "FB923C", "818CF8", "4ADE80", "F472B6",
]


def _resolve_palette(chart_spec):
    """Pick color palette from chart spec or use default dark palette."""
    colors = chart_spec.get("colors")
    if colors and len(colors) >= 2:
        return [c.lstrip("#").upper() for c in colors]
    theme = chart_spec.get("theme", "dark")
    return PALETTE_DARK if theme == "dark" else PALETTE_LIGHT


def _apply_series_colors(chart, colors):
    """Apply palette colors to chart series (or pie/doughnut points)."""
    plot = chart.plots[0]
    chart_type = chart.chart_type

    is_pie_family = chart_type in (
        XL_CHART_TYPE.PIE, XL_CHART_TYPE.PIE_EXPLODED,
        XL_CHART_TYPE.DOUGHNUT, XL_CHART_TYPE.DOUGHNUT_EXPLODED,
    )

    if is_pie_family and len(chart.series) > 0:
        series = chart.series[0]
        for idx in range(len(series.values)):
            point = series.points[idx]
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = RGBColor.from_string(
                colors[idx % len(colors)]
            )
    else:
        for idx, series in enumerate(chart.series):
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = RGBColor.from_string(
                colors[idx % len(colors)]
            )


def _apply_data_labels(chart, chart_spec):
    """Configure data label visibility and position.

    IMPORTANT: Setting data_labels.position on pie/doughnut/stacked charts
    causes corruption in MS PowerPoint (GitHub Issue #789, #692).
    Only set position on safe chart types (bar/column clustered, line).
    """
    label_conf = chart_spec.get("dataLabels", {})
    if not label_conf.get("show", False):
        return

    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels

    font_size = label_conf.get("fontSize")
    if font_size:
        dl.font.size = Pt(font_size)

    font_color = label_conf.get("color")
    if font_color:
        dl.font.color.rgb = RGBColor.from_string(font_color.lstrip("#").upper()[:6])

    dl.number_format_is_linked = True

    chart_type_key = chart_spec.get("type", "column")
    position_safe_types = {
        "bar", "bar_stacked", "column", "column_stacked",
        "line", "line_markers", "line_stacked",
    }
    if chart_type_key in position_safe_types:
        pos_map = {
            "inside_end": XL_LABEL_POSITION.INSIDE_END,
            "outside_end": XL_LABEL_POSITION.OUTSIDE_END,
            "center": XL_LABEL_POSITION.CENTER,
        }
        pos_key = label_conf.get("position", "outside_end")
        pos_val = pos_map.get(pos_key)
        if pos_val is not None:
            try:
                dl.position = pos_val
            except ValueError:
                pass


def _apply_legend(chart, chart_spec):
    """Configure chart legend."""
    legend_conf = chart_spec.get("legend", {})
    if legend_conf.get("show", True) is False:
        chart.has_legend = False
        return

    chart.has_legend = True
    pos_map = {
        "right": XL_LEGEND_POSITION.RIGHT,
        "bottom": XL_LEGEND_POSITION.BOTTOM,
        "left": XL_LEGEND_POSITION.LEFT,
        "top": XL_LEGEND_POSITION.TOP,
    }
    chart.legend.position = pos_map.get(
        legend_conf.get("position", "bottom"), XL_LEGEND_POSITION.BOTTOM
    )
    chart.legend.include_in_layout = False


def add_chart(slide, chart_spec, x_emu, y_emu, w_emu, h_emu):
    """Add an editable chart to a slide from a chart specification dict.

    chart_spec schema:
      {
        "type": "column" | "bar" | "line" | "pie" | "doughnut" | "scatter" | ...,
        "title": "Chart Title" (optional),
        "categories": ["A", "B", "C"],
        "series": [
          {"name": "S1", "values": [1, 2, 3]},
          {"name": "S2", "values": [4, 5, 6]}
        ],
        "colors": ["#22D3EE", "#3B82F6", ...] (optional),
        "theme": "dark" | "light" (optional, default "dark"),
        "dataLabels": {"show": true, "fontSize": 10, "position": "outside_end"},
        "legend": {"show": true, "position": "bottom"},
      }

    For scatter/bubble charts, series format differs:
      "series": [
        {"name": "S1", "points": [[x, y], [x, y], ...]}
      ]
    For bubble:
      "series": [
        {"name": "S1", "points": [[x, y, size], ...]}
      ]

    Returns the graphic_frame shape object.
    """
    type_key = chart_spec.get("type", "column")
    xl_type = CHART_TYPE_MAP.get(type_key, XL_CHART_TYPE.COLUMN_CLUSTERED)
    colors = _resolve_palette(chart_spec)

    is_xy = type_key.startswith("scatter")
    is_bubble = type_key.startswith("bubble")

    if is_bubble:
        chart_data = BubbleChartData()
        for s in chart_spec.get("series", []):
            series = chart_data.add_series(s.get("name", ""))
            for pt in s.get("points", []):
                if len(pt) >= 3:
                    series.add_data_point(pt[0], pt[1], pt[2])
    elif is_xy:
        chart_data = XyChartData()
        for s in chart_spec.get("series", []):
            series = chart_data.add_series(s.get("name", ""))
            for pt in s.get("points", []):
                series.add_data_point(pt[0], pt[1])
    else:
        chart_data = CategoryChartData()
        chart_data.categories = chart_spec.get("categories", [])
        for s in chart_spec.get("series", []):
            chart_data.add_series(s.get("name", ""), s.get("values", []))

    graphic_frame = slide.shapes.add_chart(
        xl_type, x_emu, y_emu, w_emu, h_emu, chart_data
    )
    chart = graphic_frame.chart

    _apply_series_colors(chart, colors)
    _apply_data_labels(chart, chart_spec)
    _apply_legend(chart, chart_spec)
    _fix_numfmt(chart)
    _fix_radar_smooth(chart, type_key)

    title_text = chart_spec.get("title")
    if title_text:
        chart.has_title = True
        chart.chart_title.text_frame.text = title_text
    else:
        chart.has_title = False

    return graphic_frame


def _fix_radar_smooth(chart, type_key):
    """Remove invalid <c:smooth> elements from radar charts.

    python-pptx may emit <c:smooth val="0"/> on radar series, but the OOXML
    schema for radarChart/ser does not allow the smooth element.
    """
    if not type_key.startswith("radar"):
        return
    chart_xml = chart._chartSpace
    ns_c = "http://schemas.openxmlformats.org/drawingml/2006/chart"
    for smooth in chart_xml.iter(qn("c:smooth")):
        parent = smooth.getparent()
        if parent is not None:
            parent.remove(smooth)


def _fix_numfmt(chart):
    """Ensure all numFmt elements have required formatCode attribute.

    python-pptx may emit <c:numFmt sourceLinked="1"/> without formatCode,
    which violates OOXML schema and triggers MS PPT repair prompt.
    Reference: ECMA-376 §21.2.2.142 (numFmt requires formatCode).
    """
    chart_xml = chart._chartSpace
    ns_c = "http://schemas.openxmlformats.org/drawingml/2006/chart"
    for nf in chart_xml.iter(qn("c:numFmt")):
        if nf.get("formatCode") is None:
            nf.set("formatCode", "General")
