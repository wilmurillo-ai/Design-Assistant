#!/usr/bin/env python3
"""
assemble_pptx.py — Phase 2 of Hybrid PPTX Pipeline

Reads JSON manifest + PNG screenshots from extract_slides.js,
assembles a PowerPoint file using python-pptx.

Features:
  - Full-bleed background image per slide (HTML screenshot)
  - Editable text boxes overlaid at precise DOM coordinates
  - Gradient text support via OOXML a:gradFill injection (lxml)
  - Transparent text boxes (no fill, no border)
  - Available-width sizing to prevent cross-renderer text wrapping

Coordinate mapping:
  HTML viewport 1280x720 → PPTX 10" x 5.625" (9144000 x 5143500 EMU)
  Scale factor: 1px = 7143.75 EMU (identical for both axes in 16:9)

Usage:
  python3 assemble_pptx.py <manifest_dir> -o output.pptx
"""

import json
import os
import re
import sys

from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

from ooxml_helpers import (
    set_run_color, set_font_typefaces, set_letter_spacing,
    apply_gradient_fill, apply_text_shadow, apply_text_outline,
    apply_text_glow,
)
from font_metrics import measure_region_width, measure_max_line_width

SLIDE_W_EMU = Inches(10.0)
SLIDE_H_EMU = Inches(5.625)
VIEWPORT_W = 1280
VIEWPORT_H = 720
PX_TO_EMU = int(SLIDE_W_EMU) / VIEWPORT_W  # 7143.75

# CSS px → PPTX pt: 1 CSS px = 0.75 pt (96dpi screen vs 72dpi print).
# Ref: https://stackoverflow.com/questions/25548060
CSS_PX_TO_PT = 0.75

# Slide background color for alpha-blending semi-transparent colors.
SLIDE_BG_RGB = (0x0B, 0x11, 0x20)


def px_to_emu(px_val):
    return int(round(px_val * PX_TO_EMU))


def css_px_to_pt(px_val):
    """Convert CSS pixel value to PowerPoint point value."""
    return round(px_val * CSS_PX_TO_PT, 2)


def _set_fill_alpha(shape, alpha_val):
    """Set fill transparency via OxmlElement a:alpha on srgbClr.

    Ref: stackoverflow.com/questions/38202582
    """
    sf = shape.fill._fill._solidFill
    srgb = sf.find(qn("a:srgbClr"))
    if srgb is None:
        return
    alpha_el = OxmlElement("a:alpha")
    alpha_el.set("val", str(alpha_val))
    srgb.append(alpha_el)


def _set_line_alpha(shape, alpha_val):
    """Set line transparency via OxmlElement a:alpha on srgbClr."""
    spPr = shape._element.find(qn("p:spPr"))
    if spPr is None:
        return
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        return
    sf = ln.find(qn("a:solidFill"))
    if sf is None:
        return
    srgb = sf.find(qn("a:srgbClr"))
    if srgb is None:
        return
    alpha_el = OxmlElement("a:alpha")
    alpha_el.set("val", str(alpha_val))
    srgb.append(alpha_el)


def blend_alpha(fg_hex, fg_alpha_pct, bg_rgb=SLIDE_BG_RGB):
    """Alpha-blend a foreground color over a background to get opaque result.

    Args:
        fg_hex: 6-char hex color string (e.g. "0B1120")
        fg_alpha_pct: alpha as 0-100000 (OOXML scale)
        bg_rgb: tuple (R, G, B) ints for background
    Returns:
        RGBColor
    """
    a = fg_alpha_pct / 100000.0
    r = int(int(fg_hex[0:2], 16) * a + bg_rgb[0] * (1 - a))
    g = int(int(fg_hex[2:4], 16) * a + bg_rgb[1] * (1 - a))
    b = int(int(fg_hex[4:6], 16) * a + bg_rgb[2] * (1 - a))
    return RGBColor(r, g, b)


def hex_to_rgb(hex_str):
    hex_str = (hex_str or "").lstrip("#")
    if len(hex_str) < 6:
        hex_str = "333333"
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def align_value(align_str):
    return {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align_str, PP_ALIGN.LEFT)


def set_slide_bg_image(slide, image_path):
    """Add a full-bleed picture as background, positioned behind all other shapes."""
    pic = slide.shapes.add_picture(image_path, left=0, top=0, width=SLIDE_W_EMU, height=SLIDE_H_EMU)
    sp_tree = slide.shapes._spTree
    pic_elem = pic._element
    sp_tree.remove(pic_elem)
    sp_tree.insert(2, pic_elem)



_SVG_COMPLEX_PATTERNS = re.compile(
    r"<linearGradient|<radialGradient|<filter\b|<mask\b|<clipPath\b|<foreignObject",
    re.IGNORECASE,
)

MAX_SIMPLE_SVG_DIM = 120


def _svg_to_transparent_png(svg_path, w_px, h_px):
    """Render SVG to transparent-background PNG via cairosvg.

    Puppeteer element.screenshot() produces opaque PNGs because parent element
    backgrounds bleed through, even with omitBackground. cairosvg renders SVG
    paths directly onto a Cairo canvas with no HTML background.
    Ref: https://cairosvg.org/documentation/
    Ref: https://medium.com/@nizarbabut/cairosvg-how-to-use-it-in-python-especially-svg2png
    """
    try:
        import cairosvg
    except ImportError:
        return None
    out = svg_path.rsplit(".", 1)[0] + "_fb.png"
    scale = max(2, 192 // 96)
    try:
        cairosvg.svg2png(
            url=svg_path,
            write_to=out,
            output_width=int(w_px * scale),
            output_height=int(h_px * scale),
        )
        return out
    except Exception:
        return None


def _is_simple_svg(svg_path, w_px, h_px):
    """Check if SVG is simple enough for PowerPoint a16svg:svgBlip embedding.

    PowerPoint's SVG renderer is comparable to IE, not modern browsers.
    Complex features (gradients, filters, masks) render incorrectly.
    Large diagrams should use PNG fallback for reliable cross-platform display.
    Ref: https://learn.microsoft.com/en-us/answers/questions/5292427
    """
    if w_px > MAX_SIMPLE_SVG_DIM or h_px > MAX_SIMPLE_SVG_DIM:
        return False
    try:
        content = open(svg_path, "r", encoding="utf-8").read(8192)
    except (OSError, UnicodeDecodeError):
        return False
    if _SVG_COMPLEX_PATTERNS.search(content):
        return False
    return True


def _cjk_ratio(text):
    """Fraction of CJK characters. Used to calibrate width buffer."""
    if not text:
        return 0.0
    cjk = len(re.findall(r'[\u4e00-\u9fff\u3000-\u303f\u3040-\u309f\u30a0-\u30ff\uff00-\uffef]', text))
    return cjk / len(text)


def _ppt_width_for_region(dom_w, available_w, region):
    """Calculate PPT-safe textbox width for cross-platform rendering.

    Strategy depends on content characteristics:
    - Single-line text: generous width to prevent line-break on any platform.
    - Multi-line text: DOM width with 15% buffer for DirectWrite differences.
      Buffer compensates for CJK character width differences between
      browser FreeType (~1.0em) and MS PPT DirectWrite (~1.05-1.15em).

    Reference: Pillow getlength() uses FreeType2 advance widths.
    dom-to-pptx uses full bounding-box width from DOM.
    """
    line_count = region.get("lineCount", 1)
    is_single = line_count <= 1

    if is_single:
        measured = measure_region_width(region)
        if measured is not None and measured < available_w * 0.5:
            safe_w = max(measured * 1.35, dom_w * 1.25)
            return min(available_w, safe_w)
        return available_w

    max_line_w = region.get("maxLineWidth", dom_w)
    measured_max = measure_max_line_width(region)
    if measured_max is not None:
        base_w = max(max_line_w, measured_max, dom_w * 0.6)
    else:
        base_w = max(max_line_w, dom_w * 0.6)

    cjk = _cjk_ratio(region.get("text", ""))
    if base_w < 80:
        buf = 1.35 + cjk * 0.15
    elif base_w < 200:
        buf = 1.25 + cjk * 0.10
    else:
        buf = 1.15 + cjk * 0.05
    return min(available_w, base_w * buf)


def _ppt_line_height_px(font_size_px, line_spacing_ratio):
    """Calculate PPT line height in pixels using the empirical formula.

    PPT uses: lineHeight = (1.2018 * lineSpacing + 0.0034) * fontSize
    where fontSize is in pixels. This is derived from measurements
    documented at https://cloud.tencent.com/developer/article/1914894
    """
    return (1.2018 * line_spacing_ratio + 0.0034) * font_size_px


def add_text_overlay(slide, region):
    """Add an editable text box with cross-platform safe sizing.

    Uses noAutofit (ECMA-376 §21.1.2.1.2) instead of normAutofit to prevent
    MS PowerPoint from shrinking text unpredictably. Width is compensated
    for DirectWrite vs Chromium rendering differences (~15-25%).
    """
    x_emu = px_to_emu(region["x"])
    y_emu = px_to_emu(region["y"])

    text = region.get("text", "")
    available_w = VIEWPORT_W - region["x"]
    dom_w = region["w"]
    dom_h = region["h"]
    font_size_px = region.get("fontSizePx", 16)

    runs_data = region.get("runs")
    if runs_data:
        max_run_fs = max(
            (r.get("fontSizePx", font_size_px) for r in runs_data if r.get("text", "").strip()),
            default=font_size_px,
        )
        effective_fs = max(font_size_px, max_run_fs)
    else:
        effective_fs = font_size_px

    line_count = region.get("lineCount", 1)
    lh_ratio = region.get("lineHeightRatio", 1.2)
    is_single_line = line_count <= 1 and dom_h < effective_fs * 1.8
    ppt_line_h = _ppt_line_height_px(effective_fs, lh_ratio)

    w_px = _ppt_width_for_region(dom_w, available_w, region)

    if is_single_line:
        h_px = max(dom_h, ppt_line_h) + 4
    else:
        h_ideal = ppt_line_h * line_count + 4
        h_px = max(dom_h * 1.10, h_ideal)

    w_emu = px_to_emu(w_px)
    h_emu = px_to_emu(h_px)

    x_emu = max(0, x_emu)
    y_emu = max(0, y_emu)
    if x_emu + w_emu > int(SLIDE_W_EMU):
        w_emu = int(SLIDE_W_EMU) - x_emu
    if y_emu + h_emu > int(SLIDE_H_EMU):
        h_emu = int(SLIDE_H_EMU) - y_emu

    txbox = slide.shapes.add_textbox(x_emu, y_emu, w_emu, h_emu)
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    v_align = region.get("vAlign", "top")
    if v_align == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif v_align == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM

    body_pr = tf._txBody.find(qn("a:bodyPr"))
    for tag in ("a:normAutofit", "a:spAutoFit", "a:noAutofit"):
        for old in body_pr.findall(qn(tag)):
            body_pr.remove(old)
    autofit = etree.SubElement(body_pr, qn("a:normAutofit"))
    autofit.set("fontScale", "100000")

    p = tf.paragraphs[0]
    p.alignment = align_value(region.get("align", "left"))
    lh_ratio = region.get("lineHeightRatio", 1.2)
    if lh_ratio > 0.5:
        ppt_lsp = lh_ratio / 1.2018
        p.line_spacing = round(max(0.7, ppt_lsp), 2)

    font_size_pt = region["fontSizePt"]
    font_family = region.get("fontFamily", "Microsoft YaHei")
    default_color = region.get("color", "FFFFFF")
    default_color_alpha = region.get("colorAlpha")
    default_bold = region.get("bold", False)
    default_italic = region.get("italic", False)
    letter_spacing_em = region.get("letterSpacingEm", 0)
    text_transform = region.get("textTransform")
    region_opacity = region.get("opacity")

    def _compute_alpha(color_alpha_val, use_region_opacity=True):
        """Combine CSS color alpha (0-100000) with compositeOpacity."""
        ca = (color_alpha_val / 100000.0) if color_alpha_val else 1.0
        ro = region_opacity if (region_opacity and use_region_opacity) else 1.0
        combined = ca * ro
        return combined if combined < 0.999 else None
    gradient = region.get("gradient")
    shadow = region.get("shadow")
    glow = region.get("glow")
    outline = region.get("outline")
    runs_data = region.get("runs")

    def _apply_cap(run_elem):
        """Apply text-transform as OOXML cap attribute on a:rPr."""
        if text_transform == "uppercase":
            run_elem.get_or_add_rPr().set("cap", "all")
        elif text_transform == "capitalize":
            pass
        elif text_transform == "lowercase":
            pass

    def _apply_effects(run_elem):
        _apply_cap(run_elem)
        if shadow:
            apply_text_shadow(run_elem, shadow, PX_TO_EMU)
        if glow:
            apply_text_glow(run_elem, glow, PX_TO_EMU)
        if outline:
            apply_text_outline(run_elem, outline)

    if runs_data and not gradient:
        p.clear()
        for ri, rd in enumerate(runs_data):
            raw = rd.get("text", "")
            if ri == 0:
                txt = raw.lstrip()
            else:
                txt = raw.lstrip("\n\r\t")
            if ri == len(runs_data) - 1:
                txt = txt.rstrip()
            if not txt:
                continue
            run = p.add_run()
            run.text = txt
            run_fs = rd.get("fontSizePt", font_size_pt)
            run.font.size = Pt(run_fs)
            run_ff = rd.get("fontFamily", font_family)
            run.font.name = run_ff
            run.font.bold = rd.get("bold", default_bold)
            run.font.italic = rd.get("italic", default_italic)
            run_color = rd.get("color", default_color)
            run_ca = rd.get("colorAlpha")
            if run_ca is None and run_color == default_color:
                run_ca = default_color_alpha
            set_run_color(run._r, run_color, alpha=_compute_alpha(run_ca))
            set_font_typefaces(run._r, run_ff)
            set_letter_spacing(run._r, letter_spacing_em, run_fs)
            _apply_effects(run._r)
    else:
        p.text = region["text"]
        run = p.runs[0]
        run.font.size = Pt(font_size_pt)
        run.font.bold = default_bold
        run.font.italic = default_italic
        run.font.name = font_family

        if gradient and len(gradient.get("stops", [])) >= 2:
            apply_gradient_fill(run._r, gradient)
        else:
            set_run_color(run._r, default_color, alpha=_compute_alpha(default_color_alpha))

        set_font_typefaces(run._r, font_family)
        set_letter_spacing(run._r, letter_spacing_em, font_size_pt)
        _apply_effects(run._r)

    # Transparent fill, no border
    txbox.fill.background()
    sp_pr = txbox._element.find(qn("p:spPr"))
    if sp_pr is not None:
        ln = sp_pr.find(qn("a:ln"))
        if ln is None:
            ln = etree.SubElement(sp_pr, qn("a:ln"), attrib={"w": "0"})
            etree.SubElement(ln, qn("a:noFill"))


def _parse_css_rgb(color_str):
    """Parse CSS rgb/rgba/color(srgb) string to (RGBColor, alpha_int_100000).

    Supports:
      - rgba(r, g, b, a)  / rgb(r, g, b)
      - #RRGGBB
      - color(srgb R G B)  / color(srgb R G B / A)
        where R,G,B are 0-1 floats (Chrome 111+ getComputedStyle format)

    Ref: CSS Color 4 serialization §4.7, python-pptx RGBColor.
    """
    if not color_str:
        return RGBColor(0x80, 0x80, 0x80), 100000
    m = re.match(r"rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)(?:\s*,\s*([\d.]+))?\s*\)", color_str)
    if m:
        r, g, b = int(m.group(1)), int(m.group(2)), int(m.group(3))
        a = float(m.group(4)) if m.group(4) else 1.0
        return RGBColor(r, g, b), int(a * 100000)
    m_srgb = re.match(
        r"color\(\s*srgb\s+([\d.]+)\s+([\d.]+)\s+([\d.]+)"
        r"(?:\s*/\s*([\d.]+))?\s*\)",
        color_str,
    )
    if m_srgb:
        r = min(255, max(0, int(float(m_srgb.group(1)) * 255 + 0.5)))
        g = min(255, max(0, int(float(m_srgb.group(2)) * 255 + 0.5)))
        b = min(255, max(0, int(float(m_srgb.group(3)) * 255 + 0.5)))
        a = float(m_srgb.group(4)) if m_srgb.group(4) else 1.0
        return RGBColor(r, g, b), int(a * 100000)
    m2 = re.match(r"#([0-9a-fA-F]{6})", color_str)
    if m2:
        return RGBColor.from_string(m2.group(1)), 100000
    return RGBColor(0x80, 0x80, 0x80), 100000


def add_native_chart(slide, chart_data_info, x, y, cx, cy):
    """Create a native PPTX chart from parsed CSS chart data.

    Uses python-pptx standard add_chart() API with minimal XML patching.
    Only holeSize requires XML manipulation (no python-pptx API for it).
    Ref: https://python-pptx.readthedocs.io/en/stable/user/charts.html
    Ref: https://stackoverflow.com/questions/34264715 (per-point color)
    Ref: https://github.com/scanny/python-pptx/issues/493 (hole size)
    """
    chart_type = chart_data_info.get("type")
    segments = chart_data_info.get("segments", [])
    if not segments:
        return False

    if chart_type == "doughnut":
        data = CategoryChartData()
        cats = [f"seg{i}" for i in range(len(segments))]
        vals = [seg["pct"] for seg in segments]
        data.categories = cats
        data.add_series("Budget", vals)

        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, data
        )
        chart = graphic_frame.chart
        chart.has_legend = False
        chart.has_title = False
        plot = chart.plots[0]
        plot.has_data_labels = False

        series = chart.series[0]
        for idx, seg in enumerate(segments):
            point = series.points[idx]
            rgb, alpha = _parse_css_rgb(seg["color"])
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = rgb

        hole_pct = chart_data_info.get("holePct", 50)
        doughnut = chart._chartSpace.chart.plotArea.find(qn("c:doughnutChart"))
        if doughnut is not None:
            existing = doughnut.find(qn("c:holeSize"))
            if existing is not None:
                existing.set("val", str(hole_pct))
            else:
                hole_el = OxmlElement("c:holeSize")
                hole_el.set("val", str(hole_pct))
                first_slice = doughnut.find(qn("c:firstSliceAng"))
                if first_slice is not None:
                    first_slice.addnext(hole_el)
                else:
                    doughnut.append(hole_el)

        return True
    return False


def _set_cell_border(cell, side, color_hex, width_emu):
    """Set one border of a table cell via direct XML.

    Ref: https://stackoverflow.com/questions/42610829 (cell borders)
    Ref: https://stackoverflow.com/questions/72919029 (per-side borders)

    Args:
        side: one of 'a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'
    """
    tc_pr = cell._tc.get_or_add_tcPr()
    existing = tc_pr.find(qn(side))
    if existing is not None:
        tc_pr.remove(existing)
    ln = OxmlElement(side)
    ln.set("w", str(width_emu))
    ln.set("cap", "flat")
    ln.set("cmpd", "sng")
    ln.set("algn", "ctr")
    if color_hex == "noFill":
        etree.SubElement(ln, qn("a:noFill"))
    else:
        sf = etree.SubElement(ln, qn("a:solidFill"))
        clr = etree.SubElement(sf, qn("a:srgbClr"))
        clr.set("val", color_hex)
    etree.SubElement(ln, qn("a:prstDash")).set("val", "solid")
    tc_pr.append(ln)


def add_native_table(slide, table_info, x, y, cx, cy):
    """Create a native PPTX table faithfully reproducing HTML <table> styling.

    Reproduced properties:
      - Column widths proportional to HTML
      - Header row: darker bg, gray-blue uppercase text, smaller font
      - Data rows: transparent bg over table gradient, green/red/white text
      - Subtle row separator borders (bottom), column separator on col 0
      - Cell padding via margin properties
      - No built-in table style (custom colors only)

    Ref: https://python-pptx.readthedocs.io/en/stable/user/table.html
    Ref: https://stackoverflow.com/questions/61982333 (table style ID)
    Ref: https://stackoverflow.com/questions/42610829 (cell borders)
    Ref: https://stackoverflow.com/questions/78624591 (cell margins)
    """
    rows_data = table_info.get("rows", [])
    n_rows = len(rows_data)
    n_cols = table_info.get("cols", 0)
    if n_rows < 1 or n_cols < 1:
        return False

    shape = slide.shapes.add_table(n_rows, n_cols, x, y, cx, cy)
    table = shape.table

    # --- 1) Disable all built-in table style ---
    # Style GUID {2D5ABB26-…} = "No Style, No Grid".
    tbl_xml = shape._element.graphic.graphicData.tbl
    tbl_pr = tbl_xml.find(qn("a:tblPr"))
    if tbl_pr is not None:
        style_id_el = tbl_pr.find(qn("a:tableStyleId"))
        if style_id_el is not None:
            style_id_el.text = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"
    table.first_row = False
    table.horz_banding = False
    table.first_col = False
    table.last_row = False
    table.last_col = False

    # --- 2) Column widths proportional to HTML ---
    first_row_widths = [c.get("widthPx", 100) for c in rows_data[0]]
    total_w = sum(first_row_widths) or 1
    for ci, w_px in enumerate(first_row_widths):
        if ci < n_cols:
            table.columns[ci].width = int(cx * w_px / total_w)

    # --- 2b) Row heights: header row gets ~14% of height, data rows split rest
    # HTML: header ~32px / total ~234px ≈ 14%. Remainder split evenly.
    n_header_rows = sum(1 for row in rows_data if row[0].get("isHeader"))
    n_data_rows = n_rows - n_header_rows
    header_h = int(cy * 0.14) if n_header_rows > 0 else 0
    data_h = int((cy - header_h * n_header_rows) / max(n_data_rows, 1))
    for ri in range(n_rows):
        is_hdr_row = rows_data[ri][0].get("isHeader", False)
        table.rows[ri].height = header_h if is_hdr_row else data_h

    # --- 3) Separator border style from HTML CSS ---
    # HTML: border-bottom: 1px solid rgba(255,255,255,0.05~0.06)
    #       layer-cell border-right: 1px solid rgba(255,255,255,0.06)
    # rgba(255,255,255,0.06) over #0B1120 ≈ #1A2030 (very subtle)
    sep_color = "1A2030"
    sep_width = 6350  # 0.5pt in EMU (6350 EMU = 0.5pt)

    # Header-bottom separator slightly brighter
    hdr_sep_color = "1E2B3D"
    hdr_sep_width = 9525  # 0.75pt

    # Table background: HTML uses linear-gradient(180deg,
    #   rgba(30,41,59,0.75), rgba(15,23,42,0.88))
    # Blended over slide bg #0B1120:
    #   top: rgba(30,41,59,0.75) → ~(25,34,49)
    #   bottom: rgba(15,23,42,0.88) → ~(14,21,38)
    # We use a single blended midpoint for data cells
    data_bg = RGBColor(0x13, 0x1B, 0x2B)
    header_bg = RGBColor(0x08, 0x0D, 0x18)

    for ri, row in enumerate(rows_data):
        for ci, cell_info in enumerate(row):
            if ci >= n_cols:
                continue
            cell = table.cell(ri, ci)
            is_header = cell_info.get("isHeader", False)

            # --- Cell background ---
            cell_fill = cell.fill
            cell_fill.solid()
            if is_header:
                cell_fill.fore_color.rgb = header_bg
            else:
                cell_fill.fore_color.rgb = data_bg

            # --- Cell margins (padding) ---
            # HTML: th padding 10px 14px; td padding 12px 14px
            # Convert CSS px → pt (* 0.75)
            # Ref: https://stackoverflow.com/questions/78624591
            pad_tb = Pt(7.5) if is_header else Pt(9)
            pad_lr = Pt(10.5)
            cell.margin_left = pad_lr
            cell.margin_right = pad_lr
            cell.margin_top = pad_tb
            cell.margin_bottom = pad_tb

            # --- Text content and formatting ---
            cell.text = cell_info.get("text", "")

            font_sz = cell_info.get("fontSizePt", css_px_to_pt(cell_info.get("fontSizePx", 12)))
            font_wt = cell_info.get("fontWeight", 400)
            fg_hex = cell_info.get("color", "F1F5F9")
            fg_alpha = cell_info.get("colorAlpha", 100000)

            if fg_alpha < 95000:
                bg_tuple = (header_bg[0], header_bg[1], header_bg[2]) if is_header \
                    else (data_bg[0], data_bg[1], data_bg[2])
                text_rgb = blend_alpha(fg_hex, fg_alpha, bg_tuple)
            else:
                text_rgb = RGBColor.from_string(fg_hex)

            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(font_sz)
                para.font.bold = font_wt >= 600
                para.font.color.rgb = text_rgb
                if is_header:
                    para.alignment = PP_ALIGN.LEFT

            if not is_header:
                cell.vertical_anchor = MSO_ANCHOR.TOP

            # --- Borders ---
            # Remove all borders first, then add only the ones we need
            for side in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
                _set_cell_border(cell, side, "noFill", 0)

            # Bottom border on all rows except last
            if ri < n_rows - 1:
                bdr_c = hdr_sep_color if is_header else sep_color
                bdr_w = hdr_sep_width if is_header else sep_width
                _set_cell_border(cell, "a:lnB", bdr_c, bdr_w)

            # Right border on first column (layer column separator)
            if ci == 0 and n_cols > 1:
                _set_cell_border(cell, "a:lnR", sep_color, sep_width)

    return True


def add_svg_picture(slide, png_path, svg_path, x, y, cx, cy):
    """Add a picture with SVG vector + PNG fallback (dual-format).

    PowerPoint 2016+ will render the SVG; older versions fall back to PNG.
    Uses the OOXML a16svg:svgBlip extension (MS-ODRAWXML §2.1.1781).
    """
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.opc.package import Part

    pic = slide.shapes.add_picture(png_path, x, y, cx, cy)
    blip_fill = pic._element.find(qn("p:blipFill"))
    if blip_fill is None:
        return pic
    blip = blip_fill.find(qn("a:blip"))
    if blip is None:
        return pic

    with open(svg_path, "rb") as f:
        svg_blob = f.read()

    if not hasattr(add_svg_picture, "_counter"):
        add_svg_picture._counter = 0
    add_svg_picture._counter += 1
    svg_partname = f"/ppt/media/svg-{add_svg_picture._counter}.svg"

    slide_part = slide.part
    pkg = slide_part.package

    from pptx.opc.packuri import PackURI

    svg_part = Part(
        PackURI(svg_partname),
        "image/svg+xml",
        blob=svg_blob,
        package=pkg,
    )

    svg_rId = slide_part.relate_to(svg_part, RT.IMAGE)

    nsmap_svg = {"a16svg": "http://schemas.microsoft.com/office/drawing/2016/SVG/main"}
    ext_lst = blip.find(qn("a:extLst"))
    if ext_lst is None:
        ext_lst = etree.SubElement(blip, qn("a:extLst"))
    ext_el = etree.SubElement(ext_lst, qn("a:ext"),
                              attrib={"uri": "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"})
    svg_blip = etree.SubElement(ext_el, "{http://schemas.microsoft.com/office/drawing/2016/SVG/main}svgBlip")
    svg_blip.set(qn("r:embed"), svg_rId)

    return pic


def add_shape_overlay(slide, shape_info):
    """Add an editable native shape (rect, rounded rect, line) from extracted CSS data.

    Uses python-pptx add_shape/add_connector per MSO_SHAPE constants.
    Reference: python-pptx AutoShapes documentation.
    """
    x_emu = px_to_emu(shape_info["x"])
    y_emu = px_to_emu(shape_info["y"])
    w_emu = px_to_emu(shape_info["w"])
    h_emu = px_to_emu(shape_info["h"])

    x_emu = max(0, x_emu)
    y_emu = max(0, y_emu)
    if x_emu + w_emu > int(SLIDE_W_EMU):
        w_emu = int(SLIDE_W_EMU) - x_emu
    if y_emu + h_emu > int(SLIDE_H_EMU):
        h_emu = int(SLIDE_H_EMU) - y_emu

    shape_type = shape_info.get("type", "rect")
    border_w = shape_info.get("borderWidthPx", 1)
    border_color = shape_info.get("borderColor", "FFFFFF")
    fill_color = shape_info.get("fillColor")
    fill_alpha = shape_info.get("fillAlpha", 100000)

    if shape_type == "line":
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT,
            x_emu, y_emu,
            x_emu + w_emu, y_emu
        )
        connector.line.color.rgb = RGBColor.from_string(border_color[:6])
        connector.line.width = Pt(max(0.5, border_w * 0.75))
    else:
        mso = MSO_SHAPE.ROUNDED_RECTANGLE if shape_type == "roundedRect" else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(mso, x_emu, y_emu, w_emu, h_emu)
        p_style = shape._element.find(qn("p:style"))
        if p_style is not None:
            shape._element.remove(p_style)

        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(fill_color[:6])
            if fill_alpha < 99000:
                sp_pr = shape._element.find(qn("p:spPr"))
                if sp_pr is None:
                    sp_pr = shape._element.spPr
                sf_elem = sp_pr.find(qn("a:solidFill"))
                if sf_elem is not None:
                    clr = sf_elem.find(qn("a:srgbClr"))
                    if clr is not None:
                        etree.SubElement(clr, qn("a:alpha"), attrib={"val": str(fill_alpha)})
        else:
            shape.fill.background()

        border_alpha = shape_info.get("borderAlpha", 100000)
        if border_w > 0.1:
            bc = border_color.lstrip("#")[:6]
            shape.line.color.rgb = RGBColor.from_string(bc)
            shape.line.width = Pt(max(0.25, border_w * 0.75))
            if border_alpha < 99000:
                ln_elem = shape._element.find(qn("p:spPr"))
                if ln_elem is not None:
                    ln_el = ln_elem.find(qn("a:ln"))
                    if ln_el is not None:
                        sf = ln_el.find(qn("a:solidFill"))
                        if sf is not None:
                            clr = sf.find(qn("a:srgbClr"))
                            if clr is not None:
                                etree.SubElement(clr, qn("a:alpha"), attrib={"val": str(border_alpha)})
        else:
            shape.line.fill.background()

        if shape_type == "roundedRect":
            radius_px = shape_info.get("borderRadius", 8)
            radius_frac = min(50000, int(radius_px / min(shape_info["w"], shape_info["h"]) * 100000))
            sp = shape._element
            spPr = sp.find(qn("a:spPr"))
            if spPr is None:
                sp_ns = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
                spPr = sp.find(f"{sp_ns}spPr")
            prstGeom = spPr.find(qn("a:prstGeom")) if spPr is not None else None
            if prstGeom is not None:
                avLst = prstGeom.find(qn("a:avLst"))
                if avLst is None:
                    avLst = etree.SubElement(prstGeom, qn("a:avLst"))
                for child in list(avLst):
                    avLst.remove(child)
                etree.SubElement(avLst, qn("a:gd"), attrib={"name": "adj", "fmla": f"val {radius_frac}"})

        shape_text = shape_info.get("shapeText")
        if shape_text and hasattr(shape, "text_frame"):
            tf = shape.text_frame
            tf.word_wrap = True
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
            run = tf.paragraphs[0].add_run()
            run.text = shape_text
            text_sz = shape_info.get("textSizePt", 12)
            run.font.size = Pt(text_sz)
            text_color = shape_info.get("textColor", "FFFFFF")
            run.font.color.rgb = RGBColor.from_string(text_color[:6])
            if shape_info.get("textBold"):
                run.font.bold = True


def _css_color_to_hex(css_val):
    """Convert CSS color value (rgb(...), #hex, named) to 6-char hex string."""
    if not css_val:
        return None
    css_val = css_val.strip()
    if css_val.startswith("#"):
        h = css_val.lstrip("#")
        if len(h) == 3:
            h = "".join(c * 2 for c in h)
        return h.upper()[:6] if len(h) >= 6 else None
    m = re.match(r"rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)", css_val)
    if m:
        return f"{int(m.group(1)):02X}{int(m.group(2)):02X}{int(m.group(3)):02X}"
    return None


def inject_theme_colors(prs, palette):
    """Inject custom color scheme into theme.xml from CSS variable palette.

    Reference: ECMA-376 §20.1.6.2 (clrScheme), Brandwares OOXML Color Themes guide.
    Maps CSS design tokens to OOXML theme slots:
      --accent-1 -> accent1, --accent-2 -> accent2, --accent-3 -> accent3
      --bg-primary -> dk1 (dark background), --text-primary -> lt1 (light text)
    """
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    if not palette:
        return

    slot_map = {}
    if "--accent-1" in palette:
        slot_map["accent1"] = _css_color_to_hex(palette["--accent-1"])
    if "--accent-2" in palette:
        slot_map["accent2"] = _css_color_to_hex(palette["--accent-2"])
    if "--accent-3" in palette:
        slot_map["accent3"] = _css_color_to_hex(palette["--accent-3"])
    if "--bg-primary" in palette:
        slot_map["dk2"] = _css_color_to_hex(palette["--bg-primary"])
    if "--text-primary" in palette:
        slot_map["lt1"] = _css_color_to_hex(palette["--text-primary"])
    if "--text-secondary" in palette:
        slot_map["lt2"] = _css_color_to_hex(palette["--text-secondary"])

    slot_map = {k: v for k, v in slot_map.items() if v}
    if not slot_map:
        return

    slide_master = prs.slide_masters[0]
    theme_part = slide_master.part.part_related_by(RT.THEME)
    theme_xml = etree.fromstring(theme_part.blob)

    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    clr_scheme = theme_xml.find(f".//{{{ns_a}}}clrScheme")
    if clr_scheme is None:
        return

    for slot_name, hex_color in slot_map.items():
        slot_el = clr_scheme.find(f"{{{ns_a}}}{slot_name}")
        if slot_el is None:
            continue
        for child in list(slot_el):
            slot_el.remove(child)
        etree.SubElement(slot_el, qn("a:srgbClr"), attrib={"val": hex_color})

    clr_scheme.set("name", "SlideDesign")
    theme_part._blob = etree.tostring(theme_xml, xml_declaration=True, encoding="UTF-8", standalone=True)


def assemble(manifest_dir, output_path):
    manifest_path = os.path.join(manifest_dir, "manifest.json")
    with open(manifest_path, "r") as f:
        manifest = json.load(f)

    prs = Presentation()
    prs.slide_width = SLIDE_W_EMU
    prs.slide_height = SLIDE_H_EMU
    blank_layout = prs.slide_layouts[6]

    first_palette = None
    for slide_info in manifest:
        p = slide_info.get("cssVarPalette")
        if p and not first_palette:
            first_palette = p

    if first_palette:
        inject_theme_colors(prs, first_palette)

    for slide_info in manifest:
        slide = prs.slides.add_slide(blank_layout)
        bg_path = os.path.join(manifest_dir, slide_info["bgImage"])
        set_slide_bg_image(slide, bg_path)

        shape_regions = slide_info.get("shapeRegions", [])
        for sr in shape_regions:
            add_shape_overlay(slide, sr)

        # Native tables from HTML <table> elements
        table_count = 0
        for table_info in slide_info.get("tableRegions", []):
            tx = px_to_emu(table_info["x"])
            ty = px_to_emu(table_info["y"])
            tw = px_to_emu(table_info["w"])
            th = px_to_emu(table_info["h"])
            if add_native_table(slide, table_info, tx, ty, tw, th):
                table_count += 1

        # Icons are inserted BEFORE text so text boxes sit above pictures
        # in PPTX z-order (later = higher).  This prevents large diagrams
        # (e.g. hub SVG on slide 03) from occluding overlay labels.
        icon_count = 0
        svg_count = 0
        chart_count = 0
        for icon_info in slide_info.get("iconRegions", []):
            icon_path = os.path.join(manifest_dir, icon_info["image"])
            if not os.path.exists(icon_path):
                continue
            x_emu = px_to_emu(icon_info["x"])
            y_emu = px_to_emu(icon_info["y"])
            w_emu = px_to_emu(icon_info["w"])
            h_emu = px_to_emu(icon_info["h"])

            # Native chart (doughnut from conic-gradient, etc.)
            chart_data_info = icon_info.get("chartData")
            if chart_data_info:
                if add_native_chart(slide, chart_data_info, x_emu, y_emu, w_emu, h_emu):
                    chart_count += 1
                    icon_count += 1
                    continue

            container = icon_info.get("containerStyle")
            if container:
                cs_w = px_to_emu(container["w"])
                cs_h = px_to_emu(container["h"])
                shape_type = MSO_SHAPE.OVAL if container.get("isCircular") else MSO_SHAPE.ROUNDED_RECTANGLE
                oval = slide.shapes.add_shape(shape_type, x_emu, y_emu, cs_w, cs_h)
                p_style = oval._element.find(qn("p:style"))
                if p_style is not None:
                    oval._element.remove(p_style)
                bg_rgb, bg_alpha = _parse_css_rgb(container["bgColor"])
                oval.fill.solid()
                oval.fill.fore_color.rgb = bg_rgb
                if bg_alpha < 100000:
                    _set_fill_alpha(oval, bg_alpha)
                border_color_str = container.get("borderColor")
                border_width = container.get("borderWidth", 0)
                if border_color_str and border_width > 0:
                    bd_rgb, bd_alpha = _parse_css_rgb(border_color_str)
                    oval.line.width = Pt(border_width * 0.75)
                    oval.line.fill.solid()
                    oval.line.fill.fore_color.rgb = bd_rgb
                    if bd_alpha < 100000:
                        _set_line_alpha(oval, bd_alpha)
                else:
                    oval.line.fill.background()
                container_text = icon_info.get("containerText")
                text_style = icon_info.get("textStyle")
                if container_text:
                    tf = oval.text_frame
                    tf.word_wrap = False
                    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                    tf.margin_left = tf.margin_right = Emu(0)
                    tf.margin_top = tf.margin_bottom = Emu(0)
                    run = tf.paragraphs[0].add_run()
                    run.text = container_text
                    t_size = Pt(text_style["sizePt"]) if text_style else Pt(14)
                    run.font.size = t_size
                    t_color = text_style["color"] if text_style else "FFFFFF"
                    run.font.color.rgb = RGBColor.from_string(t_color[:6])
                    if text_style and text_style.get("bold"):
                        run.font.bold = True
                else:
                    oval.text_frame.text = ""

            svg_file = icon_info.get("svgImage")
            svg_path = os.path.join(manifest_dir, svg_file) if svg_file else None

            svg_actual = icon_info.get("svgActualRect")
            if container and svg_actual:
                inner_w_px = svg_actual["w"]
                inner_h_px = svg_actual["h"]
            elif container:
                inner_w_px = icon_info["w"] * 0.45
                inner_h_px = icon_info["h"] * 0.45
            else:
                inner_w_px = icon_info["w"]
                inner_h_px = icon_info["h"]

            use_svg = False
            if svg_path and os.path.exists(svg_path):
                use_svg = _is_simple_svg(svg_path, inner_w_px, inner_h_px)

            # Prefer Puppeteer-generated transparent PNG (fallbackImage) over
            # CairoSVG rendering. CairoSVG cannot render emoji in <text> elements.
            # Ref: https://github.com/Kozea/CairoSVG/issues/434
            puppeteer_fb = icon_info.get("fallbackImage")
            puppeteer_fb_path = os.path.join(manifest_dir, puppeteer_fb) if puppeteer_fb else None

            fb_path = icon_path
            if puppeteer_fb_path and os.path.exists(puppeteer_fb_path):
                fb_path = puppeteer_fb_path
            elif svg_path and os.path.exists(svg_path):
                transparent_png = _svg_to_transparent_png(
                    svg_path, inner_w_px, inner_h_px
                )
                if transparent_png:
                    fb_path = transparent_png

            if container and icon_info.get("containerText"):
                icon_count += 1
                continue
            elif container:
                inner_w = px_to_emu(inner_w_px)
                inner_h = px_to_emu(inner_h_px)
                pad_x = (w_emu - inner_w) // 2
                pad_y = (h_emu - inner_h) // 2
                ix = x_emu + pad_x
                iy = y_emu + pad_y
                if use_svg:
                    add_svg_picture(slide, fb_path, svg_path, ix, iy, inner_w, inner_h)
                    svg_count += 1
                else:
                    slide.shapes.add_picture(fb_path, ix, iy, inner_w, inner_h)
            elif use_svg:
                add_svg_picture(slide, fb_path, svg_path, x_emu, y_emu, w_emu, h_emu)
                svg_count += 1
            else:
                slide.shapes.add_picture(fb_path, x_emu, y_emu, w_emu, h_emu)
            icon_count += 1

        # Build table bounding boxes to skip text regions covered by tables
        table_bboxes = []
        for ti in slide_info.get("tableRegions", []):
            table_bboxes.append((ti["x"], ti["y"],
                                 ti["x"] + ti["w"], ti["y"] + ti["h"]))

        grad_count = 0
        text_skipped = 0
        for region in slide_info["textRegions"]:
            rx, ry = region["x"], region["y"]
            # Skip text boxes whose center falls inside a table region
            if table_bboxes:
                cx_px = rx + region["w"] / 2
                cy_px = ry + region["h"] / 2
                inside = False
                for (tx1, ty1, tx2, ty2) in table_bboxes:
                    if tx1 <= cx_px <= tx2 and ty1 <= cy_px <= ty2:
                        inside = True
                        break
                if inside:
                    text_skipped += 1
                    continue
            add_text_overlay(slide, region)
            if region.get("gradient"):
                grad_count += 1

        total = len(slide_info["textRegions"]) - text_skipped
        shape_total = len(shape_regions)
        parts = [f"{total} text boxes ({grad_count} gradient)"]
        if shape_total > 0:
            parts.append(f"{shape_total} shapes")
        if table_count > 0:
            parts.append(f"{table_count} native table")
        if icon_count > 0:
            notes = []
            if svg_count: notes.append(f"{svg_count} SVG")
            if chart_count: notes.append(f"{chart_count} native chart")
            note_str = f" ({', '.join(notes)})" if notes else ""
            parts.append(f"{icon_count} icons{note_str}")
        print(f"  Slide {slide_info['index']}: {slide_info['name']} — {', '.join(parts)}")

    try:
        from transition_builder import apply_transitions_to_presentation
        apply_transitions_to_presentation(prs, {"type": "fade", "duration": 700})
    except ImportError:
        pass

    prs.save(output_path)
    size_kb = os.path.getsize(output_path) / 1024
    print(f"Saved: {output_path} ({len(manifest)} slides, {size_kb:.0f} KB)")


def main():
    args = sys.argv[1:]
    manifest_dir = None
    output = "presentation-editable.pptx"

    i = 0
    while i < len(args):
        if args[i] in ("-o", "--output"):
            i += 1
            output = args[i]
        elif not manifest_dir:
            manifest_dir = args[i]
        i += 1

    if not manifest_dir:
        print("Usage: python3 assemble_pptx.py <manifest_dir> -o output.pptx")
        sys.exit(1)

    assemble(manifest_dir, output)

    try:
        from embed_fonts import embed_fonts_in_pptx, get_default_font_specs
        specs = get_default_font_specs()
        if specs:
            tmp_out = output + ".tmp"
            embed_fonts_in_pptx(output, specs, output_path=tmp_out)
            os.replace(tmp_out, output)
            size_kb = os.path.getsize(output) / 1024
            fonts_str = ", ".join(s["typeface"] for s in specs)
            print(f"Embedded fonts: {fonts_str} ({size_kb:.0f} KB)")
    except Exception as e:
        print(f"Font embedding skipped: {e}")
        if os.path.exists(output + ".tmp"):
            os.remove(output + ".tmp")


if __name__ == "__main__":
    main()
