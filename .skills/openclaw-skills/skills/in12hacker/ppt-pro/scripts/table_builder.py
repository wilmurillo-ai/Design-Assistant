"""
table_builder.py — Native PPT table generation from structured data.

Creates editable tables with custom styling, merged cells, and border control.

python-pptx table API: https://python-pptx.readthedocs.io/en/stable/user/table.html
Border workaround: https://stackoverflow.com/questions/42610829/
"""

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Pt, Emu


TCPR_CHILD_ORDER = [
    "lnL", "lnR", "lnT", "lnB", "lnTlToBr", "lnBlToTr",
    "noFill", "solidFill", "gradFill", "blipFill", "pattFill", "grpFill",
    "cell3D", "extLst",
]


def _insert_tcPr_child(tcPr, new_elem):
    """Insert child into tcPr at the correct schema position.

    Reference: ECMA-376 §21.1.3.17 — tcPr child sequence:
    lnL, lnR, lnT, lnB, lnTlToBr, lnBlToTr, [fill], cell3D, extLst
    """
    new_local = etree.QName(new_elem).localname
    try:
        new_idx = TCPR_CHILD_ORDER.index(new_local)
    except ValueError:
        tcPr.append(new_elem)
        return
    for i, existing in enumerate(tcPr):
        existing_local = etree.QName(existing).localname
        try:
            existing_idx = TCPR_CHILD_ORDER.index(existing_local)
        except ValueError:
            continue
        if existing_idx > new_idx:
            tcPr.insert(i, new_elem)
            return
    tcPr.append(new_elem)


def _set_cell_border(cell, side, color_hex, width_pt=1.0):
    """Set border on one side of a table cell via lxml.

    side: "L", "R", "T", or "B"
    Reference: https://stackoverflow.com/questions/65072027/
    """
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    tag = qn(f"a:ln{side}")
    existing = tcPr.find(tag)
    if existing is not None:
        tcPr.remove(existing)

    ln = etree.Element(tag, attrib={"w": str(int(width_pt * 12700))})
    sf = etree.SubElement(ln, qn("a:solidFill"))
    etree.SubElement(sf, qn("a:srgbClr"), attrib={"val": color_hex[:6].upper()})
    _insert_tcPr_child(tcPr, ln)


def _set_cell_fill(cell, color_hex, alpha=None):
    """Set solid fill on a table cell."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for child in list(tcPr):
        if etree.QName(child).localname in ("solidFill", "noFill", "gradFill",
                                              "blipFill", "pattFill", "grpFill"):
            tcPr.remove(child)
    sf = etree.Element(qn("a:solidFill"))
    clr = etree.SubElement(sf, qn("a:srgbClr"), attrib={"val": color_hex[:6].upper()})
    if alpha is not None and alpha < 100000:
        etree.SubElement(clr, qn("a:alpha"), attrib={"val": str(int(alpha))})
    _insert_tcPr_child(tcPr, sf)


def add_table(slide, table_spec, x_emu, y_emu, w_emu, h_emu):
    """Add an editable table to a slide from a table specification dict.

    table_spec schema:
      {
        "headers": ["Col A", "Col B", "Col C"],
        "rows": [
          ["val1", "val2", "val3"],
          ["val4", "val5", "val6"]
        ],
        "headerBg": "#1E293B",
        "headerColor": "#FFFFFF",
        "rowBg": "#0F172A",
        "rowAltBg": "#1E293B",
        "rowColor": "#E2E8F0",
        "borderColor": "#334155",
        "borderWidth": 0.5,
        "fontSize": 10,
        "headerFontSize": 11
      }

    Returns the table shape object.
    """
    headers = table_spec.get("headers", [])
    rows_data = table_spec.get("rows", [])
    total_rows = len(rows_data) + (1 if headers else 0)
    total_cols = len(headers) if headers else (len(rows_data[0]) if rows_data else 1)

    shape = slide.shapes.add_table(total_rows, total_cols, x_emu, y_emu, w_emu, h_emu)
    table = shape.table

    header_bg = table_spec.get("headerBg", "1E293B")
    header_color = table_spec.get("headerColor", "FFFFFF")
    row_bg = table_spec.get("rowBg", "0F172A")
    row_alt_bg = table_spec.get("rowAltBg", "1E293B")
    row_color = table_spec.get("rowColor", "E2E8F0")
    border_color = table_spec.get("borderColor", "334155")
    border_w = table_spec.get("borderWidth", 0.5)
    font_size = table_spec.get("fontSize", 10)
    header_font_size = table_spec.get("headerFontSize", 11)

    def _clean_hex(c):
        return c.lstrip("#").upper()[:6] if c else "333333"

    row_offset = 0
    if headers:
        for ci, hdr in enumerate(headers):
            cell = table.cell(0, ci)
            cell.text = str(hdr)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(header_font_size)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor.from_string(_clean_hex(header_color))
            _set_cell_fill(cell, _clean_hex(header_bg))
            for side in ("L", "R", "T", "B"):
                _set_cell_border(cell, side, _clean_hex(border_color), border_w)
        row_offset = 1

    for ri, row in enumerate(rows_data):
        bg = _clean_hex(row_alt_bg if ri % 2 else row_bg)
        for ci, val in enumerate(row):
            if ci >= total_cols:
                break
            cell = table.cell(ri + row_offset, ci)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = RGBColor.from_string(_clean_hex(row_color))
            _set_cell_fill(cell, bg)
            for side in ("L", "R", "T", "B"):
                _set_cell_border(cell, side, _clean_hex(border_color), border_w)

    return shape
