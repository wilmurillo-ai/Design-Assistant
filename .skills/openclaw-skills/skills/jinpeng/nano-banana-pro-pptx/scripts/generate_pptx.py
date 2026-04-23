# /// script
# requires-python = ">=3.11"
# dependencies = [
#   "google-genai",
#   "python-pptx",
#   "Pillow",
# ]
# ///

import argparse
import base64
import io
import json
import os
import re
import sys
import time
from datetime import datetime
from pathlib import Path

from google import genai
from google.genai import types
from PIL import Image, UnidentifiedImageError
from pptx import Presentation
from pptx.util import Emu


def slugify(text: str, max_len: int = 40) -> str:
    """Convert prompt text to a filename-safe slug.

    For ASCII text, lowercases and replaces non-alphanumeric runs with hyphens.
    For non-ASCII text (e.g. Chinese), keeps the original characters up to max_len,
    replacing only filesystem-unsafe characters.
    """
    ascii_slug = text.lower()
    ascii_slug = re.sub(r"[^a-z0-9]+", "-", ascii_slug)
    ascii_slug = ascii_slug.strip("-")[:max_len].rstrip("-")
    if ascii_slug:
        return ascii_slug

    safe = re.sub(r'[\\/:*?"<>|]+', "-", text.strip())
    safe = safe[:max_len].rstrip("-").strip()
    return safe if safe else "presentation"


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Generate a PowerPoint presentation with AI-generated slide images."
    )
    parser.add_argument("--prompt", required=True, help="Topic/theme for the presentation")
    parser.add_argument("--slides", required=True, type=int, help="Number of slides (1-50)")
    parser.add_argument("--filename", default=None, help="Output directory name (default: auto-generated from prompt)")
    parser.add_argument(
        "--resolution",
        choices=["1K", "2K", "4K"],
        default="1K",
        help="Image resolution (default: 1K)",
    )
    parser.add_argument("--api-key", default=None, help="Gemini API key")
    parser.add_argument("--base-url", default=None, help="Custom Gemini API base URL")
    return parser.parse_args()


RESOLUTION_MAP = {"1K": 1024, "2K": 2048, "4K": 4096}

SLIDE_WIDTH_EMU = 9144000
SLIDE_HEIGHT_EMU = 5143500


def resolve_config(args: argparse.Namespace) -> tuple[str, str | None, Path, int]:
    """Resolve API key, base URL, output directory, and target width."""
    api_key = args.api_key or os.environ.get("GEMINI_API_KEY")
    if not api_key:
        print("Error: GEMINI_API_KEY is not set. Pass --api-key or set the environment variable.")
        sys.exit(1)

    if not (1 <= args.slides <= 50):
        print(f"Error: --slides must be between 1 and 50, got {args.slides}.")
        sys.exit(1)

    base_url = args.base_url or os.environ.get("GEMINI_BASE_URL")

    slug = slugify(args.prompt)
    timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    dir_name = args.filename if args.filename else f"{slug}-{timestamp}"
    out_dir = Path(dir_name)
    out_dir.mkdir(parents=True, exist_ok=True)

    target_width = RESOLUTION_MAP[args.resolution]
    return api_key, base_url, out_dir, target_width


def generate_slide_plans(
    client: genai.Client,
    topic: str,
    n: int,
    out_dir: Path,
) -> list[dict]:
    """Phase 1: Generate N slide plans (title + bullet points) via gemini-3.1-pro-preview."""
    system_instruction = (
        f"You are a professional presentation designer. Generate a {n}-slide PowerPoint outline "
        f"for the topic: '{topic}'.\n"
        f"Return ONLY a valid JSON array with exactly {n} objects. No markdown fences, no explanation.\n"
        'Each object must have:\n'
        '  "title": short slide title (4-8 words)\n'
        '  "points": array of 3-5 concise bullet points (each 5-15 words)\n'
        'The slides should form a coherent narrative arc from introduction to conclusion.'
    )

    def _call() -> list[dict]:
        response = client.models.generate_content(
            model="gemini-3.1-pro-preview",
            contents=system_instruction,
        )
        raw = response.text.strip()
        raw = re.sub(r"^```(?:json)?\s*", "", raw)
        raw = re.sub(r"\s*```$", "", raw)
        parsed = json.loads(raw)
        if not isinstance(parsed, list) or len(parsed) != n:
            raise ValueError(
                f"Expected list of {n} objects, got: {type(parsed).__name__} "
                f"len={len(parsed) if isinstance(parsed, list) else '?'}"
            )
        for item in parsed:
            if not isinstance(item, dict) or "title" not in item or "points" not in item:
                raise ValueError(f"Each slide must have 'title' and 'points' keys, got: {item}")
        return parsed

    print(f"Planning {n} slides for: {topic!r}")
    try:
        slides = _call()
    except Exception as e:
        print(f"Phase 1 attempt 1 failed: {e}. Retrying...")
        try:
            slides = _call()
        except Exception as e2:
            print(f"Phase 1 failed after retry: {e2}")
            sys.exit(1)

    plans_path = out_dir / "slide-plans.json"
    plans_path.write_text(
        json.dumps({"topic": topic, "slides": slides}, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    print(f"Saved slide plans to: {plans_path}")
    return slides


def build_image_prompt(slide: dict, topic: str) -> str:
    """Build a full slide design prompt from a slide plan.

    Asks the image model to generate a complete PPT slide with the title and
    bullet points visually composed into the design, banana-slides style.
    """
    points_text = "\n".join(f"• {p}" for p in slide["points"])
    return (
        f"Create a professional PowerPoint presentation slide image in 16:9 widescreen format "
        f"for a presentation about '{topic}'. "
        f"The slide must clearly display the following text content:\n\n"
        f"TITLE: {slide['title']}\n\n"
        f"KEY POINTS:\n{points_text}\n\n"
        f"Design style: visually stunning, modern presentation design. "
        f"Use a thematically relevant, high-quality background. "
        f"Title text should be large, bold, and prominent. "
        f"Bullet points should be clearly readable with high contrast. "
        f"Professional layout with clean typography and balanced composition. "
        f"The final image should look like a polished, presentation-ready slide."
    )


def generate_images(
    client: genai.Client,
    slides: list[dict],
    topic: str,
    target_width: int,
    out_dir: Path,
) -> list[Path]:
    """Phase 2: Generate one full slide design image per slide plan."""
    image_paths = []
    n = len(slides)

    # Build and save all image prompts before generating
    prompts = [
        {"slide": i, "title": s["title"], "image_prompt": build_image_prompt(s, topic)}
        for i, s in enumerate(slides, 1)
    ]
    (out_dir / "slide-prompts.json").write_text(
        json.dumps(prompts, ensure_ascii=False, indent=2), encoding="utf-8"
    )

    max_retries = 3
    retry_delay = 2  # seconds, doubles each retry
    interval = 1     # seconds between slides

    log_path = out_dir / "image-requests.log"
    log_entries = []

    for i, slide in enumerate(slides, 1):
        print(f"Generating slide {i}/{n}: {slide['title']!r}")
        img_prompt = prompts[i - 1]["image_prompt"]

        img_bytes = None
        for attempt in range(1, max_retries + 1):
            t_start = time.time()
            started_at = datetime.now().isoformat(timespec="seconds")
            status = "ok"
            error_msg = None
            try:
                response = client.models.generate_content(
                    model="gemini-3-pro-image-preview",
                    contents=img_prompt,
                    config=types.GenerateContentConfig(
                        response_modalities=["IMAGE"],
                    ),
                )
                raw_data = response.candidates[0].content.parts[0].inline_data.data
                img_bytes = raw_data if isinstance(raw_data, bytes) else base64.b64decode(raw_data)
            except Exception as e:
                status = "error"
                error_msg = str(e)
            finally:
                elapsed = round(time.time() - t_start, 2)
                entry = {
                    "slide": i,
                    "title": slide["title"],
                    "attempt": attempt,
                    "started_at": started_at,
                    "elapsed_s": elapsed,
                    "status": status,
                }
                if error_msg:
                    entry["error"] = error_msg
                log_entries.append(entry)
                log_path.write_text(
                    json.dumps(log_entries, ensure_ascii=False, indent=2), encoding="utf-8"
                )

            if status == "ok":
                break
            if attempt < max_retries:
                wait = retry_delay * (2 ** (attempt - 1))
                print(f"  Attempt {attempt} failed: {error_msg}. Retrying in {wait}s...")
                time.sleep(wait)
            else:
                print(f"Error generating image for slide {i} after {max_retries} attempts: {error_msg}")
                sys.exit(1)

        try:
            img = Image.open(io.BytesIO(img_bytes))
        except UnidentifiedImageError:
            print(f"Error: Could not identify image for slide {i}.")
            sys.exit(1)

        aspect = img.height / img.width
        new_height = int(target_width * aspect)
        img = img.resize((target_width, new_height), Image.LANCZOS)

        title_slug = slugify(slide["title"], max_len=40)
        out_path = out_dir / f"slide-{i:02d}-{title_slug}.png"
        img.save(out_path, "PNG")
        image_paths.append(out_path)

        if i < n:
            time.sleep(interval)

    return image_paths


def assemble_pptx(image_paths: list[Path], out_dir: Path, slug: str) -> Path:
    """Phase 3: Assemble full-bleed slide images into a PPTX."""
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_WIDTH_EMU)
    prs.slide_height = Emu(SLIDE_HEIGHT_EMU)

    blank_layout = prs.slide_layouts[6]

    for img_path in image_paths:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img_path),
            left=Emu(0),
            top=Emu(0),
            width=Emu(SLIDE_WIDTH_EMU),
            height=Emu(SLIDE_HEIGHT_EMU),
        )

    output_path = out_dir / f"{slug}.pptx"
    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    args = parse_args()
    api_key, base_url, out_dir, target_width = resolve_config(args)
    slug = slugify(args.prompt)

    http_options = {"base_url": base_url} if base_url else {}
    client = genai.Client(api_key=api_key, http_options=http_options)

    slides = generate_slide_plans(client, args.prompt, args.slides, out_dir)
    image_paths = generate_images(client, slides, args.prompt, target_width, out_dir)
    print("Assembling presentation...")
    pptx_path = assemble_pptx(image_paths, out_dir, slug)

    print(f"Done! Output folder: {out_dir.resolve()}")
    print(f"  Presentation: {pptx_path.name}")
    print(f"  Images: {len(image_paths)} slides")
    print(f"  Plans: slide-plans.json")
    print(f"  Image prompts: slide-prompts.json")
