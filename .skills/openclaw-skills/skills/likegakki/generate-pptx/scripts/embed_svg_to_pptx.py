#!/usr/bin/env python3
"""
Embed SVG files into a PPTX package by patching OOXML directly.
"""

import argparse
import re
import shutil
import sys
import tempfile
import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError as exc:
    raise SystemExit(
        "python-pptx is required. Run this script through run_in_skill_env.py or install it in the selected environment."
    ) from exc

SLIDE_WIDTH_EMU = Emu(9144000)
SLIDE_HEIGHT_EMU = Emu(5143500)

_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_NS_REL = "http://schemas.openxmlformats.org/package/2006/relationships"
_NS_CT = "http://schemas.openxmlformats.org/package/2006/content-types"
_NS_A14 = "http://schemas.microsoft.com/office/drawing/2010/main"
_NS_ASVG = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"

_REL_TYPE_IMAGE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
_SVG_BLIP_EXT_URI = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"
_USE_LOCAL_DPI_EXT_URI = "{28A0092B-C50C-407E-A947-70E740481C1C}"

ET.register_namespace("p", _NS_P)
ET.register_namespace("a", _NS_A)
ET.register_namespace("r", _NS_R)
ET.register_namespace("a14", _NS_A14)
ET.register_namespace("asvg", _NS_ASVG)


def _qn(namespace: str, tag: str) -> str:
    return f"{{{namespace}}}{tag}"


def _build_base_pptx(slide_count: int, output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH_EMU
    prs.slide_height = SLIDE_HEIGHT_EMU

    blank_layout = prs.slide_layouts[6]
    for _ in range(slide_count):
        prs.slides.add_slide(blank_layout)

    prs.save(str(output_path))


def _ensure_svg_content_type(content_types_path: Path) -> None:
    tree = ET.parse(content_types_path)
    root = tree.getroot()

    for node in root.findall(_qn(_NS_CT, "Default")):
        if node.attrib.get("Extension", "").lower() == "svg":
            node.attrib["ContentType"] = "image/svg+xml"
            tree.write(content_types_path, encoding="utf-8", xml_declaration=True)
            return

    ET.SubElement(
        root,
        _qn(_NS_CT, "Default"),
        {"Extension": "svg", "ContentType": "image/svg+xml"},
    )
    tree.write(content_types_path, encoding="utf-8", xml_declaration=True)


def _next_rel_id(rel_root: ET.Element) -> str:
    max_rel_id = 0
    for rel in rel_root.findall(_qn(_NS_REL, "Relationship")):
        match = re.fullmatch(r"rId(\d+)", rel.attrib.get("Id", ""))
        if match:
            max_rel_id = max(max_rel_id, int(match.group(1)))
    return f"rId{max_rel_id + 1}"


def _append_svg_relationship(rels_path: Path, target: str) -> str:
    tree = ET.parse(rels_path)
    root = tree.getroot()

    rel_id = _next_rel_id(root)
    ET.SubElement(
        root,
        _qn(_NS_REL, "Relationship"),
        {"Id": rel_id, "Type": _REL_TYPE_IMAGE, "Target": target},
    )
    tree.write(rels_path, encoding="utf-8", xml_declaration=True)
    return rel_id


def _next_shape_id(slide_root: ET.Element) -> int:
    max_shape_id = 0
    for node in slide_root.findall(f".//{_qn(_NS_P, 'cNvPr')}"):
        value = node.attrib.get("id")
        if value and value.isdigit():
            max_shape_id = max(max_shape_id, int(value))
    return max_shape_id + 1


def _build_svg_picture(shape_id: int, rel_id: str) -> ET.Element:
    pic = ET.Element(_qn(_NS_P, "pic"))

    nv_pic_pr = ET.SubElement(pic, _qn(_NS_P, "nvPicPr"))
    ET.SubElement(
        nv_pic_pr,
        _qn(_NS_P, "cNvPr"),
        {"id": str(shape_id), "name": f"SVG Picture {shape_id}"},
    )
    c_nv_pic_pr = ET.SubElement(nv_pic_pr, _qn(_NS_P, "cNvPicPr"))
    ET.SubElement(c_nv_pic_pr, _qn(_NS_A, "picLocks"), {"noChangeAspect": "1"})
    ET.SubElement(nv_pic_pr, _qn(_NS_P, "nvPr"))

    blip_fill = ET.SubElement(pic, _qn(_NS_P, "blipFill"))
    blip = ET.SubElement(blip_fill, _qn(_NS_A, "blip"), {_qn(_NS_R, "embed"): rel_id})
    ext_lst = ET.SubElement(blip, _qn(_NS_A, "extLst"))

    use_local_dpi_ext = ET.SubElement(ext_lst, _qn(_NS_A, "ext"), {"uri": _USE_LOCAL_DPI_EXT_URI})
    ET.SubElement(use_local_dpi_ext, _qn(_NS_A14, "useLocalDpi"), {"val": "0"})

    svg_ext = ET.SubElement(ext_lst, _qn(_NS_A, "ext"), {"uri": _SVG_BLIP_EXT_URI})
    ET.SubElement(svg_ext, _qn(_NS_ASVG, "svgBlip"), {_qn(_NS_R, "embed"): rel_id})

    stretch = ET.SubElement(blip_fill, _qn(_NS_A, "stretch"))
    ET.SubElement(stretch, _qn(_NS_A, "fillRect"))

    sp_pr = ET.SubElement(pic, _qn(_NS_P, "spPr"))
    xfrm = ET.SubElement(sp_pr, _qn(_NS_A, "xfrm"))
    ET.SubElement(xfrm, _qn(_NS_A, "off"), {"x": "0", "y": "0"})
    ET.SubElement(
        xfrm,
        _qn(_NS_A, "ext"),
        {"cx": str(int(SLIDE_WIDTH_EMU)), "cy": str(int(SLIDE_HEIGHT_EMU))},
    )
    prst_geom = ET.SubElement(sp_pr, _qn(_NS_A, "prstGeom"), {"prst": "rect"})
    ET.SubElement(prst_geom, _qn(_NS_A, "avLst"))
    return pic


def _append_svg_to_slide(slide_path: Path, rel_id: str) -> None:
    tree = ET.parse(slide_path)
    root = tree.getroot()
    sp_tree = root.find(f"./{_qn(_NS_P, 'cSld')}/{_qn(_NS_P, 'spTree')}")
    if sp_tree is None:
        raise ValueError(f"invalid slide XML, missing spTree: {slide_path}")

    sp_tree.append(_build_svg_picture(_next_shape_id(root), rel_id))
    tree.write(slide_path, encoding="utf-8", xml_declaration=True)


def _repack_pptx(package_dir: Path, output_path: Path) -> None:
    with zipfile.ZipFile(output_path, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        for file_path in sorted(package_dir.rglob("*")):
            if file_path.is_file():
                zf.write(file_path, file_path.relative_to(package_dir).as_posix())


def embed_svgs(svg_paths: list[Path], output_path: Path) -> None:
    if not svg_paths:
        raise ValueError("at least one SVG file is required")

    with tempfile.TemporaryDirectory() as base_tmp:
        base_pptx = Path(base_tmp) / "base.pptx"
        _build_base_pptx(len(svg_paths), base_pptx)

        with tempfile.TemporaryDirectory() as package_tmp:
            package_dir = Path(package_tmp) / "pptx"
            with zipfile.ZipFile(base_pptx, "r") as zf:
                zf.extractall(package_dir)

            _ensure_svg_content_type(package_dir / "[Content_Types].xml")

            media_dir = package_dir / "ppt" / "media"
            media_dir.mkdir(parents=True, exist_ok=True)

            for index, svg_path in enumerate(svg_paths, start=1):
                media_name = f"slide_{index:03d}.svg"
                shutil.copyfile(svg_path, media_dir / media_name)

                rels_path = package_dir / "ppt" / "slides" / "_rels" / f"slide{index}.xml.rels"
                slide_path = package_dir / "ppt" / "slides" / f"slide{index}.xml"

                rel_id = _append_svg_relationship(rels_path, f"../media/{media_name}")
                _append_svg_to_slide(slide_path, rel_id)
                print(f"embedded slide {index}: {svg_path.name}", file=sys.stderr)

            output_path.parent.mkdir(parents=True, exist_ok=True)
            _repack_pptx(package_dir, output_path)


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Embed SVG files into a PPTX package, one full-slide SVG per slide.",
    )
    parser.add_argument("svgs", nargs="+", type=Path, help="SVG file paths in slide order")
    parser.add_argument("-o", "--output", type=Path, default=Path("output.pptx"))
    args = parser.parse_args()

    missing = [path for path in args.svgs if not path.exists()]
    if missing:
        raise SystemExit(f"missing SVG files: {', '.join(str(path) for path in missing)}")

    embed_svgs(args.svgs, args.output)
    print(f"generated PPTX: {args.output}", file=sys.stderr)


if __name__ == "__main__":
    main()
