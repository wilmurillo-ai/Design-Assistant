#!/usr/bin/env python3
"""
Create PPTX file
Usage: pptx-write.py <output.pptx> [options]
"""
import sys
import argparse
from pathlib import Path

def import_pptx():
    """Import pptx with helpful error message"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        return Presentation, Inches, Pt
    except ImportError:
        print("Error: python-pptx not installed.", file=sys.stderr)
        print("Install: pip install python-pptx", file=sys.stderr)
        sys.exit(1)

def create_presentation(output_path, title=None, slides=1):
    """Create a PowerPoint presentation"""
    Presentation, Inches, Pt = import_pptx()
    
    prs = Presentation()
    
    # Title slide
    if title:
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = "Created with office-toolkit"
    
    # Add additional blank slides
    blank_layout = prs.slide_layouts[6]  # Blank layout
    for i in range(slides - 1):
        slide = prs.slides.add_slide(blank_layout)
        # Add slide number
        textbox = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(0.5)
        )
        textbox.text_frame.text = f"Slide {i + 2}"
    
    # Save
    prs.save(output_path)
    return output_path

def main():
    parser = argparse.ArgumentParser(
        description='Create PPTX file',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  pptx-write.py output.pptx --title "Presentation" --slides 5
  pptx-write.py output.pptx --title "Simple Deck"
        """
    )
    parser.add_argument('output', help='Output PPTX file')
    parser.add_argument('--title', help='Presentation title')
    parser.add_argument('--slides', type=int, default=1, help='Number of slides (default: 1)')
    
    args = parser.parse_args()
    
    # Validate output path
    output_path = Path(args.output)
    if output_path.suffix.lower() != '.pptx':
        print(f"Error: Output file must be .pptx: {output_path}", file=sys.stderr)
        sys.exit(1)
    
    if args.slides < 1:
        print("Error: Number of slides must be at least 1", file=sys.stderr)
        sys.exit(1)
    
    # Create parent directory
    output_path.parent.mkdir(parents=True, exist_ok=True)
    
    try:
        create_presentation(
            output_path,
            title=args.title,
            slides=args.slides
        )
        print(f"Created: {output_path}")
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)

if __name__ == "__main__":
    main()
