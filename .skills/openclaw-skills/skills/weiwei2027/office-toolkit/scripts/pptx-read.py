#!/usr/bin/env python3
"""
Read PPTX file and extract text content
Usage: pptx-read.py <file.pptx> [--json]
"""
import sys
import argparse
import json
from pathlib import Path

def import_pptx():
    """Import pptx with helpful error message"""
    try:
        from pptx import Presentation
        return Presentation
    except ImportError:
        print("Error: python-pptx not installed.", file=sys.stderr)
        print("Install: pip install python-pptx", file=sys.stderr)
        sys.exit(1)

def read_pptx(file_path):
    """Read PowerPoint and extract all text"""
    Presentation = import_pptx()
    
    prs = Presentation(file_path)
    content = {
        'file': str(file_path),
        'slides': [],
        'total_slides': len(prs.slides)
    }
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_content = {
            'number': slide_num,
            'texts': [],
            'tables': []
        }
        
        for shape in slide.shapes:
            # Extract text
            if hasattr(shape, "text") and shape.text.strip():
                slide_content['texts'].append(shape.text.strip())
            
            # Extract tables
            if shape.has_table:
                table_data = {
                    'rows': []
                }
                for row in shape.table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    table_data['rows'].append(row_text)
                slide_content['tables'].append(table_data)
        
        content['slides'].append(slide_content)
    
    return content

def main():
    parser = argparse.ArgumentParser(
        description='Read PPTX file and extract text content',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  pptx-read.py presentation.pptx
  pptx-read.py presentation.pptx --json
        """
    )
    parser.add_argument('file', help='Input PPTX file')
    parser.add_argument('--json', action='store_true', help='Output as JSON')
    
    args = parser.parse_args()
    
    # Validate file
    file_path = Path(args.file)
    if not file_path.exists():
        print(f"Error: File not found: {file_path}", file=sys.stderr)
        sys.exit(1)
    if not file_path.suffix.lower() == '.pptx':
        print(f"Error: Not a PPTX file: {file_path}", file=sys.stderr)
        sys.exit(1)
    
    try:
        content = read_pptx(file_path)
        
        if args.json:
            print(json.dumps(content, indent=2, ensure_ascii=False))
        else:
            # Plain text output
            print(f"File: {content['file']}")
            print(f"Total slides: {content['total_slides']}")
            print("\n" + "="*60)
            
            for slide in content['slides']:
                print(f"\n=== Slide {slide['number']} ===")
                for text in slide['texts']:
                    print(text)
                
                for i, table in enumerate(slide['tables'], 1):
                    print(f"\n[Table {i}]")
                    for row in table['rows']:
                        print(" | ".join(row))
                        
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)

if __name__ == "__main__":
    main()
