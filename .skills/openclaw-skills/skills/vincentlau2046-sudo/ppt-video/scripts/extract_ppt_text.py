#!/usr/bin/env python3
"""
PPTX 文字提取工具
从 PPTX 文件中提取每页的文字内容（标题 + 正文）
"""

import sys
import json
from pptx import Presentation

def extract_slide_text(pptx_path):
    """提取 PPTX 每页的文字内容"""
    try:
        prs = Presentation(pptx_path)
        slide_texts = []
        
        for i, slide in enumerate(prs.slides):
            text_content = []
            
            # 提取所有文本框中的文字
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text.strip())
            
            # 合并为该页的文字内容
            slide_text = ' '.join(text_content)
            slide_texts.append({
                'page': i + 1,
                'text': slide_text,
                'length': len(slide_text)
            })
        
        return slide_texts
    except Exception as e:
        print(json.dumps({'error': str(e)}), file=sys.stderr)
        return []

def extract_pdf_text(pdf_path):
    """从 PDF 提取文字（使用 pdftotext 或简单返回空）"""
    import subprocess
    try:
        # 尝试使用 pdftotext
        result = subprocess.run(
            ['pdftotext', '-layout', pdf_path, '-'],
            capture_output=True,
            text=True,
            timeout=30
        )
        if result.returncode == 0:
            content = result.stdout
            # 简单按页分割（基于换页符）
            pages = content.split('\f')
            return [
                {'page': i + 1, 'text': page.strip(), 'length': len(page)}
                for i, page in enumerate(pages) if page.strip()
            ]
    except Exception as e:
        print(json.dumps({'error': f'PDF extraction failed: {str(e)}'}), file=sys.stderr)
    return []

def main():
    if len(sys.argv) < 2:
        print(json.dumps({'error': 'No file provided'}), file=sys.stderr)
        sys.exit(1)
    
    file_path = sys.argv[1]
    
    if file_path.endswith('.pptx'):
        result = extract_slide_text(file_path)
    elif file_path.endswith('.pdf'):
        result = extract_pdf_text(file_path)
    else:
        print(json.dumps({'error': f'Unsupported file type: {file_path}'}), file=sys.stderr)
        sys.exit(1)
    
    # 输出 JSON 结果
    print(json.dumps(result, ensure_ascii=False))

if __name__ == '__main__':
    main()
