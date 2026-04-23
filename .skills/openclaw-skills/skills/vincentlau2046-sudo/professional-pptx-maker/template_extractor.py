#!/usr/bin/env python3
"""
Template Extractor - Extract layout contracts from PPTX templates (agent-slides core functionality)
"""

import json
import os
from pathlib import Path
from typing import Dict, List, Any

try:
    from pptx import Presentation
    from pptx.util import Inches
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

class TemplateExtractor:
    """Extract template contracts from PPTX files like agent-slides does."""
    
    def __init__(self):
        self.template_path = None
        self.output_dir = None
        
    def extract_template(self, template_path: str, output_dir: str) -> Dict[str, Any]:
        """Extract complete template contract from PPTX file."""
        self.template_path = template_path
        self.output_dir = output_dir
        
        # Ensure output directory exists
        os.makedirs(output_dir, exist_ok=True)
        
        # Load presentation
        prs = Presentation(template_path)
        
        # Extract all components
        template_layout = self._extract_template_layout(prs)
        content_layout = self._extract_content_layout(prs)
        archetypes = self._extract_archetypes(prs)
        theme = self._extract_theme(prs)
        resolved_manifest = self._create_resolved_manifest(
            template_layout, content_layout, archetypes, theme
        )
        
        # Create clean base template
        base_template_path = os.path.join(output_dir, "base_template.pptx")
        self._create_clean_template(template_path, base_template_path)
        
        # Save all artifacts
        self._save_artifacts(
            template_layout, content_layout, archetypes, 
            resolved_manifest, base_template_path
        )
        
        return {
            'template_layout': template_layout,
            'content_layout': content_layout,
            'archetypes': archetypes,
            'resolved_manifest': resolved_manifest,
            'base_template_path': base_template_path
        }
        
    def _extract_template_layout(self, prs: Presentation) -> Dict[str, Any]:
        """Extract physical layout families and placeholders."""
        layouts = []
        
        for i, layout in enumerate(prs.slide_layouts):
            layout_info = {
                'id': i,
                'name': layout.name or f"Layout {i}",
                'placeholders': [],
                'dimensions': {
                    'width_in': float(prs.slide_width) / 914400,
                    'height_in': float(prs.slide_height) / 914400
                }
            }
            
            # Extract placeholders
            for placeholder in layout.placeholders:
                ph_info = {
                    'idx': int(placeholder.placeholder_format.idx),
                    'type': str(placeholder.placeholder_format.type),
                    'name': placeholder.name,
                    'has_text_frame': hasattr(placeholder, 'text_frame'),
                    'position': {
                        'left_in': float(placeholder.left) / 914400 if hasattr(placeholder, 'left') else 0,
                        'top_in': float(placeholder.top) / 914400 if hasattr(placeholder, 'top') else 0,
                        'width_in': float(placeholder.width) / 914400 if hasattr(placeholder, 'width') else 0,
                        'height_in': float(placeholder.height) / 914400 if hasattr(placeholder, 'height') else 0
                    }
                }
                layout_info['placeholders'].append(ph_info)
                
            layouts.append(layout_info)
            
        return {'layouts': layouts}
        
    def _extract_content_layout(self, prs: Presentation) -> Dict[str, Any]:
        """Extract archetype-to-layout compatibility map."""
        # Define standard archetypes for financial presentations
        archetypes = {
            'title_slide': {
                'compatible_layouts': ['Title Slide'],
                'required_placeholders': ['title', 'subtitle']
            },
            'toc_slide': {
                'compatible_layouts': ['Title and Content', 'Title Only'],
                'required_placeholders': ['title', 'body']
            },
            'content_slide': {
                'compatible_layouts': ['Title and Content', 'Two Content'],
                'required_placeholders': ['title', 'body']
            },
            'chart_slide': {
                'compatible_layouts': ['Title and Content', 'Blank'],
                'required_placeholders': ['title']
            },
            'summary_slide': {
                'compatible_layouts': ['Title and Content'],
                'required_placeholders': ['title', 'body']
            },
            'section_header': {
                'compatible_layouts': ['Section Header'],
                'required_placeholders': ['title']
            }
        }
        
        return {'archetypes': archetypes}
        
    def _extract_archetypes(self, prs: Presentation) -> Dict[str, Any]:
        """Extract available archetypes with usage constraints."""
        # This would be more sophisticated in a real implementation
        # For now, we'll use the standard financial archetypes
        return self._extract_content_layout(prs)['archetypes']
        
    def _extract_theme(self, prs: Presentation) -> Dict[str, Any]:
        """Extract theme colors and typography."""
        # Default finance theme colors
        theme = {
            'palette': {
                'accent1': '007A7C',  # Finance Blue
                'accent2': 'FFC000',  # Gold
                'accent3': '4472C4',  # Corporate Blue
                'accent4': 'ED7D31',  # Orange
                'accent5': 'A5A5A5',  # Gray
                'accent6': '70AD47',  # Green
                'lt1': 'FFFFFF',      # Light text
                'dk1': '000000'       # Dark text
            },
            'font_scheme': {
                'major': 'Calibri Light',
                'minor': 'Calibri'
            },
            'background_colors': {
                'main': '0F172A',     # Dark blue background
                'secondary': '1E293B'
            }
        }
        return theme
        
    def _create_resolved_manifest(self, template_layout: Dict, content_layout: Dict, 
                                archetypes: Dict, theme: Dict) -> Dict[str, Any]:
        """Create merged contract with resolved layout bindings."""
        resolved_archetypes = {}
        
        for archetype_id, archetype in archetypes.items():
            # Find compatible layouts
            compatible_layout_names = archetype['compatible_layouts']
            resolved_layouts = []
            
            for layout_info in template_layout['layouts']:
                if layout_info['name'] in compatible_layout_names:
                    resolved_layouts.append({
                        'layout_id': layout_info['id'],
                        'layout_name': layout_info['name'],
                        'placeholders': layout_info['placeholders'],
                        'dimensions': layout_info['dimensions']
                    })
                    
            resolved_archetypes[archetype_id] = {
                'archetype_id': archetype_id,
                'resolved_layouts': resolved_layouts,
                'required_placeholders': archetype['required_placeholders']
            }
            
        manifest = {
            'theme': theme,
            'archetypes': resolved_archetypes,
            'template_dimensions': template_layout['layouts'][0]['dimensions'] if template_layout['layouts'] else {
                'width_in': 13.33,
                'height_in': 7.5
            }
        }
        
        return manifest
        
    def _create_clean_template(self, source_path: str, output_path: str):
        """Create clean template with no content slides."""
        prs = Presentation(source_path)
        
        # Remove all content slides, keep only masters and layouts
        slide_count = len(prs.slides)
        for i in range(slide_count - 1, -1, -1):
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[i])
            
        prs.save(output_path)
        
    def _save_artifacts(self, template_layout: Dict, content_layout: Dict, 
                       archetypes: Dict, resolved_manifest: Dict, base_template_path: str):
        """Save all extracted artifacts to output directory."""
        output_dir = self.output_dir
        
        # Save template layout
        with open(os.path.join(output_dir, 'template_layout.json'), 'w', encoding='utf-8') as f:
            json.dump(template_layout, f, indent=2, ensure_ascii=False)
            
        # Save content layout  
        with open(os.path.join(output_dir, 'content_layout.json'), 'w', encoding='utf-8') as f:
            json.dump(content_layout, f, indent=2, ensure_ascii=False)
            
        # Save archetypes
        with open(os.path.join(output_dir, 'archetypes.json'), 'w', encoding='utf-8') as f:
            json.dump(archetypes, f, indent=2, ensure_ascii=False)
            
        # Save resolved manifest (primary reference for build)
        with open(os.path.join(output_dir, 'resolved_manifest.json'), 'w', encoding='utf-8') as f:
            json.dump(resolved_manifest, f, indent=2, ensure_ascii=False)
            
        print(f"✅ Template extraction completed!")
        print(f"   Artifacts saved to: {output_dir}")
        print(f"   Base template: {base_template_path}")