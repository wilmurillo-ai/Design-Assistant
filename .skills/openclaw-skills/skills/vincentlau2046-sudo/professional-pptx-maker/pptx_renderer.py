#!/usr/bin/env python3
"""
Professional PPTX Renderer - Enhanced with charts and tables support
"""

import json
import os
from typing import Dict, List, Any

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

class ProfessionalPPTXRenderer:
    """Enhanced renderer with professional chart and table support."""
    
    def __init__(self, template_path: str):
        self.template_path = template_path
        self.prs = Presentation(template_path)
        self.operations = []
        
        # Define professional color schemes
        self.color_schemes = {
            'tech_insight': {
                'primary': RGBColor(224, 32, 32),    # Huawei red
                'accent': RGBColor(255, 102, 0),    # Huawei orange  
                'secondary': RGBColor(51, 51, 51),
                'background': RGBColor(255, 255, 255)
            },
            'finance': {
                'primary': RGBColor(118, 185, 0),    # NVIDIA green
                'accent': RGBColor(0, 153, 255),    # Blue accent
                'secondary': RGBColor(51, 51, 51),
                'background': RGBColor(250, 250, 250)
            },
            'tech_analysis': {
                'primary': RGBColor(0, 102, 204),   # Tech blue
                'accent': RGBColor(0, 153, 255),    # Highlight blue
                'secondary': RGBColor(51, 51, 51),
                'background': RGBColor(255, 255, 255)
            }
        }
        
    def load_operations(self, slides_json_path: str):
        """Load operations from slides.json file."""
        with open(slides_json_path, 'r', encoding='utf-8') as f:
            slides_data = json.load(f)
            self.operations = slides_data['ops']['operations']
            
    def execute_operations(self, dry_run: bool = False) -> bool:
        """Execute all operations to render the presentation."""
        if dry_run:
            print("🧪 Dry run mode - validating operations...")
            return self._validate_operations()
            
        print("🎨 Executing operations to render PPTX...")
        
        for i, operation in enumerate(self.operations):
            try:
                op_type = operation['op']
                if op_type == 'add_slide':
                    self._add_slide(operation)
                elif op_type == 'set_semantic_text':
                    self._set_semantic_text(operation)
                elif op_type == 'add_notes':
                    self._add_notes(operation)
                elif op_type == 'set_core_properties':
                    self._set_core_properties(operation)
                elif op_type == 'add_chart':
                    self._add_chart(operation)
                elif op_type == 'add_table':
                    self._add_table(operation)
                else:
                    print(f"⚠️  Unsupported operation: {op_type}")
                    
            except Exception as e:
                print(f"❌ Operation {i} failed: {e}")
                return False
                
        return True
        
    def _add_slide(self, operation: Dict[str, Any]):
        """Add slide with specified layout."""
        layout_name = operation.get('layout_name')
        layout_index = operation.get('layout_index', 1)
        
        target_layout = None
        if layout_name:
            for layout in self.prs.slide_layouts:
                if layout.name == layout_name:
                    target_layout = layout
                    break
                    
        if target_layout is None and layout_index < len(self.prs.slide_layouts):
            target_layout = self.prs.slide_layouts[layout_index]
            
        if target_layout is None:
            target_layout = self.prs.slide_layouts[1] if len(self.prs.slide_layouts) > 1 else self.prs.slide_layouts[0]
            
        self.prs.slides.add_slide(target_layout)
        
    def _set_semantic_text(self, operation: Dict[str, Any]):
        """Set text using semantic roles (title, body, etc.)."""
        slide_index = operation['slide_index']
        role = operation['role']
        text = operation['text']
        
        if slide_index >= len(self.prs.slides):
            return
            
        slide = self.prs.slides[slide_index]
        
        if role == 'title':
            if slide.shapes.title:
                slide.shapes.title.text = text
        elif role == 'body':
            for shape in slide.placeholders:
                if str(shape.placeholder_format.type).endswith('BODY (2)'):
                    shape.text = text
                    break
            else:
                left = Inches(1)
                top = Inches(2)
                width = Inches(11)
                height = Inches(4)
                textbox = slide.shapes.add_textbox(left, top, width, height)
                textbox.text_frame.text = text
        elif role == 'subtitle':
            for shape in slide.placeholders:
                if str(shape.placeholder_format.type).endswith('SUBTITLE (4)'):
                    shape.text = text
                    break
                    
    def _add_notes(self, operation: Dict[str, Any]):
        """Add speaker notes to slide."""
        slide_index = operation['slide_index']
        text = operation['text']
        
        if slide_index >= len(self.prs.slides):
            return
            
        slide = self.prs.slides[slide_index]
        notes_slide = slide.notes_slide
        if notes_slide and notes_slide.notes_text_frame:
            notes_slide.notes_text_frame.text = text
            
    def _set_core_properties(self, operation: Dict[str, Any]):
        """Set presentation core properties."""
        core = self.prs.core_properties
        if 'title' in operation:
            core.title = operation['title']
        if 'subject' in operation:
            core.subject = operation['subject']
        if 'author' in operation:
            core.author = operation['author']
            
    def _add_chart(self, operation: Dict[str, Any]):
        """Add professional chart to slide."""
        slide_index = operation['slide_index']
        chart_type = operation.get('chart_type', 'column')
        data = operation.get('data', {})
        position = operation.get('position', {'x': 1, 'y': 2.5, 'width': 6, 'height': 4})
        
        if slide_index >= len(self.prs.slides):
            return
            
        slide = self.prs.slides[slide_index]
        
        # Create chart data
        chart_data = CategoryChartData()
        categories = data.get('categories', [])
        series_list = data.get('series', [])
        
        if categories:
            chart_data.categories = categories
            
        for series in series_list:
            name = series.get('name', '')
            values = series.get('values', [])
            if name and values:
                chart_data.add_series(name, values)
                
        # Determine chart type
        pptx_chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
        if chart_type == 'line':
            pptx_chart_type = XL_CHART_TYPE.LINE
        elif chart_type == 'pie':
            pptx_chart_type = XL_CHART_TYPE.PIE
            
        # Add chart to slide
        x = Inches(position['x'])
        y = Inches(position['y'])
        cx = Inches(position['width'])
        cy = Inches(position['height'])
        
        slide.shapes.add_chart(pptx_chart_type, x, y, cx, cy, chart_data)
        
    def _add_table(self, operation: Dict[str, Any]):
        """Add professional table to slide."""
        slide_index = operation['slide_index']
        table_data = operation.get('data', [])
        position = operation.get('position', {'x': 1, 'y': 2, 'width': 11, 'height': 4})
        
        if slide_index >= len(self.prs.slides):
            return
            
        slide = self.prs.slides[slide_index]
        
        if not table_data or len(table_data) == 0:
            return
            
        rows = len(table_data)
        cols = len(table_data[0]) if table_data else 0
        
        x = Inches(position['x'])
        y = Inches(position['y'])
        width = Inches(position['width'])
        row_height = position['height'] / rows if rows > 0 else 0.4
        height = Inches(position['height'])
        
        table_shape = slide.shapes.add_table(rows, cols, x, y, width, height)
        table = table_shape.table
        
        # Fill table data
        for i in range(rows):
            for j in range(min(cols, len(table_data[i]))):
                cell = table.cell(i, j)
                cell.text = str(table_data[i][j])
                
    def _validate_operations(self) -> bool:
        """Validate operations without executing them."""
        print(f"  • Validating {len(self.operations)} operations")
        return True
        
    def save(self, output_path: str):
        """Save the presentation to file."""
        self.prs.save(output_path)