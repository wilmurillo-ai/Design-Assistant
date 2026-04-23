#!/usr/bin/env python3
"""
Create professional 16:9 finance template with bright color scheme
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_finance_template():
    """Create professional finance template with 16:9 aspect ratio and bright color scheme."""
    prs = Presentation()
    
    # Set 16:9 aspect ratio (13.33" x 7.5")
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Define bright color scheme for finance/tech theme
    COLOR_SCHEME = {
        'background': RGBColor(250, 250, 250),  # Very light gray #FAFAFA
        'primary': RGBColor(118, 185, 0),       # NVIDIA green #76B900
        'secondary': RGBColor(51, 51, 51),      # Dark gray #333333
        'light_gray': RGBColor(240, 240, 240)   # Light gray #F0F0F0
    }
    
    # Define fonts
    FONT_TITLE = 'Poppins'
    FONT_BODY = 'Roboto'
    
    # Create slide masters and layouts
    slide_master = prs.slide_masters[0]
    
    # Customize master background with bright color
    background = slide_master.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_SCHEME['background']
    
    # Title Slide Layout (Layout 0)
    title_layout = prs.slide_layouts[0]
    title_layout.name = "Title Slide"
    
    # Title and Content Layout (Layout 1)  
    content_layout = prs.slide_layouts[1]
    content_layout.name = "Title and Content"
    
    # Section Header Layout (Layout 2)
    section_layout = prs.slide_layouts[2] 
    section_layout.name = "Section Header"
    
    # Two Content Layout (Layout 3)
    two_content_layout = prs.slide_layouts[3]
    two_content_layout.name = "Two Content"
    
    # Comparison Layout (Layout 4)
    comparison_layout = prs.slide_layouts[4]
    comparison_layout.name = "Comparison"
    
    # Title Only Layout (Layout 5)
    title_only_layout = prs.slide_layouts[5]
    title_only_layout.name = "Title Only"
    
    # Blank Layout (Layout 6)
    blank_layout = prs.slide_layouts[6]
    blank_layout.name = "Blank"
    
    # Customize placeholder styles with bright color scheme
    for layout in prs.slide_layouts:
        for placeholder in layout.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                placeholder.text_frame.paragraphs[0].font.size = Pt(36)
                placeholder.text_frame.paragraphs[0].font.color.rgb = COLOR_SCHEME['secondary']
                placeholder.text_frame.paragraphs[0].font.bold = True
                placeholder.text_frame.paragraphs[0].font.name = FONT_TITLE
                placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            elif placeholder.placeholder_format.type == 2:  # Body/Content
                placeholder.text_frame.paragraphs[0].font.size = Pt(18)
                placeholder.text_frame.paragraphs[0].font.color.rgb = COLOR_SCHEME['secondary']
                placeholder.text_frame.paragraphs[0].font.name = FONT_BODY
                placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                
    return prs

if __name__ == "__main__":
    template = create_finance_template()
    template.save("finance_template.pptx")
    print("✅ Professional 16:9 bright finance template created!")