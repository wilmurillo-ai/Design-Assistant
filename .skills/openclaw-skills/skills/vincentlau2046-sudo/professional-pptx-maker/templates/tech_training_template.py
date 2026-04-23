#!/usr/bin/env python3
"""
Create professional 16:9 technology training template with white background and dark red highlight
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

def set_font(run, size, bold=False, color=None):
    """Helper function to set font properties consistently"""
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    run._element.rPr.rFonts.set(qn('w:eastAsia'), "Microsoft YaHei")

def create_tech_training_template():
    """Create professional technology training template with 16:9 aspect ratio."""
    prs = Presentation()
    
    # Set 16:9 aspect ratio (13.33" x 7.5")
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Define tech training color scheme (white background + dark red highlight)
    THEME = {
        "bg": RGBColor(255, 255, 255),      # Pure white background
        "primary": RGBColor(140, 20, 20),   # Dark red (highlight/primary)
        "secondary": RGBColor(0, 0, 0),     # Black for normal text
        "card": RGBColor(248, 248, 248),    # Light gray cards
        "line": RGBColor(210, 210, 210),    # Gray divider lines
        "gray_text": RGBColor(102, 102, 102) # Gray for secondary text
    }
    
    # Create slide masters and layouts
    slide_master = prs.slide_masters[0]
    
    # Customize master background with white color
    background = slide_master.background
    background.fill.solid()
    background.fill.fore_color.rgb = THEME["bg"]
    
    # Title Slide Layout (Layout 0)
    title_layout = prs.slide_layouts[0]
    title_layout.name = "Title Slide"
    
    # Title and Content Layout (Layout 1)  
    content_layout = prs.slide_layouts[1]
    content_layout.name = "Title and Content"
    
    # Section Header Layout (Layout 2)
    section_layout = prs.slide_layouts[2] 
    section_layout.name = "Section Header"
    
    # Two Content Layout (Layout 3)
    two_content_layout = prs.slide_layouts[3]
    two_content_layout.name = "Two Content"
    
    # Comparison Layout (Layout 4)
    comparison_layout = prs.slide_layouts[4]
    comparison_layout.name = "Comparison"
    
    # Title Only Layout (Layout 5)
    title_only_layout = prs.slide_layouts[5]
    title_only_layout.name = "Title Only"
    
    # Blank Layout (Layout 6)
    blank_layout = prs.slide_layouts[6]
    blank_layout.name = "Blank"
    
    # Customize placeholder styles with tech training color scheme
    for layout in prs.slide_layouts:
        for placeholder in layout.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                placeholder.text_frame.paragraphs[0].font.size = Pt(28)
                placeholder.text_frame.paragraphs[0].font.color.rgb = THEME["primary"]
                placeholder.text_frame.paragraphs[0].font.bold = True
                placeholder.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
                placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            elif placeholder.placeholder_format.type == 2:  # Body/Content
                placeholder.text_frame.paragraphs[0].font.size = Pt(18)
                placeholder.text_frame.paragraphs[0].font.color.rgb = THEME["secondary"]
                placeholder.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
                placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                
    return prs

if __name__ == "__main__":
    template = create_tech_training_template()
    template.save("tech_training_template.pptx")
    print("✅ Professional 16:9 technology training template created!")
