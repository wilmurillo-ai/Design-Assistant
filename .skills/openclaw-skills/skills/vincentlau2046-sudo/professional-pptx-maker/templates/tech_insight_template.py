#!/usr/bin/env python3
"""
Create professional 16:9 technology insight template with Huawei-style color scheme
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def add_content_card(slide, x, y, width, height, title_text, desc_text):
    """Helper function: Add a content card (title + description)"""
    # Add title
    title_box = slide.shapes.add_textbox(x, y, width, Inches(0.5))
    title_tf = title_box.text_frame
    title_tf.text = title_text
    title_p = title_tf.paragraphs[0]
    title_p.font.size = Pt(18)
    title_p.font.color.rgb = RGBColor(224, 32, 32)  # Huawei red
    title_p.font.bold = True
    title_p.alignment = PP_ALIGN.CENTER

    # Add description
    desc_box = slide.shapes.add_textbox(x, y + Inches(0.6), width, Inches(1))
    desc_tf = desc_box.text_frame
    desc_tf.text = desc_text
    desc_p = desc_tf.paragraphs[0]
    desc_p.font.size = Pt(12)
    desc_p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray
    desc_p.alignment = PP_ALIGN.CENTER

def create_tech_insight_template():
    """Create professional technology insight template with 16:9 aspect ratio."""
    prs = Presentation()
    
    # Set 16:9 aspect ratio (13.33" x 7.5")
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Define Huawei-style color scheme for tech insight theme
    COLOR_SCHEME = {
        'background': RGBColor(255, 255, 255),   # White #FFFFFF
        'primary_red': RGBColor(224, 32, 32),    # Huawei red #E02020
        'primary_orange': RGBColor(255, 102, 0), # Huawei orange #FF6600
        'secondary': RGBColor(51, 51, 51),       # Dark gray #333333
        'light_gray': RGBColor(240, 240, 240)    # Light gray #F0F0F0
    }
    
    # Define fonts (using system-compatible fonts)
    FONT_TITLE = 'Microsoft YaHei'  # Fallback for Source Han Sans
    FONT_BODY = 'Microsoft YaHei'   # Fallback for Source Han Sans
    
    # Set core properties
    prs.core_properties.title = "Technology Insight Template"
    prs.core_properties.subject = "Professional Tech Insight Template"
    
    # Create slide masters and layouts
    slide_master = prs.slide_masters[0]
    
    # Customize master background with white color
    background = slide_master.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_SCHEME['background']
    
    # Title Slide Layout (Layout 0)
    title_layout = prs.slide_layouts[0]
    title_layout.name = "Title Slide"
    
    # Title and Content Layout (Layout 1)  
    content_layout = prs.slide_layouts[1]
    content_layout.name = "Title and Content"
    
    # Section Header Layout (Layout 2)
    section_layout = prs.slide_layouts[2] 
    section_layout.name = "Section Header"
    
    # Two Content Layout (Layout 3)
    two_content_layout = prs.slide_layouts[3]
    two_content_layout.name = "Two Content"
    
    # Comparison Layout (Layout 4)
    comparison_layout = prs.slide_layouts[4]
    comparison_layout.name = "Comparison"
    
    # Title Only Layout (Layout 5)
    title_only_layout = prs.slide_layouts[5]
    title_only_layout.name = "Title Only"
    
    # Blank Layout (Layout 6)
    blank_layout = prs.slide_layouts[6]
    blank_layout.name = "Blank"
    
    # Customize placeholder styles with tech insight color scheme
    for layout in prs.slide_layouts:
        for placeholder in layout.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                placeholder.text_frame.paragraphs[0].font.size = Pt(36)
                placeholder.text_frame.paragraphs[0].font.color.rgb = COLOR_SCHEME['secondary']
                placeholder.text_frame.paragraphs[0].font.bold = True
                placeholder.text_frame.paragraphs[0].font.name = FONT_TITLE
                placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            elif placeholder.placeholder_format.type == 2:  # Body/Content
                placeholder.text_frame.paragraphs[0].font.size = Pt(16)
                placeholder.text_frame.paragraphs[0].font.color.rgb = COLOR_SCHEME['secondary']
                placeholder.text_frame.paragraphs[0].font.name = FONT_BODY
                placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                placeholder.text_frame.paragraphs[0].first_line_indent = Pt(0)
                
    return prs

if __name__ == "__main__":
    template = create_tech_insight_template()
    template.save("tech_insight_template.pptx")
    print("✅ Professional 16:9 technology insight template created!")