"""
演示处理模块
支持PPT演示文稿的大纲生成、智能排版和批量生成
"""

import io
import logging
from typing import Dict, List, Optional, Any, Union
from enum import Enum

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pydantic import BaseModel, Field, ConfigDict

logger = logging.getLogger(__name__)


class TemplateStyle(str, Enum):
    """模板风格"""
    BUSINESS = "business"
    CREATIVE = "creative"
    MINIMAL = "minimal"
    TECH = "tech"


class OutlineRequest(BaseModel):
    """大纲生成请求"""
    topic: str = Field(..., description="演示主题")
    slide_count: int = Field(10, description="幻灯片数量")
    target_audience: Optional[str] = Field(None, description="目标受众")
    key_points: List[str] = Field(default_factory=list, description="关键要点")


class SlideOutline(BaseModel):
    """幻灯片大纲"""
    title: str = Field(..., description="标题")
    content: List[str] = Field(default_factory=list, description="内容要点")
    layout_type: str = Field("title_content", description="布局类型")


class PresentationOutline(BaseModel):
    """演示文稿大纲"""
    title: str = Field(..., description="演示标题")
    slides: List[SlideOutline] = Field(default_factory=list, description="幻灯片列表")


class LayoutRequest(BaseModel):
    """智能排版请求"""
    outline: PresentationOutline = Field(..., description="演示大纲")
    style: TemplateStyle = Field(TemplateStyle.BUSINESS, description="模板风格")
    include_images: bool = Field(False, description="是否包含图片占位符")


class BatchGenerationRequest(BaseModel):
    """批量生成请求"""
    model_config = ConfigDict(arbitrary_types_allowed=True)
    
    data: Any = Field(..., description="数据源")
    template_slide: bytes = Field(..., description="模板幻灯片")
    field_mapping: Dict[str, str] = Field(..., description="字段映射")


class PresentationProcessor:
    """演示处理器"""
    
    def __init__(self):
        pass
        
    async def generate_outline(
        self,
        request: OutlineRequest,
    ) -> PresentationOutline:
        """
        生成演示大纲
        
        Args:
            request: 大纲生成请求
            
        Returns:
            演示大纲
        """
        logger.info(f"开始生成演示大纲: {request.topic}")
        
        return self._generate_template_outline(request)
        
    def _build_outline_prompt(self, request: OutlineRequest) -> str:
        """构建大纲生成提示词"""
        prompt = f"""请为以下主题生成一个演示文稿大纲：

主题：{request.topic}
幻灯片数量：{request.slide_count}
"""
        
        if request.target_audience:
            prompt += f"目标受众：{request.target_audience}\n"
            
        if request.key_points:
            prompt += f"关键要点：{', '.join(request.key_points)}\n"
            
        prompt += "\n请生成每张幻灯片的标题和内容要点。"
        return prompt
        
    def _generate_template_outline(self, request: OutlineRequest) -> PresentationOutline:
        """生成模板大纲"""
        slides = [
            SlideOutline(
                title=request.topic,
                content=["副标题或演讲者信息"],
                layout_type="title",
            )
        ]
        
        for i in range(1, request.slide_count - 1):
            if i <= len(request.key_points):
                slides.append(SlideOutline(
                    title=request.key_points[i - 1],
                    content=["要点1", "要点2", "要点3"],
                    layout_type="title_content",
                ))
            else:
                slides.append(SlideOutline(
                    title=f"第{i + 1}部分",
                    content=["内容要点"],
                    layout_type="title_content",
                ))
                
        slides.append(SlideOutline(
            title="谢谢",
            content=["联系方式"],
            layout_type="title",
        ))
        
        return PresentationOutline(title=request.topic, slides=slides)
        
    async def apply_layout(
        self,
        request: LayoutRequest,
    ) -> bytes:
        """
        智能排版
        
        Args:
            request: 排版请求
            
        Returns:
            生成的PPT文件字节流
        """
        logger.info(f"开始智能排版，风格: {request.style.value}")
        
        prs = Presentation()
        
        style_configs = self._get_style_config(request.style)
        
        for slide_outline in request.outline.slides:
            slide = self._create_slide(prs, slide_outline, style_configs)
            
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        logger.info(f"智能排版完成，共 {len(request.outline.slides)} 张幻灯片")
        return output.getvalue()
        
    def _get_style_config(self, style: TemplateStyle) -> Dict[str, Any]:
        """获取风格配置"""
        configs = {
            TemplateStyle.BUSINESS: {
                "title_color": RGBColor(0, 51, 102),
                "content_color": RGBColor(51, 51, 51),
                "title_size": Pt(36),
                "content_size": Pt(18),
                "background": None,
            },
            TemplateStyle.CREATIVE: {
                "title_color": RGBColor(255, 102, 0),
                "content_color": RGBColor(102, 102, 102),
                "title_size": Pt(40),
                "content_size": Pt(20),
                "background": None,
            },
            TemplateStyle.MINIMAL: {
                "title_color": RGBColor(0, 0, 0),
                "content_color": RGBColor(51, 51, 51),
                "title_size": Pt(32),
                "content_size": Pt(16),
                "background": None,
            },
            TemplateStyle.TECH: {
                "title_color": RGBColor(0, 112, 192),
                "content_color": RGBColor(89, 89, 89),
                "title_size": Pt(38),
                "content_size": Pt(18),
                "background": None,
            },
        }
        return configs.get(style, configs[TemplateStyle.BUSINESS])
        
    def _create_slide(
        self,
        prs: Presentation,
        outline: SlideOutline,
        style_config: Dict[str, Any],
    ) -> Any:
        """创建幻灯片"""
        if outline.layout_type == "title":
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            title.text = outline.title
            title.text_frame.paragraphs[0].font.size = style_config["title_size"]
            title.text_frame.paragraphs[0].font.color.rgb = style_config["title_color"]
            
        else:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            title.text = outline.title
            title.text_frame.paragraphs[0].font.size = style_config["title_size"]
            title.text_frame.paragraphs[0].font.color.rgb = style_config["title_color"]
            
            if len(slide.placeholders) > 1:
                content = slide.placeholders[1]
                text_frame = content.text_frame
                text_frame.clear()
                
                for i, point in enumerate(outline.content):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                        
                    p.text = point
                    p.font.size = style_config["content_size"]
                    p.font.color.rgb = style_config["content_color"]
                    p.level = 0
                    
        return slide
        
    async def batch_generate(
        self,
        request: BatchGenerationRequest,
    ) -> bytes:
        """
        批量生成PPT
        
        Args:
            request: 批量生成请求
            
        Returns:
            生成的PPT文件字节流
        """
        logger.info("开始批量生成PPT")
        
        df = self._load_data(request.data)
        
        template_prs = Presentation(io.BytesIO(request.template_slide))
        
        output_prs = Presentation()
        
        for _, row in df.iterrows():
            for slide in template_prs.slides:
                new_slide = output_prs.slides.add_slide(slide.slide_layout)
                
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        new_shape = new_slide.shapes.add_textbox(
                            shape.left, shape.top,
                            shape.width, shape.height
                        )
                        
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text
                            
                            for field, column in request.field_mapping.items():
                                placeholder = f"{{{{{field}}}}}"
                                if placeholder in text:
                                    text = text.replace(placeholder, str(row[column]))
                                    
                            new_para = new_shape.text_frame.add_paragraph()
                            new_para.text = text
                            new_para.font.size = paragraph.font.size
                            new_para.font.bold = paragraph.font.bold
                            
        output = io.BytesIO()
        output_prs.save(output)
        output.seek(0)
        
        logger.info(f"批量生成完成，共 {len(df)} 张幻灯片")
        return output.getvalue()
        
    def _load_data(self, data: Union[str, bytes, pd.DataFrame]) -> pd.DataFrame:
        """加载数据"""
        if isinstance(data, pd.DataFrame):
            return data.copy()
            
        if isinstance(data, bytes):
            return pd.read_excel(io.BytesIO(data))
            
        if isinstance(data, str):
            if data.endswith(('.xlsx', '.xls', '.et')):
                return pd.read_excel(data)
            elif data.endswith('.csv'):
                return pd.read_csv(data)
            else:
                raise ValueError(f"不支持的文件格式: {data}")
                
        raise ValueError("不支持的数据类型")
