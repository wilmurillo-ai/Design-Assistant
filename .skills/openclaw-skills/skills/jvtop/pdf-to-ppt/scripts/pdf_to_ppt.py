#!/usr/bin/env python3
"""
PDF to PPT Converter
Converts PDF files to PowerPoint presentations via intermediate image rendering.

Usage:
    python3 pdf_to_ppt.py input.pdf [options]

Options:
    --img-dir DIR        Directory to store intermediate images (default: <pdf_dir>/images/)
    --output PPTX        Path for output PPT file (default: <pdf_basename>.pptx in same dir as input)
    --zoom ZOOM          Zoom factor for rendering (default: 2.0)
    --slide-width-in IN  Slide width in inches (default: 13.33 for 16:9)
    --slide-height-in IN Slide height in inches (default: 7.5)
    --format {png,jpg}   Image format for intermediates (default: png)
"""

import argparse
import os
import sys
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches


def pdf_to_images(pdf_path, img_dir, zoom=2.0, format='png'):
    """Convert PDF pages to images and return list of image paths."""
    # Ensure image directory exists
    os.makedirs(img_dir, exist_ok=True)
    
    # Open PDF
    doc = fitz.open(pdf_path)
    print(f"PDF '{pdf_path}' has {doc.page_count} pages")
    
    image_paths = []
    
    # Render each page to an image
    for page_num in range(doc.page_count):
        page = doc.load_page(page_num)  # zero-based
        
        # Create transformation matrix for zoom
        mat = fitz.Matrix(zoom, zoom)
        
        # Render to pixmap (image)
        pix = page.get_pixmap(matrix=mat, alpha=False)  # alpha=False for white background
        
        # Save image
        if format.lower() == 'jpg' or format.lower() == 'jpeg':
            img_path = os.path.join(img_dir, f"page_{page_num+1:03d}.jpg")
            # Save as JPEG with quality 90
            pix.save(img_path, "jpg")
        else:
            img_path = os.path.join(img_dir, f"page_{page_num+1:03d}.png")
            pix.save(img_path)
        
        image_paths.append(img_path)
        print(f"  Saved {img_path}")
    
    doc.close()
    print(f"All {doc.page_count} pages rendered to images in '{img_dir}'")
    return image_paths


def images_to_ppt(image_paths, output_ppt, slide_width_in=13.33, slide_height_in=7.5):
    """Create a PPT from a list of image paths."""
    prs = Presentation()
    
    # Set slide size
    prs.slide_width = Inches(slide_width_in)
    prs.slide_height = Inches(slide_height_in)
    
    print(f"Creating PPT with slide size {slide_width_in}\" x {slide_height_in}\"")
    
    # Add each image as a slide
    for img_path in image_paths:
        # Add a blank slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        
        # Add picture that fills the slide
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Add picture filling entire slide
        left = top = Inches(0)
        pic = slide.shapes.add_picture(img_path, left, top, width=slide_width, height=slide_height)
        
        print(f"  Added slide with image: {os.path.basename(img_path)}")
    
    # Save PPT
    prs.save(output_ppt)
    print(f"PPT saved to '{output_ppt}' ({len(image_paths)} slides)")


def main():
    parser = argparse.ArgumentParser(
        description="Convert PDF to PPT via intermediate image rendering",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python3 pdf_to_ppt.py document.pdf
  python3 pdf_to_ppt.py document.pdf --zoom 3.0 --format jpg
  python3 pdf_to_ppt.py document.pdf --output presentation.pptx --img-dir ./tmp/images
        """
    )
    
    parser.add_argument('pdf_path', help='Path to the input PDF file')
    parser.add_argument('--img-dir', help='Directory to store intermediate images')
    parser.add_argument('--output', help='Path for output PPT file')
    parser.add_argument('--zoom', type=float, default=2.0, help='Zoom factor for rendering (default: 2.0)')
    parser.add_argument('--slide-width-in', type=float, default=13.33, help='Slide width in inches (default: 13.33)')
    parser.add_argument('--slide-height-in', type=float, default=7.5, help='Slide height in inches (default: 7.5)')
    parser.add_argument('--format', choices=['png', 'jpg', 'jpeg'], default='png', help='Image format (default: png)')
    
    args = parser.parse_args()
    
    # Validate input file
    if not os.path.isfile(args.pdf_path):
        print(f"Error: Input file '{args.pdf_path}' does not exist.", file=sys.stderr)
        sys.exit(1)
    
    # Set default paths if not provided
    pdf_dir = os.path.dirname(os.path.abspath(args.pdf_path))
    pdf_basename = os.path.splitext(os.path.basename(args.pdf_path))[0]
    
    if args.img_dir is None:
        img_dir = os.path.join(pdf_dir, 'images')
    else:
        img_dir = os.path.abspath(args.img_dir)
    
    if args.output is None:
        output_ppt = os.path.join(pdf_dir, f"{pdf_basename}.pptx")
    else:
        output_ppt = os.path.abspath(args.output)
    
    # Show configuration
    print("PDF to PPT Conversion")
    print("=" * 50)
    print(f"Input PDF:    {args.pdf_path}")
    print(f"Image dir:    {img_dir}")
    print(f"Output PPT:   {output_ppt}")
    print(f"Zoom:         {args.zoom}")
    print(f"Format:       {args.format}")
    print(f"Slide size:   {args.slide_width_in}\" x {args.slide_height_in}\"")
    print("=" * 50)
    
    try:
        # Step 1: Convert PDF to images
        image_paths = pdf_to_images(args.pdf_path, img_dir, args.zoom, args.format)
        
        # Step 2: Convert images to PPT
        images_to_ppt(image_paths, output_ppt, args.slide_width_in, args.slide_height_in)
        
        print("\nConversion completed successfully!")
        print(f"  Images stored in: {img_dir}")
        print(f"  PPT saved to:     {output_ppt}")
        
    except Exception as e:
        print(f"\nError during conversion: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()