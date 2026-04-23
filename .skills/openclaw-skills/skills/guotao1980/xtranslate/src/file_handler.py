import os
from docx import Document
from pdf2docx import Converter
import openpyxl
from pptx import Presentation
from striprtf.striprtf import rtf_to_text

class FileHandler:
    @staticmethod
    def read_txt(file_path):
        """读取 txt 文件内容"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            with open(file_path, 'r', encoding='gbk') as f:
                return f.read()

    @staticmethod
    def read_rtf(file_path):
        """读取 RTF 文件内容"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
            return rtf_to_text(content)

    @staticmethod
    def read_docx(file_path):
        """读取 docx 文件内容，返回段落列表"""
        doc = Document(file_path)
        return [p.text for p in doc.paragraphs if p.text.strip() != ""]

    @staticmethod
    def read_xlsx(file_path):
        """读取 XLSX 文件内容"""
        wb = openpyxl.load_workbook(file_path, data_only=True)
        texts = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value and isinstance(cell.value, str):
                        texts.append(cell.value)
        return texts

    @staticmethod
    def read_pptx(file_path):
        """读取 PPTX 文件内容"""
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
        return texts

    @staticmethod
    def pdf_to_docx(pdf_path, tmp_dir):
        """PDF 转 Word"""
        file_name = os.path.basename(pdf_path)
        docx_path = os.path.join(tmp_dir, file_name.replace('.pdf', '.docx'))
        
        cv = Converter(pdf_path)
        cv.convert(docx_path, start=0, end=None)
        cv.close()
        return docx_path

    @staticmethod
    def save_txt(content, output_path):
        """保存为 txt 文件"""
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(content)

    @staticmethod
    def translate_docx_in_place(file_path, translator, output_path):
        """翻译 docx 文件并保留原格式"""
        # 支持 .wps (如果是现代 wps，通常可以被 docx 读取)
        doc = Document(file_path)
        
        # 1. 收集所有需要翻译的文本对象 (段落、表格单元格)
        all_elements = []
        
        # 处理正文段落
        for p in doc.paragraphs:
            if p.text.strip():
                all_elements.append(p)
                
        # 处理表格内容
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for p in cell.paragraphs:
                        if p.text.strip():
                            all_elements.append(p)
        
        # 2. 提取文本并批量翻译
        texts_to_translate = [e.text for e in all_elements]
        translated_texts = translator.translate_list(texts_to_translate)
        
        # 3. 将翻译后的文本写回
        for element, translated_text in zip(all_elements, translated_texts):
            element.text = translated_text
            
        # 4. 保存
        doc.save(output_path)
        return output_path

    @staticmethod
    def translate_xlsx_in_place(file_path, translator, output_path):
        """翻译 XLSX 文件并保留原格式"""
        wb = openpyxl.load_workbook(file_path)
        
        for sheet in wb.worksheets:
            # 收集待翻译单元格
            cells_to_translate = []
            texts_to_translate = []
            
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value and isinstance(cell.value, str) and cell.value.strip():
                        cells_to_translate.append(cell)
                        texts_to_translate.append(cell.value)
            
            # 批量翻译并写回
            if texts_to_translate:
                translated_texts = translator.translate_list(texts_to_translate)
                for cell, translated_text in zip(cells_to_translate, translated_texts):
                    cell.value = translated_text
        
        wb.save(output_path)
        return output_path

    @staticmethod
    def translate_pptx_in_place(file_path, translator, output_path):
        """翻译 PPTX 文件并保留原格式"""
        prs = Presentation(file_path)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text.strip():
                                translated_text = translator.translate_text(run.text)
                                run.text = translated_text
        
        prs.save(output_path)
        return output_path

    @staticmethod
    def get_all_files(target_path):
        """根据路径获取所有符合要求的文件列表"""
        files = []
        # 支持更多格式：xlsx, pptx, rtf, wps, dxf, dwg
        allowed_exts = ('.txt', '.docx', '.pdf', '.xlsx', '.pptx', '.rtf', '.wps', '.dxf', '.dwg')
        if os.path.isfile(target_path):
            if target_path.lower().endswith(allowed_exts):
                files.append(target_path)
        elif os.path.isdir(target_path):
            for root, dirs, filenames in os.walk(target_path):
                for filename in filenames:
                    if filename.lower().endswith(allowed_exts):
                        files.append(os.path.join(root, filename))
        return files
