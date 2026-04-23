# -*- coding: utf-8 -*-
"""
UniversalDocProcessor - 通用文档处理核心模块
支持多格式解析、状态管理、任务执行
"""

import os
import io
import json
import chardet
from typing import Dict, List, Optional, Any, Union
from enum import Enum
from datetime import datetime

# 文件解析库（按需导入）
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    import docx
except ImportError:
    docx = None

try:
    import pandas as pd
    import openpyxl
except ImportError:
    pd = None
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


class ProcessState(Enum):
    """处理状态枚举"""
    WAITING_FOR_FILES = "waiting_for_files"
    FILES_RECEIVED = "files_received"
    EXECUTING_TASK = "executing_task"
    NEED_INFO = "need_info"


class FileInfo:
    """文件信息类"""
    def __init__(self, name: str, path: str, size: int, content: str = ""):
        self.name = name
        self.path = path
        self.size = size
        self.content = content
        self.parsed_at = datetime.now().isoformat()
        self.format = self._get_format(name)
    
    def _get_format(self, filename: str) -> str:
        ext = os.path.splitext(filename)[1].lower()
        format_map = {
            '.pdf': 'PDF',
            '.docx': 'Word',
            '.doc': 'Word',
            '.xlsx': 'Excel',
            '.xls': 'Excel',
            '.pptx': 'PPT',
            '.ppt': 'PPT',
            '.txt': 'Text',
            '.md': 'Markdown',
            '.csv': 'CSV',
            '.json': 'JSON'
        }
        return format_map.get(ext, 'Unknown')


class UniversalDocProcessor:
    """通用文档处理器"""
    
    def __init__(self):
        self.state = ProcessState.WAITING_FOR_FILES
        self.files: List[FileInfo] = []
        self.current_task: str = ""
        self.collected_info: Dict[str, Any] = {}
        self.task_history: List[Dict] = []
        self.pending_questions: List[str] = []
    
    def receive_files(self, file_list: List[Dict]) -> str:
        """接收文件并解析"""
        received_count = 0
        
        for file_data in file_list:
            try:
                content = self._parse_file(
                    file_data.get('path', ''),
                    file_data.get('name', '')
                )
                file_info = FileInfo(
                    name=file_data.get('name', 'unknown'),
                    path=file_data.get('path', ''),
                    size=file_data.get('size', 0),
                    content=content
                )
                self.files.append(file_info)
                received_count += 1
            except Exception as e:
                print(f"解析失败 {file_data.get('name')}: {e}")
        
        if received_count > 0:
            self.state = ProcessState.FILES_RECEIVED
            return f"已收到 {received_count} 个文件，等待您的具体任务指令"
        
        return "未收到有效文件，请上传需要处理的文档"
    
    def receive_task(self, task_description: str) -> str:
        """接收任务描述"""
        if self.state != ProcessState.FILES_RECEIVED:
            return "请先上传需要处理的文件"
        
        self.current_task = task_description
        self.state = ProcessState.EXECUTING_TASK
        
        # 检查信息完整性
        missing_info = self._check_info_completeness(task_description)
        
        if missing_info:
            self.state = ProcessState.NEED_INFO
            self.pending_questions = missing_info
            return self._generate_questions(missing_info)
        
        # 信息完整，执行任务
        return self._execute_task()
    
    def provide_info(self, info: str) -> str:
        """用户提供补充信息"""
        if self.state != ProcessState.NEED_INFO:
            return "当前不需要补充信息"
        
        self.collected_info['user补充'] = info
        
        # 再次检查完整性
        missing_info = self._check_info_completeness(
            self.current_task + " " + info
        )
        
        if missing_info:
            self.pending_questions = missing_info
            return self._generate_questions(missing_info)
        
        self.state = ProcessState.EXECUTING_TASK
        return self._execute_task()
    
    def _parse_file(self, file_path: str, filename: str) -> str:
        """根据文件类型解析内容"""
        ext = os.path.splitext(filename)[1].lower()
        
        parsers = {
            '.pdf': self._parse_pdf,
            '.docx': self._parse_docx,
            '.doc': self._parse_docx,
            '.xlsx': self._parse_xlsx,
            '.xls': self._parse_xlsx,
            '.pptx': self._parse_pptx,
            '.ppt': self._parse_pptx,
            '.txt': self._parse_txt,
            '.md': self._parse_txt,
            '.csv': self._parse_csv,
            '.json': self._parse_json
        }
        
        parser = parsers.get(ext, self._parse_binary)
        return parser(file_path)
    
    def _parse_pdf(self, file_path: str) -> str:
        """解析PDF"""
        if not PyPDF2:
            return self._parse_binary(file_path)
        
        try:
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text = ""
                for page in reader.pages[:50]:  # 限制页数
                    text += page.extract_text() + "\n"
                return text[:50000]  # 限制长度
        except Exception as e:
            return f"[PDF解析失败: {e}]"
    
    def _parse_docx(self, file_path: str) -> str:
        """解析Word"""
        if not docx:
            return self._parse_binary(file_path)
        
        try:
            doc = docx.Document(file_path)
            text = "\n".join([p.text for p in doc.paragraphs])
            return text[:50000]
        except Exception as e:
            return f"[Word解析失败: {e}]"
    
    def _parse_xlsx(self, file_path: str) -> str:
        """解析Excel"""
        if not pd:
            return self._parse_binary(file_path)
        
        try:
            df = pd.read_excel(file_path, nrows=1000)
            return df.to_string(max_rows=100)
        except Exception as e:
            return f"[Excel解析失败: {e}]"
    
    def _parse_pptx(self, file_path: str) -> str:
        """解析PPT"""
        if not Presentation:
            return self._parse_binary(file_path)
        
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides[:20]:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text[:30000]
        except Exception as e:
            return f"[PPT解析失败: {e}]"
    
    def _parse_txt(self, file_path: str) -> str:
        """解析文本文件"""
        try:
            with open(file_path, 'rb') as f:
                raw = f.read()
                result = chardet.detect(raw)
                encoding = result['encoding'] or 'utf-8'
                return raw.decode(encoding, errors='ignore')[:50000]
        except Exception as e:
            return f"[文本解析失败: {e}]"
    
    def _parse_csv(self, file_path: str) -> str:
        """解析CSV"""
        if not pd:
            return self._parse_binary(file_path)
        
        try:
            df = pd.read_csv(file_path, nrows=1000)
            return df.to_string(max_rows=100)
        except Exception as e:
            return f"[CSV解析失败: {e}]"
    
    def _parse_json(self, file_path: str) -> str:
        """解析JSON"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            return json.dumps(data, ensure_ascii=False, indent=2)[:50000]
        except Exception as e:
            return f"[JSON解析失败: {e}]"
    
    def _parse_binary(self, file_path: str) -> str:
        """解析二进制文件"""
        try:
            with open(file_path, 'rb') as f:
                raw = f.read(1000)
                return f"[二进制文件，{len(raw)}字节预览]\n{raw[:200]}"
        except:
            return "[无法解析的文件格式]"
    
    def _check_info_completeness(self, task: str) -> List[str]:
        """检查任务信息完整性"""
        task_lower = task.lower()
        missing = []
        
        # 检查是否需要更多信息
        vague_keywords = ['处理', '看看', '帮我', '搞一下', '弄一下']
        if any(k in task_lower for k in vague_keywords):
            missing.append("具体处理方式是什么？（分析/修改/提取/翻译/转换）")
        
        # 分析任务
        if '分析' in task or '摘要' in task:
            if '什么' in task_lower or '哪些' in task_lower or len(task) < 20:
                missing.append("需要分析哪些内容？（全文/关键信息/数据统计）")
        
        # 修改任务
        if '修改' in task or '编辑' in task:
            missing.append("需要修改的具体内容或位置？")
        
        # 翻译任务
        if '翻译' in task:
            if '英文' not in task_lower and '中文' not in task_lower and '日语' not in task_lower:
                missing.append("需要翻译成什么语言？")
        
        return missing
    
    def _generate_questions(self, missing_info: List[str]) -> str:
        """生成补充信息问题"""
        questions = "\n".join([f"{i+1}. {q}" for i, q in enumerate(missing_info)])
        return f"请补充以下信息：\n{questions}"
    
    def _execute_task(self) -> str:
        """执行任务（这里调用大模型处理）"""
        # 记录任务
        task_record = {
            "task": self.current_task,
            "files": [f.name for f in self.files],
            "collected_info": self.collected_info.copy(),
            "timestamp": datetime.now().isoformat()
        }
        self.task_history.append(task_record)
        
        # 这里应该调用OpenClaw大模型处理
        # 简化处理：返回任务确认
        result = f"[任务执行中]\n任务：{self.current_task}\n文件数：{len(self.files)}\n处理完成后的文件将直接返回给您"
        
        # 重置到FILES_RECEIVED状态，保留文件以便多轮任务
        self.state = ProcessState.FILES_RECEIVED
        
        return result
    
    def get_file_list(self) -> str:
        """获取已接收文件列表"""
        if not self.files:
            return "暂无已接收的文件"
        
        info = "已接收的文件：\n"
        for i, f in enumerate(self.files, 1):
            info += f"{i}. {f.name} ({f.format}, {f.size/1024:.1f}KB)\n"
        
        return info
    
    def reset(self):
        """重置状态"""
        self.state = ProcessState.WAITING_FOR_FILES
        self.files = []
        self.current_task = ""
        self.collected_info = {}
        self.pending_questions = []


# 测试代码
if __name__ == "__main__":
    processor = UniversalDocProcessor()
    
    # 模拟文件接收
    print("=== 测试1: 文件接收 ===")
    # 这里需要传入实际文件路径
    # result = processor.receive_files([{"name": "test.pdf", "path": "", "size": 1000}])
    # print(result)
    
    print("=== 测试2: 状态查询 ===")
    print(f"当前状态: {processor.state.value}")
    
    print("=== 测试3: 文件列表 ===")
    print(processor.get_file_list())