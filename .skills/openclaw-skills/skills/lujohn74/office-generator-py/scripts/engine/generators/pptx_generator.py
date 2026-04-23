from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, List


def _add_bullets_with_paging(prs, base_title: str, bullets: List[str]):
    from pptx.util import Inches

    max_bullets = 6
    chunks = [bullets[i:i + max_bullets] for i in range(0, len(bullets), max_bullets)] or [[]]
    slides = []
    for idx, chunk in enumerate(chunks, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if slide.shapes.title:
            title = base_title if len(chunks) == 1 else f'{base_title}（续 {idx}/{len(chunks)}）'
            slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            if chunk:
                tf.text = str(chunk[0])
                for bullet in chunk[1:]:
                    p = tf.add_paragraph()
                    p.text = str(bullet)
        slides.append(slide)
    return slides


def generate_pptx(spec: Dict[str, Any], out_path: str) -> Path:
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches
    except ImportError as exc:
        raise RuntimeError('Missing dependency: python-pptx. Install requirements first.') from exc

    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = spec.get('title', 'Untitled Presentation')
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = spec.get('purpose', '') or ''

    slides: List[Dict[str, Any]] = spec.get('contentSpec', {}).get('slides', [])
    for slide_spec in slides:
        layout_name = slide_spec.get('layout', 'content')
        if layout_name == 'table':
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            if slide.shapes.title:
                slide.shapes.title.text = slide_spec.get('title', 'Table Slide')
            table_data = slide_spec.get('table', [])
            if table_data:
                rows = len(table_data)
                cols = max(len(r) for r in table_data)
                table_shape = slide.shapes.add_table(rows, cols, Inches(0.6), Inches(1.5), Inches(8.0), Inches(3.5))
                table = table_shape.table
                for r_idx, row in enumerate(table_data):
                    for c_idx, value in enumerate(row):
                        table.cell(r_idx, c_idx).text = '' if value is None else str(value)
        elif layout_name == 'image':
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            if slide.shapes.title:
                slide.shapes.title.text = slide_spec.get('title', 'Image Slide')
            image = slide_spec.get('image', {})
            image_path = image.get('path')
            if image_path and Path(image_path).exists():
                slide.shapes.add_picture(str(image_path), Inches(0.8), Inches(1.4), width=Inches(7.0))
            bullets = slide_spec.get('bullets', [])
            if bullets:
                textbox = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(8.0), Inches(1.2))
                tf = textbox.text_frame
                tf.text = str(bullets[0])
                for bullet in bullets[1:]:
                    p = tf.add_paragraph()
                    p.text = str(bullet)
        else:
            bullets = slide_spec.get('bullets', [])
            if len(bullets) > 6:
                _add_bullets_with_paging(prs, slide_spec.get('title', 'Slide'), bullets)
            else:
                slide = prs.slides.add_slide(prs.slide_layouts[1] if layout_name != 'title' else prs.slide_layouts[0])
                if slide.shapes.title:
                    slide.shapes.title.text = slide_spec.get('title', 'Slide')
                if len(slide.placeholders) > 1:
                    tf = slide.placeholders[1].text_frame
                    if bullets:
                        tf.text = str(bullets[0])
                        for bullet in bullets[1:]:
                            p = tf.add_paragraph()
                            p.text = str(bullet)
                    else:
                        tf.text = ''

    path = Path(out_path)
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return path
